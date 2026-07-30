# --- ROBUST LOGGING SYSTEM (macOS STABILITY PATCH) ---
import os
import sys

PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

try:
    log_dir = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "scratch",
        "runtime_logs",
    )
    os.makedirs(log_dir, exist_ok=True)
    log_file_path = os.path.join(log_dir, "sperm_error_log.txt")
    f_log = open(log_file_path, 'a', encoding='utf-8')
except (PermissionError, IOError):
    # Fallback to User Home if script directory is read-only (common in macOS .app bundles)
    log_file_path = os.path.join(os.path.expanduser("~"), "sperm_error_log.txt")
    f_log = open(log_file_path, 'a', encoding='utf-8')

class Tee(object):
    def __init__(self, *files):
        self.files = files
    def write(self, obj):
        for f in self.files:
            try:
                f.write(obj)
                f.flush()
            except: pass
    def flush(self):
        for f in self.files:
            try: f.flush()
            except: pass

sys.stdout = Tee(sys.stdout, f_log)
sys.stderr = Tee(sys.stderr, f_log)
print(f"\n--- NEW SESSION STARTED (v5.7-unet-ready): {os.path.basename(__file__)} ---")
print(f"Log Path: {log_file_path}\n")
# ---------------------------------

#!/usr/bin/env python3
"""
Sperm Nucleus Segmentation & 3D Morphometrics Pipeline  -  Saturn Project
=========================================================================
A production-ready image-analysis pipeline for automated detection,
measurement, and 3D reconstruction of sperm nuclei tailored to
the *Saturn* experimental dataset acquired on the Leica confocal system.

Biological context
------------------
Spermiogenesis - the post-meiotic differentiation of round spermatids into
mature spermatozoa - involves dramatic nuclear elongation, chromatin condensation,
and apical plunging of the nucleus toward the basal lamina of the seminiferous
tubule.  Quantifying these morphological changes requires imaging the highly
elongated, condensed sperm nuclei.  This pipeline automates the process by:

1. Detecting sperm nuclei as thin, dim, ridge-like objects in each 2D
   Z-slice using a skeleton-first strategy (Frangi ridge filter -> Otsu
   binarisation -> morphological closing -> geodesic skeleton).
2. Measuring per-cell biometrics: geodesic length (um), Euclidean width,
   tortuosity, endpoint count, branch complexity, centroid, and estimated area.
3. Linking detections across Z-slices into 3D tracks using nearest-neighbor and
   overlap-aware frame-to-frame assignment.
4. Computing 3D morphometrics from tracks: estimated 3D length, Z-span,
   sampled Z-coverage, approximated volume, pitch angle, taper ratio, and
   nearest-neighbor packing density.
5. Exporting results as CSV, a multi-sheet Excel workbook, a multi-page PDF
   report, and a native PowerPoint (.pptx) dashboard with editable charts.

Saturn-specific calibration
----------------------------
Physical scale factors are derived from Leica confocal metadata for the
Saturn dataset:

- ``UM_PER_PX_XY  = 0.7568``  um/pixel  (lateral resolution)
- ``UM_PER_SLICE_Z = 1.0404``  um/slice  (z-step size)

These values are set as defaults in ``CONFIG`` and can be overridden via
the Parameter Editor GUI or by loading a JSON settings file.

v5.2 measurement notes
----------------------
``area_px`` is an estimated slender-object area, computed as 2D geodesic
length times median width. The raw skeleton-pixel count is exported separately
as ``skeleton_area_px`` for audit/debug use.

3D length is a projection-length-plus-Z-span estimate:
``sqrt(max(2D geodesic, XY displacement)^2 + z_span^2)``.  ``z_span_um`` is
endpoint-to-endpoint vertical span; ``z_covered_um`` is the sampled slab
coverage including both end slices. Volume, effective thickness, taper, and
other width/area-derived metrics are PSF- and sampling-sensitive, so use them
mainly for relative comparisons under matched imaging settings.

Pipeline architecture
---------------------
::

    Image stack (TIF / PNG / JPG)
        -> load_batch_files()
        -> process_batch() / process_one_image()
        -> segment_slice()
        -> measure_spermatids()
        -> rows_from_results()
        -> track_across_slices()
        -> generate_excel_report() / generate_batch_report() / generate_pptx_report()

Key configuration parameters (``CONFIG`` dict)
-----------------------------------------------
``UM_PER_PX_XY``       Physical pixel size in um (Leica metadata: 0.7568).
``UM_PER_SLICE_Z``     Z-step size in um (Leica metadata: 1.0404).
``FRANGI_SCALE_RANGE`` Tubeness filter scales (px); set to match sperm nucleus width.
``MIN_SKEL_LEN_PX``    Minimum skeleton length accepted (removes debris/noise).
``MAX_WIDTH_PX``       Maximum skeleton width accepted (rejects merged clusters).
``MAX_BRANCH_NODES``   Maximum branch-point count; >0 tolerates bridged tips.
``MAX_TORTUOSITY``     Maximum curvature index; rejects snaking merged networks.
``DO_TRACKING``        Enable/disable cross-slice 3D linking.
``TRACK_MAX_DIST_UM``  Maximum centroid displacement between adjacent Z-slices.

Usage
-----
Launch the interactive GUI (recommended)::

    python sperm_segmentation_saturnv5.7.py

Run a headless batch analysis::

    python sperm_segmentation_saturnv5.7.py --batch

Analyze a single slice::

    python sperm_segmentation_saturnv5.7.py --single --z 4

Dependencies
------------
numpy, scipy, scikit-image, pandas, matplotlib, tifffile, opencv-python,
Pillow, xlsxwriter, python-pptx, tkinter (stdlib)

Author
------
Dushyant Mishra  |  Findlay Lab  |  Saturn Dataset Branch
"""

import os, sys, glob, re, time, warnings, heapq, argparse, math, pathlib as pl
import time as _t
import json, webbrowser, threading, subprocess, shutil
from dataclasses import dataclass, asdict
try:
    import requests
    _HAVE_REQUESTS = True
except ImportError:
    _HAVE_REQUESTS = False
warnings.filterwarnings("ignore")

import numpy as np
import pandas as pd
import tifffile
import matplotlib
# --- macOS STABILITY PATCH: Force non-interactive backend for background reporting ---
matplotlib.use('Agg', force=True)
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.patches import Patch
from matplotlib.path import Path
print(f"[matplotlib backend: {matplotlib.get_backend()}]")

from skimage import measure, morphology, exposure, segmentation as skseg, feature
from skimage.filters import meijering, gaussian, apply_hysteresis_threshold
from skimage.morphology import skeletonize
from scipy.ndimage import distance_transform_edt, grey_dilation
from scipy.spatial import cKDTree
from scipy.optimize import linear_sum_assignment
from matplotlib.backends.backend_pdf import PdfPages

try:
    import cv2 as _cv2
    _HAVE_CV2 = True
except ImportError:
    _HAVE_CV2 = False


# =============================================================================
# CONFIG
# =============================================================================

CONFIG = {
    # ------ run mode ---------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    "RUN_MODE": "single",          # "single" | "batch"
    "ANALYSIS_MODE": "comparative",  # "comparative" | "reference_morphology" | "legacy"

    # ------ single-image selection ------------------------------------------------------------------------------------------------------------------------------------------------
    "SINGLE_IMAGE_SELECTION_MODE": "dialog",  # "path" | "z_index" | "dialog"
    "SINGLE_TEST_IMAGE": r"C:\Users\dmishra\Desktop\sperm images\Project001_Series002_z15_ch00.tif",
    "SINGLE_Z_INDEX": 15,

    # ------ input / output ------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    "INPUT_DIR":    ".",
    "OUTPUT_DIR":   "./sperm_results_saturnv5_7",
    "FILE_PATTERN": "Project001_Series002_z*_ch00.tif",
    "ROI_MASK_PATH": "",
    "PREPROCESS_MODE": "roi_adaptive",
    "LEGACY_TWO_PASS_ROI": False,
    "STACK_QC_SAMPLE_COUNT": 12,
    "NORM_LOW_PERCENTILE": 1.0,
    "NORM_HIGH_PERCENTILE": 99.5,
    "NORM_STACK_WEIGHT": 0.80,
    "CLAHE_MODE": "auto_stack",
    "CLAHE_CLIP_HIGH_CONTRAST": 0.010,
    "CLAHE_CLIP_STANDARD": 0.025,
    "CLAHE_CLIP_LOW_CONTRAST": 0.035,
    "AUTO_CONTRAST_HIGH_THRESHOLD": 0.45,
    "AUTO_CONTRAST_LOW_THRESHOLD": 0.25,
    "ROI_CROP_PADDING_PX": 16,
    "ROI_THRESHOLD_PERCENTILES_ONLY": True,
    "EXCLUSION_MASK_PATH": "",
    "ROI_BOUNDARY_SAFE_RIDGE": True,
    "ROI_THRESHOLD_EXCLUDE_BOUNDARY_PX": 4,
    "AI_PREPROCESSING_MODE": "off",

    # ------ optional U-Net 2.5D integration scaffold ---------------------------------------------------------------------------------------------------------------------------------------
    # COCO is training-only and is never read by Saturn during inference.
    # Runtime U-Net inference uses a trained checkpoint plus [z-1, z, z+1] raw planes.
    "SEGMENTATION_ENGINE": "classical_saturn",  # classical_saturn | unet_assisted | hybrid | unet_primary
    "UNET_MODEL_PATH": "",
    "UNET_THRESHOLD": 0.10,
    "UNET_THRESHOLD_MODE": "soft",
    "UNET_CANDIDATE_THRESHOLD": 0.05,
    "UNET_SEED_THRESHOLD": 0.30,
    "UNET_CONTEXT_MODE": "z_minus_z_z_plus",
    "UNET_INFERENCE_MODE": "roi_tiled",
    "UNET_TILE_SIZE": 256,
    "UNET_TILE_OVERLAP": 64,
    "UNET_TILE_BATCH_SIZE": 8,
    "UNET_ROI_PADDING_PX": 32,
    "UNET_STITCH_MODE": "weighted_average",
    "UNET_OUTSIDE_ROI_ZERO": True,
    "UNET_FAIL_HARD": True,
    "UNET_SAVE_PROBABILITY_MAPS": True,
    "UNET_CANDIDATE_ACCOUNTING": True,
    "UNET_RESCUE_ENABLE": True,
    "UNET_RESCUE_THRESHOLD": 0.60,
    "UNET_RESCUE_HYSTERESIS_ENABLE": True,
    "UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS": True,
    "UNET_RESCUE_EXCLUDE_DILATION_PX": 1,
    "UNET_RESCUE_MIN_COMPONENT_PX": 3,
    # Technical noise floor only. Do not use expected biological length as an
    # acceptance gate because genuinely short nuclei may be genotype-dependent.
    "UNET_RESCUE_MIN_SKEL_LEN_UM": 2.0,
    # Shorter centerlines remain eligible when the U-Net evidence itself is
    # strong. This avoids encoding an expected WT/mutant length into acceptance.
    "UNET_SHORT_RESCUE_MIN_MEAN_PROB": 0.85,
    "UNET_RESCUE_MAX_ADDITIONS_PER_SLICE": 0,
    "UNET_RESCUE_SPLIT_RETRY_ENABLE": True,
    "UNET_RESCUE_SPLIT_THRESHOLDS": [0.70, 0.80, 0.90],
    "UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE": True,
    "UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB": 0.85,
    "UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB": 0.75,
    "UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM": 4.0,
    "UNET_INSTANCE_SPLIT_ENABLE": True,
    "UNET_INSTANCE_SEED_THRESHOLD": 0.75,
    "UNET_INSTANCE_PEAK_MIN_DISTANCE_PX": 6,
    "UNET_INSTANCE_WATERSHED_COMPACTNESS": 0.001,
    "UNET_PRIMARY_MIN_COMPONENT_PX": 3,
    "UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE": False,
    "UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX": 2,
    "UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS": True,
    "UNET_PRIMARY_SAVE_FILLED_MASK_OVERLAY": True,
    "UNET_PRIMARY_SAVE_INSTANCE_OVERLAY": True,
    "UNET_TRACKING_SUPPORT": True,
    "ASSIGNMENT_UNET_SUPPORT_WEIGHT": 0.6,
    "ASSIGNMENT_UNET_CONTINUITY_WEIGHT": 0.25,
    "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT": 0.4,
    "UNET_DETECTED_LABEL": "Detected by U-Net",
    "UNET_RESCUED_LABEL": "U-Net rescued",
    "UNET_COMPLETED_BY_BRIDGE_LABEL": "Completed by bridge",
    "UNET_COMPLETED_BY_EXTENSION_LABEL": "Completed by extension",
    "UNET_MERGED_CANDIDATE_LABEL": "Merged candidate",
    "UNET_QC_BORDERLINE_LABEL": "QC borderline",
    "SATURN_RECOVERED_LABEL": "Recovered by Saturn",
    "FINAL_ACCEPTED_LABEL": "Final accepted",
    "EXCLUDED_FROM_MEASUREMENT_LABEL": "Excluded from measurement",

    # ------ calibration ---------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    "UM_PER_PX_XY":   0.756836,
    "UM_PER_SLICE_Z": 1.040460,

    # ------ image enhancement ---------------------------------------------------------------------------------------------------------------------------------------------------------------
    "CLAHE_CLIP":   0.025,
    "CLAHE_KERNEL": 128,
    "BG_SIGMA":     5.065,
    "BG_SIGMA_UM": 6.0,
    "DENOISE_SIGMA_UM": 0.45,
    "RIDGE_SIGMAS": [1, 2, 3, 4],
    "RIDGE_SIGMAS_UM": [0.60, 0.90, 1.20, 1.50],

    # ------ hysteresis threshold ------------------------------------------------------------------------------------------------------------------------------------------------------
    "THRESHOLD_HI": 91.0,
    "THRESHOLD_LO": 83.0,

    # ------ morphological cleanup ---------------------------------------------------------------------------------------------------------------------------------------------------
    "CLOSE_RADIUS":  0,
    "MIN_HOLE_AREA": 0,
    "MIN_OBJ_PX":    8,

    # ------ skeleton-level gap bridging ---------------------------------------------------------------------------------------------------------------------------------
    # P1: reduced from 10 -> 5 to prevent chaining distinct spermatids
    "MAX_BRIDGE_PX": 5,
    "MAX_BRIDGE_UM": 1.5,
    "MAX_BRIDGE_ANGLE_DEG": 35.0,

    # ------ branch pruning & automated splitting ------------------------------------------------------------------------------------------------------
    "MAX_BRANCH_LEN_PX": 5,   # prune spurs shorter than this before measuring
    "MAX_BRANCH_LEN_UM": 2.3,
    "BREAK_JUNCTIONS": True,  # automatically sever all branching intersections into distinct lines

    # ------ optional early mask-level shape filter ------------------------------------------------------------------------------------------------
    "USE_EARLY_SHAPE_FILTER": False,
    "MIN_ECCENTRICITY": 0.60,
    "MAX_MINOR_PX":     12.0,
    "MIN_AXIS_RATIO":   1.4,
    "MIN_MAJOR_PX":     5,

    # ------ post-skeleton filters ---------------------------------------------------------------------------------------------------------------------------------------------------
    "MIN_SKEL_LEN_PX":        6.02, # legacy JSON compatibility
    "MIN_SKEL_LEN_UM":        6.0,
    "MAX_GEODESIC_LEN_PX":    65.0, # backend actively cuts chains longer than this
    "MAX_GEODESIC_LEN_UM":    20.0,
    "MAX_WIDTH_PX":           8.94, # legacy JSON compatibility
    "MAX_WIDTH_UM":           4.2,
    "MIN_LENGTH_WIDTH_RATIO": 2.5,

    # ------ NEW topology filters ------------------------------------------------------------------------------------------------------------------------------------------------------
    "MAX_BRANCH_NODES": 0,        # 30-Gen V2 Optimized
    "MAX_TORTUOSITY": 2.5,
    "MAX_ENDPOINT_COUNT": 4,

    # N4: loops
    "ALLOW_LOOPS": False,
    "AUTO_LOCAL_REANALYSIS": False,
    "ROI_EDGE_QC_DISTANCE_UM": 1.0,

    # ------ tracking across z ---------------------------------------------------------------------------------------------------------------------------------------------------------------
    "DO_TRACKING":          True,
    "TRACKING_BACKEND":     "hybrid_repair",  # legacy, global_assignment, or hybrid_repair
    "TRACK_MAX_DIST_UM":    6.8711,
    "TRACK_MAX_GAP_SLICES": 1,
    "TRACK_BBOX_PADDING_PX": 2,
    "TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM": 15.0,

    # ------ conservative tracking stop-rules ---------------------------------------------------------------------------------------------------------------------------
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": 0.7668,      # Saturn V5 tuned default
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": 0.3134,     # Saturn V5 tuned default
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": 0.5522,       # Saturn V5 tuned default
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": 0.40,       # Absolute tortuosity jump
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": 5.2397,      # Saturn V5 tuned default

    # ------ overlap-first Stage 2 parameters ------------------------------------------------------------------------------------------------------------------------------
    "OVERLAP_STABILITY_THRESHOLD": 0.2026,            # Saturn V5 tuned default
    "OVERLAP_ORIENTATION_DEG":     27.2181,             # Saturn V5 tuned default
    "OVERLAP_MULTIPLIER":          1.6255,             # Saturn V5 tuned default
    "OVERLAP_MIN_STABLE_COUNT":    1,                # Min stable metrics required for overlap continuation
    "ASSIGNMENT_MAX_COST":          8.0,              # global-assignment tracker gate
    "ASSIGNMENT_DIST_WEIGHT":       1.0,
    "ASSIGNMENT_OVERLAP_WEIGHT":    2.0,
    "ASSIGNMENT_LENGTH_WEIGHT":     2.0,
    "ASSIGNMENT_WIDTH_WEIGHT":      1.2,
    "ASSIGNMENT_AREA_WEIGHT":       0.9,
    "ASSIGNMENT_ANGLE_WEIGHT":      0.4,
    "HYBRID_REPAIR_MAX_COST":       3.6,       # v5.5 fragment-repair gate; lower is stricter
    "HYBRID_REPAIR_MAX_GAP_SLICES": 1,
    "HYBRID_REPAIR_MAX_FRAGMENT_SLICES": 2,
    "HYBRID_REPAIR_MAX_LINK_DIST_UM": 4.8,
    "HYBRID_REPAIR_MIN_OVERLAP":    0.05,
    "HYBRID_REPAIR_MAX_FINAL_LENGTH_UM": 15.0,

    # ------ quality audit thresholds (for automated outlier filtering) ---------------------------------
    "AUDIT_MAX_LENGTH_UM":     15.0,    # Flag tracks longer than this (um)
    "AUDIT_MAX_TORTUOSITY":    1.5,     # Flag tracks more tortuous than this
    "AUDIT_MAX_THICKNESS_UM":  2.0,     # Flag tracks thicker than this (um)
    "AUDIT_MAX_TAPER_RATIO":   1.5,     # Flag tracks with taper ratio above this
    "AUDIT_EXTREME_THICKNESS_UM": 3.5,  # Hard-fail only very extreme PSF/merge thickness
    "AUDIT_EXTREME_TAPER_RATIO":  3.0,  # Hard-fail only very extreme taper instability
    "AUDIT_MIN_SLICES":        1,       # single-slice nuclei may be biologically valid at this z-step

    # ------ output / debug ------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    "SAVE_DEBUG_IMAGES":   False,
    "SAVE_MASK_TIFS":      True,
    "SAVE_LABEL_TIFS":     True,
    "SAVE_OVERLAYS":       True,
    "SAVE_DETAIL_FIGURE":  True,
    "SHOW_PREVIEW_WINDOW": True,
    "SHOW_DEBUG_PREVIEW":  True,
}


# =============================================================================
# CONFIG VALIDATION
# =============================================================================

_REQUIRED = {
    "RUN_MODE": str, "ANALYSIS_MODE": str, "SINGLE_IMAGE_SELECTION_MODE": str,
    "SINGLE_TEST_IMAGE": str, "SINGLE_Z_INDEX": int,
    "INPUT_DIR": str, "OUTPUT_DIR": str, "FILE_PATTERN": str, "ROI_MASK_PATH": str,
    "PREPROCESS_MODE": str, "LEGACY_TWO_PASS_ROI": bool,
    "STACK_QC_SAMPLE_COUNT": int, "NORM_LOW_PERCENTILE": (int, float),
    "NORM_HIGH_PERCENTILE": (int, float), "NORM_STACK_WEIGHT": (int, float),
    "CLAHE_MODE": str, "CLAHE_CLIP_HIGH_CONTRAST": (int, float),
    "CLAHE_CLIP_STANDARD": (int, float), "CLAHE_CLIP_LOW_CONTRAST": (int, float),
    "AUTO_CONTRAST_HIGH_THRESHOLD": (int, float), "AUTO_CONTRAST_LOW_THRESHOLD": (int, float),
    "ROI_CROP_PADDING_PX": int, "ROI_THRESHOLD_PERCENTILES_ONLY": bool,
    "EXCLUSION_MASK_PATH": str,
    "ROI_BOUNDARY_SAFE_RIDGE": bool, "ROI_THRESHOLD_EXCLUDE_BOUNDARY_PX": int,
    "AI_PREPROCESSING_MODE": str,
    "SEGMENTATION_ENGINE": str, "UNET_MODEL_PATH": str,
    "UNET_THRESHOLD": (int, float), "UNET_THRESHOLD_MODE": str,
    "UNET_CANDIDATE_THRESHOLD": (int, float), "UNET_SEED_THRESHOLD": (int, float),
    "UNET_CONTEXT_MODE": str, "UNET_INFERENCE_MODE": str,
    "UNET_TILE_SIZE": int, "UNET_TILE_OVERLAP": int, "UNET_TILE_BATCH_SIZE": int,
    "UNET_ROI_PADDING_PX": int,
    "UNET_STITCH_MODE": str, "UNET_OUTSIDE_ROI_ZERO": bool, "UNET_FAIL_HARD": bool,
    "UNET_SAVE_PROBABILITY_MAPS": bool, "UNET_CANDIDATE_ACCOUNTING": bool,
    "UNET_RESCUE_ENABLE": bool, "UNET_RESCUE_THRESHOLD": (int, float),
    "UNET_RESCUE_HYSTERESIS_ENABLE": bool,
    "UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS": bool,
    "UNET_RESCUE_EXCLUDE_DILATION_PX": int, "UNET_RESCUE_MIN_COMPONENT_PX": int,
    "UNET_RESCUE_MIN_SKEL_LEN_UM": (int, float),
    "UNET_SHORT_RESCUE_MIN_MEAN_PROB": (int, float),
    "UNET_RESCUE_MAX_ADDITIONS_PER_SLICE": int,
    "UNET_RESCUE_SPLIT_RETRY_ENABLE": bool, "UNET_RESCUE_SPLIT_THRESHOLDS": list,
    "UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE": bool,
    "UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB": (int, float),
    "UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB": (int, float),
    "UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM": (int, float),
    "UNET_INSTANCE_SPLIT_ENABLE": bool, "UNET_INSTANCE_SEED_THRESHOLD": (int, float),
    "UNET_INSTANCE_PEAK_MIN_DISTANCE_PX": int, "UNET_INSTANCE_WATERSHED_COMPACTNESS": (int, float),
    "UNET_PRIMARY_MIN_COMPONENT_PX": int,
    "UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE": bool,
    "UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX": int,
    "UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS": bool,
    "UNET_PRIMARY_SAVE_FILLED_MASK_OVERLAY": bool,
    "UNET_PRIMARY_SAVE_INSTANCE_OVERLAY": bool,
    "UNET_TRACKING_SUPPORT": bool,
    "ASSIGNMENT_UNET_SUPPORT_WEIGHT": (int, float),
    "ASSIGNMENT_UNET_CONTINUITY_WEIGHT": (int, float),
    "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT": (int, float),
    "UNET_DETECTED_LABEL": str, "UNET_RESCUED_LABEL": str,
    "UNET_COMPLETED_BY_BRIDGE_LABEL": str,
    "UNET_COMPLETED_BY_EXTENSION_LABEL": str, "UNET_MERGED_CANDIDATE_LABEL": str,
    "UNET_QC_BORDERLINE_LABEL": str,
    "UM_PER_PX_XY": float, "UM_PER_SLICE_Z": float,
    "CLAHE_CLIP": float, "CLAHE_KERNEL": int, "BG_SIGMA": (int, float),
    "BG_SIGMA_UM": (int, float), "DENOISE_SIGMA_UM": (int, float),
    "RIDGE_SIGMAS": list, "RIDGE_SIGMAS_UM": list,
    "THRESHOLD_HI": (int, float), "THRESHOLD_LO": (int, float),
    "CLOSE_RADIUS": int, "MIN_HOLE_AREA": int, "MIN_OBJ_PX": int,
    "MAX_BRIDGE_PX": (int, float), "MAX_BRIDGE_UM": (int, float),
    "MAX_BRIDGE_ANGLE_DEG": (int, float), "MAX_BRANCH_LEN_PX": (int, float),
    "MAX_BRANCH_LEN_UM": (int, float),
    "USE_EARLY_SHAPE_FILTER": bool,
    "MIN_SKEL_LEN_PX": (int, float), "MIN_SKEL_LEN_UM": (int, float),
    "MAX_GEODESIC_LEN_PX": (int, float), "MAX_GEODESIC_LEN_UM": (int, float),
    "MAX_WIDTH_PX": (int, float), "MAX_WIDTH_UM": (int, float),
    "MIN_LENGTH_WIDTH_RATIO": (int, float),
    "MAX_BRANCH_NODES": (int, float),
    "MAX_TORTUOSITY": (int, float),
    "MAX_ENDPOINT_COUNT": (int, float),
    "AUTO_LOCAL_REANALYSIS": bool, "ROI_EDGE_QC_DISTANCE_UM": (int, float),
    "DO_TRACKING": bool, "TRACKING_BACKEND": str, "TRACK_MAX_DIST_UM": (int, float),
    "TRACK_MAX_GAP_SLICES": int,
    "TRACK_BBOX_PADDING_PX": int,
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": float,
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": float,
    "OVERLAP_STABILITY_THRESHOLD": float,
    "OVERLAP_ORIENTATION_DEG": float,
    "OVERLAP_MULTIPLIER": float,
    "OVERLAP_MIN_STABLE_COUNT": int,
    "ASSIGNMENT_MAX_COST": (int, float),
    "ASSIGNMENT_DIST_WEIGHT": (int, float),
    "ASSIGNMENT_OVERLAP_WEIGHT": (int, float),
    "ASSIGNMENT_LENGTH_WEIGHT": (int, float),
    "ASSIGNMENT_WIDTH_WEIGHT": (int, float),
    "ASSIGNMENT_AREA_WEIGHT": (int, float),
    "ASSIGNMENT_ANGLE_WEIGHT": (int, float),
    "HYBRID_REPAIR_MAX_COST": (int, float),
    "HYBRID_REPAIR_MAX_GAP_SLICES": int,
    "HYBRID_REPAIR_MAX_FRAGMENT_SLICES": int,
    "HYBRID_REPAIR_MAX_LINK_DIST_UM": (int, float),
    "HYBRID_REPAIR_MIN_OVERLAP": (int, float),
    "HYBRID_REPAIR_MAX_FINAL_LENGTH_UM": (int, float),
    "TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM": (int, float),
    "AUDIT_MAX_LENGTH_UM": (int, float),
    "AUDIT_MAX_TORTUOSITY": (int, float),
    "AUDIT_MAX_THICKNESS_UM": (int, float),
    "AUDIT_MAX_TAPER_RATIO": (int, float),
    "AUDIT_EXTREME_THICKNESS_UM": (int, float),
    "AUDIT_EXTREME_TAPER_RATIO": (int, float),
    "AUDIT_MIN_SLICES": int,
    "SAVE_DEBUG_IMAGES": bool, "SAVE_MASK_TIFS": bool,
    "SAVE_LABEL_TIFS": bool, "SAVE_OVERLAYS": bool,
    "SAVE_DETAIL_FIGURE": bool, "SHOW_PREVIEW_WINDOW": bool,
    "SHOW_DEBUG_PREVIEW": bool,
}


def validate_config(cfg):
    """
    Validates the pipeline CONFIG dictionary for required keys, correct types, and logical consistency.

    Raises a descriptive ValueError listing ALL problems found so engineers can fix
    everything in one go rather than playing whack-a-mole with sequential errors.

    Checks performed
    ----------------
    - Every key in ``_REQUIRED`` is present in *cfg*.
    - Each value has the expected Python type (e.g., ``float`` for calibration parameters).
    - ``THRESHOLD_LO`` < ``THRESHOLD_HI`` (hysteresis thresholding only works in this order).
    - ``RUN_MODE`` is exactly one of ``'single'`` or ``'batch'``.
    - The deprecated ``'REJECT_BRANCHES'`` Boolean key is absent (replaced by ``MAX_BRANCH_NODES`` int).

    Args:
        cfg (dict): The configuration mapping to validate.

    Raises:
        ValueError: If *any* of the above checks fail. The message lists all errors.
    """
    errors = []
    for key, expected in _REQUIRED.items():
        if key not in cfg:
            errors.append(f"  MISSING: '{key}'")
        elif not isinstance(cfg[key], expected):
            errors.append(f"  WRONG TYPE '{key}': "
                          f"got {type(cfg[key]).__name__}, want {expected}")
    if cfg.get("THRESHOLD_LO", 0) >= cfg.get("THRESHOLD_HI", 100):
        errors.append("  THRESHOLD_LO must be < THRESHOLD_HI")
    engine = str(cfg.get("SEGMENTATION_ENGINE", "")).strip().lower()
    supported_engines = {
        "classical_saturn", "unet_assisted", "hybrid", "unet_primary"
    }
    if engine not in supported_engines:
        errors.append(
            "  SEGMENTATION_ENGINE must be classical_saturn, unet_assisted, "
            "hybrid, or unet_primary"
        )
    if cfg.get("UM_PER_PX_XY", 0) <= 0 or cfg.get("UM_PER_SLICE_Z", 0) <= 0:
        errors.append("  UM_PER_PX_XY and UM_PER_SLICE_Z must be positive")
    if not (0 <= cfg.get("NORM_LOW_PERCENTILE", -1) < cfg.get("NORM_HIGH_PERCENTILE", -1) <= 100):
        errors.append("  NORM_LOW_PERCENTILE must be < NORM_HIGH_PERCENTILE within [0, 100]")
    if not (0 <= cfg.get("NORM_STACK_WEIGHT", -1) <= 1):
        errors.append("  NORM_STACK_WEIGHT must be between 0 and 1")
    if not (
        0
        <= cfg.get("UNET_CANDIDATE_THRESHOLD", -1)
        < cfg.get("UNET_RESCUE_THRESHOLD", -1)
        <= 1
    ):
        errors.append(
            "  UNET_CANDIDATE_THRESHOLD must be < "
            "UNET_RESCUE_THRESHOLD within [0, 1]"
        )
    if engine == "unet_primary" and not (
        0
        <= cfg.get("UNET_CANDIDATE_THRESHOLD", -1)
        < cfg.get("UNET_SEED_THRESHOLD", -1)
        <= 1
    ):
        errors.append(
            "  unet_primary requires UNET_CANDIDATE_THRESHOLD < "
            "UNET_SEED_THRESHOLD within [0, 1]"
        )
    if cfg.get("UNET_PRIMARY_MIN_COMPONENT_PX", 0) < 1:
        errors.append("  UNET_PRIMARY_MIN_COMPONENT_PX must be at least 1")
    if cfg.get("UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX", -1) < 0:
        errors.append(
            "  UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX must be nonnegative"
        )
    if cfg.get("CLAHE_MODE") not in ("auto_stack", "auto", "no_clahe", "high_contrast", "standard", "low_signal"):
        errors.append("  CLAHE_MODE must be auto_stack/auto/no_clahe/high_contrast/standard/low_signal")
    if cfg.get("PREPROCESS_MODE") not in ("roi_adaptive", "full_frame"):
        errors.append("  PREPROCESS_MODE must be roi_adaptive or full_frame")
    for key in ("CLAHE_CLIP_HIGH_CONTRAST", "CLAHE_CLIP_STANDARD", "CLAHE_CLIP_LOW_CONTRAST",
                "DENOISE_SIGMA_UM", "MAX_BRIDGE_UM", "MIN_SKEL_LEN_UM", "MAX_WIDTH_UM",
                "MAX_GEODESIC_LEN_UM", "MAX_BRANCH_LEN_UM"):
        if cfg.get(key, 0) < 0:
            errors.append(f"  {key} must be nonnegative")
    if cfg.get("ROI_CROP_PADDING_PX", 0) < 0:
        errors.append("  ROI_CROP_PADDING_PX must be nonnegative")
    if cfg.get("BG_SIGMA_UM", 0) <= 0:
        errors.append("  BG_SIGMA_UM must be positive")
    if not cfg.get("RIDGE_SIGMAS_UM") or any(float(v) <= 0 for v in cfg.get("RIDGE_SIGMAS_UM", [])):
        errors.append("  RIDGE_SIGMAS_UM must contain positive values")
    if not (0 <= cfg.get("MAX_BRIDGE_ANGLE_DEG", -1) <= 180):
        errors.append("  MAX_BRIDGE_ANGLE_DEG must be between 0 and 180")
    if cfg.get("MAX_GEODESIC_LEN_UM", 0) <= cfg.get("MIN_SKEL_LEN_UM", 0):
        errors.append("  MAX_GEODESIC_LEN_UM must exceed MIN_SKEL_LEN_UM")
    if cfg.get("RUN_MODE", "") not in ("single", "batch"):
        errors.append("  RUN_MODE must be 'single' or 'batch'")
    if cfg.get("ANALYSIS_MODE", "") not in ("comparative", "reference_morphology", "legacy"):
        errors.append("  ANALYSIS_MODE must be comparative, reference_morphology, or legacy")
    if cfg.get("AI_PREPROCESSING_MODE", "off") != "off":
        errors.append("  AI_PREPROCESSING_MODE must be 'off' unless explicitly enabled by experimental code")
    if "REJECT_BRANCHES" in cfg:
        errors.append("  'REJECT_BRANCHES' was removed in v8. "
                      "Use MAX_BRANCH_NODES (int) instead:\n"
                      "    REJECT_BRANCHES=True  -> MAX_BRANCH_NODES=0\n"
                      "    REJECT_BRANCHES=False -> MAX_BRANCH_NODES=9999")
    if errors:
        raise ValueError("CONFIG errors:\n" + "\n".join(errors))


# =============================================================================
# UTILITIES
# =============================================================================

_N8 = [(-1,-1),(-1,0),(-1,1),(0,-1),(0,1),(1,-1),(1,0),(1,1)]


def ensure_dir(p):
    """
    Creates directory *p* (and all missing parents) if it does not already exist.

    Uses ``exist_ok=True`` to avoid TOCTOU race conditions in parallel runs.

    Args:
        p (str): Absolute or relative path of the directory to create.
    """
    os.makedirs(p, exist_ok=True)

def get_unique_batch_dir(base_dir):
    """
    Checks for 'batch_output', then 'batch_output_1', 'batch_output_2', etc.
    Returns the first available path.
    """
    candidate = os.path.join(base_dir, "batch_output")
    if not os.path.exists(candidate):
        return candidate

    counter = 1
    while True:
        candidate = os.path.join(base_dir, f"batch_output_{counter}")
        if not os.path.exists(candidate):
            return candidate
        counter += 1


def natural_sort_key(s):
    """
    Key for natural sorting (e.g., img2.tif comes before img10.tif).
    """
    import re
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split(r'(\d+)', os.path.basename(s))]

def extract_z_index(path, sequence_idx=0):
    """
    Extracts Z-index from filename, fallback to sequence index.
    """
    import re
    m = re.search(r"[zZ](\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    m = re.search(r"(\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    return sequence_idx


def ensure_2d_image(img, name="image"):
    """
    Forces an image array to be 2D grayscale.
    Handles squeezing, channel-first (Z,H,W), channel-last (H,W,C), and RGB conversion.
    """
    img = np.asarray(img)
    img = np.squeeze(img)

    if img.ndim == 2:
        return img

    if img.ndim == 3:
        # channel-last RGB/RGBA
        if img.shape[-1] in (3, 4):
            return img[..., 0]
        # channel-first
        if img.shape[0] in (1, 3, 4):
            return img[0]
        # singleton channel
        if img.shape[-1] == 1:
            return img[..., 0]
        if img.shape[0] == 1:
            return img[0]

    raise ValueError(f"{name} must be 2D after loading, got shape {img.shape}")


def robust_imread(path):
    """
    Reads image with multi-engine fallback and forced 2D grayscale enforcement.
    """
    p_lower = path.lower()

    # 1. Primary Loader: tifffile
    if p_lower.endswith(".tif") or p_lower.endswith(".tiff"):
        try:
            img = tifffile.imread(path)
            return ensure_2d_image(img, os.path.basename(path))
        except Exception:
            pass

    # 2. Fallback A: Pillow
    try:
        from PIL import Image
        img = np.array(Image.open(path))
        return ensure_2d_image(img, os.path.basename(path))
    except Exception:
        pass

    # 3. Fallback B: OpenCV
    if _HAVE_CV2:
        try:
            img = _cv2.imread(path, _cv2.IMREAD_UNCHANGED)
            if img is not None:
                return ensure_2d_image(img, os.path.basename(path))
        except Exception:
            pass

    # 4. Fallback C: Matplotlib
    try:
        img = plt.imread(path)
        return ensure_2d_image(img, os.path.basename(path))
    except Exception:
        pass

    raise RuntimeError(f"All image engines failed to read: {os.path.basename(path)}")

def normalize_display(img):
    """
    Contrast-stretches a raw image for colourmap display.

    Uses the 1st-99.5th percentile range rather than min-max to suppress hot
    pixels and dark-corner vignetting artefacts common in fluorescence microscopy.
    The small epsilon (1e-9) prevents division-by-zero on flat images.

    Args:
        img (np.ndarray): Input image of any integer or float dtype.

    Returns:
        np.ndarray: Float32 array clipped to [0, 1] ready for matplotlib display.
    """
    a = img.astype(np.float32)
    lo, hi = np.percentile(a, 1), np.percentile(a, 99.5)
    return np.clip((a - lo) / (hi - lo + 1e-9), 0, 1)


def _imwrite(path, arr_uint8, cmap="gray"):
    """
    Writes a uint8 image to disk using the best available engine.

    Engine priority:
    1. **OpenCV** - fastest; handles large images without matplotlib overhead.
       Converts RGB->BGR before writing (OpenCV internal convention).
    2. **Matplotlib** - fallback when OpenCV is absent; uses ``plt.imsave``
       which respects the *cmap* argument for single-channel saves.

    Args:
        path (str): Destination file path (extension determines format, e.g. ``.png``).
        arr_uint8 (np.ndarray): uint8 image array, shape ``(H, W)`` or ``(H, W, 3)``.
        cmap (str): Matplotlib colourmap name used only for grayscale fallback saves.
    """
    if _HAVE_CV2:
        if arr_uint8.ndim == 2:
            _cv2.imwrite(path, arr_uint8)
        else:
            # OpenCV stores BGR; convert from RGB before writing
            _cv2.imwrite(path, _cv2.cvtColor(arr_uint8, _cv2.COLOR_RGB2BGR))
    else:
        plt.imsave(path, arr_uint8,
                   cmap=(cmap if arr_uint8.ndim == 2 else None),
                   vmin=0, vmax=255)


def save_gray(path, img_float):
    """
    Saves a floating-point single-channel image as an 8-bit grayscale PNG.

    Applies a full min-max stretch (appropriate for debug images where absolute
    intensity is less important than structure visibility).

    Args:
        path (str): Output file path.
        img_float (np.ndarray): Float image of arbitrary range.
    """
    a = img_float.astype(np.float32)
    a = (a - a.min()) / (a.max() - a.min() + 1e-9)
    _imwrite(path, (a * 255).astype(np.uint8), cmap="gray")


def save_mask(path, mask_bool):
    """
    Saves a binary boolean mask as a black-and-white 8-bit PNG.

    Pixels where *mask_bool* is ``True`` are written as 255 (white);
    background pixels are 0 (black). This is the standard convention
    for binary mask overlays in ImageJ / FIJI.

    Args:
        path (str): Output file path.
        mask_bool (np.ndarray[bool]): Boolean mask array.
    """
    _imwrite(path, mask_bool.astype(np.uint8) * 255, cmap="gray")


def load_batch_files(input_dir, pattern):
    """
    Discovers and sorts all image files for a batch run.

    Uses a three-tier fallback strategy so the pipeline works even when
    users supply unusual file extensions or mixed case (.TIF vs .tif):

    1. Exact pattern match (e.g. ``Project001_Series002_z*_ch00.tif``).
    2. ``.tif`` -> ``.tiff`` substitution.
    3. Broad glob over all supported extensions (tif, tiff, png, jpg, jpeg)
       in both lower- and upper-case variants.

    After discovery, files are sorted with :func:`natural_sort_key` so that
    ``z2`` comes before ``z10`` (lexicographic sort would reverse this).  A
    Z-index is extracted from each filename or assigned by sequence position
    as a fallback.

    Args:
        input_dir (str): Directory to search for image files.
        pattern (str): Glob pattern relative to *input_dir*.

    Returns:
        tuple[list[str], list[int]]: Sorted file paths and corresponding Z-indices.

    Raises:
        FileNotFoundError: If no matching images are found after all fallbacks.
    """
    files = glob.glob(os.path.join(input_dir, pattern))
    if not files:
        # Fallback 1: .tif -> .tiff extension swap
        files = glob.glob(os.path.join(input_dir, pattern.replace(".tif", ".tiff")))
    if not files:
        # Fallback 2: Broad scan of all supported image formats
        for ext in ['*.tif', '*.tiff', '*.png', '*.jpg', '*.jpeg']:
            files.extend(glob.glob(os.path.join(input_dir, ext)))
            files.extend(glob.glob(os.path.join(input_dir, ext.upper())))
        files = list(set(files))  # remove cross-extension duplicates

    if not files:
        raise FileNotFoundError(f"No supported image files found in '{input_dir}'")

    # Natural sort preserves correct Z-stack ordering (z2 < z10)
    files = sorted(files, key=natural_sort_key)

    # Extract z-index from filename pattern; fall back to sequence index
    z_idx = [extract_z_index(f, sequence_idx=i) for i, f in enumerate(files)]

    print(f"Found {len(files)} slices: Z = {z_idx}")
    return files, z_idx


def load_roi_mask_file(path, expected_shape=None):
    """Load a boolean ROI mask from .npy or image-mask formats."""
    if not path:
        return None
    path = os.path.abspath(path)
    if not os.path.exists(path):
        raise FileNotFoundError(f"ROI mask not found: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext == ".npy":
        mask = np.load(path).astype(bool)
    else:
        mask = robust_imread(path).astype(bool)
    if mask.ndim > 2:
        mask = mask[..., 0].astype(bool)
    if expected_shape is not None and mask.shape != expected_shape:
        raise ValueError(f"ROI mask shape {mask.shape} does not match image shape {expected_shape}")
    return mask


def filter_results_to_roi(results, skel_label, roi_mask):
    """
    Keep detections whose centroid lies inside ROI and mask label pixels outside it.

    The detector can still use full-frame image statistics; this function only
    limits which detections are accepted for downstream tracking/reporting.
    """
    if roi_mask is None or not results:
        return results, skel_label
    roi_mask = roi_mask.astype(bool)
    filtered = []
    keep = []
    for r in results:
        cy = min(max(int(round(r["centroid_y"])), 0), roi_mask.shape[0] - 1)
        cx = min(max(int(round(r["centroid_x"])), 0), roi_mask.shape[1] - 1)
        if roi_mask[cy, cx]:
            filtered.append(r)
            keep.append(r["label"])
    skel_roi = np.where(np.isin(skel_label, keep) & roi_mask, skel_label, 0).astype(np.int32)
    return filtered, skel_roi


def choose_single_image(cfg):
    """
    Resolves the path to a single test image according to the configured selection mode.

    Three modes are supported (set via ``SINGLE_IMAGE_SELECTION_MODE``):

    - ``'path'``     - Use the hard-coded ``SINGLE_TEST_IMAGE`` path directly.
    - ``'z_index'``  - Find the file whose filename encodes the Z-index in
                      ``SINGLE_Z_INDEX`` (e.g. ``z15`` -> ``...z15_ch00.tif``).
    - ``'dialog'``   - Open a Tkinter file-picker dialog for interactive selection.
                      Falls back gracefully if Tkinter is unavailable.

    Args:
        cfg (dict): Pipeline configuration dictionary.

    Returns:
        str: Absolute path to the selected image file.

    Raises:
        FileNotFoundError: For ``'path'`` mode when the file does not exist, or
                           ``'z_index'`` mode when no file matches the Z-index.
        RuntimeError: For ``'dialog'`` mode when Tkinter fails or the user cancels.
        ValueError: If ``SINGLE_IMAGE_SELECTION_MODE`` is not one of the three valid values.
    """
    mode = cfg["SINGLE_IMAGE_SELECTION_MODE"].lower()
    if mode == "path":
        p = cfg["SINGLE_TEST_IMAGE"]
        if not os.path.exists(p):
            raise FileNotFoundError(f"Not found: {p}")
        return p
    if mode == "z_index":
        z = int(cfg["SINGLE_Z_INDEX"])
        files = (glob.glob(os.path.join(cfg["INPUT_DIR"], cfg["FILE_PATTERN"])) or
                 glob.glob(os.path.join(cfg["INPUT_DIR"],
                                        cfg["FILE_PATTERN"].replace(".tif", ".tiff"))))
        hits = [f for f in files if extract_z_index(f) == z]
        if not hits:
            raise FileNotFoundError(f"No file for z={z}")
        print(f"Using z={z}: {hits[0]}")
        return hits[0]
    if mode == "dialog":
        try:
            import tkinter as tk
            from tkinter import filedialog
            root = tk.Tk(); root.withdraw()  # hide the empty root window
            p = filedialog.askopenfilename(
                title="Choose image", initialdir=cfg["INPUT_DIR"],
                filetypes=[("TIFF", "*.tif *.tiff"), ("All", "*.*")])
            root.destroy()
        except Exception as e:
            raise RuntimeError(f"File dialog failed: {e}")
        if not p:
            raise RuntimeError("No file selected.")
        print(f"Selected: {p}")
        return p
    raise ValueError("SINGLE_IMAGE_SELECTION_MODE: 'path'|'z_index'|'dialog'")


# =============================================================================
# SKELETON UTILITIES
# =============================================================================

def find_endpoints(skel_bool):
    """
    Identifies all endpoint pixels (tip pixels) in a binary skeleton image.

    In an 8-connected skeleton, a pixel is an *endpoint* if exactly one of its
    eight neighbours is also part of the skeleton.  These pixels correspond to
    the physical tips of spermatid filaments.  A clean spermatid should have
    exactly 2 endpoints (head and tail end); a bridged fragment pair still has
    2 endpoints after merging.

    Args:
        skel_bool (np.ndarray[bool]): 2D binary skeleton image (True = skeleton pixel).

    Returns:
        list[tuple[int, int]]: List of ``(row, col)`` pixel coordinates of all endpoints.
    """
    H, W = skel_bool.shape
    ys, xs = np.where(skel_bool)
    sk_set = set(zip(ys.tolist(), xs.tolist()))
    return [(r, c) for r, c in sk_set
            if sum(1 for dr, dc in _N8
                   if 0 <= r+dr < H and 0 <= c+dc < W
                   and skel_bool[r+dr, c+dc]) == 1]


def _endpoint_tangent(endpoint, skel_bool, radius=5):
    r, c = endpoint
    ys, xs = np.where(skel_bool)
    if ys.size < 2:
        return None
    d2 = (ys - r) ** 2 + (xs - c) ** 2
    keep = (d2 > 0) & (d2 <= radius ** 2)
    if np.count_nonzero(keep) < 1:
        return None
    idx = int(np.argmin(d2[keep]))
    pts = np.column_stack([ys[keep], xs[keep]])
    nbr = pts[idx]
    v = np.array([float(nbr[0] - r), float(nbr[1] - c)])
    n = np.linalg.norm(v)
    return v / n if n > 0 else None


def _angle_between(v1, v2):
    if v1 is None or v2 is None:
        return 0.0
    dot = float(np.clip(np.dot(v1, v2), -1.0, 1.0))
    return math.degrees(math.acos(dot))


def bridge_skeleton_endpoints(skel_bool, skel_labeled, max_gap_px, valid_mask=None, max_angle_deg=35.0, return_stats=False):
    """
    Join endpoint pairs from DIFFERENT components within max_gap_px by a
    1-px straight line.  Preserves the original mask so width estimates
    are not inflated.
    """
    stats = {
        "skeleton_pixels_before": int(np.count_nonzero(skel_bool)),
        "skeleton_pixels_after": int(np.count_nonzero(skel_bool)),
        "proposed_bridges": 0,
        "rejected_distance": 0,
        "rejected_orientation": 0,
        "rejected_roi": 0,
        "rejected_exclusion": 0,
        "accepted": 0,
    }
    if max_gap_px <= 0:
        return (skel_bool.copy(), stats) if return_stats else skel_bool.copy()
    H, W   = skel_bool.shape
    out    = skel_bool.copy()
    if valid_mask is None:
        valid_mask = np.ones_like(skel_bool, dtype=bool)
    eps    = find_endpoints(skel_bool)
    if not eps:
        return (out, stats) if return_stats else out
    ep_arr    = np.array(eps, dtype=np.float32)
    ep_labels = skel_labeled[ep_arr[:, 0].astype(int), ep_arr[:, 1].astype(int)]
    pairs     = cKDTree(ep_arr).query_pairs(r=max_gap_px, output_type="ndarray")
    for i, j in pairs:
        stats["proposed_bridges"] += 1
        if ep_labels[i] == ep_labels[j]:
            continue
        r0, c0 = int(eps[i][0]), int(eps[i][1])
        r1, c1 = int(eps[j][0]), int(eps[j][1])
        if math.hypot(r1 - r0, c1 - c0) > max_gap_px:
            stats["rejected_distance"] += 1
            continue
        v0 = _endpoint_tangent((r0, c0), skel_bool)
        v1 = _endpoint_tangent((r1, c1), skel_bool)
        bridge_vec = np.array([float(r1 - r0), float(c1 - c0)])
        bridge_vec /= (np.linalg.norm(bridge_vec) + 1e-9)
        angle0 = _angle_between(v0, bridge_vec)
        angle1 = _angle_between(v1, -bridge_vec)
        if max(angle0, angle1) > max_angle_deg:
            stats["rejected_orientation"] += 1
            continue
        n = max(abs(r1-r0), abs(c1-c0)) + 1
        rs = np.clip(np.round(np.linspace(r0, r1, n)).astype(int), 0, H-1)
        cs = np.clip(np.round(np.linspace(c0, c1, n)).astype(int), 0, W-1)
        if not np.all(valid_mask[rs, cs]):
            stats["rejected_roi"] += 1
            continue
        out[rs, cs] = True
        stats["accepted"] += 1
    out &= valid_mask
    stats["skeleton_pixels_after"] = int(np.count_nonzero(out))
    print(f"    bridge stats: {stats}")
    return (out, stats) if return_stats else out


def prune_branches(skel_bool, max_branch_len):
    """
    Iteratively remove endpoints to shorten side-branches <= max_branch_len px.
    """
    if max_branch_len <= 0:
        return skel_bool.copy()
    H, W  = skel_bool.shape
    skel  = skel_bool.copy()
    for _ in range(int(max_branch_len)):
        eps = find_endpoints(skel)
        if not eps:
            break
        for r, c in eps:
            n = sum(1 for dr, dc in _N8
                    if 0 <= r+dr < H and 0 <= c+dc < W and skel[r+dr, c+dc])
            if n == 1:
                skel[r, c] = False
    return skel


# =============================================================================
# GEODESIC & TOPOLOGY MEASUREMENT
# =============================================================================

def _build_adj(coords, W):
    """
    Builds a lightweight adjacency list for a set of skeleton pixel coordinates.

    Uses a *linearised index* trick: each pixel ``(r, c)`` is mapped to an integer
    ``r * W + c`` so that neighbour lookup is an O(1) dictionary lookup rather than
    a 2D array access.  This is important for large skeleton components.

    Edge weights follow the Chebyshev / Euclidean convention:
    - Axis-aligned move (4-connected step): weight = 1.0
    - Diagonal move (8-connected step):     weight = sqrt2 ~ 1.41421

    Args:
        coords (np.ndarray): Shape ``(N, 2)`` integer array of ``(row, col)`` pixel positions.
        W (int): Image width in pixels (used to compute the linear index).

    Returns:
        list[list[tuple[int, float]]]: Adjacency list where ``adj[i]`` is a list of
        ``(neighbour_index, edge_weight)`` pairs.
    """
    n       = len(coords)
    # Linearise (row, col) -> single int for O(1) membership test
    lin     = coords[:, 0] * W + coords[:, 1]
    lin2idx = {int(v): i for i, v in enumerate(lin.tolist())}
    lin_set = set(lin.tolist())
    adj     = [[] for _ in range(n)]
    for i, (r, c) in enumerate(coords.tolist()):
        for dr, dc in _N8:
            lk = (r + dr) * W + (c + dc)
            if lk in lin_set:
                # Diagonal edges are longer by a factor of sqrt2
                w = 1.41421356 if (dr != 0 and dc != 0) else 1.0
                adj[i].append((lin2idx[lk], w))
    return adj


def _dijkstra(adj, src, n):
    """
    Runs Dijkstra's shortest-path algorithm from source node *src* and returns
    the *farthest* reachable node and its distance.

    This is the first BFS pass in the double-BFS algorithm for computing the
    true *geodesic diameter* (longest shortest path) of a skeleton component.
    The second BFS starts from the farthest node found here.

    Mathematical note:
        Geodesic length = shortest path distance through the skeleton graph,
        which accounts for the actual curvature of the filament rather than
        the straight-line Euclidean distance between endpoints.

    Args:
        adj (list[list[tuple[int, float]]]): Adjacency list from :func:`_build_adj`.
        src (int): Index of the source node in *adj*.
        n (int): Total number of nodes.

    Returns:
        tuple[int, float]: ``(farthest_node_index, distance_to_farthest)``
    """
    d = np.full(n, np.inf); d[src] = 0.0
    pq = [(0.0, src)]
    while pq:
        cost, u = heapq.heappop(pq)
        if cost > d[u]:
            continue  # stale entry in the heap - skip
        for v, w in adj[u]:
            nd = cost + w
            if nd < d[v]:
                d[v] = nd
                heapq.heappush(pq, (nd, v))
    far = int(np.argmax(d))
    return far, float(d[far])


def _dijkstra_with_previous(adj, src, n):
    d = np.full(n, np.inf)
    prev = np.full(n, -1, dtype=np.int32)
    d[src] = 0.0
    pq = [(0.0, src)]
    while pq:
        cost, u = heapq.heappop(pq)
        if cost > d[u]:
            continue
        for v, w in adj[u]:
            nd = cost + w
            if nd < d[v]:
                d[v] = nd
                prev[v] = u
                heapq.heappush(pq, (nd, v))
    far = int(np.argmax(d))
    return far, float(d[far]), prev


def extract_geodesic_centerline_coords(coords, W):
    """
    Return the longest simple shortest-path through a skeleton component.

    This is used only for high-confidence U-Net rescue candidates that are
    rejected for topology. It converts a branched/looped probability skeleton
    into one measurable centerline without accepting all side branches.
    """
    coords = np.asarray(coords)
    if coords.shape[0] < 2:
        return coords
    adj = _build_adj(coords, W)
    start = 0
    b, _ = _dijkstra(adj, start, len(coords))
    c, _dist, prev = _dijkstra_with_previous(adj, b, len(coords))
    path = []
    cur = c
    seen = set()
    while cur >= 0 and cur not in seen:
        seen.add(cur)
        path.append(cur)
        if cur == b:
            break
        cur = int(prev[cur])
    if not path or path[-1] != b:
        return coords
    return coords[np.asarray(path, dtype=np.int32)]


def measure_topology(coords, W, allow_loops=False):
    """
    Compute geodesic length, tortuosity, endpoint count, and branch-node count
    for one skeleton component.

    Returns
    -------
    dict with keys: geo_len, tortuosity, n_endpoints, n_branch_nodes, reason
    or None if the component is a loop and allow_loops=False.
    """
    n   = len(coords)
    adj = _build_adj(coords, W)
    deg = [len(a) for a in adj]

    n_endpoints    = sum(1 for d in deg if d == 1)
    n_branch_nodes = sum(1 for d in deg if d >  2)
    eps_idx        = [i for i, d in enumerate(deg) if d == 1]

    # ------ Loop handling ---------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    if not eps_idx:
        if not allow_loops:
            return None  # discard loop
        seen, total = set(), 0.0
        for u, nbrs in enumerate(adj):
            for v, w in nbrs:
                key = (min(u, v), max(u, v))
                if key not in seen:
                    seen.add(key); total += w
        return {"geo_len": total, "tortuosity": 1.0,
                "n_endpoints": 0, "n_branch_nodes": n_branch_nodes,
                "reason": "loop"}

    # ------ Double-BFS for geodesic ---------------------------------------------------------------------------------------------------------------------------------------------
    b, _  = _dijkstra(adj, eps_idx[0], n)
    c, gl = _dijkstra(adj, b, n)

    # ------ Tortuosity ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    p0  = coords[b]
    p1  = coords[c]
    euc = float(np.sqrt((p0[0]-p1[0])**2 + (p0[1]-p1[1])**2))
    tort = gl / (euc + 1e-9)

    return {"geo_len": float(gl), "tortuosity": tort,
            "n_endpoints": n_endpoints, "n_branch_nodes": n_branch_nodes,
            "reason": "ok"}


# =============================================================================
# SEGMENTATION PIPELINE
# =============================================================================

def apply_optional_early_shape_filter(mask, cfg):
    """
    Optionally removes non-elongated connected components from the binary mask
    *before* skeletonisation, based on region shape descriptors.

    This is an early-stage pre-filter controlled by ``USE_EARLY_SHAPE_FILTER``.
    When enabled, it eliminates round debris, fat nuclei, and large globular
    artefacts that would otherwise generate spurious skeleton trees downstream.

    Biological rationale
    --------------------
    Round debris (dead cells, lipid droplets, imaging artefacts) tends to have:
    - Low eccentricity (close to circular)
    - Short major axis (small, compact object)
    - Low axis ratio (major ~ minor, i.e., not elongated)

    Elongated spermatid nuclei, by contrast, have high eccentricity (> 0.6),
    a long major axis, and a large major/minor axis ratio.

    Args:
        mask (np.ndarray[bool]): Binary mask from hysteresis thresholding.
        cfg (dict): Pipeline config; reads ``USE_EARLY_SHAPE_FILTER``,
            ``MIN_ECCENTRICITY``, ``MAX_MINOR_PX``, ``MIN_AXIS_RATIO``,
            and ``MIN_MAJOR_PX``.

    Returns:
        np.ndarray[bool]: Filtered binary mask (unchanged if the filter is disabled).
    """
    if not cfg["USE_EARLY_SHAPE_FILTER"]:
        return mask  # filter disabled - pass through unchanged
    labeled = measure.label(mask)
    keep = [p.label for p in measure.regionprops(labeled)
            if (p.eccentricity      >= cfg["MIN_ECCENTRICITY"] and   # must be elongated
                p.minor_axis_length <= cfg["MAX_MINOR_PX"]     and   # must be narrow
                p.major_axis_length / (p.minor_axis_length + 1e-9) >= cfg["MIN_AXIS_RATIO"] and  # must be rod-like
                p.major_axis_length >= cfg["MIN_MAJOR_PX"])]
    return np.isin(labeled, keep)


def remove_objects_smaller_than(mask, min_size):
    """Compatibility wrapper preserving old remove_small_objects(min_size=...) semantics."""
    min_size = int(min_size)
    if min_size <= 1:
        return mask
    return morphology.remove_small_objects(mask, max_size=min_size - 1)


@dataclass
class StackPreprocessContext:
    normalization_low: float
    normalization_high: float
    selected_clahe_clip: float
    selected_clahe_profile: str
    contrast_score: float
    sampled_z_indices: list
    roi_percentiles: dict
    saturation_fraction: float
    slice_brightness_statistics: list
    source_dtype: str
    inferred_bit_depth: int
    resolved_pixel_parameters: dict
    configuration_provenance: dict
    image_shape: tuple
    roi_pixel_count: int
    excluded_pixel_count: int


def _json_scalar(v):
    if isinstance(v, np.generic):
        return v.item()
    if isinstance(v, np.ndarray):
        return [_json_scalar(x) for x in v.tolist()]
    if isinstance(v, dict):
        return {str(k): _json_scalar(val) for k, val in v.items()}
    if isinstance(v, (list, tuple)):
        return [_json_scalar(x) for x in v]
    if isinstance(v, pl.Path):
        return str(v)
    return v


def resolve_pixel_parameters(cfg):
    um = float(cfg.get("UM_PER_PX_XY", 1.0))
    if um <= 0:
        raise ValueError("UM_PER_PX_XY must be positive for physical-unit resolution")

    def length_px(um_key, px_key, default=0.0):
        if um_key in cfg and cfg.get(um_key) is not None:
            return float(cfg[um_key]) / um, float(cfg[um_key]), "physical"
        return float(cfg.get(px_key, default)), float(cfg.get(px_key, default)) * um, "legacy_px"

    min_skel_px, min_skel_um, min_skel_src = length_px("MIN_SKEL_LEN_UM", "MIN_SKEL_LEN_PX")
    max_geo_px, max_geo_um, max_geo_src = length_px("MAX_GEODESIC_LEN_UM", "MAX_GEODESIC_LEN_PX")
    max_width_px, max_width_um, max_width_src = length_px("MAX_WIDTH_UM", "MAX_WIDTH_PX")
    max_bridge_px_f, max_bridge_um, max_bridge_src = length_px("MAX_BRIDGE_UM", "MAX_BRIDGE_PX")
    max_branch_px_f, max_branch_um, max_branch_src = length_px("MAX_BRANCH_LEN_UM", "MAX_BRANCH_LEN_PX")
    bg_px, bg_um, bg_src = length_px("BG_SIGMA_UM", "BG_SIGMA")
    denoise_px, denoise_um, denoise_src = length_px("DENOISE_SIGMA_UM", "DENOISE_SIGMA_PX", 0.0)
    if "RIDGE_SIGMAS_UM" in cfg:
        ridge_um = [float(v) for v in cfg["RIDGE_SIGMAS_UM"]]
        ridge_px = [v / um for v in ridge_um]
        ridge_src = "physical"
    else:
        ridge_px = [float(v) for v in cfg.get("RIDGE_SIGMAS", [1, 2])]
        ridge_um = [v * um for v in ridge_px]
        ridge_src = "legacy_px"

    return {
        "physical": {
            "MIN_SKEL_LEN_UM": min_skel_um,
            "MAX_GEODESIC_LEN_UM": max_geo_um,
            "MAX_WIDTH_UM": max_width_um,
            "MAX_BRIDGE_UM": max_bridge_um,
            "MAX_BRANCH_LEN_UM": max_branch_um,
            "BG_SIGMA_UM": bg_um,
            "DENOISE_SIGMA_UM": denoise_um,
            "RIDGE_SIGMAS_UM": ridge_um,
        },
        "pixels": {
            "MIN_SKEL_LEN_PX": min_skel_px,
            "MAX_GEODESIC_LEN_PX": max_geo_px,
            "MAX_WIDTH_PX": max_width_px,
            "MAX_BRIDGE_PX": max(0, int(round(max_bridge_px_f))),
            "MAX_BRANCH_LEN_PX": max(0, int(round(max_branch_px_f))),
            "BG_SIGMA": bg_px,
            "DENOISE_SIGMA": denoise_px,
            "RIDGE_SIGMAS": ridge_px,
        },
        "source": {
            "MIN_SKEL_LEN": min_skel_src,
            "MAX_GEODESIC_LEN": max_geo_src,
            "MAX_WIDTH": max_width_src,
            "MAX_BRIDGE": max_bridge_src,
            "MAX_BRANCH_LEN": max_branch_src,
            "BG_SIGMA": bg_src,
            "DENOISE_SIGMA": denoise_src,
            "RIDGE_SIGMAS": ridge_src,
        },
    }


def cfg_with_resolved_pixels(cfg):
    out = cfg.copy()
    resolved = resolve_pixel_parameters(out)
    out.update(resolved["pixels"])
    out["_RESOLVED_PIXEL_PARAMETERS"] = resolved
    return out


def infer_bit_depth(arr):
    if np.issubdtype(arr.dtype, np.integer):
        return int(np.iinfo(arr.dtype).bits)
    finite = arr[np.isfinite(arr)]
    if finite.size == 0:
        return 32
    mx = float(np.nanmax(finite))
    if mx <= 1.0:
        return 1
    if mx <= 255:
        return 8
    if mx <= 65535:
        return 16
    return 32


def representative_indices(n, count):
    if n <= 0:
        return []
    count = max(1, min(int(count), n))
    return sorted(set(int(round(v)) for v in np.linspace(0, n - 1, count)))


def _clahe_profile_from_score(score, cfg):
    mode = cfg.get("CLAHE_MODE", "auto_stack")
    if mode in ("no_clahe", "high_contrast", "standard", "low_signal"):
        profile = mode
    elif score >= float(cfg.get("AUTO_CONTRAST_HIGH_THRESHOLD", 0.45)):
        profile = "no_clahe"
    elif score >= float(cfg.get("AUTO_CONTRAST_LOW_THRESHOLD", 0.25)):
        profile = "standard"
    else:
        profile = "low_signal"
    clip = {
        "no_clahe": 0.0,
        "high_contrast": float(cfg.get("CLAHE_CLIP_HIGH_CONTRAST", 0.010)),
        "standard": float(cfg.get("CLAHE_CLIP_STANDARD", 0.025)),
        "low_signal": float(cfg.get("CLAHE_CLIP_LOW_CONTRAST", 0.035)),
    }[profile]
    return profile, clip


def build_stack_preprocess_context(image_files, roi_mask, cfg, exclusion_mask=None):
    if not image_files:
        raise ValueError("build_stack_preprocess_context requires at least one image file")
    sample_positions = representative_indices(len(image_files), cfg.get("STACK_QC_SAMPLE_COUNT", 12))
    pooled, slice_stats = [], []
    image_shape = None
    source_dtype = None
    bit_depth = None
    saturated = 0
    total_valid = 0
    warnings_out = []

    for pos in sample_positions:
        arr = ensure_2d_image(robust_imread(image_files[pos]), os.path.basename(image_files[pos]))
        if image_shape is None:
            image_shape = arr.shape
            source_dtype = str(arr.dtype)
            bit_depth = infer_bit_depth(arr)
        elif arr.shape != image_shape:
            raise ValueError(f"Inconsistent sampled image dimensions: {arr.shape} vs {image_shape}")

        if roi_mask is None:
            valid = np.ones(arr.shape, dtype=bool)
            warnings_out.append("missing ROI: stack QC used full frame")
        else:
            if roi_mask.shape != arr.shape:
                raise ValueError(f"ROI shape {roi_mask.shape} does not match image shape {arr.shape}")
            valid = roi_mask.astype(bool).copy()
        if exclusion_mask is not None:
            if exclusion_mask.shape != arr.shape:
                raise ValueError(f"Exclusion mask shape {exclusion_mask.shape} does not match image shape {arr.shape}")
            valid &= ~exclusion_mask.astype(bool)

        vals = arr.astype(np.float64)[valid]
        vals = vals[np.isfinite(vals)]
        if vals.size == 0:
            warnings_out.append(f"empty valid pixels at sampled z position {pos}")
            continue
        pooled.append(vals)
        p95 = float(np.percentile(vals, 95))
        med = float(np.median(vals))
        slice_stats.append({"z_index": int(extract_z_index(image_files[pos])), "median": med, "p95": p95})
        total_valid += int(vals.size)
        if np.issubdtype(arr.dtype, np.integer):
            info = np.iinfo(arr.dtype)
            saturated += int(np.count_nonzero(vals >= info.max))
        else:
            saturated += int(np.count_nonzero(vals >= 1.0)) if float(np.nanmax(vals)) <= 1.5 else int(np.count_nonzero(vals >= np.nanpercentile(vals, 99.99)))

    if not pooled:
        raise ValueError("No finite stack QC pixels inside ROI after exclusion")
    pooled_vals = np.concatenate(pooled)
    pct = {f"p{p:g}": float(np.percentile(pooled_vals, p)) for p in (1, 20, 50, 95, 99.5)}
    low = float(np.percentile(pooled_vals, cfg.get("NORM_LOW_PERCENTILE", 1.0)))
    high = float(np.percentile(pooled_vals, cfg.get("NORM_HIGH_PERCENTILE", 99.5)))
    dyn = max(high - low, 1e-9)
    contrast_score = float((pct["p95"] - pct["p20"]) / dyn)
    profile, clip = _clahe_profile_from_score(contrast_score, cfg)
    medians = np.array([s["median"] for s in slice_stats], dtype=float)
    p95s = np.array([s["p95"] for s in slice_stats], dtype=float)
    cv = float(np.std(medians) / (np.mean(medians) + 1e-9)) if medians.size else 0.0
    if dyn < 1e-6:
        warnings_out.append("very low dynamic range in valid ROI pixels")
    sat_frac = float(saturated / max(total_valid, 1))
    if sat_frac > 0.01:
        warnings_out.append(f"excessive saturation fraction {sat_frac:.4f}")

    return StackPreprocessContext(
        normalization_low=low,
        normalization_high=high,
        selected_clahe_clip=clip,
        selected_clahe_profile=profile,
        contrast_score=contrast_score,
        sampled_z_indices=[int(extract_z_index(image_files[pos])) for pos in sample_positions],
        roi_percentiles=pct,
        saturation_fraction=sat_frac,
        slice_brightness_statistics=[dict(s, brightness_cv=cv, p95_mean=float(np.mean(p95s))) for s in slice_stats],
        source_dtype=source_dtype or "unknown",
        inferred_bit_depth=int(bit_depth or 0),
        resolved_pixel_parameters=resolve_pixel_parameters(cfg),
        configuration_provenance={
            "version": "v5.7-unet-ready",
            "preprocess_mode": cfg.get("PREPROCESS_MODE"),
            "legacy_two_pass_roi": bool(cfg.get("LEGACY_TWO_PASS_ROI", False)),
            "warnings": warnings_out,
        },
        image_shape=tuple(int(v) for v in image_shape),
        roi_pixel_count=int(np.count_nonzero(roi_mask)) if roi_mask is not None else int(np.prod(image_shape)),
        excluded_pixel_count=int(np.count_nonzero(exclusion_mask)) if exclusion_mask is not None else 0,
    )


def save_stack_preprocess_context(context, output_dir):
    ensure_dir(output_dir)
    data = _json_scalar(asdict(context))
    with open(os.path.join(output_dir, "stack_preprocessing_qc.json"), "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    flat = data.copy()
    flat["sampled_z_indices"] = ",".join(str(v) for v in data.get("sampled_z_indices", []))
    flat["configuration_warnings"] = "; ".join(data.get("configuration_provenance", {}).get("warnings", []))
    pd.DataFrame([flat]).to_csv(os.path.join(output_dir, "stack_preprocessing_qc.csv"), index=False)


def _full_like(shape, crop, bbox, dtype=float):
    out = np.zeros(shape, dtype=dtype)
    y0, y1, x0, x1 = bbox
    out[y0:y1, x0:x1] = crop
    return out


def _valid_bbox(mask, pad, shape):
    ys, xs = np.where(mask)
    if ys.size == 0:
        return (0, shape[0], 0, shape[1])
    return (
        max(0, int(ys.min()) - int(pad)),
        min(shape[0], int(ys.max()) + int(pad) + 1),
        max(0, int(xs.min()) - int(pad)),
        min(shape[1], int(xs.max()) + int(pad) + 1),
    )


def _normalize_with_context(img, valid_mask, cfg, preprocess_context):
    valid_vals = img[valid_mask]
    valid_vals = valid_vals[np.isfinite(valid_vals)]
    if valid_vals.size == 0:
        raise ValueError("No finite valid pixels available for slice normalization")
    slice_low = float(np.percentile(valid_vals, cfg.get("NORM_LOW_PERCENTILE", 1.0)))
    slice_high = float(np.percentile(valid_vals, cfg.get("NORM_HIGH_PERCENTILE", 99.5)))
    if preprocess_context is None:
        stack_low, stack_high = slice_low, slice_high
        profile, clip = _clahe_profile_from_score((slice_high - slice_low) / (slice_high + 1e-9), cfg)
    else:
        stack_low = float(preprocess_context.normalization_low)
        stack_high = float(preprocess_context.normalization_high)
        profile = preprocess_context.selected_clahe_profile
        clip = float(preprocess_context.selected_clahe_clip)
    w = float(cfg.get("NORM_STACK_WEIGHT", 0.80))
    blended_low = w * stack_low + (1.0 - w) * slice_low
    blended_high = w * stack_high + (1.0 - w) * slice_high
    if blended_high <= blended_low:
        blended_high = blended_low + 1e-6
    norm = np.clip((img - blended_low) / (blended_high - blended_low), 0, 1)
    return norm, {
        "stack_low": stack_low, "stack_high": stack_high,
        "slice_low": slice_low, "slice_high": slice_high,
        "blended_low": float(blended_low), "blended_high": float(blended_high),
        "profile": profile, "clip": clip,
    }


def _apply_clahe(img_norm, profile, clip, cfg):
    if profile == "no_clahe" or clip <= 0:
        return img_norm.copy()
    return exposure.equalize_adapthist(
        np.clip(img_norm, 0, 1),
        clip_limit=float(clip),
        kernel_size=int(cfg.get("CLAHE_KERNEL", 128)),
    )


def _save_v56_debug(debug_dir, z_idx, stages, debug_record):
    ensure_dir(debug_dir)
    names = [
        ("01_raw_robust_normalized", stages.get("img_norm")),
        ("02_denoised", stages.get("img_denoised")),
        ("03_clahe", stages.get("img_eq")),
        ("04_background", stages.get("background")),
        ("05_foreground", stages.get("foreground")),
        ("06_ridge", stages.get("ridge")),
        ("07_hysteresis", stages.get("mask_hyst")),
        ("08_clean", stages.get("mask_clean")),
        ("09_skeleton_clean", stages.get("skel_clean")),
        ("10_skeleton_bridged", stages.get("skel_bridged")),
        ("11_skeleton_pruned", stages.get("skel_pruned")),
        ("12_final_detections", stages.get("skel_labeled")),
        ("13_unet_probability", stages.get("unet_probability")),
        ("14_unet_candidate_mask", stages.get("unet_candidate_mask")),
        ("15_unet_seed_mask", stages.get("unet_seed_mask")),
    ]
    for name, arr in names:
        path = os.path.join(debug_dir, f"z{int(z_idx or 0):02d}_{name}.png")
        if arr is None:
            continue
        if arr.dtype == bool:
            save_mask(path, arr)
        else:
            save_gray(path, arr)
    with open(os.path.join(debug_dir, f"z{int(z_idx or 0):02d}_debug_record.json"), "w", encoding="utf-8") as f:
        json.dump(_json_scalar(debug_record), f, indent=2)


def _make_unet_context_from_paths(files_by_z, z_idx):
    """Build [z-1, z, z+1] context planes, clamping to the nearest available slice."""
    if not files_by_z:
        return None
    z_keys = sorted(int(z) for z in files_by_z)
    if not z_keys:
        return None
    z_min, z_max = z_keys[0], z_keys[-1]
    planes = []
    for zz in (int(z_idx) - 1, int(z_idx), int(z_idx) + 1):
        zz = min(max(zz, z_min), z_max)
        if zz not in files_by_z:
            nearest = min(z_keys, key=lambda k: abs(k - zz))
            zz = nearest
        arr = robust_imread(files_by_z[zz])
        planes.append(ensure_2d_image(arr, os.path.basename(files_by_z[zz])).astype(np.float32))
    return np.stack(planes, axis=0)


def _apply_unet_candidate_support(mask_hyst, ridge, valid_crop, full_shape, bbox, roi_mask_full, cfg, unet_context_stack, z_idx=None):
    engine = str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower()
    model_path = str(cfg.get("UNET_MODEL_PATH", "")).strip()
    empty = np.zeros(full_shape, dtype=np.float32)
    cache = cfg.get("_UNET_PROBABILITY_CACHE")
    has_cached_map = cache is not None and z_idx is not None and any(
        key in cache for key in (int(z_idx), str(int(z_idx)), f"z{int(z_idx):02d}")
    )
    if engine not in {"unet_assisted", "hybrid", "unet_primary"}:
        return mask_hyst, ridge, empty, empty.astype(bool), empty.astype(bool), {
            "unet_enabled": False,
            "unet_reason": "classical_engine_selected",
        }
    if not has_cached_map and not model_path:
        raise RuntimeError(
            "U-Net segmentation requires UNET_MODEL_PATH; refusing classical-only fallback"
        )
    if not has_cached_map and unet_context_stack is None:
        raise RuntimeError(
            "U-Net segmentation requires a 2.5D context stack; refusing classical-only fallback"
        )
    if not has_cached_map and not os.path.isfile(model_path):
        raise FileNotFoundError(f"U-Net checkpoint not found: {model_path}")

    try:
        unet_prob = None
        cache_hit = False
        if cache is not None and z_idx is not None:
            for key in (int(z_idx), str(int(z_idx)), f"z{int(z_idx):02d}"):
                if key in cache:
                    unet_prob = np.asarray(cache[key], dtype=np.float32)
                    cache_hit = True
                    break
            if unet_prob is not None and unet_prob.shape != full_shape:
                raise ValueError(f"cached U-Net probability shape {unet_prob.shape} does not match {full_shape}")

        if unet_prob is None:
            from utils.saturn_unet25d_bridge import predict_probability_tiled

            unet_prob = predict_probability_tiled(
                unet_context_stack,
                model_path,
                roi_mask=roi_mask_full,
                cfg=cfg,
            )
        unet_prob = np.asarray(unet_prob, dtype=np.float32)
        if unet_prob.shape != full_shape:
            raise ValueError(f"U-Net probability shape {unet_prob.shape} does not match {full_shape}")
        if not np.all(np.isfinite(unet_prob)):
            raise ValueError("U-Net probability map contains non-finite values")
        if float(np.min(unet_prob)) < 0.0 or float(np.max(unet_prob)) > 1.0:
            raise ValueError("U-Net probability map values must be within [0, 1]")
        y0, y1, x0, x1 = bbox
        prob_crop = unet_prob[y0:y1, x0:x1].astype(np.float32)
        cand_thr = float(cfg.get("UNET_CANDIDATE_THRESHOLD", cfg.get("UNET_THRESHOLD", 0.05)))
        seed_thr = float(cfg.get("UNET_SEED_THRESHOLD", max(cand_thr, 0.30)))
        threshold_mode = str(cfg.get("UNET_THRESHOLD_MODE", "soft")).strip().lower()
        candidate_crop = (prob_crop >= cand_thr) & valid_crop
        seed_crop = (prob_crop >= seed_thr) & valid_crop
        mask_action = "none"
        if engine == "unet_assisted":
            mask_hyst = candidate_crop.copy()
            ridge = prob_crop.copy()
            ridge[~valid_crop] = 0.0
            mask_action = "replace_with_unet_candidate"
        elif engine == "hybrid" and threshold_mode in {"hard", "candidate_union", "union"}:
            mask_hyst = (mask_hyst | candidate_crop) & valid_crop
            mask_action = "union_unet_candidate"

        full_valid = np.zeros(full_shape, dtype=bool)
        full_valid[y0:y1, x0:x1] = valid_crop
        full_prob = unet_prob.astype(np.float32)
        full_prob[~full_valid] = 0.0
        full_candidate = np.zeros(full_shape, dtype=bool)
        full_seed = np.zeros(full_shape, dtype=bool)
        full_candidate[y0:y1, x0:x1] = candidate_crop
        full_seed[y0:y1, x0:x1] = seed_crop
        return mask_hyst, ridge, full_prob, full_candidate, full_seed, {
            "unet_enabled": True,
            "unet_probability_source": "cache" if cache_hit else "model",
            "unet_engine": engine,
            "unet_threshold_mode": threshold_mode,
            "unet_mask_action": mask_action,
            "unet_candidate_threshold": cand_thr,
            "unet_seed_threshold": seed_thr,
            "unet_candidate_pixels": int(np.count_nonzero(full_candidate)),
            "unet_seed_pixels": int(np.count_nonzero(full_seed)),
            "unet_probability_mean_inside_roi": float(np.mean(full_prob[roi_mask_full])) if np.any(roi_mask_full) else 0.0,
            "unet_probability_max_inside_roi": float(np.max(full_prob[roi_mask_full])) if np.any(roi_mask_full) else 0.0,
        }
    except Exception as exc:
        if bool(cfg.get("UNET_FAIL_HARD", True)):
            raise RuntimeError(f"U-Net inference failed for z={z_idx}: {exc}") from exc
        print(f"  WARNING: U-Net inference failed for z={z_idx}: {exc}")
        return mask_hyst, ridge, empty, empty.astype(bool), empty.astype(bool), {
            "unet_enabled": False,
            "unet_reason": f"error: {exc}",
        }


def segment_slice(img_raw, cfg, z_idx=None, debug_dir=None, roi_mask=None, preprocess_context=None, exclusion_mask=None, unet_context_stack=None):
    """
    Executes advanced 2D multi-stage morphology segmentation to detect and isolate spermatid nuclei.

    Args:
        img_raw (np.ndarray): The raw 2D grayscale image array of the z-slice.
        cfg (dict): Pipeline hyperparameters including thresholding, CLAHE limits, and morphological radius.
        z_idx (int, optional): The current Z-index integer for debugging context.
        debug_dir (str, optional): Directory to save intermediate thresholding outputs if tracking is on.
        roi_mask (np.ndarray, optional): A boolean boolean layer representing the user's manual bounding box.

    Returns:
        np.ndarray: A labeled discrete array of contiguous pixel bodies corresponding to valid single nuclei.
    """
    cfg = cfg_with_resolved_pixels(cfg)
    img = ensure_2d_image(img_raw, f"segment_slice z={z_idx}").astype(np.float32)
    full_shape = img.shape
    if roi_mask is None:
        roi_mask_full = np.ones(full_shape, dtype=bool)
    else:
        roi_mask_full = roi_mask.astype(bool)
        if roi_mask_full.shape != full_shape:
            raise ValueError(f"roi_mask shape {roi_mask_full.shape} does not match image shape {full_shape}")
    if exclusion_mask is None:
        exclusion_full = np.zeros(full_shape, dtype=bool)
    else:
        exclusion_full = exclusion_mask.astype(bool)
        if exclusion_full.shape != full_shape:
            raise ValueError(f"exclusion_mask shape {exclusion_full.shape} does not match image shape {full_shape}")
    valid_full = roi_mask_full & ~exclusion_full
    if not np.any(valid_full):
        raise ValueError("ROI/exclusion combination leaves no valid pixels")

    bbox = _valid_bbox(roi_mask_full, cfg.get("ROI_CROP_PADDING_PX", 16), full_shape)
    y0, y1, x0, x1 = bbox
    img_crop = img[y0:y1, x0:x1].copy()
    roi_crop = roi_mask_full[y0:y1, x0:x1]
    excl_crop = exclusion_full[y0:y1, x0:x1]
    valid_crop = roi_crop & ~excl_crop

    finite_valid = img_crop[valid_crop]
    finite_valid = finite_valid[np.isfinite(finite_valid)]
    fill_value = float(np.median(finite_valid)) if finite_valid.size else float(np.nanmedian(img_crop))
    work = img_crop.copy()
    work[~roi_crop] = fill_value
    work[excl_crop] = fill_value

    img_norm, norm_record = _normalize_with_context(work, valid_crop, cfg, preprocess_context)
    boundary_safe = bool(cfg.get("ROI_BOUNDARY_SAFE_RIDGE", True))
    norm_fill = float(np.median(img_norm[valid_crop])) if np.any(valid_crop) else 0.0
    if boundary_safe:
        # Keep a continuous exterior through denoising/background/ridge filtering.
        # The exact biological ROI is applied after ridge calculation.
        img_norm[~roi_crop | excl_crop] = norm_fill
    else:
        img_norm[~roi_crop | excl_crop] = 0

    denoise_sigma = float(cfg.get("DENOISE_SIGMA", 0.0))
    img_denoised = gaussian(img_norm, sigma=denoise_sigma) if denoise_sigma > 0 else img_norm.copy()
    if not boundary_safe:
        img_denoised[~roi_crop | excl_crop] = 0
    img_eq = _apply_clahe(img_denoised, norm_record["profile"], norm_record["clip"], cfg)
    if boundary_safe:
        eq_fill = float(np.median(img_eq[valid_crop])) if np.any(valid_crop) else 0.0
        img_eq[~roi_crop | excl_crop] = eq_fill
    else:
        img_eq[~roi_crop | excl_crop] = 0

    bg  = gaussian(img_eq, sigma=float(cfg["BG_SIGMA"]))
    fg  = np.clip(img_eq - bg, 0, None)
    valid_fg = fg[valid_crop]
    hi_fg = float(np.percentile(valid_fg, 99.5)) if valid_fg.size else 1.0
    fgn = np.clip(fg / (hi_fg + 1e-9), 0, 1)
    if boundary_safe:
        fgn_fill = float(np.median(fgn[valid_crop])) if np.any(valid_crop) else 0.0
        fgn[~roi_crop | excl_crop] = fgn_fill
    else:
        fg[~roi_crop | excl_crop] = 0
        fgn[~roi_crop | excl_crop] = 0

    ridge = meijering(fgn, sigmas=cfg["RIDGE_SIGMAS"], black_ridges=False)
    ridge[~roi_crop | excl_crop] = 0

    threshold_crop = valid_crop
    if boundary_safe:
        boundary_exclude_px = max(0, int(cfg.get("ROI_THRESHOLD_EXCLUDE_BOUNDARY_PX", 4)))
        if boundary_exclude_px > 0:
            interior = distance_transform_edt(valid_crop) > boundary_exclude_px
            if np.count_nonzero(interior) > 100:
                threshold_crop = interior
    ridge_valid = ridge[threshold_crop]
    ridge_valid = ridge_valid[np.isfinite(ridge_valid)]
    if ridge_valid.size == 0:
        th_hi, th_lo = 1.0, 0.5
    else:
        if cfg["THRESHOLD_LO"] >= cfg["THRESHOLD_HI"]:
            raise ValueError("THRESHOLD_LO must be < THRESHOLD_HI")
        th_hi = float(np.percentile(ridge_valid, cfg["THRESHOLD_HI"]))
        th_lo = float(np.percentile(ridge_valid, cfg["THRESHOLD_LO"]))
    norm_record["threshold_hi"] = th_hi
    norm_record["threshold_lo"] = th_lo
    mask_hyst = apply_hysteresis_threshold(ridge, th_lo, th_hi)
    mask_hyst &= valid_crop
    classical_mask_hyst = mask_hyst.copy()
    mask_hyst, ridge, unet_prob, unet_candidate, unet_seed, unet_record = _apply_unet_candidate_support(
        mask_hyst,
        ridge,
        valid_crop,
        full_shape,
        bbox,
        roi_mask_full,
        cfg,
        unet_context_stack,
        z_idx=z_idx,
    )

    engine = str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower()
    if engine == "unet_primary":
        full = lambda crop, dtype=None: _full_like(
            full_shape, crop, bbox, dtype=(dtype or crop.dtype)
        )
        primary = _build_unet_primary_segmentation(
            unet_prob,
            valid_full,
            cfg,
            classical_mask=full(classical_mask_hyst, bool),
        )
        primary.update({
            "img_norm": full(img_norm, np.float32),
            "img_denoised": full(img_denoised, np.float32),
            "img_eq": full(img_eq, np.float32),
            "background": full(bg, np.float32),
            "foreground": full(fgn, np.float32),
            "ridge": full(ridge, np.float32),
            "roi_mask": roi_mask_full,
            "exclusion_mask": exclusion_full,
            "unet_probability": unet_prob,
            "unet_candidate_mask": unet_candidate,
            "unet_seed_mask": unet_seed,
            "unet_debug": unet_record,
            "preprocess_context": preprocess_context,
            "preprocess_debug": norm_record,
            "bridge_stats": {
                "skeleton_pixels_before": int(np.count_nonzero(
                    primary["skel_clean"]
                )),
                "skeleton_pixels_after": int(np.count_nonzero(
                    primary["skel_pruned"]
                )),
                "bridges_added": 0,
            },
            "bbox": bbox,
        })
        if cfg["SAVE_DEBUG_IMAGES"] and debug_dir and z_idx is not None:
            _save_v56_debug(
                debug_dir,
                z_idx,
                primary,
                {
                    "z_index": int(z_idx),
                    "segmentation_engine": "unet_primary",
                    "unet_debug": unet_record,
                    "unet_primary_debug": primary["unet_primary_debug"],
                    "outside_roi_skeleton_occupancy": int(np.count_nonzero(
                        primary["skel_pruned"] & ~roi_mask_full
                    )),
                    "exclusion_mask_skeleton_occupancy": int(np.count_nonzero(
                        primary["skel_pruned"] & exclusion_full
                    )),
                },
            )
        return primary

    if mask_hyst.ndim != 2:
        raise ValueError(f"mask_hyst must be 2D, got shape {mask_hyst.shape}")
    mask_clean = mask_hyst.copy()
    if int(cfg["CLOSE_RADIUS"]) > 0:
        mask_clean = morphology.binary_closing(mask_clean, morphology.disk(int(cfg["CLOSE_RADIUS"])))
        mask_clean &= valid_crop
    if int(cfg["MIN_HOLE_AREA"]) > 0:
        mask_clean = morphology.remove_small_holes(mask_clean, area_threshold=int(cfg["MIN_HOLE_AREA"]))
        mask_clean &= valid_crop
    mask_clean = remove_objects_smaller_than(mask_clean, cfg["MIN_OBJ_PX"])
    mask_clean = apply_optional_early_shape_filter(mask_clean, cfg)
    mask_clean &= valid_crop

    # Width is measured from the CLEAN (un-bridged) distance map
    dist_clean   = distance_transform_edt(mask_clean)
    skel_clean   = skeletonize(mask_clean)
    skel_clean &= valid_crop
    skel_labeled = measure.label(skel_clean)

    # Skeleton-level bridging (preserves mask / width integrity)
    skel_bridged, bridge_stats = bridge_skeleton_endpoints(
        skel_clean, skel_labeled, cfg["MAX_BRIDGE_PX"],
        valid_mask=valid_crop,
        max_angle_deg=float(cfg.get("MAX_BRIDGE_ANGLE_DEG", 35.0)),
        return_stats=True)
    # Branch pruning before measurement
    skel_pruned     = prune_branches(skel_bridged, cfg["MAX_BRANCH_LEN_PX"])
    skel_pruned &= valid_crop

    # Optional: Automatically sever complex webs into isolated individual strands
    if cfg.get("BREAK_JUNCTIONS", False):
        from scipy.ndimage import convolve
        kernel = np.array([[1, 1, 1],
                           [1, 0, 1],
                           [1, 1, 1]], dtype=np.int32)
        skel_int = skel_pruned.astype(np.int32)
        neighbors = convolve(skel_int, kernel, mode='constant', cval=0)
        # Any skeleton pixel with more than 2 neighbors is a junction
        junctions = (skel_int > 0) & (neighbors > 2)
        skel_pruned[junctions] = 0
        skel_pruned &= valid_crop

    skel_labeled_fn = measure.label(skel_pruned)

    # NEW: The Recursive Adaptive Micro-Crop Reanalyzer for dense webs and chains
    if cfg.get("MAX_GEODESIC_LEN_PX", 0) > 0 and cfg.get("AUTO_LOCAL_REANALYSIS", True):
        max_px = cfg["MAX_GEODESIC_LEN_PX"]

        # Create a rigid dict of sub-parameters to forcefully shatter the isolated components
        sub_cfg = cfg.copy()
        sub_cfg["AUTO_LOCAL_REANALYSIS"] = False # Prevent infinite recursion

        # Because the bounding box crop contains almost NO dark background,
        # a median 50% threshold perfectly isolates the brightest centers of the blobs!
        sub_cfg["THRESHOLD_HI"] = 55.0
        sub_cfg["THRESHOLD_LO"] = 45.0

        # Deactivate topological limits in the recursive sub-call so fragments aren't
        # instantly dropped before they can be spliced back into the master image.
        sub_cfg["MIN_SKEL_LEN_PX"] = 1.0
        sub_cfg["MIN_OBJ_PX"] = 3
        sub_cfg["MIN_HOLE_AREA"] = 0

        props = measure.regionprops(skel_labeled_fn)
        for sp in props:
            # Note: A dense 2D web might have len(sp.coords) = 5000 pixels.
            if len(sp.coords) > max_px:
                minr, minc, maxr, maxc = sp.bbox
                pad = 12
                minr = max(0, minr - pad)
                minc = max(0, minc - pad)
                maxr = min(img.shape[0], maxr + pad)
                maxc = min(img.shape[1], maxc + pad)

                crop_img = img[minr:maxr, minc:maxc]

                # Create a strict ROI mask over the target structure to ignore neighbors
                obj_mask = (skel_labeled_fn[minr:maxr, minc:maxc] == sp.label)
                crop_roi = morphology.dilation(obj_mask, morphology.disk(6))

                try:
                    # RECURE: Run the entire engine on the tiny isolated crop!
                    sub_seg = segment_slice(crop_img, sub_cfg, roi_mask=crop_roi)
                    sub_skel = sub_seg["skel_pruned"]
                    sub_lab = measure.label(sub_skel)

                    if sub_lab.max() > 1:
                        # SUCCESS: The adaptive threshold organically shattered the web!
                        skel_pruned[minr:maxr, minc:maxc][obj_mask] = 0
                        new_frags = (sub_skel > 0)
                        skel_pruned[minr:maxr, minc:maxc][new_frags] = 1
                    else:
                        # Failsafe: if the intensity was perfectly uniform, geometric centroid chop
                        cy, cx = sp.centroid
                        dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                        mid_idx = np.argmin(dists)
                        my, mx = sp.coords[mid_idx]
                        skel_pruned[my-1:my+2, mx-1:mx+2] = 0

                except Exception:
                    # Absolute geometric failsafe
                    cy, cx = sp.centroid
                    dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                    mid_idx = np.argmin(dists)
                    my, mx = sp.coords[mid_idx]
                    skel_pruned[my-1:my+2, mx-1:mx+2] = 0

        # Run one final relabeling array refresh after all splicing
        skel_pruned &= valid_crop
        skel_labeled_fn = measure.label(skel_pruned)

    full = lambda crop, dtype=None: _full_like(full_shape, crop, bbox, dtype=(dtype or crop.dtype))
    out = {
        "mask_hyst":    full(mask_hyst, bool),
        "mask_clean":   full(mask_clean, bool),
        "skel_clean":   full(skel_clean, bool),
        "skel_bridged": full(skel_bridged, bool),
        "skel_pruned":  full(skel_pruned, bool),
        "skel_labeled": full(skel_labeled_fn.astype(np.int32), np.int32),
        "dist_clean":   full(dist_clean, np.float32),
        "img_norm":     full(img_norm, np.float32),
        "img_denoised": full(img_denoised, np.float32),
        "img_eq":       full(img_eq, np.float32),
        "background":   full(bg, np.float32),
        "foreground":   full(fgn, np.float32),
        "ridge":        full(ridge, np.float32),
        "roi_mask":     roi_mask_full,
        "exclusion_mask": exclusion_full,
        "unet_probability": unet_prob,
        "unet_candidate_mask": unet_candidate,
        "unet_seed_mask": unet_seed,
        "unet_debug": unet_record,
        "preprocess_context": preprocess_context,
        "preprocess_debug": norm_record,
        "bridge_stats": bridge_stats,
        "bbox": bbox,
    }

    if cfg["SAVE_DEBUG_IMAGES"] and debug_dir and z_idx is not None:
        valid_full_after = roi_mask_full & ~exclusion_full
        debug_record = {
            "z_index": z_idx,
            "stack_normalization_low": norm_record["stack_low"],
            "stack_normalization_high": norm_record["stack_high"],
            "slice_normalization_low": norm_record["slice_low"],
            "slice_normalization_high": norm_record["slice_high"],
            "blended_normalization_low": norm_record["blended_low"],
            "blended_normalization_high": norm_record["blended_high"],
            "selected_clahe_profile": norm_record["profile"],
            "selected_clahe_clip": norm_record["clip"],
            "denoise_sigma_um": cfg.get("DENOISE_SIGMA_UM"),
            "denoise_sigma_px": cfg.get("DENOISE_SIGMA"),
            "background_sigma_um": cfg.get("BG_SIGMA_UM"),
            "background_sigma_px": cfg.get("BG_SIGMA"),
            "ridge_sigmas_um": cfg.get("RIDGE_SIGMAS_UM"),
            "ridge_sigmas_px": cfg.get("RIDGE_SIGMAS"),
            "numeric_ridge_high_threshold": th_hi,
            "numeric_ridge_low_threshold": th_lo,
            "valid_roi_pixel_count": int(np.count_nonzero(valid_full_after)),
            "foreground_occupancy_inside_roi": float(np.count_nonzero(out["mask_hyst"] & valid_full_after) / max(np.count_nonzero(valid_full_after), 1)),
            "foreground_occupancy_outside_roi": int(np.count_nonzero(out["mask_hyst"] & ~roi_mask_full)),
            "foreground_occupancy_inside_exclusion_mask": int(np.count_nonzero(out["mask_hyst"] & exclusion_full)),
            "unet_debug": unet_record,
            "skeleton_pixels_before_bridging": bridge_stats["skeleton_pixels_before"],
            "skeleton_pixels_after_bridging": bridge_stats["skeleton_pixels_after"],
            "bridge_inflation_fraction": float((bridge_stats["skeleton_pixels_after"] - bridge_stats["skeleton_pixels_before"]) / max(bridge_stats["skeleton_pixels_before"], 1)),
            "final_detection_count": int(np.max(out["skel_labeled"])),
            "outside_roi_skeleton_occupancy": int(np.count_nonzero(out["skel_pruned"] & ~roi_mask_full)),
            "exclusion_mask_skeleton_occupancy": int(np.count_nonzero(out["skel_pruned"] & exclusion_full)),
        }
        _save_v56_debug(debug_dir, z_idx, out, debug_record)

    return out


# =============================================================================
# MEASUREMENT  (single geodesic pass, all topology in one function)
# =============================================================================

def _expanded_component_bbox(mask, bbox=None, padding=1):
    """Return a clipped component bounding box with a small background margin."""
    h, w = mask.shape
    if bbox is None:
        ys, xs = np.nonzero(mask)
        if ys.size == 0:
            return 0, 0, 0, 0
        min_y, min_x = int(ys.min()), int(xs.min())
        max_y, max_x = int(ys.max()) + 1, int(xs.max()) + 1
    else:
        min_y, min_x, max_y, max_x = (int(v) for v in bbox)
    pad = max(0, int(padding))
    return (
        max(0, min_y - pad),
        max(0, min_x - pad),
        min(h, max_y + pad),
        min(w, max_x + pad),
    )


def _distance_transform_component(mask, bbox=None):
    """
    Compute an exact component distance transform on its local bounding box.

    A full-sized result is returned for compatibility with existing measurement
    code, but the expensive EDT itself runs only around the foreground object.
    """
    mask = np.asarray(mask, dtype=bool)
    out = np.zeros(mask.shape, dtype=np.float64)
    y0, x0, y1, x1 = _expanded_component_bbox(mask, bbox=bbox, padding=1)
    if y1 <= y0 or x1 <= x0:
        return out
    out[y0:y1, x0:x1] = distance_transform_edt(mask[y0:y1, x0:x1])
    return out


_UNET_PRIMARY_REASON_CODES = {
    "tiny_isolated_noise": 1,
    "no_high_confidence_seed": 2,
    "no_valid_centerline": 3,
    "invalid_geometry": 4,
    "unresolved_multi_instance_merge": 5,
}


def _build_unet_primary_foreground(probability, valid_mask, cfg):
    """Build seed-connected U-Net foreground without biological shape gates."""
    probability = np.asarray(probability, dtype=np.float32)
    valid_mask = np.asarray(valid_mask, dtype=bool)
    if probability.shape != valid_mask.shape:
        raise ValueError("U-Net probability and valid-mask shapes must match")
    if not np.all(np.isfinite(probability)):
        raise ValueError("U-Net probability map contains non-finite values")
    if float(np.min(probability)) < 0.0 or float(np.max(probability)) > 1.0:
        raise ValueError("U-Net probability map values must be within [0, 1]")

    low = float(cfg["UNET_CANDIDATE_THRESHOLD"])
    high = float(cfg["UNET_SEED_THRESHOLD"])
    if not 0.0 <= low < high <= 1.0:
        raise ValueError(
            "unet_primary requires candidate threshold < seed threshold in [0, 1]"
        )
    min_area = max(1, int(cfg.get("UNET_PRIMARY_MIN_COMPONENT_PX", 3)))
    probability = probability.copy()
    probability[~valid_mask] = 0.0
    low_mask = (probability >= low) & valid_mask
    seed_mask = (probability >= high) & valid_mask
    hysteresis = apply_hysteresis_threshold(probability, low, high) & valid_mask

    retained = np.zeros_like(valid_mask)
    rejected = np.zeros(valid_mask.shape, dtype=np.uint8)
    audit = []
    low_labels = measure.label(low_mask)
    for prop in measure.regionprops(low_labels):
        component = np.zeros_like(valid_mask)
        component[prop.coords[:, 0], prop.coords[:, 1]] = True
        contains_seed = bool(np.any(seed_mask & component))
        if int(prop.area) < min_area:
            reason = "tiny_isolated_noise"
        elif not contains_seed:
            reason = "no_high_confidence_seed"
        else:
            reason = ""
            retained |= component & hysteresis
        if reason:
            rejected[component] = _UNET_PRIMARY_REASON_CODES[reason]
        audit.append({
            "parent_component_id": int(prop.label),
            "marker_count": 0,
            "child_instance_count": 0,
            "child_instance_id": 0,
            "child_area": int(prop.area),
            "maximum_probability": float(np.max(probability[component])),
            "mean_probability": float(np.mean(probability[component])),
            "contains_seed": contains_seed,
            "disposition": "rejected" if reason else "foreground_retained",
            "technical_reason": reason,
        })
    retained &= valid_mask
    return retained, seed_mask, rejected, audit


def _split_unet_probability_instances(
    probability,
    foreground,
    seed_mask,
    min_component_px,
    compactness=0.0,
):
    """Split each hysteresis component with connected high-seed markers."""
    probability = np.asarray(probability, dtype=np.float32)
    foreground = np.asarray(foreground, dtype=bool)
    seed_mask = np.asarray(seed_mask, dtype=bool)
    min_area = max(1, int(min_component_px))
    output = np.zeros(foreground.shape, dtype=np.int32)
    rejected = np.zeros(foreground.shape, dtype=np.uint8)
    parent_by_instance = {}
    audit = []
    next_label = 1
    parents = measure.label(foreground)

    for parent_prop in measure.regionprops(parents):
        y0, x0, y1, x1 = _expanded_component_bbox(
            foreground, bbox=parent_prop.bbox, padding=1
        )
        component = parents[y0:y1, x0:x1] == parent_prop.label
        prob_crop = probability[y0:y1, x0:x1]
        seed_crop = seed_mask[y0:y1, x0:x1] & component
        seed_labels = measure.label(seed_crop)
        valid_markers = np.zeros(component.shape, dtype=np.int32)
        marker_count = 0
        for seed_prop in measure.regionprops(seed_labels):
            if int(seed_prop.area) < min_area:
                continue
            marker_count += 1
            valid_markers[seed_labels == seed_prop.label] = marker_count

        if marker_count == 0:
            rejected_crop = rejected[y0:y1, x0:x1]
            rejected_crop[component] = _UNET_PRIMARY_REASON_CODES[
                "no_high_confidence_seed"
            ]
            audit.append({
                "parent_component_id": int(parent_prop.label),
                "marker_count": 0,
                "child_instance_count": 0,
                "child_instance_id": 0,
                "child_area": int(parent_prop.area),
                "maximum_probability": float(np.max(prob_crop[component])),
                "mean_probability": float(np.mean(prob_crop[component])),
                "contains_seed": False,
                "disposition": "rejected",
                "technical_reason": "no_high_confidence_seed",
            })
            continue

        if marker_count == 1:
            local_labels = component.astype(np.int32)
        else:
            local_labels = skseg.watershed(
                -prob_crop,
                markers=valid_markers,
                mask=component,
                compactness=max(0.0, float(compactness)),
            ).astype(np.int32)

        accepted_children = []
        for local_id in sorted(int(v) for v in np.unique(local_labels) if v > 0):
            child = local_labels == local_id
            area = int(np.count_nonzero(child))
            contains_seed = bool(np.any(seed_crop & child))
            reason = ""
            if area < min_area:
                reason = "tiny_isolated_noise"
            elif not contains_seed:
                reason = "no_high_confidence_seed"
            if reason:
                rejected_crop = rejected[y0:y1, x0:x1]
                rejected_crop[child] = _UNET_PRIMARY_REASON_CODES[reason]
                disposition = "rejected"
                child_id = 0
            else:
                output_crop = output[y0:y1, x0:x1]
                output_crop[child] = next_label
                child_id = next_label
                parent_by_instance[next_label] = int(parent_prop.label)
                accepted_children.append(next_label)
                next_label += 1
                disposition = "accepted"
            audit.append({
                "parent_component_id": int(parent_prop.label),
                "marker_count": int(marker_count),
                "child_instance_count": 0,
                "child_instance_id": int(child_id),
                "child_area": area,
                "maximum_probability": float(np.max(prob_crop[child])),
                "mean_probability": float(np.mean(prob_crop[child])),
                "contains_seed": contains_seed,
                "disposition": disposition,
                "technical_reason": reason,
            })
        for row in audit:
            if row["parent_component_id"] == int(parent_prop.label):
                row["child_instance_count"] = len(accepted_children)

    return output, rejected, audit, parent_by_instance


def _centerline_unet_primary_instances(instance_labels):
    """Create one deterministic longest-geodesic centerline per instance."""
    instance_labels = np.asarray(instance_labels, dtype=np.int32)
    centerlines = np.zeros_like(instance_labels, dtype=np.int32)
    metadata = {}
    failures = []
    width = instance_labels.shape[1]
    for prop in measure.regionprops(instance_labels):
        instance_id = int(prop.label)
        mask = instance_labels == instance_id
        skeleton = skeletonize(mask)
        coords = np.argwhere(skeleton)
        if coords.size == 0:
            failures.append({
                "instance_id": instance_id,
                "technical_reason": "no_valid_centerline",
            })
            continue
        raw_adj = _build_adj(coords, width)
        raw_degrees = [len(neighbors) for neighbors in raw_adj]
        raw_branch_count = sum(degree > 2 for degree in raw_degrees)
        raw_endpoint_count = sum(degree == 1 for degree in raw_degrees)
        candidates = []
        for skel_prop in measure.regionprops(measure.label(skeleton)):
            component_coords = skel_prop.coords
            path_coords = extract_geodesic_centerline_coords(
                component_coords, width
            )
            path_topology = measure_topology(
                path_coords, width, allow_loops=True
            )
            path_length = (
                float(path_topology["geo_len"]) if path_topology else 0.0
            )
            first_coord = tuple(
                int(v) for v in np.min(path_coords, axis=0).tolist()
            )
            candidates.append((
                path_length,
                int(path_coords.shape[0]),
                tuple(-value for value in first_coord),
                path_coords,
            ))
        center_coords = max(candidates, key=lambda item: item[:3])[3]
        if center_coords.size == 0:
            failures.append({
                "instance_id": instance_id,
                "technical_reason": "no_valid_centerline",
            })
            continue
        centerlines[center_coords[:, 0], center_coords[:, 1]] = instance_id
        metadata[instance_id] = {
            "raw_skeleton_pixels": int(coords.shape[0]),
            "centerline_pixels": int(center_coords.shape[0]),
            "raw_branch_count": int(raw_branch_count),
            "raw_endpoint_count": int(raw_endpoint_count),
            "centerline_salvaged": bool(center_coords.shape[0] < coords.shape[0]),
        }
    return centerlines, metadata, failures


def _build_unet_primary_segmentation(
    probability,
    valid_mask,
    cfg,
    classical_mask=None,
):
    """Build filled U-Net instances and mapped centerlines for unet_primary."""
    foreground, seed_mask, rejected, foreground_audit = (
        _build_unet_primary_foreground(probability, valid_mask, cfg)
    )
    instances, split_rejected, split_audit, parent_map = (
        _split_unet_probability_instances(
            probability,
            foreground,
            seed_mask,
            cfg.get("UNET_PRIMARY_MIN_COMPONENT_PX", 3),
            compactness=cfg.get("UNET_INSTANCE_WATERSHED_COMPACTNESS", 0.0),
        )
    )
    rejected = np.maximum(rejected, split_rejected)
    centerlines, centerline_meta, failures = (
        _centerline_unet_primary_instances(instances)
    )

    for failure in failures:
        instance_id = int(failure["instance_id"])
        mask = instances == instance_id
        rejected[mask] = _UNET_PRIMARY_REASON_CODES["no_valid_centerline"]
        instances[mask] = 0
        centerlines[centerlines == instance_id] = 0
        parent_map.pop(instance_id, None)

    instance_sources = {int(v): "unet_primary" for v in np.unique(instances) if v}
    additions = 0
    if bool(cfg.get("UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE", False)):
        if classical_mask is None:
            raise ValueError(
                "classical_mask is required when U-Net-primary additions are enabled"
            )
        occupied = instances > 0
        dilation_px = max(
            0, int(cfg.get("UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX", 2))
        )
        if dilation_px:
            occupied = morphology.binary_dilation(
                occupied, morphology.disk(dilation_px)
            )
        residual = np.asarray(classical_mask, dtype=bool) & valid_mask & ~occupied
        residual = remove_objects_smaller_than(
            residual, max(1, int(cfg.get("MIN_OBJ_PX", 3)))
        )
        next_id = int(instances.max()) + 1
        residual_labels = measure.label(residual)
        for prop in measure.regionprops(residual_labels):
            mask = residual_labels == prop.label
            skeleton = skeletonize(mask)
            coords = np.argwhere(skeleton)
            if coords.size == 0:
                continue
            center_coords = extract_geodesic_centerline_coords(
                coords, instances.shape[1]
            )
            if center_coords.size == 0:
                continue
            instances[mask] = next_id
            centerlines[
                center_coords[:, 0], center_coords[:, 1]
            ] = next_id
            instance_sources[next_id] = "saturn_only_addition"
            parent_map[next_id] = 0
            centerline_meta[next_id] = {
                "raw_skeleton_pixels": int(coords.shape[0]),
                "centerline_pixels": int(center_coords.shape[0]),
                "raw_branch_count": 0,
                "raw_endpoint_count": 0,
                "centerline_salvaged": False,
            }
            additions += 1
            next_id += 1

    distance = np.zeros(instances.shape, dtype=np.float32)
    for prop in measure.regionprops(instances):
        mask = instances == prop.label
        local_distance = _distance_transform_component(mask, bbox=prop.bbox)
        distance[mask] = local_distance[mask]
    component_audit = foreground_audit + split_audit
    return {
        "mask_hyst": foreground,
        "mask_clean": instances > 0,
        "skel_clean": centerlines > 0,
        "skel_bridged": centerlines > 0,
        "skel_pruned": centerlines > 0,
        "skel_labeled": centerlines,
        "dist_clean": distance,
        "unet_primary_hysteresis_mask": foreground,
        "unet_primary_instance_labels": instances,
        "unet_primary_centerline_labels": centerlines,
        "unet_probability": np.asarray(probability, dtype=np.float32),
        "unet_seed_mask": seed_mask,
        "unet_primary_rejected_reason": rejected,
        "unet_primary_component_audit": component_audit,
        "unet_primary_debug": {
            "candidate_pixels": int(np.count_nonzero(
                (probability >= float(cfg["UNET_CANDIDATE_THRESHOLD"]))
                & valid_mask
            )),
            "seed_pixels": int(np.count_nonzero(seed_mask)),
            "hysteresis_component_count": int(measure.label(foreground).max()),
            "split_instance_count": int(
                np.count_nonzero(np.unique(instances) > 0)
            ),
            "technical_failure_count": len(failures),
            "saturn_only_additions": int(additions),
            "reason_codes": dict(_UNET_PRIMARY_REASON_CODES),
        },
        "unet_primary_parent_by_instance": parent_map,
        "unet_primary_centerline_metadata": centerline_meta,
        "unet_primary_instance_sources": instance_sources,
        "unet_primary_technical_failures": failures,
    }


def _split_unet_rescue_instances(prob, mask, min_component_px, seed_threshold, peak_min_distance, compactness):
    """
    Split a connected U-Net probability mask into putative nucleus instances.

    The U-Net probability map is trusted as the biological evidence layer; this
    helper only separates connected probability regions before Saturn measures
    length/width/topology. It does not make final accept/reject decisions.
    """
    mask = np.asarray(mask, dtype=bool)
    prob = np.asarray(prob, dtype=np.float32)
    if not np.any(mask):
        return np.zeros(mask.shape, dtype=np.int32)

    final = np.zeros(mask.shape, dtype=np.int32)
    next_label = 1
    components = measure.label(mask)
    for comp in measure.regionprops(components):
        if int(comp.area) < min_component_px:
            continue
        y0, x0, y1, x1 = _expanded_component_bbox(mask, bbox=comp.bbox, padding=1)
        comp_mask = components[y0:y1, x0:x1] == comp.label
        prob_crop = prob[y0:y1, x0:x1]

        core = comp_mask & (prob_crop >= seed_threshold)
        core = remove_objects_smaller_than(core, max(1, int(min_component_px)))
        markers = measure.label(core)

        if int(markers.max()) < 2:
            coords = feature.peak_local_max(
                prob_crop,
                labels=comp_mask.astype(np.uint8),
                min_distance=max(1, int(peak_min_distance)),
                threshold_abs=float(seed_threshold),
                exclude_border=False,
            )
            markers = np.zeros(comp_mask.shape, dtype=np.int32)
            for i, (yy, xx) in enumerate(coords, start=1):
                markers[int(yy), int(xx)] = i
            if coords.shape[0] >= 2:
                markers = morphology.dilation(markers, morphology.disk(1))
                markers *= comp_mask
            else:
                markers = measure.label(comp_mask)

        labels = skseg.watershed(
            -prob_crop,
            markers=markers,
            mask=comp_mask,
            compactness=max(0.0, float(compactness)),
        )
        for sub in measure.regionprops(labels):
            sub_mask = labels == sub.label
            if np.count_nonzero(sub_mask) < min_component_px:
                continue
            final_crop = final[y0:y1, x0:x1]
            final_crop[sub_mask] = next_label
            next_label += 1
    return final.astype(np.int32)


def _measure_unet_primary_instances(seg, cfg):
    """Measure mapped U-Net instances; morphology is warning-only."""
    instance_labels = np.asarray(
        seg["unet_primary_instance_labels"], dtype=np.int32
    )
    centerline_labels = np.asarray(
        seg["unet_primary_centerline_labels"], dtype=np.int32
    )
    probability = np.asarray(seg["unet_probability"], dtype=np.float32)
    parent_map = seg.get("unet_primary_parent_by_instance", {})
    source_map = seg.get("unet_primary_instance_sources", {})
    centerline_meta = seg.get("unet_primary_centerline_metadata", {})
    record_morphology_warnings = bool(
        cfg.get("UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS", True)
    )
    results = []
    technical_failures = list(seg.get("unet_primary_technical_failures", []))
    width_px = instance_labels.shape[1]

    for prop in measure.regionprops(instance_labels):
        instance_id = int(prop.label)
        instance_mask = instance_labels == instance_id
        center_coords = np.argwhere(centerline_labels == instance_id)
        if center_coords.size == 0:
            technical_failures.append({
                "instance_id": instance_id,
                "technical_reason": "no_valid_centerline",
            })
            continue
        topology = measure_topology(center_coords, width_px, allow_loops=True)
        if topology is None:
            technical_failures.append({
                "instance_id": instance_id,
                "technical_reason": "invalid_geometry",
            })
            continue
        distance = _distance_transform_component(instance_mask, bbox=prop.bbox)
        widths = 2.0 * distance[
            center_coords[:, 0], center_coords[:, 1]
        ]
        probability_values = probability[instance_mask]
        finite_prob = probability_values[np.isfinite(probability_values)]
        values = [
            topology["geo_len"],
            topology["tortuosity"],
            *widths.tolist(),
        ]
        if not widths.size or not np.all(np.isfinite(values)):
            technical_failures.append({
                "instance_id": instance_id,
                "technical_reason": "invalid_geometry",
            })
            continue

        geodesic = float(topology["geo_len"])
        median_width = float(np.median(widths))
        ratio = geodesic / (median_width + 1e-9)
        metadata = centerline_meta.get(instance_id, {})
        warning_reasons = []
        if geodesic < float(cfg["MIN_SKEL_LEN_PX"]):
            warning_reasons.append("short")
        if geodesic > float(cfg["MAX_GEODESIC_LEN_PX"]):
            warning_reasons.append("long_merge_review")
        if median_width > float(cfg["MAX_WIDTH_PX"]):
            warning_reasons.append("wide")
        if ratio < float(cfg["MIN_LENGTH_WIDTH_RATIO"]):
            warning_reasons.append("low_length_width_ratio")
        if (
            int(topology["n_endpoints"]) >= 2
            and float(topology["tortuosity"]) > float(cfg["MAX_TORTUOSITY"])
        ):
            warning_reasons.append("tortuous")
        if int(metadata.get("raw_branch_count", 0)) > 0:
            warning_reasons.append("branched_centerline_reduced")
        if int(topology["n_endpoints"]) > int(cfg["MAX_ENDPOINT_COUNT"]):
            warning_reasons.append("excess_endpoints")
        if not record_morphology_warnings:
            warning_reasons = []

        cy, cx = prop.centroid
        results.append({
            "label": instance_id,
            "length_px_geodesic": geodesic,
            "length_px_count": float(center_coords.shape[0]),
            "width_px": median_width,
            "length_width_ratio": ratio,
            "tortuosity": float(topology["tortuosity"]),
            "n_endpoints": int(topology["n_endpoints"]),
            "n_branch_nodes": int(metadata.get(
                "raw_branch_count", topology["n_branch_nodes"]
            )),
            "centroid_x": float(cx),
            "centroid_y": float(cy),
            "area_px": float(geodesic * median_width),
            "skeleton_area_px": float(center_coords.shape[0]),
            "instance_mask_area_px": float(prop.area),
            "bbox_min_y": float(prop.bbox[0]),
            "bbox_min_x": float(prop.bbox[1]),
            "bbox_max_y": float(prop.bbox[2]),
            "bbox_max_x": float(prop.bbox[3]),
            "orientation": float(prop.orientation),
            "unet_mean_probability": (
                float(np.mean(finite_prob)) if finite_prob.size else np.nan
            ),
            "unet_max_probability": (
                float(np.max(finite_prob)) if finite_prob.size else np.nan
            ),
            "detection_source": source_map.get(instance_id, "unet_primary"),
            "morphology_warning": bool(warning_reasons),
            "morphology_warning_reasons": ";".join(warning_reasons),
            "technical_failure": False,
            "technical_failure_reason": "",
            "parent_hysteresis_component_id": int(
                parent_map.get(instance_id, 0)
            ),
        })

    return {
        "skel_label": centerline_labels,
        "results": results,
        "unet_rescue_accepted_label": np.zeros_like(
            centerline_labels, dtype=np.int32
        ),
        "unet_rescue_rejected_reason": np.asarray(
            seg.get(
                "unet_primary_rejected_reason",
                np.zeros_like(centerline_labels, dtype=np.uint8),
            ),
            dtype=np.uint8,
        ),
        "unet_rescue_reason_codes": dict(_UNET_PRIMARY_REASON_CODES),
        "unet_rescue_rejected_counts": {},
        "unet_rescue_candidate_threshold": float(
            cfg.get("UNET_CANDIDATE_THRESHOLD", 0.0)
        ),
        "unet_rescue_seed_threshold": float(
            cfg.get("UNET_SEED_THRESHOLD", 0.0)
        ),
        "unet_primary_technical_failures": technical_failures,
    }


def measure_spermatids(seg, cfg):
    """
    Analyzes mathematically discretised nuclei arrays and derives geometric indices for individual shape profiles.

    Args:
        seg (dict): Dictionary containing labeled binary masks ("primary", "skel_pruned", "dist_clean").
        cfg (dict): Pipeline hyperparameters for absolute scaling calculations.

    Returns:
        pd.DataFrame: Table indexing each isolated object dynamically with mathematical biometrics.
                      Attributes include region area, bounding-boxes, minor/major ellipse axis mapping,
                      orientation, and geodesic structural lengths via binary skeleton processing.
    """
    cfg = cfg_with_resolved_pixels(cfg)
    if (
        str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower()
        == "unet_primary"
    ):
        return _measure_unet_primary_instances(seg, cfg)
    skel     = seg["skel_pruned"]
    dist     = seg["dist_clean"]
    skel_lab = seg["skel_labeled"]
    unet_prob = seg.get("unet_probability")
    H, W     = skel.shape

    # ------ Filter pass ---------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
    accepted_labels = []
    cache           = {}
    reasons = {"short": 0, "loop": 0, "long": 0, "wide": 0, "ratio": 0, "branches": 0, "tortuous": 0, "endpoints": 0}

    for sp in measure.regionprops(skel_lab):
        coords = sp.coords
        if coords.shape[0] < cfg["MIN_SKEL_LEN_PX"]:
            reasons["short"] += 1
            continue

        topo = measure_topology(coords, W, allow_loops=cfg.get("ALLOW_LOOPS", False))
        if topo is None:
            reasons["loop"] += 1
            continue  # loop, not allowed

        gl   = topo["geo_len"]
        tort = topo["tortuosity"]
        n_ep = topo["n_endpoints"]
        n_br = topo["n_branch_nodes"]

        if not (cfg["MIN_SKEL_LEN_PX"] <= gl <= cfg["MAX_GEODESIC_LEN_PX"]):
            if gl < cfg["MIN_SKEL_LEN_PX"]: reasons["short"] += 1
            else: reasons["long"] += 1
            continue

        width = float(np.median(2.0 * dist[coords[:, 0], coords[:, 1]]))
        if width > cfg["MAX_WIDTH_PX"]:
            reasons["wide"] += 1
            continue

        if gl / (width + 1e-9) < cfg["MIN_LENGTH_WIDTH_RATIO"]:
            reasons["ratio"] += 1
            continue

        # N1: branch-node count filter
        if n_br > cfg["MAX_BRANCH_NODES"]:
            reasons["branches"] += 1
            continue

        # N2: tortuosity filter (only for open filaments with 2+ endpoints)
        if n_ep >= 2 and tort > cfg["MAX_TORTUOSITY"]:
            reasons["tortuous"] += 1
            continue

        # N3: endpoint count filter
        if n_ep > cfg["MAX_ENDPOINT_COUNT"]:
            reasons["endpoints"] += 1
            continue

        cy, cx = sp.centroid
        unet_vals = None
        if unet_prob is not None:
            unet_vals = np.asarray(unet_prob[coords[:, 0], coords[:, 1]], dtype=np.float32)
            unet_vals = unet_vals[np.isfinite(unet_vals)]
        accepted_labels.append(sp.label)
        area_est_px = float(gl * width)
        cache[sp.label] = {
            "geo_len":            gl,
            "tortuosity":         tort,
            "n_endpoints":        n_ep,
            "n_branch_nodes":     n_br,
            "width":              width,
            "length_width_ratio": gl / (width + 1e-9),
            "length_px_count":    float(coords.shape[0]),
            "cx": cx, "cy": cy,
            "area_px": area_est_px,
            "skeleton_area_px": float(sp.area),
            "bbox_min_y": float(sp.bbox[0]),
            "bbox_min_x": float(sp.bbox[1]),
            "bbox_max_y": float(sp.bbox[2]),
            "bbox_max_x": float(sp.bbox[3]),
            "orientation": float(sp.orientation),
            "unet_mean_probability": float(np.mean(unet_vals)) if unet_vals is not None and unet_vals.size else np.nan,
            "unet_max_probability": float(np.max(unet_vals)) if unet_vals is not None and unet_vals.size else np.nan,
            "detection_source": "saturn_classical",
        }

    total_rejected = sum(reasons.values())
    if total_rejected > 0:
        print(f"    measure_spermatids rejected {total_rejected} blobs:")
        for k, v in reasons.items():
            if v > 0: print(f"      {k}: {v}")

    if accepted_labels:
        clean_skel = np.isin(skel_lab, accepted_labels)
        final_label = measure.label(clean_skel).astype(np.int32)
    else:
        clean_skel = np.zeros_like(skel, dtype=bool)
        final_label = np.zeros_like(skel_lab, dtype=np.int32)

    # ------ Re-index using cached values (no second Dijkstra pass) ---------------------------------------------
    final_results = []
    for new_i, sp in enumerate(measure.regionprops(final_label), start=1):
        old_label = skel_lab[sp.coords[0, 0], sp.coords[0, 1]]
        if old_label not in cache:
            continue
        c = cache[old_label]
        final_results.append({
            "label":               new_i,
            "length_px_geodesic":  c["geo_len"],
            "length_px_count":     c["length_px_count"],
            "width_px":            c["width"],
            "length_width_ratio":  c["length_width_ratio"],
            "tortuosity":          c["tortuosity"],
            "n_endpoints":         c["n_endpoints"],
            "n_branch_nodes":      c["n_branch_nodes"],
            "centroid_x":          c["cx"],
            "centroid_y":          c["cy"],
            "area_px":             c["area_px"],
            "skeleton_area_px":    c["skeleton_area_px"],
            "bbox_min_y":          c["bbox_min_y"],
            "bbox_min_x":          c["bbox_min_x"],
            "bbox_max_y":          c["bbox_max_y"],
            "bbox_max_x":          c["bbox_max_x"],
            "orientation":         c["orientation"],
            "unet_mean_probability": c.get("unet_mean_probability", np.nan),
            "unet_max_probability": c.get("unet_max_probability", np.nan),
            "detection_source":    c.get("detection_source", "saturn_classical"),
        })

    rescue_reasons = {"short": 0, "loop": 0, "long": 0, "wide": 0, "ratio": 0, "branches": 0, "tortuous": 0, "endpoints": 0}
    rescue_reason_codes = {
        "short": 1,
        "loop": 2,
        "long": 3,
        "wide": 4,
        "ratio": 5,
        "branches": 6,
        "tortuous": 7,
        "endpoints": 8,
    }
    rescue_rejected_reason = np.zeros_like(skel_lab, dtype=np.uint8)
    rescue_accepted_label = np.zeros_like(skel_lab, dtype=np.int32)

    def reject_rescue(reason, coords):
        rescue_reasons[reason] += 1
        rescue_rejected_reason[coords[:, 0], coords[:, 1]] = rescue_reason_codes[reason]

    if (
        cfg.get("UNET_RESCUE_ENABLE", True)
        and unet_prob is not None
        and np.any(unet_prob > 0)
        and str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower() in {"hybrid", "unet_assisted"}
    ):
        roi = seg.get("roi_mask", np.ones_like(skel, dtype=bool)).astype(bool)
        exclusion = seg.get("exclusion_mask", np.zeros_like(skel, dtype=bool)).astype(bool)
        valid = roi & ~exclusion
        rescue_thr = float(cfg.get("UNET_RESCUE_THRESHOLD", cfg.get("UNET_SEED_THRESHOLD", 0.50)))
        candidate_thr = float(
            cfg.get("UNET_CANDIDATE_THRESHOLD", min(rescue_thr, 0.05))
        )
        candidate_thr = min(candidate_thr, rescue_thr)
        use_hysteresis = bool(cfg.get("UNET_RESCUE_HYSTERESIS_ENABLE", True))
        retain_morphology_warnings = bool(
            cfg.get("UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS", True)
        )
        exclude_px = max(0, int(cfg.get("UNET_RESCUE_EXCLUDE_DILATION_PX", 3)))
        min_component = max(1, int(cfg.get("UNET_RESCUE_MIN_COMPONENT_PX", cfg.get("MIN_OBJ_PX", 3))))
        rescue_min_skel_um = float(cfg.get("UNET_RESCUE_MIN_SKEL_LEN_UM", cfg.get("MIN_SKEL_LEN_UM", 0.0)))
        rescue_min_skel_px = max(
            1.0,
            rescue_min_skel_um / max(float(cfg.get("UM_PER_PX_XY", 1.0)), 1e-9),
        )
        short_rescue_min_prob = float(
            cfg.get("UNET_SHORT_RESCUE_MIN_MEAN_PROB", 0.85)
        )
        max_additions = max(0, int(cfg.get("UNET_RESCUE_MAX_ADDITIONS_PER_SLICE", 0)))
        split_retry = bool(cfg.get("UNET_RESCUE_SPLIT_RETRY_ENABLE", True))
        split_thresholds = [
            float(v) for v in cfg.get("UNET_RESCUE_SPLIT_THRESHOLDS", [0.70, 0.80, 0.90])
            if float(v) > rescue_thr
        ]
        centerline_salvage = bool(cfg.get("UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE", True))
        centerline_min_prob = float(cfg.get("UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB", 0.85))
        low_ratio_min_prob = float(
            cfg.get("UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB", 0.75)
        )
        low_ratio_min_length_px = float(
            cfg.get("UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM", 4.0)
        ) / max(float(cfg.get("UM_PER_PX_XY", 1.0)), 1e-9)
        instance_split = bool(cfg.get("UNET_INSTANCE_SPLIT_ENABLE", True))
        instance_seed_thr = float(cfg.get("UNET_INSTANCE_SEED_THRESHOLD", max(rescue_thr, 0.75)))
        peak_min_distance = max(1, int(cfg.get("UNET_INSTANCE_PEAK_MIN_DISTANCE_PX", 6)))
        watershed_compactness = float(cfg.get("UNET_INSTANCE_WATERSHED_COMPACTNESS", 0.001))

        occupied = clean_skel.copy()
        if exclude_px > 0 and np.any(occupied):
            occupied = morphology.binary_dilation(occupied, morphology.disk(exclude_px))
        # Split complete U-Net instances before checking whether Saturn already
        # represents them. Removing occupied pixels first can turn one complete
        # nucleus into several artificial short residual fragments.
        if use_hysteresis and candidate_thr < rescue_thr:
            rescue_mask = apply_hysteresis_threshold(
                unet_prob,
                candidate_thr,
                rescue_thr,
            ) & valid
        else:
            rescue_mask = (unet_prob >= rescue_thr) & valid
        rescue_mask = remove_objects_smaller_than(rescue_mask, min_component)
        rescue_lab = _split_unet_rescue_instances(
            unet_prob,
            rescue_mask,
            min_component,
            instance_seed_thr,
            peak_min_distance,
            watershed_compactness,
        ) if instance_split else measure.label(rescue_mask)
        rescue_candidates = []

        def evaluate_rescue_coords(coords, dist_map, sp, source):
            unet_vals = np.asarray(
                unet_prob[coords[:, 0], coords[:, 1]], dtype=np.float32
            )
            unet_vals = unet_vals[np.isfinite(unet_vals)]
            mean_unet_prob = float(np.mean(unet_vals)) if unet_vals.size else 0.0
            high_confidence_short = (
                coords.shape[0] >= 1
                and mean_unet_prob >= short_rescue_min_prob
            )
            if coords.shape[0] < rescue_min_skel_px and not high_confidence_short:
                return None, "short"

            topo = measure_topology(coords, W, allow_loops=cfg.get("ALLOW_LOOPS", False))
            if topo is None:
                return None, "loop"

            gl = topo["geo_len"]
            tort = topo["tortuosity"]
            n_ep = topo["n_endpoints"]
            n_br = topo["n_branch_nodes"]
            morphology_warnings = []
            if not (rescue_min_skel_px <= gl <= cfg["MAX_GEODESIC_LEN_PX"]):
                if gl < rescue_min_skel_px and not high_confidence_short:
                    return None, "short"
                if gl > cfg["MAX_GEODESIC_LEN_PX"]:
                    return None, "long"

            width = float(np.median(2.0 * dist_map[coords[:, 0], coords[:, 1]]))
            if width > cfg["MAX_WIDTH_PX"]:
                if retain_morphology_warnings:
                    morphology_warnings.append("wide")
                else:
                    return None, "wide"
            length_width_ratio = gl / (width + 1e-9)
            low_ratio = length_width_ratio < cfg["MIN_LENGTH_WIDTH_RATIO"]
            high_confidence_low_ratio = (
                low_ratio
                and gl >= low_ratio_min_length_px
                and mean_unet_prob >= low_ratio_min_prob
            )
            if low_ratio:
                if retain_morphology_warnings:
                    morphology_warnings.append("low_length_width_ratio")
                elif not high_confidence_low_ratio:
                    return None, "ratio"
            if n_br > cfg["MAX_BRANCH_NODES"]:
                return None, "branches"
            if n_ep >= 2 and tort > cfg["MAX_TORTUOSITY"]:
                if retain_morphology_warnings:
                    morphology_warnings.append("tortuous")
                else:
                    return None, "tortuous"
            if n_ep > cfg["MAX_ENDPOINT_COUNT"]:
                return None, "endpoints"

            cy, cx = sp.centroid
            resolved_source = source
            if gl < rescue_min_skel_px and high_confidence_short:
                resolved_source = "unet_rescued_short_high_confidence"
            elif low_ratio and not retain_morphology_warnings and high_confidence_low_ratio:
                resolved_source = "unet_rescued_low_ratio_high_confidence"
            elif morphology_warnings:
                resolved_source = "unet_rescued_morphology_warning"
            return {
                "coords": coords,
                "score": mean_unet_prob,
                "result": {
                    "length_px_geodesic": gl,
                    "length_px_count": float(coords.shape[0]),
                    "width_px": width,
                    "length_width_ratio": gl / (width + 1e-9),
                    "tortuosity": tort,
                    "n_endpoints": n_ep,
                    "n_branch_nodes": n_br,
                    "centroid_x": cx,
                    "centroid_y": cy,
                    "area_px": float(gl * width),
                    "skeleton_area_px": float(sp.area),
                    "bbox_min_y": float(sp.bbox[0]),
                    "bbox_min_x": float(sp.bbox[1]),
                    "bbox_max_y": float(sp.bbox[2]),
                    "bbox_max_x": float(sp.bbox[3]),
                    "orientation": float(sp.orientation),
                    "unet_mean_probability": mean_unet_prob if unet_vals.size else np.nan,
                    "unet_max_probability": float(np.max(unet_vals)) if unet_vals.size else np.nan,
                    "detection_source": resolved_source,
                    "unet_rescue_morphology_warning": bool(morphology_warnings),
                    "unet_rescue_morphology_warning_reasons": ";".join(
                        morphology_warnings
                    ),
                },
            }, None

        def evaluate_centerline_salvage(coords, dist_map, sp, source):
            if not centerline_salvage:
                return None
            unet_vals = np.asarray(unet_prob[coords[:, 0], coords[:, 1]], dtype=np.float32)
            unet_vals = unet_vals[np.isfinite(unet_vals)]
            if not unet_vals.size or float(np.mean(unet_vals)) < centerline_min_prob:
                return None
            centerline = extract_geodesic_centerline_coords(coords, W)
            if centerline.shape[0] >= coords.shape[0]:
                return None
            candidate, _reason = evaluate_rescue_coords(centerline, dist_map, sp, source)
            return candidate

        def split_and_retry_component(component_mask):
            split_hits = []
            for level in split_thresholds:
                core = component_mask & (unet_prob >= level) & valid & ~occupied
                core = remove_objects_smaller_than(core, min_component)
                if not np.any(core):
                    continue
                core_dist = _distance_transform_component(core)
                core_skel = skeletonize(core) & valid & ~occupied
                core_lab = measure.label(core_skel)
                for sub_sp in measure.regionprops(core_lab):
                    candidate, _reason = evaluate_rescue_coords(
                        sub_sp.coords,
                        core_dist,
                        sub_sp,
                        "unet_rescued_split",
                    )
                    if candidate is not None:
                        split_hits.append(candidate)
                if split_hits:
                    return split_hits
            return split_hits

        for sp in measure.regionprops(rescue_lab):
            instance_mask = rescue_lab == sp.label
            if np.any(instance_mask & occupied):
                # This U-Net instance is already represented by an accepted
                # Saturn centerline. It is neither a rescue nor a rejection.
                continue
            instance_dist = _distance_transform_component(instance_mask, bbox=sp.bbox)
            instance_skel = skeletonize(instance_mask) & valid
            instance_skel = remove_objects_smaller_than(
                instance_skel,
                max(1, int(math.floor(rescue_min_skel_px * 0.70))),
            )
            if not np.any(instance_skel):
                reject_rescue("short", sp.coords)
                continue
            instance_skel_lab = measure.label(instance_skel)
            accepted_any = False
            rejected_reasons = []
            for skel_sp in measure.regionprops(instance_skel_lab):
                coords = skel_sp.coords
                candidate, reason = evaluate_rescue_coords(coords, instance_dist, skel_sp, "unet_rescued")
                if candidate is not None:
                    rescue_candidates.append(candidate)
                    accepted_any = True
                    continue

                if reason in {"loop", "long", "branches", "tortuous", "endpoints"}:
                    candidate = evaluate_centerline_salvage(coords, instance_dist, skel_sp, "unet_rescued")
                    if candidate is not None:
                        rescue_candidates.append(candidate)
                        accepted_any = True
                        continue

                if split_retry and reason in {"long", "branches", "loop", "tortuous", "endpoints"}:
                    # Retry against the complete probability-supported instance.
                    # Thresholding only the one-pixel skeleton cannot separate a
                    # merged U-Net region into distinct nuclei.
                    split_hits = split_and_retry_component(instance_mask)
                    if split_hits:
                        rescue_candidates.extend(split_hits)
                        accepted_any = True
                        continue
                rejected_reasons.append((reason, coords))
            if accepted_any:
                continue
            if rejected_reasons:
                reason, coords = rejected_reasons[0]
                reject_rescue(reason, coords)
            else:
                reject_rescue("short", sp.coords)

        rescue_candidates.sort(key=lambda item: item["score"], reverse=True)
        if max_additions > 0:
            rescue_candidates = rescue_candidates[:max_additions]
        next_label = int(final_label.max()) + 1
        for item in rescue_candidates:
            coords = item["coords"]
            final_label[coords[:, 0], coords[:, 1]] = next_label
            rescue_accepted_label[coords[:, 0], coords[:, 1]] = next_label
            item["result"]["label"] = next_label
            final_results.append(item["result"])
            next_label += 1

        if rescue_candidates or sum(rescue_reasons.values()) > 0:
            print(f"    U-Net rescue accepted {len(rescue_candidates)} blobs at probability >= {rescue_thr:.3f}")
            rejected = sum(rescue_reasons.values())
            if rejected:
                print(f"    U-Net rescue rejected {rejected} blobs:")
                for k, v in rescue_reasons.items():
                    if v > 0:
                        print(f"      {k}: {v}")

    return {
        "skel_label": final_label,
        "results": final_results,
        "unet_rescue_accepted_label": rescue_accepted_label,
        "unet_rescue_rejected_reason": rescue_rejected_reason,
        "unet_rescue_reason_codes": rescue_reason_codes,
        "unet_rescue_rejected_counts": dict(rescue_reasons),
        "unet_rescue_candidate_threshold": float(
            cfg.get("UNET_CANDIDATE_THRESHOLD", 0.0)
        ),
        "unet_rescue_seed_threshold": float(
            cfg.get("UNET_RESCUE_THRESHOLD", 0.0)
        ),
    }


# =============================================================================
# OVERLAY  (vectorized LUT)
# =============================================================================

_OVERLAY_DISPLAY_DILATION_SIZE = 3


def make_overlay(img_raw, skel_label):
    """
    Generates a colour-coded skeleton overlay on the grayscale raw image.

    Each detected spermatid is assigned a unique hue from the ``gist_rainbow``
    colourmap, dilated only for display visibility, and composited onto the
    contrast-stretched raw image.  Background pixels (label == 0) retain the
    original grayscale intensity.

    Implementation notes
    --------------------
    - Vectorised LUT (Look-Up Table) approach avoids per-label loop for speed.
    - ``grey_dilation`` makes thin single-pixel skeletons visible at any zoom.
    - Colour assignment is deterministic: same label ordering -> same colour.

    Args:
        img_raw (np.ndarray): Raw microscopy image (any dtype).
        skel_label (np.ndarray[int]): Integer-labelled skeleton array
            (0 = background, 1..N = individual spermatids).

    Returns:
        np.ndarray: uint8 RGB image, shape ``(H, W, 3)``, ready for ``plt.imshow``
        or saving with :func:`_imwrite`.
    """
    base = normalize_display(img_raw)
    n    = int(skel_label.max())
    if n <= 0:
        # No detections: return grayscale image as RGB
        return (np.stack([base]*3, -1) * 255).astype(np.uint8)
    # Assign one colour per label; prepend black for background (index 0)
    cols    = plt.cm.gist_rainbow(np.linspace(0, 1, n))[:, :3]
    dilated = grey_dilation(skel_label.astype(np.int32), size=_OVERLAY_DISPLAY_DILATION_SIZE)
    lut     = np.vstack([[0., 0., 0.], cols[:n]])
    rgb     = lut[dilated]
    # Restore original grayscale for background pixels
    m0      = dilated == 0
    rgb[m0, 0] = base[m0]
    rgb[m0, 1] = base[m0]
    rgb[m0, 2] = base[m0]
    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)


def make_unet_rescue_review_overlay(img_raw, skel_label, results, rescue_rejected_reason=None):
    """
    Audit overlay for v5.7 U-Net rescue review.

    Colors:
    - green: accepted Saturn classical detections
    - cyan: accepted U-Net rescued detections
    - red: U-Net candidate rejected as long/branched/loop/tortuous
    - orange: U-Net candidate rejected as wide/low ratio
    - magenta: U-Net candidate rejected as short fragment

    The dilation below is display-only and is never used for object counting,
    skeleton length, width, or 3D tracking calculations.
    """
    base = normalize_display(img_raw)
    rgb = np.stack([base, base, base], axis=-1)

    label_source = {}
    for row in results or []:
        try:
            label_source[int(row.get("label", 0))] = row.get("detection_source", "saturn_classical")
        except Exception:
            continue

    if rescue_rejected_reason is not None:
        rejected = np.asarray(rescue_rejected_reason, dtype=np.uint8)
        footprint = np.ones(
            (_OVERLAY_DISPLAY_DILATION_SIZE,) * 2,
            dtype=bool,
        )
        short = morphology.binary_dilation(
            rejected == 1,
            footprint,
        )
        severe = morphology.binary_dilation(
            np.isin(rejected, [2, 3, 6, 7, 8]),
            footprint,
        )
        shape = morphology.binary_dilation(
            np.isin(rejected, [4, 5]),
            footprint,
        )
        rgb[short] = (1.0, 0.0, 0.9)
        rgb[shape] = (1.0, 0.55, 0.0)
        rgb[severe] = (1.0, 0.0, 0.0)

    if skel_label is not None and int(np.max(skel_label)) > 0:
        unet_labels = [
            label
            for label, source in label_source.items()
            if label > 0 and str(source).startswith("unet_rescued")
        ]
        classical_labels = [
            label
            for label, source in label_source.items()
            if label > 0 and not str(source).startswith("unet_rescued")
        ]
        classical_mask = morphology.binary_dilation(
            np.isin(skel_label, classical_labels),
            footprint,
        )
        unet_mask = morphology.binary_dilation(
            np.isin(skel_label, unet_labels),
            footprint,
        )
        rgb[classical_mask] = (0.0, 1.0, 0.25)
        rgb[unet_mask] = (0.0, 0.9, 1.0)

    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)


def make_quality_overlay(img_raw, skel_label, slice_tracks, track_quality_map):
    """
    Draw an audit-coded overlay for a single Z slice.

    Colors:
    - green: accepted estimated nucleus without warnings
    - yellow/orange: accepted estimated nucleus with a morphology warning
    - red: hard-failed track
    - gray: detection has no track/audit mapping
    """
    base = normalize_display(img_raw)
    rgb = np.stack([base, base, base], axis=-1)
    if skel_label is None or int(np.max(skel_label)) <= 0:
        return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)

    label_quality = {}
    if slice_tracks is not None and not slice_tracks.empty:
        for _, row in slice_tracks.iterrows():
            try:
                label = int(row["sperm_id"])
                track_id = int(row["track_id"])
            except Exception:
                continue
            label_quality[label] = track_quality_map.get(track_id)

    dilated = grey_dilation(skel_label.astype(np.int32), size=_OVERLAY_DISPLAY_DILATION_SIZE)
    for label in np.unique(dilated):
        label = int(label)
        if label == 0:
            continue
        q = label_quality.get(label)
        if q == "candidate":
            color = np.array([0.0, 0.85, 0.25])      # candidate/pass
        elif q == "warning":
            color = np.array([1.0, 0.75, 0.05])      # warning-only
        elif q == "hard_fail":
            color = np.array([1.0, 0.18, 0.05])      # hard fail
        elif q is True:
            color = np.array([0.0, 0.85, 0.25])      # backward-compatible strict pass
        elif q is False:
            color = np.array([1.0, 0.18, 0.05])      # backward-compatible strict fail
        else:
            color = np.array([0.65, 0.65, 0.65])     # untracked/unknown
        mask = dilated == label
        rgb[mask] = 0.25 * rgb[mask] + 0.75 * color

    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)


def make_analysis_overlay(img_raw, skel_label, slice_tracks, included_track_ids):
    """Draw only observations belonging to the primary 3D analysis population."""
    base = normalize_display(img_raw)
    rgb = np.stack([base, base, base], axis=-1)
    if skel_label is None or int(np.max(skel_label)) <= 0:
        return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)

    included_labels = []
    if slice_tracks is not None and not slice_tracks.empty:
        for _, row in slice_tracks.iterrows():
            try:
                label = int(row["sperm_id"])
                track_id = int(row["track_id"])
            except Exception:
                continue
            if track_id in included_track_ids:
                included_labels.append(label)

    footprint = np.ones(
        (_OVERLAY_DISPLAY_DILATION_SIZE,) * 2,
        dtype=bool,
    )
    display_mask = morphology.binary_dilation(
        np.isin(skel_label, included_labels),
        footprint,
    )
    color = np.array([0.0, 0.85, 0.25])
    rgb[display_mask] = 0.20 * rgb[display_mask] + 0.80 * color
    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)


_QUALITY_OVERLAY_LEGEND = {
    "candidate": ("#00d940", "Included estimated nucleus"),
    "warning": ("#ffbf0d", "Included; morphology warning"),
    "hard_fail": ("#ff2e0d", "Excluded technical failure"),
    "unmapped": ("#a6a6a6", "Detected candidate without track assignment"),
}


def quality_overlay_legend_handles(active_statuses=None):
    """Return a data-dependent legend for post-tracking QC overlays."""
    statuses = (
        list(_QUALITY_OVERLAY_LEGEND)
        if active_statuses is None
        else [status for status in _QUALITY_OVERLAY_LEGEND if status in active_statuses]
    )
    return [
        Patch(
            facecolor=_QUALITY_OVERLAY_LEGEND[status][0],
            edgecolor="none",
            label=_QUALITY_OVERLAY_LEGEND[status][1],
        )
        for status in statuses
    ]


def quality_overlay_status_counts(skel_label, slice_tracks, track_quality_map):
    """Count displayed skeleton labels in each post-tracking QC category."""
    counts = {status: 0 for status in _QUALITY_OVERLAY_LEGEND}
    label_quality = {}
    if slice_tracks is not None and not slice_tracks.empty:
        for _, row in slice_tracks.iterrows():
            try:
                label = int(row["sperm_id"])
                track_id = int(row["track_id"])
            except Exception:
                continue
            label_quality[label] = track_quality_map.get(track_id)

    labels = np.unique(skel_label).astype(int) if skel_label is not None else np.array([], dtype=int)
    for label in labels:
        if label == 0:
            continue
        quality = label_quality.get(label)
        if quality in {"candidate", "warning", "hard_fail"}:
            counts[quality] += 1
        elif quality is True:
            counts["candidate"] += 1
        elif quality is False:
            counts["hard_fail"] += 1
        else:
            counts["unmapped"] += 1
    return counts


def export_quality_overlays(out_dir, slice_cache, df_tracked, track_summary):
    """
    Save per-slice and global audit-coded overlays after 3D tracking is audited.
    """
    if not slice_cache or df_tracked is None or df_tracked.empty or track_summary is None or track_summary.empty:
        return None
    if "is_quality_track" not in track_summary.columns or "track_id" not in track_summary.columns:
        return None

    quality_dir = os.path.join(out_dir, "quality_overlays")
    ensure_dir(quality_dir)

    track_quality_map = {}
    for _, row in track_summary.iterrows():
        if pd.isna(row.get("track_id")):
            continue
        tid = int(row["track_id"])
        if "technical_valid" in track_summary.columns:
            if bool(row.get("technical_valid", False)):
                track_quality_map[tid] = "warning" if bool(row.get("morphology_warning", False)) else "candidate"
            else:
                track_quality_map[tid] = "hard_fail"
        elif "is_biological_candidate" in track_summary.columns:
            if bool(row.get("is_biological_candidate")):
                track_quality_map[tid] = "warning" if bool(row.get("has_warning_only", False)) else "candidate"
            else:
                track_quality_map[tid] = "hard_fail"
        else:
            track_quality_map[tid] = bool(row.get("is_quality_track", False))

    max_proj_raw = None
    max_proj_quality = None
    status_rows = []
    for z_idx in sorted(slice_cache):
        item = slice_cache[z_idx]
        img = item["image"]
        skel_label = item["skel_label"]
        slice_tracks = df_tracked[df_tracked["z_slice"].astype(int) == int(z_idx)]
        status_rows.append({
            "z_slice": int(z_idx),
            **quality_overlay_status_counts(skel_label, slice_tracks, track_quality_map),
        })
        quality_rgb = make_quality_overlay(img, skel_label, slice_tracks, track_quality_map)

        raw_rgb = (normalize_display(img) * 255).astype(np.uint8)
        if raw_rgb.ndim == 2:
            raw_rgb = np.stack([raw_rgb] * 3, axis=-1)
        panel = np.hstack([raw_rgb, quality_rgb])
        _imwrite(os.path.join(quality_dir, f"z{int(z_idx):02d}_quality_panel.png"), panel)

        if max_proj_raw is None:
            max_proj_raw = img.copy().astype(np.float32)
            max_proj_quality = quality_rgb.copy().astype(np.float32)
        else:
            max_proj_raw = np.maximum(max_proj_raw, img.astype(np.float32))
            max_proj_quality = np.maximum(max_proj_quality, quality_rgb.astype(np.float32))

    if max_proj_raw is None:
        return None
    raw_p = (normalize_display(max_proj_raw.astype(np.uint16)) * 255).astype(np.uint8)
    if raw_p.ndim == 2:
        raw_p = np.stack([raw_p] * 3, axis=-1)
    quality_p = np.clip(max_proj_quality, 0, 255).astype(np.uint8)
    global_panel = np.hstack([raw_p, quality_p])
    out_path = os.path.join(out_dir, "quality_global_z_projection.png")
    _imwrite(out_path, global_panel)

    status_df = pd.DataFrame(status_rows)
    status_df.to_csv(os.path.join(quality_dir, "quality_overlay_counts.csv"), index=False)
    active_statuses = {
        status
        for status in _QUALITY_OVERLAY_LEGEND
        if status in status_df.columns and int(status_df[status].sum()) > 0
    }
    legend_fig, legend_ax = plt.subplots(figsize=(9, 1.2))
    legend_ax.axis("off")
    legend_ax.legend(
        handles=quality_overlay_legend_handles(active_statuses),
        loc="center",
        ncol=2,
        fontsize=10,
        frameon=True,
        title="Track-QC overlay colors",
    )
    legend_ax.text(
        0.5,
        0.02,
        "Green + amber are included. Uncolored image structures were not detected. "
        "Overlay thickness is display-only.",
        transform=legend_ax.transAxes,
        ha="center",
        va="bottom",
        fontsize=9,
    )
    legend_fig.savefig(
        os.path.join(quality_dir, "quality_overlay_legend.png"),
        dpi=180,
        bbox_inches="tight",
    )
    plt.close(legend_fig)
    return out_path


def export_analysis_overlays(out_dir, slice_cache, df_tracked, track_summary):
    """Save clean overlays containing only included estimated-nucleus observations."""
    if (
        not slice_cache
        or df_tracked is None
        or df_tracked.empty
        or track_summary is None
        or track_summary.empty
        or "track_id" not in track_summary.columns
    ):
        return None

    primary = _technical_valid_track_population(track_summary)
    included_track_ids = set(
        pd.to_numeric(primary["track_id"], errors="coerce")
        .dropna()
        .astype(int)
        .tolist()
    )
    analysis_dir = os.path.join(out_dir, "analysis_overlays")
    ensure_dir(analysis_dir)

    max_proj_raw = None
    max_proj_analysis = None
    count_rows = []
    tracked_z = pd.to_numeric(df_tracked["z_slice"], errors="coerce")
    tracked_ids = pd.to_numeric(df_tracked["track_id"], errors="coerce")
    for z_idx in sorted(slice_cache):
        item = slice_cache[z_idx]
        img = item["image"]
        skel_label = item["skel_label"]
        slice_tracks = df_tracked[tracked_z.eq(int(z_idx))].copy()
        included_count = int(
            tracked_ids.loc[slice_tracks.index].isin(included_track_ids).sum()
        )
        count_rows.append(
            {
                "z_slice": int(z_idx),
                "included_nucleus_observations": included_count,
            }
        )
        analysis_rgb = make_analysis_overlay(
            img,
            skel_label,
            slice_tracks,
            included_track_ids,
        )
        raw_rgb = (normalize_display(img) * 255).astype(np.uint8)
        if raw_rgb.ndim == 2:
            raw_rgb = np.stack([raw_rgb] * 3, axis=-1)
        panel = np.hstack([raw_rgb, analysis_rgb])
        _imwrite(
            os.path.join(analysis_dir, f"z{int(z_idx):02d}_analysis_panel.png"),
            panel,
        )

        if max_proj_raw is None:
            max_proj_raw = img.copy().astype(np.float32)
            max_proj_analysis = analysis_rgb.copy().astype(np.float32)
        else:
            max_proj_raw = np.maximum(max_proj_raw, img.astype(np.float32))
            max_proj_analysis = np.maximum(
                max_proj_analysis,
                analysis_rgb.astype(np.float32),
            )

    pd.DataFrame(count_rows).to_csv(
        os.path.join(analysis_dir, "analysis_overlay_counts.csv"),
        index=False,
    )
    if max_proj_raw is None:
        return None
    raw_projection = (
        normalize_display(max_proj_raw.astype(np.uint16)) * 255
    ).astype(np.uint8)
    if raw_projection.ndim == 2:
        raw_projection = np.stack([raw_projection] * 3, axis=-1)
    analysis_projection = np.clip(max_proj_analysis, 0, 255).astype(np.uint8)
    output_path = os.path.join(out_dir, "analysis_global_z_projection.png")
    _imwrite(
        output_path,
        np.hstack([raw_projection, analysis_projection]),
    )
    return output_path


# =============================================================================
# DISPLAY / SAVE
# =============================================================================

def _safe_show():
    """
    Calls ``plt.show(block=True)`` with a try/except guard.

    On headless systems (Linux CI servers, SSH sessions without X11 forwarding,
    or Windows without a display) matplotlib will raise a ``RuntimeError`` or
    ``_tkinter.TclError`` when trying to open an interactive window.  This
    wrapper silently swallows that error so the pipeline can continue and saves
    the image to disk as a fallback.
    """
    try:
        plt.show(block=True)
    except Exception as e:
        print(f"[WARNING] plt.show() failed: {e}")


def show_single_preview(img_raw, seg, overlay_rgb, results, z_idx, cfg):
    """
    Renders and saves an interactive preview figure for a single analysed Z-slice.

    Layout operates in two modes controlled by ``SHOW_DEBUG_PREVIEW``:

    - **Standard mode** (2 panels): Raw image | Spermatid overlay + length histogram.
    - **Debug mode** (8 panels): Adds intermediate stage images - CLAHE, ridge filter,
      hysteresis mask, cleaned mask, and pruned skeleton - for visual pipeline QC.

    After saving to ``output_dir/preview.png``, the function attempts to open the
    image with the OS default viewer (``os.startfile`` on Windows) as an immediate
    visual check, falling back silently if that fails.

    Args:
        img_raw (np.ndarray): Raw microscopy image.
        seg (dict): Segmentation dictionary from :func:`segment_slice` containing
            intermediate images (``img_eq``, ``ridge``, ``mask_hyst``, etc.).
        overlay_rgb (np.ndarray): Colour overlay from :func:`make_overlay`.
        results (list[dict]): Per-spermatid measurement dictionaries.
        z_idx (int): Z-slice index used for figure titles.
        cfg (dict): Pipeline configuration; reads ``SHOW_DEBUG_PREVIEW``,
            ``OUTPUT_DIR``, and ``UM_PER_PX_XY``.
    """
    um         = cfg["UM_PER_PX_XY"]
    lengths_um = [r["length_px_geodesic"] * um for r in results]

    nrows, ncols = (2, 4) if cfg["SHOW_DEBUG_PREVIEW"] else (1, 3)
    fig, axes = plt.subplots(nrows, ncols, figsize=(18, 5*nrows))

    def _ax(r, c):
        return axes[r, c] if nrows > 1 else axes[c]

    _ax(0,0).imshow(normalize_display(img_raw), cmap="gray")
    _ax(0,0).set_title(f"Original Z={z_idx:02d}"); _ax(0,0).axis("off")

    if cfg["SHOW_DEBUG_PREVIEW"]:
        _ax(0,1).imshow(seg["img_eq"], cmap="gray")
        _ax(0,1).set_title("CLAHE"); _ax(0,1).axis("off")
        _ax(0,2).imshow(seg["ridge"], cmap="gray")
        _ax(0,2).set_title("Ridge"); _ax(0,2).axis("off")
        _ax(0,3).imshow(seg["mask_hyst"], cmap="gray")
        _ax(0,3).set_title("Hysteresis"); _ax(0,3).axis("off")
        _ax(1,0).imshow(seg["mask_clean"], cmap="gray")
        _ax(1,0).set_title("Clean mask"); _ax(1,0).axis("off")
        _ax(1,1).imshow(seg["skel_pruned"], cmap="gray")
        _ax(1,1).set_title("Skeleton (pruned)"); _ax(1,1).axis("off")
        ov_ax   = _ax(1,2)
        hist_ax = _ax(1,3)
    else:
        ov_ax   = _ax(0,1)
        hist_ax = _ax(0,2)

    ov_ax.imshow(overlay_rgb)
    ov_ax.set_title(f"Pre-tracking 2D candidates (N={len(results)})")
    ov_ax.text(
        0.5,
        -0.025,
        "Colors separate candidate IDs; they do not indicate source or acceptance.",
        transform=ov_ax.transAxes,
        ha="center",
        va="top",
        fontsize=8,
        color="#333333",
    )
    ov_ax.axis("off")
    for r in results:
        ov_ax.text(r["centroid_x"], r["centroid_y"],
                   f"{r['length_px_geodesic']*um:.1f}",
                   color="white", fontsize=5, ha="center", va="center")

    if lengths_um:
        hist_ax.hist(lengths_um, bins=20, edgecolor="white")
        hist_ax.axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} um")
        hist_ax.set_xlabel("Geodesic length (um)"); hist_ax.legend(fontsize=8)
        hist_ax.set_title("Length distribution")
    else:
        hist_ax.text(0.5, 0.5, "No detections", ha="center", va="center")
        hist_ax.axis("off")

    plt.tight_layout()

    # Persist preview image so it can be opened by OS or reloaded by GUI
    preview_path = os.path.join(cfg["OUTPUT_DIR"], "preview.png")
    ensure_dir(cfg["OUTPUT_DIR"])
    plt.savefig(preview_path, dpi=120, bbox_inches="tight")
    print(f"  Preview saved to: {preview_path}")

    # Try interactive display
    _safe_show()
    plt.close()

    # Open with default viewer as fallback
    try:
        os.startfile(preview_path)
    except Exception:
        pass


def save_detail_figure(img_raw, overlay_rgb, results, out_path, z_idx, um):
    """
    Saves a three-panel publication-quality figure for a single Z-slice.

    Panels:
    1. **Original** - contrast-stretched raw image.
    2. **Spermatid overlay** - colour-coded detections with length labels.
    3. **Length distribution** - histogram of geodesic lengths in um with median line.

    This figure is saved as a PNG per Z-slice and is useful for visual quality
    control and inclusion in lab reports.

    Args:
        img_raw (np.ndarray): Raw microscopy image.
        overlay_rgb (np.ndarray): Colour overlay from :func:`make_overlay`.
        results (list[dict]): Per-spermatid measurement dictionaries.
        out_path (str): Destination PNG file path.
        z_idx (int): Z-slice index for figure titles.
        um (float): Microns-per-pixel scale factor (``UM_PER_PX_XY``).
    """
    lengths_um = [r["length_px_geodesic"] * um for r in results]
    fig, axes  = plt.subplots(1, 3, figsize=(18, 6))

    axes[0].imshow(normalize_display(img_raw), cmap="gray")
    axes[0].set_title(f"Z={z_idx:02d} - Original"); axes[0].axis("off")

    axes[1].imshow(overlay_rgb)
    axes[1].set_title(f"Z={z_idx:02d} - Pre-tracking 2D candidates (N={len(results)})")
    axes[1].text(
        0.5,
        -0.025,
        "Colors separate candidate IDs; they do not indicate source or acceptance.",
        transform=axes[1].transAxes,
        ha="center",
        va="top",
        fontsize=8,
        color="#333333",
    )
    axes[1].axis("off")
    for r in results:
        axes[1].text(r["centroid_x"], r["centroid_y"],
                     f"{r['length_px_geodesic']*um:.1f}",
                     color="white", fontsize=4, ha="center", va="center",
                     bbox=dict(boxstyle="round,pad=0.1", fc="black", alpha=0.4, lw=0))

    if lengths_um:
        axes[2].hist(lengths_um, bins=20, edgecolor="white")
        axes[2].axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} um")
        axes[2].set_xlabel("Geodesic length (um)"); axes[2].set_ylabel("Count")
        axes[2].set_title(f"Z={z_idx:02d} - Length distribution")
        axes[2].legend(fontsize=9)
    else:
        axes[2].text(0.5, 0.5, "No spermatids detected",
                     transform=axes[2].transAxes, ha="center", va="center")
        axes[2].axis("off")

    plt.tight_layout()
    plt.savefig(out_path, dpi=120, bbox_inches="tight")
    plt.close()


# =============================================================================
# CSV
# =============================================================================

_VERSION = "v5.7-unet-ready"


def rows_from_results(results, z_idx, um):
    """
    Converts per-spermatid measurement dictionaries into flat CSV row dictionaries.

    Applies pixel-to-micron scaling (``UM_PER_PX_XY``) for all linear dimensions
    and rounds floating-point values to 3 decimal places for clean spreadsheet output.

    Columns emitted per detection
    ------------------------------
    - ``pipeline_version``     - Code version tag for traceability.
    - ``z_slice``              - Z-plane index of this detection.
    - ``sperm_id``             - Skeleton label ID within this slice.
    - ``length_px_geodesic``   - Geodesic skeleton length in pixels.
    - ``length_um_geodesic``   - Geodesic length in micrometres.
    - ``length_px_count``      - Pixel-count skeleton length (alternative measure).
    - ``length_um_count``      - Pixel-count length in micrometres.
    - ``width_px`` / ``width_um`` - Mean width across the skeleton.
    - ``length_width_ratio``   - Elongation ratio (key morphological filter metric).
    - ``tortuosity``           - Curvature index: geodesic / Euclidean tip-to-tip.
    - ``n_endpoints``          - Number of skeleton endpoints (expect 2 for clean cells).
    - ``n_branch_nodes``       - Number of branch points (3+ neighbours); should be low.
    - ``centroid_x`` / ``centroid_y`` - Pixel centroid for overlay annotation.
    - ``area_px``              - Slender-object area estimate in pixels^2
      (geodesic length * median width), used for area-derived metrics.
    - ``skeleton_area_px``     - Raw skeleton pixel count, kept for audit/debug.

    Args:
        results (list[dict]): Output of :func:`measure_spermatids`.
        z_idx (int): Z-slice index.
        um (float): Microns-per-pixel (``UM_PER_PX_XY``).

    Returns:
        list[dict]: One flat dictionary per detected spermatid.
    """
    out_rows = []
    for i, r in enumerate(results, start=1):
        historical_area = round(float(r.get("area_px", 0.0)), 1)
        estimated_slender_area = round(float(r["length_px_geodesic"]) * float(r["width_px"]), 1)
        instance_mask_area = round(r.get("instance_mask_area_px", np.nan), 1) if np.isfinite(r.get("instance_mask_area_px", np.nan)) else np.nan
        detection_source = r.get("detection_source", "saturn_classical")

        if detection_source == "unet_primary" and np.isfinite(instance_mask_area) and instance_mask_area > 0:
            final_area = instance_mask_area
        else:
            final_area = historical_area

        out_rows.append({
            "pipeline_version":    _VERSION,
            "z_slice":             z_idx,
            "sperm_id":            int(r.get("label", i)),
            "source_instance_key": r.get("source_instance_key", ""),
            "length_px_geodesic":  round(r["length_px_geodesic"], 3),
            "length_um_geodesic":  round(r["length_px_geodesic"] * um, 3),
            "length_px_count":     round(r["length_px_count"], 1),
            "length_um_count":     round(r["length_px_count"]  * um, 3),
            "width_px":            round(r["width_px"], 2),
            "width_um":            round(r["width_px"]          * um, 3),
            "length_width_ratio":  round(r["length_width_ratio"], 3),
            "tortuosity":          round(r["tortuosity"], 3),
            "n_endpoints":         r["n_endpoints"],
            "n_branch_nodes":      r["n_branch_nodes"],
            "centroid_x":          round(r["centroid_x"], 1),
            "centroid_y":          round(r["centroid_y"], 1),
            "area_px":             final_area,
            "estimated_slender_area_px": estimated_slender_area,
            "skeleton_area_px":    round(r.get("skeleton_area_px", 0.0), 1),
            "instance_mask_area_px": instance_mask_area,
            "bbox_min_y":          r.get("bbox_min_y"),
            "bbox_min_x":          r.get("bbox_min_x"),
            "bbox_max_y":          r.get("bbox_max_y"),
            "bbox_max_x":          r.get("bbox_max_x"),
            "orientation":         round(r.get("orientation", 0.0), 3),
            "detection_source":    detection_source,
            "unet_mean_probability": round(float(r.get("unet_mean_probability", np.nan)), 4) if np.isfinite(r.get("unet_mean_probability", np.nan)) else np.nan,
            "unet_max_probability":  round(float(r.get("unet_max_probability", np.nan)), 4) if np.isfinite(r.get("unet_max_probability", np.nan)) else np.nan,
            "unet_rescue_morphology_warning": bool(
                r.get("unet_rescue_morphology_warning", False)
            ),
            "unet_rescue_morphology_warning_reasons": r.get(
                "unet_rescue_morphology_warning_reasons",
                "",
            ),
            "morphology_warning": bool(r.get("morphology_warning", False)),
            "morphology_warning_reasons": r.get(
                "morphology_warning_reasons", ""
            ),
            "technical_failure": bool(r.get("technical_failure", False)),
            "technical_failure_reason": r.get("technical_failure_reason", ""),
            "parent_hysteresis_component_id": int(
                r.get("parent_hysteresis_component_id", 0)
            ),
        })
    return out_rows


# =============================================================================
# TRACKING
# =============================================================================

def _estimated_tracking_extension_length_um(
    prev_state, candidate_x, candidate_y, candidate_z, candidate_length, cfg
):
    """Estimate the physical length created by extending one track."""
    um_xy = float(cfg["UM_PER_PX_XY"])
    um_z = float(cfg["UM_PER_SLICE_Z"])
    first_x = float(prev_state.get("first_x", prev_state["last_x"]))
    first_y = float(prev_state.get("first_y", prev_state["last_y"]))
    first_z = int(prev_state.get("first_z", prev_state["last_z"]))
    centroid_span = math.hypot(
        float(candidate_x) - first_x,
        float(candidate_y) - first_y,
    ) * um_xy
    z_span = abs(int(candidate_z) - first_z) * um_z
    finite_lengths = [
        float(value)
        for value in (
            prev_state.get("max_length_2d"),
            prev_state.get("last_length"),
            candidate_length,
            centroid_span,
        )
        if value is not None and np.isfinite(value)
    ]
    lateral_length = max(finite_lengths) if finite_lengths else centroid_span
    return math.hypot(lateral_length, z_span)


def check_extension_consistency(prev_state, candidate_detection, cfg, overlap_exists=False):
    """
    Check if extending a track with this detection would be biologically consistent.
    Implements 'Continue Unless Implausible' logic for overlapping footprints.

    Stage 2b: All thresholds are now driven by CONFIG for hyperparameter tuning.
    """
    um_xy = cfg["UM_PER_PX_XY"]

    # Read tunable Stage 2 parameters from CONFIG
    stab_thresh = cfg.get("OVERLAP_STABILITY_THRESHOLD", 0.08)
    ori_deg     = cfg.get("OVERLAP_ORIENTATION_DEG", 15.0)
    ovl_mult    = cfg.get("OVERLAP_MULTIPLIER", 1.35)
    min_stable  = cfg.get("OVERLAP_MIN_STABLE_COUNT", 1)

    # Extract previous track state
    prev_x = prev_state["last_x"]
    prev_y = prev_state["last_y"]
    prev_width = prev_state.get("last_width")
    prev_length = prev_state.get("last_length")
    prev_area = prev_state.get("last_area")
    prev_ori = prev_state.get("last_orientation")

    # Extract candidate detection features
    cand_x = candidate_detection["centroid_x"]
    cand_y = candidate_detection["centroid_y"]
    cand_width = candidate_detection.get("width_um")
    cand_length = candidate_detection.get("length_um_geodesic")
    cand_area = candidate_detection.get("area_px")
    cand_ori = candidate_detection.get("orientation")
    cand_z = candidate_detection.get("z_slice", prev_state["last_z"])

    max_joined_length = float(
        cfg.get("TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM", 15.0)
    )
    estimated_joined_length = _estimated_tracking_extension_length_um(
        prev_state,
        cand_x,
        cand_y,
        cand_z,
        cand_length,
        cfg,
    )
    if estimated_joined_length > max_joined_length:
        return False, (
            f"technical_joined_length={estimated_joined_length:.2f}um"
            f">{max_joined_length:.2f}um"
        )

    # Logic:
    # If overlap_exists, we allow the track to continue IF enough primary metrics are stable.
    # This prevents 'monster merges' where a track jumps onto a totally different cell.
    if overlap_exists:
        stable_count = 0

        # 1. Width stability
        if prev_width and cand_width:
            if (abs(cand_width - prev_width) / max(prev_width, 1e-9)) < stab_thresh:
                stable_count += 1

        # 2. Area stability
        if prev_area and cand_area:
            if (abs(cand_area - prev_area) / max(prev_area, 1e-9)) < stab_thresh:
                stable_count += 1

        # 3. Orientation stability
        if prev_ori is not None and cand_ori is not None:
            diff_rad = abs(cand_ori - prev_ori)
            if diff_rad > math.pi / 2:
                diff_rad = math.pi - diff_rad
            if diff_rad < (ori_deg * math.pi / 180):
                stable_count += 1

        # 4. Length stability
        if prev_length and cand_length:
            if (abs(cand_length - prev_length) / max(prev_length, 1e-9)) < stab_thresh:
                stable_count += 1

        # Require minimum stable metrics to continue an overlapping track
        if stable_count < min_stable:
            return False, f"overlap_but_{stable_count}_stable"

        # Even with stable metrics, still apply capped multiplier for fallback checks
        multiplier = ovl_mult
    else:
        multiplier = 1.0

    # 1. Check centroid jump
    dx = cand_x - prev_x
    dy = cand_y - prev_y
    centroid_jump_um = math.sqrt(dx*dx + dy*dy) * um_xy

    if not overlap_exists and centroid_jump_um > cfg["CONSERVATIVE_MAX_CENTROID_JUMP_UM"]:
        return False, f"centroid_jump={centroid_jump_um:.2f}um"

    # 2. Check width consistency
    if prev_width is not None and cand_width is not None:
        width_ratio = abs(cand_width - prev_width) / max(prev_width, 1e-9)
        if width_ratio > cfg["CONSERVATIVE_MAX_WIDTH_JUMP_RATIO"] * multiplier:
            return False, f"width_jump={width_ratio:.2f}"

    # 3. Check length consistency
    if prev_length is not None and cand_length is not None:
        length_ratio = abs(cand_length - prev_length) / max(prev_length, 1e-9)
        if length_ratio > cfg["CONSERVATIVE_MAX_LENGTH_JUMP_RATIO"] * multiplier:
            return False, f"length_jump={length_ratio:.2f}"

    # 4. Check area consistency
    if prev_area is not None and cand_area is not None:
        area_ratio = abs(cand_area - prev_area) / max(prev_area, 1e-9)
        if area_ratio > cfg["CONSERVATIVE_MAX_AREA_JUMP_RATIO"] * multiplier:
            return False, f"area_jump={area_ratio:.2f}"

    return True, "ok"

def _record_rejected_extension(event_map, track_id, z_index, reason):
    """Record one unique rejected candidate-extension event for a track."""
    event = f"z={int(z_index)}, reason={reason}"
    events = event_map.setdefault(int(track_id), [])
    if event not in events:
        events.append(event)


def _normalize_rejected_extension_events(event_map):
    """Normalize legacy string values and new event lists to lists of strings."""
    normalized = {}
    for track_id, raw_events in (event_map or {}).items():
        if raw_events is None:
            events = []
        elif isinstance(raw_events, str):
            events = [raw_events] if raw_events else []
        else:
            events = [str(value) for value in raw_events if str(value)]
        normalized[int(track_id)] = events
    return normalized


def _attach_tracking_audit(df, track_summary, rejected_extensions):
    """
    Attach rejected-extension history without claiming that a track stopped.

    A rejected candidate is only a near-miss considered by the tracker. The
    same track may subsequently link to a different, consistent detection.
    """
    events_by_track = _normalize_rejected_extension_events(rejected_extensions)
    counts = {track_id: len(events) for track_id, events in events_by_track.items()}
    joined = {track_id: " | ".join(events) for track_id, events in events_by_track.items()}
    first = {
        track_id: (events[0] if events else "")
        for track_id, events in events_by_track.items()
    }

    track_summary["rejected_extension_count"] = (
        track_summary["track_id"].map(counts).fillna(0).astype(int)
    )
    track_summary["has_rejected_extension"] = track_summary["rejected_extension_count"] > 0
    track_summary["rejected_extension_reasons"] = track_summary["track_id"].map(joined).fillna("")
    track_summary["first_rejected_extension_reason"] = track_summary["track_id"].map(first).fillna("")

    # Keep the old column for file compatibility, but do not populate it with
    # unconfirmed stop claims.
    track_summary["track_stop_reason"] = ""

    df["track_rejected_extension_count"] = (
        df["track_id"].map(counts).fillna(0).astype(int)
    )
    df["track_has_rejected_extension"] = df["track_rejected_extension_count"] > 0
    return df, track_summary


def track_across_slices_legacy(detections_df, cfg):
    """
    Conservative tracking natively: stop tracks when consistency breaks.
    """
    if detections_df.empty:
        detections_df = detections_df.copy()
        detections_df["track_id"] = pd.Series(dtype=int)
        return detections_df, pd.DataFrame()

    max_dist_px = cfg["TRACK_MAX_DIST_UM"] / (cfg["UM_PER_PX_XY"] + 1e-9)
    df = (detections_df.copy()
                       .sort_values(["z_slice", "sperm_id"])
                       .reset_index(drop=True))

    next_tid = 1
    active = {}
    track_ids = [-1] * len(df)
    link_methods = ["new"] * len(df)
    link_distances_um = [np.nan] * len(df)
    link_gap_slices = [0] * len(df)
    rejected_extensions = {}

    rows_by_z = {z: df.index[df["z_slice"] == z].to_numpy()
                 for z in sorted(df["z_slice"].unique())}

    for z, idxs in rows_by_z.items():
        # Get detection features for this slice
        xs = df.loc[idxs, "centroid_x"].to_numpy(float)
        ys = df.loc[idxs, "centroid_y"].to_numpy(float)

        # Extract morphological features (with fallbacks)
        widths = df.loc[idxs, "width_um"].to_numpy(float) if "width_um" in df.columns else np.full(len(idxs), np.nan)
        lengths = df.loc[idxs, "length_um_geodesic"].to_numpy(float) if "length_um_geodesic" in df.columns else np.full(len(idxs), np.nan)
        areas = df.loc[idxs, "area_px"].to_numpy(float) if "area_px" in df.columns else np.full(len(idxs), np.nan)
        oris = df.loc[idxs, "orientation"].to_numpy(float) if "orientation" in df.columns else np.full(len(idxs), np.nan)

        # Extract Bounding Box Arrays for Overlap-First Algorithm
        bbox_min_ys = df.loc[idxs, "bbox_min_y"].to_numpy(float) if "bbox_min_y" in df.columns else np.full(len(idxs), np.nan)
        bbox_min_xs = df.loc[idxs, "bbox_min_x"].to_numpy(float) if "bbox_min_x" in df.columns else np.full(len(idxs), np.nan)
        bbox_max_ys = df.loc[idxs, "bbox_max_y"].to_numpy(float) if "bbox_max_y" in df.columns else np.full(len(idxs), np.nan)
        bbox_max_xs = df.loc[idxs, "bbox_max_x"].to_numpy(float) if "bbox_max_x" in df.columns else np.full(len(idxs), np.nan)

        # Find candidate tracks from previous slices
        cand_tracks = [t for t, st in active.items()
                       if 1 <= z - st["last_z"] <= cfg["TRACK_MAX_GAP_SLICES"] + 1]

        used_det, used_trk = set(), set()

        if cand_tracks:
            candidates = []
            pad = cfg.get("TRACK_BBOX_PADDING_PX", 5.0)

            for k, (x, y) in enumerate(zip(xs, ys)):
                # Prepare Candidate State
                det_min_y, det_min_x = bbox_min_ys[k], bbox_min_xs[k]
                det_max_y, det_max_x = bbox_max_ys[k], bbox_max_xs[k]
                has_bbox = np.isfinite(det_min_y)

                cand_det = {
                    "z_slice": int(z),
                    "centroid_x": float(x),
                    "centroid_y": float(y),
                    "width_um": float(widths[k]) if np.isfinite(widths[k]) else None,
                    "length_um_geodesic": float(lengths[k]) if np.isfinite(lengths[k]) else None,
                    "area_px": float(areas[k]) if np.isfinite(areas[k]) else None,
                    "orientation": float(oris[k]) if np.isfinite(oris[k]) else None,
                }

                for j, tid in enumerate(cand_tracks):
                    trk_st = active[tid]

                    # Compute standard spatial distance
                    dx = float(x) - trk_st["last_x"]
                    dy = float(y) - trk_st["last_y"]
                    d_val = np.sqrt(dx*dx + dy*dy)

                    # Perform Overlap-First physical footprint collision check
                    overlap_exists = False
                    if has_bbox and "last_bbox" in trk_st and trk_st["last_bbox"] is not None:
                        t_min_y, t_min_x, t_max_y, t_max_x = trk_st["last_bbox"]
                        # Bounding box intersection formula with padding forgiveness
                        if not (det_max_y + pad < t_min_y or det_min_y - pad > t_max_y or
                                det_max_x + pad < t_min_x or det_min_x - pad > t_max_x):
                            overlap_exists = True

                    # Accept candidate if it physically overlaps, OR if it's within pure centroid limits (implausible fallback)
                    if overlap_exists or d_val <= max_dist_px:
                        # Continue unless implausible bounds
                        is_consistent, reason = check_extension_consistency(
                            trk_st, cand_det, cfg, overlap_exists=overlap_exists
                        )

                        if is_consistent:
                            # Massive scoring favor (+10,000) for overlaps so they greedily override standard loose distance matches
                            score = float(d_val) if not overlap_exists else (float(d_val) - 10000.0)
                            method = "overlap" if overlap_exists else "centroid"
                            gap_slices = int(z - trk_st["last_z"])
                            candidates.append((score, k, j, method, float(d_val), gap_slices))
                        else:
                            _record_rejected_extension(
                                rejected_extensions,
                                tid,
                                z,
                                reason,
                            )

            # Sort by score and assign greedily
            candidates.sort(key=lambda x: x[0])
            for score, det_k, trk_j, method, d_val_px, gap_slices in candidates:
                if det_k in used_det or trk_j in used_trk:
                    continue

                used_det.add(det_k)
                used_trk.add(trk_j)
                tid = cand_tracks[trk_j]
                row_idx = int(idxs[det_k])
                track_ids[row_idx] = tid
                link_methods[row_idx] = method
                link_distances_um[row_idx] = d_val_px * cfg["UM_PER_PX_XY"]
                link_gap_slices[row_idx] = gap_slices

                # Update track state
                previous = active[tid]
                active[tid] = {
                    "first_z": previous.get("first_z", previous["last_z"]),
                    "first_x": previous.get("first_x", previous["last_x"]),
                    "first_y": previous.get("first_y", previous["last_y"]),
                    "last_z": int(z),
                    "last_x": float(xs[det_k]),
                    "last_y": float(ys[det_k]),
                    "last_width": float(widths[det_k]) if np.isfinite(widths[det_k]) else None,
                    "last_length": float(lengths[det_k]) if np.isfinite(lengths[det_k]) else None,
                    "last_area": float(areas[det_k]) if np.isfinite(areas[det_k]) else None,
                    "last_orientation": float(oris[det_k]) if np.isfinite(oris[det_k]) else None,
                    "last_bbox": (bbox_min_ys[det_k], bbox_min_xs[det_k], bbox_max_ys[det_k], bbox_max_xs[det_k]) if np.isfinite(bbox_min_ys[det_k]) else None,
                    "max_length_2d": max(
                        float(previous.get("max_length_2d") or 0.0),
                        float(lengths[det_k]) if np.isfinite(lengths[det_k]) else 0.0,
                    ),
                }

        # Create new tracks for unmatched detections
        for det_k in range(len(idxs)):
            if track_ids[int(idxs[det_k])] == -1:
                track_ids[int(idxs[det_k])] = next_tid
                active[next_tid] = {
                    "first_z": int(z),
                    "first_x": float(xs[det_k]),
                    "first_y": float(ys[det_k]),
                    "last_z": int(z),
                    "last_x": float(xs[det_k]),
                    "last_y": float(ys[det_k]),
                    "last_width": float(widths[det_k]) if np.isfinite(widths[det_k]) else None,
                    "last_length": float(lengths[det_k]) if np.isfinite(lengths[det_k]) else None,
                    "last_area": float(areas[det_k]) if np.isfinite(areas[det_k]) else None,
                    "last_orientation": float(oris[det_k]) if np.isfinite(oris[det_k]) else None,
                    "last_bbox": (bbox_min_ys[det_k], bbox_min_xs[det_k], bbox_max_ys[det_k], bbox_max_xs[det_k]) if np.isfinite(bbox_min_ys[det_k]) else None,
                    "max_length_2d": float(lengths[det_k]) if np.isfinite(lengths[det_k]) else 0.0,
                }
                next_tid += 1

        # Remove stale tracks (exceeded gap)
        for tid in [t for t, st in active.items()
                    if z - st["last_z"] > cfg["TRACK_MAX_GAP_SLICES"] + 1]:
            del active[tid]

    df["track_id"] = track_ids
    df["track_link_method"] = link_methods
    df["track_link_distance_um"] = np.round(link_distances_um, 3)
    df["track_link_gap_slices"] = link_gap_slices

    # Print tracking stats
    print(
        "  Conservative tracking: "
        f"{len(rejected_extensions)} tracks encountered rejected candidate extensions"
    )

    # Inject maximum 2D Euclidean distance of the physical shape prior to grouping
    if "tortuosity" in df.columns:
        df["euc_um_2d"] = df["length_um_geodesic"] / df["tortuosity"]
    else:
        df["euc_um_2d"] = df["length_um_geodesic"]
    if "length_width_ratio" not in df.columns:
        width = pd.to_numeric(df.get("width_um"), errors="coerce")
        length = pd.to_numeric(df.get("length_um_geodesic"), errors="coerce")
        df["length_width_ratio"] = length / width.clip(lower=1e-9)

    g = df.groupby("track_id", as_index=False)
    ts = g.agg(
        n_slices        = ("z_slice",            "count"),
        z_start         = ("z_slice",            "min"),
        z_end           = ("z_slice",            "max"),
        max_length_2d   = ("length_um_geodesic", "max"),
        median_width_2d = ("width_um",            "median"),
        median_length_width_ratio_2d = ("length_width_ratio", "median"),
        max_euc_2d      = ("euc_um_2d",          "max"),
        sum_area_px     = ("area_px",            "sum"),
        min_area_px     = ("area_px",            "min"),
        max_area_px     = ("area_px",            "max"),
        area_start      = ("area_px",            "first"),
        area_end        = ("area_px",            "last"),
        x_mean          = ("centroid_x",         "mean"),
        y_mean          = ("centroid_y",         "mean"),
        z_mean          = ("z_slice",            "mean"),
        x_start         = ("centroid_x",         "first"),
        y_start         = ("centroid_y",         "first"),
        x_end           = ("centroid_x",         "last"),
        y_end           = ("centroid_y",         "last"),
    )

    um_xy = cfg["UM_PER_PX_XY"]
    um_z  = cfg["UM_PER_SLICE_Z"]

    # 1. Z span and sampled Z coverage are distinct:
    # span is endpoint-to-endpoint displacement; covered includes slice thickness.
    z_span = (ts["z_end"] - ts["z_start"]) * um_z
    z_covered = (ts["z_end"] - ts["z_start"] + 1) * um_z
    ts["z_extent_um"] = z_span
    ts["z_span_um"] = z_span
    ts["z_covered_um"] = z_covered

    # Geodesic vertical displacement for Euclidean distance is purely centroid-to-centroid
    dz_euc = z_span

    # 2. Total 3D Length (Lateral-Corrected Hypotenuse)
    # The physical 3D shape arc relies on the maximum length of its 2D projection
    euc_2d_centroid = np.sqrt((ts["x_end"] - ts["x_start"])**2 + (ts["y_end"] - ts["y_start"])**2) * um_xy
    lat_geodesic = np.maximum(ts["max_length_2d"], euc_2d_centroid)
    l3d = np.sqrt(lat_geodesic**2 + z_span**2)
    ts["total_3d_length_um"] = l3d

    # 3. 3D Volume (sum of per-slice projected area estimates * Z_step)
    ts["volume_um3"] = ts["sum_area_px"] * (um_xy**2) * um_z

    # 4. 3D Tortuosity (Total 3D Geodesic Length / 3D End-To-End Euclidean Distance)
    euc_3d = np.sqrt(ts["max_euc_2d"]**2 + dz_euc**2)
    safe_euc = np.maximum(euc_3d, 0.1)
    tort_raw = l3d / safe_euc
    ts["tortuosity_3d"] = np.minimum(tort_raw, 20.0)

    # 5. Taper Ratio (max/min area across the full track)
    ts["taper_ratio"] = ts["max_area_px"] / np.maximum(ts["min_area_px"], 0.001)

    # 6. Effective Thickness / Diameter
    cross_area = ts["volume_um3"] / np.maximum(ts["total_3d_length_um"], 0.1)
    ts["thickness_um"] = 2 * np.sqrt(cross_area / np.pi)

    # 7. Orientation Angles (Pitch and Yaw)
    dx = (ts["x_end"] - ts["x_start"]) * um_xy
    dy = (ts["y_end"] - ts["y_start"]) * um_xy
    v_mag = np.sqrt(dx**2 + dy**2 + dz_euc**2)
    safe_v = np.maximum(v_mag, 1e-9)
    ts["pitch_deg"] = np.abs(np.arcsin(dz_euc / safe_v)) * (180.0 / np.pi)
    ts["yaw_deg"] = np.arctan2(dy, dx) * (180.0 / np.pi)

    if len(ts) > 1:
        centers = np.column_stack((ts["x_mean"] * um_xy, ts["y_mean"] * um_xy, ts["z_mean"] * um_z))
        tree = cKDTree(centers)
        dists, _ = tree.query(centers, k=2)
        ts["nearest_neighbor_um"] = dists[:, 1]
    else:
        ts["nearest_neighbor_um"] = np.nan

    cols_ordered = [
        "track_id", "total_3d_length_um", "z_extent_um", "z_span_um", "z_covered_um", "volume_um3", "tortuosity_3d",
        "thickness_um", "pitch_deg", "yaw_deg", "taper_ratio", "nearest_neighbor_um",
        "n_slices", "z_start", "z_end", "max_length_2d",
        "median_width_2d", "median_length_width_ratio_2d", "sum_area_px",
        "min_area_px", "max_area_px", "area_start", "area_end"
    ]
    ts = ts[cols_ordered]
    return _attach_tracking_audit(df, ts, rejected_extensions)


def _angle_diff_deg(a, b):
    if a is None or b is None or not np.isfinite(a) or not np.isfinite(b):
        return 0.0
    d = abs(float(a) - float(b)) % 180.0
    return min(d, 180.0 - d)


def _relative_change(a, b):
    if a is None or b is None or not np.isfinite(a) or not np.isfinite(b):
        return 0.0
    return abs(float(a) - float(b)) / max(abs(float(b)), 1e-9)


def _bbox_overlap_fraction(box_a, box_b):
    if box_a is None or box_b is None:
        return 0.0
    ay0, ax0, ay1, ax1 = box_a
    by0, bx0, by1, bx1 = box_b
    vals = [ay0, ax0, ay1, ax1, by0, bx0, by1, bx1]
    if not all(np.isfinite(v) for v in vals):
        return 0.0
    inter_y = max(0.0, min(ay1, by1) - max(ay0, by0))
    inter_x = max(0.0, min(ax1, bx1) - max(ax0, bx0))
    inter = inter_y * inter_x
    area_a = max(0.0, (ay1 - ay0) * (ax1 - ax0))
    area_b = max(0.0, (by1 - by0) * (bx1 - bx0))
    denom = max(min(area_a, area_b), 1e-9)
    return float(min(inter / denom, 1.0))


def _finite_float_or_none(value):
    try:
        value = float(value)
    except (TypeError, ValueError):
        return None
    if not np.isfinite(value):
        return None
    return value


def _read_unet_probability(row):
    """
    Read optional U-Net support from a detection row.

    v5.7 treats U-Net fields as optional evidence. Classical Saturn detections
    do not need these columns, and missing fields produce no tracking penalty.
    """
    for col in (
        "unet_mean_probability",
        "unet_max_probability",
        "unet_seed_probability",
        "unet_probability",
    ):
        if col in row.index:
            value = _finite_float_or_none(row[col])
            if value is not None:
                return float(np.clip(value, 0.0, 1.0))
    return None


def _unet_tracking_enabled(cfg):
    engine = str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower()
    return bool(cfg.get("UNET_TRACKING_SUPPORT", True)) and engine in (
        "unet_assisted", "hybrid", "unet_primary"
    )


def _unet_link_cost_terms(det_prob, prev_prob, cfg, repair=False):
    if (not _unet_tracking_enabled(cfg)) or (det_prob is None and prev_prob is None):
        return 0.0, 0.0, 0.0

    support_weight = float(
        cfg.get(
            "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT" if repair else "ASSIGNMENT_UNET_SUPPORT_WEIGHT",
            0.0,
        )
    )
    continuity_weight = float(cfg.get("ASSIGNMENT_UNET_CONTINUITY_WEIGHT", 0.0))

    probs = [p for p in (det_prob, prev_prob) if p is not None]
    mean_prob = float(np.mean(probs)) if probs else 0.0
    support_term = support_weight * (1.0 - mean_prob)
    continuity_term = 0.0
    if det_prob is not None and prev_prob is not None:
        continuity_term = continuity_weight * abs(float(det_prob) - float(prev_prob))
    return support_term + continuity_term, support_term, continuity_term


def _summarize_tracked_detections(df, rejected_extensions, cfg):
    if df.empty:
        return df, pd.DataFrame()

    if "tortuosity" in df.columns:
        df["euc_um_2d"] = df["length_um_geodesic"] / df["tortuosity"]
    else:
        df["euc_um_2d"] = df["length_um_geodesic"]
    if "length_width_ratio" not in df.columns:
        width = pd.to_numeric(df.get("width_um"), errors="coerce")
        length = pd.to_numeric(df.get("length_um_geodesic"), errors="coerce")
        df["length_width_ratio"] = length / width.clip(lower=1e-9)

    g = df.groupby("track_id", as_index=False)
    ts = g.agg(
        n_slices        = ("z_slice",            "count"),
        z_start         = ("z_slice",            "min"),
        z_end           = ("z_slice",            "max"),
        max_length_2d   = ("length_um_geodesic", "max"),
        median_width_2d = ("width_um",            "median"),
        median_length_width_ratio_2d = ("length_width_ratio", "median"),
        max_euc_2d      = ("euc_um_2d",          "max"),
        sum_area_px     = ("area_px",            "sum"),
        min_area_px     = ("area_px",            "min"),
        max_area_px     = ("area_px",            "max"),
        area_start      = ("area_px",            "first"),
        area_end        = ("area_px",            "last"),
        x_mean          = ("centroid_x",         "mean"),
        y_mean          = ("centroid_y",         "mean"),
        z_mean          = ("z_slice",            "mean"),
        x_start         = ("centroid_x",         "first"),
        y_start         = ("centroid_y",         "first"),
        x_end           = ("centroid_x",         "last"),
        y_end           = ("centroid_y",         "last"),
    )

    um_xy = cfg["UM_PER_PX_XY"]
    um_z = cfg["UM_PER_SLICE_Z"]
    z_span = (ts["z_end"] - ts["z_start"]) * um_z
    z_covered = (ts["z_end"] - ts["z_start"] + 1) * um_z
    ts["z_extent_um"] = z_span
    ts["z_span_um"] = z_span
    ts["z_covered_um"] = z_covered

    euc_2d_centroid = np.sqrt((ts["x_end"] - ts["x_start"])**2 + (ts["y_end"] - ts["y_start"])**2) * um_xy
    lat_geodesic = np.maximum(ts["max_length_2d"], euc_2d_centroid)
    l3d = np.sqrt(lat_geodesic**2 + z_span**2)
    ts["total_3d_length_um"] = l3d
    ts["volume_um3"] = ts["sum_area_px"] * (um_xy**2) * um_z

    euc_3d = np.sqrt(ts["max_euc_2d"]**2 + z_span**2)
    safe_euc = np.maximum(euc_3d, 0.1)
    ts["tortuosity_3d"] = np.minimum(l3d / safe_euc, 20.0)
    ts["taper_ratio"] = ts["max_area_px"] / np.maximum(ts["min_area_px"], 0.001)

    cross_area = ts["volume_um3"] / np.maximum(ts["total_3d_length_um"], 0.1)
    ts["thickness_um"] = 2 * np.sqrt(cross_area / np.pi)

    dx = (ts["x_end"] - ts["x_start"]) * um_xy
    dy = (ts["y_end"] - ts["y_start"]) * um_xy
    v_mag = np.sqrt(dx**2 + dy**2 + z_span**2)
    safe_v = np.maximum(v_mag, 1e-9)
    ts["pitch_deg"] = np.abs(np.arcsin(z_span / safe_v)) * (180.0 / np.pi)
    ts["yaw_deg"] = np.arctan2(dy, dx) * (180.0 / np.pi)

    if len(ts) > 1:
        centers = np.column_stack((ts["x_mean"] * um_xy, ts["y_mean"] * um_xy, ts["z_mean"] * um_z))
        tree = cKDTree(centers)
        dists, _ = tree.query(centers, k=2)
        ts["nearest_neighbor_um"] = dists[:, 1]
    else:
        ts["nearest_neighbor_um"] = np.nan

    unet_summary_cols = []
    for prob_col in (
        "unet_mean_probability",
        "unet_max_probability",
        "unet_seed_probability",
        "unet_probability",
    ):
        if prob_col in df.columns:
            vals = pd.to_numeric(df[prob_col], errors="coerce")
            if vals.notna().any():
                df[prob_col] = vals
                mean_by_track = df.groupby("track_id")[prob_col].mean()
                max_by_track = df.groupby("track_id")[prob_col].max()
                ts[f"track_mean_{prob_col}"] = ts["track_id"].map(mean_by_track)
                ts[f"track_max_{prob_col}"] = ts["track_id"].map(max_by_track)
                unet_summary_cols.extend([f"track_mean_{prob_col}", f"track_max_{prob_col}"])

    cols_ordered = [
        "track_id", "total_3d_length_um", "z_extent_um", "z_span_um", "z_covered_um", "volume_um3", "tortuosity_3d",
        "thickness_um", "pitch_deg", "yaw_deg", "taper_ratio", "nearest_neighbor_um",
        "n_slices", "z_start", "z_end", "max_length_2d",
        "median_width_2d", "median_length_width_ratio_2d", "sum_area_px",
        "min_area_px", "max_area_px", "area_start", "area_end"
    ] + unet_summary_cols
    ts = ts[cols_ordered]
    return _attach_tracking_audit(df, ts, rejected_extensions)


def track_across_slices_global_assignment(detections_df, cfg):
    """
    V5.6 ROI-ADAPTIVE tracker: link each slice with a global assignment cost matrix.

    This is intentionally self-contained. It gives us a LapTrack-like assignment
    backend without requiring a new dependency while we evaluate whether global
    assignment improves the fragmentation/over-linking tradeoff.
    """
    if detections_df.empty:
        detections_df = detections_df.copy()
        detections_df["track_id"] = pd.Series(dtype=int)
        return detections_df, pd.DataFrame()

    df = (detections_df.copy()
                       .sort_values(["z_slice", "sperm_id"])
                       .reset_index(drop=True))
    track_ids = [-1] * len(df)
    link_methods = ["new"] * len(df)
    link_distances_um = [np.nan] * len(df)
    link_gap_slices = [0] * len(df)
    rejected_extensions = {}

    um_xy = cfg["UM_PER_PX_XY"]
    max_dist_um = float(cfg.get("TRACK_MAX_DIST_UM", 7.0))
    max_cost = float(cfg.get("ASSIGNMENT_MAX_COST", 8.0))
    weights = {
        "dist": float(cfg.get("ASSIGNMENT_DIST_WEIGHT", 1.0)),
        "overlap": float(cfg.get("ASSIGNMENT_OVERLAP_WEIGHT", 2.0)),
        "length": float(cfg.get("ASSIGNMENT_LENGTH_WEIGHT", 2.0)),
        "width": float(cfg.get("ASSIGNMENT_WIDTH_WEIGHT", 1.2)),
        "area": float(cfg.get("ASSIGNMENT_AREA_WEIGHT", 0.9)),
        "angle": float(cfg.get("ASSIGNMENT_ANGLE_WEIGHT", 0.4)),
    }

    next_tid = 1
    active = {}
    rows_by_z = {z: df.index[df["z_slice"] == z].to_numpy()
                 for z in sorted(df["z_slice"].unique())}

    for z, idxs in rows_by_z.items():
        dets = []
        for row_idx in idxs:
            row = df.loc[row_idx]
            bbox = None
            if {"bbox_min_y", "bbox_min_x", "bbox_max_y", "bbox_max_x"}.issubset(df.columns):
                bbox = (row["bbox_min_y"], row["bbox_min_x"], row["bbox_max_y"], row["bbox_max_x"])
            dets.append({
                "row_idx": int(row_idx),
                "z": int(z),
                "x": float(row["centroid_x"]),
                "y": float(row["centroid_y"]),
                "width": float(row["width_um"]) if "width_um" in df.columns and np.isfinite(row["width_um"]) else None,
                "length": float(row["length_um_geodesic"]) if "length_um_geodesic" in df.columns and np.isfinite(row["length_um_geodesic"]) else None,
                "area": float(row["area_px"]) if "area_px" in df.columns and np.isfinite(row["area_px"]) else None,
                "orientation": float(row["orientation"]) if "orientation" in df.columns and np.isfinite(row["orientation"]) else None,
                "unet_probability": _read_unet_probability(row),
                "bbox": bbox,
            })

        cand_tids = [tid for tid, st in active.items()
                     if 1 <= z - st["last_z"] <= cfg["TRACK_MAX_GAP_SLICES"] + 1]
        assigned_dets = set()
        assigned_tracks = set()

        if cand_tids and dets:
            cost = np.full((len(cand_tids), len(dets)), 1e9, dtype=float)
            dist_cache = {}
            method_cache = {}
            for ti, tid in enumerate(cand_tids):
                st = active[tid]
                for di, det in enumerate(dets):
                    estimated_joined_length = _estimated_tracking_extension_length_um(
                        st,
                        det["x"],
                        det["y"],
                        det["z"],
                        det["length"],
                        cfg,
                    )
                    if estimated_joined_length > float(
                        cfg.get("TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM", 15.0)
                    ):
                        continue
                    dx = det["x"] - st["last_x"]
                    dy = det["y"] - st["last_y"]
                    dist_um = math.sqrt(dx * dx + dy * dy) * um_xy
                    overlap = _bbox_overlap_fraction(det["bbox"], st.get("last_bbox"))
                    if dist_um > max_dist_um and overlap <= 0:
                        continue

                    dist_term = dist_um / max(max_dist_um, 1e-9)
                    overlap_term = 1.0 - overlap
                    length_term = _relative_change(det["length"], st.get("last_length"))
                    width_term = _relative_change(det["width"], st.get("last_width"))
                    area_term = _relative_change(det["area"], st.get("last_area"))
                    angle_term = _angle_diff_deg(det["orientation"], st.get("last_orientation")) / 90.0
                    unet_term, _, _ = _unet_link_cost_terms(
                        det["unet_probability"],
                        st.get("last_unet_probability"),
                        cfg,
                        repair=False,
                    )
                    c = (
                        weights["dist"] * dist_term
                        + weights["overlap"] * overlap_term
                        + weights["length"] * length_term
                        + weights["width"] * width_term
                        + weights["area"] * area_term
                        + weights["angle"] * angle_term
                        + unet_term
                    )
                    cost[ti, di] = c
                    dist_cache[(ti, di)] = dist_um
                    method_cache[(ti, di)] = "assignment_overlap" if overlap > 0 else "assignment_cost"

            row_ind, col_ind = linear_sum_assignment(cost)
            for ti, di in zip(row_ind, col_ind):
                if cost[ti, di] > max_cost:
                    continue
                tid = cand_tids[int(ti)]
                det = dets[int(di)]
                row_idx = det["row_idx"]
                track_ids[row_idx] = tid
                link_methods[row_idx] = method_cache.get((int(ti), int(di)), "assignment_cost")
                link_distances_um[row_idx] = dist_cache.get((int(ti), int(di)), np.nan)
                link_gap_slices[row_idx] = int(z - active[tid]["last_z"])
                assigned_dets.add(int(di))
                assigned_tracks.add(tid)
                previous = active[tid]
                active[tid] = {
                    "first_z": previous.get("first_z", previous["last_z"]),
                    "first_x": previous.get("first_x", previous["last_x"]),
                    "first_y": previous.get("first_y", previous["last_y"]),
                    "last_z": int(z),
                    "last_x": det["x"],
                    "last_y": det["y"],
                    "last_width": det["width"],
                    "last_length": det["length"],
                    "last_area": det["area"],
                    "last_orientation": det["orientation"],
                    "last_unet_probability": det["unet_probability"],
                    "last_bbox": det["bbox"],
                    "max_length_2d": max(
                        float(previous.get("max_length_2d") or 0.0),
                        float(det["length"] or 0.0),
                    ),
                }

            for ti, tid in enumerate(cand_tids):
                if tid not in assigned_tracks and np.isfinite(cost[ti]).any():
                    best = float(np.min(cost[ti]))
                    if best <= max_cost * 1.5:
                        _record_rejected_extension(
                            rejected_extensions,
                            tid,
                            z,
                            f"assignment_unmatched_cost={best:.2f}",
                        )

        for di, det in enumerate(dets):
            row_idx = det["row_idx"]
            if di in assigned_dets or track_ids[row_idx] != -1:
                continue
            track_ids[row_idx] = next_tid
            active[next_tid] = {
                "first_z": int(z),
                "first_x": det["x"],
                "first_y": det["y"],
                "last_z": int(z),
                "last_x": det["x"],
                "last_y": det["y"],
                "last_width": det["width"],
                "last_length": det["length"],
                "last_area": det["area"],
                "last_orientation": det["orientation"],
                "last_unet_probability": det["unet_probability"],
                "last_bbox": det["bbox"],
                "max_length_2d": float(det["length"] or 0.0),
            }
            next_tid += 1

        for tid in [t for t, st in active.items()
                    if z - st["last_z"] > cfg["TRACK_MAX_GAP_SLICES"] + 1]:
            del active[tid]

    df["track_id"] = track_ids
    df["track_link_method"] = link_methods
    df["track_link_distance_um"] = np.round(link_distances_um, 3)
    df["track_link_gap_slices"] = link_gap_slices
    print(
        "  Global-assignment tracking: "
        f"{len(rejected_extensions)} tracks encountered rejected assignment extensions"
    )
    return _summarize_tracked_detections(df, rejected_extensions, cfg)


def _row_endpoint(row):
    bbox = None
    if {"bbox_min_y", "bbox_min_x", "bbox_max_y", "bbox_max_x"}.issubset(row.index):
        bbox = (row["bbox_min_y"], row["bbox_min_x"], row["bbox_max_y"], row["bbox_max_x"])
    return {
        "z": int(row["z_slice"]),
        "x": float(row["centroid_x"]),
        "y": float(row["centroid_y"]),
        "width": float(row["width_um"]) if "width_um" in row.index and np.isfinite(row["width_um"]) else None,
        "length": float(row["length_um_geodesic"]) if "length_um_geodesic" in row.index and np.isfinite(row["length_um_geodesic"]) else None,
        "area": float(row["area_px"]) if "area_px" in row.index and np.isfinite(row["area_px"]) else None,
        "orientation": float(row["orientation"]) if "orientation" in row.index and np.isfinite(row["orientation"]) else None,
        "unet_probability": _read_unet_probability(row),
        "bbox": bbox,
    }


def _hybrid_repair_cost(src_end, dst_start, cfg):
    um_xy = cfg["UM_PER_PX_XY"]
    max_dist_um = float(cfg.get("HYBRID_REPAIR_MAX_LINK_DIST_UM", cfg.get("TRACK_MAX_DIST_UM", 6.0)))
    weights = {
        "dist": float(cfg.get("ASSIGNMENT_DIST_WEIGHT", 1.0)),
        "overlap": float(cfg.get("ASSIGNMENT_OVERLAP_WEIGHT", 2.0)),
        "length": float(cfg.get("ASSIGNMENT_LENGTH_WEIGHT", 2.0)),
        "width": float(cfg.get("ASSIGNMENT_WIDTH_WEIGHT", 1.2)),
        "area": float(cfg.get("ASSIGNMENT_AREA_WEIGHT", 0.9)),
        "angle": float(cfg.get("ASSIGNMENT_ANGLE_WEIGHT", 0.4)),
    }
    dx = dst_start["x"] - src_end["x"]
    dy = dst_start["y"] - src_end["y"]
    dist_um = math.sqrt(dx * dx + dy * dy) * um_xy
    overlap = _bbox_overlap_fraction(dst_start["bbox"], src_end["bbox"])
    if dist_um > max_dist_um and overlap <= 0:
        return np.inf, dist_um, overlap

    dist_term = dist_um / max(max_dist_um, 1e-9)
    overlap_term = 1.0 - overlap
    unet_term, _, _ = _unet_link_cost_terms(
        dst_start.get("unet_probability"),
        src_end.get("unet_probability"),
        cfg,
        repair=True,
    )
    cost = (
        weights["dist"] * dist_term
        + weights["overlap"] * overlap_term
        + weights["length"] * _relative_change(dst_start["length"], src_end["length"])
        + weights["width"] * _relative_change(dst_start["width"], src_end["width"])
        + weights["area"] * _relative_change(dst_start["area"], src_end["area"])
        + weights["angle"] * (_angle_diff_deg(dst_start["orientation"], src_end["orientation"]) / 90.0)
        + unet_term
    )
    return cost, dist_um, overlap


def _estimated_merged_length_um(df, track_ids, cfg):
    sub = df[df["track_id"].isin(track_ids)].sort_values(["z_slice", "sperm_id"])
    if sub.empty:
        return np.inf
    um_xy = cfg["UM_PER_PX_XY"]
    um_z = cfg["UM_PER_SLICE_Z"]
    first = sub.iloc[0]
    last = sub.iloc[-1]
    z_span = (float(sub["z_slice"].max()) - float(sub["z_slice"].min())) * um_z
    centroid_span = math.sqrt(
        (float(last["centroid_x"]) - float(first["centroid_x"])) ** 2
        + (float(last["centroid_y"]) - float(first["centroid_y"])) ** 2
    ) * um_xy
    lat_len = max(float(sub["length_um_geodesic"].max()), centroid_span)
    return math.sqrt(lat_len * lat_len + z_span * z_span)


def track_across_slices_hybrid_repair(detections_df, cfg):
    """
    V5.6 ROI-ADAPTIVE tracker: keep the conservative legacy tracker, then run a
    narrow global-assignment-style repair only on short fragments.

    The repair pass is deliberately asymmetric: it can merge obvious fragments,
    but ambiguous candidates remain split so the audit does not inherit long
    false-positive chains.
    """
    df, ts = track_across_slices_legacy(detections_df, cfg)
    if df.empty or ts.empty or len(ts) < 2:
        return df, ts

    max_gap = int(cfg.get("HYBRID_REPAIR_MAX_GAP_SLICES", cfg.get("TRACK_MAX_GAP_SLICES", 1)))
    max_frag_slices = int(cfg.get("HYBRID_REPAIR_MAX_FRAGMENT_SLICES", 2))
    max_cost = float(cfg.get("HYBRID_REPAIR_MAX_COST", 3.6))
    min_overlap = float(cfg.get("HYBRID_REPAIR_MIN_OVERLAP", 0.05))
    max_final_length = min(
        float(cfg.get("HYBRID_REPAIR_MAX_FINAL_LENGTH_UM", 15.0)),
        float(cfg.get("TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM", 15.0)),
    )

    endpoints = {}
    for tid, grp in df.sort_values(["z_slice", "sperm_id"]).groupby("track_id"):
        endpoints[int(tid)] = {
            "first": _row_endpoint(grp.iloc[0]),
            "last": _row_endpoint(grp.iloc[-1]),
        }

    ts_by_tid = ts.set_index("track_id")
    candidates = []
    tids = list(endpoints.keys())
    starts_by_z = {}
    for tid in tids:
        starts_by_z.setdefault(endpoints[tid]["first"]["z"], []).append(tid)

    for src_tid in tids:
        src_end = endpoints[src_tid]["last"]
        src_n = int(ts_by_tid.loc[src_tid, "n_slices"])
        possible_dsts = []
        for gap in range(max_gap + 1):
            possible_dsts.extend(starts_by_z.get(src_end["z"] + gap + 1, []))
        for dst_tid in possible_dsts:
            if src_tid == dst_tid:
                continue
            dst_start = endpoints[dst_tid]["first"]
            dst_n = int(ts_by_tid.loc[dst_tid, "n_slices"])
            z_gap = dst_start["z"] - src_end["z"] - 1
            if z_gap < 0 or z_gap > max_gap:
                continue
            if src_n > max_frag_slices and dst_n > max_frag_slices:
                continue
            cost, dist_um, overlap = _hybrid_repair_cost(src_end, dst_start, cfg)
            if not np.isfinite(cost) or cost > max_cost:
                continue
            if overlap < min_overlap and dist_um > float(cfg.get("HYBRID_REPAIR_MAX_LINK_DIST_UM", cfg.get("TRACK_MAX_DIST_UM", 6.0))) * 0.6:
                continue
            est_len = _estimated_merged_length_um(df, [src_tid, dst_tid], cfg)
            if est_len > max_final_length:
                continue
            candidates.append((cost, dist_um, overlap, src_tid, dst_tid))

    parent = {tid: tid for tid in tids}

    def find(tid):
        while parent[tid] != tid:
            parent[tid] = parent[parent[tid]]
            tid = parent[tid]
        return tid

    repair_candidate_count = len(candidates)
    repair_evaluated_count = 0
    repair_accepted_count = 0
    repair_rejected_same_z_count = 0
    repair_rejected_length_count = 0
    repair_skipped_same_root_count = 0
    repair_skipped_target_repaired_count = 0
    repaired_targets = set()
    for (
        cost,
        dist_um,
        overlap,
        src_tid,
        dst_tid,
    ) in sorted(candidates, key=lambda item: item[0]):
        src_root = find(src_tid)
        dst_root = find(dst_tid)

        if src_root == dst_root:
            repair_skipped_same_root_count += 1
            continue

        if dst_root in repaired_targets:
            repair_skipped_target_repaired_count += 1
            continue

        repair_evaluated_count += 1

        src_members = [tid for tid in tids if find(tid) == src_root]
        dst_members = [tid for tid in tids if find(tid) == dst_root]

        src_z = set(df[df["track_id"].isin(src_members)]["z_slice"])
        dst_z = set(df[df["track_id"].isin(dst_members)]["z_slice"])

        if src_z & dst_z:
            repair_rejected_same_z_count += 1
            continue

        merged_members = src_members + dst_members
        if _estimated_merged_length_um(df, merged_members, cfg) > max_final_length:
            repair_rejected_length_count += 1
            continue
        parent[dst_root] = src_root
        repaired_targets.add(dst_root)
        repair_accepted_count += 1

        dst_start_z = endpoints[dst_tid]["first"]["z"]
        first_dst_mask = (df["track_id"] == dst_tid) & (df["z_slice"] == dst_start_z)
        first_dst_idx = df[first_dst_mask].index[:1]
        if len(first_dst_idx):
            idx = first_dst_idx[0]
            df.loc[idx, "track_link_method"] = "hybrid_repair"
            df.loc[idx, "track_link_distance_um"] = round(float(dist_um), 3)
            df.loc[idx, "track_link_gap_slices"] = int(endpoints[dst_tid]["first"]["z"] - endpoints[src_tid]["last"]["z"])

    if repair_accepted_count:
        df["track_id"] = df["track_id"].map(lambda tid: find(int(tid)))

    rejected_extensions = {}
    for _, row in ts.iterrows():
        old_track_id = int(row["track_id"])
        final_track_id = find(old_track_id)
        raw_reasons = row.get("rejected_extension_reasons", "")
        raw_reasons = "" if pd.isna(raw_reasons) else str(raw_reasons)
        events = [value for value in raw_reasons.split(" | ") if value]
        rejected_extensions.setdefault(final_track_id, []).extend(events)

    print(f"  Hybrid repair tracking: {repair_accepted_count} conservative fragment merges accepted")

    if not df.empty and df.groupby(["track_id", "z_slice"]).size().max() > 1:
        raise RuntimeError("Final track contains duplicate observations at the same z_slice.")

    final_df, final_ts = _summarize_tracked_detections(df, rejected_extensions, cfg)

    members_per_final = {}
    for old_track_id in tids:
        final_track_id = find(old_track_id)
        members_per_final[final_track_id] = members_per_final.get(final_track_id, 0) + 1
    repair_merges = {
        track_id: max(0, member_count - 1)
        for track_id, member_count in members_per_final.items()
    }
    final_ts["hybrid_repair_merge_count"] = (
        final_ts["track_id"].map(repair_merges).fillna(0).astype(int)
    )
    final_df["track_hybrid_repair_merge_count"] = (
        final_df["track_id"].map(repair_merges).fillna(0).astype(int)
    )

    audit_columns = {
        "track_hybrid_repair_candidate_count":
            repair_candidate_count,
        "track_hybrid_repair_evaluated_count":
            repair_evaluated_count,
        "track_hybrid_repair_accepted_count":
            repair_accepted_count,
        "track_hybrid_repair_rejected_same_z_count":
            repair_rejected_same_z_count,
        "track_hybrid_repair_rejected_length_count":
            repair_rejected_length_count,
        "track_hybrid_repair_skipped_same_root_count":
            repair_skipped_same_root_count,
        "track_hybrid_repair_skipped_target_repaired_count":
            repair_skipped_target_repaired_count,
    }

    for column, value in audit_columns.items():
        final_df[column] = int(value)
        final_ts[column] = int(value)

    return final_df, final_ts

def track_across_slices(detections_df, cfg):
    backend = str(cfg.get("TRACKING_BACKEND", "legacy")).strip().lower()
    if backend in ("hybrid_repair", "hybrid", "repair"):
        return track_across_slices_hybrid_repair(detections_df, cfg)
    if backend in ("global_assignment", "assignment", "hungarian"):
        return track_across_slices_global_assignment(detections_df, cfg)
    return track_across_slices_legacy(detections_df, cfg)


# =============================================================================
# QUALITY AUDIT & OUTLIER FLAGGING
# =============================================================================

def _join_flag_names(mask_list, n_rows):
    strings = [""] * n_rows
    any_mask = np.zeros(n_rows, dtype=bool)
    for flag_name, mask in mask_list:
        mask_arr = mask.values if hasattr(mask, "values") else np.array(mask)
        mask_arr = mask_arr.astype(bool)
        any_mask |= mask_arr
        for i in range(n_rows):
            if mask_arr[i]:
                strings[i] = f"{strings[i]},{flag_name}" if strings[i] else flag_name
    return strings, any_mask


def _series_false(df_tracks):
    return pd.Series(False, index=df_tracks.index)


def _comparative_audit_masks(df_tracks, cfg):
    """Return technical-failure and morphology-warning masks for comparative analysis."""
    technical_masks = []
    morphology_masks = []
    reference_masks = []
    n = len(df_tracks)
    false = _series_false(df_tracks)

    for col in ("track_id", "centroid_x", "centroid_y"):
        if col in df_tracks.columns:
            bad = ~np.isfinite(pd.to_numeric(df_tracks[col], errors="coerce"))
            technical_masks.append((f"invalid_{col}", bad))

    if "total_3d_length_um" in df_tracks.columns:
        length = pd.to_numeric(df_tracks["total_3d_length_um"], errors="coerce")
        invalid_length = (~np.isfinite(length)) | (length <= 0)
        technical_masks.append(("invalid_length", invalid_length))
        morphology_masks.extend([
            ("long", length > cfg.get("AUDIT_MAX_LENGTH_UM", 15.0)),
            ("short", length < cfg.get("MIN_SKEL_LEN_UM", 0.0)),
            ("extreme_component_length", length > cfg.get("MAX_GEODESIC_LEN_UM", 20.0) * 2.0),
        ])
        reference_masks.extend([
            ("long", length > cfg.get("AUDIT_MAX_LENGTH_UM", 15.0)),
            ("short", length < cfg.get("MIN_SKEL_LEN_UM", 0.0)),
        ])

    if "tortuosity_3d" in df_tracks.columns:
        tort = pd.to_numeric(df_tracks["tortuosity_3d"], errors="coerce")
        morphology_masks.append(("high_tortuosity", tort > cfg.get("AUDIT_MAX_TORTUOSITY", 1.5)))
        reference_masks.append(("high_tortuosity", tort > cfg.get("AUDIT_MAX_TORTUOSITY", 1.5)))

    if "thickness_um" in df_tracks.columns:
        thick = pd.to_numeric(df_tracks["thickness_um"], errors="coerce")
        morphology_masks.extend([
            ("wide", thick > cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0)),
            ("thin", thick < max(0.0, cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0) * 0.20)),
        ])
        reference_masks.append(("wide", thick > cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0)))

    if "taper_ratio" in df_tracks.columns:
        taper = pd.to_numeric(df_tracks["taper_ratio"], errors="coerce")
        morphology_masks.extend([
            ("high_taper", taper > cfg.get("AUDIT_MAX_TAPER_RATIO", 1.5)),
            ("low_taper", taper < 1.0),
        ])
        reference_masks.append(("high_taper", taper > cfg.get("AUDIT_MAX_TAPER_RATIO", 1.5)))

    if "length_width_ratio" in df_tracks.columns:
        lwr = pd.to_numeric(df_tracks["length_width_ratio"], errors="coerce")
        morphology_masks.append(("low_length_width_ratio", lwr < cfg.get("MIN_LENGTH_WIDTH_RATIO", 2.5)))
        reference_masks.append(("low_length_width_ratio", lwr < cfg.get("MIN_LENGTH_WIDTH_RATIO", 2.5)))

    if "pitch_deg" in df_tracks.columns:
        pitch = pd.to_numeric(df_tracks["pitch_deg"], errors="coerce")
        morphology_masks.append(("unusual_pitch", (pitch < 5.0) | (pitch > 175.0)))

    if "volume_um3" in df_tracks.columns:
        volume = pd.to_numeric(df_tracks["volume_um3"], errors="coerce")
        morphology_masks.append(("unusual_volume", volume <= 0))

    z_col = next((c for c in ("z_span_um", "z_covered_um") if c in df_tracks.columns), None)
    if z_col:
        zspan = pd.to_numeric(df_tracks[z_col], errors="coerce")
        morphology_masks.append(("unusual_z_span", zspan < 0))

    if "nearest_neighbor_um" in df_tracks.columns:
        nn = pd.to_numeric(df_tracks["nearest_neighbor_um"], errors="coerce")
        morphology_masks.append(("unusual_nearest_neighbor_distance", (nn >= 0) & (nn < cfg.get("UM_PER_PX_XY", 1.0))))

    branch_col = next((c for c in ("max_branch_nodes", "branch_nodes", "n_branch_nodes") if c in df_tracks.columns), None)
    if branch_col:
        branches = pd.to_numeric(df_tracks[branch_col], errors="coerce").fillna(0)
        technical_masks.append(("gross_branched_tissue_network", branches > max(3, cfg.get("MAX_BRANCH_NODES", 0) + 3)))

    if "segmentation_leakage" in df_tracks.columns:
        technical_masks.append(("segmentation_leakage", df_tracks["segmentation_leakage"].astype(bool)))
    if "outside_roi_pixel_count" in df_tracks.columns:
        technical_masks.append(("outside_roi", pd.to_numeric(df_tracks["outside_roi_pixel_count"], errors="coerce").fillna(0) > 0))
    if "exclusion_mask_overlap_count" in df_tracks.columns:
        technical_masks.append(("exclusion_mask_overlap", pd.to_numeric(df_tracks["exclusion_mask_overlap_count"], errors="coerce").fillna(0) > 0))
    if "suspected_multi_object_merge" in df_tracks.columns:
        technical_masks.append(("clear_multi_object_connected_component", df_tracks["suspected_multi_object_merge"].astype(bool)))
    if "unrecoverable_tracking_inconsistency" in df_tracks.columns:
        technical_masks.append(("unrecoverable_label_or_tracking_inconsistency", df_tracks["unrecoverable_tracking_inconsistency"].astype(bool)))

    if not technical_masks:
        technical_masks.append(("none", false))
    if not morphology_masks:
        morphology_masks.append(("none", false))
    if not reference_masks:
        reference_masks.append(("none", false))
    return technical_masks, morphology_masks, reference_masks


def flag_quality_tracks(df_tracks, cfg):
    """
    Annotate completed tracks without deleting data.

    In comparative mode, morphology outliers are warnings and remain in the
    primary ``technical_valid`` population. Reference-morphology filtering is a
    diagnostic subset only. Legacy columns are still emitted for older reports.
    """
    if df_tracks.empty:
        df_tracks["is_quality_track"] = pd.Series(dtype=bool)
        df_tracks["quality_flags"] = pd.Series(dtype=str)
        df_tracks["is_biological_candidate"] = pd.Series(dtype=bool)
        df_tracks["hard_flags"] = pd.Series(dtype=str)
        df_tracks["warning_flags"] = pd.Series(dtype=str)
        df_tracks["technical_valid"] = pd.Series(dtype=bool)
        df_tracks["technical_failure_reasons"] = pd.Series(dtype=str)
        df_tracks["morphology_warning"] = pd.Series(dtype=bool)
        df_tracks["morphology_warning_reasons"] = pd.Series(dtype=str)
        df_tracks["reference_morphology_pass"] = pd.Series(dtype=bool)
        df_tracks["is_reference_morphology_track"] = pd.Series(dtype=bool)
        df_tracks["is_warning_free_track"] = pd.Series(dtype=bool)
        df_tracks["segmentation_parameter_set"] = pd.Series(dtype=str)
        df_tracks["preprocessing_profile"] = pd.Series(dtype=str)
        df_tracks["analysis_mode"] = pd.Series(dtype=str)
        return df_tracks

    mode = str(cfg.get("ANALYSIS_MODE", "comparative")).strip().lower()
    parameter_set = str(cfg.get("SEGMENTATION_PARAMETER_SET", "")).strip()
    if not parameter_set:
        engine = str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn")).strip().lower()
        parameter_set = f"{engine}_{_VERSION}"
    preprocessing_profile = str(
        cfg.get(
            "RESOLVED_PREPROCESSING_PROFILE",
            cfg.get("CLAHE_MODE", cfg.get("PREPROCESS_MODE", "unspecified")),
        )
    )
    n = len(df_tracks)
    if mode == "comparative":
        technical_masks, morphology_masks, reference_masks = _comparative_audit_masks(df_tracks, cfg)
        technical_strs, any_technical = _join_flag_names(
            [(name, mask) for name, mask in technical_masks if name != "none"], n)
        morphology_strs, any_morphology = _join_flag_names(
            [(name, mask) for name, mask in morphology_masks if name != "none"], n)
        reference_strs, any_reference_fail = _join_flag_names(
            [(name, mask) for name, mask in reference_masks if name != "none"], n)

        df_tracks["technical_valid"] = ~any_technical
        df_tracks["technical_failure_reasons"] = technical_strs
        df_tracks["morphology_warning"] = any_morphology & (~any_technical)
        df_tracks["morphology_warning_reasons"] = morphology_strs
        df_tracks["reference_morphology_pass"] = (~any_technical) & (~any_reference_fail)
        df_tracks["is_reference_morphology_track"] = df_tracks["reference_morphology_pass"].astype(bool)
        df_tracks["is_warning_free_track"] = (
            df_tracks["technical_valid"] & (~df_tracks["morphology_warning"])
        )
        df_tracks["analysis_mode"] = mode
        df_tracks["segmentation_parameter_set"] = parameter_set
        df_tracks["preprocessing_profile"] = preprocessing_profile

        df_tracks["quality_flags"] = [
            ",".join([x for x in [technical_strs[i], morphology_strs[i]] if x])
            for i in range(n)
        ]
        # Legacy compatibility alias. New reports use the explicit reference
        # and warning-free columns above.
        df_tracks["is_quality_track"] = df_tracks["is_reference_morphology_track"]
        df_tracks["hard_flags"] = technical_strs
        df_tracks["warning_flags"] = morphology_strs
        df_tracks["is_biological_candidate"] = df_tracks["technical_valid"].astype(bool)
        df_tracks["has_warning_only"] = df_tracks["morphology_warning"].astype(bool)

        print(
            "  Comparative audit: "
            f"{int(df_tracks['technical_valid'].sum())} technical-valid, "
            f"{int(df_tracks['reference_morphology_pass'].sum())} reference-morphology, "
            f"{int(df_tracks['morphology_warning'].sum())} morphology-warning, "
            f"{int((~df_tracks['technical_valid']).sum())} technical failures out of {n} total"
        )
        return df_tracks

    strict_masks = []
    hard_masks = []
    warning_masks = []

    # Length
    if "total_3d_length_um" in df_tracks.columns:
        long_mask = df_tracks["total_3d_length_um"] > cfg.get("AUDIT_MAX_LENGTH_UM", 15.0)
        strict_masks.append(("long", long_mask))
        hard_masks.append(("long", long_mask))

    # Tortuosity
    if "tortuosity_3d" in df_tracks.columns:
        tort_mask = df_tracks["tortuosity_3d"] > cfg.get("AUDIT_MAX_TORTUOSITY", 1.5)
        strict_masks.append(("tortuous", tort_mask))
        hard_masks.append(("tortuous", tort_mask))

    # Thickness
    if "thickness_um" in df_tracks.columns:
        thick_mask = df_tracks["thickness_um"] > cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0)
        extreme_thick_mask = df_tracks["thickness_um"] > cfg.get("AUDIT_EXTREME_THICKNESS_UM", 3.5)
        strict_masks.append(("thick", thick_mask))
        warning_masks.append(("thick", thick_mask & ~extreme_thick_mask))
        hard_masks.append(("extreme_thick", extreme_thick_mask))

    # Taper
    if "taper_ratio" in df_tracks.columns:
        taper_mask = df_tracks["taper_ratio"] > cfg.get("AUDIT_MAX_TAPER_RATIO", 1.5)
        extreme_taper_mask = df_tracks["taper_ratio"] > cfg.get("AUDIT_EXTREME_TAPER_RATIO", 3.0)
        strict_masks.append(("taper", taper_mask))
        warning_masks.append(("taper", taper_mask & ~extreme_taper_mask))
        hard_masks.append(("extreme_taper", extreme_taper_mask))

    # Single-slice / shallow
    if "n_slices" in df_tracks.columns:
        shallow_mask = df_tracks["n_slices"] < cfg.get("AUDIT_MIN_SLICES", 2)
        strict_masks.append(("single_slice", shallow_mask))
        hard_masks.append(("single_slice", shallow_mask))

    flag_strs, any_flagged = _join_flag_names(strict_masks, n)
    hard_strs, any_hard = _join_flag_names(hard_masks, n)
    warning_strs, any_warning = _join_flag_names(warning_masks, n)

    df_tracks["quality_flags"] = flag_strs
    df_tracks["is_quality_track"] = ~any_flagged
    df_tracks["hard_flags"] = hard_strs
    df_tracks["warning_flags"] = warning_strs
    df_tracks["is_biological_candidate"] = ~any_hard
    df_tracks["has_warning_only"] = any_warning & ~any_hard
    df_tracks["technical_valid"] = ~any_hard
    df_tracks["technical_failure_reasons"] = hard_strs
    df_tracks["morphology_warning"] = any_flagged & ~any_hard
    df_tracks["morphology_warning_reasons"] = warning_strs
    df_tracks["reference_morphology_pass"] = ~any_flagged
    df_tracks["is_reference_morphology_track"] = df_tracks["reference_morphology_pass"].astype(bool)
    df_tracks["is_warning_free_track"] = (~any_flagged).astype(bool)
    df_tracks["segmentation_parameter_set"] = parameter_set
    df_tracks["preprocessing_profile"] = preprocessing_profile
    df_tracks["analysis_mode"] = mode

    n_quality = int((~any_flagged).sum())
    n_flagged = int(any_flagged.sum())
    n_candidates = int((~any_hard).sum())
    n_hard = int(any_hard.sum())
    n_warning_only = int((any_warning & ~any_hard).sum())
    print(
        f"  Audit: {n_candidates} estimated nuclei, {n_hard} technical failures, "
        f"{n_warning_only} morphology warnings, {n_quality} warning-free out of {n} reconstructed tracks"
    )

    return df_tracks


def export_comparative_track_tables(out_dir, track_summary, version_label=None):
    """
    Save exceptional comparative-mode audit populations.

    The canonical biological population is already written by
    :func:`export_biologist_results`. The master track summary contains all
    flags, so this function only creates a separate table when technical
    failures actually exist.
    """
    if track_summary is None or track_summary.empty:
        return {}
    ensure_dir(out_dir)
    suffix = f"_{version_label}" if version_label else ""
    technical_valid = track_summary["technical_valid"].astype(bool) if "technical_valid" in track_summary.columns else pd.Series(True, index=track_summary.index)
    paths = {}
    technical_failures = track_summary[~technical_valid].copy()
    if not technical_failures.empty:
        name = f"track_summary_technical_failures{suffix}.csv"
        path = os.path.join(out_dir, name)
        technical_failures.to_csv(path, index=False)
        paths[name] = path
    note = (
        "Primary sample-level table: biologist_results/sample_summary*.csv\n"
        "Primary nucleus-level table: biologist_results/nuclei_for_analysis*.csv\n"
        "Master audit table: track_summary*.csv\n"
        "Morphology warnings remain in the primary population because they may "
        "represent genuine genotype-dependent phenotypes.\n"
        "The legacy is_biological_candidate flag is identical to technical_valid "
        "and does not define a separate population.\n"
    )
    with open(os.path.join(out_dir, f"comparative_population_note{suffix}.txt"), "w", encoding="utf-8") as f:
        f.write(note)
    return paths


def export_biologist_results(out_dir, track_summary, version_label=None):
    """Export one canonical analysis population with a compact metric set."""
    if track_summary is None or track_summary.empty:
        return {}

    result_dir = os.path.join(out_dir, "biologist_results")
    ensure_dir(result_dir)
    suffix = f"_{version_label}" if version_label else ""
    primary = (
        track_summary[track_summary["technical_valid"].astype(bool)].copy()
        if "technical_valid" in track_summary.columns
        else track_summary.copy()
    )

    column_map = {
        "track_id": "estimated_nucleus_id",
        "total_3d_length_um": "length_3d_um",
        "max_length_2d": "maximum_2d_length_um",
        "median_width_2d": "median_2d_width_um",
        "median_length_width_ratio_2d": "median_2d_length_width_ratio",
        "thickness_um": "effective_thickness_um_psf_sensitive",
        "tortuosity_3d": "tortuosity_3d",
        "z_span_um": "z_span_um",
        "n_slices": "slices_detected",
        "morphology_warning": "morphology_warning_for_review",
    }
    available = [column for column in column_map if column in primary.columns]
    nuclei = primary[available].rename(columns=column_map)
    nuclei_path = os.path.join(result_dir, f"nuclei_for_analysis{suffix}.csv")
    nuclei.to_csv(nuclei_path, index=False)

    def median(column):
        if column not in primary.columns or primary.empty:
            return np.nan
        return float(pd.to_numeric(primary[column], errors="coerce").median())

    summary = pd.DataFrame([{
        "analysis_population": "included estimated nuclei",
        "estimated_unique_nuclei": int(len(primary)),
        "median_3d_length_um": median("total_3d_length_um"),
        "median_maximum_2d_length_um": median("max_length_2d"),
        "median_2d_width_um": median("median_width_2d"),
        "median_2d_length_width_ratio": median(
            "median_length_width_ratio_2d"
        ),
        "median_effective_thickness_um_psf_sensitive": median("thickness_um"),
        "median_3d_tortuosity": median("tortuosity_3d"),
        "median_z_span_um": median("z_span_um"),
        "median_slices_detected": median("n_slices"),
    }])
    summary_path = os.path.join(result_dir, f"sample_summary{suffix}.csv")
    summary.to_csv(summary_path, index=False)

    readme_path = os.path.join(result_dir, "README.txt")
    with open(readme_path, "w", encoding="utf-8") as handle:
        handle.write(
            "BIOLOGIST RESULTS\n"
            "==================\n"
            "Use sample_summary*.csv for sample-level comparisons.\n"
            "Use nuclei_for_analysis*.csv for nucleus-level statistics.\n\n"
            "Primary population: included estimated nuclei reconstructed in 3D.\n"
            "Morphology-warning tracks remain included because they may represent real biology.\n"
            "Effective thickness is PSF-sensitive and should be compared only between matched acquisitions.\n\n"
            "Do not use raw 2D detections, U-Net contribution counts, warning-free counts,\n"
            "reference-morphology counts, or rejected-extension counts as the biological nucleus count.\n"
        )
    return {
        "summary": summary_path,
        "nuclei": nuclei_path,
        "readme": readme_path,
    }


def _technical_valid_track_population(track_summary):
    """Return the one canonical 3D population used for biological summaries."""
    if track_summary is None or track_summary.empty:
        return pd.DataFrame()
    if "technical_valid" in track_summary.columns:
        return track_summary[_study_series_bool(track_summary["technical_valid"])].copy()
    return track_summary.copy()


def _unet_detection_accounting(detections):
    """Return one consistent U-Net provenance contract for every run mode."""
    frame = detections if detections is not None else pd.DataFrame()
    source = (
        frame["detection_source"].fillna("saturn_classical").astype(str)
        if not frame.empty and "detection_source" in frame.columns
        else pd.Series(["saturn_classical"] * len(frame), index=frame.index, dtype=str)
    )
    unet = source.str.startswith("unet_rescued")
    split = source == "unet_rescued_split"
    short = source == "unet_rescued_short_high_confidence"
    low_ratio = source == "unet_rescued_low_ratio_high_confidence"
    direct = source == "unet_rescued"
    known = direct | split | short | low_ratio
    probability_supported = pd.Series(False, index=frame.index)
    median_probability = np.nan
    if not frame.empty and "unet_mean_probability" in frame.columns:
        probability = pd.to_numeric(frame["unet_mean_probability"], errors="coerce")
        probability_supported = probability.notna() & np.isfinite(probability)
        accepted_probability = probability[unet].dropna()
        if not accepted_probability.empty:
            median_probability = float(accepted_probability.median())
    return {
        "saturn_classical_2d_count": int((source == "saturn_classical").sum()),
        "unet_rescued_2d_count": int(unet.sum()),
        "unet_rescued_direct_2d_count": int(direct.sum()),
        "unet_rescued_split_2d_count": int(split.sum()),
        "unet_rescued_short_high_confidence_2d_count": int(short.sum()),
        "unet_rescued_low_ratio_high_confidence_2d_count": int(low_ratio.sum()),
        "unet_rescued_other_2d_count": int((unet & ~known).sum()),
        "unet_probability_supported_2d_count": int(probability_supported.sum()),
        "median_probability_of_unet_rescues": median_probability,
        "unet_rescue_fraction_of_2d_detections": float(unet.sum() / max(len(frame), 1)),
    }


def build_analysis_summary(
    df=None,
    track_summary=None,
    run_scope="full_stack_3d",
    z_index=None,
    cfg=None,
):
    """Build one concise result contract shared by CLI, GUI batch, and slice preview."""
    detections = df.copy() if df is not None else pd.DataFrame()
    tracks = track_summary.copy() if track_summary is not None else pd.DataFrame()
    unet_accounting = _unet_detection_accounting(detections)
    engine = str((cfg or {}).get("SEGMENTATION_ENGINE", "unknown"))
    checkpoint = str((cfg or {}).get("UNET_MODEL_PATH", ""))

    def median(frame, column):
        if frame.empty or column not in frame.columns:
            return np.nan
        values = pd.to_numeric(frame[column], errors="coerce").dropna()
        return float(values.median()) if not values.empty else np.nan

    if run_scope == "single_slice_preview":
        return {
            "run_scope": "single_slice_preview",
            "analysis_population": "2D candidate detections in one preview slice",
            "biological_count_available": False,
            "estimated_unique_nuclei": np.nan,
            "z_index": z_index,
            "candidate_2d_detection_count": int(len(detections)),
            "median_2d_length_um": median(detections, "length_um_geodesic"),
            "median_2d_width_um": median(detections, "width_um"),
            "median_2d_length_width_ratio": median(
                detections,
                "length_width_ratio",
            ),
            "segmentation_engine": engine,
            "unet_checkpoint": checkpoint,
            **unet_accounting,
            "interpretation": (
                "Preview candidates are not unique nuclei. Run the complete stack "
                "with 3D tracking for biological counts and 3D morphology."
            ),
        }

    tracking_completed = run_scope == "full_stack_3d"
    primary = _technical_valid_track_population(tracks)
    morphology_warning_count = 0
    if not primary.empty and "morphology_warning" in primary.columns:
        morphology_warning_count = int(
            _study_series_bool(primary["morphology_warning"]).sum()
        )
    return {
        "run_scope": "full_stack_3d" if tracking_completed else "stack_2d_only",
        "analysis_population": (
            "included estimated nuclei"
            if tracking_completed
            else "2D candidate detections; no unique-nucleus estimate"
        ),
        "biological_count_available": bool(tracking_completed),
        "estimated_unique_nuclei": int(len(primary)) if tracking_completed else np.nan,
        "median_3d_length_um": median(primary, "total_3d_length_um"),
        "median_maximum_2d_length_um": median(primary, "max_length_2d"),
        "median_2d_width_um": median(primary, "median_width_2d"),
        "median_2d_length_width_ratio": median(
            primary,
            "median_length_width_ratio_2d",
        ),
        "median_effective_thickness_um_psf_sensitive": median(primary, "thickness_um"),
        "median_3d_tortuosity": median(primary, "tortuosity_3d"),
        "median_z_span_um": median(primary, "z_span_um"),
        "raw_2d_detection_count_qc": int(len(detections)),
        "reconstructed_track_count_qc": int(len(tracks)),
        "technical_failure_track_count_qc": int(len(tracks) - len(primary)),
        "morphology_review_note_count_qc": morphology_warning_count,
        "segmentation_engine": engine,
        "unet_checkpoint": checkpoint,
        **unet_accounting,
        "interpretation": (
            "Use estimated_unique_nuclei and the technical-valid morphology medians "
            "for specimen summaries. Fields ending in _qc are technical diagnostics, "
            "not alternative biological populations."
            if tracking_completed
            else "Run with 3D tracking before interpreting a biological nucleus count."
        ),
    }


def export_analysis_summary(
    out_dir,
    df=None,
    track_summary=None,
    run_scope="full_stack_3d",
    z_index=None,
    cfg=None,
):
    """Write the canonical concise summary used by every execution path."""
    ensure_dir(out_dir)
    summary = build_analysis_summary(
        df=df,
        track_summary=track_summary,
        run_scope=run_scope,
        z_index=z_index,
        cfg=cfg,
    )
    probability_map_count = 0
    if os.path.isdir(out_dir):
        probability_map_count = len(
            [
                name
                for name in os.listdir(out_dir)
                if name.endswith("_unet_probability.tif")
            ]
        )
    summary["saved_unet_probability_map_count"] = int(probability_map_count)

    primary_keys = [
        "run_scope",
        "analysis_population",
        "biological_count_available",
        "estimated_unique_nuclei",
    ]
    if summary.get("run_scope") == "single_slice_preview":
        primary_keys.extend(
            [
                "z_index",
                "candidate_2d_detection_count",
                "median_2d_length_um",
                "median_2d_width_um",
                "median_2d_length_width_ratio",
            ]
        )
    else:
        primary_keys.extend(
            [
                "median_3d_length_um",
                "median_maximum_2d_length_um",
                "median_2d_width_um",
                "median_2d_length_width_ratio",
                "median_effective_thickness_um_psf_sensitive",
                "median_3d_tortuosity",
                "median_z_span_um",
            ]
        )
    primary_keys.append("interpretation")
    primary_summary = {
        key: summary.get(key)
        for key in primary_keys
        if key in summary
    }

    csv_path = os.path.join(out_dir, "analysis_summary.csv")
    json_path = os.path.join(out_dir, "analysis_summary.json")
    pd.DataFrame([primary_summary]).to_csv(csv_path, index=False)
    pd.DataFrame([summary]).to_csv(
        os.path.join(out_dir, "technical_qc_summary.csv"),
        index=False,
    )

    def json_safe(payload):
        cleaned = {}
        for key, value in payload.items():
            if isinstance(value, (np.integer, np.floating)):
                value = value.item()
            if isinstance(value, float) and not np.isfinite(value):
                value = None
            cleaned[key] = value
        return cleaned

    json_summary = json_safe(primary_summary)
    technical_json_summary = json_safe(summary)
    with open(json_path, "w", encoding="utf-8") as handle:
        json.dump(json_summary, handle, indent=2)
    with open(
        os.path.join(out_dir, "technical_qc_summary.json"),
        "w",
        encoding="utf-8",
    ) as handle:
        json.dump(technical_json_summary, handle, indent=2)
    return summary


def summarize_comparative_population(df):
    """Return population-level summary metrics for sensitivity and blinded reports."""
    if df is None or df.empty:
        return {
            "total_technical_valid_count": 0,
            "technical_failure_fraction": 0.0,
            "morphology_warning_fraction": 0.0,
        }
    technical_valid = df["technical_valid"].astype(bool) if "technical_valid" in df.columns else pd.Series(True, index=df.index)
    morphology_warning = df["morphology_warning"].astype(bool) if "morphology_warning" in df.columns else pd.Series(False, index=df.index)
    valid = df[technical_valid].copy()
    denom = max(len(df), 1)
    out = {
        "total_technical_valid_count": int(len(valid)),
        "technical_failure_fraction": float((~technical_valid).sum() / denom),
        "morphology_warning_fraction": float(morphology_warning.sum() / denom),
    }
    metric_cols = {
        "total_3d_length_um": "length",
        "thickness_um": "width",
        "taper_ratio": "taper",
        "tortuosity_3d": "tortuosity",
        "volume_um3": "volume",
        "z_span_um": "z_span",
        "pitch_deg": "pitch",
        "nearest_neighbor_um": "nearest_neighbor_distance",
    }
    for col, label in metric_cols.items():
        if col in valid.columns and not valid.empty:
            vals = pd.to_numeric(valid[col], errors="coerce").dropna()
            out[f"median_{label}"] = float(vals.median()) if not vals.empty else np.nan
            out[f"mean_{label}"] = float(vals.mean()) if not vals.empty else np.nan
        else:
            out[f"median_{label}"] = np.nan
            out[f"mean_{label}"] = np.nan
    for flag_col, label in [
        ("technical_failure_reasons", "technical_failure"),
        ("morphology_warning_reasons", "morphology_warning"),
    ]:
        if flag_col in df.columns:
            reasons = df[flag_col].fillna("").astype(str)
            for token in sorted(set(",".join(reasons).split(",")) - {""}):
                out[f"{label}_{token}_fraction"] = float(reasons.map(lambda s, t=token: t in [p for p in s.split(",") if p]).sum() / denom)
    return out


def compare_preset_track_summaries(preset_tables):
    """
    Compare already-generated track summaries from comparative presets.

    ``preset_tables`` maps preset names to annotated track-summary DataFrames.
    Object identity is approximated by rounded centroid and z-span fields when
    available; this function never asserts which biological distribution is
    correct.
    """
    summaries = []
    identity_sets = {}
    for name, df in preset_tables.items():
        rec = {"preset": name, **summarize_comparative_population(df)}
        summaries.append(rec)
        if df is None or df.empty:
            identity_sets[name] = set()
            continue
        cols = [c for c in ("centroid_x", "centroid_y", "z_min", "z_max") if c in df.columns]
        if len(cols) >= 2:
            technical_valid = df["technical_valid"].astype(bool) if "technical_valid" in df.columns else pd.Series(True, index=df.index)
            valid = df[technical_valid].copy()
            keys = set(tuple(round(float(row[c]), 1) for c in cols) for _, row in valid.iterrows())
        elif "track_id" in df.columns:
            technical_valid = df["technical_valid"].astype(bool) if "technical_valid" in df.columns else pd.Series(True, index=df.index)
            keys = set(df.loc[technical_valid, "track_id"].astype(str))
        else:
            keys = set(df.index.astype(str))
        identity_sets[name] = keys
    all_names = list(preset_tables.keys())
    shared = set.intersection(*(identity_sets[n] for n in all_names)) if all_names else set()
    permissive = identity_sets.get("permissive", set())
    conservative = identity_sets.get("conservative", set())
    selected = identity_sets.get("selected", set())
    sensitivity = {
        "detected_by_all_presets": len(shared),
        "detected_only_by_permissive": len(permissive - set.union(*(identity_sets[n] for n in all_names if n != "permissive"))) if permissive and len(all_names) > 1 else 0,
        "lost_only_by_conservative": len((selected | permissive | identity_sets.get("intermediate", set())) - conservative),
        "preset_count": len(all_names),
    }
    return pd.DataFrame(summaries), sensitivity


def assign_blinded_dataset_ids(manifest_df, seed=560123):
    """
    Return a manifest copy with anonymized dataset IDs.

    Genotype labels are preserved only in a separate reveal table. Segmentation
    code should consume the blinded manifest and not the reveal table.
    """
    rng = np.random.default_rng(seed)
    manifest = manifest_df.copy()
    order = np.arange(len(manifest))
    rng.shuffle(order)
    blinded_ids = [f"anon_{i + 1:03d}" for i in range(len(manifest))]
    manifest["blinded_dataset_id"] = ""
    for blinded_id, idx in zip(blinded_ids, order):
        manifest.loc[manifest.index[idx], "blinded_dataset_id"] = blinded_id
    reveal_cols = [c for c in ("blinded_dataset_id", "genotype", "group", "dataset_label") if c in manifest.columns]
    reveal = manifest[reveal_cols].copy()
    blinded = manifest.drop(columns=[c for c in ("genotype", "group", "dataset_label") if c in manifest.columns])
    return blinded, reveal


def make_blinded_review_sheet(crop_records, candidates=None):
    """Create blank manual-review rows without exposing genotype labels."""
    candidates = candidates or ["comparative_selected"]
    rows = []
    manual_cols = [
        "true_detection", "missed_nucleus", "split_nucleus", "merged_nuclei",
        "tissue_edge_false_positive", "puncta_ring_false_positive", "uncertain",
        "reviewer_notes",
    ]
    for crop in crop_records:
        for cand in candidates:
            row = {
                "blinded_dataset_id": crop.get("blinded_dataset_id", ""),
                "crop_id": crop.get("crop_id", ""),
                "candidate": cand,
                "z_index": crop.get("z_index", ""),
                "x0": crop.get("x0", ""),
                "y0": crop.get("y0", ""),
                "x1": crop.get("x1", ""),
                "y1": crop.get("y1", ""),
            }
            row.update({col: "" for col in manual_cols})
            rows.append(row)
    return pd.DataFrame(rows)


def differential_error_indicators(summary_by_group, warning_threshold=0.15):
    """
    Compare segmentation error indicators between anonymized groups.

    The function reports differences only; it does not correct or normalize
    biological distributions to make groups agree.
    """
    df = pd.DataFrame(summary_by_group).copy()
    if df.empty or "group" not in df.columns:
        return df, []
    indicators = [
        "technical_failure_fraction", "morphology_warning_fraction",
        "short_fragment_fraction", "suspected_merge_fraction",
        "branch_network_fraction", "roi_edge_fraction",
        "permissive_only_detection_fraction", "conservative_loss_fraction",
    ]
    warnings = []
    for col in indicators:
        if col not in df.columns:
            continue
        vals = pd.to_numeric(df[col], errors="coerce")
        if vals.notna().sum() >= 2 and float(vals.max() - vals.min()) >= warning_threshold:
            warnings.append(f"{col} differs by {float(vals.max() - vals.min()):.3f} across anonymized groups")
    return df, warnings


def export_outlier_audit(out_dir, df_tracks, cfg):
    """
    Generate the outlier_audit/ subfolder with per-class CSVs and a summary,
    mirroring the output of audit_sperm_outliers.py but done automatically.
    """
    if df_tracks.empty:
        return

    audit_dir = os.path.join(out_dir, "outlier_audit")
    ensure_dir(audit_dir)

    # Robust column matching
    length_col = next((c for c in ["total_3d_length_um", "length_3d_um_est", "max_length_um"] if c in df_tracks.columns), None)
    tort_col = next((c for c in ["tortuosity_3d", "tortuosity"] if c in df_tracks.columns), None)
    thick_col = next((c for c in ["thickness_um", "effective_thickness_um", "median_width_um"] if c in df_tracks.columns), None)
    taper_col = next((c for c in ["taper_ratio", "morphological_taper_ratio"] if c in df_tracks.columns), None)
    nslices_col = next((c for c in ["n_slices", "n_detections"] if c in df_tracks.columns), None)

    # Apply thresholds
    thresh_len = cfg.get("AUDIT_MAX_LENGTH_UM", 15.0)
    thresh_tort = cfg.get("AUDIT_MAX_TORTUOSITY", 1.5)
    thresh_thick = cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0)
    thresh_taper = cfg.get("AUDIT_MAX_TAPER_RATIO", 1.5)
    thresh_slices = cfg.get("AUDIT_MIN_SLICES", 2)

    outlier_sets = {}

    if length_col:
        long_df = df_tracks[df_tracks[length_col] > thresh_len].sort_values(length_col, ascending=False)
        long_df.to_csv(os.path.join(audit_dir, "outliers_long.csv"), index=False)
        outlier_sets["long"] = len(long_df)

    if tort_col:
        tort_df = df_tracks[df_tracks[tort_col] > thresh_tort].sort_values(tort_col, ascending=False)
        tort_df.to_csv(os.path.join(audit_dir, "outliers_tortuous.csv"), index=False)
        outlier_sets["tortuous"] = len(tort_df)

    if thick_col:
        thick_df = df_tracks[df_tracks[thick_col] > thresh_thick].sort_values(thick_col, ascending=False)
        thick_df.to_csv(os.path.join(audit_dir, "outliers_thick.csv"), index=False)
        outlier_sets["thick"] = len(thick_df)

    if taper_col:
        taper_df = df_tracks[df_tracks[taper_col] > thresh_taper].sort_values(taper_col, ascending=False)
        taper_df.to_csv(os.path.join(audit_dir, "outliers_taper.csv"), index=False)
        outlier_sets["taper"] = len(taper_df)

    if nslices_col:
        single_df = df_tracks[df_tracks[nslices_col] < thresh_slices].sort_values(nslices_col, ascending=True)
        single_df.to_csv(os.path.join(audit_dir, "outliers_single_slice.csv"), index=False)
        outlier_sets["single_slice"] = len(single_df)

    if "technical_valid" in df_tracks.columns:
        tech_fail_df = df_tracks[~df_tracks["technical_valid"]].copy()
        morph_warn_df = df_tracks[df_tracks["morphology_warning"]].copy() if "morphology_warning" in df_tracks.columns else pd.DataFrame()
        ref_df = df_tracks[df_tracks["reference_morphology_pass"]].copy() if "reference_morphology_pass" in df_tracks.columns else pd.DataFrame()
        if not tech_fail_df.empty:
            tech_fail_df.to_csv(os.path.join(audit_dir, "technical_failures.csv"), index=False)
        morph_warn_df.to_csv(os.path.join(audit_dir, "morphology_warnings.csv"), index=False)
        ref_df.to_csv(os.path.join(audit_dir, "reference_morphology.csv"), index=False)
        outlier_sets["technical_failures"] = len(tech_fail_df)
        outlier_sets["morphology_warnings"] = len(morph_warn_df)
        outlier_sets["reference_morphology"] = len(ref_df)

    # All flagged
    if "is_quality_track" in df_tracks.columns:
        flagged = df_tracks[~df_tracks["is_quality_track"]].copy()
    else:
        flagged = pd.DataFrame()

    if not flagged.empty:
        flagged.to_csv(os.path.join(audit_dir, "outliers_all_flagged.csv"), index=False)

    # Summary
    lines = []
    lines.append("OUTLIER AUDIT SUMMARY (Auto-generated)")
    lines.append("=" * 60)
    lines.append(f"Tracks total: {len(df_tracks)}")
    if "is_quality_track" in df_tracks.columns:
        lines.append(f"Reference-morphology diagnostic tracks: {int(df_tracks['is_quality_track'].sum())}")
        lines.append(f"Outside reference-morphology diagnostic: {int((~df_tracks['is_quality_track']).sum())}")
    if "has_warning_only" in df_tracks.columns:
        lines.append(f"Warning-only tracks: {int(df_tracks['has_warning_only'].sum())}")
    if "technical_valid" in df_tracks.columns:
        lines.append(f"Primary estimated nuclei: {int(df_tracks['technical_valid'].sum())}")
        lines.append(f"Technical-failure tracks: {int((~df_tracks['technical_valid']).sum())}")
    if "morphology_warning" in df_tracks.columns:
        lines.append(f"Morphology-warning tracks retained in comparative population: {int(df_tracks['morphology_warning'].sum())}")
    if "reference_morphology_pass" in df_tracks.columns:
        lines.append(f"Reference-morphology diagnostic subset: {int(df_tracks['reference_morphology_pass'].sum())}")
    lines.append("Morphology warnings are retained in the comparative population because they may represent genuine genotype-dependent phenotypes.")
    lines.append("")
    lines.append("Thresholds:")
    lines.append(f"  length > {thresh_len}")
    lines.append(f"  tortuosity > {thresh_tort}")
    lines.append(f"  thickness > {thresh_thick}")
    lines.append(f"  taper > {thresh_taper}")
    lines.append(f"  extreme thickness > {cfg.get('AUDIT_EXTREME_THICKNESS_UM', 3.5)}")
    lines.append(f"  extreme taper > {cfg.get('AUDIT_EXTREME_TAPER_RATIO', 3.0)}")
    lines.append(f"  n_slices < {thresh_slices}")
    lines.append("")
    lines.append("Counts:")
    for name, count in outlier_sets.items():
        lines.append(f"  {name:20s}: {count}")
    lines.append(f"  {'all_flagged':20s}: {len(flagged)}")

    with open(os.path.join(audit_dir, "outlier_summary.txt"), "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"  Outlier audit exported to: {audit_dir}")


def export_post_detection_qc(out_dir, df_detections, df_tracks):
    """
    Save compact diagnostics for the post-detection 3D processing stage.

    These metrics make GUI and CLI 3D runs easier to audit: how detections were
    linked, how many tracks are single-slice, and which quality flags dominate.
    """
    try:
        ensure_dir(out_dir)
        lines = []
        lines.append("POST-DETECTION 3D QC")
        lines.append("=" * 72)

        n_det = int(len(df_detections)) if df_detections is not None else 0
        n_tracks = int(len(df_tracks)) if df_tracks is not None else 0
        lines.append(f"2D detections: {n_det}")
        lines.append(f"3D tracks: {n_tracks}")

        if df_detections is not None and not df_detections.empty and "track_link_method" in df_detections.columns:
            method_counts = df_detections["track_link_method"].fillna("unknown").value_counts()
            lines.append("\nDetection link methods:")
            for method, count in method_counts.items():
                lines.append(f"  {method}: {int(count)}")
            if "track_link_distance_um" in df_detections.columns:
                linked_dist = pd.to_numeric(df_detections["track_link_distance_um"], errors="coerce").dropna()
                if not linked_dist.empty:
                    lines.append(f"Median link distance um: {float(linked_dist.median()):.3f}")
                    lines.append(f"95th percentile link distance um: {float(linked_dist.quantile(0.95)):.3f}")

        if df_tracks is not None and not df_tracks.empty:
            primary = _technical_valid_track_population(df_tracks)
            technical_valid_count = int(len(primary))
            lines.append(
                f"\nPRIMARY BIOLOGICAL POPULATION (technical-valid 3D tracks): "
                f"{technical_valid_count}"
            )
            if "n_slices" in primary.columns and not primary.empty:
                single_frac = float((primary["n_slices"] <= 1).mean() * 100)
                lines.append(f"\nSingle-slice tracks: {single_frac:.1f}%")
                lines.append(f"Median n_slices: {float(primary['n_slices'].median()):.2f}")
            if "total_3d_length_um" in primary.columns and not primary.empty:
                lines.append(f"Median 3D length um: {float(primary['total_3d_length_um'].median()):.3f}")
            if "z_span_um" in primary.columns and not primary.empty:
                lines.append(f"Median Z-span um: {float(primary['z_span_um'].median()):.3f}")
            if "technical_valid" in df_tracks.columns:
                lines.append("\nFINAL ANALYSIS POPULATION:")
                lines.append(f"  estimated_unique_nuclei: {technical_valid_count}")
                lines.append(f"  technical_failures: {int(len(df_tracks) - technical_valid_count)}")
            if "morphology_warning" in primary.columns:
                lines.append(
                    "  nuclei_with_morphology_review_note: "
                    f"{int(_study_series_bool(primary['morphology_warning']).sum())}"
                )
                lines.append(
                    "  interpretation: review notes do not remove nuclei and do "
                    "not define another biological population"
                )
            lines.append(
                "\nDetailed segmentation-source and rejected-link diagnostics "
                "remain in the master audit tables; they are not biological counts."
            )

        qc_path = os.path.join(out_dir, "post_detection_qc.txt")
        with open(qc_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines) + "\n")
        print(f"  Post-detection QC exported to: {qc_path}")
    except Exception as e:
        print(f"  WARNING: Could not export post-detection QC: {e}")


def get_audit_flag_counts(df_tracks, cfg=None):
    """Return consistent audit category counts for reports, regardless of GUI/report settings."""
    counts = {
        "long": 0, "tortuous": 0, "thick": 0, "taper": 0, "single_slice": 0,
        "all_flagged": 0, "quality": 0, "shape_outliers": 0,
        "estimated_nuclei": 0, "technical_failures": 0, "warning_only": 0,
    }
    if df_tracks is None or df_tracks.empty:
        return counts

    if "quality_flags" in df_tracks.columns:
        flags = df_tracks["quality_flags"].fillna("").astype(str)
        split_flags = flags.str.split(',')
        token_map = {
            "long": "long",
            "tortuous": "high_tortuosity",
            "thick": "wide",
            "taper": "high_taper",
            "single_slice": "single_slice",
        }
        for key, token in token_map.items():
            counts[key] = int(
                split_flags.apply(
                    lambda items, expected=token: expected in items
                    if isinstance(items, list)
                    else False
                ).sum()
            )
    elif cfg is not None:
        if "total_3d_length_um" in df_tracks.columns:
            counts["long"] = int((df_tracks["total_3d_length_um"] > cfg.get("AUDIT_MAX_LENGTH_UM", 15.0)).sum())
        if "tortuosity_3d" in df_tracks.columns:
            counts["tortuous"] = int((df_tracks["tortuosity_3d"] > cfg.get("AUDIT_MAX_TORTUOSITY", 1.5)).sum())
        if "thickness_um" in df_tracks.columns:
            counts["thick"] = int((df_tracks["thickness_um"] > cfg.get("AUDIT_MAX_THICKNESS_UM", 2.0)).sum())
        if "taper_ratio" in df_tracks.columns:
            counts["taper"] = int((df_tracks["taper_ratio"] > cfg.get("AUDIT_MAX_TAPER_RATIO", 1.5)).sum())
        if "n_slices" in df_tracks.columns:
            counts["single_slice"] = int((df_tracks["n_slices"] < cfg.get("AUDIT_MIN_SLICES", 2)).sum())

    if "is_quality_track" in df_tracks.columns:
        counts["quality"] = int(df_tracks["is_quality_track"].sum())
        counts["all_flagged"] = int((~df_tracks["is_quality_track"]).sum())
    else:
        counts["quality"] = len(df_tracks)
        counts["all_flagged"] = 0

    counts["shape_outliers"] = max(counts["all_flagged"] - counts["single_slice"], 0)
    if "technical_valid" in df_tracks.columns:
        counts["estimated_nuclei"] = int(df_tracks["technical_valid"].sum())
        counts["technical_failures"] = int((~df_tracks["technical_valid"]).sum())
    else:
        counts["estimated_nuclei"] = len(df_tracks)
        counts["technical_failures"] = 0
    if "has_warning_only" in df_tracks.columns:
        counts["warning_only"] = int(df_tracks["has_warning_only"].sum())
    return counts


# =============================================================================
# PROCESS ONE IMAGE
# =============================================================================

def process_one_image(image_path, cfg, output_dir):
    """
    Runs the complete segmentation-and-measurement pipeline on a single Z-slice image.

    This function is the fundamental unit of the batch processing loop.  For each
    input image it:

    1. Loads the raw image (multi-format with TIFF / OpenCV fallback).
    2. Selects the correct Z-slice from a multi-plane volume if needed.
    3. Calls :func:`segment_slice` to produce the binary mask, ridge map, and skeleton.
    4. Calls :func:`measure_spermatids` to extract per-cell biometrics.
    5. Optionally validates or post-filters results with shape-quality checks.
    6. Saves the colour overlay PNG, detail figure, mask TIFs, and per-slice CSV.
    7. Optionally opens the preview window for interactive QC.

    Args:
        image_path (str): Absolute path to the input image file.
        cfg (dict): Full pipeline configuration dictionary.  Key fields:
            - ``UM_PER_PX_XY`` - physical scale factor.
            - ``Z_INDEX`` - which plane to extract from a Z-stack.
            - ``OUTPUT_DIR`` - top-level output folder.
            - ``SHOW_PREVIEW_WINDOW`` - whether to open the preview GUI.
            - ``SAVE_MASK_TIFS``, ``SAVE_LABEL_TIFS`` - optional TIF outputs.
            - ``SAVE_DETAIL_FIGURE`` - whether to save 3-panel figure.
        output_dir (str): Destination directory for this image's outputs.

    Returns:
        tuple[list[dict], dict]:
            - ``results`` - list of per-spermatid measurement dicts.
            - ``seg``     - full segmentation dict from :func:`segment_slice`
              including ``mask_clean``, ``skel_pruned``, ``skel_label``, etc.
    """
    ensure_dir(output_dir)
    overlay_dir = os.path.join(output_dir, "overlays")
    debug_dir   = os.path.join(output_dir, "debug")
    rescue_review_dir = os.path.join(
        output_dir,
        "technical_qc",
        "unet_rescue_overlays",
    )
    ensure_dir(overlay_dir)
    ensure_dir(rescue_review_dir)
    if cfg["SAVE_DEBUG_IMAGES"]:
        ensure_dir(debug_dir)

    z_idx   = extract_z_index(image_path)
    img_raw = robust_imread(image_path)
    img_2d  = ensure_2d_image(img_raw, os.path.basename(image_path))
    roi_mask = None
    exclusion_mask = None
    if cfg.get("ROI_MASK_PATH"):
        roi_mask = load_roi_mask_file(cfg["ROI_MASK_PATH"], expected_shape=img_2d.shape)
    if cfg.get("EXCLUSION_MASK_PATH"):
        exclusion_mask = load_roi_mask_file(cfg["EXCLUSION_MASK_PATH"], expected_shape=img_2d.shape)
    preprocess_context = build_stack_preprocess_context([image_path], roi_mask, cfg, exclusion_mask=exclusion_mask)
    save_stack_preprocess_context(preprocess_context, output_dir)
    print(f"\nProcessing: {os.path.basename(image_path)}")

    t0      = time.time()
    unet_context = np.stack([img_2d.astype(np.float32)] * 3, axis=0)
    seg     = segment_slice(img_2d, cfg, z_idx=z_idx,
                            debug_dir=debug_dir if cfg["SAVE_DEBUG_IMAGES"] else None,
                            roi_mask=roi_mask,
                            preprocess_context=preprocess_context,
                            exclusion_mask=exclusion_mask,
                            unet_context_stack=unet_context)
    meas    = measure_spermatids(seg, cfg)
    results = meas["results"]
    elapsed = time.time() - t0

    um = cfg["UM_PER_PX_XY"]
    print(f"  2D preview candidates: {len(results)}  ({elapsed:.1f}s)")
    if results:
        ls = [r["length_px_geodesic"]*um for r in results]
        print(f"  Geodesic length um: median={np.median(ls):.2f}  max={max(ls):.2f}")

    overlay_rgb = make_overlay(img_raw, meas["skel_label"])

    if cfg["SAVE_OVERLAYS"]:
        _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_overlay.png"), overlay_rgb)
        rescue_rgb = make_unet_rescue_review_overlay(
            img_raw,
            meas["skel_label"],
            results,
            meas.get("unet_rescue_rejected_reason"),
        )
        _imwrite(
            os.path.join(
                rescue_review_dir,
                f"z{z_idx:02d}_unet_rescue_review.png",
            ),
            rescue_rgb,
        )
    if cfg["SAVE_DETAIL_FIGURE"]:
        save_detail_figure(img_raw, overlay_rgb, results,
                           os.path.join(overlay_dir, f"z{z_idx:02d}_detail.png"),
                           z_idx, um)
    if cfg["SAVE_MASK_TIFS"]:
        tifffile.imwrite(os.path.join(output_dir, f"z{z_idx:02d}_mask.tif"),
                         seg["mask_clean"].astype(np.uint8) * 255)
    if cfg.get("UNET_SAVE_PROBABILITY_MAPS", True):
        if seg.get("unet_probability") is not None and np.any(seg.get("unet_probability")):
            tifffile.imwrite(os.path.join(output_dir, f"z{z_idx:02d}_unet_probability.tif"),
                             seg["unet_probability"].astype(np.float32))
    if cfg["SAVE_LABEL_TIFS"]:
        tifffile.imwrite(os.path.join(output_dir, f"z{z_idx:02d}_skel_labels.tif"),
                         meas["skel_label"].astype(np.uint16))

    result_frame = pd.DataFrame(rows_from_results(results, z_idx, um))
    result_frame.to_csv(
        os.path.join(output_dir, f"single_measurements_{_VERSION}.csv"), index=False)
    analysis_summary = export_analysis_summary(
        output_dir,
        df=result_frame,
        run_scope="single_slice_preview",
        z_index=z_idx,
        cfg=cfg,
    )
    print(
        "  Detection sources: "
        f"{analysis_summary['saturn_classical_2d_count']} Saturn classical, "
        f"{analysis_summary['unet_rescued_2d_count']} U-Net rescued"
    )

    if cfg["SHOW_PREVIEW_WINDOW"]:
        show_single_preview(img_raw, seg, overlay_rgb, results, z_idx, cfg)

    print("  Preview only: these are not estimated unique nuclei.")
    print(f"Saved to: {output_dir}")
    return results, seg


# =============================================================================
# PROCESS BATCH
# =============================================================================

def process_batch(cfg):
    """
    Orchestrates the entire end-to-end biological data processing engine for batch image iterations.

    Iterates over all `.tif`/`.tiff` files across defined spatial bounds, extracting single
    segmentation matrices and triggering the 3D concatenation models upon directory exhaustion.

    Args:
        cfg (dict): Active session configuration parameter dictionary.
    """
    ensure_dir(cfg["OUTPUT_DIR"])
    overlay_dir = os.path.join(cfg["OUTPUT_DIR"], "overlays")
    debug_dir   = os.path.join(cfg["OUTPUT_DIR"], "debug")
    rescue_review_dir = os.path.join(
        cfg["OUTPUT_DIR"],
        "technical_qc",
        "unet_rescue_overlays",
    )
    ensure_dir(overlay_dir)
    ensure_dir(rescue_review_dir)
    if cfg["SAVE_DEBUG_IMAGES"]:
        ensure_dir(debug_dir)

    files, z_indices = load_batch_files(cfg["INPUT_DIR"], cfg["FILE_PATTERN"])
    files_by_z = {int(z): f for f, z in zip(files, z_indices)}
    um         = cfg["UM_PER_PX_XY"]
    roi_mask = None
    exclusion_mask = None
    first_img = ensure_2d_image(robust_imread(files[0]), os.path.basename(files[0]))
    if cfg.get("ROI_MASK_PATH"):
        roi_mask = load_roi_mask_file(cfg["ROI_MASK_PATH"], expected_shape=first_img.shape)
        roi_out = os.path.join(cfg["OUTPUT_DIR"], "roi_mask_used.tif")
        tifffile.imwrite(roi_out, roi_mask.astype(np.uint8) * 255)
        print(f"Using ROI mask: {cfg['ROI_MASK_PATH']}")
        print(f"Saved ROI copy: {roi_out}")
    if cfg.get("EXCLUSION_MASK_PATH"):
        exclusion_mask = load_roi_mask_file(cfg["EXCLUSION_MASK_PATH"], expected_shape=first_img.shape)
        excl_out = os.path.join(cfg["OUTPUT_DIR"], "exclusion_mask_used.tif")
        tifffile.imwrite(excl_out, exclusion_mask.astype(np.uint8) * 255)
        print(f"Using exclusion mask: {cfg['EXCLUSION_MASK_PATH']}")
        print(f"Saved exclusion copy: {excl_out}")
    preprocess_context = build_stack_preprocess_context(files, roi_mask, cfg, exclusion_mask=exclusion_mask)
    save_stack_preprocess_context(preprocess_context, cfg["OUTPUT_DIR"])
    cfg["RESOLVED_PREPROCESSING_PROFILE"] = preprocess_context.selected_clahe_profile
    print(f"Stack preprocessing: profile={preprocess_context.selected_clahe_profile}, clip={preprocess_context.selected_clahe_clip}, norm=({preprocess_context.normalization_low:.3f}, {preprocess_context.normalization_high:.3f}), sampled_z={preprocess_context.sampled_z_indices}")
    all_rows   = []
    summaries  = []
    t_batch    = time.time()
    t_slices   = []
    max_proj_raw = None
    max_proj_ov = None
    slice_cache = {}

    for idx_i, (fpath, z_idx) in enumerate(zip(files, z_indices)):
        t0 = time.time()
        print(f"\n[{idx_i+1}/{len(files)}]  Z={z_idx:02d}  {os.path.basename(fpath)}")
        img_raw = robust_imread(fpath)
        img_2d  = ensure_2d_image(img_raw, os.path.basename(fpath))
        unet_context = _make_unet_context_from_paths(files_by_z, z_idx)
        seg     = segment_slice(img_2d, cfg, z_idx=z_idx,
                                debug_dir=debug_dir if cfg["SAVE_DEBUG_IMAGES"] else None,
                                roi_mask=roi_mask,
                                preprocess_context=preprocess_context,
                                exclusion_mask=exclusion_mask,
                                unet_context_stack=unet_context)
        meas    = measure_spermatids(seg, cfg)
        results = meas["results"]
        skel_label = meas["skel_label"]
        ls_um   = [r["length_px_geodesic"]*um for r in results]
        ws_um   = [r["width_px"]*um for r in results]

        t_s = time.time() - t0
        t_slices.append(t_s)
        eta = (len(files) - idx_i - 1) * float(np.mean(t_slices))
        print(f"  N={len(results)}", end="")
        if ls_um:
            print(f"  med_len={np.median(ls_um):.2f}um", end="")
        print(f"  {t_s:.1f}s  ETA {eta:.0f}s")

        all_rows.extend(rows_from_results(results, z_idx, um))
        summaries.append({
            "z_slice":          z_idx,
            "n_spermatids":     len(results),
            "mean_length_um":   round(float(np.mean(ls_um)),   3) if ls_um else 0,
            "median_length_um": round(float(np.median(ls_um)), 3) if ls_um else 0,
            "mean_width_um":    round(float(np.mean(ws_um)),   3) if ws_um  else 0,
        })

        overlay_rgb = make_overlay(img_2d, skel_label)
        if cfg["SAVE_OVERLAYS"]:
            # Create side-by-side panel: [Original | Overlay]
            orig_rgb = (normalize_display(img_2d) * 255).astype(np.uint8)
            # if grayscale, convert to RGB for hstack
            if orig_rgb.ndim == 2:
                orig_rgb = np.stack([orig_rgb]*3, axis=-1)

            panel = np.hstack([orig_rgb, overlay_rgb])
            _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_panel.png"), panel)
            rescue_rgb = make_unet_rescue_review_overlay(
                img_2d,
                skel_label,
                results,
                meas.get("unet_rescue_rejected_reason"),
            )
            _imwrite(
                os.path.join(
                    rescue_review_dir,
                    f"z{z_idx:02d}_unet_rescue_review.png",
                ),
                rescue_rgb,
            )

            if max_proj_raw is None:
                max_proj_raw = img_2d.copy().astype(np.float32)
                max_proj_ov = overlay_rgb.copy().astype(np.float32)
            else:
                max_proj_raw = np.maximum(max_proj_raw, img_2d.astype(np.float32))
                max_proj_ov = np.maximum(max_proj_ov, overlay_rgb.astype(np.float32))
            if cfg.get("DO_TRACKING", True):
                slice_cache[int(z_idx)] = {
                    "image": img_2d.copy(),
                    "skel_label": skel_label.copy().astype(np.int32),
                }

        if cfg["SAVE_DETAIL_FIGURE"]:
            save_detail_figure(img_2d, overlay_rgb, results,
                               os.path.join(overlay_dir, f"z{z_idx:02d}_detail.png"),
                               z_idx, um)
        if cfg["SAVE_MASK_TIFS"]:
            tifffile.imwrite(os.path.join(cfg["OUTPUT_DIR"], f"z{z_idx:02d}_mask.tif"),
                             (seg["mask_clean"] & roi_mask if roi_mask is not None else seg["mask_clean"]).astype(np.uint8) * 255)
        if cfg.get("UNET_SAVE_PROBABILITY_MAPS", True):
            if seg.get("unet_probability") is not None and np.any(seg.get("unet_probability")):
                tifffile.imwrite(os.path.join(cfg["OUTPUT_DIR"], f"z{z_idx:02d}_unet_probability.tif"),
                                 seg["unet_probability"].astype(np.float32))
        if cfg["SAVE_LABEL_TIFS"]:
            tifffile.imwrite(os.path.join(cfg["OUTPUT_DIR"], f"z{z_idx:02d}_skel_labels.tif"),
                             skel_label.astype(np.uint16))

    df     = pd.DataFrame(all_rows)
    df_sum = pd.DataFrame(summaries)
    if cfg["SAVE_OVERLAYS"] and max_proj_raw is not None and max_proj_ov is not None:
        raw_p = (normalize_display(max_proj_raw.astype(np.uint16)) * 255).astype(np.uint8)
        if raw_p.ndim == 2:
            raw_p = np.stack([raw_p] * 3, axis=-1)
        ov_p = np.clip(max_proj_ov, 0, 255).astype(np.uint8)
        global_panel = np.hstack([raw_p, ov_p])
        _imwrite(os.path.join(cfg["OUTPUT_DIR"], "global_z_projection.png"), global_panel)

    df.to_csv(    os.path.join(cfg["OUTPUT_DIR"], f"spermatid_measurements_{_VERSION}.csv"), index=False)
    df_sum.to_csv(os.path.join(cfg["OUTPUT_DIR"], f"slice_summary_{_VERSION}.csv"), index=False)

    # Robust initialization for reporting
    df_trk = None
    ts = None

    if cfg["DO_TRACKING"] and not df.empty:
        df_trk, ts = track_across_slices(df, cfg)
        df_trk.to_csv(
            os.path.join(cfg["OUTPUT_DIR"], f"measurements_with_tracks_{_VERSION}.csv"),
            index=False)

        # Auto quality audit
        ts = flag_quality_tracks(ts, cfg)
        ts.to_csv(
            os.path.join(cfg["OUTPUT_DIR"], f"track_summary_{_VERSION}.csv"),
            index=False)

        export_comparative_track_tables(cfg["OUTPUT_DIR"], ts, _VERSION)
        export_biologist_results(cfg["OUTPUT_DIR"], ts, _VERSION)

        # Quality-coded overlays: green = audit-passed, red = audit-flagged.
        if cfg["SAVE_OVERLAYS"]:
            export_quality_overlays(cfg["OUTPUT_DIR"], slice_cache, df_trk, ts)
            export_analysis_overlays(
                cfg["OUTPUT_DIR"],
                slice_cache,
                df_trk,
                ts,
            )

        # Export outlier_audit/ folder
        export_outlier_audit(cfg["OUTPUT_DIR"], ts, cfg)
        export_post_detection_qc(cfg["OUTPUT_DIR"], df_trk, ts)

    analysis_summary = export_analysis_summary(
        cfg["OUTPUT_DIR"],
        df=df,
        track_summary=ts,
        run_scope="full_stack_3d" if cfg["DO_TRACKING"] else "stack_2d_only",
        cfg=cfg,
    )

    # --- Reporting Phase (CLI/Batch) ---
    print(f"\nGenerating final reports in {cfg['OUTPUT_DIR']}...")
    generate_batch_report(
        cfg["OUTPUT_DIR"], df, df_sum, um, ts, df_tracked=df_trk)
    generate_excel_report(cfg["OUTPUT_DIR"], df, df_sum, ts)

    total = time.time() - t_batch
    print(f"\n{'='*55}")
    print(f"{_VERSION} DONE | {len(files)} slices | {total:.1f}s")
    if analysis_summary["biological_count_available"]:
        print(
            "Primary result | estimated unique nuclei: "
            f"{analysis_summary['estimated_unique_nuclei']} | "
            "median 3D length: "
            f"{analysis_summary['median_3d_length_um']:.3f} um"
        )
    else:
        print("No biological nucleus count was produced because 3D tracking was unavailable.")
    print(f"Saved to: {cfg['OUTPUT_DIR']}")
    print("Concise result: analysis_summary.csv")
    print("Technical diagnostics: technical_qc_summary.csv")


def write_error_log(out_dir, component, message):
    """
    Writes a persistent error log to report_generation_errors.txt in the output directory.
    """
    try:
        import os as _os
        import time as _time
        from datetime import datetime as _dt
        log_path = _os.path.join(out_dir, "report_generation_errors.txt")
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"\n[{_dt.now().strftime('%Y-%m-%d %H:%M:%S')}] COMPONENT: {component}\n")
            f.write(f"MESSAGE:\n{message}\n")
            f.write("-" * 80 + "\n")
    except Exception:
        pass


def summarize_unet_rescue_for_reports(df, out_dir=None):
    """
    Return report-ready v5.7 U-Net rescue accounting from raw 2D detections.

    The summary is intentionally based on exported measurement columns, not
    overlay pixels. Display dilation in PNG overlays is cosmetic only.
    """
    if df is None or df.empty or "detection_source" not in df.columns:
        return None

    source = df["detection_source"].fillna("saturn_classical").astype(str)
    total = int(len(df))
    classical = int((source == "saturn_classical").sum())
    unet_mask = source.str.startswith("unet_rescued")
    unet_split_mask = source == "unet_rescued_split"
    unet_direct = int((unet_mask & ~unet_split_mask).sum())
    unet_split = int((source == "unet_rescued_split").sum())
    unet_total = int(unet_mask.sum())

    def pct(n):
        return 100.0 * float(n) / max(total, 1)

    rows = [
        {"Population": "All 2D detections", "Count": total, "Percent": 100.0, "Median U-Net probability": np.nan},
        {"Population": "Saturn classical", "Count": classical, "Percent": pct(classical), "Median U-Net probability": np.nan},
        {"Population": "U-Net rescued", "Count": unet_direct, "Percent": pct(unet_direct), "Median U-Net probability": np.nan},
        {"Population": "U-Net rescued after split/centerline", "Count": unet_split, "Percent": pct(unet_split), "Median U-Net probability": np.nan},
        {"Population": "All U-Net rescued", "Count": unet_total, "Percent": pct(unet_total), "Median U-Net probability": np.nan},
    ]

    if "unet_mean_probability" in df.columns:
        probs = pd.to_numeric(df["unet_mean_probability"], errors="coerce")
        for row in rows:
            if row["Population"] == "Saturn classical":
                mask = source == "saturn_classical"
            elif row["Population"] == "U-Net rescued":
                mask = unet_mask & ~unet_split_mask
            elif row["Population"] == "U-Net rescued after split/centerline":
                mask = source == "unet_rescued_split"
            elif row["Population"] == "All U-Net rescued":
                mask = source.str.startswith("unet_rescued")
            else:
                mask = pd.Series([True] * len(df), index=df.index)
            vals = probs[mask].dropna()
            if not vals.empty:
                row["Median U-Net probability"] = float(vals.median())

    overlay_count = 0
    probability_map_count = 0
    if out_dir:
        overlay_directories = [
            os.path.join(out_dir, "technical_qc", "unet_rescue_overlays"),
            os.path.join(out_dir, "overlays"),
        ]
        for overlay_dir in overlay_directories:
            if os.path.isdir(overlay_dir):
                overlay_count += len(
                    [
                        name
                        for name in os.listdir(overlay_dir)
                        if name.endswith("_unet_rescue_review.png")
                    ]
                )
        if os.path.isdir(out_dir):
            probability_map_count = len([p for p in os.listdir(out_dir) if p.endswith("_unet_probability.tif")])

    return {
        "enabled": unet_total > 0 or probability_map_count > 0,
        "total_2d": total,
        "saturn_classical": classical,
        "unet_rescued": unet_direct,
        "unet_rescued_split": unet_split,
        "unet_total_rescued": unet_total,
        "unet_rescue_fraction": float(unet_total / max(total, 1)),
        "overlay_count": int(overlay_count),
        "probability_map_count": int(probability_map_count),
        "table": pd.DataFrame(rows),
    }


def generate_excel_report(out_dir, df, df_summary, df_tracks=None):
    """
    Generates a multi-tab Excel workbook with formatted data, summary statistics,
    embedded chart images, and source-data hyperlinks.

    Workbook structure
    ------------------
    - **Batch_Summary** - one row per Z-slice with detection counts, mean/median
      length, and total area.  Includes an embedded histogram image and conditional
      formatting for high/low detection slices.
    - **3D_Track_Audit** - complete track-level provenance and 3D metrics from
      :func:`track_across_slices`, with an embedded 3D length distribution plot.
    - **Raw_2D_Detections** - full per-spermatid measurement table matching the CSV
      export, with number formatting and frozen header pane.
    - **Statistics_Summary** - descriptive statistics (mean, median, std, IQR,
      percentiles) for the primary biometric columns.

    All sheets include a hyperlink from cell A1 back to the source Excel file so
    that clicking within PowerPoint linked charts opens the correct workbook row.

    Biological interpretation
    -------------------------
    The statistical summary sheet is designed to be directly copy-pasteable into
    lab reports.  The IQR (Interquartile Range = Q75 - Q25) is reported because
    spermatid length distributions are often right-skewed and IQR is more robust
    than standard deviation in these cases.

    Args:
        out_dir (str): Top-level analysis output directory.  The workbook is saved
            as ``<out_dir>/batch_analysis_results_<ver>.xlsx``.
        df (pd.DataFrame): Per-spermatid measurement table (all Z-slices combined).
        df_summary (pd.DataFrame): Per-slice summary statistics.
        df_tracks (pd.DataFrame, optional): 3D track table from
            :func:`track_across_slices`.  ``None`` if tracking was not run.
    """
    excel_path = os.path.join(out_dir, f"batch_analysis_results_{_VERSION}.xlsx")
    plot_dir = os.path.join(out_dir, "summary_plots")
    print(f"Generating Interactive Excel Audit: {excel_path} ...")

    try:
        with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
            workbook  = writer.book
            unet_report = summarize_unet_rescue_for_reports(df, out_dir)

            def excel_col_letter(frame, column_name):
                """Return the Excel column letter for a DataFrame column."""
                idx = frame.columns.get_loc(column_name)
                letters = ""
                idx += 1
                while idx:
                    idx, rem = divmod(idx - 1, 26)
                    letters = chr(65 + rem) + letters
                return letters

            # Formats
            bold = workbook.add_format({'bold': True, 'bg_color': '#D7E4BC', 'border': 1})
            num_fmt = workbook.add_format({'num_format': '0.00'})

            # --- Sheet 1: one canonical biologist-facing result set ---
            ws_bio = workbook.add_worksheet('Biologist_Results')
            ws_bio.write(0, 0, "BIOLOGIST RESULTS - USE THIS SHEET", bold)
            if df_tracks is not None and not df_tracks.empty:
                primary = (
                    df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                    if "technical_valid" in df_tracks.columns
                    else df_tracks.copy()
                )
                def primary_median(column):
                    if column not in primary.columns:
                        return np.nan
                    return float(
                        pd.to_numeric(
                            primary[column],
                            errors="coerce",
                        ).median()
                    )

                biologist_metrics = [
                    ("Analysis population", "Included estimated nuclei"),
                    ("Estimated unique nuclei", int(len(primary))),
                    ("Median 3D length (um)", primary_median("total_3d_length_um")),
                    ("Median maximum 2D length (um)", primary_median("max_length_2d")),
                    ("Median 2D width (um)", primary_median("median_width_2d")),
                    ("Median 2D length / width", primary_median("median_length_width_ratio_2d")),
                    ("Median effective thickness (um; PSF-sensitive)", primary_median("thickness_um")),
                    ("Median 3D tortuosity", primary_median("tortuosity_3d")),
                    ("Median Z-span (um)", primary_median("z_span_um")),
                    ("Median slices detected", primary_median("n_slices")),
                ]
                for row_index, (label, value) in enumerate(biologist_metrics, start=2):
                    ws_bio.write(row_index, 0, label)
                    if isinstance(value, float) and not np.isfinite(value):
                        ws_bio.write_blank(row_index, 1, None)
                    else:
                        ws_bio.write(
                            row_index,
                            1,
                            value,
                            num_fmt if isinstance(value, float) else None,
                        )
                ws_bio.write(
                    14,
                    0,
                    "Raw 2D, U-Net, warning-free, reference-morphology, and rejected-extension counts are technical QC only.",
                )
                ws_bio.write(
                    15,
                    0,
                    "Use effective thickness only for samples acquired with matched microscope settings.",
                )
            ws_bio.set_column('A:A', 72)
            ws_bio.set_column('B:B', 24)

            # --- Technical sheet: population summary with dynamic formulas ---
            ws_sum = workbook.add_worksheet('Technical_QC')
            headers = ["Metric", "Average (Formula)", "Median (Formula)", "Std Dev (Formula)"]
            for col, h in enumerate(headers):
                ws_sum.write(0, col, h, bold)

            row = 1
            # 2D Length from Raw_2D_Detections length_um_geodesic column.
            if not df.empty:
                n_2d = len(df) + 1
                length_col = excel_col_letter(df, "length_um_geodesic")
                ws_sum.write(row, 0, "2D Geodesic Length (um)")
                ws_sum.write_formula(row, 1, f"=AVERAGE('Raw_2D_Detections'!{length_col}2:{length_col}{n_2d})", num_fmt)
                ws_sum.write_formula(row, 2, f"=MEDIAN('Raw_2D_Detections'!{length_col}2:{length_col}{n_2d})", num_fmt)
                ws_sum.write_formula(row, 3, f"=STDEV.P('Raw_2D_Detections'!{length_col}2:{length_col}{n_2d})", num_fmt)
                row += 1

            # 3D metrics in the master track-audit table.
            if df_tracks is not None and not df_tracks.empty:
                n_3d = len(df_tracks) + 1 if not df_tracks.empty else 2  # Default to 2 to avoid #DIV/0! bounds
                metrics_3d = [
                    ("3D Geodesic Length (um)", "total_3d_length_um"),
                    ("3D Z-Span (um)", "z_span_um"),
                    ("3D Z-Covered (um)", "z_covered_um"),
                    ("3D Volume (um3)", "volume_um3"),
                    ("3D Tortuosity", "tortuosity_3d")
                ]
                for m_name, col_name in metrics_3d:
                    if col_name not in df_tracks.columns:
                        continue
                    col_letter = excel_col_letter(df_tracks, col_name)
                    ws_sum.write(row, 0, m_name)
                    ws_sum.write_formula(row, 1, f"=AVERAGE('3D_Track_Audit'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 2, f"=MEDIAN('3D_Track_Audit'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 3, f"=STDEV.P('3D_Track_Audit'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    row += 1

            ws_sum.set_column('A:A', 30)

            # --- Sheet 2: complete 3D track audit ---
            if df_tracks is not None and not df_tracks.empty:
                has_reference = "is_reference_morphology_track" in df_tracks.columns or "reference_morphology_pass" in df_tracks.columns
                df_tracks.to_excel(writer, sheet_name='3D_Track_Audit', index=False)
                ws_3d = writer.sheets['3D_Track_Audit']
                ws_3d.set_column('A:Z', 15)
                # Insert the primary estimated-nuclei 3D distribution graph.
                p_3d = os.path.join(plot_dir, "3d_length_distribution.png")
                if os.path.exists(p_3d):
                    ws_3d.insert_image('K2', p_3d, {'x_scale': 0.4, 'y_scale': 0.4})

                # Reference morphology is diagnostic and remains separate from
                # warning-free and biological-candidate populations.
                if has_reference:
                    reference_col = (
                        "is_reference_morphology_track"
                        if "is_reference_morphology_track" in df_tracks.columns
                        else "reference_morphology_pass"
                    )
                    df_reference = df_tracks[df_tracks[reference_col]].copy()
                    df_reference.to_excel(writer, sheet_name='3D_Reference_Morphology', index=False)
                    ws_reference = writer.sheets['3D_Reference_Morphology']
                    ws_reference.set_column('A:Z', 15)
            # --- Sheet 3: Raw 2D Detections ---
            if not df.empty:
                df.to_excel(writer, sheet_name='Raw_2D_Detections', index=False)
                ws_2d = writer.sheets['Raw_2D_Detections']
                ws_2d.set_column('A:Z', 15)

            # --- Sheet 4: Slice Summary ---
            if not df_summary.empty:
                df_summary.to_excel(writer, sheet_name='Slice_Summary', index=False)

            # --- Sheet 4b: v5.7 U-Net Rescue Audit ---
            if unet_report and unet_report["enabled"]:
                ws_unet = workbook.add_worksheet('U-Net_Rescue_Audit')
                ws_unet.write(0, 0, "Saturn v5.7 U-Net Rescue Audit", bold)
                ws_unet.write(1, 0, "All counts below are measurement-table counts, not overlay-pixel counts.")
                ws_unet.write(2, 0, "Overlay dilation is display-only and does not affect count, length, width, or 3D tracking.")
                ws_unet.write(4, 0, "Probability maps saved")
                ws_unet.write(4, 1, unet_report["probability_map_count"])
                ws_unet.write(5, 0, "U-Net rescue review overlays saved")
                ws_unet.write(5, 1, unet_report["overlay_count"])
                ws_unet.write(6, 0, "U-Net rescued fraction of 2D detections")
                ws_unet.write(6, 1, unet_report["unet_rescue_fraction"], num_fmt)
                table_df = unet_report["table"]
                for c_idx, col_name in enumerate(table_df.columns):
                    ws_unet.write(8, c_idx, col_name, bold)
                for r_idx, row_data in enumerate(table_df.itertuples(index=False), start=9):
                    for c_idx, val in enumerate(row_data):
                        if isinstance(val, (int, float, np.integer, np.floating)) and not pd.isna(val):
                            ws_unet.write(r_idx, c_idx, float(val), num_fmt if c_idx >= 2 else None)
                        else:
                            ws_unet.write(r_idx, c_idx, "" if pd.isna(val) else val)
                ws_unet.set_column('A:A', 36)
                ws_unet.set_column('B:D', 18)

            ws_sum.set_column('B:D', 18)

            # Insert Histograms into Summary
            p_hist = os.path.join(plot_dir, "global_histograms.png")
            p_slice = os.path.join(plot_dir, "length_by_slice.png")
            if os.path.exists(p_hist):
                ws_sum.insert_image('F2', p_hist, {'x_scale': 0.5, 'y_scale': 0.5})
            if os.path.exists(p_slice):
                ws_sum.insert_image('F25', p_slice, {'x_scale': 0.5, 'y_scale': 0.5})

            # --- Sheet 5: Methods Dictionary ---
            dictionary_data = [
                ["Metric", "Formula / Definition", "Biological Interpretation"],
                ["2D Geodesic Length", "Shortest-path length along a 2D skeleton", "Curved fragment length within a single optical slice."],
                ["Total 3D Geodesic Length", "sqrt(max(2D geodesic, XY displacement)^2 + z_span^2)", "Projection-length plus Z-span estimate of whole-nucleus 3D length."],
                ["3D Euclidean Distance", "sqrt(XY displacement^2 + z_span^2)", "Straight-line span used as the tortuosity denominator."],
                ["3D Tortuosity", "Total 3D length / 3D Euclidean distance", "Curvature or over-merge index. Values near 1 indicate straighter nuclei."],
                ["Z-Span (Vertical Span)", "(max_z - min_z) * UM_PER_SLICE_Z", "Endpoint-to-endpoint vertical span; single-slice tracks have span 0."],
                ["Z-Covered", "(max_z - min_z + 1) * UM_PER_SLICE_Z", "Sampled slab thickness covered by the detections."],
                ["3D Volume (um3) *", "sum(area_est_slice * XY_pixel_area * Z_step)", "PSF- and voxel-sensitive Riemann-sum approximation of nuclear volume. Use mainly for relative comparisons acquired under matched imaging settings."],
                ["Effective Diameter Proxy (um) *", "2 * sqrt((V_3D / L_3D) / pi)", "PSF-sensitive cylinder-equivalent diameter. Comparative metric only; do not interpret as literal physical diameter."],
                ["Pitch Angle (degrees)", "abs(arcsin(z_span / Euclidean_3D)) * 180/pi", "Absolute plunge angle relative to the imaging plane."],
                ["Taper Ratio *", "max(area_est across track) / min(area_est across track)", "PSF-sensitive area-derived metric. Useful for relative comparison and instability screening, not as a literal anatomical ratio."],
                ["Nearest Neighbor (um)", "Nearest 3D centroid-to-centroid distance", "Simple local packing-density readout."],
                ["Quality Audit", "Strict post-tracking flag based on length, tortuosity, thickness, taper, and minimum slice count", "Strict audit labels completed tracks after tracking. It does not change segmentation; it defines the no-warning subset used for conservative summaries."],
                ["Estimated Unique Nucleus", "A technical-valid 3D track; morphology warnings remain included", "This is the one nucleus population used for biological analysis."],
                ["v5.7 U-Net Rescue", "U-Net probability maps add a rescue lane for detections missed by classical Saturn; accepted detections are labeled by detection_source.", "Use source counts to audit how much the AI lane contributed. The U-Net lane does not replace downstream biological QC."],
                ["U-Net Rescue Overlay", "Green = Saturn classical; cyan = accepted U-Net rescue; magenta/orange/red = U-Net-positive candidates rejected by rescue gates.", "Overlay line thickness is display-only and never used for count, length, width, or tracking calculations."],
                ["Parameter Tuning Guidance", "Candidate audit first -> Tracking second -> Segmentation last", "Change audit interpretation when summaries are too strict; change tracking when tracks are fragmented or over-merged; change segmentation only when the raw 2D detections themselves are wrong."],
                ["Standard Deviation", "Std Dev", "Population spread around the mean."]
            ]
            pd.DataFrame(dictionary_data[1:], columns=dictionary_data[0]).to_excel(writer, sheet_name='Methods_Dictionary', index=False)
            ws_dict = writer.sheets['Methods_Dictionary']

            # Manually write the dictionary to avoid any encoding or header issues
            for r_idx, r_data in enumerate(dictionary_data):
                for c_idx, val in enumerate(r_data):
                    ws_dict.write(r_idx, c_idx, val, bold if r_idx==0 else None)

            ws_dict.set_column('A:A', 25)
            ws_dict.set_column('B:B', 40)
            ws_dict.set_column('C:C', 60)
            p_guide = os.path.join(plot_dir, "methods_guide.png")
            if os.path.exists(p_guide):
                ws_dict.insert_image('A12', p_guide, {'x_scale': 0.6, 'y_scale': 0.6})

            print(f"Interactive Excel report successfully saved to {excel_path}")

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating Excel report: {e}")
        write_error_log(out_dir, "Excel Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"Excel Report failed to generate completely.\n{e}")
        except Exception:
            pass



def generate_batch_report(
        out_dir, df, df_summary, um, df_tracks=None, gui_callback=None,
        generate_pptx=True, df_tracked=None):
    """
    Compiles standard summary global output architectures including histograms, mathematical summaries,
    biological methodology pages, and graphical slice overlays natively to a `.pdf` file.

    Args:
        out_dir (str): Root export directory.
        df (pd.DataFrame): Flat 2D analysis parameters.
        df_summary (pd.DataFrame): Top-level slice aggregation tracking.
        um (dict): User-defined pixel to micron mapping ratios.
        df_tracks (pd.DataFrame, optional): Unified 3D tracking geometries array.
        gui_callback (function, optional): Live variable passing to front-end dashboards elements.
        df_tracked (pd.DataFrame, optional): Per-slice measurements with assigned
            track IDs. Used to keep detailed overlay statistics aligned with the
            technical-valid biological population.

    Returns:
        None (Saves directly to absolute path PDF)
    """
    pdf_path = os.path.join(out_dir, f"batch_report_{_VERSION}.pdf")
    print(f"Generating high-res PDF report: {pdf_path} ...")

    # Create directory for high-res standalone plots (for easy copy-pasting into papers/presentations)
    plot_dir = os.path.join(out_dir, "summary_plots")
    os.makedirs(plot_dir, exist_ok=True)

    try:
        with PdfPages(pdf_path) as pdf:
            unet_report = summarize_unet_rescue_for_reports(df, out_dir)
            # --- PAGE 1: GLOBAL SUMMARY ---
            fig_sum = plt.figure(figsize=(11, 8.5))
            fig_sum.suptitle(
                f"Technical QC: Detection and Tracking Provenance - {_VERSION}\n"
                f"Do not use these counts as the biological population | Location: {out_dir}",
                fontsize=13,
                fontweight='bold',
            )
            df_candidate_tracks = (
                df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                if df_tracks is not None and not df_tracks.empty and "technical_valid" in df_tracks.columns
                else (df_tracks.copy() if df_tracks is not None else pd.DataFrame())
            )
            has_candidate_data = not df_candidate_tracks.empty

            # Global Z projection for the one primary analysis population.
            analysis_z_proj_path = os.path.join(
                out_dir,
                "analysis_global_z_projection.png",
            )
            raw_z_proj_path = os.path.join(out_dir, "global_z_projection.png")
            z_proj_path = (
                analysis_z_proj_path
                if os.path.exists(analysis_z_proj_path)
                else raw_z_proj_path
            )
            if os.path.exists(z_proj_path):
                ax_z = fig_sum.add_axes([0.15, 0.62, 0.7, 0.28]) # [left, bottom, width, height]
                ax_z.imshow(plt.imread(z_proj_path))
                if z_proj_path == analysis_z_proj_path:
                    ax_z.set_title(
                        "Included Estimated Nuclei Across the Z Stack",
                        fontsize=10,
                    )
                else:
                    ax_z.set_title("Global Z-Projection (Composite [Original | Overlay])", fontsize=10)
                ax_z.axis('off')

            # Plot 1: Counts per slice
            ax1 = fig_sum.add_subplot(2, 2, 3)
            if has_candidate_data and "z_start" in df_candidate_tracks.columns:
                q_counts = (
                    df_candidate_tracks["z_start"]
                    .astype(int)
                    .value_counts()
                    .reindex(df_summary["z_slice"].astype(int), fill_value=0)
                    .sort_index()
                )
                ax1.plot(q_counts.index, q_counts.values, 'go-', markersize=4, linewidth=1.5, label='Estimated nuclei by start Z')
                ax1.set_title("Estimated Unique Nuclei by Z-Start")
            else:
                ax1.plot(
                    df_summary['z_slice'],
                    df_summary['n_spermatids'],
                    color='forestgreen',
                    marker='o',
                    markersize=3,
                    linewidth=1,
                    label='2D preview candidates',
                )
                ax1.set_title("Raw 2D Detections per Z-Slice")
            ax1.set_xlabel("Z-Index")
            ax1.set_ylabel("Count")
            ax1.grid(True, alpha=0.3)
            ax1.legend(fontsize=7)

            # Plot 2: primary estimated-nuclei length distribution.
            ax2 = fig_sum.add_subplot(2, 2, 4)
            if has_candidate_data and "total_3d_length_um" in df_candidate_tracks.columns:
                vals = df_candidate_tracks['total_3d_length_um'].dropna()
                ax2.hist(vals, bins=25, color='darkorange', edgecolor='black', alpha=0.75, label='Estimated nuclei')
                m_med = vals.median()
                m_avg = vals.mean()
                ax2.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax2.axvline(m_avg, color='black', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax2.set_title("Estimated-Nuclei 3D Length Distribution")
                ax2.set_xlabel("Total 3D Length (um)")
                ax2.set_ylabel("Frequency")
                ax2.legend(fontsize=8)
            elif not df.empty:
                vals = df['length_um_geodesic']
                ax2.hist(vals, bins=25, color='forestgreen', edgecolor='black', alpha=0.7)
                m_med = vals.median()
                m_avg = vals.mean()
                ax2.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax2.axvline(m_avg, color='orange', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax2.set_title("Raw 2D Length Distribution (No Candidate Tracks)")
                ax2.set_xlabel("2D Geodesic Length (um)")
                ax2.set_ylabel("Frequency")
                ax2.legend(fontsize=8)

            fig_sum.savefig(os.path.join(plot_dir, "global_summary.png"), dpi=300, bbox_inches='tight')
            # NEW: Explicitly save global_histograms.png for Excel embedding
            fig_sum.savefig(os.path.join(plot_dir, "global_histograms.png"), dpi=300, bbox_inches='tight')

            # Save length_by_slice.png
            fig_l_slice = plt.figure(figsize=(10, 5))
            ax_ls = fig_l_slice.add_subplot(1, 1, 1)
            ax_ls.plot(df_summary['z_slice'], df_summary['median_length_um'], 'go-', label='Median Length')
            ax_ls.set_title("Median Length by Slice")
            ax_ls.set_xlabel("Z-Slice")
            ax_ls.set_ylabel("Length (um)")
            ax_ls.grid(True, alpha=0.3)
            fig_l_slice.savefig(os.path.join(plot_dir, "length_by_slice.png"), dpi=300, bbox_inches='tight')
            plt.close(fig_l_slice)

            # --- PAGE 1.5: BIOLOGIST-FACING PRIMARY RESULTS ---
            if df_tracks is not None and not df_tracks.empty:
                fig_dyn = plt.figure(figsize=(11, 8.5))
                fig_dyn.suptitle("Biologist Results: Primary 3D Population", fontsize=16, fontweight='bold')
                primary = (
                    df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                    if "technical_valid" in df_tracks.columns
                    else df_tracks.copy()
                )
                def report_median(column):
                    if column not in primary.columns:
                        return np.nan
                    return float(
                        pd.to_numeric(
                            primary[column],
                            errors="coerce",
                        ).median()
                    )

                ax_metrics = fig_dyn.add_subplot(1, 2, 1)
                ax_metrics.axis("off")
                metric_lines = [
                    ("Estimated unique nuclei", f"{len(primary):,}"),
                    ("Median 3D length", f"{report_median('total_3d_length_um'):.2f} um"),
                    ("Median maximum 2D length", f"{report_median('max_length_2d'):.2f} um"),
                    ("Median 2D width", f"{report_median('median_width_2d'):.2f} um"),
                    ("Median 2D length/width", f"{report_median('median_length_width_ratio_2d'):.2f}"),
                    ("Median effective thickness*", f"{report_median('thickness_um'):.2f} um"),
                ]
                y = 0.88
                for label, value in metric_lines:
                    ax_metrics.text(0.02, y, label, fontsize=10, color="#444444", va="top")
                    ax_metrics.text(0.02, y - 0.045, value, fontsize=18, fontweight="bold", color="#145a32", va="top")
                    y -= 0.135
                ax_metrics.text(
                    0.02,
                    0.02,
                    "*PSF-sensitive; compare only between samples acquired with matched microscope settings.",
                    fontsize=8.5,
                    color="#555555",
                    wrap=True,
                )

                ax_length = fig_dyn.add_subplot(1, 2, 2)
                ax_length.hist(
                    primary["total_3d_length_um"].dropna(),
                    bins=25,
                    color="#2ca02c",
                    edgecolor="black",
                    alpha=0.8,
                )
                ax_length.axvline(
                    primary["total_3d_length_um"].median(),
                    color="black",
                    linestyle="--",
                    linewidth=1.5,
                    label=f"Median {primary['total_3d_length_um'].median():.2f} um",
                )
                ax_length.set_title("Primary Population: 3D Length")
                ax_length.set_xlabel("3D length (um)")
                ax_length.set_ylabel("Estimated unique nuclei")
                ax_length.legend(fontsize=9)
                ax_length.grid(True, alpha=0.2)

                fig_dyn.text(
                    0.5,
                    0.035,
                    "Use this included estimated-nucleus population for biological comparisons. "
                    "Detection provenance and excluded technical failures are stored separately "
                    "and are not alternative biological counts.",
                    ha="center",
                    fontsize=9.5,
                    fontweight="bold",
                    color="#333333",
                    wrap=True,
                )
                plt.tight_layout(rect=[0, 0.08, 1, 0.94])
                fig_dyn.savefig(os.path.join(plot_dir, "population_consolidation.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_dyn, dpi=300, bbox_inches='tight')
                plt.close(fig_dyn)

            # Keep detection/tracking provenance as a standalone technical-QC
            # image. It is intentionally excluded from the biologist-facing PDF.
            plt.close(fig_sum)

            # Save the detailed U-Net audit as a standalone technical-QC image,
            # not as another competing population in the final PDF.
            if unet_report and unet_report["enabled"]:
                fig_unet = plt.figure(figsize=(11, 8.5))
                fig_unet.suptitle("Saturn v5.7 U-Net Rescue Audit", fontsize=15, fontweight='bold')
                gs = fig_unet.add_gridspec(2, 2, height_ratios=[2.0, 1.2])
                ax_bar = fig_unet.add_subplot(gs[0, 0])
                table = unet_report["table"].copy()
                plot_table = table[table["Population"].isin([
                    "Saturn classical",
                    "U-Net rescued",
                    "U-Net rescued after split/centerline",
                ])]
                ax_bar.bar(
                    plot_table["Population"],
                    plot_table["Count"],
                    color=["#2ca02c", "#17becf", "#00a6c8"],
                    edgecolor="black",
                )
                ax_bar.set_title("2D Detection Source Accounting")
                ax_bar.set_ylabel("Detection count")
                ax_bar.tick_params(axis='x', labelrotation=20)
                for idx, val in enumerate(plot_table["Count"]):
                    ax_bar.text(idx, val + max(plot_table["Count"].max() * 0.02, 1), f"{int(val):,}", ha='center', fontsize=10, fontweight='bold')

                ax_text = fig_unet.add_subplot(gs[0, 1])
                ax_text.axis('off')
                note = (
                    f"Total 2D detections: {unet_report['total_2d']:,}\n"
                    f"Saturn classical: {unet_report['saturn_classical']:,}\n"
                    f"U-Net rescued total: {unet_report['unet_total_rescued']:,} "
                    f"({unet_report['unet_rescue_fraction']*100:.1f}%)\n"
                    f"Probability maps saved: {unet_report['probability_map_count']:,}\n"
                    f"Rescue-review overlays saved: {unet_report['overlay_count']:,}\n\n"
                    "Interpretation\n"
                    "Cyan rescue detections are U-Net-supported objects accepted by the same measurement table used for CSV/Excel exports. "
                    "Red/orange/magenta review-overlay marks show U-Net-positive candidates that were not accepted by the rescue gates.\n\n"
                    "Overlay display note\n"
                    "Overlay dilation is cosmetic only. It is not used for counts, skeleton length, width, or 3D tracking."
                )
                ax_text.text(0, 1, note, transform=ax_text.transAxes, va='top', fontsize=11, linespacing=1.35)

                ax_table = fig_unet.add_subplot(gs[1, :])
                ax_table.axis('off')
                display_table = table.copy()
                display_table["Percent"] = display_table["Percent"].map(lambda v: f"{v:.1f}%")
                display_table["Median U-Net probability"] = display_table["Median U-Net probability"].map(
                    lambda v: "" if pd.isna(v) else f"{v:.3f}"
                )
                tab = ax_table.table(
                    cellText=display_table.values,
                    colLabels=display_table.columns,
                    loc='center',
                    cellLoc='center',
                )
                tab.auto_set_font_size(False)
                tab.set_fontsize(9)
                tab.scale(1, 1.35)
                fig_unet.tight_layout(rect=[0, 0.03, 1, 0.93])
                fig_unet.savefig(os.path.join(plot_dir, "unet_rescue_audit.png"), dpi=300, bbox_inches='tight')
                plt.close(fig_unet)

            # --- PAGE 2: PRIMARY 3D MORPHOMETRICS SUMMARY ---
            if df_tracks is not None and not df_tracks.empty:
                fig_3d = plt.figure(figsize=(11, 8.5))

                df_q = (
                    df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                    if "technical_valid" in df_tracks.columns
                    else df_tracks.copy()
                )
                fig_3d.suptitle(
                    f"3D Statistics: {len(df_q):,} Estimated Unique Nuclei",
                    fontsize=14,
                    fontweight='bold',
                )

                # 3D Length
                ax3d_1 = fig_3d.add_subplot(2, 2, 1)
                vals_q = df_q['total_3d_length_um']
                ax3d_1.hist(vals_q, bins=20, color='darkorange', edgecolor='black', alpha=0.7, label='Estimated nuclei')
                stats_len = vals_q
                m_med = stats_len.median()
                m_avg = stats_len.mean()
                ax3d_1.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax3d_1.axvline(m_avg, color='black', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax3d_1.set_title("Total 3D Geodesic Length")
                ax3d_1.set_xlabel("Length (um)")
                ax3d_1.set_ylabel("Frequency")
                ax3d_1.legend(fontsize=7)

                # Save the primary-population distribution for Excel embedding.
                fig_3d_len = plt.figure(figsize=(6, 4))
                ax_3dl = fig_3d_len.add_subplot(1, 1, 1)
                ax_3dl.hist(vals_q, bins=20, color='steelblue', edgecolor='black', alpha=0.75)
                ax_3dl.set_title("3D Length Distribution (Estimated Nuclei)")
                fig_3d_len.savefig(os.path.join(plot_dir, "3d_length_distribution.png"), dpi=300, bbox_inches='tight')
                plt.close(fig_3d_len)

                # 3D Tortuosity
                ax3d_2 = fig_3d.add_subplot(2, 2, 2)
                vt_q = df_q['tortuosity_3d']
                vt_q_viz = vt_q[(vt_q >= 0.95) & (vt_q <= 3.0)]
                ax3d_2.hist(vt_q_viz, bins=25, color='purple', edgecolor='black', alpha=0.6, label='Estimated nuclei')
                stats_tort = vt_q
                ax3d_2.axvline(stats_tort.median(), color='red', linestyle='-', label=f"Median: {stats_tort.median():.2f}")
                ax3d_2.axvline(stats_tort.mean(), color='black', linestyle='--', label=f"Mean: {stats_tort.mean():.2f}")
                ax3d_2.set_xlim(0.95, 3.0)
                ax3d_2.set_title("3D Tortuosity (Curvature)")
                ax3d_2.set_xlabel("Ratio (Length / Distance)")
                ax3d_2.set_ylabel("Frequency")
                ax3d_2.legend(fontsize=7)

                # Vertical Extent
                ax3d_3 = fig_3d.add_subplot(2, 2, 3)
                z_col = "z_span_um" if "z_span_um" in df_tracks.columns else "z_extent_um"
                ve_q = df_q[z_col]
                ax3d_3.hist(ve_q, bins=15, color='teal', edgecolor='black', alpha=0.7, label='Estimated nuclei')
                stats_z = ve_q
                ax3d_3.axvline(stats_z.median(), color='red', linestyle='-', label=f"Median: {stats_z.median():.1f}")
                ax3d_3.axvline(stats_z.mean(), color='black', linestyle='--', label=f"Mean: {stats_z.mean():.1f}")
                ax3d_3.set_title("Z-Span (Vertical Span)")
                ax3d_3.set_xlabel("Vertical Span (um)")
                ax3d_3.set_ylabel("Frequency")
                ax3d_3.legend(fontsize=7)

                # Volume
                ax3d_4 = fig_3d.add_subplot(2, 2, 4)
                vv_q = df_q['volume_um3']
                ax3d_4.hist(vv_q, bins=20, color='gray', edgecolor='black', alpha=0.7, label='Estimated nuclei')
                stats_vol = vv_q
                ax3d_4.axvline(stats_vol.median(), color='red', linestyle='-', label=f"Median: {stats_vol.median():.0f}")
                ax3d_4.axvline(stats_vol.mean(), color='black', linestyle='--', label=f"Mean: {stats_vol.mean():.0f}")
                ax3d_4.set_title("Approximated 3D Volume")
                ax3d_4.set_xlabel("Volume (um\u00b3)")
                ax3d_4.set_ylabel("Frequency")
                ax3d_4.legend(fontsize=7)

                plt.tight_layout(rect=[0, 0.03, 1, 0.95])
                fig_3d.savefig(os.path.join(plot_dir, "3d_population_stats.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_3d, dpi=300, bbox_inches='tight')
                plt.close(fig_3d)

            # Methods Guide securely moved down past the Advanced Biometrics block.

            # --- PAGE 4: ADVANCED 3D BIOMETRICS FOR ESTIMATED NUCLEI ---
            if df_tracks is not None and not df_tracks.empty:
                fig_adv = plt.figure(figsize=(11, 8.5))
                df_q = (
                    df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                    if "technical_valid" in df_tracks.columns
                    else df_tracks.copy()
                )
                fig_adv.suptitle(
                    f"Advanced 3D Biometrics: {len(df_q):,} Estimated Unique Nuclei",
                    fontsize=16,
                    fontweight='bold',
                    y=0.96,
                )

                # Helper for Mean/Median
                def add_stats_lines(ax, data_series):
                    if data_series.empty or data_series.isna().all(): return
                    m = data_series.mean()
                    med = data_series.median()
                    ax.axvline(med, color='red', linestyle='--', linewidth=1.5, label=f'Median: {med:.2f}')
                    ax.axvline(m, color='green', linestyle=':', linewidth=2, label=f'Mean: {m:.2f}')
                    ax.legend(fontsize=8)

                def dual_hist(ax, col, title, xlabel, color_q, bins=30):
                    vals_q = df_q[col].dropna() if col in df_q.columns else pd.Series(dtype=float)
                    if not vals_q.empty:
                        sns.histplot(vals_q, bins=bins, ax=ax, color=color_q, edgecolor='black', alpha=0.7, label='Estimated nuclei')
                        add_stats_lines(ax, vals_q)
                    ax.set_title(title)
                    ax.set_xlabel(xlabel)
                    ax.set_ylabel("Frequency")
                    if len(ax.get_legend_handles_labels()[0]) > 0:
                        ax.legend(fontsize=7)

                # 4 panels: Pitch, Thickness, Taper, Nearest Neighbor
                ax_p = fig_adv.add_subplot(2, 2, 1)
                dual_hist(ax_p, 'pitch_deg', "Pitch Angle (Vertical Plunge)", "Degrees (0=Flat, 90=Vertical)", 'orange')

                ax_th = fig_adv.add_subplot(2, 2, 2)
                dual_hist(ax_th, 'thickness_um', "Effective Nucleus Thickness", "Average Diameter (\u00b5m)", '#17becf')

                ax_ta = fig_adv.add_subplot(2, 2, 3)
                dual_hist(ax_ta, 'taper_ratio', "Morphological Taper Ratio", "Max Area / Min Area", 'purple')

                ax_nn = fig_adv.add_subplot(2, 2, 4)
                dual_hist(ax_nn, 'nearest_neighbor_um', "Spatial Packing Density", "Distance to Nearest Neighbor (\u00b5m)", 'brown')

                plt.tight_layout(rect=[0, 0.03, 1, 0.93])
                fig_adv.savefig(os.path.join(plot_dir, "advanced_biometrics.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_adv, dpi=300, bbox_inches='tight')
                plt.close(fig_adv)

            # --- PAGE 5: METHODS & INTERPRETATION GUIDE ---
            fig_guide = plt.figure(figsize=(11, 8.5))
            ax_g = fig_guide.add_axes([0.05, 0.05, 0.9, 0.9])
            ax_g.axis('off')
            guide_full = (
                "METHODS, FORMULAE, AND AUDIT GUIDE\n"
                f"{'='*80}\n\n"
                "1. Total 3D Geodesic Length (um)\n"
                "   Formula: L_3D = sqrt(max(2D geodesic, XY displacement)^2 + z_span^2)\n"
                "   Meaning: Projection-length plus Z-span estimate of whole-nucleus 3D length.\n\n"
                "2. 3D Euclidean Distance (um)\n"
                "   Formula: D_3D = sqrt(XY displacement^2 + z_span^2)\n"
                "   Meaning: Straight-line span used as the tortuosity denominator.\n\n"
                "3. 3D Tortuosity\n"
                "   Formula: T = L_3D / D_3D\n"
                "   Meaning: Curvature index. Values near 1 are straighter; high values suggest bent or fused tracks.\n\n"
                "4. Z-Span and Z-Covered\n"
                "   Formula: z_span = (max_z - min_z) * UM_PER_SLICE_Z; z_covered = (max_z - min_z + 1) * UM_PER_SLICE_Z\n"
                "   Meaning: Z-span is endpoint-to-endpoint displacement; Z-covered is sampled slab thickness.\n\n"
                "5. Approximated 3D Volume (um3)\n"
                "   Formula: V_3D = sum(area_est_slice * XY_pixel_area * Z_step)\n"
                "   Meaning: PSF- and voxel-sensitive volume estimate. Best used for relative comparisons between datasets acquired with matched imaging settings, not as a literal absolute volume.\n\n"
                "6. Effective Diameter Proxy (Average Diameter, um)\n"
                "   Formula: D_avg = 2 * sqrt((V_3D / L_3D) / pi)\n"
                "   Meaning: PSF-sensitive cylinder-equivalent diameter. Useful for relative comparisons under matched imaging settings, not as a literal physical diameter.\n\n"
                "7. Pitch Angle (Plunge Vector, Degrees)\n"
                "   Formula: theta = abs(arcsin(z_span / D_3D)) * 180/pi\n"
                "   Meaning: Absolute plunge angle relative to the imaging plane.\n\n"
                "8. Taper Ratio\n"
                "   Formula: R = max(area_est across track) / min(area_est across track)\n"
                "   Meaning: PSF-sensitive area-derived instability metric. Large values can reflect abrupt narrowing, fusion, or segmentation instability; use mainly for relative comparison.\n\n"
                "9. Spatial Packing Density (Nearest Neighbor Dist, um)\n"
                "   Formula: nearest 3D centroid-to-centroid distance\n"
                "   Meaning: Simple local packing-density readout.\n\n"
                "10. Technical Track Audit (post-tracking)\n"
                "   Hard-fail rules: length > AUDIT_MAX_LENGTH_UM, tortuosity > AUDIT_MAX_TORTUOSITY, extreme thickness > AUDIT_EXTREME_THICKNESS_UM, extreme taper > AUDIT_EXTREME_TAPER_RATIO, and n_slices < AUDIT_MIN_SLICES.\n"
                "   Warning-only rules: thickness > AUDIT_MAX_THICKNESS_UM and taper > AUDIT_MAX_TAPER_RATIO. These PSF-sensitive flags remain in the estimated-nuclei population for review.\n"
                "   Interpretation: Technical-valid tracks define the one estimated-nuclei population. Reference morphology and warning-free flags are QC annotations, not separate biological populations.\n"
                "\n11. v5.7 U-Net Rescue Lane\n"
                "   U-Net probability maps can add a rescue lane for classical-Saturn misses. Accepted rescues retain detailed detection_source labels for direct, split, short high-confidence, and low-ratio high-confidence recovery.\n"
                "   The rescue-review overlay uses green for Saturn classical detections, cyan for accepted U-Net rescues, and magenta/orange/red for U-Net-positive candidates rejected by rescue gates.\n"
                "   Overlay dilation is display-only and is never used for count, length, width, or 3D tracking calculations.\n\n"
                "12. PSF-sensitive metrics note\n"
                "   Volume, effective thickness, taper, and other width/area-derived values are broadened by microscope PSF and voxel sampling. Use them mainly for relative comparisons between matched WT and mutant datasets, not as literal physical dimensions.\n"
            )
            ax_g.text(0, 1, guide_full, transform=ax_g.transAxes, fontsize=10, family='monospace', verticalalignment='top', linespacing=1.3)
            fig_guide.savefig(os.path.join(plot_dir, "methods_guide.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_guide, dpi=300, bbox_inches='tight')
            plt.close(fig_guide)

            # --- PAGE 6: GLOBAL STATISTICS TABLE (after Methods) ---
            fig_tab = plt.figure(figsize=(11, 8.5))
            ax_t = fig_tab.add_subplot(1, 1, 1)
            ax_t.axis('off')
            ax_t.set_title("Global Population Statistics Summary", fontsize=14, fontweight='bold', pad=20)

            stats_rows = []
            if df_tracks is not None and not df_tracks.empty:
                primary = (
                    df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                    if "technical_valid" in df_tracks.columns
                    else df_tracks.copy()
                )
                stats_rows.append(["--- ESTIMATED UNIQUE NUCLEI ---", f"N={len(primary)}", "", ""])
                l3d = primary['total_3d_length_um']
                z_col = "z_span_um" if "z_span_um" in primary.columns else "z_extent_um"
                ze = primary[z_col]
                vo = primary['volume_um3']
                to = primary['tortuosity_3d']
                th = primary['thickness_um']
                stats_rows.append(["3D Length (um)", f"{l3d.mean():.2f}", f"{l3d.median():.2f}", f"{l3d.std():.2f}"])
                if "median_width_2d" in primary.columns:
                    width_2d = primary["median_width_2d"]
                    stats_rows.append(["2D Width (um)", f"{width_2d.mean():.2f}", f"{width_2d.median():.2f}", f"{width_2d.std():.2f}"])
                if "median_length_width_ratio_2d" in primary.columns:
                    ratio_2d = primary["median_length_width_ratio_2d"]
                    stats_rows.append(["2D Length / Width", f"{ratio_2d.mean():.2f}", f"{ratio_2d.median():.2f}", f"{ratio_2d.std():.2f}"])
                stats_rows.append(["3D Z-Span (um)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
                stats_rows.append(["3D Volume (um3)*", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
                stats_rows.append(["3D Tortuosity", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
                stats_rows.append(["3D Thickness (um)*", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])
                pitch = primary['pitch_deg']
                taper = primary['taper_ratio']
                neighbor = primary['nearest_neighbor_um'].dropna()
                stats_rows.append(["3D Pitch (degrees)", f"{pitch.mean():.1f}", f"{pitch.median():.1f}", f"{pitch.std():.1f}"])
                stats_rows.append(["3D Taper Ratio*", f"{taper.mean():.2f}", f"{taper.median():.2f}", f"{taper.std():.2f}"])
                if not neighbor.empty:
                    stats_rows.append(["Nearest Neighbor (um)", f"{neighbor.mean():.1f}", f"{neighbor.median():.1f}", f"{neighbor.std():.1f}"])

            if stats_rows:
                table = ax_t.table(
                    cellText=stats_rows,
                    colLabels=["Metric", "Mean", "Median", "Std Dev"],
                    loc='center', cellLoc='center',
                    colWidths=[0.38, 0.2, 0.2, 0.2]
                )
                table.auto_set_font_size(False)
                table.set_fontsize(10)
                table.scale(1.2, 2.5)

            fig_tab.savefig(os.path.join(plot_dir, "global_statistics_table.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_tab, dpi=300, bbox_inches='tight')
            plt.close(fig_tab)

            # --- SUBSEQUENT PAGES: PER-SLICE DETAILS (2 panels: [Panel] | [Histogram]) ---
            overlay_dir = os.path.join(out_dir, "overlays")
            analysis_dir = os.path.join(out_dir, "analysis_overlays")
            accepted_track_ids = set()
            if (
                    df_candidate_tracks is not None
                    and not df_candidate_tracks.empty
                    and "track_id" in df_candidate_tracks.columns):
                accepted_track_ids = set(
                    pd.to_numeric(df_candidate_tracks["track_id"], errors="coerce")
                    .dropna()
                    .astype(int)
                    .tolist()
                )
            for idx_p, (row_idx, row) in enumerate(df_summary.iterrows()):
                z = int(row['z_slice'])
                analysis_panel_path = os.path.join(
                    analysis_dir,
                    f"z{z:02d}_analysis_panel.png",
                )
                raw_panel_path = os.path.join(overlay_dir, f"z{z:02d}_panel.png")
                panel_path = (
                    analysis_panel_path
                    if os.path.exists(analysis_panel_path)
                    else raw_panel_path
                )
                uses_analysis_overlay = panel_path == analysis_panel_path

                if not os.path.exists(panel_path):
                    continue

                fig_slice = plt.figure(figsize=(18, 7))
                fig_slice.suptitle(
                    f"Z-Slice {z:02d} [Original | Included Nuclei | Length Distribution]"
                    if uses_analysis_overlay else
                    f"Z-Slice {z:02d} Pre-Tracking Review [Original | Candidate Overlay | Distribution]",
                    fontsize=12,
                    fontweight='bold',
                )

                # Panel: Side-by-Side (Original | Overlay)
                ax_panel = fig_slice.add_subplot(1, 2, 1)
                ax_panel.imshow(plt.imread(panel_path))
                if (
                        uses_analysis_overlay
                        and df_tracked is not None
                        and not df_tracked.empty
                        and {"z_slice", "track_id"}.issubset(df_tracked.columns)):
                    tracked_z = pd.to_numeric(df_tracked["z_slice"], errors="coerce")
                    tracked_ids = pd.to_numeric(df_tracked["track_id"], errors="coerce")
                    slice_data = df_tracked[
                        tracked_z.eq(z) & tracked_ids.isin(accepted_track_ids)
                    ].copy()
                    panel_count_label = f"Included track observations: N={len(slice_data)}"
                else:
                    slice_data = df[df['z_slice'] == z].copy()
                    panel_count_label = f"Pre-tracking 2D candidates: N={len(slice_data)}"
                ax_panel.set_title(panel_count_label)
                if uses_analysis_overlay:
                    ax_panel.legend(
                        handles=[
                            Patch(
                                facecolor="#00d940",
                                edgecolor="none",
                                label="Included estimated nucleus observation",
                            )
                        ],
                        loc="lower center",
                        bbox_to_anchor=(0.5, -0.12),
                        ncol=1,
                        fontsize=8,
                        frameon=True,
                    )
                else:
                    ax_panel.text(
                        0.5,
                        -0.04,
                        "Rainbow colors separate candidate IDs; they are not QC categories.",
                        transform=ax_panel.transAxes,
                        ha="center",
                        va="top",
                        fontsize=8,
                    )
                ax_panel.axis('off')

                # Plot: Stats
                ax_hist = fig_slice.add_subplot(1, 2, 2)
                if not slice_data.empty:
                    ax_hist.hist(slice_data['length_um_geodesic'], bins=15, color='skyblue', edgecolor='black')
                    ax_hist.set_title(
                        f"Z={z} Included-Observation 2D Lengths"
                        if uses_analysis_overlay else
                        f"Z={z} Pre-Tracking Candidate 2D Lengths"
                    )
                    ax_hist.set_xlabel("2D geodesic length (um)")
                    ax_hist.set_ylabel("Observations")

                    m_med = slice_data['length_um_geodesic'].median()
                    m_avg = slice_data['length_um_geodesic'].mean()
                    ax_hist.axvline(m_med, color='red', linestyle='-', alpha=0.7, label=f"Median: {m_med:.1f}")
                    ax_hist.axvline(m_avg, color='orange', linestyle='--', alpha=0.7, label=f"Mean: {m_avg:.1f}")
                    ax_hist.legend(fontsize=9)
                else:
                    ax_hist.text(
                        0.5,
                        0.5,
                        "No included track observations"
                        if uses_analysis_overlay else "No detections",
                        ha='center',
                        va='center',
                    )

                pdf.savefig(fig_slice, dpi=300)
                plt.close(fig_slice)

                if gui_callback:
                    gui_callback(int(80 + (20 * (idx_p+1) / len(df_summary))))

        print(f"Report successfully saved to {pdf_path}")

        # --- GENERATE POWERPOINT REPORT ---
        try:
            if generate_pptx:
                generate_pptx_report(out_dir, df, df_summary, um, df_tracks)
        except Exception as e:
            import traceback
            err_msg = traceback.format_exc()
            print(f"PPTX Report failed: {e}")
            write_error_log(out_dir, "PowerPoint Generator (via Batch)", err_msg)

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating PDF report: {e}")
        write_error_log(out_dir, "PDF Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"PDF Report failed to generate completely.\n{e}")
        except Exception:
            pass



def generate_pptx_report(out_dir, df, df_summary, um, df_tracks=None):
    """
    Generates a native Microsoft Office PowerPoint presentation (.pptx) with
    fully editable, data-embedded charts for each key biometric.

    Each chart is a *true* Office Open XML chart object (not a rasterised image),
    so the researcher can re-style, re-colour, and re-export from PowerPoint without
    any additional software.  The chart data tables are embedded inside the PPTX
    file itself, enabling standalone sharing.

    Slide structure
    ---------------
    1. **Global Population Analytics** - side-by-side column histograms of:
       - 2D geodesic length distribution (from the current batch).
       - 3D estimated track length distribution (if tracking was run).
    2. **Population Tracking Consolidation** (if tracking was run) - pie chart
       showing the fraction of single-slice vs. multi-slice reassigned tracks.
    3. **3D Biometrics Dashboard** - scatter plot of tortuosity vs. 3D length,
       plus a bar chart of pitch angle distribution.
    4. **Methods & Calculation Details** - text slide with biological justification
       and mathematical formulae for all metrics.

    Hyperlink behaviour
    -------------------
    The title text on each slide contains a hyperlink to the source Excel workbook.
    Clicking the title in Slide Show mode opens the associated Excel file so the
    viewer can inspect the raw numbers behind any chart.

    Biological context of each chart
    ---------------------------------
    - **Length histogram**: A Gaussian distribution centred near the species-expected
      nuclear length (about 10 um for D. melanogaster in this workflow) supports
      mature elongation. Bimodal distributions may indicate simultaneous maturation
      cohorts, heteromorphic sperm classes, or segmentation/tracking mixtures.
    - **Pie chart**: Quantifies how many tracked objects span multiple Z-slices
      (multi-slice continuity) vs. transient single-slice events (possibly debris).
    - **Tortuosity vs. 3D length scatter**: Healthy spermatids cluster in the
      low-tortuosity, high-length quadrant.  High-tortuosity outliers indicate
      coiled or bent morphologies.
    - **Pitch angle**: A right-skewed distribution toward 90 degrees is expected during
      the apical plunging phase of spermatid elongation.

    Args:
        out_dir (str): Top-level analysis output directory.  The presentation is
            saved as ``<out_dir>/spermatid_analysis_report.pptx``.
        df (pd.DataFrame): Per-spermatid 2D measurement table.
        df_summary (pd.DataFrame): Per-slice summary statistics table.
        um (float): Microns-per-pixel scale factor (``UM_PER_PX_XY``).
        df_tracks (pd.DataFrame, optional): 3D track table.  ``None`` skips
            tracking-specific slides (slides 2 and 3).
    """
    try:
        import os as _os
        import numpy as _np
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from tkinter import messagebox

        print("PPTX: Starting report generation...")
        # We need a fallback blank pptx
        prs = Presentation()
        unet_report = summarize_unet_rescue_for_reports(df, out_dir)

        # Safe blank layout (often index 6 or 5)
        try:
            blank_slide_layout = prs.slide_layouts[6]
        except Exception:
            blank_slide_layout = prs.slide_layouts[0]

        def add_hyperlink(slide, sheet_name="Population_Summary"):
            excel_name = f"batch_analysis_results_{_VERSION}.xlsx"
            # PowerPoint/Excel deep-link fragment: filename.xlsx#Sheet!A1
            target = f"{excel_name}#'{sheet_name}'!A1"

            txBox = slide.shapes.add_textbox(Inches(0.2), Inches(7.1), Inches(9.5), Inches(0.4))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"Click to View Detailed Data: {excel_name} [{sheet_name}]"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0, 0, 255) # Blue
            run.font.underline = True
            try:
                # Shape-level link
                txBox.click_action.hyperlink.address = target
                # Text-run level link
                run.hyperlink.address = target
            except Exception:
                pass

        def add_line_chart(slide, x_data, y_data, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = list(x_data)
            chart_data.add_series('Count', list(y_data))

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.has_major_gridlines = True

        def add_horizontal_bar_chart(slide, categories, values, colors, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', values)

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.size = Pt(8)

            # Show data labels for bar values
            plot = chart.plots[0]
            plot.has_data_labels = True
            for i, point in enumerate(plot.series[0].points):
                point.data_label.font.size = Pt(10)
                point.data_label.font.bold = True

        def add_histogram(slide, data_series, left, top, width, height, title, bins=20):
            if data_series is None or data_series.empty or data_series.isna().all():
                return

            clean_data = data_series.dropna()
            counts, bin_edges = _np.histogram(clean_data, bins=bins)
            avg = clean_data.mean()
            med = clean_data.median()

            # Create category names from bin boundaries
            categories = [f"{bin_edges[i]:.1f}-{bin_edges[i+1]:.1f}" for i in range(len(counts))]

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', list(counts))

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = f"{title} | Mean {avg:.2f}, Median {med:.2f}"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

            # Reduce axis label crowding
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

        # --- Slide 1: canonical biologist-facing results ---
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Biologist Results: Primary 3D Population"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True

        if df_tracks is not None and not df_tracks.empty:
            primary = (
                df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                if "technical_valid" in df_tracks.columns
                else df_tracks.copy()
            )
            metrics_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.0), Inches(5.8))
            metrics_frame = metrics_box.text_frame
            metrics_frame.text = (
                f"Estimated unique nuclei\n{len(primary):,}\n\n"
                f"Median 3D length\n{primary['total_3d_length_um'].median():.2f} um\n\n"
                f"Median maximum 2D length\n{primary['max_length_2d'].median():.2f} um\n\n"
                f"Median effective thickness*\n{primary['thickness_um'].median():.2f} um\n\n"
                f"Median 3D tortuosity\n{primary['tortuosity_3d'].median():.3f}"
            )
            metrics_frame.paragraphs[0].font.size = Pt(16)
            add_histogram(
                slide1,
                primary['total_3d_length_um'],
                Inches(4.8),
                Inches(1.0),
                Inches(4.8),
                Inches(5.2),
                "Primary Population 3D Length",
            )
            note_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9.0), Inches(0.5))
            note_box.text_frame.text = (
                "Primary population: technical-valid 3D tracks. "
                "*Thickness is PSF-sensitive; compare only matched acquisitions."
            )
            note_box.text_frame.paragraphs[0].font.size = Pt(9)

        add_hyperlink(slide1)

        # --- Slide 2: Population Consolidation ---
        if df_tracks is not None and not df_tracks.empty:
            slide2 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "Technical QC - Do Not Use as the Biological Population"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True

            # Left: Reduction Bar Chart (PDF Parity)
            total_2d = len(df)
            total_3d = len(df_tracks)
            n_candidate = (
                int(df_tracks["technical_valid"].astype(bool).sum())
                if "technical_valid" in df_tracks.columns
                else total_3d
            )
            n_warning_only = int(df_tracks["has_warning_only"].sum()) if "has_warning_only" in df_tracks.columns else 0
            n_hard_fail = total_3d - n_candidate

            add_horizontal_bar_chart(slide2,
                                     ['Estimated Unique Nuclei', 'Raw 2D Detections'],
                                     [n_candidate, total_2d],
                                     None, Inches(0.2), Inches(1.5), Inches(4.5), Inches(4.5), "Primary Count and 2D Provenance")

            # Right: Composition Pie Chart
            n_candidate_clean = max(n_candidate - n_warning_only, 0)
            pie_sizes = [n_candidate_clean, n_warning_only, n_hard_fail]
            pie_labels = ['Warning Free', 'Morphology Warning', 'Technical Failure']

            # Filter zero values for pptx chart
            valid_idx = [i for i, v in enumerate(pie_sizes) if v > 0]
            pie_sizes = [pie_sizes[i] for i in valid_idx]
            pie_labels = [pie_labels[i] for i in valid_idx]

            chart_data = CategoryChartData()
            chart_data.categories = pie_labels
            chart_data.add_series('Population', pie_sizes)

            chart2 = slide2.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(4.6), Inches(1.2), Inches(5.2), Inches(5.2), chart_data
            ).chart
            chart2.has_legend = True
            chart2.legend.position = XL_LEGEND_POSITION.CORNER
            chart2.legend.font.size = Pt(8)
            chart2.chart_title.text_frame.text = f"Technical Audit of {total_3d:,} Reconstructed Tracks"
            chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

            plot = chart2.plots[0]
            plot.has_data_labels = True
            total = sum(pie_sizes) if sum(pie_sizes) > 0 else 1
            for i, point in enumerate(plot.series[0].points):
                val = pie_sizes[i]
                pct = (val / total) * 100
                label_text = f"{val:,}\n({pct:.1f}%)"
                point.data_label.text_frame.text = label_text
                point.data_label.font.size = Pt(9)
                point.data_label.font.bold = True

            add_hyperlink(slide2)

        # ---------------------------------------------------------------------
        # SLIDE 3: Advanced 3D biometrics for estimated nuclei
        # ---------------------------------------------------------------------
        if df_tracks is not None and not df_tracks.empty:
            slide3 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            plot_df = (
                df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                if "technical_valid" in df_tracks.columns
                else df_tracks.copy()
            )
            tf.text = f"Advanced 3D Biometrics: {len(plot_df):,} Estimated Unique Nuclei"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True

            add_histogram(slide3, plot_df['pitch_deg'], Inches(0.2), Inches(0.8), Inches(4.5), Inches(2.9), "Pitch Angle (Degrees)", bins=20)
            add_histogram(slide3, plot_df['thickness_um'], Inches(5.0), Inches(0.8), Inches(4.5), Inches(2.9), "Effective Diameter Proxy (\u00b5m)", bins=20)
            add_histogram(slide3, plot_df['taper_ratio'], Inches(0.2), Inches(3.8), Inches(4.5), Inches(2.9), "Morphological Taper Ratio", bins=20)
            add_histogram(slide3, plot_df['nearest_neighbor_um'].dropna(), Inches(5.0), Inches(3.8), Inches(4.5), Inches(2.9), "Nearest-Neighbor Distance (\u00b5m)", bins=20)

            add_hyperlink(slide3, "3D_Track_Audit")

        # ---------------------------------------------------------------------
        # SLIDE 3b: Saturn v5.7 U-Net Rescue Audit
        # ---------------------------------------------------------------------
        if unet_report and unet_report["enabled"]:
            slide_unet = prs.slides.add_slide(blank_slide_layout)
            txBox = slide_unet.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "Saturn v5.7 U-Net Rescue Audit"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True

            table_rows = [
                ["Source", "Count", "Percent"],
                ["Saturn classical", f"{unet_report['saturn_classical']:,}", f"{100.0 * unet_report['saturn_classical'] / max(unet_report['total_2d'], 1):.1f}%"],
                ["U-Net rescued", f"{unet_report['unet_rescued']:,}", f"{100.0 * unet_report['unet_rescued'] / max(unet_report['total_2d'], 1):.1f}%"],
                ["U-Net split/centerline", f"{unet_report['unet_rescued_split']:,}", f"{100.0 * unet_report['unet_rescued_split'] / max(unet_report['total_2d'], 1):.1f}%"],
                ["All U-Net rescued", f"{unet_report['unet_total_rescued']:,}", f"{unet_report['unet_rescue_fraction'] * 100.0:.1f}%"],
            ]
            table_shape = slide_unet.shapes.add_table(len(table_rows), 3, Inches(0.5), Inches(1.0), Inches(4.4), Inches(2.4))
            table = table_shape.table
            for r, row_vals in enumerate(table_rows):
                for c, val in enumerate(row_vals):
                    cell = table.cell(r, c)
                    cell.text_frame.text = val
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if r == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        cell.text_frame.paragraphs[0].font.bold = True

            note_box = slide_unet.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.1), Inches(4.8))
            tf_note = note_box.text_frame
            tf_note.word_wrap = True
            tf_note.text = (
                f"Probability maps saved: {unet_report['probability_map_count']:,}\n"
                f"Rescue-review overlays saved: {unet_report['overlay_count']:,}\n\n"
                "Green overlay: Saturn classical\n"
                "Cyan overlay: accepted U-Net rescue\n"
                "Magenta/orange/red: U-Net-positive candidates rejected by rescue gates\n\n"
                "Overlay dilation is display-only. Counts, lengths, widths, and 3D tracking use the measurement tables, not overlay pixels."
            )
            for p in tf_note.paragraphs:
                p.font.size = Pt(11)

            add_hyperlink(slide_unet, "U-Net_Rescue_Audit")

        # ---------------------------------------------------------------------
        # SLIDE 4: Global Population Statistics Summary Table
        # ---------------------------------------------------------------------
        slide4 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Global Population Statistics Summary"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True

        # Prepare statistics for the one primary estimated-nuclei population.
        stats_rows = [["Metric", "Mean", "Median", "Std Dev"]]
        if not df.empty:
            l2d = df['length_um_geodesic']
            stats_rows.append(["2D Fragment Length (\u00b5m)", f"{l2d.mean():.2f}", f"{l2d.median():.2f}", f"{l2d.std():.2f}"])

        if df_tracks is not None and not df_tracks.empty:
            df_q = (
                df_tracks[df_tracks["technical_valid"].astype(bool)].copy()
                if "technical_valid" in df_tracks.columns
                else df_tracks.copy()
            )
            stats_rows.append(["--- ESTIMATED UNIQUE NUCLEI ---", f"N={len(df_q)}", "", ""])

            def add_pop_rows(pop_df, prefix=""):
                l3 = pop_df['total_3d_length_um']
                z_col = "z_span_um" if "z_span_um" in pop_df.columns else "z_extent_um"
                ze = pop_df[z_col]
                vo = pop_df['volume_um3']
                to = pop_df['tortuosity_3d']
                th = pop_df['thickness_um']
                stats_rows.append([f"{prefix}3D Length (\u00b5m)", f"{l3.mean():.2f}", f"{l3.median():.2f}", f"{l3.std():.2f}"])
                stats_rows.append([f"{prefix}3D Z-Span (\u00b5m)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
                stats_rows.append([f"{prefix}3D Volume (\u00b5m\u00b3)", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
                stats_rows.append([f"{prefix}3D Tortuosity", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
                stats_rows.append([f"{prefix}3D Thickness (\u00b5m)", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])

            add_pop_rows(df_q)

        if len(stats_rows) > 1:
            rows = len(stats_rows)
            cols = 4
            table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(6.0))
            table = table_shape.table

            # Header styling
            for c in range(cols):
                cell = table.cell(0, c)
                cell.text_frame.text = stats_rows[0][c]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196) # Standard Blue
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.bold = True

            # Body styling
            for r in range(1, rows):
                is_separator = "---" in stats_rows[r][0]
                for c in range(cols):
                    cell = table.cell(r, c)
                    cell.text_frame.text = stats_rows[r][c]
                    cell.text_frame.paragraphs[0].font.size = Pt(8)
                    if is_separator:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                        cell.text_frame.paragraphs[0].font.bold = True
                    elif c == 0:
                        cell.text_frame.paragraphs[0].font.bold = True

        add_hyperlink(slide4, "Population_Summary")

        # ---------------------------------------------------------------------
        # SLIDE 5: Methods & Interpretation Guide (Exact PDF Synchronization)
        # ---------------------------------------------------------------------
        slide5 = prs.slides.add_slide(blank_slide_layout)
        txBox_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.text = "Methods, Formulae, Parameter Guide & Audit Details"
        tf_title.paragraphs[0].font.size = Pt(20)
        tf_title.paragraphs[0].font.bold = True

        # Synchronized text from the PDF report version
        methods_items = [
            ("1. Total 3D Geodesic Length (um)", [
                ("Formula: ", "L_3D = sqrt(max(2D geodesic, XY displacement)^2 + z_span^2)"),
                ("Meaning: ", "Projection-length plus Z-span estimate of whole-nucleus 3D length.")
            ]),
            ("2. 3D Euclidean Distance (um)", [
                ("Formula: ", "D_3D = sqrt(XY displacement^2 + z_span^2)"),
                ("Meaning: ", "Straight-line span used as the tortuosity denominator.")
            ]),
            ("3. 3D Tortuosity", [
                ("Formula: ", "T = L_3D / D_3D"),
                ("Meaning: ", "Curvature or over-merge index. Values near 1 are straighter; high values suggest bent or fused tracks.")
            ]),
            ("4. Z-Span and Z-Covered", [
                ("Formula: ", "z_span = (max_z - min_z) * UM_PER_SLICE_Z; z_covered = (max_z - min_z + 1) * UM_PER_SLICE_Z"),
                ("Meaning: ", "Z-span is endpoint-to-endpoint displacement; Z-covered is sampled slab thickness.")
            ]),
            ("5. Approximated 3D Volume (um3)", [
                ("Formula: ", "V_3D = sum(area_est_slice * XY_pixel_area * Z_step)"),
                ("Meaning: ", "PSF- and voxel-sensitive volume estimate. Use mainly for relative comparisons between matched datasets, not as a literal absolute volume.")
            ]),
            ("6. Effective Diameter Proxy (Average Diameter, um)", [
                ("Formula: ", "D_avg = 2 * sqrt((V_3D / L_3D) / pi)"),
                ("Meaning: ", "PSF-sensitive cylinder-equivalent diameter. Useful for relative comparisons under matched imaging settings, not as a literal physical diameter.")
            ]),
            ("7. Pitch Angle (Plunge Vector, Degrees)", [
                ("Formula: ", "theta = abs(arcsin(z_span / D_3D)) * 180/pi"),
                ("Meaning: ", "Absolute plunge angle relative to the imaging plane.")
            ]),
            ("8. Taper Ratio", [
                ("Formula: ", "R = max(area_est across track) / min(area_est across track)"),
                ("Meaning: ", "PSF-sensitive area-derived instability metric. Large values can reflect abrupt narrowing, fusion, or segmentation instability; use mainly for relative comparison.")
            ]),
            ("9. Spatial Packing Density (Nearest Neighbor Dist, um)", [
                ("Formula: ", "nearest 3D centroid-to-centroid distance"),
                ("Meaning: ", "Simple local packing-density readout.")
            ]),
            ("10. Tracking Parameter Provenance", [
                ("Source: ", "Several overlap and conservative tracking thresholds were originally selected by an evolutionary tuning script."),
                ("Goal: ", "The optimizer rewarded multi-slice continuity while penalizing fragmentation, implausible merges, and biology-aware outlier categories."),
                ("Practical use: ", "Treat tuned tracking values as strong starting points. Adjust them only if a new dataset clearly shows fragmentation or over-merging.")
            ]),
            ("11. Technical Track Audit (post-tracking)", [
                ("Audit rules: ", "hard-fail tracks when 3D length > AUDIT_MAX_LENGTH_UM, 3D tortuosity > AUDIT_MAX_TORTUOSITY, effective thickness > AUDIT_EXTREME_THICKNESS_UM, extreme taper > AUDIT_EXTREME_TAPER_RATIO, or track slices < AUDIT_MIN_SLICES. Ordinary thick/taper tracks remain warning-only."),
                ("Meaning: ", "Audit does not change raw detection or tracking. Technical-valid tracks form the estimated-nuclei population; technical failures are excluded."),
                ("Practical use: ", "Use estimated unique nuclei as the one analysis population. Morphology-warning, warning-free, and reference-morphology flags are QC annotations only."),
                ("Biology note: ", "At this acquisition z-step (~1.04 um), single-slice nuclei can be biologically valid because the true mature Drosophila sperm nucleus is much thinner in z than the optical sampling.")
            ]),
            ("12. v5.7 U-Net rescue lane", [
                ("Source labels: ", "Raw 2D detections are labeled as saturn_classical, unet_rescued, or unet_rescued_split in the measurement table."),
                ("Overlay colors: ", "Green is Saturn classical; cyan is accepted U-Net rescue; magenta/orange/red are U-Net-positive candidates rejected by rescue gates."),
                ("Display note: ", "Overlay dilation is cosmetic only and is never used for count, length, width, or 3D tracking calculations.")
            ]),
            ("13. PSF-sensitive metrics note", [
                ("Important: ", "Volume, effective thickness, taper, and other width/area-derived values are broadened by microscope PSF and voxel sampling."),
                ("Use them for: ", "relative comparison between WT and mutant datasets acquired with matched settings."),
                ("Do not use them as: ", "literal physical dimensions or absolute biophysical ground truth.")
            ])
        ]

        top_y = 1.0
        for title, content_list in methods_items:
            tb = slide5.shapes.add_textbox(Inches(0.5), Inches(top_y), Inches(9), Inches(0.72))
            tf = tb.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(10)

            for label, text in content_list:
                p2 = tf.add_paragraph()
                p2.font.size = Pt(8.5)
                if label:
                    run1 = p2.add_run()
                    run1.text = label
                    run1.font.bold = True
                run2 = p2.add_run()
                run2.text = text
                run2.font.bold = False

            top_y += 0.74

        add_hyperlink(slide5, "Population_Summary")

        # Save the presentation matching Batch Name
        final_pptx_name = f"batch_analysis_results_{_VERSION}.pptx"
        output_path = os.path.join(out_dir, final_pptx_name)
        try:
            print(f"PPTX: Saving to {output_path}...")
            prs.save(output_path)
            print(f"PPTX Report successfully saved to: {output_path}")
            return True
        except PermissionError:
            print(f"CRITICAL ERROR: Could not save PowerPoint because the file '{final_pptx_name}' is currently open!")
            return False

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"Failed to generate PPTX report: {e}")
        write_error_log(out_dir, "PowerPoint Generator", err_msg)
        return False

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"Failed to generate PPTX report: {e}")
        write_error_log(out_dir, "PowerPoint Generator", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showerror("Reporting Error", f"Failed to generate PowerPoint Report:\n{e}\n\nSee report_generation_errors.txt for details.")
        except Exception:
            pass
        return False




# =============================================================================
# ROI GUI (single-file v9)
# =============================================================================
import traceback
try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
    from matplotlib.figure import Figure
    _TK_AVAILABLE = True
except Exception:
    _TK_AVAILABLE = False

PARAM_SECTIONS = {
    "Calibration & Scale": [
        "UM_PER_PX_XY", "UM_PER_SLICE_Z"
    ],
    "2.5D U-Net Integration": [
        "SEGMENTATION_ENGINE", "UNET_MODEL_PATH", "UNET_THRESHOLD_MODE",
        "UNET_CANDIDATE_THRESHOLD", "UNET_SEED_THRESHOLD", "UNET_CONTEXT_MODE",
        "UNET_INFERENCE_MODE", "UNET_TILE_SIZE", "UNET_TILE_OVERLAP",
        "UNET_TILE_BATCH_SIZE", "UNET_ROI_PADDING_PX", "UNET_FAIL_HARD",
        "UNET_RESCUE_ENABLE", "UNET_RESCUE_THRESHOLD",
        "UNET_RESCUE_HYSTERESIS_ENABLE",
        "UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS",
        "UNET_RESCUE_EXCLUDE_DILATION_PX",
        "UNET_RESCUE_MIN_COMPONENT_PX", "UNET_RESCUE_MIN_SKEL_LEN_UM",
        "UNET_SHORT_RESCUE_MIN_MEAN_PROB",
        "UNET_RESCUE_MAX_ADDITIONS_PER_SLICE",
        "UNET_RESCUE_SPLIT_RETRY_ENABLE", "UNET_RESCUE_SPLIT_THRESHOLDS",
        "UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE", "UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB",
        "UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB", "UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM",
        "UNET_INSTANCE_SPLIT_ENABLE", "UNET_INSTANCE_SEED_THRESHOLD",
        "UNET_INSTANCE_PEAK_MIN_DISTANCE_PX", "UNET_INSTANCE_WATERSHED_COMPACTNESS",
        "UNET_PRIMARY_MIN_COMPONENT_PX",
        "UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE",
        "UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX",
        "UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS",
        "UNET_PRIMARY_SAVE_FILLED_MASK_OVERLAY",
        "UNET_PRIMARY_SAVE_INSTANCE_OVERLAY",
        "UNET_TRACKING_SUPPORT", "ASSIGNMENT_UNET_SUPPORT_WEIGHT",
        "ASSIGNMENT_UNET_CONTINUITY_WEIGHT", "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT"
    ],
    "Image Enhancement": [
        "CLAHE_CLIP", "CLAHE_KERNEL", "BG_SIGMA", "THRESHOLD_HI", "THRESHOLD_LO"
    ],
    "Binary Mask Cleanup": [
        "CLOSE_RADIUS", "MIN_HOLE_AREA", "MIN_OBJ_PX"
    ],
    "Skeleton Cleanup & Splitting": [
        "MAX_BRIDGE_PX", "MAX_BRANCH_LEN_PX", "BREAK_JUNCTIONS"
    ],
    "Early Shape Filter": [
        "USE_EARLY_SHAPE_FILTER", "MIN_ECCENTRICITY", "MAX_MINOR_PX", "MIN_AXIS_RATIO", "MIN_MAJOR_PX"
    ],
    "Post-Skeleton Acceptance Filters": [
        "MIN_SKEL_LEN_PX", "MAX_GEODESIC_LEN_PX", "MAX_WIDTH_PX", "MIN_LENGTH_WIDTH_RATIO",
        "MAX_BRANCH_NODES", "MAX_TORTUOSITY", "MAX_ENDPOINT_COUNT", "ALLOW_LOOPS"
    ],
    "3D Tracking": [
        "DO_TRACKING", "TRACKING_BACKEND", "TRACK_MAX_DIST_UM", "TRACK_MAX_GAP_SLICES", "TRACK_BBOX_PADDING_PX",
        "TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM",
        "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO", "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO",
        "CONSERVATIVE_MAX_AREA_JUMP_RATIO", "CONSERVATIVE_MAX_TORTUOSITY_JUMP",
        "CONSERVATIVE_MAX_CENTROID_JUMP_UM"
    ],
    "Overlap-First Continuation": [
        "OVERLAP_STABILITY_THRESHOLD", "OVERLAP_ORIENTATION_DEG", "OVERLAP_MULTIPLIER", "OVERLAP_MIN_STABLE_COUNT"
    ],
    "3D Tracking - Global Assignment Prototype": [
        "ASSIGNMENT_MAX_COST", "ASSIGNMENT_DIST_WEIGHT", "ASSIGNMENT_OVERLAP_WEIGHT",
        "ASSIGNMENT_LENGTH_WEIGHT", "ASSIGNMENT_WIDTH_WEIGHT", "ASSIGNMENT_AREA_WEIGHT",
        "ASSIGNMENT_ANGLE_WEIGHT"
    ],
    "3D Tracking - Hybrid Fragment Repair": [
        "HYBRID_REPAIR_MAX_COST", "HYBRID_REPAIR_MAX_GAP_SLICES", "HYBRID_REPAIR_MAX_FRAGMENT_SLICES",
        "HYBRID_REPAIR_MAX_LINK_DIST_UM", "HYBRID_REPAIR_MIN_OVERLAP", "HYBRID_REPAIR_MAX_FINAL_LENGTH_UM"
    ],
    "Candidate Audit (post-tracking)": [
        "AUDIT_MAX_LENGTH_UM", "AUDIT_MAX_TORTUOSITY", "AUDIT_MAX_THICKNESS_UM", "AUDIT_MAX_TAPER_RATIO",
        "AUDIT_EXTREME_THICKNESS_UM", "AUDIT_EXTREME_TAPER_RATIO", "AUDIT_MIN_SLICES"
    ],
}

PARAM_TITLES = {
    "UM_PER_PX_XY": "XY pixel size (um / pixel)",
    "UM_PER_SLICE_Z": "Z step size (um / slice)",
    "CLAHE_CLIP": "CLAHE contrast limit",
    "CLAHE_KERNEL": "CLAHE tile size",
    "BG_SIGMA": "Background subtraction sigma",
    "THRESHOLD_HI": "High ridge threshold percentile",
    "THRESHOLD_LO": "Low ridge threshold percentile",
    "CLOSE_RADIUS": "Mask closing radius",
    "MIN_HOLE_AREA": "Fill holes up to this size (px)",
    "MIN_OBJ_PX": "Minimum object size (px)",
    "MAX_BRIDGE_PX": "Maximum skeleton bridge gap (px)",
    "MAX_BRANCH_LEN_PX": "Maximum spur length to prune (px)",
    "BREAK_JUNCTIONS": "Break skeleton junctions",
    "USE_EARLY_SHAPE_FILTER": "Use early blob / round-shape filter",
    "MIN_ECCENTRICITY": "Minimum eccentricity",
    "MAX_MINOR_PX": "Maximum minor axis (px)",
    "MIN_AXIS_RATIO": "Minimum major/minor ratio",
    "MIN_MAJOR_PX": "Minimum major axis (px)",
    "MIN_SKEL_LEN_PX": "Minimum 2D skeleton length (px)",
    "MAX_GEODESIC_LEN_PX": "Maximum 2D skeleton length (px)",
    "MAX_WIDTH_PX": "Maximum 2D width (px)",
    "MIN_LENGTH_WIDTH_RATIO": "Minimum length/width ratio",
    "MAX_BRANCH_NODES": "Maximum branch points",
    "MAX_TORTUOSITY": "Maximum 2D tortuosity",
    "MAX_ENDPOINT_COUNT": "Maximum endpoints",
    "ALLOW_LOOPS": "Allow looped skeletons",
    "DO_TRACKING": "Enable 3D tracking",
    "TRACKING_BACKEND": "Tracking backend",
    "TRACK_MAX_DIST_UM": "Maximum slice-to-slice centroid jump (um)",
    "TRACK_MAX_GAP_SLICES": "Maximum missing slices in one track",
    "TRACK_BBOX_PADDING_PX": "Bounding-box overlap padding (px)",
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": "Maximum width change ratio",
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": "Maximum length change ratio",
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": "Maximum area change ratio",
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": "Maximum tortuosity change",
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": "Hard centroid jump stop-rule (um)",
    "OVERLAP_STABILITY_THRESHOLD": "Overlap stability tolerance",
    "OVERLAP_ORIENTATION_DEG": "Maximum stable orientation difference (deg)",
    "OVERLAP_MULTIPLIER": "Overlap relaxation multiplier",
    "OVERLAP_MIN_STABLE_COUNT": "Minimum stable features for overlap continuation",
    "ASSIGNMENT_MAX_COST": "Assignment maximum accepted cost",
    "ASSIGNMENT_DIST_WEIGHT": "Assignment centroid-distance weight",
    "ASSIGNMENT_OVERLAP_WEIGHT": "Assignment non-overlap penalty weight",
    "ASSIGNMENT_LENGTH_WEIGHT": "Assignment length-change weight",
    "ASSIGNMENT_WIDTH_WEIGHT": "Assignment width-change weight",
    "ASSIGNMENT_AREA_WEIGHT": "Assignment area-change weight",
    "ASSIGNMENT_ANGLE_WEIGHT": "Assignment orientation-change weight",
    "HYBRID_REPAIR_MAX_COST": "Hybrid repair maximum accepted cost",
    "HYBRID_REPAIR_MAX_GAP_SLICES": "Hybrid repair maximum missing slices",
    "HYBRID_REPAIR_MAX_FRAGMENT_SLICES": "Hybrid repair maximum fragment size",
    "HYBRID_REPAIR_MAX_LINK_DIST_UM": "Hybrid repair maximum link distance (um)",
    "HYBRID_REPAIR_MIN_OVERLAP": "Hybrid repair minimum bounding-box overlap",
    "HYBRID_REPAIR_MAX_FINAL_LENGTH_UM": "Hybrid repair maximum merged 3D length (um)",
    "TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM": "Maximum physically joined nucleus length (um)",
    "SEGMENTATION_ENGINE": "Segmentation engine",
    "UNET_MODEL_PATH": "U-Net checkpoint path",
    "UNET_THRESHOLD_MODE": "U-Net threshold mode",
    "UNET_CANDIDATE_THRESHOLD": "U-Net candidate threshold",
    "UNET_SEED_THRESHOLD": "U-Net seed threshold",
    "UNET_CONTEXT_MODE": "U-Net 2.5D context",
    "UNET_INFERENCE_MODE": "U-Net inference mode",
    "UNET_TILE_SIZE": "U-Net tile size (px)",
    "UNET_TILE_OVERLAP": "U-Net tile overlap (px)",
    "UNET_TILE_BATCH_SIZE": "U-Net tile batch size",
    "UNET_ROI_PADDING_PX": "U-Net ROI padding (px)",
    "UNET_FAIL_HARD": "Stop run if U-Net inference fails",
    "UNET_RESCUE_ENABLE": "Enable U-Net rescue lane",
    "UNET_RESCUE_THRESHOLD": "U-Net rescue seed threshold",
    "UNET_RESCUE_HYSTERESIS_ENABLE": "Connect low-confidence support to rescue seeds",
    "UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS": "Retain unusual U-Net morphology for measurement",
    "UNET_RESCUE_EXCLUDE_DILATION_PX": "Rescue exclusion around Saturn hits (px)",
    "UNET_RESCUE_MIN_COMPONENT_PX": "Minimum rescue component size (px)",
    "UNET_RESCUE_MIN_SKEL_LEN_UM": "Minimum resolvable U-Net centerline (um)",
    "UNET_SHORT_RESCUE_MIN_MEAN_PROB": "Short-nucleus U-Net confidence",
    "UNET_RESCUE_MAX_ADDITIONS_PER_SLICE": "Maximum rescued detections per slice",
    "UNET_RESCUE_SPLIT_RETRY_ENABLE": "Retry splitting rejected U-Net candidates",
    "UNET_RESCUE_SPLIT_THRESHOLDS": "U-Net split retry thresholds",
    "UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE": "Salvage red U-Net centerlines",
    "UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB": "Minimum salvage mean probability",
    "UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB": "Low-ratio rescue confidence",
    "UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM": "Low-ratio rescue minimum length (um)",
    "UNET_INSTANCE_SPLIT_ENABLE": "Split U-Net probability instances first",
    "UNET_INSTANCE_SEED_THRESHOLD": "U-Net instance seed threshold",
    "UNET_INSTANCE_PEAK_MIN_DISTANCE_PX": "Minimum seed peak spacing (px)",
    "UNET_INSTANCE_WATERSHED_COMPACTNESS": "Watershed compactness",
    "UNET_PRIMARY_MIN_COMPONENT_PX": "U-Net-primary minimum component size (px)",
    "UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE": "Add residual Saturn-only detections",
    "UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX": "Saturn residual exclusion dilation (px)",
    "UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS": "Retain U-Net-primary morphology warnings",
    "UNET_PRIMARY_SAVE_FILLED_MASK_OVERLAY": "Save U-Net-primary filled-mask overlay",
    "UNET_PRIMARY_SAVE_INSTANCE_OVERLAY": "Save U-Net-primary instance overlay",
    "UNET_TRACKING_SUPPORT": "Use U-Net evidence for 3D linking",
    "ASSIGNMENT_UNET_SUPPORT_WEIGHT": "Assignment U-Net support penalty",
    "ASSIGNMENT_UNET_CONTINUITY_WEIGHT": "Assignment U-Net continuity penalty",
    "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT": "Hybrid repair U-Net support penalty",
    "AUDIT_MAX_LENGTH_UM": "Audit: maximum 3D length (um)",
    "AUDIT_MAX_TORTUOSITY": "Audit: maximum 3D tortuosity",
    "AUDIT_MAX_THICKNESS_UM": "Audit: maximum effective thickness (um, PSF-sensitive)",
    "AUDIT_MAX_TAPER_RATIO": "Audit: maximum taper ratio (PSF-sensitive)",
    "AUDIT_EXTREME_THICKNESS_UM": "Candidate hard-fail: extreme thickness (um)",
    "AUDIT_EXTREME_TAPER_RATIO": "Candidate hard-fail: extreme taper ratio",
    "AUDIT_MIN_SLICES": "Audit: minimum slices required (1 recommended here)"
}

PARAM_DESCRIPTIONS = {
    "UM_PER_PX_XY": "What it affects: all lateral measurements. Increase only if the microscope pixel size really is larger. Wrong values scale every reported x/y length, width, and distance.",
    "UM_PER_SLICE_Z": "What it affects: z-span, 3D length, pitch angle, and volume. Increase only if the physical slice spacing is larger. This does not change segmentation, only measurement scaling.",
    "CLAHE_CLIP": "What it affects: local contrast before thresholding. Increase to pull out dim nuclei, but too high amplifies noise and halos. Decrease if background texture starts getting segmented.",
    "CLAHE_KERNEL": "What it affects: how local the contrast correction is. Smaller values react to finer local contrast; larger values give smoother correction across the field.",
    "BG_SIGMA": "What it affects: background haze removal. Increase to subtract broader blur/haze; decrease if real diffuse signal is being removed together with background.",
    "THRESHOLD_HI": "What it affects: the strongest ridge pixels used as confident seeds. Raise it for stricter detection; lower it to keep weaker nuclei. High impact on raw detection counts.",
    "THRESHOLD_LO": "What it affects: how far masks expand outward from the confident seeds. Raise it for tighter masks; lower it for more inclusive masks. Must stay below THRESHOLD_HI.",
    "CLOSE_RADIUS": "What it affects: tiny breaks in the thresholded mask. Increase slightly to reconnect small gaps; too high can fuse neighboring nuclei.",
    "MIN_HOLE_AREA": "What it affects: small dark holes inside detected objects. Increase to fill more internal holes; decrease if thin structures are being overfilled.",
    "MIN_OBJ_PX": "What it affects: debris removal. Increase to discard more tiny specks; decrease only if real nuclei are being lost because they are very small.",
    "MAX_BRIDGE_PX": "What it affects: skeleton fragmentation. Increase to reconnect nearby broken skeleton ends; too high can join neighboring nuclei and create long false chains.",
    "MAX_BRANCH_LEN_PX": "What it affects: tiny branch spurs. Increase to prune more side spikes; too high can shorten real tips.",
    "BREAK_JUNCTIONS": "What it affects: dense skeleton webs. True = cut branch intersections so complex tangles split into simpler strands. This is high impact and can over-fragment if used too aggressively.",
    "USE_EARLY_SHAPE_FILTER": "What it affects: pre-skeleton removal of round blobs. Turn on only if round debris is a major problem. It does not fix tracking; it changes what enters tracking.",
    "MIN_ECCENTRICITY": "Early shape filter only. Higher values keep only more elongated objects; lower values allow rounder objects to survive.",
    "MAX_MINOR_PX": "Early shape filter only. Lower values enforce thinner objects; higher values allow thicker blobs through.",
    "MIN_AXIS_RATIO": "Early shape filter only. Higher values require stronger elongation; lower values allow shorter/fatter objects.",
    "MIN_MAJOR_PX": "Early shape filter only. Rejects very short components before skeletonization.",
    "MIN_SKEL_LEN_PX": "What it affects: short fragment rejection in 2D. Increase to remove tiny fragments and debris; decrease if real short nuclei fragments are being lost.",
    "MAX_GEODESIC_LEN_PX": "What it affects: very long 2D fragments. Lower values reject long merged chains; higher values keep more long objects but risk allowing fused webs.",
    "MAX_WIDTH_PX": "What it affects: thick 2D fragments. Lower values remove broad fused objects; higher values allow thicker candidates to pass.",
    "MIN_LENGTH_WIDTH_RATIO": "What it affects: how rod-like an object must be. Increase to reject blobs/fusions; decrease if real nuclei are slightly thicker or shorter than expected.",
    "MAX_BRANCH_NODES": "What it affects: branching complexity. Lower values reject tangled skeletons; higher values tolerate more branching and may admit fused objects.",
    "MAX_TORTUOSITY": "What it affects: very curved 2D fragments. Lower values are stricter; higher values tolerate more bending. Useful for removing tangled chains.",
    "MAX_ENDPOINT_COUNT": "What it affects: cluttered skeleton topology. Lower values reject objects with many tips; higher values tolerate more fragmented-looking shapes.",
    "ALLOW_LOOPS": "If True, looped skeletons can still be measured. Usually safe to leave on unless loops are a known artifact source.",
    "DO_TRACKING": "What it affects: whether 2D detections are linked into 3D tracks. Turn off only for debugging pure per-slice detection.",
    "TRACKING_BACKEND": "Which tracking engine to use. v5.5 default is hybrid_repair: legacy tracking first, followed by a conservative fragment-repair pass. Other options are legacy and global_assignment.",
    "TRACK_MAX_DIST_UM": "What it affects: fragmentation versus false merging across z. Tuned by the evolutionary optimizer. Increase if the same nucleus is failing to link between slices; decrease if neighboring nuclei are being fused into one track. This is one of the most important tracking parameters.",
    "TRACK_MAX_GAP_SLICES": "What it affects: tolerance for missed slices. Increase to allow tracks to bridge one or more blank slices; too high can stitch unrelated objects.",
    "TRACK_BBOX_PADDING_PX": "What it affects: overlap matching sensitivity. Tuned by the evolutionary optimizer. Increase to count near-touching detections as overlapping; too high can make crowded nuclei look like one track.",
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": "What it affects: whether sudden width changes are allowed inside a track. Tuned by the evolutionary optimizer. Increase if real tracks are breaking on modest width changes; decrease if fused tracks are slipping through.",
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": "What it affects: whether sudden length changes are allowed inside a track. Tuned by the evolutionary optimizer. Lower values reduce false merges; higher values reduce fragmentation.",
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": "What it affects: how much 2D area can change from one slice to the next. Tuned by the evolutionary optimizer. Useful for preventing track jumps in crowded regions.",
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": "What it affects: abrupt curvature change between linked detections. Increase only if real nuclei bend more than expected across slices.",
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": "Hard stop-rule for large centroid jumps when overlap evidence is weak. Tuned by the evolutionary optimizer. Increase carefully if tracking is fragmented; too high can cause track hopping.",
    "OVERLAP_STABILITY_THRESHOLD": "What it affects: how similar width/area/length must remain when overlap suggests the same object. Tuned by the evolutionary optimizer. Higher values reduce fragmentation; lower values are stricter.",
    "OVERLAP_ORIENTATION_DEG": "What it affects: how much the local orientation can rotate between slices while still counting as stable. Tuned by the evolutionary optimizer. Increase if slight angular jitter is breaking tracks.",
    "OVERLAP_MULTIPLIER": "What it affects: how much continuation rules are relaxed when overlap is strong and the object looks stable. Tuned by the evolutionary optimizer. Higher values are more forgiving, but can increase false continuation.",
    "OVERLAP_MIN_STABLE_COUNT": "What it affects: how many stable features must agree before an overlapping candidate continues a track. Lower values reduce fragmentation; higher values are stricter.",
    "ASSIGNMENT_MAX_COST": "Global-assignment tracker only. Candidate links above this total cost are rejected. Lower values reduce false links; higher values reduce fragmentation.",
    "ASSIGNMENT_DIST_WEIGHT": "Global-assignment and hybrid repair weight for centroid displacement between slices.",
    "ASSIGNMENT_OVERLAP_WEIGHT": "Global-assignment and hybrid repair penalty for weak/no bounding-box overlap.",
    "ASSIGNMENT_LENGTH_WEIGHT": "Global-assignment and hybrid repair penalty for sudden 2D length changes.",
    "ASSIGNMENT_WIDTH_WEIGHT": "Global-assignment and hybrid repair penalty for sudden width changes.",
    "ASSIGNMENT_AREA_WEIGHT": "Global-assignment and hybrid repair penalty for sudden area changes.",
    "ASSIGNMENT_ANGLE_WEIGHT": "Global-assignment and hybrid repair penalty for orientation changes between slices.",
    "HYBRID_REPAIR_MAX_COST": "V5.6 ROI-ADAPTIVE repair only. Candidate fragment merges above this cost are rejected. Lower values are safer and leave more fragments split.",
    "HYBRID_REPAIR_MAX_GAP_SLICES": "V5.6 ROI-ADAPTIVE repair only. Maximum blank slices allowed between two fragments being considered for repair.",
    "HYBRID_REPAIR_MAX_FRAGMENT_SLICES": "V5.6 ROI-ADAPTIVE repair only. At least one side of a repaired link must have this many slices or fewer, so the pass targets fragments rather than rewriting stable tracks.",
    "HYBRID_REPAIR_MAX_LINK_DIST_UM": "V5.6 ROI-ADAPTIVE repair only. Hard distance gate for fragment merges unless there is direct bounding-box overlap evidence.",
    "HYBRID_REPAIR_MIN_OVERLAP": "V5.6 ROI-ADAPTIVE repair only. Preferred minimum bounding-box overlap for a repair link. Very close non-overlap links can still pass if the total cost is low.",
    "HYBRID_REPAIR_MAX_FINAL_LENGTH_UM": "V5.6 ROI-ADAPTIVE repair only. Rejects a proposed merge if the estimated merged 3D length would exceed this value.",
    "TRACK_TECHNICAL_MAX_JOINED_LENGTH_UM": "Hard physical guard used by every tracking backend. A proposed cross-slice link is rejected when it would create a reconstructed nucleus longer than this value; individual long observations remain visible for review.",
    "SEGMENTATION_ENGINE": "classical_saturn preserves Saturn segmentation; hybrid and unet_assisted add U-Net support; unet_primary makes seed-connected U-Net probability instances authoritative.",
    "UNET_MODEL_PATH": "Path to a trained 2.5D U-Net checkpoint. Leave blank to keep Saturn fully classical.",
    "UNET_THRESHOLD_MODE": "soft keeps U-Net output as probability evidence; hard turns the probability map into a binary candidate mask. Soft is safer while the model is still being validated.",
    "UNET_CANDIDATE_THRESHOLD": "Low probability cutoff for inclusive U-Net candidate support. Lower values recover faint nuclei but produce more candidates for downstream biological QC.",
    "UNET_SEED_THRESHOLD": "Higher probability cutoff for confident U-Net seed support. This should usually stay above the candidate threshold.",
    "UNET_CONTEXT_MODE": "How neighboring z-slices are presented to the model. z_minus_z_z_plus uses previous/current/next slices as a 2.5D input.",
    "UNET_INFERENCE_MODE": "roi_tiled runs U-Net only on ROI-aware tiles and stitches probabilities back into full-frame coordinates.",
    "UNET_TILE_SIZE": "Tile width/height sent to the U-Net. Smaller tiles zoom into local detail less by themselves, but reduce memory and keep nuclei prominent in each crop.",
    "UNET_TILE_OVERLAP": "Overlap between tiles. More overlap reduces edge artifacts during stitched inference but costs more GPU time.",
    "UNET_TILE_BATCH_SIZE": "How many ROI tiles are sent through the U-Net at once. Increase for faster GPU inference if memory allows; lower it if you hit CUDA memory errors.",
    "UNET_ROI_PADDING_PX": "Extra context around the selected ROI when preparing U-Net tiles. Helps avoid boundary artifacts without letting off-ROI tissue influence output.",
    "UNET_FAIL_HARD": "Required safety gate for hybrid analysis. If model loading or inference fails, stop the run instead of silently substituting classical-only segmentation.",
    "UNET_RESCUE_ENABLE": "If enabled, U-Net high-probability regions not already covered by accepted Saturn detections are skeletonized and measured as a separate rescue lane.",
    "UNET_RESCUE_THRESHOLD": "High-confidence seed cutoff for the rescue lane. With hysteresis enabled, connected support extends down to UNET_CANDIDATE_THRESHOLD.",
    "UNET_RESCUE_HYSTERESIS_ENABLE": "If enabled, U-Net candidate pixels at the low candidate threshold are retained only when connected to a high-confidence rescue seed. This recovers faint boundaries without admitting isolated low-probability noise.",
    "UNET_RESCUE_RETAIN_MORPHOLOGY_WARNINGS": "If enabled, unusual width, length-to-width ratio, or tortuosity is retained and reported as morphology for biological comparison. Technical failures such as unresolved loops, branches, excess endpoints, and objects above the 20 um guard remain excluded.",
    "UNET_RESCUE_EXCLUDE_DILATION_PX": "How far to dilate accepted Saturn skeletons before searching for U-Net-only missed detections. Increase to avoid duplicate detections around existing nuclei.",
    "UNET_RESCUE_MIN_COMPONENT_PX": "Minimum binary component size before U-Net rescue skeletonization. Increase to suppress tiny U-Net specks; decrease to recover very faint fragments.",
    "UNET_RESCUE_MIN_SKEL_LEN_UM": "Technical resolution floor for U-Net-only centerlines. Keep this well below expected biological nucleus lengths so genuinely short WT or mutant nuclei are retained; increase only when confirmed pixel-scale specks are being accepted.",
    "UNET_SHORT_RESCUE_MIN_MEAN_PROB": "Mean U-Net probability required to retain a centerline below the technical length floor. This confidence exception prevents expected WT or mutant length from becoming an acceptance rule.",
    "UNET_RESCUE_MAX_ADDITIONS_PER_SLICE": "Optional cap on rescued objects per slice. Set to 0 for no cap. Use only if visual review shows the rescue lane is too permissive.",
    "UNET_RESCUE_SPLIT_RETRY_ENABLE": "If enabled, U-Net rescue candidates rejected as long, branched, looped, tortuous, or endpoint-heavy are retried at stricter probability-core thresholds before final rejection.",
    "UNET_RESCUE_SPLIT_THRESHOLDS": "Probability core thresholds used during split retry. Higher thresholds can separate connected U-Net regions into cleaner individual nuclei before biological QC.",
    "UNET_RESCUE_CENTERLINE_SALVAGE_ENABLE": "If enabled, high-confidence red U-Net candidates rejected for topology are reduced to their longest simple centerline and measured once more.",
    "UNET_RESCUE_CENTERLINE_MIN_MEAN_PROB": "Minimum mean U-Net probability required before topology-rejected red candidates can be centerline-salvaged.",
    "UNET_LOW_RATIO_RESCUE_MIN_MEAN_PROB": "Allows a low length-to-width U-Net candidate only when its mean model probability reaches this value. This avoids globally weakening the morphology rule.",
    "UNET_LOW_RATIO_RESCUE_MIN_LENGTH_UM": "Technical minimum centerline length for the high-confidence low-ratio exception. The default 4 um excludes tiny fragments while remaining below expected WT and mutant nucleus lengths.",
    "UNET_INSTANCE_SPLIT_ENABLE": "If enabled, connected U-Net rescue probability regions are split into putative instances before skeletonization and measurement.",
    "UNET_INSTANCE_SEED_THRESHOLD": "Probability threshold for watershed/core seeds used to split connected U-Net rescue regions. Higher values create cleaner but fewer seeds.",
    "UNET_INSTANCE_PEAK_MIN_DISTANCE_PX": "Minimum distance between fallback U-Net probability peaks used as instance seeds. Lower values can split crowded regions more aggressively.",
    "UNET_INSTANCE_WATERSHED_COMPACTNESS": "Compactness term for watershed instance splitting. Larger values favor compact regions; near-zero follows probability topology more closely.",
    "UNET_PRIMARY_MIN_COMPONENT_PX": "Technical pixel-noise floor used before U-Net-primary instance measurement. It is not a biological length or shape gate.",
    "UNET_PRIMARY_CLASSICAL_ADDITIONS_ENABLE": "When enabled, Saturn may add detections only in residual space outside accepted U-Net-primary masks. It cannot remove or alter U-Net instances.",
    "UNET_PRIMARY_CLASSICAL_EXCLUDE_DILATION_PX": "Display-independent exclusion margin around accepted U-Net-primary masks before optional Saturn-only additions are searched.",
    "UNET_PRIMARY_RETAIN_MORPHOLOGY_WARNINGS": "Keep short, wide, low-ratio, curved, or tortuous U-Net-primary instances and report those properties as warnings.",
    "UNET_PRIMARY_SAVE_FILLED_MASK_OVERLAY": "Save a filled-mask review overlay for U-Net-primary output.",
    "UNET_PRIMARY_SAVE_INSTANCE_OVERLAY": "Save a uniquely colored filled-instance review overlay for U-Net-primary output.",
    "UNET_TRACKING_SUPPORT": "If enabled in hybrid/U-Net mode, per-detection U-Net probabilities can reduce confidence in weak 3D links and favor links with consistent model support.",
    "ASSIGNMENT_UNET_SUPPORT_WEIGHT": "Global-assignment penalty for linking detections with weak U-Net support. Set to 0 to ignore U-Net evidence during assignment.",
    "ASSIGNMENT_UNET_CONTINUITY_WEIGHT": "Global-assignment penalty for abrupt U-Net probability changes across adjacent slices.",
    "HYBRID_REPAIR_UNET_SUPPORT_WEIGHT": "Hybrid fragment-repair penalty for repairing tracks through weak U-Net support. Set to 0 if U-Net evidence is too conservative.",
    "AUDIT_MAX_LENGTH_UM": "Audit only. Tracks longer than this are flagged after tracking. This does not change segmentation or tracking itself; it changes the audit-passed subset.",
    "AUDIT_MAX_TORTUOSITY": "Audit only. Flags unusually curved 3D tracks. Lower values make the quality set stricter; higher values keep more bent nuclei.",
    "AUDIT_MAX_THICKNESS_UM": "Audit only. Flags tracks that look too thick for a nucleus. PSF-sensitive: use mainly for relative WT-versus-mutant comparison under matched imaging settings rather than as a literal physical diameter cutoff.",
    "AUDIT_MAX_TAPER_RATIO": "Audit only. Flags tracks with extreme change from thickest to thinnest slice. PSF-sensitive and area-derived: useful for instability screening and relative comparison, not as a literal anatomical ratio.",
    "AUDIT_EXTREME_THICKNESS_UM": "Biological-candidate audit only. Tracks above this very high effective-thickness threshold hard-fail the candidate tier; ordinary thick tracks remain warning-only.",
    "AUDIT_EXTREME_TAPER_RATIO": "Biological-candidate audit only. Tracks above this very high taper threshold hard-fail the candidate tier; ordinary taper tracks remain warning-only.",
    "AUDIT_MIN_SLICES": "Audit only. Minimum number of slices required to avoid the single-slice flag. For this Leica SP8 stack (z-step ~1.04 um), set to 1 because mature Drosophila sperm nuclei can be biologically valid even when they appear in a single optical slice."
}

PARAM_TITLES.update({
    "UM_PER_PX_XY": "XY pixel size (um / pixel)",
    "UM_PER_SLICE_Z": "Z step size (um / slice)",
    "TRACK_MAX_DIST_UM": "Maximum slice-to-slice centroid jump (um)",
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": "Hard centroid jump stop-rule (um)",
    "AUDIT_MAX_LENGTH_UM": "Audit: maximum 3D length (um)",
    "AUDIT_MAX_THICKNESS_UM": "Audit: maximum effective thickness (um, PSF-sensitive)",
    "AUDIT_EXTREME_THICKNESS_UM": "Candidate hard-fail: extreme thickness (um)",
    "AUDIT_EXTREME_TAPER_RATIO": "Candidate hard-fail: extreme taper ratio",
})

PARAM_DESCRIPTIONS.update({
    "UM_PER_SLICE_Z": (
        "What it affects: z-span, z-covered, 3D length, pitch angle, and volume. "
        "Increase only if the physical slice spacing is larger. This does not change "
        "segmentation, only measurement scaling."
    ),
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": (
        "What it affects: how much estimated 2D area can change from one slice to the next. "
        "Area is estimated as geodesic length times median width, not raw skeleton pixel count. "
        "Tuned by the evolutionary optimizer and useful for preventing track jumps in crowded regions."
    ),
    "OVERLAP_STABILITY_THRESHOLD": (
        "What it affects: how similar width, estimated area, and length must remain when overlap "
        "suggests the same object. Higher values reduce fragmentation; lower values are stricter."
    ),
    "AUDIT_MAX_LENGTH_UM": (
        "Audit only. Tracks longer than this projection-length-plus-Z 3D estimate are flagged after "
        "tracking. This does not change segmentation or tracking itself; it changes the audit-passed subset."
    ),
    "AUDIT_MAX_THICKNESS_UM": (
        "Audit only. Flags tracks with a large cylinder-equivalent thickness derived from estimated volume "
        "and 3D length. This is PSF-sensitive, so use mainly for relative comparisons under matched imaging "
        "settings rather than as a literal physical diameter cutoff."
    ),
    "AUDIT_MAX_TAPER_RATIO": (
        "Audit only. Flags tracks with extreme max/min estimated area across slices. Area is derived from "
        "length times median width, so this is PSF-sensitive and best used for instability screening."
    ),
    "AUDIT_MIN_SLICES": (
        "Audit only. Minimum number of detected slices required to avoid the single-slice flag. For this "
        "Leica SP8 stack (z-step ~1.04 um), set to 1 because mature Drosophila sperm nuclei can be "
        "biologically valid even when they appear in a single optical slice."
    ),
    "AUDIT_EXTREME_THICKNESS_UM": (
        "Candidate tier only. Effective thickness above this high threshold becomes a hard fail. "
        "Values above AUDIT_MAX_THICKNESS_UM but below this remain warning-only because thickness is PSF-sensitive."
    ),
    "AUDIT_EXTREME_TAPER_RATIO": (
        "Candidate tier only. Taper above this high threshold becomes a hard fail. "
        "Values above AUDIT_MAX_TAPER_RATIO but below this remain warning-only because taper is PSF-sensitive."
    ),
})

TRACKING_TUNER_EXPLANATION = (
    "Several tracking parameters in this section were originally selected by an evolutionary tuning script. "
    "That script rewarded multi-slice continuity while penalizing fragmentation, implausible merges, and several biology-aware outlier categories. "
    "These values are strong starting points, but they can still be adjusted for a new dataset."
)

AUDIT_EXPLANATION = (
    "The technical audit is applied after 3D tracking and does not change raw segmentation or linkage. "
    "Technical-valid tracks form the one estimated-nuclei analysis population. "
    "Morphology-warning, warning-free, and reference-morphology flags are QC annotations, not alternative populations."
)

PARAM_ENUM_OPTIONS = {
    "SEGMENTATION_ENGINE": (
        "classical_saturn", "hybrid", "unet_assisted", "unet_primary"
    ),
    "UNET_THRESHOLD_MODE": ("soft", "hard"),
    "UNET_CONTEXT_MODE": ("z_minus_z_z_plus",),
    "UNET_INFERENCE_MODE": ("roi_tiled",),
    "TRACKING_BACKEND": ("legacy", "global_assignment", "hybrid_repair"),
}

PARAM_EDITOR_VALUE_TYPES = (int, float, bool, str, list)


def _parameter_editor_can_display(key, cfg):
    """Return whether a documented configuration value can be edited safely."""
    return (
        key in cfg
        and key in PARAM_DESCRIPTIONS
        and isinstance(cfg[key], PARAM_EDITOR_VALUE_TYPES)
    )


def _coerce_parameter_editor_value(value, expected_type):
    """Convert a Tk variable value back to its CONFIG-compatible type."""
    if expected_type is bool:
        if isinstance(value, bool):
            return value
        return str(value).strip().lower() in {"true", "1", "t", "y", "yes", "on"}
    if expected_type is str:
        return str(value).strip()
    if expected_type is list:
        if isinstance(value, list):
            return value
        parsed = json.loads(str(value).strip())
        if not isinstance(parsed, list):
            raise ValueError("list parameters must use JSON list syntax, for example [0.7, 0.8, 0.9]")
        return parsed
    return expected_type(str(value).strip())



class ParameterEditor(tk.Toplevel):
    """Interactive configuration editor with grouped sections and mouse-wheel scrolling."""
    def __init__(self, parent, current_config, default_config, apply_callback):
        super().__init__(parent)
        self.title("Parameter Configuration")
        self.geometry("1180x760")
        self.minsize(980, 620)
        self.current_config = current_config
        self.default_config = default_config
        self.apply_callback = apply_callback
        self.entries = {}

        tools = tk.Frame(self, bg='#e0e0e0', padx=10, pady=10)
        tools.pack(side='top', fill='x')

        tk.Button(tools, text="Apply to Session", command=self.apply, bg="#d4edda", font=("Arial", 10, "bold")).pack(side="left", padx=5)
        tk.Button(tools, text="Reset to Defaults", command=self.reset_defaults, bg="#f8d7da").pack(side="left", padx=5)
        tk.Button(tools, text="Load JSON", command=self.load_json).pack(side="left", padx=5)
        tk.Button(tools, text="Save JSON", command=self.save_json).pack(side="left", padx=5)
        tk.Label(tools, text="Grouped sections with mouse-wheel / trackpad scrolling. Audit rules are listed separately.", bg='#e0e0e0', fg='dimgray').pack(side="right", padx=8)

        self.canvas = tk.Canvas(self, highlightthickness=0)
        self.v_scrollbar = ttk.Scrollbar(self, orient="vertical", command=self.canvas.yview)
        self.h_scrollbar = ttk.Scrollbar(self, orient="horizontal", command=self.canvas.xview)
        self.scrollable_frame = tk.Frame(self.canvas)
        self.window_id = self.canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=self.v_scrollbar.set, xscrollcommand=self.h_scrollbar.set)

        self.canvas.pack(side="left", fill="both", expand=True)
        self.v_scrollbar.pack(side="right", fill="y")
        self.h_scrollbar.pack(side="bottom", fill="x")

        self.scrollable_frame.bind("<Configure>", self._on_frame_configure)
        self.canvas.bind("<Configure>", self._on_canvas_configure)
        self.canvas.bind("<Enter>", self._bind_mousewheel)
        self.canvas.bind("<Leave>", self._unbind_mousewheel)

        self.populate_form(self.current_config)

    def _on_frame_configure(self, event=None):
        self.canvas.configure(scrollregion=self.canvas.bbox("all"))

    def _on_canvas_configure(self, event):
        self.canvas.itemconfigure(self.window_id, width=max(event.width - 4, 900))

    def _bind_mousewheel(self, event=None):
        self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)
        self.canvas.bind_all("<Shift-MouseWheel>", self._on_shift_mousewheel)
        self.canvas.bind_all("<Button-4>", self._on_linux_scroll_up)
        self.canvas.bind_all("<Button-5>", self._on_linux_scroll_down)

    def _unbind_mousewheel(self, event=None):
        self.canvas.unbind_all("<MouseWheel>")
        self.canvas.unbind_all("<Shift-MouseWheel>")
        self.canvas.unbind_all("<Button-4>")
        self.canvas.unbind_all("<Button-5>")

    def _on_mousewheel(self, event):
        delta = event.delta
        if delta == 0:
            return
        step = -1 * int(delta / 120) if abs(delta) >= 120 else (-1 if delta > 0 else 1)
        self.canvas.yview_scroll(step, "units")

    def _on_shift_mousewheel(self, event):
        delta = event.delta
        if delta == 0:
            return
        step = -1 * int(delta / 120) if abs(delta) >= 120 else (-1 if delta > 0 else 1)
        self.canvas.xview_scroll(step, "units")

    def _on_linux_scroll_up(self, event):
        self.canvas.yview_scroll(-1, "units")

    def _on_linux_scroll_down(self, event):
        self.canvas.yview_scroll(1, "units")

    def populate_form(self, cfg):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()

        self.entries = {}
        row = 0

        intro = tk.Label(
            self.scrollable_frame,
            text=(
                "Edit numeric and boolean parameters below. Descriptions explain what each threshold does biologically or computationally.\n"
                + AUDIT_EXPLANATION
            ),
            justify="left", anchor="w", fg="#333333", wraplength=1040, pady=8
        )
        intro.grid(row=row, column=0, columnspan=3, sticky="ew", padx=12, pady=(8, 12))
        row += 1

        active_engine = str(cfg.get("SEGMENTATION_ENGINE", "classical_saturn"))
        active_checkpoint = str(cfg.get("UNET_MODEL_PATH", "")).strip()
        checkpoint_label = os.path.basename(active_checkpoint) if active_checkpoint else "none selected"
        checkpoint_state = (
            "found"
            if active_checkpoint and os.path.isfile(active_checkpoint)
            else ("missing" if active_checkpoint else "not required for classical mode")
        )
        active_summary = tk.Label(
            self.scrollable_frame,
            text=(
                f"Active segmentation engine: {active_engine}    |    "
                f"U-Net checkpoint: {checkpoint_label} ({checkpoint_state})"
            ),
            justify="left",
            anchor="w",
            bg="#ecfdf5" if checkpoint_state != "missing" else "#fef2f2",
            fg="#166534" if checkpoint_state != "missing" else "#991b1b",
            padx=10,
            pady=7,
        )
        active_summary.grid(row=row, column=0, columnspan=3, sticky="ew", padx=10, pady=(0, 8))
        row += 1

        displayed = set()
        for section, keys in PARAM_SECTIONS.items():
            section_frame = tk.Frame(self.scrollable_frame, bg="#f6f8fb", bd=1, relief="groove")
            section_frame.grid(row=row, column=0, columnspan=3, sticky="ew", padx=10, pady=(4, 10))
            section_frame.grid_columnconfigure(2, weight=1)

            tk.Label(section_frame, text=section, font=("Arial", 11, "bold"), bg="#dce6f7", anchor="w", padx=8, pady=6).grid(row=0, column=0, columnspan=3, sticky="ew")

            if section == "3D Tracking":
                track_note = tk.Label(
                    section_frame,
                    text="These settings control how 2D detections are linked into 3D nuclei. Several of them were originally selected by an evolutionary tuning script, so treat them as tuned starting points rather than arbitrary defaults.",
                    fg="dimgray", bg="#f6f8fb", justify="left", anchor="w", wraplength=980
                )
                track_note.grid(row=1, column=0, columnspan=3, sticky="ew", padx=8, pady=(6, 8))
                local_row = 2
            elif section == "Candidate Audit (post-tracking)":
                audit_note = tk.Label(
                    section_frame,
                    text="The audit labels reconstructed tracks after tracking. Technical-valid tracks form the one estimated-nuclei population; the remaining flags are QC annotations. Audit thresholds do not change segmentation or 3D linkage.",
                    fg="dimgray", bg="#f6f8fb", justify="left", anchor="w", wraplength=980
                )
                audit_note.grid(row=1, column=0, columnspan=3, sticky="ew", padx=8, pady=(6, 8))
                local_row = 2
            else:
                local_row = 1

            for k in keys:
                if not _parameter_editor_can_display(k, cfg):
                    continue
                displayed.add(k)
                v = cfg[k]
                label_txt = f"{PARAM_TITLES.get(k, k)}\n[{k}]"
                tk.Label(section_frame, text=label_txt, font=("Arial", 10, "bold"), width=34, anchor="e", justify="right", bg="#f6f8fb").grid(row=local_row, column=0, padx=(8, 10), pady=4, sticky="e")
                if isinstance(v, bool):
                    var = tk.BooleanVar(value=v)
                    control = tk.Checkbutton(
                        section_frame,
                        variable=var,
                        bg="#f6f8fb",
                        activebackground="#f6f8fb",
                    )
                    control.grid(row=local_row, column=1, padx=(0, 10), pady=4, sticky="w")
                elif k in PARAM_ENUM_OPTIONS:
                    var = tk.StringVar(value=str(v))
                    control = ttk.Combobox(
                        section_frame,
                        textvariable=var,
                        values=PARAM_ENUM_OPTIONS[k],
                        state="readonly",
                        width=24,
                    )
                    control.grid(row=local_row, column=1, padx=(0, 10), pady=4, sticky="w")
                else:
                    display_value = json.dumps(v) if isinstance(v, list) else str(v)
                    var = tk.StringVar(value=display_value)
                    control_frame = tk.Frame(section_frame, bg="#f6f8fb")
                    control_frame.grid(row=local_row, column=1, padx=(0, 10), pady=4, sticky="w")
                    width = 52 if k == "UNET_MODEL_PATH" else 24
                    tk.Entry(control_frame, textvariable=var, width=width).pack(side="left")
                    if k == "UNET_MODEL_PATH":
                        tk.Button(
                            control_frame,
                            text="Browse...",
                            command=lambda target=var: self._browse_checkpoint(target),
                        ).pack(side="left", padx=(5, 0))
                self.entries[k] = (var, type(v))
                desc = tk.Label(section_frame, text=PARAM_DESCRIPTIONS[k], fg="dimgray", bg="#f6f8fb", anchor="w", justify="left", wraplength=760)
                desc.grid(row=local_row, column=2, padx=(0, 10), pady=4, sticky="w")
                local_row += 1

            row += 1

        remaining = [k for k, v in cfg.items() if k not in displayed and isinstance(v, (int, float, bool))]
        if remaining:
            extra = tk.Frame(self.scrollable_frame, bg="#fdfdfd", bd=1, relief="groove")
            extra.grid(row=row, column=0, columnspan=3, sticky="ew", padx=10, pady=(4, 10))
            extra.grid_columnconfigure(2, weight=1)
            tk.Label(extra, text="Other Numeric / Boolean Parameters", font=("Arial", 11, "bold"), bg="#ececec", anchor="w", padx=8, pady=6).grid(row=0, column=0, columnspan=3, sticky="ew")
            local_row = 1
            for k in remaining:
                v = cfg[k]
                label_txt = f"{PARAM_TITLES.get(k, k)}\n[{k}]"
                tk.Label(extra, text=label_txt, font=("Arial", 10, "bold"), width=34, anchor="e", justify="right", bg="#fdfdfd").grid(row=local_row, column=0, padx=(8, 10), pady=4, sticky="e")
                var = tk.StringVar(value=str(v))
                ent = tk.Entry(extra, textvariable=var, width=16)
                ent.grid(row=local_row, column=1, padx=(0, 10), pady=4, sticky="w")
                self.entries[k] = (var, type(v))
                tk.Label(extra, text="No extended description registered yet. Use caution and compare with nearby parameters in the same stage.", fg="dimgray", bg="#fdfdfd", anchor="w", justify="left", wraplength=760).grid(row=local_row, column=2, padx=(0, 10), pady=4, sticky="w")
                local_row += 1

    def _browse_checkpoint(self, target_var):
        current = str(target_var.get()).strip()
        initial_dir = os.path.dirname(current) if current and os.path.isdir(os.path.dirname(current)) else os.getcwd()
        selected = filedialog.askopenfilename(
            title="Select v5.7 U-Net Checkpoint",
            initialdir=initial_dir,
            filetypes=[
                ("PyTorch checkpoints", "*.pt *.pth"),
                ("All files", "*.*"),
            ],
        )
        if selected:
            target_var.set(os.path.abspath(selected))

    def apply(self):
        new_cfg = self.current_config.copy()
        try:
            for k, (var, t) in self.entries.items():
                new_cfg[k] = _coerce_parameter_editor_value(var.get(), t)
        except (TypeError, ValueError, json.JSONDecodeError) as e:
            messagebox.showerror("Validation Error", f"Invalid input format: {e}")
            return

        self.apply_callback(new_cfg)
        self.destroy()

    def reset_defaults(self):
        self.populate_form(self.default_config)

    def load_json(self):
        fpath = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                with open(fpath, 'r') as f:
                    loaded = json.load(f)
                filtered = {k: v for k, v in loaded.items() if k in self.current_config}
                temp = self.current_config.copy()
                temp.update(filtered)
                self.populate_form(temp)
            except Exception as e:
                messagebox.showerror("Load Error", str(e))

    def save_json(self):
        fpath = filedialog.asksaveasfilename(defaultextension=".json", filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                new_cfg = self.current_config.copy()
                for k, (var, t) in self.entries.items():
                    new_cfg[k] = _coerce_parameter_editor_value(var.get(), t)
                with open(fpath, 'w') as f:
                    json.dump(new_cfg, f, indent=4)
                messagebox.showinfo("Saved", f"Parameters saved to {os.path.basename(fpath)}")
            except Exception as e:
                messagebox.showerror("Save Error", str(e))

# =============================================================================
# V5.2 ADVANCED AI BIOLOGICAL INTERPRETATION ENGINE
# =============================================================================

SPECIES_PROFILES = {
    "D. melanogaster": {
        "length_target": "10.0 um",
        "thickness": "0.3-0.4 um",
        "context": "Standard model species. Needle-like, highly condensed nucleus.",
        "heteromorphism": False
    },
    "D. simulans": {
        "length_target": "9.5-10.5 um",
        "thickness": "0.3 um",
        "context": "Close relative of D. mel, very similar nuclear profile.",
        "heteromorphism": False
    },
    "D. yakuba": {
        "length_target": "10.0-11.0 um",
        "thickness": "0.35 um",
        "context": "Close relative of D. mel, slightly different condensation timing.",
        "heteromorphism": False
    },
    "D. ananassae": {
        "length_target": "8.0-10.0 um",
        "thickness": "0.3-0.4 um",
        "context": "Melanogaster subgroup. Moderate nuclear elongation with distinct chromatin packaging.",
        "heteromorphism": False
    },
    "D. pseudoobscura (Dpse)": {
        "length_target": "Variable (Heteromorphic)",
        "thickness": "Variable",
        "context": "Produces both fertile 'eusperm' and shorter, non-fertilizing 'parasperm'.",
        "heteromorphism": True
    },
    "D. virilis (Dvir)": {
        "length_target": "15.0-18.0 um",
        "thickness": "0.4 um",
        "context": "Large species with robust, stable nuclear morphology.",
        "heteromorphism": False
    },
    "General / Evolutionary": {
        "length_target": "Unknown",
        "thickness": "Unknown",
        "context": "Perform comparative discovery mode to infer species strategy.",
        "heteromorphism": "Possible"
    }
}

SPECIES_PROFILES.update({
    "D. melanogaster": {
        "length_target": "10.0 um",
        "thickness": "0.3-0.4 um anatomical; reported thickness is PSF-sensitive",
        "context": "Standard model species. Needle-like, highly condensed nucleus.",
        "heteromorphism": False
    },
    "D. simulans": {
        "length_target": "9.5-10.5 um",
        "thickness": "0.3 um anatomical; reported thickness is PSF-sensitive",
        "context": "Close relative of D. mel, very similar nuclear profile.",
        "heteromorphism": False
    },
    "D. yakuba": {
        "length_target": "10.0-11.0 um",
        "thickness": "0.35 um anatomical; reported thickness is PSF-sensitive",
        "context": "Close relative of D. mel, slightly different condensation timing.",
        "heteromorphism": False
    },
    "D. ananassae": {
        "length_target": "8.0-10.0 um",
        "thickness": "0.3-0.4 um anatomical; reported thickness is PSF-sensitive",
        "context": "Melanogaster subgroup. Moderate nuclear elongation with distinct chromatin packaging.",
        "heteromorphism": False
    },
    "D. virilis (Dvir)": {
        "length_target": "15.0-18.0 um",
        "thickness": "0.4 um anatomical; reported thickness is PSF-sensitive",
        "context": "Large species with robust, stable nuclear morphology.",
        "heteromorphism": False
    },
})

def get_ai_biological_interpretation(csv_summary_str, species, folder_name, model_id="gemini-2.5-pro"):
    """Calls Gemini API for biological narrative. Priority: Local File > Env Var."""
    profile = SPECIES_PROFILES.get(species, SPECIES_PROFILES["General / Evolutionary"])

    if not _HAVE_REQUESTS:
        return "AI ANALYSIS SKIPPED: 'requests' library not found. Please run 'pip install requests'."

    # 1. Try local file path first
    key_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gemini_api_key.txt")
    api_key = ""
    if os.path.exists(key_file):
        with open(key_file, 'r') as f:
            api_key = f.read().strip()

    # 2. Fallback to Env Var
    if not api_key:
        api_key = os.environ.get("GEMINI_API_KEY", "")

    if not api_key:
        return "AI ANALYSIS SKIPPED: No Gemini API Key found. Use 'Set API Key' in the GUI (Free at aistudio.google.com)."

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent?key={api_key}"
    system_prompt = f"""You are a world-class Drosophila reproductive biologist and evolutionary morphologist.
    You are interpreting 3D-tracked sperm nuclei data from a confocal Z-stack batch analysis.

    BIOLOGICAL CONTEXT:
    - Target Species: {species}
    - Baseline Profile: {profile['context']}
    - Expected Length Class: {profile['length_target']}
    - Thickness Reference: {profile['thickness']}
    - Heteromorphism (Multiple Morphs): {'Yes' if profile.get('heteromorphism') else 'No'}
    - Source Data Folder: {folder_name}
    - Measurement Caveat: reported 3D length is a projection-length-plus-Z-span estimate. Volume,
      effective thickness, taper, and area-derived metrics are PSF- and sampling-sensitive; use them
      mainly for relative comparisons under matched imaging settings, not literal anatomical dimensions.

    INPUT DATA (CSV Summary):
    {csv_summary_str}

    YOUR TASK:
    1. Determine the 'Morphological Class' of this population. Does it align with the expected {species} profile?
    2. Analyze maturation: Are these likely mature motile sperm or transitionary elongating spermatids based primarily on length and track morphology, with thickness treated as PSF-sensitive supporting evidence?
    3. Identify Anomalies: Highlight outliers in tortuosity, estimated thickness, taper, or abrupt area change that may indicate fixation artifacts, segmentation instability, over-merging, or developmental defects.
    4. Evolutionary Insight: Provide 2-3 sentences on the evolutionary context of this sperm morphology.
    5. Formatting: Use professional, high-density scientific language. Use Markdown formatting.
    """
    payload = {"contents": [{"parts": [{"text": system_prompt}]}]}
    try:
        response = requests.post(url, json=payload, timeout=300)
        if response.status_code == 200:
            return response.json()['candidates'][0]['content']['parts'][0]['text']
        else: return f"AI API Error ({response.status_code}): {response.text}"
    except Exception as e: return f"Failed to connect to AI Service: {str(e)}"

def generate_ai_html_report(out_dir, ai_text, stats_summary, species):
    """Generates premium HTML dashboard."""
    html_path = os.path.join(out_dir, f"AI_Biological_Analysis_{_VERSION}.html")
    try:
        import markdown
        ai_html = markdown.markdown(ai_text)
    except ImportError:
        ai_html = ai_text.replace("\n", "<br>")

    stats_html = "".join([f'<div class="stat-box"><div class="stat-val">{v}</div><div class="stat-label">{k.replace("_", " ")}</div></div>' for k,v in stats_summary.items()])

    html_content = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <title>Sperm AI Analysis - {_VERSION}</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap" rel="stylesheet">
        <style>
            :root {{ --primary: #2563eb; --bg: #f8fafc; --card: rgba(255, 255, 255, 0.8); }}
            body {{ font-family: 'Inter', sans-serif; background: var(--bg); color: #1e293b; line-height: 1.6; padding: 40px; }}
            .container {{ max-width: 900px; margin: 0 auto; }}
            .badge {{ background: #dbeafe; color: #1e40af; padding: 4px 12px; border-radius: 99px; font-size: 0.8rem; font-weight: bold; }}
            .glass-card {{ background: var(--card); backdrop-filter: blur(10px); border: 1px solid rgba(255,255,255,0.3);
                          border-radius: 16px; padding: 30px; box-shadow: 0 4px 20px rgba(0,0,0,0.05); margin-bottom: 30px; }}
            .stats-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 20px; margin-bottom: 30px; }}
            .stat-box {{ background: white; padding: 15px; border-radius: 12px; text-align: center; border: 1px solid #e2e8f0; }}
            .stat-val {{ font-size: 1.5rem; font-weight: 800; color: var(--primary); }}
            .stat-label {{ font-size: 0.75rem; text-transform: uppercase; color: #64748b; }}
            .ai-content h2 {{ color: var(--primary); border-bottom: 1px solid #e2e8f0; padding-bottom: 8px; }}
        </style>
    </head>
    <body>
        <div class="container">
            <span class="badge">Drosophila Morphometrics {_VERSION}</span>
            <h1>AI Biological Interpretation</h1>
            <p>Advanced Analysis for <strong>{species}</strong></p>
            <div class="stats-grid">{stats_html}</div>
            <div class="glass-card">
                <div class="ai-content">{ai_html}</div>
            </div>
            <footer>Generated | {time.strftime('%Y-%m-%d %H:%M:%S')} | {_VERSION} AI Layer</footer>
        </div>
    </body>
    </html>
    """
    with open(html_path, "w", encoding="utf-8") as f: f.write(html_content)
    return html_path


# =============================================================================
# MULTI-SAMPLE STUDY MANAGEMENT
# =============================================================================

STUDY_MANIFEST_COLUMNS = [
    "include",
    "sample_id",
    "group",
    "input_dir",
    "roi_path",
    "file_pattern",
    "slice_count",
    "z_min",
    "z_max",
    "xy_um_per_pixel",
    "z_um_per_slice",
    "acquisition_class",
    "status",
    "message",
    "output_dir",
]

_STUDY_SOURCE_RE = re.compile(
    r"^Project(?P<project>\d*)_Series(?P<series>\d+)_z(?P<z>\d+)_ch00\.tif{1,2}$",
    re.IGNORECASE,
)
_STUDY_EXPLICIT_Z_RE = re.compile(
    r"^(?P<prefix>.+?)(?P<z_sep>[_ .-]?)(?P<z_token>z)(?P<z>\d+)"
    r"(?P<channel_part>[_ .-](?P<channel_tag>ch|c)(?P<channel>\d+))?"
    r"\.(?P<extension>tif{1,2})$",
    re.IGNORECASE,
)
_STUDY_TRAILING_INDEX_RE = re.compile(
    r"^(?P<prefix>.+?)(?P<index_sep>[_ .-])(?P<z>\d+)"
    r"\.(?P<extension>tif{1,2})$",
    re.IGNORECASE,
)
_STUDY_EXCLUDED_DIR_NAMES = {
    "overlays",
    "quality_overlays",
    "plots",
    "debug",
    "debug_images",
    "masks",
    "labels",
    "segmentation_outputs",
    "biologist_results",
    "parameter_tuning_results",
}


def _study_parse_source_name(filename):
    """Return a conservative stack identity and Z index for a source TIFF name."""
    match = _STUDY_SOURCE_RE.match(filename)
    if match:
        project = match.group("project") or ""
        series = int(match.group("series"))
        project_token = f"Project{project}" if project else "Project"
        return {
            "kind": "leica",
            "stack_key": ("leica", project.lower(), series),
            "project": project,
            "series": series,
            "z": int(match.group("z")),
            "file_pattern": f"{project_token}_Series{series:03d}_z*_ch00.tif",
            "label": f"{project_token}_Series{series:03d}",
        }

    match = _STUDY_EXPLICIT_Z_RE.match(filename)
    if match:
        channel = match.group("channel")
        if channel is not None and int(channel) != 0:
            return None
        prefix = match.group("prefix")
        z_sep = match.group("z_sep")
        z_token = match.group("z_token")
        channel_part = match.group("channel_part") or ""
        extension = match.group("extension")
        identity = f"{prefix}{channel_part}".lower()
        return {
            "kind": "explicit_z",
            "stack_key": ("explicit_z", identity),
            "project": "",
            "series": None,
            "z": int(match.group("z")),
            "file_pattern": (
                f"{prefix}{z_sep}{z_token}[0-9]*{channel_part}.{extension}"
            ),
            "label": prefix.rstrip("_ .-") or "Z_stack",
        }

    match = _STUDY_TRAILING_INDEX_RE.match(filename)
    if match:
        prefix = match.group("prefix")
        separator = match.group("index_sep")
        extension = match.group("extension")
        return {
            "kind": "trailing_index",
            "stack_key": ("trailing_index", prefix.lower()),
            "project": "",
            "series": None,
            "z": int(match.group("z")),
            "file_pattern": f"{prefix}{separator}[0-9]*.{extension}",
            "label": prefix.rstrip("_ .-") or "Indexed_stack",
        }
    return None


def _study_is_output_directory(folder, study_root):
    """Exclude known generated-output trees from flexible TIFF discovery."""
    try:
        parts = folder.resolve().relative_to(study_root.resolve()).parts
    except Exception:
        parts = folder.parts
    for part in parts:
        normalized = part.lower()
        if normalized in _STUDY_EXCLUDED_DIR_NAMES:
            return True
        if normalized.startswith(("batch_output", "attempt_", "v5_6_", "v5_7_")):
            return True
    return False


def _study_bool(value):
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on", "include"}


def _study_series_bool(series):
    if series.empty:
        return pd.Series(dtype=bool)
    if pd.api.types.is_bool_dtype(series):
        return series.fillna(False)
    return series.fillna(False).map(_study_bool)


def _study_safe_id(value):
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", str(value).strip()).strip("._-")
    return cleaned or "Sample"


def _study_group_from_folder(folder, study_root):
    try:
        rel = pl.Path(folder).resolve().relative_to(pl.Path(study_root).resolve())
        candidate = rel.parts[0] if len(rel.parts) > 1 else pl.Path(folder).parent.name
    except Exception:
        candidate = pl.Path(folder).parent.name
    candidate = re.sub(r"\s+Test\s+SV$", "", candidate, flags=re.IGNORECASE).strip()
    normalized = candidate.lower()
    if normalized.startswith(("kj ", "lkj ")):
        return "KJ"
    if normalized.startswith("w1118 "):
        return "WT"
    return candidate or "Unassigned"


def _study_parse_leica_metadata(
    sample_dir,
    project_number,
    series_number,
    fallback_xy,
    fallback_z,
):
    """Read the physical calibration and acquisition signature from Leica XML."""
    import xml.etree.ElementTree as ET

    project_token = f"Project{project_number}" if project_number else "Project"
    xml_path = (
        pl.Path(sample_dir)
        / "MetaData"
        / f"{project_token}_Series{int(series_number):03d}.xml"
    )
    result = {
        "xy_um_per_pixel": float(fallback_xy),
        "z_um_per_slice": float(fallback_z),
        "acquisition_class": "metadata unavailable",
        "metadata_path": str(xml_path) if xml_path.exists() else "",
    }
    if not xml_path.exists():
        return result

    try:
        root = ET.parse(xml_path).getroot()
        dimensions = {}
        setting_candidates = []
        detector_candidates = []
        for elem in root.iter():
            tag = elem.tag.split("}")[-1]
            if tag == "DimensionDescription" and "DimID" in elem.attrib:
                dimensions[int(elem.attrib["DimID"])] = elem.attrib
            elif tag == "ATLConfocalSettingDefinition":
                setting_candidates.append(elem.attrib)
            elif tag == "Detector" and elem.attrib.get("IsActive") == "1":
                detector_candidates.append(elem.attrib)

        def setting_score(candidate):
            fields = (
                "Begin",
                "End",
                "Sections",
                "ObjectiveName",
                "Magnification",
                "NumericalAperture",
                "Zoom",
                "MicroscopeModel",
            )
            return sum(bool(str(candidate.get(field, "")).strip()) for field in fields)

        settings = max(setting_candidates, key=setting_score, default={})
        detector = next(
            (
                candidate
                for candidate in detector_candidates
                if abs(float(candidate.get("Gain", 0) or 0)) > 0
            ),
            detector_candidates[0] if detector_candidates else {},
        )

        dim_x = dimensions.get(1, {})
        if dim_x.get("NumberOfElements") and dim_x.get("Length"):
            result["xy_um_per_pixel"] = (
                abs(float(dim_x["Length"])) * 1_000_000.0 / int(dim_x["NumberOfElements"])
            )

        dim_z = dimensions.get(3, {})
        n_z = int(dim_z.get("NumberOfElements", 0) or 0)
        begin = settings.get("Begin")
        end = settings.get("End")
        if n_z > 1 and begin is not None and end is not None:
            result["z_um_per_slice"] = abs(float(end) - float(begin)) * 1_000_000.0 / (n_z - 1)
        elif n_z > 1 and dim_z.get("Length"):
            result["z_um_per_slice"] = (
                abs(float(dim_z["Length"])) * 1_000_000.0 / (n_z - 1)
            )

        objective = settings.get("ObjectiveName", "unknown").strip()
        zoom = settings.get("Zoom", "unknown")
        magnification = settings.get("Magnification", "unknown")
        numerical_aperture = settings.get("NumericalAperture", "unknown")
        microscope = settings.get("MicroscopeModel", "unknown")
        gain = detector.get("Gain", "unknown")
        time_gate = detector.get("IsTimeGateActivated", "unknown")
        result["acquisition_class"] = (
            f"microscope={microscope}; objective={objective}; "
            f"magnification={magnification}; NA={numerical_aperture}; "
            f"zoom={zoom}; gain={gain}; time_gate={time_gate}"
        )
    except Exception as exc:
        result["acquisition_class"] = f"metadata parse warning: {exc}"
    return result


def discover_multisample_study(study_root, roi_filename="analysis_roi_v5_7.npy", base_cfg=None):
    """Discover top-level source Z-stacks below a study root without reading outputs."""
    study_root = pl.Path(study_root).resolve()
    if not study_root.is_dir():
        raise FileNotFoundError(f"Study root does not exist: {study_root}")
    cfg = CONFIG if base_cfg is None else base_cfg
    rows = []
    used_ids = set()

    candidate_dirs = [study_root]
    candidate_dirs.extend(path for path in study_root.rglob("*") if path.is_dir())
    for folder in candidate_dirs:
        if _study_is_output_directory(folder, study_root):
            continue
        sources = []
        for path in folder.iterdir():
            if not path.is_file():
                continue
            parsed = _study_parse_source_name(path.name)
            if parsed:
                sources.append((parsed, path))
        if not sources:
            continue

        stack_keys = sorted({parsed["stack_key"] for parsed, _path in sources})
        for stack_key in stack_keys:
            series_sources = sorted(
                [
                    (parsed["z"], path, parsed)
                    for parsed, path in sources
                    if parsed["stack_key"] == stack_key
                ],
                key=lambda item: item[0],
            )
            z_values = [z for z, _path, _parsed in series_sources]
            stack_info = series_sources[0][2]
            base_id = _study_safe_id(folder.name)
            if len(stack_keys) > 1:
                base_id = f"{base_id}_{_study_safe_id(stack_info['label'])}"
            sample_id = base_id
            suffix = 2
            while sample_id.lower() in used_ids:
                sample_id = f"{base_id}_{suffix}"
                suffix += 1
            used_ids.add(sample_id.lower())

            if stack_info["kind"] == "leica":
                metadata = _study_parse_leica_metadata(
                    folder,
                    stack_info["project"],
                    stack_info["series"],
                    cfg.get("UM_PER_PX_XY", 1.0),
                    cfg.get("UM_PER_SLICE_Z", 1.0),
                )
            else:
                metadata = {
                    "xy_um_per_pixel": float(cfg.get("UM_PER_PX_XY", 1.0)),
                    "z_um_per_slice": float(cfg.get("UM_PER_SLICE_Z", 1.0)),
                    "acquisition_class": (
                        f"generic filename ({stack_info['kind']}); "
                        "calibration inherited from current settings"
                    ),
                }
            roi_path = folder / roi_filename
            if not roi_path.is_file():
                roi_candidates = sorted(
                    path
                    for path in folder.glob("*.npy")
                    if path.is_file() and path.name.lower().startswith("roi")
                )
                if len(roi_candidates) == 1:
                    roi_path = roi_candidates[0]
            rows.append(
                {
                    "include": True,
                    "sample_id": sample_id,
                    "group": _study_group_from_folder(folder, study_root),
                    "input_dir": str(folder),
                    "roi_path": str(roi_path),
                    "file_pattern": stack_info["file_pattern"],
                    "slice_count": len(series_sources),
                    "z_min": min(z_values),
                    "z_max": max(z_values),
                    "xy_um_per_pixel": float(metadata["xy_um_per_pixel"]),
                    "z_um_per_slice": float(metadata["z_um_per_slice"]),
                    "acquisition_class": metadata["acquisition_class"],
                    "status": "pending",
                    "message": "",
                    "output_dir": "",
                }
            )
    return rows


def save_multisample_manifest(rows, path):
    path = pl.Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    frame = pd.DataFrame(rows)
    for column in STUDY_MANIFEST_COLUMNS:
        if column not in frame.columns:
            frame[column] = ""
    frame[STUDY_MANIFEST_COLUMNS].to_csv(path, index=False)
    return str(path)


def load_multisample_manifest(path):
    frame = pd.read_csv(path, keep_default_na=False)
    rows = []
    for record in frame.to_dict(orient="records"):
        row = {column: record.get(column, "") for column in STUDY_MANIFEST_COLUMNS}
        row["include"] = _study_bool(row["include"])
        for key in ("slice_count", "z_min", "z_max"):
            row[key] = int(float(row[key])) if str(row[key]).strip() else 0
        for key in ("xy_um_per_pixel", "z_um_per_slice"):
            row[key] = float(row[key]) if str(row[key]).strip() else 0.0
        rows.append(row)
    return rows


def organize_multisample_study_copy(
    rows,
    output_root,
    progress_callback=None,
    copy_metadata=True,
):
    """Create a canonical, non-destructive copy of discovered study stacks."""
    from datetime import datetime

    if not rows:
        raise ValueError("No study rows were supplied for organization.")

    output_root = pl.Path(output_root).resolve()
    input_roots = {
        pl.Path(str(row.get("input_dir", ""))).resolve()
        for row in rows
        if str(row.get("input_dir", "")).strip()
    }
    common_source_root = pl.Path(
        os.path.commonpath([str(path) for path in input_roots])
    ).resolve()
    try:
        output_root.relative_to(common_source_root)
        raise ValueError(
            "The organized output must be outside the discovered source tree."
        )
    except ValueError as exc:
        if str(exc).startswith("The organized output"):
            raise
    for input_root in input_roots:
        if output_root == input_root:
            raise ValueError("The organized output cannot be an input specimen folder.")
        try:
            output_root.relative_to(input_root)
            raise ValueError(
                "The organized output must be outside every source specimen folder."
            )
        except ValueError as exc:
            if str(exc).startswith("The organized output"):
                raise

    marker_path = output_root / "organization_summary.json"
    if output_root.exists():
        existing = list(output_root.iterdir())
        if existing and not marker_path.is_file():
            raise ValueError(
                "Choose an empty output folder. Existing non-organizer content was found."
            )
    output_root.mkdir(parents=True, exist_ok=True)

    validated, _errors = validate_multisample_manifest(rows)
    organized_rows = []
    mapping_records = []
    copied_files = 0
    reused_files = 0
    copied_bytes = 0
    total = sum(_study_bool(row.get("include", True)) for row in validated)
    position = 0

    for row in validated:
        if not _study_bool(row.get("include", True)):
            continue
        position += 1
        sample_id = _study_safe_id(row.get("sample_id", "Sample"))
        group = _study_safe_id(row.get("group", "Unassigned"))
        source_dir = pl.Path(str(row.get("input_dir", ""))).resolve()
        pattern = str(row.get("file_pattern", "")).strip()

        blocking_messages = [
            message.strip()
            for message in str(row.get("message", "")).split(";")
            if message.strip()
            and message.strip() != "ROI file missing"
            and not message.strip().startswith("unusual Z spacing")
        ]
        if blocking_messages:
            raise ValueError(
                f"{sample_id} cannot be organized: {'; '.join(blocking_messages)}"
            )

        source_members = []
        for source_path in source_dir.glob(pattern):
            parsed = _study_parse_source_name(source_path.name)
            if parsed:
                source_members.append((int(parsed["z"]), source_path))
        source_members.sort(key=lambda item: item[0])
        if not source_members:
            raise ValueError(f"{sample_id} has no source TIFF files matching {pattern!r}.")

        z_values = [z_index for z_index, _path in source_members]
        if len(z_values) != len(set(z_values)):
            raise ValueError(f"{sample_id} has duplicate Z indices.")
        missing = sorted(set(range(min(z_values), max(z_values) + 1)) - set(z_values))
        if missing:
            raise ValueError(f"{sample_id} has missing Z indices: {missing}")

        destination_dir = output_root / group / sample_id
        destination_dir.mkdir(parents=True, exist_ok=True)
        for z_index, source_path in source_members:
            destination_name = f"{sample_id}_z{z_index:04d}_ch00.tif"
            destination_path = destination_dir / destination_name
            if destination_path.exists():
                if destination_path.stat().st_size != source_path.stat().st_size:
                    raise ValueError(
                        f"Existing organized file differs: {destination_path}"
                    )
                action = "reused"
                reused_files += 1
            else:
                shutil.copy2(source_path, destination_path)
                action = "copied"
                copied_files += 1
                copied_bytes += int(source_path.stat().st_size)
            mapping_records.append(
                {
                    "sample_id": sample_id,
                    "group": row.get("group", ""),
                    "z_index": z_index,
                    "source_path": str(source_path),
                    "organized_path": str(destination_path),
                    "action": action,
                }
            )

        source_roi = pl.Path(str(row.get("roi_path", "")))
        destination_roi = destination_dir / "analysis_roi_v5_7.npy"
        if source_roi.is_file():
            if destination_roi.exists():
                if destination_roi.stat().st_size != source_roi.stat().st_size:
                    raise ValueError(
                        f"Existing organized ROI differs: {destination_roi}"
                    )
            else:
                shutil.copy2(source_roi, destination_roi)

        metadata_source = source_dir / "MetaData"
        metadata_destination = destination_dir / "MetaData"
        if copy_metadata and metadata_source.is_dir():
            shutil.copytree(
                metadata_source,
                metadata_destination,
                dirs_exist_ok=True,
                copy_function=shutil.copy2,
            )

        organized_rows.append(
            {
                "include": True,
                "sample_id": sample_id,
                "group": row.get("group", ""),
                "input_dir": str(destination_dir),
                "roi_path": str(destination_roi),
                "file_pattern": f"{sample_id}_z[0-9]*_ch00.tif",
                "slice_count": len(source_members),
                "z_min": min(z_values),
                "z_max": max(z_values),
                "xy_um_per_pixel": float(row["xy_um_per_pixel"]),
                "z_um_per_slice": float(row["z_um_per_slice"]),
                "acquisition_class": row.get("acquisition_class", ""),
                "status": "pending",
                "message": "" if destination_roi.is_file() else "ROI file missing",
                "output_dir": "",
            }
        )
        if progress_callback:
            progress_callback(
                {
                    "sample_id": sample_id,
                    "position": position,
                    "total": total,
                    "message": f"organized {len(source_members)} slices",
                }
            )

    manifest_path = output_root / "organized_study_manifest.csv"
    mapping_path = output_root / "source_file_mapping.csv"
    save_multisample_manifest(organized_rows, manifest_path)
    pd.DataFrame(mapping_records).to_csv(mapping_path, index=False)
    summary = {
        "pipeline_version": _VERSION,
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "mode": "non-destructive copy",
        "sample_count": len(organized_rows),
        "copied_files": copied_files,
        "reused_files": reused_files,
        "copied_bytes": copied_bytes,
        "samples_missing_roi": [
            row["sample_id"]
            for row in organized_rows
            if not pl.Path(row["roi_path"]).is_file()
        ],
        "manifest_path": str(manifest_path),
        "mapping_path": str(mapping_path),
    }
    _study_atomic_json(marker_path, summary)
    return organized_rows, summary


def validate_multisample_manifest(rows):
    """Validate study rows and return copied rows plus a flat error list."""
    from PIL import Image

    validated = []
    errors = []
    seen_ids = set()
    for index, original in enumerate(rows):
        row = dict(original)
        row["include"] = _study_bool(row.get("include", True))
        row_errors = []
        row_warnings = []
        sample_id = _study_safe_id(row.get("sample_id", ""))
        row["sample_id"] = sample_id
        if sample_id.lower() in seen_ids:
            row_errors.append("duplicate sample ID")
        seen_ids.add(sample_id.lower())
        if not str(row.get("group", "")).strip():
            row_errors.append("group is blank")

        folder = pl.Path(str(row.get("input_dir", "")))
        roi_path = pl.Path(str(row.get("roi_path", "")))
        pattern = str(row.get("file_pattern", "")).strip()
        source_files = []
        z_values = []
        if not folder.is_dir():
            row_errors.append("input directory missing")
        else:
            for path in folder.glob(pattern):
                parsed = _study_parse_source_name(path.name)
                if parsed:
                    source_files.append(path)
                    z_values.append(int(parsed["z"]))
            order = np.argsort(z_values) if z_values else []
            source_files = [source_files[int(i)] for i in order]
            z_values = [z_values[int(i)] for i in order]
            if not source_files:
                row_errors.append("no exact source TIFF files")
            elif len(z_values) != len(set(z_values)):
                row_errors.append("duplicate Z indices")
            else:
                expected = set(range(min(z_values), max(z_values) + 1))
                missing = sorted(expected - set(z_values))
                if missing:
                    row_errors.append(f"missing Z indices: {missing}")

        image_shape = None
        if source_files:
            shapes = set()
            try:
                for path in source_files:
                    with Image.open(path) as image:
                        shapes.add((int(image.height), int(image.width)))
                if len(shapes) != 1:
                    row_errors.append(f"inconsistent image dimensions: {sorted(shapes)}")
                else:
                    image_shape = next(iter(shapes))
            except Exception as exc:
                row_errors.append(f"could not inspect TIFF dimensions: {exc}")

        if not roi_path.is_file():
            row_errors.append("ROI file missing")
        else:
            try:
                roi = np.load(roi_path, mmap_mode="r")
                if image_shape is not None and tuple(roi.shape) != tuple(image_shape):
                    row_errors.append(f"ROI shape {tuple(roi.shape)} != image shape {image_shape}")
                if not np.any(roi):
                    row_errors.append("ROI is empty")
            except Exception as exc:
                row_errors.append(f"ROI could not be read: {exc}")

        try:
            xy = float(row.get("xy_um_per_pixel", 0))
            if xy <= 0:
                row_errors.append("XY calibration must be positive")
            row["xy_um_per_pixel"] = xy
        except Exception:
            row_errors.append("XY calibration is invalid")
        try:
            z_step = float(row.get("z_um_per_slice", 0))
            if z_step <= 0:
                row_errors.append("Z calibration must be positive")
            if z_step > 1.5:
                row_warnings.append(f"unusual Z spacing {z_step:.4f} um")
            row["z_um_per_slice"] = z_step
        except Exception:
            row_errors.append("Z calibration is invalid")

        row["slice_count"] = len(source_files)
        row["z_min"] = min(z_values) if z_values else 0
        row["z_max"] = max(z_values) if z_values else 0
        if not row["include"]:
            row["status"] = "excluded"
        elif row_errors:
            row["status"] = "invalid"
        else:
            row["status"] = "validated"
        row["message"] = "; ".join(row_errors + row_warnings)
        validated.append(row)
        for message in row_errors:
            errors.append(f"row {index + 1} ({sample_id}): {message}")
    return validated, errors


def _study_config_fingerprint(cfg):
    import hashlib

    excluded = {"INPUT_DIR", "OUTPUT_DIR", "ROI_MASK_PATH", "EXCLUSION_MASK_PATH", "FILE_PATTERN"}
    serializable = {
        str(key): _json_scalar(value)
        for key, value in cfg.items()
        if key not in excluded and not str(key).startswith("_")
    }
    payload = json.dumps(serializable, sort_keys=True, default=str).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _study_atomic_json(path, data):
    path = pl.Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    temp_path = path.with_suffix(path.suffix + ".tmp")
    with temp_path.open("w", encoding="utf-8") as handle:
        json.dump(_json_scalar(data), handle, indent=2)
    os.replace(temp_path, path)


def _study_next_attempt_dir(sample_root):
    sample_root = pl.Path(sample_root)
    sample_root.mkdir(parents=True, exist_ok=True)
    existing = []
    for path in sample_root.glob("attempt_*"):
        try:
            existing.append(int(path.name.rsplit("_", 1)[-1]))
        except ValueError:
            continue
    number = max(existing, default=0) + 1
    return sample_root / f"attempt_{number:03d}"


def _study_find_output_csv(output_dir, stem, exclude_tokens=()):
    candidates = []
    for path in pl.Path(output_dir).glob(f"{stem}*.csv"):
        lower = path.name.lower()
        if any(token.lower() in lower for token in exclude_tokens):
            continue
        candidates.append(path)
    return max(candidates, key=lambda path: path.stat().st_mtime) if candidates else None


def _study_find_primary_track_summary(output_dir):
    """Select the complete track table, never a filtered derivative."""
    filtered_tokens = (
        "_all", "_quality", "_biological", "_candidate", "_technical",
        "_reference", "_morphology", "_failure",
    )
    candidates = []
    for path in pl.Path(output_dir).glob("track_summary*.csv"):
        lower = path.stem.lower()
        if any(token in lower for token in filtered_tokens):
            continue
        candidates.append(path)
    return max(candidates, key=lambda path: path.stat().st_mtime) if candidates else None


def _study_scaled_rate(count, exposure, scale):
    exposure = float(exposure)
    if not np.isfinite(exposure) or exposure <= 0:
        return np.nan
    return float(count) * float(scale) / exposure


def _study_exposure_metrics(row, roi_pixels_from_qc=0):
    """Return geometric sampling exposure without changing image processing."""
    roi_pixels = int(roi_pixels_from_qc or 0)
    frame_pixels = 0
    roi_path = pl.Path(str(row.get("roi_path", "")))
    if roi_path.is_file():
        try:
            roi = np.asarray(np.load(roi_path, mmap_mode="r"), dtype=bool)
            roi_pixels_file = int(np.count_nonzero(roi))
            frame_pixels = int(roi.size)
            if roi_pixels_file > 0:
                roi_pixels = roi_pixels_file
        except Exception:
            pass

    xy = float(row.get("xy_um_per_pixel", 0) or 0)
    z_step = float(row.get("z_um_per_slice", 0) or 0)
    slices = int(row.get("slice_count", 0) or 0)
    area = float(roi_pixels) * xy * xy
    sampled_depth = float(slices) * z_step
    stack_span = float(max(slices - 1, 0)) * z_step
    sampled_volume = area * sampled_depth
    return {
        "roi_pixel_count": roi_pixels,
        "frame_pixel_count": frame_pixels,
        "roi_fraction_of_frame": float(roi_pixels / frame_pixels) if frame_pixels > 0 else np.nan,
        "roi_area_um2": area,
        "included_slice_count": slices,
        "z_spacing_um": z_step,
        "sampled_depth_um": sampled_depth,
        "stack_span_um": stack_span,
        "sampled_roi_volume_um3": sampled_volume,
    }


def _study_track_source_sets(output_dir):
    tracked_path = _study_find_output_csv(output_dir, "measurements_with_tracks")
    if tracked_path is None:
        return set(), set()
    tracked = pd.read_csv(tracked_path)
    if tracked.empty or "track_id" not in tracked.columns or "detection_source" not in tracked.columns:
        return set(), set()
    source = tracked["detection_source"].fillna("saturn_classical").astype(str)
    tracked = tracked.assign(_is_unet=source.str.startswith("unet_rescued"))
    unet_by_track = tracked.groupby("track_id")["_is_unet"].any()
    unet_tracks = set(unet_by_track[unet_by_track].index.tolist())
    classical_only_tracks = set(unet_by_track[~unet_by_track].index.tolist())
    return classical_only_tracks, unet_tracks


def summarize_study_sample(row, output_dir):
    """Create one specimen-level summary with raw and exposure-normalized counts."""
    output_dir = pl.Path(output_dir)
    measurements_path = _study_find_output_csv(output_dir, "spermatid_measurements")
    tracks_path = _study_find_primary_track_summary(output_dir)
    detections = pd.read_csv(measurements_path) if measurements_path else pd.DataFrame()
    tracks = pd.read_csv(tracks_path) if tracks_path else pd.DataFrame()
    unet_accounting = _unet_detection_accounting(detections)

    def median(frame, column):
        if column not in frame.columns or frame.empty:
            return np.nan
        return float(pd.to_numeric(frame[column], errors="coerce").median())

    roi_pixels = 0
    qc_path = output_dir / "stack_preprocessing_qc.json"
    if qc_path.exists():
        try:
            with qc_path.open("r", encoding="utf-8") as handle:
                roi_pixels = int(json.load(handle).get("roi_pixel_count", 0))
        except Exception:
            roi_pixels = 0

    if "technical_valid" in tracks:
        analysis_mask = _study_series_bool(tracks["technical_valid"])
    elif "is_biological_candidate" in tracks:
        analysis_mask = _study_series_bool(tracks["is_biological_candidate"])
    else:
        analysis_mask = pd.Series(True, index=tracks.index)
    analysis_tracks = tracks.loc[analysis_mask].copy()
    estimated_nuclei = int(analysis_mask.sum())
    technical_failures = int((~analysis_mask).sum())
    exposure = _study_exposure_metrics(row, roi_pixels)
    roi_area = exposure["roi_area_um2"]
    roi_volume = exposure["sampled_roi_volume_um3"]
    n_slices = exposure["included_slice_count"]

    analysis_ids = set(tracks.loc[analysis_mask, "track_id"].tolist()) if "track_id" in tracks else set()
    classical_track_ids, unet_track_ids = _study_track_source_sets(output_dir)
    analysis_unet_track_ids = unet_track_ids & analysis_ids
    analysis_classical_track_ids = classical_track_ids & analysis_ids

    tracked_measurements_path = _study_find_output_csv(output_dir, "measurements_with_tracks")
    tracked_measurements = (
        pd.read_csv(tracked_measurements_path)
        if tracked_measurements_path is not None
        else pd.DataFrame()
    )
    per_track_2d = pd.DataFrame()
    if (
        analysis_ids
        and not tracked_measurements.empty
        and "track_id" in tracked_measurements.columns
    ):
        accepted_detections = tracked_measurements[
            tracked_measurements["track_id"].isin(analysis_ids)
        ].copy()
        aggregations = {}
        if "length_um_geodesic" in accepted_detections.columns:
            aggregations["maximum_2d_length_um"] = ("length_um_geodesic", "max")
        if "width_um" in accepted_detections.columns:
            aggregations["median_2d_width_um"] = ("width_um", "median")
        if "length_width_ratio" in accepted_detections.columns:
            aggregations["median_2d_length_width_ratio"] = (
                "length_width_ratio",
                "median",
            )
        if aggregations:
            per_track_2d = accepted_detections.groupby("track_id").agg(**aggregations)

    z_min = int(row.get("z_min", 0) or 0)
    z_max = int(row.get("z_max", max(n_slices - 1, 0)) or 0)
    lower_boundary = int((pd.to_numeric(tracks["z_start"], errors="coerce") <= z_min).sum()) if "z_start" in tracks else 0
    upper_boundary = int((pd.to_numeric(tracks["z_end"], errors="coerce") >= z_max).sum()) if "z_end" in tracks else 0
    if "z_start" in tracks and "z_end" in tracks:
        boundary_mask = (
            (pd.to_numeric(tracks["z_start"], errors="coerce") <= z_min)
            | (pd.to_numeric(tracks["z_end"], errors="coerce") >= z_max)
        )
        boundary_count = int(boundary_mask.sum())
    else:
        boundary_count = 0

    if "z_slice" in detections and not detections.empty:
        detection_z = pd.to_numeric(detections["z_slice"], errors="coerce").dropna()
        positive_slice_count = int(detection_z.nunique())
        positive_z_min = int(detection_z.min()) if not detection_z.empty else np.nan
        positive_z_max = int(detection_z.max()) if not detection_z.empty else np.nan
    else:
        positive_slice_count = 0
        positive_z_min = np.nan
        positive_z_max = np.nan

    boundary_fraction = float(boundary_count / len(tracks)) if len(tracks) > 0 else 0.0
    normalization_warnings = []
    if roi_area <= 0 or roi_volume <= 0:
        normalization_warnings.append("invalid ROI exposure")
    if boundary_fraction > 0.20:
        normalization_warnings.append(f"high Z-boundary track fraction ({boundary_fraction:.1%})")
    if positive_slice_count > 0 and (positive_z_min <= z_min or positive_z_max >= z_max):
        normalization_warnings.append("detections reach an acquisition Z boundary")

    summary = {
        "sample_id": row["sample_id"],
        "group": row.get("group", ""),
        "status": "complete",
        "input_dir": row.get("input_dir", ""),
        "output_dir": str(output_dir),
        "slice_count": n_slices,
        "xy_um_per_pixel": float(row.get("xy_um_per_pixel", 0)),
        "z_um_per_slice": float(row.get("z_um_per_slice", 0)),
        **exposure,
        "detection_positive_slice_count": positive_slice_count,
        "detection_positive_z_min": positive_z_min,
        "detection_positive_z_max": positive_z_max,
        "detection_positive_slice_fraction": float(positive_slice_count / n_slices) if n_slices > 0 else np.nan,
        "raw_2d_detection_count": int(len(detections)),
        "raw_2d_detections_per_1000_um2_per_slice": _study_scaled_rate(len(detections), roi_area * n_slices, 1_000.0),
        "saturn_classical_count": unet_accounting["saturn_classical_2d_count"],
        "unet_rescued_count": unet_accounting["unet_rescued_2d_count"],
        "unet_rescue_fraction": unet_accounting["unet_rescue_fraction_of_2d_detections"],
        **unet_accounting,
        "estimated_unique_nuclei": estimated_nuclei,
        "estimated_unique_nuclei_classical_only": int(len(analysis_classical_track_ids)),
        "estimated_unique_nuclei_with_unet_evidence": int(len(analysis_unet_track_ids)),
        "estimated_unique_nuclei_unet_fraction": float(
            len(analysis_unet_track_ids) / max(estimated_nuclei, 1)
        ),
        "estimated_nuclei_per_1000_um2": _study_scaled_rate(estimated_nuclei, roi_area, 1_000.0),
        "estimated_nuclei_per_100000_um3": _study_scaled_rate(estimated_nuclei, roi_volume, 100_000.0),
        "qc_technical_failure_track_count": technical_failures,
        "qc_classical_only_3d_track_count": int(len(classical_track_ids)),
        "qc_unet_associated_3d_track_count": int(len(unet_track_ids)),
        "qc_analysis_population_unet_track_count": int(len(analysis_unet_track_ids)),
        "lower_z_boundary_track_count": lower_boundary,
        "upper_z_boundary_track_count": upper_boundary,
        "z_boundary_track_count": boundary_count,
        "z_boundary_track_fraction": boundary_fraction,
        "z_coverage_censored": bool(boundary_count > 0),
        "normalization_valid": bool(roi_area > 0 and roi_volume > 0),
        "normalization_warning": "; ".join(normalization_warnings),
        "median_2d_length_um": (
            median(per_track_2d, "maximum_2d_length_um")
            if "maximum_2d_length_um" in per_track_2d
            else median(analysis_tracks, "max_length_2d")
        ),
        "median_2d_width_um": (
            median(per_track_2d, "median_2d_width_um")
            if "median_2d_width_um" in per_track_2d
            else median(analysis_tracks, "median_width_2d")
        ),
        "median_2d_length_width_ratio": (
            median(per_track_2d, "median_2d_length_width_ratio")
            if "median_2d_length_width_ratio" in per_track_2d
            else median(analysis_tracks, "median_length_width_ratio_2d")
        ),
        "median_3d_length_um": median(analysis_tracks, "total_3d_length_um"),
        "median_3d_tortuosity": median(analysis_tracks, "tortuosity_3d"),
        "median_3d_thickness_um": median(analysis_tracks, "thickness_um"),
        "median_3d_volume_um3": median(analysis_tracks, "volume_um3"),
        "median_3d_z_span_um": median(analysis_tracks, "z_span_um"),
        "acquisition_class": row.get("acquisition_class", ""),
    }
    return summary


def _study_group_summary(specimen_frame):
    """Summarize groups with specimens, rather than tracks, as replicates."""
    if specimen_frame.empty or "group" not in specimen_frame.columns or "status" not in specimen_frame.columns:
        return pd.DataFrame()
    complete = specimen_frame[specimen_frame["status"] == "complete"].copy()
    if complete.empty:
        return pd.DataFrame()
    metrics = [
        "estimated_unique_nuclei",
        "estimated_nuclei_per_1000_um2",
        "estimated_nuclei_per_100000_um3",
        "median_2d_length_um",
        "median_2d_width_um",
        "median_2d_length_width_ratio",
        "median_3d_length_um",
        "median_3d_tortuosity",
        "median_3d_thickness_um",
        "median_3d_volume_um3",
        "median_3d_z_span_um",
        "roi_area_um2",
        "sampled_roi_volume_um3",
        "z_boundary_track_fraction",
    ]
    records = []
    for group, frame in complete.groupby("group", dropna=False):
        record = {"group": group, "specimen_count": int(len(frame))}
        for metric in metrics:
            if metric not in frame.columns:
                continue
            values = pd.to_numeric(frame[metric], errors="coerce").dropna()
            record[f"{metric}_mean"] = float(values.mean()) if not values.empty else np.nan
            record[f"{metric}_median"] = float(values.median()) if not values.empty else np.nan
            record[f"{metric}_std"] = float(values.std(ddof=1)) if len(values) > 1 else np.nan
        records.append(record)
    return pd.DataFrame(records)


_STUDY_COMPARISON_METRICS = {
    "estimated_nuclei_per_1000_um2": "Estimated nuclei per 1,000 um2",
    "estimated_nuclei_per_100000_um3": "Estimated nuclei per 100,000 um3",
    "median_2d_length_um": "Specimen median 2D length (um)",
    "median_2d_width_um": "Specimen median 2D width (um)",
    "median_2d_length_width_ratio": "Specimen median 2D length / width",
    "median_3d_length_um": "Specimen median 3D length (um)",
    "median_3d_tortuosity": "Specimen median 3D tortuosity",
    "median_3d_thickness_um": "Specimen median effective thickness (um)",
    "median_3d_volume_um3": "Specimen median volume (um3)",
    "median_3d_z_span_um": "Specimen median Z span (um)",
}


def _study_reference_group(groups):
    """Prefer a clearly named WT group as the comparison reference."""
    groups = sorted(str(group) for group in groups)
    for group in groups:
        label = group.lower()
        if "wild" in label or "w1118" in label or re.search(r"(^|[^a-z])wt([^a-z]|$)", label):
            return group
    return groups[0] if groups else ""


def _study_cliffs_delta(reference, comparison):
    """Return Cliff's delta with positive values favoring the comparison group."""
    reference = np.asarray(reference, dtype=float)
    comparison = np.asarray(comparison, dtype=float)
    if reference.size == 0 or comparison.size == 0:
        return np.nan
    differences = comparison[:, None] - reference[None, :]
    return float((np.count_nonzero(differences > 0) - np.count_nonzero(differences < 0)) / differences.size)


def _study_bh_qvalues(p_values):
    """Benjamini-Hochberg correction while preserving missing values."""
    values = np.asarray(p_values, dtype=float)
    adjusted = np.full(values.shape, np.nan, dtype=float)
    valid_indices = np.flatnonzero(np.isfinite(values))
    if valid_indices.size == 0:
        return adjusted
    order = valid_indices[np.argsort(values[valid_indices])]
    ranked = values[order] * valid_indices.size / np.arange(1, valid_indices.size + 1)
    ranked = np.minimum.accumulate(ranked[::-1])[::-1]
    adjusted[order] = np.clip(ranked, 0.0, 1.0)
    return adjusted


def _study_specimen_group_comparisons(
    specimen_frame,
    random_seed=57057,
    bootstrap_resamples=5000,
    permutation_resamples=9999,
):
    """Compare two groups using specimens, never individual nuclei, as replicates."""
    complete = specimen_frame.copy()
    if "status" in complete.columns:
        complete = complete[complete["status"] == "complete"].copy()
    if "group" not in complete.columns:
        complete["group"] = ""
    complete["group"] = complete["group"].fillna("").astype(str).str.strip()
    complete = complete[complete["group"] != ""]
    groups = sorted(complete["group"].unique().tolist())
    qc = {
        "analysis_unit": "biological specimen",
        "nested_nucleus_records_role": "descriptive and audit only; not independent replicates",
        "random_seed": int(random_seed),
        "bootstrap_resamples": int(bootstrap_resamples),
        "permutation_resamples": int(permutation_resamples),
        "groups": groups,
        "specimen_counts": {
            group: int((complete["group"] == group).sum()) for group in groups
        },
        "comparison_status": "not_run",
        "warnings": [],
    }
    if len(groups) != 2:
        qc["warnings"].append(
            f"Specimen comparison requires exactly two non-empty groups; found {len(groups)}."
        )
        return pd.DataFrame(
            columns=[
                "metric",
                "metric_label",
                "analysis_unit",
                "reference_group",
                "comparison_group",
                "reference_n",
                "comparison_n",
                "reference_median",
                "comparison_median",
                "median_difference_comparison_minus_reference",
                "cliffs_delta_comparison_minus_reference",
                "bootstrap_95ci_low",
                "bootstrap_95ci_high",
                "permutation_p_value",
                "bh_fdr_q_value",
                "inference_status",
            ]
        ), qc

    reference_group = _study_reference_group(groups)
    comparison_group = next(group for group in groups if group != reference_group)
    qc["reference_group"] = reference_group
    qc["comparison_group"] = comparison_group
    qc["effect_direction"] = f"{comparison_group} minus {reference_group}"
    qc["comparison_status"] = "exploratory"
    if min(qc["specimen_counts"].values()) < 3:
        qc["warnings"].append(
            "At least one group has fewer than three specimens; inferential estimates are not reported."
        )
    elif min(qc["specimen_counts"].values()) < 5:
        qc["warnings"].append(
            "At least one group has fewer than five specimens; intervals and p-values are highly uncertain."
        )

    from scipy.stats import permutation_test

    records = []
    for metric_index, (metric, label) in enumerate(_STUDY_COMPARISON_METRICS.items()):
        if metric not in complete.columns:
            continue
        reference = pd.to_numeric(
            complete.loc[complete["group"] == reference_group, metric],
            errors="coerce",
        ).dropna().to_numpy(dtype=float)
        comparison = pd.to_numeric(
            complete.loc[complete["group"] == comparison_group, metric],
            errors="coerce",
        ).dropna().to_numpy(dtype=float)
        if reference.size == 0 and comparison.size == 0:
            continue

        reference_median = float(np.median(reference)) if reference.size else np.nan
        comparison_median = float(np.median(comparison)) if comparison.size else np.nan
        median_difference = comparison_median - reference_median
        record = {
            "metric": metric,
            "metric_label": label,
            "analysis_unit": "biological specimen",
            "reference_group": reference_group,
            "comparison_group": comparison_group,
            "reference_n": int(reference.size),
            "comparison_n": int(comparison.size),
            "reference_mean": float(np.mean(reference)) if reference.size else np.nan,
            "comparison_mean": float(np.mean(comparison)) if comparison.size else np.nan,
            "reference_median": reference_median,
            "comparison_median": comparison_median,
            "median_difference_comparison_minus_reference": median_difference,
            "median_percent_difference": (
                float(100.0 * median_difference / reference_median)
                if np.isfinite(reference_median) and reference_median != 0
                else np.nan
            ),
            "cliffs_delta_comparison_minus_reference": _study_cliffs_delta(
                reference, comparison
            ),
            "bootstrap_95ci_low": np.nan,
            "bootstrap_95ci_high": np.nan,
            "permutation_p_value": np.nan,
            "inference_status": "insufficient_specimens",
        }

        if reference.size >= 3 and comparison.size >= 3:
            rng = np.random.default_rng(random_seed + metric_index)
            bootstrap_differences = np.empty(bootstrap_resamples, dtype=float)
            for index in range(bootstrap_resamples):
                reference_sample = rng.choice(reference, size=reference.size, replace=True)
                comparison_sample = rng.choice(comparison, size=comparison.size, replace=True)
                bootstrap_differences[index] = (
                    np.median(comparison_sample) - np.median(reference_sample)
                )
            record["bootstrap_95ci_low"], record["bootstrap_95ci_high"] = [
                float(value)
                for value in np.quantile(bootstrap_differences, [0.025, 0.975])
            ]
            permutation = permutation_test(
                (reference, comparison),
                lambda ref, comp: np.median(comp) - np.median(ref),
                permutation_type="independent",
                vectorized=False,
                n_resamples=permutation_resamples,
                alternative="two-sided",
                rng=np.random.default_rng(random_seed + 1000 + metric_index),
            )
            record["permutation_p_value"] = float(permutation.pvalue)
            record["inference_status"] = (
                "exploratory_small_sample"
                if min(reference.size, comparison.size) < 5
                else "exploratory"
            )
        records.append(record)

    result = pd.DataFrame(records)
    if not result.empty:
        result["bh_fdr_q_value"] = _study_bh_qvalues(result["permutation_p_value"])
    return result, qc


def _write_study_specimen_comparison_plot(specimen_frame, comparison_frame, output_path):
    """Plot every specimen; do not substitute pooled nucleus-level distributions."""
    if comparison_frame.empty:
        return False
    metrics = comparison_frame["metric"].tolist()
    groups = [
        comparison_frame.iloc[0]["reference_group"],
        comparison_frame.iloc[0]["comparison_group"],
    ]
    panel_count = len(metrics)
    columns = 3
    rows = int(math.ceil(panel_count / columns))
    figure, axes = plt.subplots(rows, columns, figsize=(13.5, 3.8 * rows), squeeze=False)
    rng = np.random.default_rng(57057)
    colors = ["#2878B5", "#D1495B"]
    for axis, metric in zip(axes.flat, metrics):
        for group_index, group in enumerate(groups):
            values = pd.to_numeric(
                specimen_frame.loc[
                    (specimen_frame["status"] == "complete")
                    & (specimen_frame["group"].astype(str) == group),
                    metric,
                ],
                errors="coerce",
            ).dropna().to_numpy(dtype=float)
            x_values = group_index + rng.uniform(-0.075, 0.075, size=values.size)
            axis.scatter(
                x_values,
                values,
                s=42,
                color=colors[group_index],
                edgecolor="white",
                linewidth=0.6,
                alpha=0.9,
                zorder=3,
            )
            if values.size:
                median_value = float(np.median(values))
                axis.plot(
                    [group_index - 0.18, group_index + 0.18],
                    [median_value, median_value],
                    color="#202020",
                    linewidth=2.2,
                    zorder=4,
                )
        label = _STUDY_COMPARISON_METRICS.get(metric, metric)
        axis.set_title(label, fontsize=10)
        axis.set_xticks([0, 1], groups)
        axis.grid(axis="y", color="#D8D8D8", linewidth=0.7, alpha=0.8)
        axis.spines[["top", "right"]].set_visible(False)
    for axis in axes.flat[panel_count:]:
        axis.set_visible(False)
    figure.suptitle(
        "Specimen-level group comparison\nEach point is one biological specimen; bars show medians",
        fontsize=14,
        fontweight="bold",
    )
    figure.tight_layout(rect=(0, 0, 1, 0.95))
    figure.savefig(output_path, dpi=180, bbox_inches="tight")
    plt.close(figure)
    return True


def _study_normalization_qc(specimen_frame):
    """Audit whether specimen exposures are sufficiently comparable to interpret."""
    complete = specimen_frame.copy()
    if "status" in complete.columns:
        complete = complete[complete["status"] == "complete"]
    warnings = []

    def exposure_range(column):
        if column not in complete.columns:
            return np.nan, np.nan, np.nan
        values = pd.to_numeric(complete[column], errors="coerce")
        values = values[np.isfinite(values) & (values > 0)]
        if values.empty:
            return np.nan, np.nan, np.nan
        minimum = float(values.min())
        maximum = float(values.max())
        return minimum, maximum, float(maximum / minimum)

    area_min, area_max, area_ratio = exposure_range("roi_area_um2")
    volume_min, volume_max, volume_ratio = exposure_range("sampled_roi_volume_um3")
    if np.isfinite(area_ratio) and area_ratio > 4.0:
        warnings.append(f"ROI areas differ by {area_ratio:.2f}x across specimens")
    if np.isfinite(volume_ratio) and volume_ratio > 4.0:
        warnings.append(f"sampled ROI volumes differ by {volume_ratio:.2f}x across specimens")

    invalid_count = 0
    if "normalization_valid" in complete.columns:
        invalid_count = int((~_study_series_bool(complete["normalization_valid"])).sum())
        if invalid_count:
            warnings.append(f"{invalid_count} specimens have invalid normalization exposure")
    high_boundary_count = 0
    if "z_boundary_track_fraction" in complete.columns:
        boundary = pd.to_numeric(complete["z_boundary_track_fraction"], errors="coerce").fillna(0)
        high_boundary_count = int((boundary > 0.20).sum())
        if high_boundary_count:
            warnings.append(f"{high_boundary_count} specimens exceed 20% Z-boundary tracks")

    return {
        "specimen_count": int(len(complete)),
        "roi_area_um2_min": area_min,
        "roi_area_um2_max": area_max,
        "roi_area_max_min_ratio": area_ratio,
        "sampled_roi_volume_um3_min": volume_min,
        "sampled_roi_volume_um3_max": volume_max,
        "sampled_roi_volume_max_min_ratio": volume_ratio,
        "invalid_normalization_specimen_count": invalid_count,
        "high_z_boundary_specimen_count": high_boundary_count,
        "normalization_review_required": bool(warnings),
        "warnings": warnings,
    }


def _write_study_aggregates(output_root, rows, state):
    output_root = pl.Path(output_root)
    summaries = []
    track_frames = []
    for row in rows:
        record = state.get("samples", {}).get(row["sample_id"], {})
        output_dir = record.get("output_dir", "")
        if record.get("status") != "complete" or not output_dir:
            summaries.append(
                {
                    "sample_id": row["sample_id"],
                    "group": row.get("group", ""),
                    "status": record.get("status", row.get("status", "pending")),
                    "input_dir": row.get("input_dir", ""),
                    "output_dir": output_dir,
                    "message": record.get("message", ""),
                }
            )
            continue
        summary = summarize_study_sample(row, output_dir)
        summary["message"] = record.get("message", "")
        summaries.append(summary)

        tracks_path = _study_find_primary_track_summary(output_dir)
        if tracks_path:
            tracks = pd.read_csv(tracks_path)
            if not tracks.empty:
                tracks.insert(0, "group", row.get("group", ""))
                tracks.insert(0, "sample_id", row["sample_id"])
                if "track_id" in tracks.columns:
                    tracks.insert(
                        2,
                        "study_track_id",
                        row["sample_id"] + ":" + tracks["track_id"].astype(str),
                    )
                track_frames.append(tracks)

    specimen_frame = pd.DataFrame(summaries)
    biological_columns = [
        "sample_id",
        "group",
        "status",
        "input_dir",
        "output_dir",
        "acquisition_class",
        "slice_count",
        "xy_um_per_pixel",
        "z_um_per_slice",
        "roi_area_um2",
        "sampled_roi_volume_um3",
        "estimated_unique_nuclei",
        "estimated_nuclei_per_1000_um2",
        "estimated_nuclei_per_100000_um3",
        "median_2d_length_um",
        "median_2d_width_um",
        "median_2d_length_width_ratio",
        "median_3d_length_um",
        "median_3d_tortuosity",
        "median_3d_thickness_um",
        "median_3d_volume_um3",
        "median_3d_z_span_um",
        "normalization_warning",
    ]
    specimen_frame[
        [column for column in biological_columns if column in specimen_frame.columns]
    ].to_csv(output_root / "specimen_summary.csv", index=False)
    specimen_frame.to_csv(output_root / "specimen_technical_qc.csv", index=False)
    _study_group_summary(specimen_frame).to_csv(output_root / "group_summary.csv", index=False)
    _study_atomic_json(output_root / "normalization_qc.json", _study_normalization_qc(specimen_frame))
    comparison_frame, comparison_qc = _study_specimen_group_comparisons(specimen_frame)
    comparison_frame.to_csv(
        output_root / "specimen_group_comparisons.csv",
        index=False,
    )
    _study_atomic_json(
        output_root / "specimen_group_comparison_qc.json",
        comparison_qc,
    )
    _write_study_specimen_comparison_plot(
        specimen_frame,
        comparison_frame,
        output_root / "specimen_group_comparison.pdf",
    )
    if track_frames:
        pd.concat(track_frames, ignore_index=True).to_csv(output_root / "study_track_records.csv", index=False)
    return summaries


def run_multisample_study(
    rows,
    output_root,
    base_cfg=None,
    progress_callback=None,
    resume=True,
    batch_runner=None,
    stop_requested=None,
):
    """Run validated specimens sequentially and resume completed sample attempts."""
    from datetime import datetime

    def should_stop():
        if stop_requested is None:
            return False
        if callable(stop_requested):
            return bool(stop_requested())
        is_set = getattr(stop_requested, "is_set", None)
        return bool(is_set()) if callable(is_set) else bool(stop_requested)

    validated, errors = validate_multisample_manifest(rows)
    included_errors = []
    for index, row in enumerate(validated, start=1):
        if row["include"] and row["status"] == "invalid":
            included_errors.append(f"row {index} ({row['sample_id']}): {row['message']}")
    if included_errors:
        raise ValueError("Study manifest validation failed:\n" + "\n".join(included_errors))

    output_root = pl.Path(output_root).resolve()
    output_root.mkdir(parents=True, exist_ok=True)
    save_multisample_manifest(validated, output_root / "study_manifest.csv")
    cfg_template = CONFIG.copy() if base_cfg is None else dict(base_cfg)
    config_hash = _study_config_fingerprint(cfg_template)
    state_path = output_root / "study_run_state.json"
    state = {
        "pipeline_version": _VERSION,
        "config_hash": config_hash,
        "updated_at": datetime.now().isoformat(timespec="seconds"),
        "samples": {},
    }
    if resume and state_path.exists():
        try:
            with state_path.open("r", encoding="utf-8") as handle:
                prior = json.load(handle)
            if prior.get("config_hash") == config_hash:
                state = prior
        except Exception:
            pass
    state["updated_at"] = datetime.now().isoformat(timespec="seconds")
    state["run_status"] = "running"
    _study_atomic_json(output_root / "runtime_parameters.json", {
        key: value for key, value in cfg_template.items() if not str(key).startswith("_")
    })
    _study_atomic_json(state_path, state)

    runner = process_batch if batch_runner is None else batch_runner
    included = [row for row in validated if row["include"]]
    total = len(included)
    for position, row in enumerate(included, start=1):
        if should_stop():
            state["run_status"] = "stopped"
            state["updated_at"] = datetime.now().isoformat(timespec="seconds")
            _study_atomic_json(state_path, state)
            if progress_callback:
                progress_callback(
                    {
                        "event": "stopped",
                        "sample_id": "",
                        "position": position - 1,
                        "total": total,
                        "message": "stopped before starting the next specimen",
                    }
                )
            break
        sample_id = row["sample_id"]
        prior = state.get("samples", {}).get(sample_id, {})
        marker = pl.Path(prior.get("output_dir", "")) / "sample_complete.json" if prior.get("output_dir") else None
        if resume and prior.get("status") == "complete" and marker is not None and marker.exists():
            if progress_callback:
                progress_callback({"event": "skipped", "sample_id": sample_id, "position": position, "total": total, "message": "already complete"})
            continue

        attempt_dir = _study_next_attempt_dir(output_root / "samples" / sample_id)
        attempt_dir.mkdir(parents=True, exist_ok=True)
        state.setdefault("samples", {})[sample_id] = {
            "status": "running",
            "group": row.get("group", ""),
            "output_dir": str(attempt_dir),
            "message": "",
            "started_at": datetime.now().isoformat(timespec="seconds"),
        }
        state["updated_at"] = datetime.now().isoformat(timespec="seconds")
        _study_atomic_json(state_path, state)
        if progress_callback:
            progress_callback({"event": "started", "sample_id": sample_id, "position": position, "total": total, "message": "processing"})

        sample_cfg = cfg_template.copy()
        sample_cfg.update(
            {
                "RUN_MODE": "batch",
                "INPUT_DIR": row["input_dir"],
                "OUTPUT_DIR": str(attempt_dir),
                "FILE_PATTERN": row["file_pattern"],
                "ROI_MASK_PATH": row["roi_path"],
                "UM_PER_PX_XY": float(row["xy_um_per_pixel"]),
                "UM_PER_SLICE_Z": float(row["z_um_per_slice"]),
                "SHOW_PREVIEW_WINDOW": False,
                "DO_TRACKING": True,
            }
        )
        try:
            runner(sample_cfg)
            summary = summarize_study_sample(row, attempt_dir)
            marker_data = {
                "sample_id": sample_id,
                "group": row.get("group", ""),
                "pipeline_version": _VERSION,
                "config_hash": config_hash,
                "completed_at": datetime.now().isoformat(timespec="seconds"),
                "summary": summary,
            }
            _study_atomic_json(attempt_dir / "sample_complete.json", marker_data)
            state["samples"][sample_id].update(
                {
                    "status": "complete",
                    "message": "",
                    "completed_at": marker_data["completed_at"],
                }
            )
            event = "complete"
            message = "complete"
        except Exception as exc:
            state["samples"][sample_id].update(
                {
                    "status": "failed",
                    "message": f"{type(exc).__name__}: {exc}",
                    "completed_at": datetime.now().isoformat(timespec="seconds"),
                }
            )
            event = "failed"
            message = state["samples"][sample_id]["message"]
        state["updated_at"] = datetime.now().isoformat(timespec="seconds")
        _study_atomic_json(state_path, state)
        _write_study_aggregates(output_root, validated, state)
        if progress_callback:
            progress_callback({"event": event, "sample_id": sample_id, "position": position, "total": total, "message": message})
        if should_stop():
            state["run_status"] = "stopped"
            state["updated_at"] = datetime.now().isoformat(timespec="seconds")
            _study_atomic_json(state_path, state)
            if progress_callback:
                progress_callback(
                    {
                        "event": "stopped",
                        "sample_id": sample_id,
                        "position": position,
                        "total": total,
                        "message": "stopped after the current specimen",
                    }
                )
            break

    if state.get("run_status") != "stopped":
        state["run_status"] = "complete"
        state["updated_at"] = datetime.now().isoformat(timespec="seconds")
        _study_atomic_json(state_path, state)
    summaries = _write_study_aggregates(output_root, validated, state)
    return state, pd.DataFrame(summaries)

# =============================================================================

class SpermGUI:
    """
    The primary Tkinter-based graphical user interface for the Sperm Segmentation ROI Tool.

    The GUI provides a two-panel layout:
    - **Left sidebar** - controls for directory loading, Z-slice navigation,
      tool selection, ROI drawing, single-slice analysis, batch analysis,
      and two progress bars (2D segmentation and post-analysis reporting).
    - **Right canvas** - a matplotlib ``FigureCanvasTkAgg`` that renders the
      currently selected Z-slice image (raw, overlay, or debug) and accepts
      polygon ROI drawing interactions.

    Key interaction modes (controlled by ``mode_var``)
    ---------------------------------------------------
    - ``'view'``   - pans and inspects the raw image.
    - ``'review'`` - displays the saved overlay PNG for the current slice (requires
      a prior batch run).
    - ``'roi'``    - left-click to add polygon vertices; right-click to undo the
      last vertex; call *Finalize Polygon* to close and rasterise the mask.

    Thread architecture
    -------------------
    Batch processing (:meth:`run_batch_analysis`) runs in a background ``threading.Thread``
    so the Tkinter event loop remains responsive.  Progress bar updates are marshalled
    back to the main thread via ``root.after()`` calls.
    """
    def open_parameter_editor(self):
        """
        Opens the :class:`ParameterEditor` ``Toplevel`` window.

        Defines an ``on_apply`` callback that merges the edited values into the
        global ``CONFIG`` dict and updates the status label so the researcher
        knows the new parameters are active.
        """
        def on_apply(new_cfg):
            CONFIG.update(new_cfg)
            self.lbl_roi.config(text="Parameters updated in memory.")

        editor = ParameterEditor(self.root, CONFIG, self.default_config, on_apply)

    def _load_tuned_params(self):
        """Load a tuned parameters JSON file and merge into CONFIG."""
        from tkinter import filedialog, messagebox
        filepath = filedialog.askopenfilename(
            title="Select Tuned Parameters JSON",
            filetypes=[("JSON Files", "*.json"), ("All Files", "*.*")],
            initialdir=os.path.dirname(os.path.abspath(__file__))
        )
        if not filepath:
            return
        try:
            with open(filepath, 'r') as f:
                tuned = json.load(f)

            # Only update keys that exist in CONFIG (safety check)
            applied = []
            for key, value in tuned.items():
                if key in CONFIG:
                    old_val = CONFIG[key]
                    CONFIG[key] = value
                    applied.append(f"  {key}: {old_val} -> {value}")

            if applied:
                n = len(applied)
                short_name = os.path.basename(filepath)
                self.lbl_params_status.config(
                    text=f'OK Loaded {n} params from {short_name}',
                    fg='green'
                )
                detail = "\n".join(applied)
                messagebox.showinfo(
                    "Parameters Loaded",
                    f"Loaded {n} parameters from:\n{short_name}\n\n"
                    f"Changes applied:\n{detail}\n\n"
                    f"Use 'Revert Defaults' to undo."
                )
            else:
                messagebox.showwarning("No Matching Keys",
                    f"No recognized CONFIG keys found in {os.path.basename(filepath)}")
        except Exception as e:
            messagebox.showerror("Load Error", f"Failed to load parameters:\n{e}")

    def _launch_parameter_tuner(self, mode):
        """Launch the external V5.6 ROI-ADAPTIVE tuner without blocking the GUI."""
        from tkinter import messagebox, simpledialog

        if not self.input_dir:
            messagebox.showwarning(
                "Load Directory First",
                "Load the image directory before launching the tuner."
            )
            return

        default_slices = "0-12"
        if self.files:
            z_values = []
            for f in self.files:
                m = re.search(r"z(\d+)", os.path.basename(f), re.IGNORECASE)
                if m:
                    z_values.append(int(m.group(1)))
            if z_values:
                z_values = sorted(set(z_values))
                start = z_values[0]
                end = z_values[min(len(z_values) - 1, 12)]
                default_slices = f"{start}-{end}"

        slice_str = simpledialog.askstring(
            "Tuner Slices",
            "Enter representative z-slices for tuning (for example 0-12 or 10,15,20,25):",
            initialvalue=default_slices,
            parent=self.root
        )
        if not slice_str:
            return

        tuner_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "utils", "tune_parameters_Saturnv5_7.py")
        if not os.path.exists(tuner_path):
            messagebox.showerror("Tuner Missing", f"Could not find tuner:\n{tuner_path}")
            return

        outdir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "parameter_tuning_results")
        ensure_dir(outdir)
        try:
            roi_mask = self.build_roi_mask()
            if roi_mask is not None:
                roi_path = os.path.join(outdir, "last_drawn_roi_saturnv5_7_tune.tif")
                tifffile.imwrite(roi_path, roi_mask.astype(np.uint8) * 255)
        except Exception as e:
            print(f"Could not export current ROI for tuner: {e}")

        cmd = [
            sys.executable,
            tuner_path,
            "--mode", mode,
            "--dir", self.input_dir,
            "--slices", slice_str,
            "--outdir", outdir,
        ]

        if mode == "segmentation":
            cmd.extend(["--maxiter", "6", "--popsize", "6", "--review-candidates", "8"])

        try:
            creationflags = getattr(subprocess, "CREATE_NEW_CONSOLE", 0)
            subprocess.Popen(cmd, cwd=os.path.dirname(os.path.abspath(__file__)), creationflags=creationflags)
            self.lbl_params_status.config(
                text=f"Started {mode} tuner for z={slice_str}",
                fg="#0f766e"
            )
            messagebox.showinfo(
                "Tuner Started",
                f"Started {mode} tuning in a separate process.\n\n"
                f"Results will be saved to:\n{outdir}\n\n"
                "When it finishes, use 'Load Tuned Params' to apply a JSON candidate."
            )
        except Exception as e:
            messagebox.showerror("Tuner Launch Error", f"Could not launch tuner:\n{e}")

    def _on_sidebar_frame_configure(self, event=None):
        self.sidebar_canvas.configure(scrollregion=self.sidebar_canvas.bbox('all'))

    def _on_sidebar_canvas_configure(self, event):
        target_width = max(event.width - 4, 320)
        self.sidebar_canvas.itemconfigure(self.sidebar_window, width=target_width)

    def _bind_sidebar_widget_tree(self, widget):
        """Make mouse wheel scrolling work over buttons, labels, entries, and frames."""
        widget.bind('<MouseWheel>', self._on_sidebar_mousewheel, add='+')
        widget.bind('<Shift-MouseWheel>', self._on_sidebar_shift_mousewheel, add='+')
        widget.bind('<Button-4>', self._on_sidebar_linux_scroll_up, add='+')
        widget.bind('<Button-5>', self._on_sidebar_linux_scroll_down, add='+')
        widget.bind('<Home>', self._sidebar_scroll_top, add='+')
        widget.bind('<End>', self._sidebar_scroll_bottom, add='+')
        widget.bind('<Prior>', self._sidebar_page_up, add='+')
        widget.bind('<Next>', self._sidebar_page_down, add='+')
        for child in widget.winfo_children():
            self._bind_sidebar_widget_tree(child)

    def _bind_sidebar_mousewheel(self, event=None):
        self.sidebar_canvas.bind_all('<MouseWheel>', self._on_sidebar_mousewheel)
        self.sidebar_canvas.bind_all('<Shift-MouseWheel>', self._on_sidebar_shift_mousewheel)
        self.sidebar_canvas.bind_all('<Button-4>', self._on_sidebar_linux_scroll_up)
        self.sidebar_canvas.bind_all('<Button-5>', self._on_sidebar_linux_scroll_down)

    def _unbind_sidebar_mousewheel(self, event=None):
        self.sidebar_canvas.unbind_all('<MouseWheel>')
        self.sidebar_canvas.unbind_all('<Shift-MouseWheel>')
        self.sidebar_canvas.unbind_all('<Button-4>')
        self.sidebar_canvas.unbind_all('<Button-5>')

    def _on_sidebar_mousewheel(self, event):
        delta = getattr(event, 'delta', 0)
        if delta == 0:
            return
        step = -1 * int(delta / 120) if abs(delta) >= 120 else (-1 if delta > 0 else 1)
        self.sidebar_canvas.yview_scroll(step * 4, 'units')
        return "break"

    def _on_sidebar_shift_mousewheel(self, event):
        delta = getattr(event, 'delta', 0)
        if delta == 0:
            return
        step = -1 * int(delta / 120) if abs(delta) >= 120 else (-1 if delta > 0 else 1)
        self.sidebar_canvas.xview_scroll(step * 4, 'units')
        return "break"

    def _on_sidebar_linux_scroll_up(self, event):
        self.sidebar_canvas.yview_scroll(-4, 'units')
        return "break"

    def _on_sidebar_linux_scroll_down(self, event):
        self.sidebar_canvas.yview_scroll(4, 'units')
        return "break"

    def _sidebar_scroll_top(self, event=None):
        self.sidebar_canvas.yview_moveto(0.0)
        return "break"

    def _sidebar_scroll_bottom(self, event=None):
        self.sidebar_canvas.yview_moveto(1.0)
        return "break"

    def _sidebar_page_up(self, event=None):
        self.sidebar_canvas.yview_scroll(-1, 'pages')
        return "break"

    def _sidebar_page_down(self, event=None):
        self.sidebar_canvas.yview_scroll(1, 'pages')
        return "break"

    def _make_sidebar_section(self, parent, title, default_open=True, accent="#dbeafe"):
        """Create a compact collapsible section for the main sidebar."""
        outer = tk.Frame(parent, bg="#f0f0f0", bd=1, relief="groove")
        outer.pack(fill="x", padx=6, pady=(6, 2))

        body = tk.Frame(outer, bg="#f7f7f7")
        state = {"open": bool(default_open)}

        def sync_header():
            marker = "-" if state["open"] else "+"
            header.config(text=f"{marker} {title}")

        def toggle():
            state["open"] = not state["open"]
            if state["open"]:
                body.pack(fill="x", padx=0, pady=(0, 4))
            else:
                body.pack_forget()
            sync_header()
            self._on_sidebar_frame_configure()

        header = tk.Button(
            outer,
            command=toggle,
            anchor="w",
            bg=accent,
            relief="flat",
            font=("Arial", 9, "bold"),
            padx=8,
            pady=4
        )
        header.pack(fill="x")
        if state["open"]:
            body.pack(fill="x", padx=0, pady=(0, 4))
        sync_header()
        return body

    def _revert_to_defaults(self):
        """Revert CONFIG back to the original defaults captured at startup."""
        from tkinter import messagebox
        CONFIG.update(self.default_config)
        self.lbl_params_status.config(
            text='Using default parameters',
            fg='#555'
        )
        messagebox.showinfo("Reverted", "All parameters have been reverted to their original defaults.")

    def __init__(self, root):
        """
        Initialises the main application window, all sidebar controls, the matplotlib
        canvas, and mouse/key event bindings.

        Args:
            root (tk.Tk): The root Tkinter window created by :func:`launch_gui`.
        """
        self.root = root
        self.root.title(f'Sperm Segmentation ROI Tool - Saturn Project')
        self.root.geometry('1450x920')

        self.input_dir = ''
        self.files = []
        self.current_idx = 0
        self.current_img = None
        self.last_out_dir = ""

        self.study_rows = []
        self.study_root_dir = ""
        self.study_output_dir = ""
        self.study_window = None
        self.study_tree = None
        self.study_run_button = None
        self.study_stop_button = None
        self.study_progress_bar = None
        self.study_status_var = tk.StringVar(value="No study loaded")
        self.study_output_var = tk.StringVar(value="Output: not selected")
        self.study_progress_var = tk.DoubleVar(value=0)
        self.study_progress_text_var = tk.StringVar(value="Progress: 0 / 0 specimens")
        self._study_running = False
        self._study_stop_event = threading.Event()

        self.roi_points = []
        self.drawing = False
        self.roi_active = False
        self._loaded_roi_mask = None

        # --- AI ANALYSIS CONTROLS ---
        self.species_var = tk.StringVar(value='D. melanogaster')
        self.model_var = tk.StringVar(value='Gemini 2.5 Pro (Thinking)')
        self._last_batch_ts = None
        self._last_batch_out_dir = None

        # --- Main GUI sidebar: flexibly scrollable with mouse wheel, touchpad, and horizontal support ---
        self.sidebar_container = tk.Frame(root, width=340, bg='#f0f0f0')
        self.sidebar_container.pack(side='left', fill='y')
        self.sidebar_container.pack_propagate(False)

        self.sidebar_canvas = tk.Canvas(self.sidebar_container, bg='#f0f0f0', highlightthickness=0)
        self.sidebar_scrollbar = ttk.Scrollbar(self.sidebar_container, orient='vertical', command=self.sidebar_canvas.yview)
        self.sidebar_hscrollbar = ttk.Scrollbar(self.sidebar_container, orient='horizontal', command=self.sidebar_canvas.xview)
        self.sidebar_canvas.configure(yscrollcommand=self.sidebar_scrollbar.set, xscrollcommand=self.sidebar_hscrollbar.set)

        self.sidebar_quick_actions = tk.Frame(self.sidebar_container, bg='#e5e7eb', bd=1, relief='ridge')
        self.sidebar_quick_actions.pack(side='bottom', fill='x')
        tk.Label(
            self.sidebar_quick_actions,
            text='Quick Actions',
            bg='#e5e7eb',
            fg='#374151',
            font=('Arial', 8, 'bold')
        ).pack(anchor='w', padx=6, pady=(4, 0))
        quick_row1 = tk.Frame(self.sidebar_quick_actions, bg='#e5e7eb')
        quick_row1.pack(fill='x', padx=6, pady=(3, 2))
        tk.Button(
            quick_row1,
            text='Run Slice',
            command=self.run_analysis_slice,
            width=10
        ).pack(side='left', expand=True, fill='x', padx=(0, 2))
        tk.Button(
            quick_row1,
            text='Run Batch',
            command=self.run_batch_analysis,
            bg='#d4edda',
            font=('Arial', 9, 'bold'),
            width=10
        ).pack(side='right', expand=True, fill='x', padx=(2, 0))
        quick_row2 = tk.Frame(self.sidebar_quick_actions, bg='#e5e7eb')
        quick_row2.pack(fill='x', padx=6, pady=(2, 6))
        tk.Button(quick_row2, text='Top', command=self._sidebar_scroll_top, width=6).pack(side='left', fill='x', padx=(0, 2))
        tk.Button(quick_row2, text='Bottom', command=self._sidebar_scroll_bottom, width=7).pack(side='left', fill='x', padx=(2, 2))
        tk.Button(quick_row2, text='Save ROI', command=self.save_roi_mask, width=8).pack(side='left', expand=True, fill='x', padx=(2, 2))
        tk.Button(quick_row2, text='Load ROI', command=self.load_roi_mask, width=8).pack(side='right', expand=True, fill='x', padx=(2, 0))

        self.sidebar_canvas.pack(side='left', fill='both', expand=True)
        self.sidebar_scrollbar.pack(side='right', fill='y')
        self.sidebar_hscrollbar.pack(side='bottom', fill='x')

        self.sidebar = tk.Frame(self.sidebar_canvas, bg='#f0f0f0')
        self.sidebar_window = self.sidebar_canvas.create_window((0, 0), window=self.sidebar, anchor='nw')

        self.sidebar.bind("<Configure>", self._on_sidebar_frame_configure)
        self.sidebar_canvas.bind("<Configure>", self._on_sidebar_canvas_configure)
        self.sidebar_canvas.bind('<Enter>', self._bind_sidebar_mousewheel)
        self.sidebar_canvas.bind('<Leave>', self._unbind_sidebar_mousewheel)
        self.sidebar_container.bind('<Enter>', self._bind_sidebar_mousewheel)
        self.sidebar_container.bind('<Leave>', self._unbind_sidebar_mousewheel)
        # -----------------------------------------------------------------------------------------------

        self.default_config = CONFIG.copy()
        self.mode_var = tk.StringVar(value='view')

        data_section = self._make_sidebar_section(self.sidebar, "Data", default_open=True, accent="#dbeafe")
        tk.Button(data_section, text='Load Directory', command=self.load_directory, height=2).pack(fill='x', padx=8, pady=(8, 4))
        self.lbl_status = tk.Label(data_section, text='No directory loaded', wraplength=260, justify='left', bg="#f7f7f7", fg="#374151")
        self.lbl_status.pack(fill='x', padx=8, pady=(0, 8))

        study_section = self._make_sidebar_section(self.sidebar, "Multi-Sample Study", default_open=False, accent="#ccfbf1")
        tk.Button(
            study_section,
            text="Open Study Manager",
            command=self.open_multisample_study,
            bg="#ccfbf1",
            font=("Arial", 9, "bold"),
        ).pack(fill="x", padx=8, pady=(8, 4))
        tk.Label(
            study_section,
            textvariable=self.study_status_var,
            wraplength=260,
            justify="left",
            bg="#f7f7f7",
            fg="#374151",
            font=("Arial", 8),
        ).pack(fill="x", padx=8, pady=(0, 8))

        params_section = self._make_sidebar_section(self.sidebar, "Parameters", default_open=True, accent="#e2e3e5")
        tk.Button(params_section, text='Configure Parameters', command=self.open_parameter_editor, bg='#e2e3e5').pack(fill='x', padx=8, pady=(8, 4))
        params_frame = tk.Frame(params_section, bg='#f7f7f7')
        params_frame.pack(fill='x', padx=8, pady=(0, 4))
        tk.Button(params_frame, text='Load Tuned Params', command=self._load_tuned_params, bg='#d4edda', width=16).pack(side='left', expand=True, fill='x', padx=(0, 2))
        tk.Button(params_frame, text='Revert Defaults', command=self._revert_to_defaults, bg='#f8d7da', width=14).pack(side='right', expand=True, fill='x', padx=(2, 0))
        self.lbl_params_status = tk.Label(params_section, text='Using default parameters', wraplength=260, justify='left', fg='#555', bg="#f7f7f7", font=('Arial', 8))
        self.lbl_params_status.pack(fill='x', padx=8, pady=(0, 8))

        tuning_section = self._make_sidebar_section(self.sidebar, "Tuning", default_open=True, accent="#cffafe")
        tuner_frame = tk.Frame(tuning_section, bg='#f7f7f7')
        tuner_frame.pack(fill='x', padx=8, pady=8)
        tk.Button(
            tuner_frame,
            text='Tune Segmentation',
            command=lambda: self._launch_parameter_tuner("segmentation"),
            bg='#cffafe',
            width=16
        ).pack(side='left', expand=True, fill='x', padx=(0, 2))
        tk.Button(
            tuner_frame,
            text='Tune Tracking',
            command=lambda: self._launch_parameter_tuner("tracking"),
            bg='#e0e7ff',
            width=14
        ).pack(side='right', expand=True, fill='x', padx=(2, 0))

        nav_section = self._make_sidebar_section(self.sidebar, "Z Navigation", default_open=True, accent="#fef3c7")
        self.scale_z = tk.Scale(nav_section, from_=0, to=0, orient='horizontal', command=self.on_slide_change, bg="#f7f7f7", highlightthickness=0)
        self.scale_z.pack(fill='x', padx=8, pady=(8, 2))
        self.lbl_z = tk.Label(nav_section, text='Z: 0 / 0', bg="#f7f7f7")
        self.lbl_z.pack(pady=(0, 8))

        tools_section = self._make_sidebar_section(self.sidebar, "View / ROI Drawing", default_open=True, accent="#fde68a")
        tk.Radiobutton(tools_section, text='View/Nav (Raw Image)', variable=self.mode_var, value='view', command=self.render, bg="#f7f7f7").pack(anchor='w', padx=10, pady=(6, 0))
        tk.Radiobutton(tools_section, text='Review Overlays (After Batch)', variable=self.mode_var, value='review', command=self.render, bg="#f7f7f7").pack(anchor='w', padx=10)
        tk.Radiobutton(tools_section, text='Draw ROI (Polygon)', variable=self.mode_var, value='roi', command=self.render, bg="#f7f7f7").pack(anchor='w', padx=10)
        tk.Label(tools_section, text='Left-click points, right-click undo', font=('Arial', 8, 'italic'), fg='dimgray', bg="#f7f7f7").pack(fill='x', padx=10, pady=(2, 4))
        tk.Button(tools_section, text='Finalize Polygon', command=self.finalize_roi, bg='#ffeeba').pack(fill='x', padx=20, pady=(0, 8))

        roi_section = self._make_sidebar_section(self.sidebar, "ROI Files", default_open=True, accent="#dcfce7")
        roi_buttons = tk.Frame(roi_section, bg="#f7f7f7")
        roi_buttons.pack(fill='x', padx=8, pady=(8, 4))
        tk.Button(roi_buttons, text='Reset ROI', command=self.reset_roi).pack(side='left', expand=True, fill='x', padx=(0, 2))
        tk.Button(roi_buttons, text='Save ROI', command=self.save_roi_mask).pack(side='left', expand=True, fill='x', padx=2)
        tk.Button(roi_buttons, text='Load ROI', command=self.load_roi_mask).pack(side='right', expand=True, fill='x', padx=(2, 0))
        self.lbl_roi = tk.Label(roi_section, text='ROI: none', wraplength=260, justify='left', bg="#f7f7f7", fg="#374151")
        self.lbl_roi.pack(fill='x', padx=8, pady=(0, 8))

        analysis_section = self._make_sidebar_section(self.sidebar, "Analysis", default_open=True, accent="#d4edda")
        tk.Button(analysis_section, text='Run Analysis on Slice', command=self.run_analysis_slice).pack(fill='x', padx=8, pady=(8, 4))
        tk.Button(analysis_section, text='Run Batch (All Slices + 3D Track)', command=self.run_batch_analysis, bg='#d4edda', font=('Arial', 10, 'bold')).pack(fill='x', padx=8, pady=(0, 8))

        ai_section = self._make_sidebar_section(self.sidebar, f'{_VERSION} AI Biological Analysis', default_open=False, accent="#ede9fe")
        species_list = [
            'D. melanogaster', 'D. simulans', 'D. yakuba',
            'D. ananassae', 'D. pseudoobscura (Dpse)',
            'D. virilis (Dvir)', 'General / Evolutionary'
        ]
        self.species_dropdown = ttk.Combobox(ai_section, textvariable=self.species_var, values=species_list, state='readonly')
        self.species_dropdown.pack(fill='x', padx=8, pady=(8, 5))
        tk.Label(ai_section, text='AI Model:', font=('Arial', 8), bg="#f7f7f7").pack(anchor='w', padx=8)
        models = ['Gemini 2.5 Pro (Thinking)', 'Gemini 2.5 Flash (Fast)']
        self.model_dropdown = ttk.Combobox(ai_section, textvariable=self.model_var, values=models, state='readonly')
        self.model_dropdown.pack(fill='x', padx=8, pady=(0, 5))
        tk.Button(ai_section, text='Set API Key (Free)', command=self.set_ai_key, font=('Arial', 8)).pack(fill='x', padx=8, pady=(0, 5))
        self.btn_ai = tk.Button(ai_section, text='Run AI Analysis', command=self.run_ai_analysis,
                                 bg='#8b5cf6', fg='white', font=('Arial', 9, 'bold'), state='disabled')
        self.btn_ai.pack(fill='x', padx=8, pady=(0, 8))

        self._bind_sidebar_widget_tree(self.sidebar)
        self._bind_sidebar_widget_tree(self.sidebar_quick_actions)
        self.sidebar_canvas.bind('<Button-1>', lambda event: self.sidebar_canvas.focus_set(), add='+')

        self.canvas_frame = tk.Frame(root, bg='black')
        self.canvas_frame.pack(side='right', expand=True, fill='both')

        # --- TOP STATUS BAR (For Progress Visibility) ---
        self.top_status_frame = tk.Frame(self.canvas_frame, bg='#f0f0f0', height=80)
        self.top_status_frame.pack(side='top', fill='x')

        # --- Sub-frame for Status (Left) ---
        self.status_sub_frame = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.status_sub_frame.pack(side='left', padx=10, fill='y')

        # --- Dynamic Status Label ---
        self.lbl_batch_op = tk.Label(self.status_sub_frame, text='GUI Ready', font=('Arial', 10, 'bold'), fg='#2c3e50', bg='#f0f0f0')
        self.lbl_batch_op.pack(side='left', pady=5)

        # --- PROGRESS BOX: Sequential Progress Bars (At Top Right) ---
        self.p_container = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.p_container.pack(side='right', padx=20, pady=5)

        # Frame 1: 2D Batch Segmentation
        self.batch_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        self.batch_p_frame.pack(fill='x')
        tk.Label(self.batch_p_frame, text='Batch Progress (2D)', font=('Arial', 9, 'bold'), bg='#f0f0f0').pack(side='left', padx=10)
        self.progress = ttk.Progressbar(self.batch_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress.pack(side='left', padx=5, pady=5)
        self.lbl_progress_val = tk.Label(self.batch_p_frame, text='0%', font=('Arial', 10, 'bold'), fg='blue', bg='#f0f0f0')
        self.lbl_progress_val.pack(side='left', padx=10)

        # Frame 2: Post-Analysis (Initially Hidden)
        self.post_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        # We don't pack self.post_p_frame here
        tk.Label(self.post_p_frame, text='Post-Analysis Progress', font=('Arial', 9, 'bold'), bg="#f0f0f0").pack(side='left', padx=10)
        self.progress_post = ttk.Progressbar(self.post_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress_post.pack(side='left', padx=5, pady=5)
        self.lbl_post_progress_val = tk.Label(self.post_p_frame, text='Waiting...', font=('Arial', 10, 'bold'), fg='dimgray', bg="#f0f0f0")
        self.lbl_post_progress_val.pack(side='left', padx=10)

        self.fig = Figure(figsize=(8, 8), dpi=100)
        self.ax = self.fig.add_subplot(111)
        self.ax.axis('off')
        self.canvas = FigureCanvasTkAgg(self.fig, master=self.canvas_frame)
        self.canvas.get_tk_widget().pack(fill='both', expand=True)

        self.canvas.mpl_connect('button_press_event', self.on_click)
        self.canvas.mpl_connect('key_press_event', self.on_key)


    def open_multisample_study(self):
        """Open the manifest-backed study manager for independent specimens."""
        if self.study_window is not None and self.study_window.winfo_exists():
            self.study_window.deiconify()
            self.study_window.lift()
            self.study_window.focus_force()
            return

        window = tk.Toplevel(self.root)
        self.study_window = window
        window.title(f"Saturn {_VERSION} Multi-Sample Study Manager")
        window.geometry("1220x700")
        window.minsize(900, 520)
        window.protocol("WM_DELETE_WINDOW", self._study_close_window)

        toolbar = tk.Frame(window, bg="#e5e7eb", bd=1, relief="ridge")
        toolbar.pack(fill="x")
        for text, command in (
            ("Discover Root", self._study_discover_root),
            ("Load Manifest", self._study_load_manifest),
            ("Save Manifest", self._study_save_manifest),
            ("Assign Group", self._study_assign_group),
            ("Organize Dataset Copy", self._study_organize_copy),
            ("Select Output", self._study_select_output),
            ("Validate", self._study_validate),
        ):
            tk.Button(toolbar, text=text, command=command).pack(side="left", padx=4, pady=6)
        self.study_run_button = tk.Button(
            toolbar,
            text="Run / Resume Study",
            command=self._study_run,
            bg="#bbf7d0",
            font=("Arial", 9, "bold"),
        )
        self.study_run_button.pack(side="right", padx=6, pady=6)
        self.study_stop_button = tk.Button(
            toolbar,
            text="Stop After Current Sample",
            command=self._study_request_stop,
            bg="#fecaca",
            state="disabled",
        )
        self.study_stop_button.pack(side="right", padx=4, pady=6)

        status = tk.Frame(window, bg="#f8fafc")
        status.pack(fill="x", padx=8, pady=(6, 2))
        tk.Label(status, textvariable=self.study_status_var, anchor="w", bg="#f8fafc").pack(fill="x")
        tk.Label(status, textvariable=self.study_output_var, anchor="w", bg="#f8fafc", fg="#475569").pack(fill="x")
        progress_row = tk.Frame(status, bg="#f8fafc")
        progress_row.pack(fill="x", pady=(5, 1))
        self.study_progress_bar = ttk.Progressbar(
            progress_row,
            orient="horizontal",
            mode="determinate",
            variable=self.study_progress_var,
            maximum=1,
        )
        self.study_progress_bar.pack(side="left", fill="x", expand=True)
        tk.Label(
            progress_row,
            textvariable=self.study_progress_text_var,
            anchor="e",
            bg="#f8fafc",
            fg="#475569",
            width=27,
        ).pack(side="right", padx=(8, 0))

        table_frame = tk.Frame(window)
        table_frame.pack(fill="both", expand=True, padx=8, pady=6)
        columns = (
            "include", "sample_id", "group", "slices", "z_range", "roi",
            "xy", "z_step", "status", "message",
        )
        tree = ttk.Treeview(table_frame, columns=columns, show="headings", selectmode="extended")
        self.study_tree = tree
        headings = {
            "include": "Include", "sample_id": "Sample ID", "group": "Group",
            "slices": "Slices", "z_range": "Z range", "roi": "ROI",
            "xy": "XY um/px", "z_step": "Z um/slice", "status": "Status",
            "message": "Message",
        }
        widths = {
            "include": 60, "sample_id": 150, "group": 120, "slices": 55,
            "z_range": 70, "roi": 65, "xy": 80, "z_step": 85,
            "status": 85, "message": 330,
        }
        for column in columns:
            tree.heading(column, text=headings[column])
            tree.column(column, width=widths[column], minwidth=45, stretch=(column == "message"))
        yscroll = ttk.Scrollbar(table_frame, orient="vertical", command=tree.yview)
        xscroll = ttk.Scrollbar(table_frame, orient="horizontal", command=tree.xview)
        tree.configure(yscrollcommand=yscroll.set, xscrollcommand=xscroll.set)
        tree.grid(row=0, column=0, sticky="nsew")
        yscroll.grid(row=0, column=1, sticky="ns")
        xscroll.grid(row=1, column=0, sticky="ew")
        table_frame.rowconfigure(0, weight=1)
        table_frame.columnconfigure(0, weight=1)
        tree.bind("<Double-1>", self._study_edit_cell)
        tree.bind("<space>", self._study_toggle_selected)

        footer = tk.Label(
            window,
            text=(
                "Select one or more rows to Assign Group. Double-click Include, "
                "Sample ID, Group, XY, or Z spacing to edit. Space toggles Include."
            ),
            anchor="w",
            fg="#475569",
        )
        footer.pack(fill="x", padx=10, pady=(0, 8))
        self._study_refresh_tree()

    def _study_close_window(self):
        if self._study_running:
            messagebox.showinfo("Study Running", "The study manager remains open while processing is active.")
            return
        if self.study_window is not None:
            self.study_window.destroy()
        self.study_window = None
        self.study_tree = None
        self.study_run_button = None
        self.study_stop_button = None
        self.study_progress_bar = None

    def _study_refresh_tree(self):
        tree = self.study_tree
        if tree is None or not tree.winfo_exists():
            return
        selected = tree.selection()
        selected_id = selected[0] if selected else None
        tree.delete(*tree.get_children())
        for index, row in enumerate(self.study_rows):
            iid = f"row_{index}"
            roi_ok = pl.Path(str(row.get("roi_path", ""))).is_file()
            tree.insert(
                "",
                "end",
                iid=iid,
                values=(
                    "Yes" if _study_bool(row.get("include", True)) else "No",
                    row.get("sample_id", ""),
                    row.get("group", ""),
                    row.get("slice_count", 0),
                    f"{row.get('z_min', 0)}-{row.get('z_max', 0)}",
                    "Ready" if roi_ok else "Missing",
                    f"{float(row.get('xy_um_per_pixel', 0) or 0):.6g}",
                    f"{float(row.get('z_um_per_slice', 0) or 0):.6g}",
                    row.get("status", "pending"),
                    row.get("message", ""),
                ),
            )
        if selected_id and tree.exists(selected_id):
            tree.selection_set(selected_id)
            tree.see(selected_id)

    def _study_selected_index(self):
        if self.study_tree is None:
            return None
        selected = self.study_tree.selection()
        if not selected:
            return None
        try:
            return int(selected[0].rsplit("_", 1)[-1])
        except (ValueError, IndexError):
            return None

    def _study_selected_indices(self):
        if self.study_tree is None:
            return []
        indices = []
        for item in self.study_tree.selection():
            try:
                indices.append(int(item.rsplit("_", 1)[-1]))
            except (ValueError, IndexError):
                continue
        return sorted(set(indices))

    def _study_toggle_selected(self, event=None):
        indices = self._study_selected_indices()
        if not indices or self._study_running:
            return "break"
        for index in indices:
            self.study_rows[index]["include"] = not _study_bool(
                self.study_rows[index].get("include", True)
            )
            self.study_rows[index]["status"] = (
                "pending" if self.study_rows[index]["include"] else "excluded"
            )
        self._study_refresh_tree()
        return "break"

    def _study_edit_cell(self, event):
        if self._study_running or self.study_tree is None:
            return
        region = self.study_tree.identify_region(event.x, event.y)
        item = self.study_tree.identify_row(event.y)
        column_id = self.study_tree.identify_column(event.x)
        if region != "cell" or not item or not column_id:
            return
        self.study_tree.selection_set(item)
        index = self._study_selected_index()
        if index is None:
            return
        column_index = int(column_id[1:]) - 1
        keys = ("include", "sample_id", "group", "slice_count", "z_range", "roi", "xy_um_per_pixel", "z_um_per_slice", "status", "message")
        key = keys[column_index]
        if key == "include":
            self._study_toggle_selected()
            return
        if key not in {"sample_id", "group", "xy_um_per_pixel", "z_um_per_slice"}:
            return

        from tkinter import simpledialog
        current = self.study_rows[index].get(key, "")
        label = {
            "sample_id": "Sample ID",
            "group": "Biological group",
            "xy_um_per_pixel": "XY calibration (um/pixel)",
            "z_um_per_slice": "Z calibration (um/slice)",
        }[key]
        value = simpledialog.askstring("Edit Study Row", label, initialvalue=str(current), parent=self.study_window)
        if value is None:
            return
        if key in {"xy_um_per_pixel", "z_um_per_slice"}:
            try:
                value = float(value)
            except ValueError:
                messagebox.showerror("Invalid Calibration", f"{label} must be a number.", parent=self.study_window)
                return
        elif not value.strip():
            messagebox.showerror("Invalid Value", f"{label} cannot be blank.", parent=self.study_window)
            return
        self.study_rows[index][key] = value
        self.study_rows[index]["status"] = "pending"
        self.study_rows[index]["message"] = ""
        self._study_refresh_tree()

    def _study_discover_root(self):
        if self._study_running:
            return
        root_dir = filedialog.askdirectory(
            title="Select the folder containing all biological samples",
            initialdir=self.study_root_dir or CONFIG.get("INPUT_DIR", os.getcwd()),
            parent=self.study_window,
        )
        if not root_dir:
            return
        try:
            rows = discover_multisample_study(root_dir, base_cfg=CONFIG)
            if not rows:
                raise ValueError("No exact ch00 Z-stacks were found below the selected folder.")
            self.study_root_dir = root_dir
            self.study_rows = rows
            self.study_status_var.set(f"Discovered {len(rows)} samples in {root_dir}")
            self._study_refresh_tree()
        except Exception as exc:
            messagebox.showerror("Study Discovery Failed", str(exc), parent=self.study_window)

    def _study_load_manifest(self):
        if self._study_running:
            return
        path = filedialog.askopenfilename(
            title="Load study manifest",
            filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
            parent=self.study_window,
        )
        if not path:
            return
        try:
            self.study_rows = load_multisample_manifest(path)
            input_dirs = [pl.Path(str(row["input_dir"])) for row in self.study_rows if row.get("input_dir")]
            if input_dirs:
                self.study_root_dir = str(pl.Path(os.path.commonpath([str(path) for path in input_dirs])))
            self.study_status_var.set(f"Loaded {len(self.study_rows)} samples from {path}")
            self._study_refresh_tree()
        except Exception as exc:
            messagebox.showerror("Manifest Load Failed", str(exc), parent=self.study_window)

    def _study_save_manifest(self):
        if not self.study_rows:
            messagebox.showwarning("No Study", "Discover or load a study first.", parent=self.study_window)
            return
        path = filedialog.asksaveasfilename(
            title="Save study manifest",
            defaultextension=".csv",
            initialfile="study_manifest_v5_7.csv",
            filetypes=[("CSV files", "*.csv")],
            parent=self.study_window,
        )
        if not path:
            return
        try:
            save_multisample_manifest(self.study_rows, path)
            self.study_status_var.set(f"Saved manifest: {path}")
        except Exception as exc:
            messagebox.showerror("Manifest Save Failed", str(exc), parent=self.study_window)

    def _study_assign_group(self):
        if self._study_running:
            return
        indices = self._study_selected_indices()
        if not indices:
            messagebox.showwarning(
                "No Samples Selected",
                "Select one or more specimen rows first.",
                parent=self.study_window,
            )
            return
        from tkinter import simpledialog

        current_groups = {
            str(self.study_rows[index].get("group", "")).strip()
            for index in indices
        }
        initial = next(iter(current_groups)) if len(current_groups) == 1 else ""
        group = simpledialog.askstring(
            "Assign Biological Group",
            f"Group label for {len(indices)} selected specimen(s):",
            initialvalue=initial,
            parent=self.study_window,
        )
        if group is None:
            return
        group = group.strip()
        if not group:
            messagebox.showerror(
                "Invalid Group",
                "The biological group cannot be blank.",
                parent=self.study_window,
            )
            return
        for index in indices:
            self.study_rows[index]["group"] = group
            self.study_rows[index]["status"] = "pending"
            self.study_rows[index]["message"] = ""
        self.study_status_var.set(
            f"Assigned {len(indices)} specimen(s) to group {group!r}"
        )
        self._study_refresh_tree()

    def _study_organize_copy(self):
        if self._study_running:
            return
        if not self.study_rows or not self.study_root_dir:
            messagebox.showwarning(
                "Discover Study First",
                "Use Discover Root, review the specimens, and assign biological "
                "groups before organizing the dataset.",
                parent=self.study_window,
            )
            return
        source_root = self.study_root_dir
        source_rows = [dict(row) for row in self.study_rows]

        output_root = filedialog.askdirectory(
            title="Select an EMPTY folder for the organized dataset copy",
            initialdir=str(pl.Path(source_root).resolve().parent),
            mustexist=False,
            parent=self.study_window,
        )
        if not output_root:
            return

        slice_count = sum(int(row.get("slice_count", 0)) for row in source_rows)
        proceed = messagebox.askyesno(
            "Organize Dataset Copy",
            "Create a non-destructive canonical copy?\n\n"
            f"Discovered specimens: {len(source_rows)}\n"
            f"TIFF planes: {slice_count}\n\n"
            "The biological group labels currently shown in the table will "
            "define the organized group folders.\n\n"
            f"Source:\n{source_root}\n\n"
            f"Organized copy:\n{output_root}\n\n"
            "Original files will not be renamed, moved, or deleted.",
            parent=self.study_window,
        )
        if not proceed:
            return

        self._study_running = True
        if self.study_run_button is not None:
            self.study_run_button.config(state="disabled", text="Organizing...")
        self.study_status_var.set(
            f"Preparing canonical copy of {len(source_rows)} specimens"
        )

        def progress(event):
            self.root.after(
                0,
                lambda item=dict(event): self._study_organization_progress(item),
            )

        def worker():
            try:
                organized_rows, summary = organize_multisample_study_copy(
                    source_rows,
                    output_root,
                    progress_callback=progress,
                )
                self.root.after(
                    0,
                    lambda: self._study_organization_finished(
                        output_root,
                        organized_rows,
                        summary,
                    ),
                )
            except Exception as exc:
                detail = f"{type(exc).__name__}: {exc}"
                self.root.after(
                    0,
                    lambda: self._study_organization_failed(detail),
                )

        threading.Thread(target=worker, daemon=True).start()

    def _study_organization_progress(self, event):
        self.study_status_var.set(
            f"[{event.get('position', 0)}/{event.get('total', 0)}] "
            f"{event.get('sample_id', '')}: {event.get('message', '')}"
        )

    def _study_organization_finished(self, output_root, organized_rows, summary):
        self._study_running = False
        if self.study_run_button is not None and self.study_run_button.winfo_exists():
            self.study_run_button.config(state="normal", text="Run / Resume Study")
        self.study_root_dir = str(pl.Path(output_root).resolve())
        self.study_rows = organized_rows
        missing_roi = len(summary.get("samples_missing_roi", []))
        self.study_status_var.set(
            f"Organized {len(organized_rows)} samples; {missing_roi} need ROI"
        )
        self._study_refresh_tree()
        messagebox.showinfo(
            "Dataset Organization Complete",
            f"Samples organized: {len(organized_rows)}\n"
            f"Files copied: {summary.get('copied_files', 0)}\n"
            f"Files reused: {summary.get('reused_files', 0)}\n"
            f"Samples missing ROI: {missing_roi}\n\n"
            f"Manifest:\n{summary.get('manifest_path', '')}\n\n"
            "The organized study has been loaded into the manager.",
            parent=self.study_window,
        )

    def _study_organization_failed(self, detail):
        self._study_running = False
        if self.study_run_button is not None and self.study_run_button.winfo_exists():
            self.study_run_button.config(state="normal", text="Run / Resume Study")
        self.study_status_var.set("Dataset organization failed")
        messagebox.showerror(
            "Dataset Organization Failed",
            detail,
            parent=self.study_window,
        )

    def _study_select_output(self):
        if self._study_running:
            return
        output_dir = filedialog.askdirectory(
            title="Select or create a separate study output folder",
            initialdir=self.study_output_dir or os.path.dirname(os.path.abspath(__file__)),
            mustexist=False,
            parent=self.study_window,
        )
        if output_dir:
            self.study_output_dir = output_dir
            self.study_output_var.set(f"Output: {output_dir}")

    def _study_validate(self, show_dialog=True):
        if not self.study_rows:
            if show_dialog:
                messagebox.showwarning("No Study", "Discover or load a study first.", parent=self.study_window)
            return False
        self.study_rows, errors = validate_multisample_manifest(self.study_rows)
        included = [row for row in self.study_rows if row["include"]]
        invalid = [row for row in included if row["status"] == "invalid"]
        warning_count = sum("unusual Z spacing" in row.get("message", "") for row in included)
        self._study_refresh_tree()
        if invalid:
            self.study_status_var.set(f"Validation failed: {len(invalid)} of {len(included)} included samples are invalid")
            if show_dialog:
                details = "\n".join(errors[:15])
                if len(errors) > 15:
                    details += f"\n... and {len(errors) - 15} more"
                messagebox.showerror("Study Validation Failed", details, parent=self.study_window)
            return False
        self.study_status_var.set(f"Validated {len(included)} included samples; {warning_count} calibration warnings")
        if show_dialog:
            messagebox.showinfo(
                "Study Validated",
                f"{len(included)} included samples are ready.\nCalibration warnings: {warning_count}",
                parent=self.study_window,
            )
        return bool(included)

    def _study_output_is_separate(self):
        if not self.study_output_dir:
            return False, "Select a study output folder first."
        output = pl.Path(self.study_output_dir).resolve()
        if self.study_root_dir:
            study_root = pl.Path(self.study_root_dir).resolve()
            try:
                output.relative_to(study_root)
                return False, "The output folder must be outside the source study folder."
            except ValueError:
                pass
        for row in self.study_rows:
            sample_dir = pl.Path(str(row.get("input_dir", ""))).resolve()
            if output == sample_dir:
                return False, f"The output folder is the input folder for {row.get('sample_id', 'a sample')}."
        return True, ""

    def _study_run(self):
        if self._study_running:
            return
        if not self._study_validate(show_dialog=False):
            messagebox.showerror("Study Not Ready", "Correct the invalid study rows before running.", parent=self.study_window)
            return
        separate, reason = self._study_output_is_separate()
        if not separate:
            messagebox.showerror("Invalid Output Folder", reason, parent=self.study_window)
            return
        included = sum(_study_bool(row.get("include", True)) for row in self.study_rows)
        proceed = messagebox.askyesno(
            "Run Multi-Sample Study",
            f"Run or resume {included} independent samples?\n\nOutput:\n{self.study_output_dir}",
            parent=self.study_window,
        )
        if not proceed:
            return

        self._study_running = True
        self._study_stop_event.clear()
        if self.study_run_button is not None:
            self.study_run_button.config(state="disabled", text="Study Running...")
        if self.study_stop_button is not None:
            self.study_stop_button.config(
                state="normal",
                text="Stop After Current Sample",
            )
        self.study_progress_var.set(0)
        if self.study_progress_bar is not None:
            self.study_progress_bar.config(maximum=max(included, 1))
        self.study_progress_text_var.set(f"Progress: 0 / {included} specimens")
        self.study_status_var.set(f"Starting {included}-sample study")
        rows = [dict(row) for row in self.study_rows]
        output_dir = self.study_output_dir
        cfg = CONFIG.copy()

        def progress(event):
            self.root.after(0, lambda item=dict(event): self._study_progress_event(item))

        def worker():
            try:
                state, summary = run_multisample_study(
                    rows,
                    output_dir,
                    base_cfg=cfg,
                    progress_callback=progress,
                    resume=True,
                    stop_requested=self._study_stop_event,
                )
                self.root.after(0, lambda: self._study_run_finished(state, summary))
            except Exception as exc:
                detail = f"{type(exc).__name__}: {exc}"
                self.root.after(0, lambda: self._study_run_failed(detail))

        threading.Thread(target=worker, daemon=True).start()

    def _study_request_stop(self):
        if not self._study_running or self._study_stop_event.is_set():
            return
        self._study_stop_event.set()
        if self.study_stop_button is not None:
            self.study_stop_button.config(
                state="disabled",
                text="Stop Requested",
            )
        self.study_status_var.set(
            "Stop requested; the current specimen will finish, then the study will pause."
        )

    def _study_progress_event(self, event):
        sample_id = event.get("sample_id", "")
        status = event.get("event", "running")
        message = event.get("message", "")
        if sample_id and status != "stopped":
            for row in self.study_rows:
                if row.get("sample_id") == sample_id:
                    row["status"] = "complete" if status in {"complete", "skipped"} else status
                    row["message"] = message
                    break
        position = int(event.get("position", 0) or 0)
        total = int(event.get("total", 0) or 0)
        if status == "started":
            completed_position = max(position - 1, 0)
        else:
            completed_position = position
        self.study_progress_var.set(completed_position)
        if self.study_progress_bar is not None:
            self.study_progress_bar.config(maximum=max(total, 1))
        self.study_progress_text_var.set(
            f"Progress: {completed_position} / {total} specimens"
        )
        if status == "stopped":
            self.study_status_var.set(message)
        else:
            self.study_status_var.set(
                f"[{position}/{total}] {sample_id}: {message}"
            )
        self._study_refresh_tree()

    def _study_run_finished(self, state, summary):
        self._study_running = False
        self._study_stop_event.clear()
        if self.study_run_button is not None and self.study_run_button.winfo_exists():
            self.study_run_button.config(state="normal", text="Run / Resume Study")
        if self.study_stop_button is not None and self.study_stop_button.winfo_exists():
            self.study_stop_button.config(
                state="disabled",
                text="Stop After Current Sample",
            )
        completed = sum(record.get("status") == "complete" for record in state.get("samples", {}).values())
        failed = sum(record.get("status") == "failed" for record in state.get("samples", {}).values())
        for row in self.study_rows:
            record = state.get("samples", {}).get(row.get("sample_id"), {})
            if record:
                row["status"] = record.get("status", row.get("status", "pending"))
                row["message"] = record.get("message", "")
                row["output_dir"] = record.get("output_dir", "")
        total = sum(_study_bool(row.get("include", True)) for row in self.study_rows)
        self.study_progress_var.set(min(completed + failed, total))
        self.study_progress_text_var.set(
            f"Progress: {min(completed + failed, total)} / {total} specimens"
        )
        stopped = state.get("run_status") == "stopped"
        self.study_status_var.set(
            (
                f"Study paused: {completed} complete, {failed} failed"
                if stopped
                else f"Study finished: {completed} complete, {failed} failed"
            )
        )
        self._study_refresh_tree()
        messagebox.showinfo(
            "Study Paused" if stopped else "Study Run Finished",
            f"Complete: {completed}\nFailed: {failed}\n"
            + (
                "Pending specimens can be continued with Run / Resume Study.\n\n"
                if stopped
                else "\n"
            )
            + f"Specimen analysis table:\n{pl.Path(self.study_output_dir) / 'specimen_summary.csv'}\n\n"
            + f"Two-group comparison:\n{pl.Path(self.study_output_dir) / 'specimen_group_comparisons.csv'}",
            parent=self.study_window,
        )

    def _study_run_failed(self, detail):
        self._study_running = False
        self._study_stop_event.clear()
        if self.study_run_button is not None and self.study_run_button.winfo_exists():
            self.study_run_button.config(state="normal", text="Run / Resume Study")
        if self.study_stop_button is not None and self.study_stop_button.winfo_exists():
            self.study_stop_button.config(
                state="disabled",
                text="Stop After Current Sample",
            )
        self.study_status_var.set(f"Study stopped: {detail}")
        messagebox.showerror("Study Run Failed", detail, parent=self.study_window)


    def set_ai_key(self):
        """Opens a small dialog to save the Gemini API Key to a local file."""
        win = tk.Toplevel(self.root)
        win.title("AI Key Management")
        win.geometry("420x260")
        win.grab_set()

        tk.Label(win, text="Gemini AI Key Setup", font=('Arial', 12, 'bold')).pack(pady=(15,5))

        link = tk.Label(win, text="Get your FREE Key at Google AI Studio", fg='blue', cursor="hand2", font=('Arial', 9, 'underline'))
        link.pack()
        link.bind("<Button-1>", lambda e: webbrowser.open("https://aistudio.google.com/app/apikey"))

        tk.Label(win, text="(No credit card required for Free Tier)", font=('Arial', 8), fg='#64748b').pack(pady=(0, 10))

        tk.Label(win, text="Paste API Key here:", font=('Arial', 9)).pack()
        entry = tk.Entry(win, show='*', width=45)
        entry.pack(pady=5)

        # Load current key if exists
        key_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gemini_api_key.txt")
        exists = os.path.exists(key_file)
        if exists:
            with open(key_file, 'r') as f:
                entry.insert(0, f.read().strip())

        def save():
            k = entry.get().strip()
            if not k:
                messagebox.showwarning("Warning", "Key cannot be empty.")
                return
            with open(key_file, 'w') as f:
                f.write(k)
            status = "updated" if exists else "saved"
            messagebox.showinfo("Success", f"API Key {status} successfully!\nPortable analysis is ready.")
            win.destroy()

        btn_text = "Update API Key" if exists else "Save Key Permanently"
        tk.Button(win, text=btn_text, command=save, bg='#8b5cf6', fg='white', font=('Arial', 9, 'bold'), padx=20).pack(pady=15)

    def run_ai_analysis(self):
        """Runs AI biological interpretation on the last completed batch data."""
        if self._last_batch_ts is None or self._last_batch_out_dir is None:
            messagebox.showinfo("AI Analysis", "Please run a batch analysis first before requesting AI interpretation.")
            return

        ts = self._last_batch_ts
        out_dir = self._last_batch_out_dir

        if ts.empty or len(ts) == 0:
            messagebox.showinfo("AI Analysis", "No valid 3D tracks found in the last batch. AI requires track data.")
            return

        self.lbl_batch_op.config(text='AI Analyzing Biological Data...', fg='#8b5cf6')
        self.btn_ai.config(state='disabled', text='AI Working...')
        self.root.update()

        # Send only the canonical estimated-nuclei population.
        ts_main = (
            ts[ts["technical_valid"].astype(bool)]
            if "technical_valid" in ts.columns
            else ts
        )
        csv_summary_str = ts_main.head(100).to_csv(index=False)
        species = self.species_var.get()
        folder_name = os.path.basename(self.input_dir or "Current Project")

        # Model Selection Mapping
        m_map = {
            'Gemini 2.5 Pro (Thinking)': 'gemini-2.5-pro',
            'Gemini 2.5 Flash (Fast)': 'gemini-2.5-flash'
        }
        model_id = m_map.get(self.model_var.get(), 'gemini-2.5-pro')

        def _ai_thread():
            try:
                print(f"AI START: Interpreting {len(ts)} tracks for {species} via {model_id}...")
                ai_text = get_ai_biological_interpretation(csv_summary_str, species, folder_name, model_id=model_id)

                if not ai_text or "AI ANALYSIS SKIPPED" in ai_text or "AI API Error" in ai_text:
                    err_hint = ""
                    if "(429)" in ai_text:
                        err_hint = "\n\nTIP: Gemini 2.5 Pro has lower free-tier quotas. Try switching to 'Gemini 2.5 Flash' in the sidebar or wait a minute."
                    print(f"AI FAIL: {ai_text}")
                    self.root.after(0, lambda: messagebox.showwarning("AI Analysis", f"The AI analysis could not proceed:\n\n{ai_text}{err_hint}"))
                    self.root.after(0, lambda: self.btn_ai.config(state='normal', text='Run AI Analysis'))
                    return

                stats_summary = {
                    "Median_Length_um": f"{ts['total_3d_length_um'].median():.2f}",
                    "Median_Z_Span_um": f"{ts['z_span_um'].median():.2f}" if "z_span_um" in ts.columns else "NA",
                    "Median_Effective_Thickness_um_PSF_sensitive": f"{ts['thickness_um'].median():.2f}" if "thickness_um" in ts.columns else "NA",
                    "Median_Taper_Ratio_PSF_sensitive": f"{ts['taper_ratio'].median():.2f}" if "taper_ratio" in ts.columns else "NA",
                    "Avg_Tortuosity": f"{ts['tortuosity_3d'].mean():.2f}",
                    "Track_Count": f"{len(ts)}",
                    "Species": species
                }
                report_path = generate_ai_html_report(out_dir, ai_text, stats_summary, species)

                if os.path.exists(report_path):
                    print(f"AI SUCCESS: Report generated at {report_path}")
                    abs_report = os.path.abspath(report_path)

                    def _open_report():
                        try:
                            if os.name == 'nt':
                                os.startfile(abs_report)
                            else:
                                webbrowser.open("file:///" + abs_report.replace("\\", "/"))
                        except Exception as oe:
                            print(f"AI OPEN ERROR: {oe}")

                    self.root.after(0, lambda: self.lbl_batch_op.config(text='AI REPORT READY OK', fg='green'))
                    self.root.after(500, _open_report)
                    self.root.after(1200, lambda: messagebox.showinfo("AI Analysis Complete",
                        f"Biological interpretation report generated and opened:\n\n{os.path.basename(report_path)}"))
                else:
                    raise FileNotFoundError(f"Report file was not created at {report_path}")

            except Exception as e:
                import traceback
                print(f"AI EXCEPTION: {traceback.format_exc()}")
                err_msg = f"Failed during biological interpretation:\n\n{str(e)}"
                self.root.after(0, lambda m=err_msg: messagebox.showerror("AI Error", m))
            finally:
                self.root.after(0, lambda: self.btn_ai.config(state='normal', text='Run AI Analysis'))

        threading.Thread(target=_ai_thread, daemon=True).start()

    def load_directory(self):
        """
        Opens a file-picker dialog so the user can select any image in a Z-stack folder.

        Discovers all supported image files (``*.tif``, ``*.tiff``, ``*.png``, ``*.jpg``,
        ``*.jpeg``) in the same directory as the selected file, applies natural sort order
        (so ``z01`` comes before ``z10``), synchronises the Z-slice slider to the selected
        file, and calls :meth:`load_image` to display the first slice.
        """
        initial = CONFIG.get('INPUT_DIR', os.getcwd())

        # Use file picker so user can see images
        fpath = filedialog.askopenfilename(
            initialdir=str(initial),
            title="Select any image in the stack",
            filetypes=[("Image files", "*.tif *.tiff *.png *.jpg *.jpeg"), ("All files", "*.*")]
        )

        if not fpath:
            return

        selected_file = pl.Path(fpath)
        p = selected_file.parent
        self.input_dir = str(p)

        exts = ['.tif', '.tiff', '.png', '.jpg', '.jpeg']
        found_files_path = []

        # 1. Discover all images in the same folder
        for ext in exts:
            found_files_path.extend(list(p.glob(f"*{ext}")))
            found_files_path.extend(list(p.glob(f"*{ext.upper()}")))

        # 2. Recursive fallback (if needed, though file picker implies they are in the right spot)
        if not found_files_path:
            for ext in exts:
                found_files_path.extend(list(p.rglob(f"*{ext}")))
                found_files_path.extend(list(p.rglob(f"*{ext.upper()}")))

        if not found_files_path:
            messagebox.showerror('Error', f"No supported images found in:\n{p}")
            return

        # Standardize and Natural Sort
        found_files_str = [os.path.abspath(str(f)) for f in found_files_path]
        unique_files = list(set(found_files_str))
        self.files = sorted(unique_files, key=natural_sort_key)

        # Sync to the selected file
        try:
            self.current_idx = self.files.index(os.path.abspath(fpath))
        except ValueError:
            self.current_idx = 0

        self.scale_z.config(to=len(self.files) - 1)
        self.scale_z.set(self.current_idx)
        self.reset_roi(redraw=False)
        self.load_image()
        self.lbl_status.config(text=f'Opened: {selected_file.name}\n({len(self.files)} slices in folder)', fg='blue')
        self.root.update()

    def load_image(self):
        if not self.files:
            return
        try:
            self.current_img = robust_imread(self.files[self.current_idx])
            self.lbl_z.config(text=f'Z: {self.current_idx} / {len(self.files)-1}')
            self.render()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load image:\n{self.files[self.current_idx]}\n\nError: {e}")

    def on_slide_change(self, val):
        self.current_idx = int(val)
        self.load_image()

    def render(self):
        self.ax.clear()
        self.ax.axis('off')

        # --- NEW: Review Overlays Logic ---
        if self.mode_var.get() == 'review':
            try:
                if not hasattr(self, 'last_out_dir') or not self.last_out_dir:
                    self.ax.text(0.5, 0.5, "No Batch Analysis Results Found.\nRun Batch First.",
                                 ha='center', va='center', color='red', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return

                z_idx = extract_z_index(self.files[self.current_idx])
                quality_panel_path = os.path.join(
                    self.last_out_dir, "quality_overlays", f"z{z_idx:02d}_quality_panel.png")
                raw_panel_path = os.path.join(
                    self.last_out_dir, "overlays", f"z{z_idx:02d}_panel.png")
                panel_path = (
                    quality_panel_path
                    if os.path.exists(quality_panel_path)
                    else raw_panel_path
                )

                if os.path.exists(panel_path):
                    if _HAVE_CV2:
                        img = _cv2.imread(panel_path)
                        img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
                    else:
                        img = plt.imread(panel_path)
                    self.ax.imshow(img)
                    if panel_path == quality_panel_path:
                        active_statuses = {"candidate", "warning", "hard_fail"}
                        quality_counts_path = os.path.join(
                            self.last_out_dir,
                            "quality_overlays",
                            "quality_overlay_counts.csv",
                        )
                        if os.path.exists(quality_counts_path):
                            quality_counts = pd.read_csv(quality_counts_path)
                            status_row = (
                                quality_counts[
                                    pd.to_numeric(
                                        quality_counts["z_slice"],
                                        errors="coerce",
                                    ).eq(z_idx)
                                ]
                                if "z_slice" in quality_counts.columns
                                else pd.DataFrame()
                            )
                            if not status_row.empty:
                                active_statuses = {
                                    status
                                    for status in _QUALITY_OVERLAY_LEGEND
                                    if status in status_row.columns
                                    and int(status_row.iloc[0][status]) > 0
                                }
                        self.ax.set_title(
                            "Track-QC overlay (green/amber are included)",
                            fontsize=9,
                        )
                        self.ax.legend(
                            handles=quality_overlay_legend_handles(active_statuses),
                            loc="lower center",
                            bbox_to_anchor=(0.5, -0.12),
                            ncol=2,
                            fontsize=7,
                            frameon=True,
                        )
                    else:
                        self.ax.set_title(
                            "Pre-tracking candidates; colors identify candidate IDs",
                            fontsize=9,
                        )
                    self.canvas.draw_idle()
                    return
                else:
                    self.ax.text(0.5, 0.5, f"Overlay not found for Z={z_idx:02d}\n{os.path.basename(panel_path)}",
                                 ha='center', va='center', color='orange', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return
            except Exception as e:
                print(f"Overlay Render Error: {e}")

        if self.current_img is not None and isinstance(self.current_img, np.ndarray):
            img = self.current_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            self.ax.imshow(disp, cmap='gray')
            if len(self.roi_points) > 0:
                pts = np.array(self.roi_points)
                # If active (closed), draw it closed. Otherwise, draw the open line and endpoints.
                if self.roi_active:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=2)
                else:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=1.5)
                    self.ax.plot(pts[:,0], pts[:,1], 'ro', markersize=4)
            # If a loaded mask exists, always redraw the red contour so it persists across slices
            elif self._loaded_roi_mask is not None:
                self.ax.contour(self._loaded_roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def on_click(self, event):
        if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
            return
        if self.mode_var.get() == 'roi':
            # Left click = add point
            if event.button == 1:
                # If they start clicking after already having a finished ROI, wipe it to start fresh
                if self.roi_active:
                    self.reset_roi(redraw=False)
                self.roi_points.append([event.xdata, event.ydata])
                self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                self.render()

            # Right click = undo last point
            elif event.button == 3:
                if not self.roi_active and len(self.roi_points) > 0:
                    self.roi_points.pop()
                    self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                    self.render()

    def finalize_roi(self):
        if not self.roi_active and len(self.roi_points) > 2:
            self.roi_points.append(self.roi_points[0])
            self.roi_active = True
            self.lbl_roi.config(text=f'ROI: active ({len(self.roi_points)-1} points)')
            self.render()
        elif not self.roi_active and len(self.roi_points) > 0:
            messagebox.showwarning('Draw ROI', 'Please place at least 3 points before finalizing.')

    def on_key(self, event):
        if self.mode_var.get() == 'roi' and event.key == 'enter':
            self.finalize_roi()

    def reset_roi(self, redraw=True):
        self.roi_points = []
        self.roi_active = False
        self._loaded_roi_mask = None
        self.lbl_roi.config(text='ROI: none')
        if redraw:
            self.render()

    def build_roi_mask(self):
        # If a mask was loaded from file, use it directly
        if self._loaded_roi_mask is not None and self.current_img is not None:
            return self._loaded_roi_mask
        if not self.roi_active or len(self.roi_points) < 4 or self.current_img is None:
            return None
        h, w = self.current_img.shape
        yy, xx = np.mgrid[:h, :w]
        pts = np.column_stack((xx.ravel(), yy.ravel()))
        path = Path(self.roi_points)
        return path.contains_points(pts).reshape(h, w)

    def save_roi_mask(self):
        mask = self.build_roi_mask()
        if mask is None:
            messagebox.showinfo('Save ROI', 'No active ROI to save.')
            return
        default = f'roi_z{self.current_idx:02d}.npy'
        path = filedialog.asksaveasfilename(defaultextension='.npy', initialfile=default,
                                            filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        np.save(path, mask.astype(np.uint8))
        messagebox.showinfo('Save ROI', f'Saved ROI mask to:\n{path}')

    def load_roi_mask(self):
        if self.current_img is None:
            messagebox.showinfo('Load ROI', 'Load an image first.')
            return
        path = filedialog.askopenfilename(filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        mask = np.load(path).astype(bool)
        if mask.shape != self.current_img.shape:
            messagebox.showerror('Load ROI', f'Mask shape {mask.shape} does not match image shape {self.current_img.shape}')
            return
        # Convert mask boundary to points for display
        ys, xs = np.where(mask)
        if len(xs) == 0:
            messagebox.showerror('Load ROI', 'Loaded ROI mask is empty.')
            return
        self.roi_points = []
        self.roi_active = True
        self._loaded_roi_mask = mask  # Store the loaded mask for build_roi_mask!
        self.lbl_roi.config(text=f'ROI: loaded mask\n{os.path.basename(path)}')
        self.render()
        self.ax.contour(mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def run_analysis_slice(self):
        if self.current_img is None:
            messagebox.showinfo('Info', 'No image loaded. Load a directory first.')
            return

        self.lbl_roi.config(text='Running analysis...')
        self.root.update_idletasks()

        try:
            import time as _t
            log_lines = []
            def log(msg):
                log_lines.append(msg)
                print(msg)
                sys.stdout.flush()

            log(f"\n--- GUI Analysis: slice {self.current_idx} ---")
            log(f"  File: {os.path.basename(self.files[self.current_idx]) if self.files else 'N/A'}")
            log(f"  Image shape: {self.current_img.shape}, dtype: {self.current_img.dtype}")

            params = CONFIG.copy()
            params['SAVE_DEBUG_IMAGES'] = False
            roi_mask = self.build_roi_mask()
            preview_z = extract_z_index(
                self.files[self.current_idx],
                sequence_idx=self.current_idx,
            )

            full_img = self.current_img
            crop_offset_y, crop_offset_x = 0, 0
            crop_roi = roi_mask

            t0 = _t.time()
            log("  v5.7 U-Net-ready single-pass analysis...")
            preview_context = build_stack_preprocess_context(
                self.files if self.files else [self.files[self.current_idx]],
                roi_mask,
                params,
                exclusion_mask=None,
            )
            log(f"  Temporary preview context: profile={preview_context.selected_clahe_profile}, sampled_z={preview_context.sampled_z_indices}")
            seg1 = segment_slice(
                full_img,
                params,
                roi_mask=roi_mask,
                preprocess_context=preview_context,
                exclusion_mask=None,
                z_idx=preview_z,
            )
            meas1 = measure_spermatids(seg1, params)
            results = meas1['results']
            skel_label_full = meas1['skel_label']

            elapsed = _t.time() - t0
            log(f"  RESULT: {len(results)} 2D preview candidates ({elapsed:.1f}s)")

            preview_output_dir = params["OUTPUT_DIR"]
            ensure_dir(preview_output_dir)
            preview_frame = pd.DataFrame(
                rows_from_results(results, preview_z, params["UM_PER_PX_XY"])
            )
            preview_frame.to_csv(
                os.path.join(
                    preview_output_dir,
                    f"single_measurements_z{preview_z:03d}_{_VERSION}.csv",
                ),
                index=False,
            )
            export_analysis_summary(
                preview_output_dir,
                df=preview_frame,
                run_scope="single_slice_preview",
                z_index=preview_z,
                cfg=params,
            )

            overlay = make_overlay(full_img, skel_label_full)

            # Write log to file
            try:
                log_path = os.path.join(CONFIG['OUTPUT_DIR'], 'gui_analysis_log.txt')
                ensure_dir(CONFIG['OUTPUT_DIR'])
                with open(log_path, 'a', encoding='utf-8') as f:
                    f.write('\n'.join(log_lines) + '\n\n')
            except Exception:
                pass

            # Show results popup
            top = tk.Toplevel(self.root)
            top.title(f'2D Preview Z={preview_z} - {len(results)} candidates')
            top.geometry('1200x650')

            fig = Figure(figsize=(14, 6))
            ax1 = fig.add_subplot(121)
            ax2 = fig.add_subplot(122)

            img = full_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            ax1.imshow(disp, cmap='gray')
            if roi_mask is not None:
                ax1.contour(roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.2)
            ax1.set_title('Original + ROI')
            ax1.axis('off')

            ax2.imshow(overlay)
            ax2.set_title(f'2D candidate overlay (N={len(results)})')
            ax2.axis('off')

            um = params['UM_PER_PX_XY']
            for r in results:
                ax2.text(r['centroid_x'], r['centroid_y'],
                         f"{r['length_px_geodesic'] * um:.1f}",
                         color='white', fontsize=5, ha='center', va='center')

            fig.tight_layout()
            can = FigureCanvasTkAgg(fig, master=top)
            can.get_tk_widget().pack(fill='both', expand=True)
            can.draw()

            if results:
                lengths = [r['length_px_geodesic'] * um for r in results]
                text = (
                    f'2D preview candidates: {len(results)} | median 2D length '
                    f'{np.median(lengths):.2f} um ({elapsed:.1f}s)\n'
                    'Not a unique-nucleus count; run the complete stack for 3D results.'
                )
            else:
                text = (
                    f'2D preview candidates: 0 ({elapsed:.1f}s) - '
                    'see gui_analysis_log.txt for diagnostics'
                )
            lbl_stats = tk.Label(top, text=text, font=('Arial', 11))
            lbl_stats.pack(pady=4)

            lbl_tool = tk.Label(top, text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue', font=('Arial', 10, 'bold'))
            lbl_tool.pack(pady=2)

            self.lbl_roi.config(text=f'Preview done: {len(results)} 2D candidates')

            # ------ INTERACTIVE MANUAL CORRECTION LOGIC ------
            class ManualCorrector:
                def __init__(self, canvas, ax_overlay, seg_data, prms, crop_oy, crop_ox, fimg):
                    self.canvas = canvas
                    self.ax = ax_overlay
                    self.seg = seg_data
                    self.params = prms
                    self.crop_oy = crop_oy
                    self.crop_ox = crop_ox
                    self.fimg = fimg

                    self.active_tool = None
                    self.cid_press = self.canvas.mpl_connect('button_press_event', self.on_click)
                    self.cid_key = self.canvas.mpl_connect('key_press_event', self.on_key)

                    self.overlay_imshow = None
                    self.text_artists = []

                def on_key(self, event):
                    if event.key == 'e':
                        self.active_tool = 'erase'
                        lbl_tool.config(text="Active Tool: ERASE (Click a colored spermatid to delete it)", fg='red')
                    elif event.key == 's':
                        self.active_tool = 'split'
                        lbl_tool.config(text="Active Tool: SPLIT (Click on a skeleton branch to sever it)", fg='orange')
                    elif event.key == 'escape':
                        self.active_tool = None
                        lbl_tool.config(text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue')

                def on_click(self, event):
                    if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
                        return
                    if not self.active_tool:
                        return

                    # Map global click to cropped coordinate space
                    x_crop = int(round(event.xdata)) - self.crop_ox
                    y_crop = int(round(event.ydata)) - self.crop_oy

                    ch, cw = self.seg['skel_pruned'].shape
                    if not (0 <= x_crop < cw and 0 <= y_crop < ch):
                        return # Clicked outside the ROI working area

                    modified = False
                    if self.active_tool == 'erase':
                        # Find the label at this pixel in the current skeleton
                        lab = self.seg['skel_labeled'][y_crop, x_crop]
                        # If the user clicked slightly off, search a 3x3 neighborhood
                        if lab == 0:
                            for dy in (-1, 0, 1):
                                for dx in (-1, 0, 1):
                                    yy, xx = y_crop+dy, x_crop+dx
                                    if 0 <= yy < ch and 0 <= xx < cw and self.seg['skel_labeled'][yy, xx] > 0:
                                        lab = self.seg['skel_labeled'][yy, xx]
                                        break
                        if lab > 0:
                            self.seg['skel_pruned'][self.seg['skel_labeled'] == lab] = False
                            modified = True

                    elif self.active_tool == 'split':
                        # Draw a small black circle to sever the topological skeleton connection
                        from skimage.draw import disk
                        rr, cc = disk((y_crop, x_crop), radius=2.5, shape=self.seg['skel_pruned'].shape)
                        self.seg['skel_pruned'][rr, cc] = False
                        modified = True

                    if modified:
                        self.recalculate_and_redraw()

                def recalculate_and_redraw(self):
                    # Re-label the modified skeleton
                    self.seg['skel_labeled'] = measure.label(self.seg['skel_pruned'])

                    # Re-run measurement filters
                    new_meas = measure_spermatids(self.seg, self.params)
                    new_results = new_meas['results']

                    # Filter by ROI inside centroid logic
                    if crop_roi is not None:
                        filtered = []
                        keep_labels = []
                        for r in new_results:
                            c_y = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0] - 1)
                            c_x = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1] - 1)
                            if crop_roi[c_y, c_x]:
                                filtered.append(r)
                                keep_labels.append(r['label'])
                        new_results = filtered
                        skel_l = np.where(np.isin(new_meas['skel_label'], keep_labels), new_meas['skel_label'], 0).astype(np.int32)
                    else:
                        skel_l = new_meas['skel_label']

                    # Map back to full image
                    H, W = self.fimg.shape
                    skel_label_full = np.zeros((H, W), dtype=np.int32)
                    ch, cw = skel_l.shape
                    skel_label_full[self.crop_oy:self.crop_oy+ch, self.crop_ox:self.crop_ox+cw] = skel_l

                    # Redraw Overlay
                    new_overlay = make_overlay(self.fimg, skel_label_full)

                    self.ax.clear()
                    self.ax.imshow(new_overlay)
                    self.ax.set_title(f'2D candidate overlay (N={len(new_results)}) - Manual Corrections Applied')
                    self.ax.axis('off')

                    # Redraw Text
                    _um = self.params['UM_PER_PX_XY']
                    for r in new_results:
                        self.ax.text(r['centroid_x'] + self.crop_ox, r['centroid_y'] + self.crop_oy,
                                 f"{r['length_px_geodesic'] * _um:.1f}",
                                 color='white', fontsize=5, ha='center', va='center')

                    self.canvas.draw()

                    # Update Stats Label
                    if new_results:
                        lengths = [r['length_px_geodesic'] * _um for r in new_results]
                        lbl_stats.config(
                            text=(
                                f'Corrected 2D candidates: {len(new_results)} | '
                                f'median 2D length {np.median(lengths):.2f} um | '
                                'not a unique-nucleus count'
                            )
                        )
                    else:
                        lbl_stats.config(text='Corrected 2D candidates: 0')

                    corrected_frame = pd.DataFrame(
                        rows_from_results(new_results, preview_z, _um)
                    )
                    corrected_frame.to_csv(
                        os.path.join(
                            preview_output_dir,
                            f"single_measurements_z{preview_z:03d}_{_VERSION}.csv",
                        ),
                        index=False,
                    )
                    export_analysis_summary(
                        preview_output_dir,
                        df=corrected_frame,
                        run_scope="single_slice_preview",
                        z_index=preview_z,
                        cfg=params,
                    )

            # Attach to popup so it doesn't get garbage collected
            top.corrector = ManualCorrector(can, ax2, seg1, params, crop_offset_y, crop_offset_x, full_img)

        except Exception as e:
            import traceback
            traceback.print_exc()
            sys.stdout.flush()
            self.lbl_roi.config(text=f'Analysis error: {e}')
            messagebox.showerror('Analysis Error', f'{type(e).__name__}: {e}')

    def run_batch_analysis(self):
        if not self.files:
            messagebox.showinfo('Info', 'No directory loaded.')
            return

        # Auto-incremental output directory inside selected folder
        out_dir = get_unique_batch_dir(self.input_dir)
        self.last_out_dir = out_dir

        # EXPLICIT CONFIRMATION: Show the user where the data will go
        confirm = messagebox.askokcancel("Confirm Output",
            f"Results (Excel, PDF, CSV) will be saved to:\n\n{out_dir}\n\nContinue?")
        if not confirm:
            return

        ensure_dir(out_dir)
        overlay_dir = os.path.join(out_dir, "overlays")
        rescue_review_dir = os.path.join(
            out_dir,
            "technical_qc",
            "unet_rescue_overlays",
        )
        ensure_dir(overlay_dir)
        ensure_dir(rescue_review_dir)

        params = CONFIG.copy()
        params['OUTPUT_DIR'] = out_dir
        params['SAVE_DEBUG_IMAGES'] = False
        params['DO_TRACKING'] = True

        roi_mask = self.build_roi_mask()
        exclusion_mask = None

        self.lbl_roi.config(text="Processing... See Top Bar")
        self.root.update_idletasks()

        try:
            import time as _t
            t_batch = _t.time()
            self.lbl_batch_op.config(text=f"Batch Segmenting: 0 / {len(self.files)} slices...", fg='blue')
            self.root.update()

            all_rows = []
            summaries = []

            # Z-Projection accumulation
            max_proj_raw = None
            max_proj_ov = None
            slice_cache = {}

            # Robust initialization
            ts = None
            df_trk = None
            first_img = ensure_2d_image(robust_imread(self.files[0]), os.path.basename(self.files[0]))
            if roi_mask is not None and roi_mask.shape != first_img.shape:
                raise ValueError(f"ROI shape {roi_mask.shape} does not match image shape {first_img.shape}")
            files_by_z = {
                int(extract_z_index(fpath, sequence_idx=i)): fpath
                for i, fpath in enumerate(self.files)
            }
            preprocess_context = build_stack_preprocess_context(self.files, roi_mask, params, exclusion_mask=exclusion_mask)
            save_stack_preprocess_context(preprocess_context, out_dir)
            if roi_mask is not None:
                tifffile.imwrite(os.path.join(out_dir, "roi_mask_used.tif"), roi_mask.astype(np.uint8) * 255)

            self.progress['value'] = 0
            self.progress['maximum'] = len(self.files)
            self.progress_post['value'] = 0
            self.progress_post['maximum'] = 100
            self.lbl_post_progress_val.config(text="0%", fg='orange')

            for idx, fpath in enumerate(self.files):
                z_idx = extract_z_index(fpath, sequence_idx=idx)

                pct = int(((idx + 1) / len(self.files)) * 100)
                self.progress['value'] = idx + 1
                self.lbl_progress_val.config(text=f"{pct}%")
                self.root.update()

                print(f"[{idx+1}/{len(self.files)}] Processing Z={z_idx:02d}...")

                img_raw = robust_imread(fpath)
                process_img = ensure_2d_image(img_raw, os.path.basename(fpath))
                print(f"DEBUG batch image {os.path.basename(fpath)} shape: {process_img.shape}, ndim={process_img.ndim}")

                full_img = process_img
                crop_oy, crop_ox = 0, 0
                unet_context = _make_unet_context_from_paths(files_by_z, z_idx)
                seg = segment_slice(process_img, params, z_idx=z_idx,
                                    roi_mask=roi_mask,
                                    preprocess_context=preprocess_context,
                                    exclusion_mask=exclusion_mask,
                                    unet_context_stack=unet_context)
                meas = measure_spermatids(seg, params)
                res = meas['results']
                sl_full = meas['skel_label']

                for r in res:
                    r['centroid_x'] += crop_ox
                    r['centroid_y'] += crop_oy

                um = params['UM_PER_PX_XY']
                ls_um = [r['length_px_geodesic']*um for r in res]
                ws_um = [r['width_px']*um for r in res]

                all_rows.extend(rows_from_results(res, z_idx, um))
                summaries.append({
                    "z_slice": z_idx,
                    "n_spermatids": len(res),
                    "mean_length_um": round(float(np.mean(ls_um)), 3) if ls_um else 0,
                    "median_length_um": round(float(np.median(ls_um)), 3) if ls_um else 0,
                    "mean_width_um": round(float(np.mean(ws_um)), 3) if ws_um else 0,
                })

                if params['SAVE_OVERLAYS']:
                    ov = make_overlay(full_img, sl_full)
                    # Create side-by-side panel
                    orig_rgb = (normalize_display(full_img) * 255).astype(np.uint8)
                    if orig_rgb.ndim == 2:
                        orig_rgb = np.stack([orig_rgb]*3, axis=-1)
                    panel = np.hstack([orig_rgb, ov])
                    _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_panel.png"), panel)
                    rescue_rgb = make_unet_rescue_review_overlay(
                        full_img,
                        sl_full,
                        res,
                        meas.get("unet_rescue_rejected_reason"),
                    )
                    _imwrite(
                        os.path.join(
                            rescue_review_dir,
                            f"z{z_idx:02d}_unet_rescue_review.png",
                        ),
                        rescue_rgb,
                    )

                    # ---- LIVE GUI UPDATE ----
                    # Show the side-by-side segmentation panel during batch execution
                    if hasattr(self, 'ax') and hasattr(self, 'canvas'):
                        try:
                            self.ax.clear()
                            self.ax.axis('off')
                            self.ax.imshow(panel)
                            self.canvas.draw()
                        except Exception:
                            pass

                    # Update Z-Projections
                    if max_proj_raw is None:
                        max_proj_raw = full_img.copy().astype(np.float32)
                        max_proj_ov = ov.copy().astype(np.float32)
                    else:
                        max_proj_raw = np.maximum(max_proj_raw, full_img.astype(np.float32))
                        max_proj_ov = np.maximum(max_proj_ov, ov.astype(np.float32))
                    slice_cache[int(z_idx)] = {
                        "image": full_img.copy(),
                        "skel_label": sl_full.copy().astype(np.int32),
                    }

                if params["SAVE_MASK_TIFS"]:
                    tifffile.imwrite(os.path.join(out_dir, f"z{z_idx:02d}_mask.tif"),
                                     (seg["mask_clean"] & roi_mask if roi_mask is not None else seg["mask_clean"]).astype(np.uint8) * 255)
                if params.get("UNET_SAVE_PROBABILITY_MAPS", True):
                    if seg.get("unet_probability") is not None and np.any(seg.get("unet_probability")):
                        tifffile.imwrite(os.path.join(out_dir, f"z{z_idx:02d}_unet_probability.tif"),
                                         seg["unet_probability"].astype(np.float32))

                # Update Progress Bar for each slice
                self.lbl_batch_op.config(text=f"Batch Segmenting: {idx+1} / {len(self.files)} slices...", fg='blue')
                self.progress['value'] = idx + 1
                self.root.update()

            df = pd.DataFrame(all_rows)
            df_sum = pd.DataFrame(summaries)
            df.to_csv(os.path.join(out_dir, "spermatid_measurements.csv"), index=False)
            df_sum.to_csv(os.path.join(out_dir, "slice_summary.csv"), index=False)

            if not df.empty:
                self.lbl_batch_op.config(text='Running 3D Tracking & Morphometrics...', fg='#e67e22') # Orange-ish

                # --- Sequential Progress Bar Swap ---
                # Hide 2D progress, Show Post-Analysis progress
                self.batch_p_frame.pack_forget()
                self.post_p_frame.pack(fill='x')

                self.progress_post['value'] = 25
                self.lbl_post_progress_val.config(text="25%")
                self.root.update()

                df_trk, ts = track_across_slices(df, params)

                # --- Advanced 3D Morphometrics ---
                self.lbl_batch_op.config(text='Calculating Advanced 3D Metrics...', fg='#e67e22')
                # Metrics are natively generated correctly and safely in track_across_slices

                self.progress_post['value'] = 60
                self.lbl_post_progress_val.config(text="60%")
                self.root.update()

                df_trk.to_csv(os.path.join(out_dir, "measurements_with_tracks.csv"), index=False)

                # ------ AUTO CANDIDATE AUDIT ----------------------------------------------------------------------------------------------------------------
                self.lbl_batch_op.config(text='Running Candidate Audit...', fg='#e67e22')
                self.root.update()

                ts = flag_quality_tracks(ts, params)

                # Save annotated track summary with explicit population flags.
                ts.to_csv(os.path.join(out_dir, "track_summary.csv"), index=False)

                export_comparative_track_tables(out_dir, ts, None)
                export_biologist_results(out_dir, ts, None)

                # Generate candidate-coded overlays after audit.
                if params['SAVE_OVERLAYS']:
                    export_quality_overlays(out_dir, slice_cache, df_trk, ts)
                    export_analysis_overlays(
                        out_dir,
                        slice_cache,
                        df_trk,
                        ts,
                    )

                # Generate outlier_audit/ subfolder automatically
                export_outlier_audit(out_dir, ts, params)
                export_post_detection_qc(out_dir, df_trk, ts)

                n_candidates = int(ts["technical_valid"].sum()) if "technical_valid" in ts.columns else len(ts)
                primary_tracks = _technical_valid_track_population(ts)
                median_3d_length = (
                    float(primary_tracks["total_3d_length_um"].median())
                    if not primary_tracks.empty and "total_3d_length_um" in primary_tracks.columns
                    else np.nan
                )
                self.lbl_batch_op.config(
                    text=(
                        f'Primary result: {n_candidates} estimated unique nuclei | '
                        f'median 3D length {median_3d_length:.2f} um'
                    ),
                    fg='#27ae60')
                self.root.update()

                # Save Global Z-Projection
                if max_proj_raw is not None:
                    self.lbl_batch_op.config(text='Generating Global Z-Projection...', fg='#e67e22')
                    raw_p = (normalize_display(max_proj_raw.astype(np.uint16)) * 255).astype(np.uint8)
                    if raw_p.ndim == 2: raw_p = np.stack([raw_p]*3, axis=-1)
                    ov_p = max_proj_ov.astype(np.uint8)
                    global_panel = np.hstack([raw_p, ov_p])
                    _imwrite(os.path.join(out_dir, "global_z_projection.png"), global_panel)

                self.progress_post['value'] = 80
                self.lbl_post_progress_val.config(text="80%")
                self.root.update()

            analysis_summary = export_analysis_summary(
                out_dir,
                df=df,
                track_summary=ts,
                run_scope="full_stack_3d",
                cfg=params,
            )

            elapsed = _t.time() - t_batch
            if analysis_summary["biological_count_available"]:
                msg = (
                    f"Batch complete in {elapsed:.1f}s.\n\n"
                    f"Estimated unique nuclei: {analysis_summary['estimated_unique_nuclei']}\n"
                    f"Median 3D length: {analysis_summary['median_3d_length_um']:.2f} um\n\n"
                    f"Primary sample summary:\n"
                    f"{os.path.join(out_dir, 'biologist_results', 'sample_summary.csv')}\n\n"
                    f"Saved to:\n{out_dir}"
                )
            else:
                msg = (
                    f"Batch complete in {elapsed:.1f}s, but no 3D biological "
                    f"population was produced.\n\nSaved to:\n{out_dir}"
                )
            print(msg)
            self.lbl_batch_op.config(text='Batch Analysis Complete.', fg='green')
            self.lbl_progress_val.config(text="100% - Done", fg='green')
            self.root.update()

            # Generate High-Res Graphical Report and Excel Audit
            self.lbl_batch_op.config(text='Generating PDF & Excel Reports...', fg='#8e44ad') # Purple

            def update_cb(v):
                self.progress_post['value'] = v
                self.lbl_post_progress_val.config(text=f"{v}%")
                self.root.update()

            generate_batch_report(
                out_dir,
                df,
                df_sum,
                um,
                ts if not df.empty else None,
                update_cb,
                generate_pptx=False,
                df_tracked=df_trk,
            )
            generate_excel_report(out_dir, df, df_sum, ts if not df.empty else None)

            # --- Store batch data for AI button ---
            primary_for_ai = _technical_valid_track_population(ts)
            if not primary_for_ai.empty:
                self._last_batch_ts = primary_for_ai
                self._last_batch_out_dir = out_dir
                self.btn_ai.config(state='normal')
                print(
                    "AI READY: "
                    f"{len(self._last_batch_ts)} technical-valid tracks stored. "
                    "Click 'Run AI Analysis' to interpret."
                )
            else:
                self._last_batch_ts = None
                self._last_batch_out_dir = None
                self.btn_ai.config(state='disabled')

            self.lbl_batch_op.config(text='Generating PowerPoint Dashboard...', fg='#c0392b') # Red-ish
            self.root.update()
            try:
                ok = generate_pptx_report(out_dir, df, df_sum, um, ts if not df.empty else None)
                if not ok:
                    print("WARNING: PPTX generation returned False - check console for traceback.")
            except Exception as pptx_err:
                import traceback as _tb
                print(f"ERROR generating PPTX: {pptx_err}")
                _tb.print_exc()

            print(msg)
            self.lbl_batch_op.config(text='ALL OPERATIONS COMPLETE', fg='green')
            self.lbl_progress_val.config(text="100%", fg='green')
            self.lbl_post_progress_val.config(text="100% - DONE", fg='green')
            self.progress_post['value'] = 100
            self.root.update()

            # --- macOS STABILITY PATCH: Ensure all plot resources are freed before UI popup ---
            try:
                import matplotlib.pyplot as plt
                plt.close('all')
            except: pass

            messagebox.showinfo('Batch Complete', msg)

        except Exception as e:
            import traceback
            traceback.print_exc()
            self.lbl_batch_op.config(text=f'Batch error: {e}', fg='red')
            messagebox.showerror('Batch Error', str(e))



def launch_gui():
    if not _TK_AVAILABLE:
        raise RuntimeError("Tkinter GUI components are not available in this Python environment.")
    root = tk.Tk()
    app = SpermGUI(root)
    root.mainloop()


# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    import multiprocessing
    multiprocessing.freeze_support()
    print(f"\n{'='*60}")
    print(f" SPERMATID ANALYSIS PIPELINE - v{_VERSION} - INITIALIZING...")
    print(f" Mode: {'GUI' if '--gui' in sys.argv or len(sys.argv)<=1 else 'CLI'}")
    print(f" Features: [Incremental Folders] [Multi-Tab Excel] [3D Stats]")
    print(f"{'='*60}\n")

    ap = argparse.ArgumentParser(description=f"Spermatid segmentation {_VERSION} / ROI GUI")
    ap.add_argument("--batch",  action="store_true", help="Run batch processing")
    ap.add_argument("--single", action="store_true", help="Run a single-image analysis")
    ap.add_argument("--gui",    action="store_true", help="Launch ROI GUI")
    ap.add_argument("--z",      type=int, default=None, help="Choose z-index in single mode")
    ap.add_argument("--params", type=str, default=None, help="Path to tuned parameters JSON file to override CONFIG")
    ap.add_argument("--roi-mask", type=str, default=None, help="Optional .npy/.tif ROI mask for batch processing")
    args = ap.parse_args()

    # Load tuned parameters from JSON if provided
    if args.params:
        import json as _json
        params_path = os.path.abspath(args.params)
        if os.path.exists(params_path):
            with open(params_path, 'r') as _pf:
                tuned = _json.load(_pf)
            applied = 0
            for key, value in tuned.items():
                if key in CONFIG:
                    CONFIG[key] = value
                    applied += 1
            print(f"  Loaded {applied} tuned parameters from: {os.path.basename(params_path)}")
        else:
            print(f"  WARNING: Params file not found: {params_path}")

    if args.roi_mask:
        CONFIG["ROI_MASK_PATH"] = os.path.abspath(args.roi_mask)
        print(f"  Loaded ROI mask path: {CONFIG['ROI_MASK_PATH']}")

    validate_config(CONFIG)

    # Launch GUI by default if no explicit CLI flags are provided
    if args.gui or not (args.batch or args.single or args.z is not None):
        launch_gui()
        raise SystemExit

    if args.batch:
        CONFIG["RUN_MODE"] = "batch"
        # CLI Incremental Folder Logic - Anchor to INPUT_DIR
        base_parent = CONFIG["INPUT_DIR"]
        if not os.path.isabs(base_parent):
            base_parent = os.path.abspath(base_parent)

        CONFIG["OUTPUT_DIR"] = get_unique_batch_dir(base_parent)
        ensure_dir(CONFIG["OUTPUT_DIR"])
        print(f"CLI BATCH MODE: Results will be saved inside input folder: {CONFIG['OUTPUT_DIR']}")
    if args.single:
        CONFIG["RUN_MODE"] = "single"
    if args.z is not None:
        CONFIG["SINGLE_IMAGE_SELECTION_MODE"] = "z_index"
        CONFIG["SINGLE_Z_INDEX"] = args.z

    if CONFIG["RUN_MODE"] == "single":
        process_one_image(choose_single_image(CONFIG), CONFIG, CONFIG["OUTPUT_DIR"])
    else:
        process_batch(CONFIG)
