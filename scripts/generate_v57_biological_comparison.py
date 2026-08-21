"""Generate a specimen-level biological comparison package from a v5.7 study."""

import argparse
import json
import math
import shutil
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use("Agg", force=True)
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.backends.backend_pdf import PdfPages
from scipy.stats import mannwhitneyu, permutation_test, spearmanr, ttest_ind


METRICS = {
    "estimated_nuclei_per_1000_um2": {
        "label": "Estimated nuclei per 1,000 um2",
        "short": "Area density",
        "question": "How many nuclei are present per unit of tissue area?",
        "meaning": (
            "An XY-area density. Higher values indicate more reconstructed nuclei "
            "inside the ROI footprint, but this measure does not correct for stack depth."
        ),
        "role": "count_exploratory",
    },
    "estimated_nuclei_per_100000_um3": {
        "label": "Estimated nuclei per 100,000 um3",
        "short": "Volume density",
        "question": "How many nuclei are present per unit of sampled 3D volume?",
        "meaning": (
            "A volume-normalized density. It adjusts for nominal sampled depth but "
            "can remain sensitive to stack boundaries and cross-slice tracking."
        ),
        "role": "count_exploratory",
    },
    "median_2d_length_um": {
        "label": "Specimen median 2D length (um)",
        "short": "2D length",
        "question": "Are nuclei typically longer in the image plane?",
        "meaning": (
            "The specimen median of each reconstructed nucleus's maximum calibrated "
            "2D centerline length. Higher values indicate longer projected nuclei."
        ),
        "role": "morphology",
    },
    "median_body_width_um": {
        "label": "Specimen median apparent body width (um)",
        "short": "Body width",
        "question": "Are nuclei typically broader or thinner?",
        "meaning": (
            "The specimen median of representative-plane subpixel perpendicular "
            "body chords. Higher values indicate broader apparent masks; width "
            "remains sensitive to PSF, focus, and mask boundaries."
        ),
        "role": "morphology",
    },
    "median_length_body_width_ratio": {
        "label": "Specimen median length / apparent body width",
        "short": "Length / width",
        "question": "Are nuclei more elongated or more rounded?",
        "meaning": (
            "A shape ratio. Higher values indicate longer, more slender nuclei; "
            "interpret it together with length and width."
        ),
        "role": "morphology",
    },
    "median_3d_length_um": {
        "label": "Specimen median 3D length (um)",
        "short": "3D length",
        "question": "Are nuclei longer after accounting for Z orientation?",
        "meaning": (
            "A calibrated projection-plus-Z estimate combining maximum lateral "
            "length with Z span. It is not a surface-mesh trace."
        ),
        "role": "morphology",
    },
    "median_3d_tortuosity": {
        "label": "Specimen median 3D tortuosity",
        "short": "3D tortuosity",
        "question": "Are nuclei straighter or more curved?",
        "meaning": (
            "Estimated path length divided by end-to-end distance. Values near 1 "
            "are straight; larger values indicate increasing curvature."
        ),
        "role": "morphology",
    },
    "median_3d_thickness_um": {
        "label": "Specimen median effective thickness (um)",
        "short": "Effective thickness",
        "question": "Is the reconstructed nucleus effectively thicker?",
        "meaning": (
            "A diameter proxy derived from volume divided by length. It is not a "
            "direct width and is PSF- and segmentation-sensitive."
        ),
        "role": "morphology_psf_sensitive",
    },
    "median_3d_volume_um3": {
        "label": "Specimen median volume (um3)",
        "short": "Volume",
        "question": "How much calibrated 3D mask volume does a typical nucleus occupy?",
        "meaning": (
            "Filled-mask area accumulated through Z. It can reflect length and "
            "thickness but also depends on thresholds, Z sampling, and PSF."
        ),
        "role": "morphology_psf_sensitive",
    },
    "median_3d_z_span_um": {
        "label": "Specimen median Z span (um)",
        "short": "Z span",
        "question": "Through how much optical depth is a nucleus observed?",
        "meaning": (
            "The calibrated distance from first to last linked Z plane. A "
            "single-slice track has zero endpoint-to-endpoint Z span."
        ),
        "role": "tracking_sensitive",
    },
}

BIOLOGICAL_METRICS = tuple(
    metric
    for metric, definition in METRICS.items()
    if definition["role"] in {"morphology", "morphology_psf_sensitive"}
)
QC_METRICS = tuple(metric for metric in METRICS if metric not in BIOLOGICAL_METRICS)


def reference_group(groups):
    groups = sorted(str(group) for group in groups)
    for group in groups:
        label = group.lower()
        if any(token in label for token in ("control", "reference", "wild", "w1118")):
            return group
        if label in {"wt", "ctrl"}:
            return group
    return groups[0]


def bh_qvalues(values):
    values = np.asarray(values, dtype=float)
    result = np.full(values.shape, np.nan)
    finite = np.flatnonzero(np.isfinite(values))
    if not finite.size:
        return result
    order = finite[np.argsort(values[finite])]
    ranked = values[order]
    adjusted = ranked * len(ranked) / np.arange(1, len(ranked) + 1)
    adjusted = np.minimum.accumulate(adjusted[::-1])[::-1]
    result[order] = np.minimum(adjusted, 1.0)
    return result


def cliffs_delta(reference, comparison):
    differences = comparison[:, None] - reference[None, :]
    return float(
        (
            np.count_nonzero(differences > 0)
            - np.count_nonzero(differences < 0)
        )
        / differences.size
    )


def hedges_g(reference, comparison):
    n_reference = len(reference)
    n_comparison = len(comparison)
    if min(n_reference, n_comparison) < 2:
        return np.nan
    variance = (
        (n_reference - 1) * np.var(reference, ddof=1)
        + (n_comparison - 1) * np.var(comparison, ddof=1)
    ) / (n_reference + n_comparison - 2)
    if not np.isfinite(variance) or variance <= 0:
        return 0.0 if np.mean(reference) == np.mean(comparison) else np.nan
    d_value = (np.mean(comparison) - np.mean(reference)) / math.sqrt(variance)
    correction = 1 - 3 / (4 * (n_reference + n_comparison) - 9)
    return float(d_value * correction)


def boolean_series(series):
    if pd.api.types.is_bool_dtype(series):
        return series.fillna(False)
    return (
        series.fillna("")
        .astype(str)
        .str.strip()
        .str.lower()
        .isin({"true", "1", "yes"})
    )


MIN_INFERENCE_SPECIMENS_PER_GROUP = 3


def compute_statistics(specimens, reference, comparison, seed=57057):
    records = []
    for index, (metric, definition) in enumerate(METRICS.items()):
        if metric not in specimens.columns:
            continue
        reference_values = pd.to_numeric(
            specimens.loc[specimens["group"] == reference, metric],
            errors="coerce",
        ).dropna().to_numpy()
        comparison_values = pd.to_numeric(
            specimens.loc[specimens["group"] == comparison, metric],
            errors="coerce",
        ).dropna().to_numpy()
        if not len(reference_values) or not len(comparison_values):
            continue
        reference_median = float(np.median(reference_values))
        comparison_median = float(np.median(comparison_values))
        median_difference = comparison_median - reference_median
        inference_available = (
            len(reference_values) >= MIN_INFERENCE_SPECIMENS_PER_GROUP
            and len(comparison_values) >= MIN_INFERENCE_SPECIMENS_PER_GROUP
        )
        if inference_available:
            rng = np.random.default_rng(seed + index)
            bootstrap = np.asarray(
                [
                    np.median(rng.choice(comparison_values, len(comparison_values), replace=True))
                    - np.median(rng.choice(reference_values, len(reference_values), replace=True))
                    for _ in range(5_000)
                ]
            )
            ci_low = float(np.quantile(bootstrap, 0.025))
            ci_high = float(np.quantile(bootstrap, 0.975))
            permutation_p = float(permutation_test(
                (reference_values, comparison_values),
                lambda ref, comp: np.median(comp) - np.median(ref),
                permutation_type="independent",
                vectorized=False,
                n_resamples=9_999,
                alternative="two-sided",
                rng=np.random.default_rng(seed + 1_000 + index),
            ).pvalue)
            mann_whitney = mannwhitneyu(
                reference_values, comparison_values,
                alternative="two-sided", method="auto",
            )
            welch = ttest_ind(
                comparison_values, reference_values,
                equal_var=False, nan_policy="omit",
            )
            mann_whitney_u = float(mann_whitney.statistic)
            mann_whitney_p = float(mann_whitney.pvalue)
            welch_t = float(welch.statistic)
            welch_p = float(welch.pvalue)
            delta = cliffs_delta(reference_values, comparison_values)
            hedges = hedges_g(reference_values, comparison_values)
            inference_status = "exploratory_specimen_level_inference"
        else:
            ci_low = ci_high = permutation_p = np.nan
            mann_whitney_u = mann_whitney_p = np.nan
            welch_t = welch_p = delta = hedges = np.nan
            inference_status = "insufficient_specimens"
        records.append(
            {
                "metric": metric,
                "metric_label": definition["label"],
                "analysis_role": definition["role"],
                "analysis_unit": "biological specimen",
                "inference_status": inference_status,
                "minimum_specimens_per_group_for_inference": (
                    MIN_INFERENCE_SPECIMENS_PER_GROUP
                ),
                "reference_group": reference,
                "comparison_group": comparison,
                "reference_n": len(reference_values),
                "comparison_n": len(comparison_values),
                "reference_mean": np.mean(reference_values),
                "comparison_mean": np.mean(comparison_values),
                "reference_sd": (
                    np.std(reference_values, ddof=1)
                    if len(reference_values) > 1 else np.nan
                ),
                "comparison_sd": (
                    np.std(comparison_values, ddof=1)
                    if len(comparison_values) > 1 else np.nan
                ),
                "reference_q1": np.quantile(reference_values, 0.25),
                "reference_median": reference_median,
                "reference_q3": np.quantile(reference_values, 0.75),
                "comparison_q1": np.quantile(comparison_values, 0.25),
                "comparison_median": comparison_median,
                "comparison_q3": np.quantile(comparison_values, 0.75),
                "median_difference_comparison_minus_reference": median_difference,
                "median_percent_difference": (
                    100 * median_difference / reference_median
                    if reference_median != 0
                    else np.nan
                ),
                "bootstrap_median_difference_95ci_low": ci_low,
                "bootstrap_median_difference_95ci_high": ci_high,
                "cliffs_delta_comparison_minus_reference": delta,
                "permutation_median_test_p": permutation_p,
                "mann_whitney_u": mann_whitney_u,
                "mann_whitney_p": mann_whitney_p,
                "welch_t": welch_t,
                "welch_t_p": welch_p,
                "hedges_g_comparison_minus_reference": hedges,
            }
        )
    result = pd.DataFrame(records)
    for p_column, q_column in (
        ("permutation_median_test_p", "permutation_bh_fdr_q"),
        ("mann_whitney_p", "mann_whitney_bh_fdr_q"),
        ("welch_t_p", "welch_t_bh_fdr_q"),
    ):
        result[q_column] = bh_qvalues(result[p_column])
    return result


def morphology_proportions(tracks):
    valid = tracks[boolean_series(tracks["technical_valid"])].copy()
    valid["morphology_warning"] = boolean_series(valid["morphology_warning"])
    warning_text = (
        valid["morphology_warning_reasons"].fillna("").astype(str).str.lower()
    )
    rows = []
    for (sample_id, group), frame in valid.groupby(["sample_id", "group"]):
        length = pd.to_numeric(frame["total_3d_length_um"], errors="coerce")
        warnings = warning_text.loc[frame.index]
        denominator = max(len(frame), 1)
        rows.append(
            {
                "sample_id": sample_id,
                "group": group,
                "technical_valid_nuclei": len(frame),
                "fraction_below_2_um": float((length < 2).sum() / denominator),
                "fraction_2_to_15_um": float(
                    ((length >= 2) & (length < 15)).sum() / denominator
                ),
                "fraction_15_to_20_um": float(
                    ((length >= 15) & (length <= 20)).sum() / denominator
                ),
                "fraction_single_slice": float(
                    (pd.to_numeric(frame["n_slices"], errors="coerce") == 1).sum()
                    / denominator
                ),
                "fraction_morphology_warning": float(
                    frame["morphology_warning"].sum() / denominator
                ),
                "fraction_short_warning": float(
                    warnings.str.contains(r"(?:^|,)short(?:,|$)").sum()
                    / denominator
                ),
                "fraction_long_warning": float(
                    warnings.str.contains(r"(?:^|,)long(?:,|$)").sum()
                    / denominator
                ),
                "fraction_wide_warning": float(
                    warnings.str.contains(r"(?:^|,)wide(?:,|$)").sum()
                    / denominator
                ),
                "fraction_tortuous_warning": float(
                    warnings.str.contains("tortuous").sum() / denominator
                ),
            }
        )
    return pd.DataFrame(rows)


def count_depth_correlations(specimens):
    records = []
    for metric in (
        "estimated_nuclei_per_1000_um2",
        "estimated_nuclei_per_100000_um3",
    ):
        for exposure in (
            "slice_count",
            "roi_area_um2",
            "sampled_roi_volume_um3",
        ):
            rho, p_value = spearmanr(specimens[metric], specimens[exposure])
            records.append(
                {
                    "count_metric": metric,
                    "exposure_metric": exposure,
                    "spearman_rho": rho,
                    "spearman_p": p_value,
                    "interpretation": (
                        "Residual exposure association; count comparison requires caution"
                        if abs(rho) >= 0.3
                        else "Weak exposure association in this dataset"
                    ),
                }
            )
    result = pd.DataFrame(records)
    result["spearman_bh_fdr_q"] = bh_qvalues(result["spearman_p"])
    return result


def save_figure(figure, figure_dir, stem):
    png_path = figure_dir / f"{stem}.png"
    figure.savefig(png_path, dpi=220, bbox_inches="tight")
    return png_path


def specimen_overview_figure(
    specimens,
    groups,
    metric_keys=None,
    title="Specimen-level biological comparison",
):
    metric_keys = tuple(metric_keys or METRICS)
    columns = 3
    rows = math.ceil(len(metric_keys) / columns)
    figure, axes = plt.subplots(
        rows,
        columns,
        figsize=(13.5, max(5.2, rows * 3.5)),
        squeeze=False,
    )
    colors = ["#2878B5", "#D1495B"]
    rng = np.random.default_rng(57057)
    for axis, metric in zip(axes.flat, metric_keys):
        definition = METRICS[metric]
        for group_index, group in enumerate(groups):
            values = pd.to_numeric(
                specimens.loc[specimens["group"] == group, metric],
                errors="coerce",
            ).dropna().to_numpy()
            x_values = group_index + rng.uniform(-0.075, 0.075, len(values))
            axis.scatter(
                x_values,
                values,
                s=42,
                color=colors[group_index],
                edgecolor="white",
                linewidth=0.6,
            )
            axis.hlines(
                np.median(values),
                group_index - 0.18,
                group_index + 0.18,
                color="#202020",
                linewidth=2.2,
            )
        axis.set_title(definition["short"], fontsize=10)
        axis.set_xticks([0, 1], groups)
        axis.grid(axis="y", color="#D8D8D8", linewidth=0.7)
        axis.spines[["top", "right"]].set_visible(False)
    for axis in axes.flat[len(metric_keys) :]:
        axis.set_visible(False)
    figure.suptitle(
        f"{title}\n"
        "Each point is one specimen; black bars show medians",
        fontsize=16,
        fontweight="bold",
    )
    figure.tight_layout(rect=(0, 0, 1, 0.96))
    return figure


def forest_figure(statistics):
    frame = statistics.iloc[::-1].reset_index(drop=True)
    reference = frame["reference_group"].iloc[0]
    comparison = frame["comparison_group"].iloc[0]
    denominator = frame["reference_median"].replace(0, np.nan)
    effect = frame["median_percent_difference"]
    low = 100 * frame["bootstrap_median_difference_95ci_low"] / denominator
    high = 100 * frame["bootstrap_median_difference_95ci_high"] / denominator
    colors = [
        "#C98B2E" if role == "count_exploratory" else "#3A7D70"
        for role in frame["analysis_role"]
    ]
    figure, axis = plt.subplots(figsize=(11, 7.5))
    y_values = np.arange(len(frame))
    for y_value, estimate, lower, upper, color in zip(
        y_values, effect, low, high, colors
    ):
        axis.errorbar(
            estimate,
            y_value,
            xerr=[[estimate - lower], [upper - estimate]],
            fmt="none",
            ecolor=color,
            elinewidth=2,
            capsize=4,
        )
    axis.scatter(effect, y_values, color=colors, s=58, zorder=3)
    axis.axvline(0, color="#303030", linestyle="--", linewidth=1)
    axis.set_yticks(y_values, frame["metric_label"])
    axis.set_xlabel(
        f"Median percent difference: {comparison} minus {reference}"
    )
    axis.set_title(
        "Specimen-level effect estimates\n"
        "Points are median differences; bars are bootstrap 95% intervals",
        fontweight="bold",
    )
    axis.grid(axis="x", alpha=0.25)
    axis.spines[["top", "right"]].set_visible(False)
    figure.tight_layout()
    return figure


def length_width_figure(specimens, groups):
    colors = ["#2878B5", "#D1495B"]
    figure, axis = plt.subplots(figsize=(8, 6))
    for index, group in enumerate(groups):
        frame = specimens[specimens["group"] == group]
        axis.scatter(
            frame["median_2d_length_um"],
            frame["median_body_width_um"],
            s=68,
            color=colors[index],
            edgecolor="white",
            linewidth=0.7,
            label=f"{group} specimens",
        )
        axis.scatter(
            frame["median_2d_length_um"].median(),
            frame["median_body_width_um"].median(),
            marker="X",
            s=180,
            color=colors[index],
            edgecolor="#202020",
            linewidth=0.8,
        )
    axis.set_xlabel("Specimen median 2D length (um)")
    axis.set_ylabel("Specimen median apparent body width (um)")
    axis.set_title(
        "Length-width relationship\nLarge X symbols show group medians",
        fontweight="bold",
    )
    axis.legend()
    axis.grid(alpha=0.25)
    axis.spines[["top", "right"]].set_visible(False)
    figure.tight_layout()
    return figure


def morphology_figure(proportions, groups):
    panels = {
        "fraction_below_2_um": "Below 2 um",
        "fraction_2_to_15_um": "2-15 um",
        "fraction_15_to_20_um": "15-20 um review band",
        "fraction_single_slice": "Single-slice tracks",
        "fraction_morphology_warning": "Any morphology warning",
        "fraction_tortuous_warning": "Tortuous warning",
    }
    colors = ["#2878B5", "#D1495B"]
    rng = np.random.default_rng(57057)
    figure, axes = plt.subplots(2, 3, figsize=(12, 7))
    for axis, (column, title) in zip(axes.flat, panels.items()):
        for index, group in enumerate(groups):
            values = proportions.loc[proportions["group"] == group, column] * 100
            x = index + rng.uniform(-0.07, 0.07, len(values))
            axis.scatter(
                x,
                values,
                color=colors[index],
                edgecolor="white",
                linewidth=0.6,
                s=40,
            )
            axis.hlines(
                values.median(),
                index - 0.18,
                index + 0.18,
                color="#202020",
                linewidth=2,
            )
        axis.set_title(title, fontsize=10)
        axis.set_xticks([0, 1], groups)
        axis.set_ylabel("Nuclei per specimen (%)")
        axis.grid(axis="y", alpha=0.25)
        axis.spines[["top", "right"]].set_visible(False)
    figure.suptitle(
        "Per-specimen morphology and tracking proportions\n"
        "Categories annotate technical-valid nuclei; they are not rejection populations",
        fontsize=14,
        fontweight="bold",
    )
    figure.tight_layout(rect=(0, 0, 1, 0.93))
    return figure


def count_depth_figure(specimens, common_depth, groups):
    colors = ["#2878B5", "#D1495B"]
    figure, axes = plt.subplots(2, 2, figsize=(11, 8))
    panels = [
        (
            "slice_count",
            "estimated_nuclei_per_1000_um2",
            "Area density versus stack depth",
        ),
        (
            "slice_count",
            "estimated_nuclei_per_100000_um3",
            "Volume density versus stack depth",
        ),
        (
            "sampled_roi_volume_um3",
            "estimated_unique_nuclei",
            "Raw estimated nuclei versus sampled volume",
        ),
    ]
    for axis, (x_column, y_column, title) in zip(axes.flat[:3], panels):
        for index, group in enumerate(groups):
            frame = specimens[specimens["group"] == group]
            axis.scatter(
                frame[x_column],
                frame[y_column],
                color=colors[index],
                s=48,
                edgecolor="white",
                linewidth=0.6,
                label=group,
            )
        rho, p_value = spearmanr(specimens[x_column], specimens[y_column])
        axis.set_title(f"{title}\nSpearman rho={rho:.2f}, p={p_value:.3g}")
        axis.set_xlabel(x_column.replace("_", " "))
        axis.set_ylabel(y_column.replace("_", " "))
        axis.grid(alpha=0.25)
        axis.spines[["top", "right"]].set_visible(False)
    axis = axes.flat[3]
    if common_depth is not None and not common_depth.empty:
        rng = np.random.default_rng(57057)
        column = "central_window_nuclei_per_100000_um3"
        for index, group in enumerate(groups):
            values = common_depth.loc[common_depth["group"] == group, column]
            x = index + rng.uniform(-0.07, 0.07, len(values))
            axis.scatter(
                x,
                values,
                color=colors[index],
                s=48,
                edgecolor="white",
                linewidth=0.6,
            )
            axis.hlines(
                values.median(),
                index - 0.18,
                index + 0.18,
                color="#202020",
                linewidth=2,
            )
        depth_slices = int(common_depth["common_depth_slices"].iloc[0])
        axis.set_title(f"Fixed central {depth_slices}-slice sensitivity")
        axis.set_xticks([0, 1], groups)
        axis.set_ylabel("Estimated nuclei per 100,000 um3")
        axis.grid(axis="y", alpha=0.25)
        axis.spines[["top", "right"]].set_visible(False)
    figure.suptitle(
        "Count and acquisition-depth diagnostics\nCount endpoints remain exploratory",
        fontsize=14,
        fontweight="bold",
    )
    handles, labels = axes.flat[0].get_legend_handles_labels()
    figure.legend(handles, labels, loc="lower center", ncol=2)
    figure.tight_layout(rect=(0, 0.05, 1, 0.94))
    return figure


def statistics_table_figure(statistics):
    frame = statistics.copy()
    reference = frame["reference_group"].iloc[0]
    comparison = frame["comparison_group"].iloc[0]
    display = pd.DataFrame(
        {
            "Measure": frame["metric_label"],
            f"{reference}\nmedian": frame["reference_median"].map(
                lambda x: f"{x:.3g}"
            ),
            f"{comparison}\nmedian": frame["comparison_median"].map(
                lambda x: f"{x:.3g}"
            ),
            "Median\nchange (%)": frame["median_percent_difference"].map(
                lambda x: f"{x:+.1f}%"
            ),
            "Effect size\nCliff's delta": frame[
                "cliffs_delta_comparison_minus_reference"
            ].map(lambda x: f"{x:+.2f}"),
            "Median test\nFDR q": frame["permutation_bh_fdr_q"].map(
                lambda x: f"{x:.3g}"
            ),
            "Rank test\nFDR q": frame["mann_whitney_bh_fdr_q"].map(
                lambda x: f"{x:.3g}"
            ),
            "Mean test\nFDR q": frame["welch_t_bh_fdr_q"].map(
                lambda x: f"{x:.3g}"
            ),
        }
    )
    figure, axis = plt.subplots(figsize=(13.5, 7.5))
    axis.axis("off")
    table = axis.table(
        cellText=display.values,
        colLabels=display.columns,
        bbox=[0.01, 0.28, 0.98, 0.58],
        cellLoc="center",
        colLoc="center",
        colWidths=[0.29, 0.10, 0.10, 0.10, 0.11, 0.10, 0.10, 0.10],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1, 1.6)
    for column in range(len(display.columns)):
        table[(0, column)].set_facecolor("#DCE8EE")
        table[(0, column)].set_text_props(weight="bold")
    axis.set_title(
        "Specimen-level statistical summary\n"
        "Each value is calculated across biological specimens, not pooled nuclei",
        fontsize=15,
        fontweight="bold",
        pad=20,
    )
    interpretation = (
        f"Medians: the typical specimen value in each group.  "
        f"Median change: ({comparison} median - {reference} median) / "
        f"{reference} median x 100; positive means higher in {comparison}.\n"
        "Cliff's delta: direction and separation, from -1 to +1; positive means "
        f"{comparison} specimens tend to have higher values.  "
        "FDR q: p-value adjusted for testing multiple measurements; q < 0.05 "
        "is evidence of a group difference.\n"
        "Median test: permutation comparison of specimen medians (primary).  "
        "Rank test: Mann-Whitney comparison of specimen ordering.  "
        "Mean test: Welch comparison of specimen means."
    )
    axis.text(
        0.01,
        0.20,
        interpretation,
        transform=axis.transAxes,
        fontsize=9.5,
        va="top",
        linespacing=1.45,
        bbox={
            "boxstyle": "round,pad=0.6",
            "facecolor": "#F3F6F7",
            "edgecolor": "#B8C6CC",
        },
    )
    figure.tight_layout()
    return figure


def text_page(title, paragraphs):
    figure = plt.figure(figsize=(13.5, 10.5))
    figure.patch.set_facecolor("white")
    figure.suptitle(title, fontsize=19, fontweight="bold", y=0.965)
    y = 0.89
    for heading, text in paragraphs:
        figure.text(0.06, y, heading, fontsize=13, fontweight="bold", va="top")
        wrapped = "\n".join(textwrap.wrap(text, width=145))
        figure.text(0.06, y - 0.035, wrapped, fontsize=11, va="top")
        y -= 0.035 + 0.031 * (wrapped.count("\n") + 1) + 0.055
    return figure


def write_biological_excel(
    path,
    specimens,
    group_summary,
    statistics,
    reference,
    comparison,
):
    definitions = pd.DataFrame(
        [
            {
                "metric": metric,
                "display_name": definition["label"],
                "analysis_role": definition["role"],
                "biological_question": definition["question"],
                "meaning": definition["meaning"],
            }
            for metric, definition in METRICS.items()
            if metric in BIOLOGICAL_METRICS
        ]
    )
    methods = pd.DataFrame(
        [
            {
                "item": "Analysis unit",
                "explanation": (
                    "One biological specimen. Individual nuclei are nested "
                    "measurements and are not independent replicates."
                ),
            },
            {
                "item": "Primary test",
                "explanation": (
                    "Two-sided permutation test of the difference in specimen "
                    "medians, with 9,999 permutations."
                ),
            },
            {
                "item": "Mann-Whitney U",
                "explanation": (
                    "Secondary rank-based sensitivity test. It tests whether "
                    "values from one group tend to rank above values from the other."
                ),
            },
            {
                "item": "Welch t-test",
                "explanation": (
                    "Secondary mean-based sensitivity test that does not assume "
                    "equal group variances."
                ),
            },
            {
                "item": "Why not two-group ANOVA?",
                "explanation": (
                    "With exactly two groups, ordinary one-way ANOVA and the "
                    "equal-variance two-sample t-test are mathematically equivalent. "
                    "Welch's t-test is more useful here because equal variance is "
                    "not assumed."
                ),
            },
            {
                "item": "Multiple testing",
                "explanation": (
                    "Benjamini-Hochberg FDR q-values are calculated separately "
                    "for permutation, Mann-Whitney, and Welch test families."
                ),
            },
            {
                "item": "Effect sizes",
                "explanation": (
                    "Cliff's delta is rank-based; Hedges g is standardized "
                    "mean difference. Positive values indicate comparison above reference."
                ),
            },
        ]
    )
    readme = pd.DataFrame(
        {
            "Item": [
                "Reference group",
                "Comparison group",
                "Primary biological table",
                "Primary statistical table",
                "QC location",
            ],
            "Value": [
                reference,
                comparison,
                "Specimen_Data",
                "Statistical_Tests",
                "../02_quality_control",
            ],
        }
    )
    with pd.ExcelWriter(path, engine="xlsxwriter") as writer:
        sheets = {
            "README": readme,
            "Specimen_Data": specimens,
            "Group_Descriptives": group_summary,
            "Statistical_Tests": statistics,
            "Metric_Definitions": definitions,
            "Statistical_Methods": methods,
        }
        for sheet_name, frame in sheets.items():
            frame.to_excel(writer, sheet_name=sheet_name, index=False)
            worksheet = writer.sheets[sheet_name]
            worksheet.freeze_panes(1, 0)
            if not frame.empty:
                worksheet.autofilter(0, 0, len(frame), len(frame.columns) - 1)
            for index, column in enumerate(frame.columns):
                sample_lengths = [
                    len(str(value)) for value in frame[column].head(200).tolist()
                ]
                width = min(
                    max([len(str(column)), *sample_lengths]) + 2,
                    55,
                )
                worksheet.set_column(index, index, width)


def write_qc_excel(
    path,
    specimen_qc,
    qc_statistics,
    proportions,
    correlations,
    common_depth,
    source_inventory,
    reference,
    comparison,
):
    readme = pd.DataFrame(
        {
            "Item": [
                "Reference group",
                "Comparison group",
                "Purpose",
                "Primary biological results",
            ],
            "Value": [
                reference,
                comparison,
                (
                    "Technical and exploratory diagnostics. These sheets support "
                    "interpretation but are not the primary biological result."
                ),
                "../01_biological_results",
            ],
        }
    )
    definitions = pd.DataFrame(
        [
            {
                "metric": metric,
                "display_name": METRICS[metric]["label"],
                "analysis_role": METRICS[metric]["role"],
                "meaning": METRICS[metric]["meaning"],
            }
            for metric in QC_METRICS
        ]
    )
    sheets = {
        "README": readme,
        "Specimen_Technical_QC": specimen_qc,
        "Exploratory_Statistics": qc_statistics,
        "Morphology_Warnings": proportions,
        "Count_Depth_Diagnostics": correlations,
        "Common_Depth_Sensitivity": (
            common_depth if common_depth is not None else pd.DataFrame()
        ),
        "QC_Metric_Definitions": definitions,
        "Source_File_Inventory": source_inventory,
    }
    with pd.ExcelWriter(path, engine="xlsxwriter") as writer:
        for sheet_name, frame in sheets.items():
            frame.to_excel(writer, sheet_name=sheet_name, index=False)
            worksheet = writer.sheets[sheet_name]
            worksheet.freeze_panes(1, 0)
            if not frame.empty:
                worksheet.autofilter(0, 0, len(frame), len(frame.columns) - 1)
            for index, column in enumerate(frame.columns):
                sample_lengths = [
                    len(str(value)) for value in frame[column].head(200).tolist()
                ]
                width = min(
                    max([len(str(column)), *sample_lengths]) + 2,
                    55,
                )
                worksheet.set_column(index, index, width)


def add_ppt_title(slide, text):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.55))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True


def write_powerpoint(path, title, figures, explanation_pages):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(2))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(30)
    paragraph.font.bold = True
    paragraph.alignment = 1
    for figure_title, figure_path in figures:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_ppt_title(slide, figure_title)
        slide.shapes.add_picture(
            str(figure_path),
            Inches(0.55),
            Inches(0.85),
            width=Inches(12.2),
            height=Inches(6.25),
        )
    for page_title, bullets in explanation_pages:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        add_ppt_title(slide, page_title)
        box = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(11.9), Inches(5.9))
        frame = box.text_frame
        frame.word_wrap = True
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(17)
            paragraph.space_after = Pt(11)
    presentation.save(path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--study-output", required=True)
    parser.add_argument("--output-folder", default="")
    args = parser.parse_args()

    study_root = Path(args.study_output).resolve()
    output = (
        Path(args.output_folder).resolve()
        if args.output_folder
        else study_root / "between_sample_analysis"
    )
    biological_dir = output / "01_biological_results"
    qc_dir = output / "02_quality_control"
    biological_figure_dir = biological_dir / "figures"
    biological_data_dir = biological_dir / "data"
    qc_figure_dir = qc_dir / "figures"
    qc_data_dir = qc_dir / "data"
    qc_source_dir = qc_dir / "source_files"
    for directory in (
        biological_figure_dir,
        biological_data_dir,
        qc_figure_dir,
        qc_data_dir,
        qc_source_dir,
    ):
        directory.mkdir(parents=True, exist_ok=True)

    specimens = pd.read_csv(study_root / "specimen_summary.csv")
    specimens = specimens[specimens["status"] == "complete"].copy()
    tracks = pd.read_csv(study_root / "study_track_records.csv")
    technical_qc_path = study_root / "specimen_technical_qc.csv"
    specimen_qc = (
        pd.read_csv(technical_qc_path)
        if technical_qc_path.exists()
        else specimens.copy()
    )
    common_path = study_root / "common_depth_sensitivity.csv"
    common_depth = pd.read_csv(common_path) if common_path.exists() else None
    groups = sorted(specimens["group"].dropna().astype(str).unique())
    if len(groups) != 2:
        raise ValueError(
            f"This pairwise report requires exactly two groups; found {groups}"
        )
    reference = reference_group(groups)
    comparison = next(group for group in groups if group != reference)
    groups = [reference, comparison]

    all_statistics = compute_statistics(specimens, reference, comparison)
    inference_available = bool(
        len(all_statistics)
        and (all_statistics["inference_status"] != "insufficient_specimens").all()
    )
    inference_guidance = (
        "Use specimen-level medians, effect sizes, confidence intervals, and "
        "permutation-test FDR q-values as exploratory comparisons."
        if inference_available
        else (
            "Descriptive comparison only. Inferential statistics are unavailable "
            f"until each group has at least {MIN_INFERENCE_SPECIMENS_PER_GROUP} "
            "biological specimens."
        )
    )
    biological_statistics = all_statistics[
        all_statistics["metric"].isin(BIOLOGICAL_METRICS)
    ].copy()
    qc_statistics = all_statistics[
        all_statistics["metric"].isin(QC_METRICS)
    ].copy()
    proportions = morphology_proportions(tracks)
    correlations = count_depth_correlations(specimens)
    biological_group_summary = specimens.groupby("group").agg(
        specimen_count=("sample_id", "count"),
        **{
            f"{metric}_median": (metric, "median")
            for metric in BIOLOGICAL_METRICS
            if metric in specimens.columns
        },
    ).reset_index()
    biological_columns = [
        column
        for column in (
            "sample_id",
            "group",
            "acquisition_class",
            *BIOLOGICAL_METRICS,
        )
        if column in specimens.columns
    ]
    biological_specimens = specimens[biological_columns].copy()

    biological_statistics.to_csv(
        biological_data_dir / "biological_statistical_tests.csv",
        index=False,
    )
    biological_specimens.to_csv(
        biological_data_dir / "specimen_biological_measurements.csv",
        index=False,
    )
    biological_group_summary.to_csv(
        biological_data_dir / "biological_group_descriptives.csv",
        index=False,
    )
    qc_statistics.to_csv(
        qc_data_dir / "exploratory_qc_statistics.csv",
        index=False,
    )
    specimen_qc.to_csv(
        qc_data_dir / "specimen_technical_qc.csv",
        index=False,
    )
    proportions.to_csv(
        qc_data_dir / "morphology_warning_proportions.csv",
        index=False,
    )
    correlations.to_csv(
        qc_data_dir / "count_depth_correlations.csv",
        index=False,
    )
    if common_depth is not None:
        common_depth.to_csv(
            qc_data_dir / "common_depth_sensitivity.csv",
            index=False,
        )

    source_names = (
        "audit_rebuild_summary.json",
        "common_depth_sensitivity.json",
        "common_depth_sensitivity.pdf",
        "group_summary.csv",
        "normalization_qc.json",
        "runtime_parameters.json",
        "specimen_group_comparison.pdf",
        "specimen_group_comparisons.csv",
        "specimen_group_comparison_qc.json",
        "study_manifest.csv",
        "study_run_state.json",
    )
    source_records = []
    for source_name in source_names:
        source = study_root / source_name
        if not source.exists():
            continue
        destination_name = (
            "legacy_specimen_group_comparison.pdf"
            if source_name == "specimen_group_comparison.pdf"
            else source_name
        )
        destination = qc_source_dir / destination_name
        shutil.copy2(source, destination)
        source_records.append(
            {
                "source_file": source_name,
                "copied_to": str(destination.relative_to(output)),
                "purpose": (
                    "Superseded source report retained for provenance; its useful "
                    "biological content is incorporated in the primary report."
                    if source_name == "specimen_group_comparison.pdf"
                    else "Study-level QC or provenance source"
                ),
            }
        )
    source_records.append(
        {
            "source_file": "study_track_records.csv",
            "copied_to": "",
            "purpose": (
                "Large track-level audit remains in the study root to avoid a "
                "104 MB duplicate; derived QC tables are included here."
            ),
        }
    )
    legacy_package = qc_source_dir / "legacy_unsplit_comparison_package"
    if legacy_package.is_dir():
        source_records.append(
            {
                "source_file": (
                    "legacy biologically_meaningful_between_sample_comparisons"
                ),
                "copied_to": str(legacy_package.relative_to(output)),
                "purpose": (
                    "Superseded unsplit package retained for provenance only."
                ),
            }
        )
    source_inventory = pd.DataFrame(source_records)
    source_inventory.to_csv(
        qc_data_dir / "source_file_inventory.csv",
        index=False,
    )

    biological_figures = [
        (
            "Specimen-level overview",
            "01_specimen_overview",
            specimen_overview_figure(
                specimens,
                groups,
                BIOLOGICAL_METRICS,
                "Specimen-level biological measurements",
            ),
        ),
        (
            "Effect-size forest plot",
            "02_effect_size_forest",
            forest_figure(biological_statistics),
        ),
        (
            "Length-width relationship",
            "03_length_width_relationship",
            length_width_figure(specimens, groups),
        ),
        (
            "Biological statistical summary",
            "04_biological_statistical_summary",
            statistics_table_figure(biological_statistics),
        ),
    ]
    biological_figure_paths = []
    for title, stem, figure in biological_figures:
        path = save_figure(figure, biological_figure_dir, stem)
        biological_figure_paths.append((title, path))

    qc_figures = [
        (
            "Exploratory count and tracking endpoints",
            "01_qc_endpoint_overview",
            specimen_overview_figure(
                specimens,
                groups,
                QC_METRICS,
                "Exploratory count and tracking endpoints",
            ),
        ),
        (
            "Morphology-warning and tracking proportions",
            "02_morphology_warning_proportions",
            morphology_figure(proportions, groups),
        ),
        (
            "Count and stack-depth diagnostics",
            "03_count_depth_diagnostics",
            count_depth_figure(specimens, common_depth, groups),
        ),
        (
            "Exploratory QC statistical summary",
            "04_qc_statistical_summary",
            statistics_table_figure(qc_statistics),
        ),
    ]
    qc_figure_paths = []
    for title, stem, figure in qc_figures:
        path = save_figure(figure, qc_figure_dir, stem)
        qc_figure_paths.append((title, path))

    methods_page = text_page(
        "Statistical methods and interpretation",
        [
            (
                "Inference availability",
                inference_guidance,
            ),
            (
                "Analysis unit",
                "Each point and test uses one biological specimen. Individual "
                "nuclei are nested measurements and are not counted as independent replicates.",
            ),
            (
                "How to read the summary table",
                "The two group medians describe the typical specimen in each group. "
                "Median change is the comparison-minus-reference difference expressed "
                "as a percentage of the reference median. Cliff's delta describes "
                "direction and group separation from -1 to +1. Each FDR q-value is "
                "the corresponding test's p-value after correction for testing "
                "multiple measurements; q below 0.05 is the conventional evidence threshold.",
            ),
            (
                "Primary inference",
                "The two-sided permutation test compares the difference in specimen "
                "medians. Bootstrap intervals describe uncertainty in that median difference.",
            ),
            (
                "Mann-Whitney U",
                "A secondary rank-based sensitivity test. It asks whether values "
                "from one group tend to rank above values from the other; it is "
                "not strictly a test of medians unless distribution shapes are similar.",
            ),
            (
                "Welch's t-test and ANOVA",
                "Welch's t-test is a secondary mean-based sensitivity test and "
                "does not require equal variances. With exactly two groups, ordinary "
                "one-way ANOVA adds no independent information because it is "
                "equivalent to the equal-variance two-sample t-test.",
            ),
            (
                "Multiple tests",
                "Benjamini-Hochberg FDR q-values are reported separately for each "
                "test family. Effect sizes and confidence intervals should be "
                "interpreted alongside q-values.",
            ),
        ],
    )
    meaning_page = text_page(
        "Biological meaning of the primary measurements",
        [
            (
                definition["label"],
                f"{definition['question']} {definition['meaning']}",
            )
            for metric, definition in METRICS.items()
            if metric in BIOLOGICAL_METRICS
        ],
    )

    biological_pdf = biological_dir / "Biological_Comparison_Report.pdf"
    with PdfPages(biological_pdf) as pdf:
        cover = text_page(
            "Biologically meaningful between-sample comparison",
            [
                (
                    "Groups",
                    f"Reference: {reference}. Comparison: {comparison}. "
                    f"Specimens: {len(specimens[specimens.group == reference])} "
                    f"and {len(specimens[specimens.group == comparison])}.",
                ),
                (
                    "What to use",
                    inference_guidance,
                ),
                (
                    "Scope",
                    "This report contains biological morphology comparisons only. "
                    "Count, acquisition-depth, tracking, warning-category, and audit "
                    "diagnostics are in the separate quality-control package.",
                ),
                (
                    "Supporting files",
                    "The accompanying Excel workbook contains every specimen-level "
                    "value and every derived table plotted in this report.",
                ),
            ],
        )
        pdf.savefig(cover, bbox_inches="tight")
        plt.close(cover)
        for _, _, figure in biological_figures:
            pdf.savefig(figure, bbox_inches="tight")
            plt.close(figure)
        pdf.savefig(methods_page, bbox_inches="tight")
        pdf.savefig(meaning_page, bbox_inches="tight")
    plt.close(methods_page)
    plt.close(meaning_page)

    biological_excel = biological_dir / "Biological_Comparison_Data.xlsx"
    write_biological_excel(
        biological_excel,
        biological_specimens,
        biological_group_summary,
        biological_statistics,
        reference,
        comparison,
    )
    biological_ppt = biological_dir / "Biological_Comparison_Presentation.pptx"
    write_powerpoint(
        biological_ppt,
        f"Biological comparison: {reference} versus {comparison}",
        biological_figure_paths,
        [
            (
                "Statistical methods",
                [
                    "Biological specimen is the replicate; nuclei are nested measurements.",
                    inference_guidance,
                    "When available, the permutation median test is primary; Mann-Whitney and Welch tests are sensitivity analyses.",
                    "Nuclei are never treated as independent biological replicates.",
                ],
            ),
            (
                "How to interpret the measurements",
                [
                    f"{definition['short']}: {definition['question']} "
                    f"{definition['meaning']}"
                    for metric, definition in METRICS.items()
                    if metric in BIOLOGICAL_METRICS
                ],
            ),
        ],
    )

    qc_explanation_page = text_page(
        "How to use the quality-control package",
        [
            (
                "Not the primary biological result",
                "These panels assess acquisition depth, count normalization, "
                "tracking span, morphology-warning frequencies, and processing "
                "provenance. They should not replace the biological morphology report.",
            ),
            (
                "Count endpoints",
                "Area- and volume-normalized counts remain exploratory when they "
                "retain an association with slice count or sampled volume. The "
                "fixed-depth sensitivity panel helps assess this dependence.",
            ),
            (
                "Morphology warnings",
                "Warnings annotate unusual but technical-valid nuclei. They are "
                "not rejection populations and must not be used to force either "
                "group toward reference morphology.",
            ),
            (
                "Tracking Z span",
                "Z span is sensitive to acquisition spacing and cross-slice linking. "
                "A single-slice nucleus can still be biologically genuine.",
            ),
        ],
    )
    qc_pdf = qc_dir / "Quality_Control_Report.pdf"
    with PdfPages(qc_pdf) as pdf:
        cover = text_page(
            "Between-sample quality control and exploratory diagnostics",
            [
                (
                    "Groups",
                    f"Reference: {reference}. Comparison: {comparison}.",
                ),
                (
                    "Purpose",
                    "Use this package to judge acquisition comparability, "
                    "normalization, tracking behavior, and warning frequencies.",
                ),
                (
                    "Primary results",
                    "Biological conclusions belong in ../01_biological_results.",
                ),
            ],
        )
        pdf.savefig(cover, bbox_inches="tight")
        plt.close(cover)
        for _, _, figure in qc_figures:
            pdf.savefig(figure, bbox_inches="tight")
            plt.close(figure)
        pdf.savefig(qc_explanation_page, bbox_inches="tight")
    plt.close(qc_explanation_page)

    qc_excel = qc_dir / "Quality_Control_Data.xlsx"
    write_qc_excel(
        qc_excel,
        specimen_qc,
        qc_statistics,
        proportions,
        correlations,
        common_depth,
        source_inventory,
        reference,
        comparison,
    )
    qc_ppt = qc_dir / "Quality_Control_Presentation.pptx"
    write_powerpoint(
        qc_ppt,
        f"Quality control: {reference} versus {comparison}",
        qc_figure_paths,
        [
            (
                "How to use these diagnostics",
                [
                    "These are technical and exploratory diagnostics, not the primary biological result.",
                    "Count endpoints require acquisition-depth review.",
                    "Morphology warnings annotate unusual technical-valid nuclei; they do not reject them.",
                    "Tracking Z span is acquisition- and linking-sensitive.",
                    "The primary biological report is in ../01_biological_results.",
                ],
            )
        ],
    )

    metadata = {
        "study_output": str(study_root),
        "reference_group": reference,
        "comparison_group": comparison,
        "specimen_counts": specimens.groupby("group").size().to_dict(),
        "analysis_unit": "biological specimen",
        "inference_status": (
            "exploratory_specimen_level_inference"
            if inference_available else "insufficient_specimens"
        ),
        "minimum_specimens_per_group_for_inference": MIN_INFERENCE_SPECIMENS_PER_GROUP,
        "primary_test": (
            "two-sided permutation test of specimen median difference"
            if inference_available else "unavailable"
        ),
        "secondary_tests": (
            ["Mann-Whitney U", "Welch t-test"]
            if inference_available else []
        ),
        "random_seed": 57057,
        "biological_results": {
            "pdf": str(biological_pdf),
            "powerpoint": str(biological_ppt),
            "excel": str(biological_excel),
        },
        "quality_control": {
            "pdf": str(qc_pdf),
            "powerpoint": str(qc_ppt),
            "excel": str(qc_excel),
        },
    }
    (output / "report_metadata.json").write_text(
        json.dumps(metadata, indent=2),
        encoding="utf-8",
    )
    (output / "README.md").write_text(
        "\n".join(
            [
                "# Between-sample analysis package",
                "",
                f"Reference group: {reference}",
                f"Comparison group: {comparison}",
                "",
                "Use `01_biological_results` for biological interpretation.",
                "Use `02_quality_control` for acquisition, normalization, tracking, and audit diagnostics.",
                "",
                "The biological specimen is the analysis unit.",
                "The previous specimen-group PDF is retained only as a provenance source in the QC folder.",
            ]
        ),
        encoding="utf-8",
    )
    print(f"Created {biological_pdf}")
    print(f"Created {biological_ppt}")
    print(f"Created {biological_excel}")
    print(f"Created {qc_pdf}")
    print(f"Created {qc_ppt}")
    print(f"Created {qc_excel}")


if __name__ == "__main__":
    main()
