#!/usr/bin/env python3
"""
Sperm Nucleus Segmentation & 3D Morphometrics Pipeline  —  v11 (Combined)
=========================================================================
A production-ready image-analysis pipeline for automated detection,
measurement, and 3D reconstruction of sperm nuclei in
fluorescence z-stack microscopy images.

Biological context
------------------
Spermiogenesis — the post-meiotic differentiation of round spermatids into
mature spermatozoa — involves dramatic nuclear elongation, chromatin condensation,
and apical plunging of the nucleus toward the basal lamina of the seminiferous
tubule.  Quantifying these morphological changes across large image stacks is
labour-intensive when done manually.  This pipeline automates the process by:

1. Detecting spermatid nuclei as thin, dim, ridge-like objects in each 2D
   Z-slice using a skeleton-first strategy (Frangi ridge filter → Otsu
   binarisation → morphological closing → geodesic skeleton).
2. Measuring per-cell biometrics: geodesic length (µm), Euclidean width,
   tortuosity, endpoint count, branch complexity, centroid, and area.
3. Linking detections across Z-slices into 3D tracks using a nearest-neighbour
   frame-to-frame assignment algorithm.
4. Computing 3D morphometrics from tracks: true 3D length, Z-extent,
   approximated volume, pitch angle, taper ratio, and nearest-neighbour
   packing density.
5. Exporting results as CSV, a multi-sheet Excel workbook, a multi-page PDF
   report, and a native PowerPoint (.pptx) dashboard with editable charts.

Pipeline architecture
---------------------
::

    Image stack (TIF / PNG / JPG)
          │
          ▼
    load_batch_files()  ────────────────────── natural-sorted file list
          │
    ┌─────┴──────────────────────────────────────────────────────────┐
    │  Per-slice loop  (process_batch → process_one_image)           │
    │                                                                 │
    │  robust_imread()   →  segment_slice()  →  measure_spermatids() │
    │       │                    │                     │              │
    │   raw ndarray        mask / skel / ridge    list[dict]           │
    │                            │                     │              │
    │                     make_overlay()        rows_from_results()   │
    │                            │                     │              │
    │                      overlay PNG              per-slice CSV     │
    └─────────────────────────────────────────────────────────────────┘
          │
    track_across_slices()  ──── 3D track DataFrame
          │
    ┌─────┴──────────────────────────────────────────────────────────┐
    │  Reporting                                                      │
    │  generate_excel_report()  →  multi-sheet XLSX workbook         │
    │  generate_batch_report()  →  multi-page PDF (matplotlib)       │
    │  generate_pptx_report()   →  native Office PPTX dashboard      │
    └─────────────────────────────────────────────────────────────────┘

Key configuration parameters (``CONFIG`` dict)
-----------------------------------------------
``UM_PER_PX_XY``       Physical pixel size in µm (from microscope metadata).
``UM_PER_SLICE_Z``     Z-step size in µm (from microscope metadata).
``FRANGI_SCALE_RANGE`` Tubeness filter scales (px); set to match sperm nucleus width.
``MIN_SKEL_LEN_PX``    Minimum skeleton length accepted (removes debris/noise).
``MAX_WIDTH_PX``       Maximum skeleton width accepted (rejects merged clusters).
``MAX_BRANCH_NODES``   Maximum branch-point count; >0 tolerates bridged tips.
``MAX_TORTUOSITY``     Maximum curvature index; rejects snaking merged networks.
``DO_TRACKING``        Enable/disable cross-slice 3D linking.
``TRACK_MAX_DIST_UM``  Maximum centroid displacement between adjacent Z-slices.

Usage
-----
Launch the interactive GUI (recommended)::

    python sperm_segmentation_v11_combined.py

Run a headless batch analysis::

    python sperm_segmentation_v11_combined.py --batch

Analyse a single slice::

    python sperm_segmentation_v11_combined.py --single --z 4

Dependencies
------------
numpy, scipy, scikit-image, pandas, matplotlib, tifffile, opencv-python,
Pillow, xlsxwriter, python-pptx, tkinter (stdlib)

Author
------
Dushyant Mishra  |  Findlay Lab
"""

import os, sys, glob, re, time, warnings, heapq, argparse, pathlib as pl
warnings.filterwarnings("ignore")

import numpy as np
import pandas as pd
import tifffile
import matplotlib
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.path import Path
print(f"[matplotlib backend: {matplotlib.get_backend()}]")

from skimage import measure, morphology, exposure
from skimage.filters import meijering, gaussian, apply_hysteresis_threshold
from skimage.morphology import skeletonize
from scipy.ndimage import distance_transform_edt, grey_dilation
from scipy.spatial import cKDTree
from matplotlib.backends.backend_pdf import PdfPages

try:
    import cv2 as _cv2
    _HAVE_CV2 = True
except ImportError:
    _HAVE_CV2 = False


# =============================================================================
# CONFIG
# =============================================================================

CONFIG = {
    # ── run mode ─────────────────────────────────────────────────────────────
    "RUN_MODE": "single",          # "single" | "batch"

    # ── single-image selection ────────────────────────────────────────────────
    "SINGLE_IMAGE_SELECTION_MODE": "dialog",  # "path" | "z_index" | "dialog"
    "SINGLE_TEST_IMAGE": r"C:\Users\dmishra\Desktop\sperm images\Project001_Series002_z15_ch00.tif",
    "SINGLE_Z_INDEX": 15,

    # ── input / output ────────────────────────────────────────────────────────
    "INPUT_DIR":    r"C:\Users\dmishra\Desktop\sperm images",
    "OUTPUT_DIR":   r"C:\Users\dmishra\Desktop\sperm images\sperm_results_v8",
    "FILE_PATTERN": "Project001_Series002_z*_ch00.tif",

    # ── calibration ───────────────────────────────────────────────────────────
    "UM_PER_PX_XY":   0.378788,
    "UM_PER_SLICE_Z": 0.346184,

    # ── image enhancement ─────────────────────────────────────────────────────
    "CLAHE_CLIP":   0.120,   # Raised significantly to pull out dim edges inside solid center cysts
    "CLAHE_KERNEL": 64,
    "BG_SIGMA":     10.9,    # Widened background subtraction to flatten out the solid centers
    "RIDGE_SIGMAS": [1, 2, 3, 4],

    # ── hysteresis threshold ──────────────────────────────────────────────────
    "THRESHOLD_HI": 90.1,    # High percentile suppresses all background noise cleanly
    "THRESHOLD_LO": 81.0,    # Firm lower bound stops fragmentation

    # ── morphological cleanup ─────────────────────────────────────────────────
    "CLOSE_RADIUS":  1,
    "MIN_HOLE_AREA": 10,
    "MIN_OBJ_PX":    5,

    # ── skeleton-level gap bridging ───────────────────────────────────────────
    # P1: reduced from 10 → 5 to prevent chaining distinct spermatids
    "MAX_BRIDGE_PX": 5,

    # ── branch pruning & automated splitting ──────────────────────────────────
    "MAX_BRANCH_LEN_PX": 5,   # prune spurs shorter than this before measuring
    "BREAK_JUNCTIONS": True,  # automatically sever all branching intersections into distinct lines

    # ── optional early mask-level shape filter ────────────────────────────────
    "USE_EARLY_SHAPE_FILTER": False,
    "MIN_ECCENTRICITY": 0.60,
    "MAX_MINOR_PX":     12.0,
    "MIN_AXIS_RATIO":   1.4,
    "MIN_MAJOR_PX":     5,

    # ── post-skeleton filters ─────────────────────────────────────────────────
    "MIN_SKEL_LEN_PX":        13.9, # min geodesic length (px)
    "MAX_GEODESIC_LEN_PX":    65.0, # (~13.0um); backend actively cuts chains longer than this
    "MAX_WIDTH_PX":           25.2, # Wide enough to allow multi-stranded clusters to survive
    "MIN_LENGTH_WIDTH_RATIO": 1.6,  # Relaxed to allow clustered blobs to pass

    # ── NEW topology filters ──────────────────────────────────────────────────
    "MAX_BRANCH_NODES": 872,      # Allowed to be massive so user can Split the resulting thick webs
    "MAX_TORTUOSITY": 840.6,      # Relaxed to let extremely zig-zagging clusters live safely
    "MAX_ENDPOINT_COUNT": 734,    # Relaxed so massive webs aren't dropped outright

    # N4: loops
    "ALLOW_LOOPS": True,     # don't reject spermatids that cross themselves  # moderate

    # ── tracking across z ─────────────────────────────────────────────────────
    "DO_TRACKING":          True,
    "TRACK_MAX_DIST_UM":    4.0,
    "TRACK_MAX_GAP_SLICES": 1,

    # ── output / debug ────────────────────────────────────────────────────────
    "SAVE_DEBUG_IMAGES":   True,
    "SAVE_MASK_TIFS":      True,
    "SAVE_LABEL_TIFS":     True,
    "SAVE_OVERLAYS":       True,
    "SAVE_DETAIL_FIGURE":  True,
    "SHOW_PREVIEW_WINDOW": True,
    "SHOW_DEBUG_PREVIEW":  True,
}


# =============================================================================
# CONFIG VALIDATION
# =============================================================================

_REQUIRED = {
    "RUN_MODE": str, "SINGLE_IMAGE_SELECTION_MODE": str,
    "SINGLE_TEST_IMAGE": str, "SINGLE_Z_INDEX": int,
    "INPUT_DIR": str, "OUTPUT_DIR": str, "FILE_PATTERN": str,
    "UM_PER_PX_XY": float, "UM_PER_SLICE_Z": float,
    "CLAHE_CLIP": float, "CLAHE_KERNEL": int, "BG_SIGMA": (int, float),
    "RIDGE_SIGMAS": list,
    "THRESHOLD_HI": (int, float), "THRESHOLD_LO": (int, float),
    "CLOSE_RADIUS": int, "MIN_HOLE_AREA": int, "MIN_OBJ_PX": int,
    "MAX_BRIDGE_PX": (int, float), "MAX_BRANCH_LEN_PX": (int, float),
    "USE_EARLY_SHAPE_FILTER": bool,
    "MIN_SKEL_LEN_PX": (int, float), "MAX_GEODESIC_LEN_PX": (int, float),
    "MAX_WIDTH_PX": (int, float), "MIN_LENGTH_WIDTH_RATIO": (int, float),
    "MAX_BRANCH_NODES": (int, float),
    "MAX_TORTUOSITY": (int, float),
    "MAX_ENDPOINT_COUNT": (int, float),
    "DO_TRACKING": bool, "TRACK_MAX_DIST_UM": (int, float),
    "TRACK_MAX_GAP_SLICES": int,
    "SAVE_DEBUG_IMAGES": bool, "SAVE_MASK_TIFS": bool,
    "SAVE_LABEL_TIFS": bool, "SAVE_OVERLAYS": bool,
    "SAVE_DETAIL_FIGURE": bool, "SHOW_PREVIEW_WINDOW": bool,
    "SHOW_DEBUG_PREVIEW": bool,
}


def validate_config(cfg):
    """
    Validates the pipeline CONFIG dictionary for required keys, correct types, and logical consistency.

    Raises a descriptive ValueError listing ALL problems found so engineers can fix
    everything in one go rather than playing whack-a-mole with sequential errors.

    Checks performed
    ----------------
    - Every key in ``_REQUIRED`` is present in *cfg*.
    - Each value has the expected Python type (e.g., ``float`` for calibration parameters).
    - ``THRESHOLD_LO`` < ``THRESHOLD_HI`` (hysteresis thresholding only works in this order).
    - ``RUN_MODE`` is exactly one of ``'single'`` or ``'batch'``.
    - The deprecated ``'REJECT_BRANCHES'`` Boolean key is absent (replaced by ``MAX_BRANCH_NODES`` int).

    Args:
        cfg (dict): The configuration mapping to validate.

    Raises:
        ValueError: If *any* of the above checks fail. The message lists all errors.
    """
    errors = []
    for key, expected in _REQUIRED.items():
        if key not in cfg:
            errors.append(f"  MISSING: '{key}'")
        elif not isinstance(cfg[key], expected):
            errors.append(f"  WRONG TYPE '{key}': "
                          f"got {type(cfg[key]).__name__}, want {expected}")
    if cfg.get("THRESHOLD_LO", 0) >= cfg.get("THRESHOLD_HI", 100):
        errors.append("  THRESHOLD_LO must be < THRESHOLD_HI")
    if cfg.get("RUN_MODE", "") not in ("single", "batch"):
        errors.append("  RUN_MODE must be 'single' or 'batch'")
    if "REJECT_BRANCHES" in cfg:
        errors.append("  'REJECT_BRANCHES' was removed in v8. "
                      "Use MAX_BRANCH_NODES (int) instead:\n"
                      "    REJECT_BRANCHES=True  → MAX_BRANCH_NODES=0\n"
                      "    REJECT_BRANCHES=False → MAX_BRANCH_NODES=9999")
    if errors:
        raise ValueError("CONFIG errors:\n" + "\n".join(errors))


# =============================================================================
# UTILITIES
# =============================================================================

_N8 = [(-1,-1),(-1,0),(-1,1),(0,-1),(0,1),(1,-1),(1,0),(1,1)]


def ensure_dir(p):
    """
    Creates directory *p* (and all missing parents) if it does not already exist.

    Uses ``exist_ok=True`` to avoid TOCTOU race conditions in parallel runs.

    Args:
        p (str): Absolute or relative path of the directory to create.
    """
    os.makedirs(p, exist_ok=True)

def get_unique_batch_dir(base_dir):
    """
    Checks for 'batch_output', then 'batch_output_1', 'batch_output_2', etc.
    Returns the first available path.
    """
    candidate = os.path.join(base_dir, "batch_output")
    if not os.path.exists(candidate):
        return candidate
    
    counter = 1
    while True:
        candidate = os.path.join(base_dir, f"batch_output_{counter}")
        if not os.path.exists(candidate):
            return candidate
        counter += 1


def natural_sort_key(s):
    """
    Key for natural sorting (e.g., img2.tif comes before img10.tif).
    """
    import re
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split(r'(\d+)', os.path.basename(s))]

def extract_z_index(path, sequence_idx=0):
    """
    Extracts Z-index from filename, fallback to sequence index.
    """
    import re
    m = re.search(r"[zZ](\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    m = re.search(r"(\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    return sequence_idx

def robust_imread(path):
    """
    Reads image with multi-engine fallback for maximum compatibility.
    Handles TIF (compressed/uncompressed), PNG, and JPG.
    """
    p_lower = path.lower()
    
    # 1. Primary Loader: tifffile (Best for 16-bit/Microscopy TIFFs)
    if p_lower.endswith('.tif') or p_lower.endswith('.tiff'):
        try:
            return tifffile.imread(path)
        except Exception:
            # If tifffile fails (e.g. missing LZW codecs), proceed to fallbacks
            pass

    # 2. Fallback A: Pillow (Very robust for standard formats + LZW)
    try:
        from PIL import Image
        return np.array(Image.open(path))
    except Exception:
        pass
        
    # 3. Fallback B: OpenCV (Fast, handles many JPG/PNG variants)
    if _HAVE_CV2:
        try:
            img = _cv2.imread(path, _cv2.IMREAD_UNCHANGED)
            if img is not None:
                # CV2 loads BGR by default, but for grayscale sperm images it's fine
                # If 3 channels, convert to RGB
                if len(img.shape) == 3:
                    return _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
                return img
        except Exception:
            pass
            
    # 4. Last resort: matplotlib
    try:
        return plt.imread(path)
    except Exception:
        pass
            
    # Final default failure (will be caught by UI)
    raise RuntimeError(f"All image engines failed to read: {os.path.basename(path)}")

def normalize_display(img):
    """
    Contrast-stretches a raw image for colourmap display.

    Uses the 1st–99.5th percentile range rather than min–max to suppress hot
    pixels and dark-corner vignetting artefacts common in fluorescence microscopy.
    The small epsilon (1e-9) prevents division-by-zero on flat images.

    Args:
        img (np.ndarray): Input image of any integer or float dtype.

    Returns:
        np.ndarray: Float32 array clipped to [0, 1] ready for matplotlib display.
    """
    a = img.astype(np.float32)
    lo, hi = np.percentile(a, 1), np.percentile(a, 99.5)
    return np.clip((a - lo) / (hi - lo + 1e-9), 0, 1)


def _imwrite(path, arr_uint8, cmap="gray"):
    """
    Writes a uint8 image to disk using the best available engine.

    Engine priority:
    1. **OpenCV** — fastest; handles large images without matplotlib overhead.
       Converts RGB→BGR before writing (OpenCV internal convention).
    2. **Matplotlib** — fallback when OpenCV is absent; uses ``plt.imsave``
       which respects the *cmap* argument for single-channel saves.

    Args:
        path (str): Destination file path (extension determines format, e.g. ``.png``).
        arr_uint8 (np.ndarray): uint8 image array, shape ``(H, W)`` or ``(H, W, 3)``.
        cmap (str): Matplotlib colourmap name used only for grayscale fallback saves.
    """
    if _HAVE_CV2:
        if arr_uint8.ndim == 2:
            _cv2.imwrite(path, arr_uint8)
        else:
            # OpenCV stores BGR; convert from RGB before writing
            _cv2.imwrite(path, _cv2.cvtColor(arr_uint8, _cv2.COLOR_RGB2BGR))
    else:
        plt.imsave(path, arr_uint8,
                   cmap=(cmap if arr_uint8.ndim == 2 else None),
                   vmin=0, vmax=255)


def save_gray(path, img_float):
    """
    Saves a floating-point single-channel image as an 8-bit grayscale PNG.

    Applies a full min–max stretch (appropriate for debug images where absolute
    intensity is less important than structure visibility).

    Args:
        path (str): Output file path.
        img_float (np.ndarray): Float image of arbitrary range.
    """
    a = img_float.astype(np.float32)
    a = (a - a.min()) / (a.max() - a.min() + 1e-9)
    _imwrite(path, (a * 255).astype(np.uint8), cmap="gray")


def save_mask(path, mask_bool):
    """
    Saves a binary boolean mask as a black-and-white 8-bit PNG.

    Pixels where *mask_bool* is ``True`` are written as 255 (white);
    background pixels are 0 (black). This is the standard convention
    for binary mask overlays in ImageJ / FIJI.

    Args:
        path (str): Output file path.
        mask_bool (np.ndarray[bool]): Boolean mask array.
    """
    _imwrite(path, mask_bool.astype(np.uint8) * 255, cmap="gray")


def load_batch_files(input_dir, pattern):
    """
    Discovers and sorts all image files for a batch run.

    Uses a three-tier fallback strategy so the pipeline works even when
    users supply unusual file extensions or mixed case (.TIF vs .tif):

    1. Exact pattern match (e.g. ``Project001_Series002_z*_ch00.tif``).
    2. ``.tif`` → ``.tiff`` substitution.
    3. Broad glob over all supported extensions (tif, tiff, png, jpg, jpeg)
       in both lower- and upper-case variants.

    After discovery, files are sorted with :func:`natural_sort_key` so that
    ``z2`` comes before ``z10`` (lexicographic sort would reverse this).  A
    Z-index is extracted from each filename or assigned by sequence position
    as a fallback.

    Args:
        input_dir (str): Directory to search for image files.
        pattern (str): Glob pattern relative to *input_dir*.

    Returns:
        tuple[list[str], list[int]]: Sorted file paths and corresponding Z-indices.

    Raises:
        FileNotFoundError: If no matching images are found after all fallbacks.
    """
    files = glob.glob(os.path.join(input_dir, pattern))
    if not files:
        # Fallback 1: .tif → .tiff extension swap
        files = glob.glob(os.path.join(input_dir, pattern.replace(".tif", ".tiff")))
    if not files:
        # Fallback 2: Broad scan of all supported image formats
        for ext in ['*.tif', '*.tiff', '*.png', '*.jpg', '*.jpeg']:
            files.extend(glob.glob(os.path.join(input_dir, ext)))
            files.extend(glob.glob(os.path.join(input_dir, ext.upper())))
        files = list(set(files))  # remove cross-extension duplicates

    if not files:
        raise FileNotFoundError(f"No supported image files found in '{input_dir}'")

    # Natural sort preserves correct Z-stack ordering (z2 < z10)
    files = sorted(files, key=natural_sort_key)

    # Extract z-index from filename pattern; fall back to sequence index
    z_idx = [extract_z_index(f, sequence_idx=i) for i, f in enumerate(files)]

    print(f"Found {len(files)} slices: Z = {z_idx}")
    return files, z_idx


def choose_single_image(cfg):
    """
    Resolves the path to a single test image according to the configured selection mode.

    Three modes are supported (set via ``SINGLE_IMAGE_SELECTION_MODE``):

    - ``'path'``     — Use the hard-coded ``SINGLE_TEST_IMAGE`` path directly.
    - ``'z_index'``  — Find the file whose filename encodes the Z-index in
                      ``SINGLE_Z_INDEX`` (e.g. ``z15`` → ``...z15_ch00.tif``).
    - ``'dialog'``   — Open a Tkinter file-picker dialog for interactive selection.
                      Falls back gracefully if Tkinter is unavailable.

    Args:
        cfg (dict): Pipeline configuration dictionary.

    Returns:
        str: Absolute path to the selected image file.

    Raises:
        FileNotFoundError: For ``'path'`` mode when the file does not exist, or
                           ``'z_index'`` mode when no file matches the Z-index.
        RuntimeError: For ``'dialog'`` mode when Tkinter fails or the user cancels.
        ValueError: If ``SINGLE_IMAGE_SELECTION_MODE`` is not one of the three valid values.
    """
    mode = cfg["SINGLE_IMAGE_SELECTION_MODE"].lower()
    if mode == "path":
        p = cfg["SINGLE_TEST_IMAGE"]
        if not os.path.exists(p):
            raise FileNotFoundError(f"Not found: {p}")
        return p
    if mode == "z_index":
        z = int(cfg["SINGLE_Z_INDEX"])
        files = (glob.glob(os.path.join(cfg["INPUT_DIR"], cfg["FILE_PATTERN"])) or
                 glob.glob(os.path.join(cfg["INPUT_DIR"],
                                        cfg["FILE_PATTERN"].replace(".tif", ".tiff"))))
        hits = [f for f in files if extract_z_index(f) == z]
        if not hits:
            raise FileNotFoundError(f"No file for z={z}")
        print(f"Using z={z}: {hits[0]}")
        return hits[0]
    if mode == "dialog":
        try:
            import tkinter as tk
            from tkinter import filedialog
            root = tk.Tk(); root.withdraw()  # hide the empty root window
            p = filedialog.askopenfilename(
                title="Choose image", initialdir=cfg["INPUT_DIR"],
                filetypes=[("TIFF", "*.tif *.tiff"), ("All", "*.*")])
            root.destroy()
        except Exception as e:
            raise RuntimeError(f"File dialog failed: {e}")
        if not p:
            raise RuntimeError("No file selected.")
        print(f"Selected: {p}")
        return p
    raise ValueError("SINGLE_IMAGE_SELECTION_MODE: 'path'|'z_index'|'dialog'")


# =============================================================================
# SKELETON UTILITIES
# =============================================================================

def find_endpoints(skel_bool):
    """
    Identifies all endpoint pixels (tip pixels) in a binary skeleton image.

    In an 8-connected skeleton, a pixel is an *endpoint* if exactly one of its
    eight neighbours is also part of the skeleton.  These pixels correspond to
    the physical tips of spermatid filaments.  A clean spermatid should have
    exactly 2 endpoints (head and tail end); a bridged fragment pair still has
    2 endpoints after merging.

    Args:
        skel_bool (np.ndarray[bool]): 2D binary skeleton image (True = skeleton pixel).

    Returns:
        list[tuple[int, int]]: List of ``(row, col)`` pixel coordinates of all endpoints.
    """
    H, W = skel_bool.shape
    ys, xs = np.where(skel_bool)
    sk_set = set(zip(ys.tolist(), xs.tolist()))
    return [(r, c) for r, c in sk_set
            if sum(1 for dr, dc in _N8
                   if 0 <= r+dr < H and 0 <= c+dc < W
                   and skel_bool[r+dr, c+dc]) == 1]


def bridge_skeleton_endpoints(skel_bool, skel_labeled, max_gap_px):
    """
    Join endpoint pairs from DIFFERENT components within max_gap_px by a
    1-px straight line.  Preserves the original mask so width estimates
    are not inflated.
    """
    if max_gap_px <= 0:
        return skel_bool.copy()
    H, W   = skel_bool.shape
    out    = skel_bool.copy()
    eps    = find_endpoints(skel_bool)
    if not eps:
        return out
    ep_arr    = np.array(eps, dtype=np.float32)
    ep_labels = skel_labeled[ep_arr[:, 0].astype(int), ep_arr[:, 1].astype(int)]
    pairs     = cKDTree(ep_arr).query_pairs(r=max_gap_px, output_type="ndarray")
    for i, j in pairs:
        if ep_labels[i] == ep_labels[j]:
            continue
        r0, c0 = int(eps[i][0]), int(eps[i][1])
        r1, c1 = int(eps[j][0]), int(eps[j][1])
        n = max(abs(r1-r0), abs(c1-c0)) + 1
        rs = np.clip(np.round(np.linspace(r0, r1, n)).astype(int), 0, H-1)
        cs = np.clip(np.round(np.linspace(c0, c1, n)).astype(int), 0, W-1)
        out[rs, cs] = True
    return out


def prune_branches(skel_bool, max_branch_len):
    """
    Iteratively remove endpoints to shorten side-branches ≤ max_branch_len px.
    """
    if max_branch_len <= 0:
        return skel_bool.copy()
    H, W  = skel_bool.shape
    skel  = skel_bool.copy()
    for _ in range(int(max_branch_len)):
        eps = find_endpoints(skel)
        if not eps:
            break
        for r, c in eps:
            n = sum(1 for dr, dc in _N8
                    if 0 <= r+dr < H and 0 <= c+dc < W and skel[r+dr, c+dc])
            if n == 1:
                skel[r, c] = False
    return skel


# =============================================================================
# GEODESIC & TOPOLOGY MEASUREMENT
# =============================================================================

def _build_adj(coords, W):
    """
    Builds a lightweight adjacency list for a set of skeleton pixel coordinates.

    Uses a *linearised index* trick: each pixel ``(r, c)`` is mapped to an integer
    ``r * W + c`` so that neighbour lookup is an O(1) dictionary lookup rather than
    a 2D array access.  This is important for large skeleton components.

    Edge weights follow the Chebyshev / Euclidean convention:
    - Axis-aligned move (4-connected step): weight = 1.0
    - Diagonal move (8-connected step):     weight = √2 ≈ 1.41421

    Args:
        coords (np.ndarray): Shape ``(N, 2)`` integer array of ``(row, col)`` pixel positions.
        W (int): Image width in pixels (used to compute the linear index).

    Returns:
        list[list[tuple[int, float]]]: Adjacency list where ``adj[i]`` is a list of
        ``(neighbour_index, edge_weight)`` pairs.
    """
    n       = len(coords)
    # Linearise (row, col) -> single int for O(1) membership test
    lin     = coords[:, 0] * W + coords[:, 1]
    lin2idx = {int(v): i for i, v in enumerate(lin.tolist())}
    lin_set = set(lin.tolist())
    adj     = [[] for _ in range(n)]
    for i, (r, c) in enumerate(coords.tolist()):
        for dr, dc in _N8:
            lk = (r + dr) * W + (c + dc)
            if lk in lin_set:
                # Diagonal edges are longer by a factor of √2
                w = 1.41421356 if (dr != 0 and dc != 0) else 1.0
                adj[i].append((lin2idx[lk], w))
    return adj


def _dijkstra(adj, src, n):
    """
    Runs Dijkstra's shortest-path algorithm from source node *src* and returns
    the *farthest* reachable node and its distance.

    This is the first BFS pass in the double-BFS algorithm for computing the
    true *geodesic diameter* (longest shortest path) of a skeleton component.
    The second BFS starts from the farthest node found here.

    Mathematical note:
        Geodesic length = shortest path distance through the skeleton graph,
        which accounts for the actual curvature of the filament rather than
        the straight-line Euclidean distance between endpoints.

    Args:
        adj (list[list[tuple[int, float]]]): Adjacency list from :func:`_build_adj`.
        src (int): Index of the source node in *adj*.
        n (int): Total number of nodes.

    Returns:
        tuple[int, float]: ``(farthest_node_index, distance_to_farthest)``
    """
    d = np.full(n, np.inf); d[src] = 0.0
    pq = [(0.0, src)]
    while pq:
        cost, u = heapq.heappop(pq)
        if cost > d[u]:
            continue  # stale entry in the heap — skip
        for v, w in adj[u]:
            nd = cost + w
            if nd < d[v]:
                d[v] = nd
                heapq.heappush(pq, (nd, v))
    far = int(np.argmax(d))
    return far, float(d[far])


def measure_topology(coords, W, allow_loops=False):
    """
    Compute geodesic length, tortuosity, endpoint count, and branch-node count
    for one skeleton component.

    Returns
    -------
    dict with keys: geo_len, tortuosity, n_endpoints, n_branch_nodes, reason
    or None if the component is a loop and allow_loops=False.
    """
    n   = len(coords)
    adj = _build_adj(coords, W)
    deg = [len(a) for a in adj]

    n_endpoints    = sum(1 for d in deg if d == 1)
    n_branch_nodes = sum(1 for d in deg if d >  2)
    eps_idx        = [i for i, d in enumerate(deg) if d == 1]

    # ── Loop handling ─────────────────────────────────────────────────────────
    if not eps_idx:
        if not allow_loops:
            return None  # discard loop
        seen, total = set(), 0.0
        for u, nbrs in enumerate(adj):
            for v, w in nbrs:
                key = (min(u, v), max(u, v))
                if key not in seen:
                    seen.add(key); total += w
        return {"geo_len": total, "tortuosity": 1.0,
                "n_endpoints": 0, "n_branch_nodes": n_branch_nodes,
                "reason": "loop"}

    # ── Double-BFS for geodesic ───────────────────────────────────────────────
    b, _  = _dijkstra(adj, eps_idx[0], n)
    c, gl = _dijkstra(adj, b, n)

    # ── Tortuosity ────────────────────────────────────────────────────────────
    p0  = coords[b]
    p1  = coords[c]
    euc = float(np.sqrt((p0[0]-p1[0])**2 + (p0[1]-p1[1])**2))
    tort = gl / (euc + 1e-9)

    return {"geo_len": float(gl), "tortuosity": tort,
            "n_endpoints": n_endpoints, "n_branch_nodes": n_branch_nodes,
            "reason": "ok"}


# =============================================================================
# SEGMENTATION PIPELINE
# =============================================================================

def apply_optional_early_shape_filter(mask, cfg):
    """
    Optionally removes non-elongated connected components from the binary mask
    *before* skeletonisation, based on region shape descriptors.

    This is an early-stage pre-filter controlled by ``USE_EARLY_SHAPE_FILTER``.
    When enabled, it eliminates round debris, fat nuclei, and large globular
    artefacts that would otherwise generate spurious skeleton trees downstream.

    Biological rationale
    --------------------
    Round debris (dead cells, lipid droplets, imaging artefacts) tends to have:
    - Low eccentricity (close to circular)
    - Short major axis (small, compact object)
    - Low axis ratio (major ≈ minor, i.e., not elongated)

    Elongated spermatid nuclei, by contrast, have high eccentricity (> 0.6),
    a long major axis, and a large major/minor axis ratio.

    Args:
        mask (np.ndarray[bool]): Binary mask from hysteresis thresholding.
        cfg (dict): Pipeline config; reads ``USE_EARLY_SHAPE_FILTER``,
            ``MIN_ECCENTRICITY``, ``MAX_MINOR_PX``, ``MIN_AXIS_RATIO``,
            and ``MIN_MAJOR_PX``.

    Returns:
        np.ndarray[bool]: Filtered binary mask (unchanged if the filter is disabled).
    """
    if not cfg["USE_EARLY_SHAPE_FILTER"]:
        return mask  # filter disabled — pass through unchanged
    labeled = measure.label(mask)
    keep = [p.label for p in measure.regionprops(labeled)
            if (p.eccentricity      >= cfg["MIN_ECCENTRICITY"] and   # must be elongated
                p.minor_axis_length <= cfg["MAX_MINOR_PX"]     and   # must be narrow
                p.major_axis_length / (p.minor_axis_length + 1e-9) >= cfg["MIN_AXIS_RATIO"] and  # must be rod-like
                p.major_axis_length >= cfg["MIN_MAJOR_PX"])]
    return np.isin(labeled, keep)


def segment_slice(img_raw, cfg, z_idx=None, debug_dir=None, roi_mask=None):
    """
    Executes advanced 2D multi-stage morphology segmentation to detect and isolate spermatid nuclei.
    
    Args:
        img_raw (np.ndarray): The raw 2D grayscale image array of the z-slice.
        cfg (dict): Pipeline hyperparameters including thresholding, CLAHE limits, and morphological radius.
        z_idx (int, optional): The current Z-index integer for debugging context.
        debug_dir (str, optional): Directory to save intermediate thresholding outputs if tracking is on.
        roi_mask (np.ndarray, optional): A boolean boolean layer representing the user's manual bounding box.

    Returns:
        np.ndarray: A labeled discrete array of contiguous pixel bodies corresponding to valid single nuclei.
    """
    img = img_raw.astype(np.float32)
    if roi_mask is not None:
        roi_mask = roi_mask.astype(bool)
        if roi_mask.shape != img.shape:
            raise ValueError(f"roi_mask shape {roi_mask.shape} does not match image shape {img.shape}")
        # NOTE: do NOT zero out pixels here — that destroys image statistics
        # (CLAHE, percentiles, ridge detection).  Instead, process full image
        # and apply roi_mask only at the mask/skeleton stages below.

    # 1. Intensity Normalization
    img_norm = (img - img.min()) / (img.max() - img.min() + 1e-9)

    # 2. Adaptive Contrast Enhancement (CLAHE) - Targets local variance over global washouts
    img_eq = exposure.equalize_adapthist(
        img_norm, clip_limit=cfg["CLAHE_CLIP"], kernel_size=cfg["CLAHE_KERNEL"])

    # 3. Background Subtraction
    bg  = gaussian(img_eq, sigma=cfg["BG_SIGMA"])
    fg  = np.clip(img_eq - bg, 0, None)
    fgn = fg / (fg.max() + 1e-9)

    # 4. Neurite/Ridge Enhanced Topological Edge Detection (Meijering Filter)
    ridge = meijering(fgn, sigmas=cfg["RIDGE_SIGMAS"], black_ridges=False)

    th_hi     = np.percentile(ridge, cfg["THRESHOLD_HI"])
    th_lo     = np.percentile(ridge, cfg["THRESHOLD_LO"])
    mask_hyst = apply_hysteresis_threshold(ridge, th_lo, th_hi)
    if roi_mask is not None:
        mask_hyst &= roi_mask

    mask_clean = morphology.binary_closing(mask_hyst, morphology.disk(cfg["CLOSE_RADIUS"]))
    mask_clean = morphology.remove_small_holes(mask_clean, area_threshold=cfg["MIN_HOLE_AREA"])
    mask_clean = morphology.remove_small_objects(mask_clean, min_size=cfg["MIN_OBJ_PX"])
    mask_clean = apply_optional_early_shape_filter(mask_clean, cfg)
    if roi_mask is not None:
        mask_clean &= roi_mask

    # Width is measured from the CLEAN (un-bridged) distance map
    dist_clean   = distance_transform_edt(mask_clean)
    skel_clean   = skeletonize(mask_clean)
    skel_labeled = measure.label(skel_clean)

    # Skeleton-level bridging (preserves mask / width integrity)
    skel_bridged    = bridge_skeleton_endpoints(
        skel_clean, skel_labeled, cfg["MAX_BRIDGE_PX"])
    if roi_mask is not None:
        skel_bridged &= roi_mask
    # Branch pruning before measurement
    skel_pruned     = prune_branches(skel_bridged, cfg["MAX_BRANCH_LEN_PX"])
    
    # Optional: Automatically sever complex webs into isolated individual strands
    if cfg.get("BREAK_JUNCTIONS", False):
        from scipy.ndimage import convolve
        kernel = np.array([[1, 1, 1],
                           [1, 0, 1],
                           [1, 1, 1]], dtype=np.int32)
        skel_int = skel_pruned.astype(np.int32)
        neighbors = convolve(skel_int, kernel, mode='constant', cval=0)
        # Any skeleton pixel with more than 2 neighbors is a junction
        junctions = (skel_int > 0) & (neighbors > 2)
        skel_pruned[junctions] = 0
        
    if roi_mask is not None:
        skel_pruned &= roi_mask
    skel_labeled_fn = measure.label(skel_pruned)

    # NEW: The Recursive Adaptive Micro-Crop Reanalyzer for dense webs and chains
    if cfg.get("MAX_GEODESIC_LEN_PX", 0) > 0 and cfg.get("AUTO_LOCAL_REANALYSIS", True):
        max_px = cfg["MAX_GEODESIC_LEN_PX"]
        
        # Create a rigid dict of sub-parameters to forcefully shatter the isolated components
        sub_cfg = cfg.copy()
        sub_cfg["AUTO_LOCAL_REANALYSIS"] = False # Prevent infinite recursion
        
        # Because the bounding box crop contains almost NO dark background, 
        # a median 50% threshold perfectly isolates the brightest centers of the blobs!
        sub_cfg["THRESHOLD_HI"] = 55.0
        sub_cfg["THRESHOLD_LO"] = 45.0
        
        # Deactivate topological limits in the recursive sub-call so fragments aren't 
        # instantly dropped before they can be spliced back into the master image.
        sub_cfg["MIN_SKEL_LEN_PX"] = 1.0
        sub_cfg["MIN_OBJ_PX"] = 3
        sub_cfg["MIN_HOLE_AREA"] = 0

        props = measure.regionprops(skel_labeled_fn)
        for sp in props:
            # Note: A dense 2D web might have len(sp.coords) = 5000 pixels.
            if len(sp.coords) > max_px:
                minr, minc, maxr, maxc = sp.bbox
                pad = 12
                minr = max(0, minr - pad)
                minc = max(0, minc - pad)
                maxr = min(img.shape[0], maxr + pad)
                maxc = min(img.shape[1], maxc + pad)
                
                crop_img = img[minr:maxr, minc:maxc]
                
                # Create a strict ROI mask over the target structure to ignore neighbors
                obj_mask = (skel_labeled_fn[minr:maxr, minc:maxc] == sp.label)
                crop_roi = morphology.dilation(obj_mask, morphology.disk(6))
                
                try:
                    # RECURE: Run the entire engine on the tiny isolated crop!
                    sub_seg = segment_slice(crop_img, sub_cfg, roi_mask=crop_roi)
                    sub_skel = sub_seg["skel_pruned"]
                    sub_lab = measure.label(sub_skel)
                    
                    if sub_lab.max() > 1:
                        # SUCCESS: The adaptive threshold organically shattered the web!
                        skel_pruned[minr:maxr, minc:maxc][obj_mask] = 0
                        new_frags = (sub_skel > 0)
                        skel_pruned[minr:maxr, minc:maxc][new_frags] = 1
                    else:
                        # Failsafe: if the intensity was perfectly uniform, geometric centroid chop
                        cy, cx = sp.centroid
                        dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                        mid_idx = np.argmin(dists)
                        my, mx = sp.coords[mid_idx]
                        skel_pruned[my-1:my+2, mx-1:mx+2] = 0
                        
                except Exception:
                    # Absolute geometric failsafe
                    cy, cx = sp.centroid
                    dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                    mid_idx = np.argmin(dists)
                    my, mx = sp.coords[mid_idx]
                    skel_pruned[my-1:my+2, mx-1:mx+2] = 0

        # Run one final relabeling array refresh after all splicing
        skel_labeled_fn = measure.label(skel_pruned)

    out = {
        "mask_hyst":    mask_hyst,
        "mask_clean":   mask_clean,
        "skel_clean":   skel_clean,
        "skel_bridged": skel_bridged,
        "skel_pruned":  skel_pruned,
        "skel_labeled": skel_labeled_fn,
        "dist_clean":   dist_clean,
        "img_eq":       img_eq,
        "ridge":        ridge,
        "roi_mask":     roi_mask,
    }

    if cfg["SAVE_DEBUG_IMAGES"] and debug_dir and z_idx is not None:
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_01_norm.png"),          img_norm)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_02_clahe.png"),         img_eq)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_03_fg.png"),            fgn)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_04_ridge.png"),         ridge)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_05_hysteresis.png"),    mask_hyst)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_06_clean.png"),         mask_clean)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_07_skel_clean.png"),    skel_clean)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_08_skel_bridged.png"),  skel_bridged)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_09_skel_pruned.png"),   skel_pruned)

    return out


# =============================================================================
# MEASUREMENT  (single geodesic pass, all topology in one function)
# =============================================================================

def measure_spermatids(seg, cfg):
    """
    Analyzes mathematically discretised nuclei arrays and derives geometric indices for individual shape profiles.
    
    Args:
        seg (dict): Dictionary containing labeled binary masks ("primary", "skel_pruned", "dist_clean").
        cfg (dict): Pipeline hyperparameters for absolute scaling calculations.
        
    Returns:
        pd.DataFrame: Table indexing each isolated object dynamically with mathematical biometrics.
                      Attributes include region area, bounding-boxes, minor/major ellipse axis mapping,
                      orientation, and geodesic structural lengths via binary skeleton processing.
    """
    skel     = seg["skel_pruned"]
    dist     = seg["dist_clean"]
    skel_lab = seg["skel_labeled"]
    H, W     = skel.shape

    # ── Filter pass ───────────────────────────────────────────────────────────
    accepted_labels = []
    cache           = {}
    reasons = {"short": 0, "loop": 0, "long": 0, "wide": 0, "ratio": 0, "branches": 0, "tortuous": 0, "endpoints": 0}

    for sp in measure.regionprops(skel_lab):
        coords = sp.coords
        if coords.shape[0] < cfg["MIN_SKEL_LEN_PX"]:
            reasons["short"] += 1
            continue

        topo = measure_topology(coords, W, allow_loops=cfg.get("ALLOW_LOOPS", False))
        if topo is None:
            reasons["loop"] += 1
            continue  # loop, not allowed

        gl   = topo["geo_len"]
        tort = topo["tortuosity"]
        n_ep = topo["n_endpoints"]
        n_br = topo["n_branch_nodes"]

        if not (cfg["MIN_SKEL_LEN_PX"] <= gl <= cfg["MAX_GEODESIC_LEN_PX"]):
            if gl < cfg["MIN_SKEL_LEN_PX"]: reasons["short"] += 1
            else: reasons["long"] += 1
            continue

        width = float(np.median(2.0 * dist[coords[:, 0], coords[:, 1]]))
        if width > cfg["MAX_WIDTH_PX"]:
            reasons["wide"] += 1
            continue

        if gl / (width + 1e-9) < cfg["MIN_LENGTH_WIDTH_RATIO"]:
            reasons["ratio"] += 1
            continue

        # N1: branch-node count filter
        if n_br > cfg["MAX_BRANCH_NODES"]:
            reasons["branches"] += 1
            continue

        # N2: tortuosity filter (only for open filaments with 2+ endpoints)
        if n_ep >= 2 and tort > cfg["MAX_TORTUOSITY"]:
            reasons["tortuous"] += 1
            continue

        # N3: endpoint count filter
        if n_ep > cfg["MAX_ENDPOINT_COUNT"]:
            reasons["endpoints"] += 1
            continue

        cy, cx = sp.centroid
        accepted_labels.append(sp.label)
        cache[sp.label] = {
            "geo_len":            gl,
            "tortuosity":         tort,
            "n_endpoints":        n_ep,
            "n_branch_nodes":     n_br,
            "width":              width,
            "length_width_ratio": gl / (width + 1e-9),
            "length_px_count":    float(coords.shape[0]),
            "cx": cx, "cy": cy,
            "area_px": float(sp.area),
        }

    total_rejected = sum(reasons.values())
    if total_rejected > 0:
        print(f"    measure_spermatids rejected {total_rejected} blobs:")
        for k, v in reasons.items():
            if v > 0: print(f"      {k}: {v}")

    if not accepted_labels:
        return {"skel_label": np.zeros_like(skel_lab, dtype=np.int32),
                "results": []}

    clean_skel  = np.isin(skel_lab, accepted_labels)
    final_label = measure.label(clean_skel).astype(np.int32)

    # ── Re-index using cached values (no second Dijkstra pass) ───────────────
    final_results = []
    for new_i, sp in enumerate(measure.regionprops(final_label), start=1):
        old_label = skel_lab[sp.coords[0, 0], sp.coords[0, 1]]
        if old_label not in cache:
            continue
        c = cache[old_label]
        final_results.append({
            "label":               new_i,
            "length_px_geodesic":  c["geo_len"],
            "length_px_count":     c["length_px_count"],
            "width_px":            c["width"],
            "length_width_ratio":  c["length_width_ratio"],
            "tortuosity":          c["tortuosity"],
            "n_endpoints":         c["n_endpoints"],
            "n_branch_nodes":      c["n_branch_nodes"],
            "centroid_x":          c["cx"],
            "centroid_y":          c["cy"],
            "area_px":             c["area_px"],
        })

    return {"skel_label": final_label, "results": final_results}


# =============================================================================
# OVERLAY  (vectorized LUT)
# =============================================================================

def make_overlay(img_raw, skel_label):
    """
    Generates a colour-coded skeleton overlay on the grayscale raw image.

    Each detected spermatid is assigned a unique hue from the ``gist_rainbow``
    colourmap, dilated by 3 pixels for visibility, and composited onto the
    contrast-stretched raw image.  Background pixels (label == 0) retain the
    original grayscale intensity.

    Implementation notes
    --------------------
    - Vectorised LUT (Look-Up Table) approach avoids per-label loop for speed.
    - ``grey_dilation`` makes thin single-pixel skeletons visible at any zoom.
    - Colour assignment is deterministic: same label ordering → same colour.

    Args:
        img_raw (np.ndarray): Raw microscopy image (any dtype).
        skel_label (np.ndarray[int]): Integer-labelled skeleton array
            (0 = background, 1..N = individual spermatids).

    Returns:
        np.ndarray: uint8 RGB image, shape ``(H, W, 3)``, ready for ``plt.imshow``
        or saving with :func:`_imwrite`.
    """
    base = normalize_display(img_raw)
    n    = int(skel_label.max())
    if n <= 0:
        # No detections: return grayscale image as RGB
        return (np.stack([base]*3, -1) * 255).astype(np.uint8)
    # Assign one colour per label; prepend black for background (index 0)
    cols    = plt.cm.gist_rainbow(np.linspace(0, 1, n))[:, :3]
    dilated = grey_dilation(skel_label.astype(np.int32), size=3)
    lut     = np.vstack([[0., 0., 0.], cols[:n]])
    rgb     = lut[dilated]
    # Restore original grayscale for background pixels
    m0      = dilated == 0
    rgb[m0, 0] = base[m0]
    rgb[m0, 1] = base[m0]
    rgb[m0, 2] = base[m0]
    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)


# =============================================================================
# DISPLAY / SAVE
# =============================================================================

def _safe_show():
    """
    Calls ``plt.show(block=True)`` with a try/except guard.

    On headless systems (Linux CI servers, SSH sessions without X11 forwarding,
    or Windows without a display) matplotlib will raise a ``RuntimeError`` or
    ``_tkinter.TclError`` when trying to open an interactive window.  This
    wrapper silently swallows that error so the pipeline can continue and saves
    the image to disk as a fallback.
    """
    try:
        plt.show(block=True)
    except Exception as e:
        print(f"[WARNING] plt.show() failed: {e}")


def show_single_preview(img_raw, seg, overlay_rgb, results, z_idx, cfg):
    um         = cfg["UM_PER_PX_XY"]
    lengths_um = [r["length_px_geodesic"] * um for r in results]

    nrows, ncols = (2, 4) if cfg["SHOW_DEBUG_PREVIEW"] else (1, 3)
    fig, axes = plt.subplots(nrows, ncols, figsize=(18, 5*nrows))

    def _ax(r, c):
        return axes[r, c] if nrows > 1 else axes[c]

    _ax(0,0).imshow(normalize_display(img_raw), cmap="gray")
    _ax(0,0).set_title(f"Original Z={z_idx:02d}"); _ax(0,0).axis("off")

    if cfg["SHOW_DEBUG_PREVIEW"]:
        _ax(0,1).imshow(seg["img_eq"], cmap="gray")
        _ax(0,1).set_title("CLAHE"); _ax(0,1).axis("off")
        _ax(0,2).imshow(seg["ridge"], cmap="gray")
        _ax(0,2).set_title("Ridge"); _ax(0,2).axis("off")
        _ax(0,3).imshow(seg["mask_hyst"], cmap="gray")
        _ax(0,3).set_title("Hysteresis"); _ax(0,3).axis("off")
        _ax(1,0).imshow(seg["mask_clean"], cmap="gray")
        _ax(1,0).set_title("Clean mask"); _ax(1,0).axis("off")
        _ax(1,1).imshow(seg["skel_pruned"], cmap="gray")
        _ax(1,1).set_title("Skeleton (pruned)"); _ax(1,1).axis("off")
        ov_ax   = _ax(1,2)
        hist_ax = _ax(1,3)
    else:
        ov_ax   = _ax(0,1)
        hist_ax = _ax(0,2)

    ov_ax.imshow(overlay_rgb)
    ov_ax.set_title(f"Overlay N={len(results)}"); ov_ax.axis("off")
    for r in results:
        ov_ax.text(r["centroid_x"], r["centroid_y"],
                   f"{r['length_px_geodesic']*um:.1f}",
                   color="white", fontsize=5, ha="center", va="center")

    if lengths_um:
        hist_ax.hist(lengths_um, bins=20, edgecolor="white")
        hist_ax.axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} µm")
        hist_ax.set_xlabel("Geodesic length (µm)"); hist_ax.legend(fontsize=8)
        hist_ax.set_title("Length distribution")
    else:
        hist_ax.text(0.5, 0.5, "No detections", ha="center", va="center")
        hist_ax.axis("off")

    plt.tight_layout()

    # Save preview to file and open with default viewer
    preview_path = os.path.join(cfg["OUTPUT_DIR"], "preview.png")
    ensure_dir(cfg["OUTPUT_DIR"])
    plt.savefig(preview_path, dpi=120, bbox_inches="tight")
    print(f"  Preview saved to: {preview_path}")

    # Try interactive display
    _safe_show()
    plt.close()

    # Open with default viewer as fallback
    try:
        os.startfile(preview_path)
    except Exception:
        pass


def save_detail_figure(img_raw, overlay_rgb, results, out_path, z_idx, um):
    lengths_um = [r["length_px_geodesic"] * um for r in results]
    fig, axes  = plt.subplots(1, 3, figsize=(18, 6))

    axes[0].imshow(normalize_display(img_raw), cmap="gray")
    axes[0].set_title(f"Z={z_idx:02d} — Original"); axes[0].axis("off")

    axes[1].imshow(overlay_rgb)
    axes[1].set_title(f"Z={z_idx:02d} — Spermatids (N={len(results)})")
    axes[1].axis("off")
    for r in results:
        axes[1].text(r["centroid_x"], r["centroid_y"],
                     f"{r['length_px_geodesic']*um:.1f}",
                     color="white", fontsize=4, ha="center", va="center",
                     bbox=dict(boxstyle="round,pad=0.1", fc="black", alpha=0.4, lw=0))

    if lengths_um:
        axes[2].hist(lengths_um, bins=20, edgecolor="white")
        axes[2].axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} µm")
        axes[2].set_xlabel("Geodesic length (µm)"); axes[2].set_ylabel("Count")
        axes[2].set_title(f"Z={z_idx:02d} — Length distribution")
        axes[2].legend(fontsize=9)
    else:
        axes[2].text(0.5, 0.5, "No spermatids detected",
                     transform=axes[2].transAxes, ha="center", va="center")
        axes[2].axis("off")

    plt.tight_layout()
    plt.savefig(out_path, dpi=120, bbox_inches="tight")
    plt.close()


# =============================================================================
# CSV
# =============================================================================

_VERSION = "v11"


def rows_from_results(results, z_idx, um):
    """
    Converts per-spermatid measurement dictionaries into flat CSV row dictionaries.

    Applies pixel-to-micron scaling (``UM_PER_PX_XY``) for all linear dimensions
    and rounds floating-point values to 3 decimal places for clean spreadsheet output.

    Columns emitted per detection
    ------------------------------
    - ``pipeline_version``     — Code version tag for traceability.
    - ``z_slice``              — Z-plane index of this detection.
    - ``sperm_id``             — 1-based integer ID within this slice.
    - ``length_px_geodesic``   — Geodesic skeleton length in pixels.
    - ``length_um_geodesic``   — Geodesic length in micrometres.
    - ``length_px_count``      — Pixel-count skeleton length (alternative measure).
    - ``length_um_count``      — Pixel-count length in micrometres.
    - ``width_px`` / ``width_um`` — Mean width across the skeleton.
    - ``length_width_ratio``   — Elongation ratio (key morphological filter metric).
    - ``tortuosity``           — Curvature index: geodesic / Euclidean tip-to-tip.
    - ``n_endpoints``          — Number of skeleton endpoints (expect 2 for clean cells).
    - ``n_branch_nodes``       — Number of branch points (3+ neighbours); should be low.
    - ``centroid_x`` / ``centroid_y`` — Pixel centroid for overlay annotation.
    - ``area_px``              — Mask area in pixels^2 (used for volume estimation).

    Args:
        results (list[dict]): Output of :func:`measure_spermatids`.
        z_idx (int): Z-slice index.
        um (float): Microns-per-pixel (``UM_PER_PX_XY``).

    Returns:
        list[dict]: One flat dictionary per detected spermatid.
    """
    return [{
        "pipeline_version":    _VERSION,
        "z_slice":             z_idx,
        "sperm_id":            i,
        "length_px_geodesic":  round(r["length_px_geodesic"], 3),
        "length_um_geodesic":  round(r["length_px_geodesic"] * um, 3),
        "length_px_count":     round(r["length_px_count"], 1),
        "length_um_count":     round(r["length_px_count"]  * um, 3),
        "width_px":            round(r["width_px"], 2),
        "width_um":            round(r["width_px"]          * um, 3),
        "length_width_ratio":  round(r["length_width_ratio"], 3),
        "tortuosity":          round(r["tortuosity"], 3),
        "n_endpoints":         r["n_endpoints"],
        "n_branch_nodes":      r["n_branch_nodes"],
        "centroid_x":          round(r["centroid_x"], 1),
        "centroid_y":          round(r["centroid_y"], 1),
        "area_px":             round(r["area_px"], 1),
    } for i, r in enumerate(results, start=1)]


# =============================================================================
# TRACKING
# =============================================================================

def track_across_slices(detections_df, cfg):
    """
    Links separate discrete mathematically proven 2D slices into unified 3D physical continuous objects.
    Constructs accurate physical geometric bounding matrices and extrapolates z-plane depths across continuous Z-slices.
    
    Algorithm:
        DBSCAN density clustering determines neighbor linkage across Z distances, mapped mathematically based
        on XY drift coefficients. 
    
    Args:
        detections_df (pd.DataFrame): Target frame containing absolute xy indices mapped to variable slice limits.
        cfg (dict): Master parameters including thresholding spatial mappings and slice distances.
        
    Returns:
        pd.DataFrame: Finalized unified track summaries dictating 3D metrics (Volume, True Length, Tortuosity, Pitch)
    """
    if detections_df.empty:
        detections_df = detections_df.copy()
        detections_df["track_id"] = pd.Series(dtype=int)
        return detections_df, pd.DataFrame()

    max_dist_px = cfg["TRACK_MAX_DIST_UM"] / (cfg["UM_PER_PX_XY"] + 1e-9)
    df = (detections_df.copy()
                       .sort_values(["z_slice","sperm_id"])
                       .reset_index(drop=True))

    next_tid  = 1
    active    = {}
    track_ids = [-1] * len(df)
    rows_by_z = {z: df.index[df["z_slice"] == z].to_numpy()
                 for z in sorted(df["z_slice"].unique())}

    for z, idxs in rows_by_z.items():
        xs = df.loc[idxs, "centroid_x"].to_numpy(float)
        ys = df.loc[idxs, "centroid_y"].to_numpy(float)

        cand_tracks = [t for t, st in active.items()
                       if 1 <= z - st["last_z"] <= cfg["TRACK_MAX_GAP_SLICES"] + 1]
        cand_pos    = [(active[t]["last_x"], active[t]["last_y"]) for t in cand_tracks]

        used_det, used_trk = set(), set()
        if cand_tracks:
            tree = cKDTree(np.array(cand_pos, float))
            candidates = []
            for k, (x, y) in enumerate(zip(xs, ys)):
                d_val, j = tree.query([x, y])
                if np.isfinite(d_val) and d_val <= max_dist_px:
                    candidates.append((float(d_val), k, int(j)))
            for d_val, det_k, trk_j in sorted(candidates):
                if det_k in used_det or trk_j in used_trk:
                    continue
                used_det.add(det_k); used_trk.add(trk_j)
                tid = cand_tracks[trk_j]
                track_ids[int(idxs[det_k])] = tid
                active[tid] = {"last_z": int(z), "last_x": float(xs[det_k]),
                               "last_y": float(ys[det_k])}

        for det_k in range(len(idxs)):
            if track_ids[int(idxs[det_k])] == -1:
                track_ids[int(idxs[det_k])] = next_tid
                active[next_tid] = {"last_z": int(z), "last_x": float(xs[det_k]),
                                    "last_y": float(ys[det_k])}
                next_tid += 1

        for tid in [t for t, st in active.items()
                    if z - st["last_z"] > cfg["TRACK_MAX_GAP_SLICES"] + 1]:
            del active[tid]

    df["track_id"] = track_ids
    
    # Inject maximum 2D Euclidean distance of the physical shape prior to grouping
    if "tortuosity" in df.columns:
        df["euc_um_2d"] = df["length_um_geodesic"] / df["tortuosity"]
    else:
        df["euc_um_2d"] = df["length_um_geodesic"]
        
    g = df.groupby("track_id", as_index=False)
    ts = g.agg(
        n_slices        = ("z_slice",            "count"),
        z_start         = ("z_slice",            "min"),
        z_end           = ("z_slice",            "max"),
        max_length_2d   = ("length_um_geodesic", "max"),
        max_euc_2d      = ("euc_um_2d",          "max"),
        sum_area_px     = ("area_px",            "sum"),
        area_start      = ("area_px",            "first"),
        area_end        = ("area_px",            "last"),
        x_mean          = ("centroid_x",         "mean"),
        y_mean          = ("centroid_y",         "mean"),
        z_mean          = ("z_slice",            "mean"),
        # Centroids of the start/end for bounding verification (optional)
        x_start         = ("centroid_x",         "first"),
        y_start         = ("centroid_y",         "first"),
        x_end           = ("centroid_x",         "last"),
        y_end           = ("centroid_y",         "last"),
    )
    
    um_xy = cfg["UM_PER_PX_XY"]
    um_z  = cfg["UM_PER_SLICE_Z"]
    
    # 1. Z-Extent (Vertical Span includes +1 for slice thickness, matching PDF)
    z_extent = (ts["z_end"] - ts["z_start"] + 1) * um_z
    ts["z_extent_um"] = z_extent
    
    # Geodesic vertical displacement for Euclidean distance is purely centroid-to-centroid
    dz_euc = (ts["z_end"] - ts["z_start"]) * um_z
    
    # 2. Total 3D Length (Lateral-Corrected Hypotenuse)
    # The physical 3D shape arc relies on the maximum length of its 2D projection
    euc_2d_centroid = np.sqrt((ts["x_end"] - ts["x_start"])**2 + (ts["y_end"] - ts["y_start"])**2) * um_xy
    lat_geodesic = np.maximum(ts["max_length_2d"], euc_2d_centroid)
    l3d = np.sqrt(lat_geodesic**2 + z_extent**2)
    ts["total_3d_length_um"] = l3d
    
    # 3. 3D Volume (Sum Area * Z_step)
    ts["volume_um3"] = ts["sum_area_px"] * (um_xy**2) * um_z
    
    # 4. 3D Tortuosity (Total 3D Geodesic Length / 3D End-To-End Euclidean Distance)
    # The true 3D shape euclidean distance uses the maximum 2D shape displacement, 
    # not the tracking centroid displacement!
    euc_3d = np.sqrt(ts["max_euc_2d"]**2 + dz_euc**2)
    
    # Numerical safety
    safe_euc = np.maximum(euc_3d, 0.1)
    tort_raw = l3d / safe_euc
    ts["tortuosity_3d"] = np.minimum(tort_raw, 20.0)
    
    # 5. Taper Ratio (Max Area / Min Area between first and last z-slice)
    ts["taper_ratio"] = np.maximum(ts["area_start"], ts["area_end"]) / np.maximum(np.minimum(ts["area_start"], ts["area_end"]), 0.001)
    
    # 6. Effective Thickness / Diameter
    cross_area = ts["volume_um3"] / np.maximum(ts["total_3d_length_um"], 0.1)
    ts["thickness_um"] = 2 * np.sqrt(cross_area / np.pi)
    
    # 7. Orientation Angles (Pitch and Yaw)
    dx = (ts["x_end"] - ts["x_start"]) * um_xy
    dy = (ts["y_end"] - ts["y_start"]) * um_xy
    v_mag = np.sqrt(dx**2 + dy**2 + dz_euc**2)
    safe_v = np.maximum(v_mag, 1e-9)
    ts["pitch_deg"] = np.abs(np.arcsin(dz_euc / safe_v)) * (180.0 / np.pi)
    ts["yaw_deg"] = np.arctan2(dy, dx) * (180.0 / np.pi)
    
    # 8. Spatial Packing Density (Nearest Neighbor Distance in 3D)
    if len(ts) > 1:
        centers = np.column_stack((ts["x_mean"] * um_xy, ts["y_mean"] * um_xy, ts["z_mean"] * um_z))
        tree = cKDTree(centers)
        dists, _ = tree.query(centers, k=2)
        ts["nearest_neighbor_um"] = dists[:, 1]
    else:
        ts["nearest_neighbor_um"] = np.nan
        
    # FINAL COLUMN ORDERING (Extended)
    cols_ordered = [
        "track_id", "total_3d_length_um", "z_extent_um", "volume_um3", "tortuosity_3d",
        "thickness_um", "pitch_deg", "yaw_deg", "taper_ratio", "nearest_neighbor_um",
        "n_slices", "z_start", "z_end", "max_length_2d", "sum_area_px"
    ]
    ts = ts[cols_ordered]
    
    return df, ts


# =============================================================================
# PROCESS ONE IMAGE
# =============================================================================

def process_one_image(image_path, cfg, output_dir):
    """
    Runs the complete segmentation-and-measurement pipeline on a single Z-slice image.

    This function is the fundamental unit of the batch processing loop.  For each
    input image it:

    1. Loads the raw image (multi-format with TIFF / OpenCV fallback).
    2. Selects the correct Z-slice from a multi-plane volume if needed.
    3. Calls :func:`segment_slice` to produce the binary mask, ridge map, and skeleton.
    4. Calls :func:`measure_spermatids` to extract per-cell biometrics.
    5. Optionally validates or post-filters results with shape-quality checks.
    6. Saves the colour overlay PNG, detail figure, mask TIFs, and per-slice CSV.
    7. Optionally opens the preview window for interactive QC.

    Args:
        image_path (str): Absolute path to the input image file.
        cfg (dict): Full pipeline configuration dictionary.  Key fields:
            - ``UM_PER_PX_XY`` — physical scale factor.
            - ``Z_INDEX`` — which plane to extract from a Z-stack.
            - ``OUTPUT_DIR`` — top-level output folder.
            - ``SHOW_PREVIEW_WINDOW`` — whether to open the preview GUI.
            - ``SAVE_MASK_TIFS``, ``SAVE_LABEL_TIFS`` — optional TIF outputs.
            - ``SAVE_DETAIL_FIGURE`` — whether to save 3-panel figure.
        output_dir (str): Destination directory for this image's outputs.

    Returns:
        tuple[list[dict], dict]:
            - ``results`` — list of per-spermatid measurement dicts.
            - ``seg``     — full segmentation dict from :func:`segment_slice`
              including ``mask_clean``, ``skel_pruned``, ``skel_label``, etc.
    """
    ensure_dir(output_dir)
    overlay_dir = os.path.join(output_dir, "overlays")
    debug_dir   = os.path.join(output_dir, "debug")
    ensure_dir(overlay_dir)
    if cfg["SAVE_DEBUG_IMAGES"]:
        ensure_dir(debug_dir)

    z_idx   = extract_z_index(image_path)
    img_raw = robust_imread(image_path)
    print(f"\nProcessing: {os.path.basename(image_path)}")

    t0      = time.time()
    seg     = segment_slice(img_raw, cfg, z_idx=z_idx,
                            debug_dir=debug_dir if cfg["SAVE_DEBUG_IMAGES"] else None)
    meas    = measure_spermatids(seg, cfg)
    results = meas["results"]
    elapsed = time.time() - t0

    um = cfg["UM_PER_PX_XY"]
    print(f"  Detected: {len(results)} spermatids  ({elapsed:.1f}s)")
    if results:
        ls = [r["length_px_geodesic"]*um for r in results]
        print(f"  Geodesic length µm: median={np.median(ls):.2f}  max={max(ls):.2f}")

    overlay_rgb = make_overlay(img_raw, meas["skel_label"])

    if cfg["SAVE_OVERLAYS"]:
        _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_overlay.png"), overlay_rgb)
    if cfg["SAVE_DETAIL_FIGURE"]:
        save_detail_figure(img_raw, overlay_rgb, results,
                           os.path.join(overlay_dir, f"z{z_idx:02d}_detail.png"),
                           z_idx, um)
    if cfg["SAVE_MASK_TIFS"]:
        tifffile.imwrite(os.path.join(output_dir, f"z{z_idx:02d}_mask.tif"),
                         seg["mask_clean"].astype(np.uint8) * 255)
    if cfg["SAVE_LABEL_TIFS"]:
        tifffile.imwrite(os.path.join(output_dir, f"z{z_idx:02d}_skel_labels.tif"),
                         meas["skel_label"].astype(np.uint16))

    pd.DataFrame(rows_from_results(results, z_idx, um)).to_csv(
        os.path.join(output_dir, f"single_measurements_{_VERSION}.csv"), index=False)

    if cfg["SHOW_PREVIEW_WINDOW"]:
        show_single_preview(img_raw, seg, overlay_rgb, results, z_idx, cfg)

    print(f"Saved to: {output_dir}")


# =============================================================================
# PROCESS BATCH
# =============================================================================

def process_batch(cfg):
    """
    Orchestrates the entire end-to-end biological data processing engine for batch image iterations.
    
    Iterates over all `.tif`/`.tiff` files across defined spatial bounds, extracting single
    segmentation matrices and triggering the 3D concatenation models upon directory exhaustion.
    
    Args:
        cfg (dict): Active session configuration parameter dictionary.
    """
    ensure_dir(cfg["OUTPUT_DIR"])
    overlay_dir = os.path.join(cfg["OUTPUT_DIR"], "overlays")
    debug_dir   = os.path.join(cfg["OUTPUT_DIR"], "debug")
    ensure_dir(overlay_dir)
    if cfg["SAVE_DEBUG_IMAGES"]:
        ensure_dir(debug_dir)

    files, z_indices = load_batch_files(cfg["INPUT_DIR"], cfg["FILE_PATTERN"])
    um         = cfg["UM_PER_PX_XY"]
    all_rows   = []
    summaries  = []
    t_batch    = time.time()
    t_slices   = []

    for idx_i, (fpath, z_idx) in enumerate(zip(files, z_indices)):
        t0 = time.time()
        print(f"\n[{idx_i+1}/{len(files)}]  Z={z_idx:02d}  {os.path.basename(fpath)}")
        img_raw = robust_imread(fpath)
        seg     = segment_slice(img_raw, cfg, z_idx=z_idx,
                                debug_dir=debug_dir if cfg["SAVE_DEBUG_IMAGES"] else None)
        meas    = measure_spermatids(seg, cfg)
        results = meas["results"]
        ls_um   = [r["length_px_geodesic"]*um for r in results]
        ws_um   = [r["width_px"]*um for r in results]

        t_s = time.time() - t0
        t_slices.append(t_s)
        eta = (len(files) - idx_i - 1) * float(np.mean(t_slices))
        print(f"  N={len(results)}", end="")
        if ls_um:
            print(f"  med_len={np.median(ls_um):.2f}µm", end="")
        print(f"  {t_s:.1f}s  ETA {eta:.0f}s")

        all_rows.extend(rows_from_results(results, z_idx, um))
        summaries.append({
            "z_slice":          z_idx,
            "n_spermatids":     len(results),
            "mean_length_um":   round(float(np.mean(ls_um)),   3) if ls_um else 0,
            "median_length_um": round(float(np.median(ls_um)), 3) if ls_um else 0,
            "mean_width_um":    round(float(np.mean(ws_um)),   3) if ws_um  else 0,
        })

        overlay_rgb = make_overlay(img_raw, meas["skel_label"])
        if cfg["SAVE_OVERLAYS"]:
            # Create side-by-side panel: [Original | Overlay]
            orig_rgb = (normalize_display(img_raw) * 255).astype(np.uint8)
            # if grayscale, convert to RGB for hstack
            if orig_rgb.ndim == 2:
                orig_rgb = np.stack([orig_rgb]*3, axis=-1)
            
            panel = np.hstack([orig_rgb, overlay_rgb])
            _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_panel.png"), panel)

        if cfg["SAVE_DETAIL_FIGURE"]:
            save_detail_figure(img_raw, overlay_rgb, results,
                               os.path.join(overlay_dir, f"z{z_idx:02d}_detail.png"),
                               z_idx, um)
        if cfg["SAVE_MASK_TIFS"]:
            tifffile.imwrite(os.path.join(cfg["OUTPUT_DIR"], f"z{z_idx:02d}_mask.tif"),
                             seg["mask_clean"].astype(np.uint8) * 255)
        if cfg["SAVE_LABEL_TIFS"]:
            tifffile.imwrite(os.path.join(cfg["OUTPUT_DIR"], f"z{z_idx:02d}_skel_labels.tif"),
                             meas["skel_label"].astype(np.uint16))

    df     = pd.DataFrame(all_rows)
    df_sum = pd.DataFrame(summaries)
    df.to_csv(    os.path.join(cfg["OUTPUT_DIR"], f"spermatid_measurements_{_VERSION}.csv"), index=False)
    df_sum.to_csv(os.path.join(cfg["OUTPUT_DIR"], f"slice_summary_{_VERSION}.csv"), index=False)

    # Robust initialization for reporting
    df_trk = None
    ts = None

    if cfg["DO_TRACKING"] and not df.empty:
        df_trk, ts = track_across_slices(df, cfg)
        df_trk.to_csv(
            os.path.join(cfg["OUTPUT_DIR"], f"measurements_with_tracks_{_VERSION}.csv"),
            index=False)
        ts.to_csv(
            os.path.join(cfg["OUTPUT_DIR"], f"track_summary_{_VERSION}.csv"),
            index=False)
    
    # --- Reporting Phase (CLI/Batch) ---
    print(f"\nGenerating final reports in {cfg['OUTPUT_DIR']}...")
    generate_batch_report(cfg["OUTPUT_DIR"], df, df_sum, um, ts)
    generate_excel_report(cfg["OUTPUT_DIR"], df, df_sum, ts)

    total = time.time() - t_batch
    print(f"\n{'='*55}")
    print(f"v11 DONE | {len(files)} slices | {total:.1f}s")
    print(f"Saved to: {cfg['OUTPUT_DIR']}")
    print(df_sum.to_string(index=False))


def write_error_log(out_dir, component, message):
    """
    Writes a persistent error log to report_generation_errors.txt in the output directory.
    """
    try:
        import os as _os
        from datetime import datetime as _dt
        log_path = _os.path.join(out_dir, "report_generation_errors.txt")
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"\n[{_dt.now().strftime('%Y-%m-%d %H:%M:%S')}] COMPONENT: {component}\n")
            f.write(f"MESSAGE:\n{message}\n")
            f.write("-" * 80 + "\n")
    except Exception:
        pass


def generate_excel_report(out_dir, df, df_summary, df_tracks=None):
    """
    Generates a multi-tab Excel workbook with formatted data, summary statistics,
    embedded chart images, and source-data hyperlinks.

    Workbook structure
    ------------------
    - **Batch_Summary** — one row per Z-slice with detection counts, mean/median
      length, and total area.  Includes an embedded histogram image and conditional
      formatting for high/low detection slices.
    - **3D_Morphometrics** — track-level 3D metrics exported from
      :func:`track_across_slices`, with an embedded 3D length distribution plot.
    - **Raw_2D_Detections** — full per-spermatid measurement table matching the CSV
      export, with number formatting and frozen header pane.
    - **Statistics_Summary** — descriptive statistics (mean, median, std, IQR,
      percentiles) for the primary biometric columns.

    Biological interpretation
    -------------------------
    The statistical summary sheet is designed to be directly copy-pasteable into
    lab reports.  The IQR (Interquartile Range = Q75 - Q25) is reported because
    spermatid length distributions are often right-skewed and IQR is more robust
    than standard deviation in these cases.

    Args:
        out_dir (str): Top-level analysis output directory.  The workbook is saved
            as ``<out_dir>/batch_analysis_results_<ver>.xlsx``.
        df (pd.DataFrame): Per-spermatid measurement table (all Z-slices combined).
        df_summary (pd.DataFrame): Per-slice summary statistics.
        df_tracks (pd.DataFrame, optional): 3D track table from
            :func:`track_across_slices`.  ``None`` if tracking was not run.
    """
    excel_path = os.path.join(out_dir, f"batch_analysis_results_{_VERSION}.xlsx")
    plot_dir = os.path.join(out_dir, "summary_plots")
    print(f"Generating Interactive Excel Audit: {excel_path} ...")
    
    try:
        with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
            workbook  = writer.book
            
            # Formats
            bold = workbook.add_format({'bold': True, 'bg_color': '#D7E4BC', 'border': 1})
            num_fmt = workbook.add_format({'num_format': '0.00'})

            # --- Sheet 1: Population Summary (DYNAMIC FORMULAS) --- [FIRST TAB]
            ws_sum = workbook.add_worksheet('Population_Summary')
            headers = ["Metric", "Average (Formula)", "Median (Formula)", "Std Dev (Formula)"]
            for col, h in enumerate(headers):
                ws_sum.write(0, col, h, bold)
            
            row = 1
            # 2D Length (from Raw_2D_Detections column G - length_um_geodesic)
            if not df.empty:
                n_2d = len(df) + 1
                ws_sum.write(row, 0, "2D Geodesic Length (um)")
                ws_sum.write_formula(row, 1, f"=AVERAGE('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                ws_sum.write_formula(row, 2, f"=MEDIAN('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                ws_sum.write_formula(row, 3, f"=STDEV.P('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                row += 1
            
            # 3D Metrics
            if df_tracks is not None and not df_tracks.empty:
                n_3d = len(df_tracks) + 1
                metrics_3d = [
                    ("3D Geodesic Length (um)", "B"),
                    ("3D Z-Extent (um)", "C"),
                    ("3D Volume (um3)", "D"),
                    ("3D Tortuosity", "E")
                ]
                for m_name, col_letter in metrics_3d:
                    ws_sum.write(row, 0, m_name)
                    ws_sum.write_formula(row, 1, f"=AVERAGE('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 2, f"=MEDIAN('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 3, f"=STDEV.P('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    row += 1

            ws_sum.set_column('A:A', 30)

            # --- Sheet 2: 3D Morphometrics ---
            if df_tracks is not None and not df_tracks.empty:
                df_tracks.to_excel(writer, sheet_name='3D_Morphometrics', index=False)
                ws_3d = writer.sheets['3D_Morphometrics']
                ws_3d.set_column('A:Z', 15)
                # Insert 3D Distribution Graph
                p_3d = os.path.join(plot_dir, "3d_length_distribution.png")
                if os.path.exists(p_3d):
                    ws_3d.insert_image('K2', p_3d, {'x_scale': 0.4, 'y_scale': 0.4})

            # --- Sheet 3: Raw 2D Detections ---
            if not df.empty:
                df.to_excel(writer, sheet_name='Raw_2D_Detections', index=False)
                ws_2d = writer.sheets['Raw_2D_Detections']
                ws_2d.set_column('A:Z', 15)

            # --- Sheet 4: Slice Summary ---
            if not df_summary.empty:
                df_summary.to_excel(writer, sheet_name='Slice_Summary', index=False)

            ws_sum.set_column('B:D', 18)

            # Insert Histograms into Summary
            p_hist = os.path.join(plot_dir, "global_histograms.png")
            p_slice = os.path.join(plot_dir, "length_by_slice.png")
            if os.path.exists(p_hist):
                ws_sum.insert_image('F2', p_hist, {'x_scale': 0.5, 'y_scale': 0.5})
            if os.path.exists(p_slice):
                ws_sum.insert_image('F25', p_slice, {'x_scale': 0.5, 'y_scale': 0.5})

            # --- Sheet 5: Methods Dictionary ---
            dictionary_data = [
                ["Metric", "Formula / Definition", "Biological Interpretation"],
                ["2D Geodesic Length", "Dijkstra path length on skeleton", "The curved length of a spermatid fragment within a single 2D slice."],
                ["Total 3D Geodesic Length", "sqrt( max(L_2d, L_drift)^2 + Z_extent^2 )", "Robust 3D reconstruction accounting for lateral drift across slices."],
                ["3D Tortuosity", "3D Geodesic Length / 3D Euclidean distance", "Curvature Index. 1.0 = Perfectly straight."],
                ["3D Euclidean Distance", "sqrt( dx^2 + dy^2 + dz^2 )", "Straight line between centroids."],
                ["Z-Extent (Vertical Span)", "Vertical Span (slices * um/slice)", "Total vertical depth covered."],
                ["3D Volume (um3)", "sum( Area_i * Z_step )", "Approximated Volumetric Mass."],
                ["Standard Deviation", "Std Dev", "Population variance."]
            ]
            pd.DataFrame(dictionary_data[1:], columns=dictionary_data[0]).to_excel(writer, sheet_name='Methods_Dictionary', index=False)
            ws_dict = writer.sheets['Methods_Dictionary']
            
            # Manually write the dictionary to avoid any encoding or header issues
            for r_idx, r_data in enumerate(dictionary_data):
                for c_idx, val in enumerate(r_data):
                    ws_dict.write(r_idx, c_idx, val, bold if r_idx==0 else None)
            
            ws_dict.set_column('A:A', 25)
            ws_dict.set_column('B:B', 40)
            ws_dict.set_column('C:C', 60)
            p_guide = os.path.join(plot_dir, "methods_guide.png")
            if os.path.exists(p_guide):
                ws_dict.insert_image('A12', p_guide, {'x_scale': 0.6, 'y_scale': 0.6})

            print(f"Interactive Excel report successfully saved to {excel_path}")

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating Excel report: {e}")
        write_error_log(out_dir, "Excel Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"Excel Report failed to generate completely.\n{e}")
        except Exception:
            pass



def generate_batch_report(out_dir, df, df_summary, um, df_tracks=None, gui_callback=None):
    """
    Compiles standard summary global output architectures including histograms, mathematical summaries,
    biological methodology pages, and graphical slice overlays natively to a `.pdf` file.
    
    Args:
        out_dir (str): Root export directory.
        df (pd.DataFrame): Flat 2D analysis parameters.
        df_summary (pd.DataFrame): Top-level slice aggregation tracking.
        um (dict): User-defined pixel to micron mapping ratios.
        df_tracks (pd.DataFrame, optional): Unified 3D tracking geometries array.
        gui_callback (function, optional): Live variable passing to front-end dashboards elements.
        
    Returns:
        None (Saves directly to absolute path PDF)
    """
    pdf_path = os.path.join(out_dir, f"batch_report_{_VERSION}.pdf")
    print(f"Generating high-res PDF report: {pdf_path} ...")
    
    # Create directory for high-res standalone plots (for easy copy-pasting into papers/presentations)
    plot_dir = os.path.join(out_dir, "summary_plots")
    os.makedirs(plot_dir, exist_ok=True)
    
    try:
        with PdfPages(pdf_path) as pdf:
            # --- PAGE 1: GLOBAL SUMMARY ---
            fig_sum = plt.figure(figsize=(11, 8.5))
            fig_sum.suptitle(f"Spermatid Analysis Batch Summary - {_VERSION}\nLocation: {out_dir}", fontsize=14, fontweight='bold')
            
            # Global Z-Projection Image (Top Center)
            z_proj_path = os.path.join(out_dir, "global_z_projection.png")
            if os.path.exists(z_proj_path):
                ax_z = fig_sum.add_axes([0.15, 0.62, 0.7, 0.28]) # [left, bottom, width, height]
                ax_z.imshow(plt.imread(z_proj_path))
                ax_z.set_title("Global Z-Projection (Composite [Original | Overlay])", fontsize=10)
                ax_z.axis('off')
            
            # Plot 1: Counts per slice
            ax1 = fig_sum.add_subplot(2, 2, 3)
            ax1.plot(df_summary['z_slice'], df_summary['n_spermatids'], 'bo-', markersize=4)
            ax1.set_title("Detections per Z-Slice")
            ax1.set_xlabel("Z-Index")
            ax1.set_ylabel("Count")
            ax1.grid(True, alpha=0.3)
            
            # Plot 2: Length Distribution (Global)
            ax2 = fig_sum.add_subplot(2, 2, 4)
            if not df.empty:
                vals = df['length_um_geodesic']
                ax2.hist(vals, bins=25, color='forestgreen', edgecolor='black', alpha=0.7)
                m_med = vals.median()
                m_avg = vals.mean()
                ax2.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax2.axvline(m_avg, color='orange', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax2.set_title("Global 2D Length Distribution")
                ax2.set_xlabel("Geodesic Length (um)")
                ax2.set_ylabel("Frequency")
                ax2.legend(fontsize=8)
            
            fig_sum.savefig(os.path.join(plot_dir, "global_summary.png"), dpi=300, bbox_inches='tight')
            # NEW: Explicitly save global_histograms.png for Excel embedding
            fig_sum.savefig(os.path.join(plot_dir, "global_histograms.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_sum, dpi=300, bbox_inches='tight')
            plt.close(fig_sum)
            
            # Save length_by_slice.png
            fig_l_slice = plt.figure(figsize=(10, 5))
            ax_ls = fig_l_slice.add_subplot(1, 1, 1)
            ax_ls.plot(df_summary['z_slice'], df_summary['median_length_um'], 'go-', label='Median Length')
            ax_ls.set_title("Median Length by Slice")
            ax_ls.set_xlabel("Z-Slice")
            ax_ls.set_ylabel("Length (um)")
            ax_ls.grid(True, alpha=0.3)
            fig_l_slice.savefig(os.path.join(plot_dir, "length_by_slice.png"), dpi=300, bbox_inches='tight')
            plt.close(fig_l_slice)

            # --- PAGE 1.5: POPULATION CONSOLIDATION ---
            if df_tracks is not None and not df_tracks.empty:
                fig_dyn = plt.figure(figsize=(11, 8.5))
                fig_dyn.suptitle("3D Population Tracking & Consolidation", fontsize=15, fontweight='bold')
                
                total_2d = len(df)
                total_3d = len(df_tracks)
                
                # A) Side-by-side Horizontal Bar
                ax_bar = fig_dyn.add_subplot(1, 2, 1)
                y_pos = [1, 0]
                counts = [total_2d, total_3d]
                colors = ['coral', 'steelblue']
                labels = ['Raw 2D Detections\n(All Fragments)', 'True 3D Spermatid Nuclei\n(Consolidated)']
                
                bars = ax_bar.barh(y_pos, counts, color=colors, edgecolor='black', height=0.6)
                
                # Dynamic x-limit for labels
                ax_bar.set_xlim(0, max(counts) * 1.3)
                
                for i, v in enumerate(counts):
                    ax_bar.text(v + (max(counts)*0.02), y_pos[i], f"{v:,}", va='center', fontweight='bold', fontsize=12)
                    
                ax_bar.set_yticks(y_pos)
                ax_bar.set_yticklabels(labels, fontsize=11, fontweight='bold')
                ax_bar.set_xlabel("Total Count", fontsize=12)
                ax_bar.set_title("Tracking Reduction Effect", fontsize=13, fontweight='bold')
                ax_bar.spines['top'].set_visible(False)
                ax_bar.spines['right'].set_visible(False)
                
                # B) Donut Chart for Composition
                ax_pie = fig_dyn.add_subplot(1, 2, 2)
                
                # Calculate Single vs Multi slice
                n_single = len(df_tracks[df_tracks['n_slices'] == 1])
                n_multi = len(df_tracks[df_tracks['n_slices'] > 1])
                
                pie_labels = [f"Multi-Slice\n({n_multi:,})", f"Single-Slice Counts\n({n_single:,})"]
                pie_sizes = [n_multi, n_single]
                pie_colors = ['#2ca02c', '#d62728'] # Green and Red-ish matching standard dists
                
                wedge_props = dict(width=0.45, edgecolor='white', linewidth=2)
                wedges, texts, autotexts = ax_pie.pie(
                    pie_sizes, labels=None, colors=pie_colors,
                    autopct='%1.1f%%',
                    startangle=90,
                    pctdistance=0.75,
                    wedgeprops=wedge_props,
                    textprops={'fontsize': 12, 'fontweight': 'bold'}
                )
                
                # Place percentage text inside wedges with high contrast
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontsize(13)
                    autotext.set_fontweight('bold')

                # Place external labels with leader lines
                ax_pie.legend(
                    wedges, pie_labels,
                    title="Track Type",
                    loc="lower center",
                    bbox_to_anchor=(0.5, -0.15),
                    ncol=2,
                    fontsize=11,
                    frameon=False
                )
                    
                ax_pie.set_title(f"Composition of {total_3d:,} True 3D Spermatid Nuclei", fontsize=13, fontweight='bold')
                
                plt.tight_layout(rect=[0, 0.03, 1, 0.95])
                fig_dyn.savefig(os.path.join(plot_dir, "population_consolidation.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_dyn, dpi=300, bbox_inches='tight')
                plt.close(fig_dyn)

            # --- PAGE 2: 3D MORPHOMETRICS SUMMARY ---
            if df_tracks is not None and not df_tracks.empty:
                fig_3d = plt.figure(figsize=(11, 8.5))
                fig_3d.suptitle("3D Population Statistics (Tracked Spermatid Nuclei)", fontsize=14, fontweight='bold')
                
                # 3D Length
                ax3d_1 = fig_3d.add_subplot(2, 2, 1)
                vals_3d = df_tracks['total_3d_length_um']
                ax3d_1.hist(vals_3d, bins=20, color='darkorange', edgecolor='black', alpha=0.7)
                m3d_med = vals_3d.median()
                m3d_avg = vals_3d.mean()
                ax3d_1.axvline(m3d_med, color='red', linestyle='-', label=f"Median: {m3d_med:.1f}")
                ax3d_1.axvline(m3d_avg, color='black', linestyle='--', label=f"Mean: {m3d_avg:.1f}")
                ax3d_1.set_title("Total 3D Geodesic Length")
                ax3d_1.set_xlabel("Length (um)")
                ax3d_1.set_ylabel("Frequency")
                ax3d_1.legend(fontsize=8)
                
                # NEW: Save 3d_length_distribution.png for Excel embedding
                fig_3d_len = plt.figure(figsize=(6, 4))
                ax_3dl = fig_3d_len.add_subplot(1, 1, 1)
                ax_3dl.hist(vals_3d, bins=20, color='darkorange', edgecolor='black', alpha=0.7)
                ax_3dl.set_title("3D Length Distribution")
                fig_3d_len.savefig(os.path.join(plot_dir, "3d_length_distribution.png"), dpi=300, bbox_inches='tight')
                plt.close(fig_3d_len)

                # 3D Tortuosity
                ax3d_2 = fig_3d.add_subplot(2, 2, 2)
                vt = df_tracks['tortuosity_3d']
                # Filter out extreme collision outliers strictly for visualization resolution
                vt_viz = vt[(vt >= 0.95) & (vt <= 3.0)]
                ax3d_2.hist(vt_viz, bins=25, color='purple', edgecolor='black', alpha=0.6)
                ax3d_2.axvline(vt.median(), color='red', linestyle='-', label=f"Median: {vt.median():.2f}")
                ax3d_2.axvline(vt.mean(), color='black', linestyle='--', label=f"Mean: {vt.mean():.2f}")
                ax3d_2.set_xlim(0.95, 3.0)
                ax3d_2.set_title("3D Tortuosity (Curvature)")
                ax3d_2.set_xlabel("Ratio (Length / Distance)")
                ax3d_2.set_ylabel("Frequency")
                ax3d_2.legend(fontsize=8)

                # Vertical Extent
                ax3d_3 = fig_3d.add_subplot(2, 2, 3)
                ve = df_tracks['z_extent_um']
                ax3d_3.hist(ve, bins=15, color='teal', edgecolor='black', alpha=0.7)
                ax3d_3.axvline(ve.median(), color='red', linestyle='-', label=f"Median: {ve.median():.1f}")
                ax3d_3.axvline(ve.mean(), color='black', linestyle='--', label=f"Mean: {ve.mean():.1f}")
                ax3d_3.set_title("Z-Extent (Vertical Span)")
                ax3d_3.set_xlabel("Vertical Height (um)")
                ax3d_3.set_ylabel("Frequency")
                ax3d_3.legend(fontsize=8)

                # Volume
                ax3d_4 = fig_3d.add_subplot(2, 2, 4)
                vv = df_tracks['volume_um3']
                ax3d_4.hist(vv, bins=20, color='gray', edgecolor='black', alpha=0.7)
                ax3d_4.axvline(vv.median(), color='red', linestyle='-', label=f"Median: {vv.median():.0f}")
                ax3d_4.axvline(vv.mean(), color='black', linestyle='--', label=f"Mean: {vv.mean():.0f}")
                ax3d_4.set_title("Approximated 3D Volume")
                ax3d_4.set_xlabel("Volume (um³)")
                ax3d_4.set_ylabel("Frequency")
                ax3d_4.legend(fontsize=8)

                plt.tight_layout(rect=[0, 0.03, 1, 0.95])
                fig_3d.savefig(os.path.join(plot_dir, "3d_population_stats.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_3d, dpi=300, bbox_inches='tight')
                plt.close(fig_3d)

            # --- PAGE 4: ADVANCED 3D BIOMETRICS ---
            if df_tracks is not None and not df_tracks.empty:
                fig_adv = plt.figure(figsize=(11, 8.5))
                fig_adv.suptitle("Advanced 3D Biometrics Dashboard", fontsize=16, fontweight='bold', y=0.96)
                
                # Helper for Mean/Median
                def add_stats_lines(ax, data_series):
                    if data_series.empty or data_series.isna().all(): return
                    m = data_series.mean()
                    med = data_series.median()
                    ax.axvline(med, color='red', linestyle='--', linewidth=1.5, label=f'Median: {med:.2f}')
                    ax.axvline(m, color='green', linestyle=':', linewidth=2, label=f'Mean: {m:.2f}')
                    ax.legend(fontsize=8)

                # 4 panels: Pitch, Thickness, Taper, Nearest Neighbor
                ax_p = fig_adv.add_subplot(2, 2, 1)
                sns.histplot(data=df_tracks, x='pitch_deg', bins=30, ax=ax_p, color='orange')
                ax_p.set_title("Pitch Angle (Vertical Plunge)")
                ax_p.set_xlabel("Degrees (0=Flat, 90=Vertical)")
                add_stats_lines(ax_p, df_tracks['pitch_deg'])
                
                ax_th = fig_adv.add_subplot(2, 2, 2)
                sns.histplot(data=df_tracks, x='thickness_um', bins=30, ax=ax_th, color='#17becf')
                ax_th.set_title("Effective Nucleus Thickness")
                ax_th.set_xlabel("Average Diameter (µm)")
                add_stats_lines(ax_th, df_tracks['thickness_um'])
                
                ax_ta = fig_adv.add_subplot(2, 2, 3)
                sns.histplot(data=df_tracks, x='taper_ratio', bins=30, ax=ax_ta, color='purple')
                ax_ta.set_title("Morphological Taper Ratio")
                ax_ta.set_xlabel("Max Area / Min Area")
                add_stats_lines(ax_ta, df_tracks['taper_ratio'])
                
                ax_nn = fig_adv.add_subplot(2, 2, 4)
                valid_nn = df_tracks['nearest_neighbor_um'].dropna()
                if not valid_nn.empty:
                    sns.histplot(valid_nn, bins=30, ax=ax_nn, color='brown')
                    add_stats_lines(ax_nn, valid_nn)
                ax_nn.set_title("Spatial Packing Density")
                ax_nn.set_xlabel("Distance to Nearest Neighbor (µm)")
                
                plt.tight_layout(rect=[0, 0.03, 1, 0.93])
                fig_adv.savefig(os.path.join(plot_dir, "advanced_biometrics.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_adv, dpi=300, bbox_inches='tight')
                plt.close(fig_adv)

            # --- PAGE 5: METHODS & INTERPRETATION GUIDE ---
            fig_guide = plt.figure(figsize=(11, 8.5))
            ax_g = fig_guide.add_axes([0.05, 0.05, 0.9, 0.9])
            ax_g.axis('off')
            guide_full = (
                "METHODS & CALCULATION DETAILS (Biometrics & Tissue Mechanics)\n"
                f"{'='*80}\n\n"
                "1. Total 3D Geodesic Length (µm)\n"
                "   Formula: L_3D = Σ sqrt((Δx)² + (Δy)² + (Δz)²)\n"
                "   Biology: Measures the true multi-slice structural backbone of the spermatid,\n"
                "   tracking along its curvature rather than assuming rigid linearity.\n\n"
                "2. 3D Tortuosity / Curvature\n"
                "   Formula: T = (Total Geodesic Length) / (Straight-Line Point-to-Point Distance)\n"
                "   Biology: A morphological metric of acrosomal or tail curvature. Normal spermatids\n"
                "   should be moderately rigid (~1.0 - 1.2). Extreme values (>2.0) generally flag\n"
                "   tissue deformation, fixation artifacts, or algorithmic tracking fusions.\n\n"
                "3. Z-Extent (Vertical Tissue Span)\n"
                "   Biology: Captures how deep the nucleus penetrates through the epithelial layer.\n"
                "   Plunging nuclei indicate distinct maturation stages closer to the lumen boundary.\n\n"
                "4. Approximated 3D Volume (µm³)\n"
                "   Calculated via Riemann sum of 2D segmented pixel area × Z-slice thickness.\n"
                "   Biology: Directly correlates to chromatin condensation levels. Aberrantly\n"
                "   large volumes may signify poor condensation or double-nucleation.\n\n"
                "5. Effective Thickness (Average Diameter, µm)\n"
                "   Formula: D_avg = 2 * sqrt( (V_3D / L_3D) / π )\n"
                "   Biology: Models the nucleus as a cylinder to quantify chromatin thickness.\n"
                "   Particularly useful for isolating swelling or decompaction pathologies.\n\n"
                "6. Pitch Angle (Plunge Vector, Degrees)\n"
                "   Formula: θ = |arcsin(ΔZ / Length_Euclidean)| * (180/π)\n"
                "   Biology: Quantifies the absolute orientation of the elongating nucleus relative\n"
                "   to the basal lamina plane. Values near 90° indicate perpendicular apical plunging.\n\n"
                "7. Taper Ratio\n"
                "   Formula: R = (Maximum Cross-sectional Area) / (Minimum Cross-sectional Area)\n"
                "   Biology: Identifies tapering anomalies near the acrosome versus the tail insertion.\n"
                "   Extreme variance implies severe morphological irregularity or fusion errors.\n\n"
                "8. Spatial Packing Density (Nearest Neighbor Dist, µm)\n"
                "   Calculated: 3D Euclidean distance from this track's centroid to its closest sibling.\n"
                "   Biology: Measures tubule capacity and cyst grouping density during spermiation."
            )
            ax_g.text(0, 1, guide_full, transform=ax_g.transAxes, fontsize=10, family='monospace', verticalalignment='top', linespacing=1.3)
            fig_guide.savefig(os.path.join(plot_dir, "methods_guide.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_guide, dpi=300, bbox_inches='tight')
            plt.close(fig_guide)

            # --- PAGE 6: GLOBAL STATISTICS TABLE (after Methods) ---
            # (Table was moved here from before Methods for logical document flow)
            fig_tab = plt.figure(figsize=(11, 8.5))
            ax_t = fig_tab.add_subplot(1, 1, 1)
            ax_t.axis('off')
            ax_t.set_title("Global Population Statistics Summary", fontsize=14, fontweight='bold', pad=20)
            
            stats_rows = []
            if not df.empty:
                l2d = df['length_um_geodesic']
                stats_rows.append(["2D Fragment Geodesic Length (µm)", f"{l2d.mean():.2f}", f"{l2d.median():.2f}", f"{l2d.std():.2f}"])
            
            if df_tracks is not None and not df_tracks.empty:
                l3d = df_tracks['total_3d_length_um']
                ze  = df_tracks['z_extent_um']
                vo  = df_tracks['volume_um3']
                to  = df_tracks['tortuosity_3d']
                th  = df_tracks['thickness_um']
                pi  = df_tracks['pitch_deg']
                ta  = df_tracks['taper_ratio']
                nn  = df_tracks['nearest_neighbor_um'].dropna()
                stats_rows.append(["3D Nucleus Geodesic Length (µm)", f"{l3d.mean():.2f}", f"{l3d.median():.2f}", f"{l3d.std():.2f}"])
                stats_rows.append(["3D Nucleus Z-Extent (Vertical Depth, µm)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
                stats_rows.append(["3D Nucleus Approximated Volume (µm³)", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
                stats_rows.append(["3D Nucleus Tortuosity (Curvature Ratio)", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
                stats_rows.append(["3D Nucleus Thickness (Average Diameter, µm)", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])
                stats_rows.append(["3D Nucleus Pitch (Plunge Angle, degrees)", f"{pi.mean():.1f}", f"{pi.median():.1f}", f"{pi.std():.1f}"])
                stats_rows.append(["3D Nucleus Morphological Taper (Ratio)", f"{ta.mean():.2f}", f"{ta.median():.2f}", f"{ta.std():.2f}"])
                if not nn.empty:
                    stats_rows.append(["3D Nucleus Spatial Density (Nearest Neighbour, µm)", f"{nn.mean():.1f}", f"{nn.median():.1f}", f"{nn.std():.1f}"])
            
            if stats_rows:
                table = ax_t.table(
                    cellText=stats_rows,
                    colLabels=["Metric", "Mean", "Median", "Std Dev"],
                    loc='center', cellLoc='center',
                    colWidths=[0.38, 0.2, 0.2, 0.2]
                )
                table.auto_set_font_size(False)
                table.set_fontsize(10)
                table.scale(1.2, 2.5)
            
            fig_tab.savefig(os.path.join(plot_dir, "global_statistics_table.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_tab, dpi=300, bbox_inches='tight')
            plt.close(fig_tab)

            
            # --- SUBSEQUENT PAGES: PER-SLICE DETAILS (2 panels: [Panel] | [Histogram]) ---
            overlay_dir = os.path.join(out_dir, "overlays")
            for idx_p, (row_idx, row) in enumerate(df_summary.iterrows()):
                z = int(row['z_slice'])
                panel_path = os.path.join(overlay_dir, f"z{z:02d}_panel.png")
                
                if not os.path.exists(panel_path):
                    continue
                
                fig_slice = plt.figure(figsize=(18, 7))
                fig_slice.suptitle(f"Z-Slice {z:02d} Analysis [Original | Overlay | Distribution]", fontsize=12, fontweight='bold')
                
                # Panel: Side-by-Side (Original | Overlay)
                ax_panel = fig_slice.add_subplot(1, 2, 1)
                ax_panel.imshow(plt.imread(panel_path))
                ax_panel.set_title(f"Visual Verification (N={int(row['n_spermatids'])})")
                ax_panel.axis('off')
                
                # Plot: Stats
                ax_hist = fig_slice.add_subplot(1, 2, 2)
                slice_data = df[df['z_slice'] == z]
                if not slice_data.empty:
                    ax_hist.hist(slice_data['length_um_geodesic'], bins=15, color='skyblue', edgecolor='black')
                    ax_hist.set_title(f"Z={z} Length Distribution")
                    ax_hist.set_xlabel("Spermatid Nucleus Length (µm)")
                    ax_hist.set_ylabel("Frequency (Count)")
                    
                    m_med = slice_data['length_um_geodesic'].median()
                    m_avg = slice_data['length_um_geodesic'].mean()
                    ax_hist.axvline(m_med, color='red', linestyle='-', alpha=0.7, label=f"Median: {m_med:.1f}")
                    ax_hist.axvline(m_avg, color='orange', linestyle='--', alpha=0.7, label=f"Mean: {m_avg:.1f}")
                    ax_hist.legend(fontsize=9)
                else:
                    ax_hist.text(0.5, 0.5, "No Detections", ha='center', va='center')
                
                pdf.savefig(fig_slice, dpi=300)
                plt.close(fig_slice)
                
                if gui_callback:
                    gui_callback(int(80 + (20 * (idx_p+1) / len(df_summary))))
                
        print(f"Report successfully saved to {pdf_path}")
        
        # --- GENERATE POWERPOINT REPORT ---
        try:
            generate_pptx_report(out_dir, df, df_summary, um, df_tracks)
        except Exception as e:
            import traceback
            err_msg = traceback.format_exc()
            print(f"PPTX Report failed: {e}")
            write_error_log(out_dir, "PowerPoint Generator (via Batch)", err_msg)
            
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating PDF report: {e}")
        write_error_log(out_dir, "PDF Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"PDF Report failed to generate completely.\n{e}")
        except Exception:
            pass



def generate_pptx_report(out_dir, df, df_summary, um, df_tracks=None):
    """
    Generates a native Microsoft Office PowerPoint presentation (.pptx) with
    fully editable, data-embedded charts for each key biometric.

    Each chart is a *true* Office Open XML chart object (not a rasterised image),
    so the researcher can re-style, re-colour, and re-export from PowerPoint without
    any additional software.

    Slide structure
    ---------------
    1. **Global Population Analytics** — column histograms of 2D geodesic length
       and 3D track length (if tracking was run).
    2. **Population Tracking Consolidation** (if tracking) — pie chart of
       single-slice vs. multi-slice track fractions.
    3. **3D Biometrics Dashboard** — tortuosity vs. 3D length scatter and
       pitch angle distribution.
    4. **Methods & Calculation Details** — text slide with biological and
       mathematical formulae for all reported metrics.

    Biological context of each chart
    ---------------------------------
    - **Length histogram**: A Gaussian centred near the species-expected length
      indicates healthy elongation; bimodal distributions suggest two cohorts.
    - **Pie chart**: Fraction of multi-slice tracks indicates true 3D continuity.
    - **Tortuosity scatter**: Healthy cells cluster in the low-tortuosity region.
    - **Pitch angle**: Right-skewed towards 90° during the apical plunging phase.

    Args:
        out_dir (str): Directory where ``spermatid_analysis_report.pptx`` is saved.
        df (pd.DataFrame): Per-spermatid 2D measurement table.
        df_summary (pd.DataFrame): Per-slice summary statistics.
        um (float): Microns-per-pixel (``UM_PER_PX_XY``).
        df_tracks (pd.DataFrame, optional): 3D track table; ``None`` skips
            tracking-specific slides.
    """
    try:
        import os as _os
        import numpy as _np
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from tkinter import messagebox
        
        print("PPTX: Starting report generation...")
        # We need a fallback blank pptx
        prs = Presentation()
        
        # Safe blank layout (often index 6 or 5)
        try:
            blank_slide_layout = prs.slide_layouts[6] 
        except Exception:
            blank_slide_layout = prs.slide_layouts[0]
            
        def add_hyperlink(slide, sheet_name="Population_Summary"):
            excel_name = "batch_analysis_results_v11.xlsx"
            # PowerPoint/Excel deep-link fragment: filename.xlsx#Sheet!A1
            target = f"{excel_name}#'{sheet_name}'!A1"
            
            txBox = slide.shapes.add_textbox(Inches(0.2), Inches(7.1), Inches(9.5), Inches(0.4))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"Click to View Detailed Data: {excel_name} [{sheet_name}]"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0, 0, 255) # Blue
            run.font.underline = True
            try:
                # Shape-level link
                txBox.click_action.hyperlink.address = target
                # Text-run level link
                run.hyperlink.address = target
            except Exception:
                pass

        def add_line_chart(slide, x_data, y_data, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = list(x_data)
            chart_data.add_series('Count', list(y_data))
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.has_major_gridlines = True

        def add_horizontal_bar_chart(slide, categories, values, colors, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', values)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.size = Pt(8)
            
            # Show data labels for bar values
            plot = chart.plots[0]
            plot.has_data_labels = True
            for i, point in enumerate(plot.series[0].points):
                point.data_label.font.size = Pt(10)
                point.data_label.font.bold = True

        def add_histogram(slide, data_series, left, top, width, height, title, bins=20):
            if data_series is None or data_series.empty or data_series.isna().all():
                return
            
            clean_data = data_series.dropna()
            counts, bin_edges = _np.histogram(clean_data, bins=bins)
            avg = clean_data.mean()
            med = clean_data.median()
            
            # Create category names from bin boundaries
            categories = [f"{bin_edges[i]:.1f}-{bin_edges[i+1]:.1f}" for i in range(len(counts))]
            
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', list(counts))
            # ADD DUMMY SERIES FOR LEGEND STATS
            chart_data.add_series(f'Mean: {avg:.2f}', [])
            chart_data.add_series(f'Median: {med:.2f}', [])
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.CORNER
            chart.legend.font.size = Pt(8)
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Reduce axis label crowding
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

        # --- Slide 1: Global Analytics Overview ---
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Spermatid Population Overview"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        
        # Bottom Left: Detections per Slice (Line Chart) - Identical to PDF Page 1
        if not df_summary.empty:
            add_line_chart(slide1, df_summary['z_slice'], df_summary['n_spermatids'], Inches(0.2), Inches(4.0), Inches(4.5), Inches(3.0), "Detections per Z-Slice (Raw)")

        # Top Center Left: 2D Length Dist
        if not df.empty:
            add_histogram(slide1, df['length_um_geodesic'], Inches(0.2), Inches(0.8), Inches(4.5), Inches(3.0), "Global 2D Geodesic Length Distribution")
        
            add_histogram(slide1, df_tracks['total_3d_length_um'], Inches(5.0), Inches(0.8), Inches(4.5), Inches(3.0), "Global 3D Solid Nuclei Length")
            
        add_hyperlink(slide1, "Population_Summary")
        
        # --- Slide 2: Population Consolidation ---
        if df_tracks is not None and not df_tracks.empty:
            slide2 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "3D Population Tracking & Consolidation"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            
            # Left: Reduction Bar Chart (PDF Parity)
            total_2d = len(df)
            total_3d = len(df_tracks)
            add_horizontal_bar_chart(slide2, 
                                     ['3D Consolidated Nuclei', 'Raw 2D Detections'], 
                                     [total_3d, total_2d], 
                                     None, Inches(0.2), Inches(1.5), Inches(4.5), Inches(4.5), "Tracking Reduction Effect")

            # Right: Composition Pie Chart
            n_single = len(df_tracks[df_tracks['n_slices'] == 1])
            n_multi = len(df_tracks[df_tracks['n_slices'] > 1])
            
            chart_data = CategoryChartData()
            chart_data.categories = ['Multi-Slice Tracks', 'Single-Slice Counts']
            pie_values = [n_multi, n_single]
            chart_data.add_series('Population', pie_values)
            
            chart2 = slide2.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(4.6), Inches(1.2), Inches(5.2), Inches(5.2), chart_data
            ).chart
            chart2.has_legend = True
            chart2.legend.position = XL_LEGEND_POSITION.CORNER
            chart2.legend.font.size = Pt(8)
            chart2.chart_title.text_frame.text = f"Composition of {total_3d:,} True 3D Nuclei"
            chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            
            plot = chart2.plots[0]
            plot.has_data_labels = True
            total = sum(pie_values) if sum(pie_values) > 0 else 1
            for i, point in enumerate(plot.series[0].points):
                val = pie_values[i]
                pct = (val / total) * 100
                label_text = f"{val:,}\n({pct:.1f}%)"
                point.data_label.text_frame.text = label_text
                point.data_label.font.size = Pt(9)
                point.data_label.font.bold = True
            
            add_hyperlink(slide2, "Population_Summary")
            
        # ---------------------------------------------------------------------
        # SLIDE 3: Advanced 3D Biometrics Dashboard
        # ---------------------------------------------------------------------
        if df_tracks is not None and not df_tracks.empty:
            slide3 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "Advanced 3D Biometrics Dashboard"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            
            add_histogram(slide3, df_tracks['pitch_deg'], Inches(0.2), Inches(0.8), Inches(4.5), Inches(2.9), "Pitch Angle (Degrees)", bins=20)
            add_histogram(slide3, df_tracks['thickness_um'], Inches(5.0), Inches(0.8), Inches(4.5), Inches(2.9), "Effective Thickness (µm)", bins=20)
            add_histogram(slide3, df_tracks['taper_ratio'], Inches(0.2), Inches(3.8), Inches(4.5), Inches(2.9), "Morphological Taper Ratio", bins=20)
            add_histogram(slide3, df_tracks['nearest_neighbor_um'], Inches(5.0), Inches(3.8), Inches(4.5), Inches(2.9), "Nearest Neighbor Density (µm)", bins=20)
            
            add_hyperlink(slide3, "3D_Track_Summary")
            
        # ---------------------------------------------------------------------
        # SLIDE 4: Global Population Statistics Summary Table
        # ---------------------------------------------------------------------
        slide4 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Global Population Statistics Summary"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        
        # Prepare stats matching PDF
        stats_rows = [["Metric", "Mean", "Median", "Std Dev"]]
        if not df.empty:
            l2d = df['length_um_geodesic']
            stats_rows.append(["2D Fragment Length (µm)", f"{l2d.mean():.2f}", f"{l2d.median():.2f}", f"{l2d.std():.2f}"])
        
        if df_tracks is not None and not df_tracks.empty:
            l3d = df_tracks['total_3d_length_um']
            ze  = df_tracks['z_extent_um']
            vo  = df_tracks['volume_um3']
            to  = df_tracks['tortuosity_3d']
            th  = df_tracks['thickness_um']
            pi  = df_tracks['pitch_deg']
            ta  = df_tracks['taper_ratio']
            nn  = df_tracks['nearest_neighbor_um'].dropna()
            
            stats_rows.append(["3D Geodesic Length (µm)", f"{l3d.mean():.2f}", f"{l3d.median():.2f}", f"{l3d.std():.2f}"])
            stats_rows.append(["3D Z-Extent (Height, µm)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
            stats_rows.append(["3D Volume (µm³)", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
            stats_rows.append(["3D Tortuosity (Ratio)", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
            stats_rows.append(["3D Thickness (µm)", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])
            stats_rows.append(["3D Pitch (Degrees)", f"{pi.mean():.1f}", f"{pi.median():.1f}", f"{pi.std():.1f}"])
            stats_rows.append(["3D Taper Ratio", f"{ta.mean():.2f}", f"{ta.median():.2f}", f"{ta.std():.2f}"])
            if not nn.empty:
                stats_rows.append(["Nearest Neighbor (µm)", f"{nn.mean():.1f}", f"{nn.median():.1f}", f"{nn.std():.1f}"])
        
        if len(stats_rows) > 1:
            rows = len(stats_rows)
            cols = 4
            table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
            table = table_shape.table
            
            # Header styling
            for c in range(cols):
                cell = table.cell(0, c)
                cell.text_frame.text = stats_rows[0][c]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196) # Standard Blue
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True
                
            # Body styling
            for r in range(1, rows):
                for c in range(cols):
                    cell = table.cell(r, c)
                    cell.text_frame.text = stats_rows[r][c]
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    if c == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
        
        add_hyperlink(slide4, "Population_Summary")

        # ---------------------------------------------------------------------
        # SLIDE 5: Methods & Interpretation Guide (Exact PDF Synchronization)
        # ---------------------------------------------------------------------
        slide5 = prs.slides.add_slide(blank_slide_layout)
        txBox_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.text = "Methods & Calculation Details"
        tf_title.paragraphs[0].font.size = Pt(20)
        tf_title.paragraphs[0].font.bold = True

        # Synchronized text from the PDF report version
        methods_items = [
            ("1. Total 3D Geodesic Length (µm)", [
                ("Formula: ", "L_3D = Σ sqrt((Δx)² + (Δy)² + (Δz)²)"),
                ("Biology: ", "Measures the true multi-slice structural backbone of the spermatid, tracking along its curvature rather than assuming rigid linearity.")
            ]),
            ("2. 3D Tortuosity / Curvature", [
                ("Formula: ", "T = (Total Geodesic Length) / (Straight-Line Point-to-Point Distance)"),
                ("Biology: ", "A morphological metric of acrosomal or tail curvature. Normal spermatids should be moderately rigid (~1.0 - 1.2). Extreme values (>2.0) generally flag tissue deformation, fixation artifacts, or algorithmic tracking fusions.")
            ]),
            ("3. Z-Extent (Vertical Tissue Span)", [
                ("Biology: ", "Captures how deep the nucleus penetrates through the epithelial layer. Plunging nuclei indicate distinct maturation stages closer to the lumen boundary.")
            ]),
            ("4. Approximated 3D Volume (µm³)", [
                ("", "Calculated via Riemann sum of 2D segmented pixel area × Z-slice thickness."),
                ("Biology: ", "Directly correlates to chromatin condensation levels. Aberrantly large volumes may signify poor condensation or double-nucleation.")
            ]),
            ("5. Effective Thickness (Average Diameter, µm)", [
                ("Formula: ", "D_avg = 2 * sqrt( (V_3D / L_3D) / π )"),
                ("Biology: ", "Models the nucleus as a cylinder to quantify chromatin thickness. Particularly useful for isolating swelling or decompaction pathologies.")
            ]),
            ("6. Pitch Angle (Plunge Vector, Degrees)", [
                ("Formula: ", "θ = |arcsin(ΔZ / Length_Euclidean)| * (180/π)"),
                ("Biology: ", "Quantifies the absolute orientation of the elongating nucleus relative to the basal lamina plane. Values near 90° indicate perpendicular apical plunging.")
            ]),
            ("7. Taper Ratio", [
                ("Formula: ", "R = (Maximum Cross-sectional Area) / (Minimum Cross-sectional Area)"),
                ("Biology: ", "Identifies tapering anomalies near the acrosome versus the tail insertion. Extreme variance implies severe morphological irregularity or fusion errors.")
            ]),
            ("8. Spatial Packing Density (Nearest Neighbor Dist, µm)", [
                ("Calculated: ", "3D Euclidean distance from this track's centroid to its closest sibling."),
                ("Biology: ", "Measures tubule capacity and cyst grouping density during spermiation.")
            ])
        ]

        top_y = 1.0
        for title, content_list in methods_items:
            tb = slide5.shapes.add_textbox(Inches(0.5), Inches(top_y), Inches(9), Inches(0.72))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(10)
            
            for label, text in content_list:
                p2 = tf.add_paragraph()
                p2.font.size = Pt(8.5)
                if label:
                    run1 = p2.add_run()
                    run1.text = label
                    run1.font.bold = True
                run2 = p2.add_run()
                run2.text = text
                run2.font.bold = False
            
            top_y += 0.74

        add_hyperlink(slide5, "Population_Summary")

        # Build and Save
        pptx_name = "batch_analysis_results_v11.pptx"
        pptx_path = os.path.join(out_dir, pptx_name)
        try:
            print(f"PPTX: Saving to {pptx_path}...")
            prs.save(pptx_path)
            print(f"PPTX Report successfully saved to: {pptx_path}")
            return True
        except PermissionError:
            print(f"CRITICAL ERROR: File '{os.path.basename(pptx_path)}' is currently open!")
            return False
            
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"Failed to generate PPTX report: {e}")
        # Standardize out_dir for logging
        log_dir = os.path.dirname(pptx_path) if pptx_path else "."
        write_error_log(log_dir, "PowerPoint Generator", err_msg)
        return False


# =============================================================================
# ROI GUI (single-file v9)
# =============================================================================
import traceback
try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
    from matplotlib.figure import Figure
    _TK_AVAILABLE = True
except Exception:
    _TK_AVAILABLE = False

PARAM_DESCRIPTIONS = {
    "UM_PER_PX_XY": "Microns per pixel in XY (Calibration)",
    "UM_PER_SLICE_Z": "Microns per Z-slice step (Calibration)",
    "CLAHE_CLIP": "Local contrast limit for CLAHE enhancement",
    "CLAHE_KERNEL": "Grid size for CLAHE enhancement",
    "BG_SIGMA": "Gaussian blur radius for background subtraction",
    "THRESHOLD_HI": "Intensity threshold for main spermatid seeds",
    "THRESHOLD_LO": "Intensity threshold for trailing edges",
    "CLOSE_RADIUS": "Radius for morphological closing (px)",
    "MIN_HOLE_AREA": "Smallest hole area to fill (px)",
    "MIN_OBJ_PX": "Absolute minimum foreground size (px)",
    "MAX_BRIDGE_PX": "Max gap to bridge in skeleton (px)",
    "MAX_BRANCH_LEN_PX": "Prune skeleton spurs shorter than this (px)",
    "BREAK_JUNCTIONS": "Automatically sever intersecting skeletons (True/False)",
    "USE_EARLY_SHAPE_FILTER": "Apply basic eccentricity filters early (True/False)",
    "MIN_ECCENTRICITY": "Minimum eccentricity (if early filtered)",
    "MAX_MINOR_PX": "Maximum minor axis length (px)",
    "MIN_AXIS_RATIO": "Minimum major/minor axis ratio",
    "MIN_MAJOR_PX": "Minimum major axis length (px)",
    "MIN_SKEL_LEN_PX": "Minimum valid skeleton length (px)",
    "MAX_GEODESIC_LEN_PX": "Maximum valid skeleton length (px)",
    "MAX_WIDTH_PX": "Maximum allowed thickness (px)",
    "MIN_LENGTH_WIDTH_RATIO": "Minimum required elongation ratio",
    "MAX_BRANCH_NODES": "Maximum allowed branch intersections",
    "MAX_TORTUOSITY": "Maximum allowed curvature index",
    "MAX_ENDPOINT_COUNT": "Maximum allowed skeleton endpoints",
    "ALLOW_LOOPS": "Allow topological loops in targets (True/False)",
    "DO_TRACKING": "Enable 3D tracking across Z-slices (True/False)",
    "TRACK_MAX_DIST_UM": "Maximum allowed tracking displacement (um)",
    "TRACK_MAX_GAP_SLICES": "Maximum allowed missing Z-slices in a track"
}

class ParameterEditor(tk.Toplevel):
    """
    A Tkinter ``Toplevel`` window providing an interactive editor for all pipeline
    configuration parameters.

    The editor auto-generates entry widgets per config key with human-readable labels
    from ``PARAM_DESCRIPTIONS``.  Supports:
    - **Apply to Session** — updates the live ``CONFIG`` dict without restart.
    - **Reset to Defaults** — restores factory defaults.
    - **Load JSON** — loads parameters from an external JSON file.
    - **Save JSON** — serialises current values to disk for reuse.

    The window is scrollable so all parameters are accessible on any screen size.
    """
    def __init__(self, parent, current_config, default_config, apply_callback):
        super().__init__(parent)
        self.title("Parameter Configuration")
        self.geometry("900x700")
        self.current_config = current_config
        self.default_config = default_config
        self.apply_callback = apply_callback
        
        self.entries = {}
        
        # Tools Frame (Top)
        tools = tk.Frame(self, bg='#e0e0e0', padx=10, pady=10)
        tools.pack(side='top', fill='x')
        
        tk.Button(tools, text="Apply to Session", command=self.apply, bg="#d4edda", font=("Arial", 10, "bold")).pack(side="left", padx=5)
        tk.Button(tools, text="Reset to Defaults", command=self.reset_defaults, bg="#f8d7da").pack(side="left", padx=5)
        tk.Button(tools, text="Load JSON", command=self.load_json).pack(side="left", padx=5)
        tk.Button(tools, text="Save JSON", command=self.save_json).pack(side="left", padx=5)
        
        # Scrollable Canvas
        canvas = tk.Canvas(self)
        scrollbar = ttk.Scrollbar(self, orient="vertical", command=canvas.yview)
        self.scrollable_frame = tk.Frame(canvas)
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        self.populate_form(self.current_config)

    def populate_form(self, cfg):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        self.entries = {}
        row = 0
        for k, v in cfg.items():
            if k in PARAM_DESCRIPTIONS and isinstance(v, (int, float, bool)):
                # Label Key
                tk.Label(self.scrollable_frame, text=k, font=("Arial", 10, "bold"), width=30, anchor="e").grid(row=row, column=0, padx=10, pady=5, sticky="e")
                
                # Entry value
                var = tk.StringVar(value=str(v))
                ent = tk.Entry(self.scrollable_frame, textvariable=var, width=15)
                ent.grid(row=row, column=1, padx=10, pady=5)
                self.entries[k] = (var, type(v))
                
                # Description
                tk.Label(self.scrollable_frame, text=PARAM_DESCRIPTIONS[k], fg="dimgray", anchor="w").grid(row=row, column=2, padx=10, pady=5, sticky="w")
                
                row += 1

    def apply(self):
        new_cfg = self.current_config.copy()
        try:
            for k, (var, t) in self.entries.items():
                val = var.get()
                if t == bool:
                    new_cfg[k] = val.lower() in ['true', '1', 't', 'y', 'yes']
                else:
                    new_cfg[k] = t(val)
        except ValueError as e:
            messagebox.showerror("Validation Error", f"Invalid input format: {e}")
            return
            
        self.apply_callback(new_cfg)
        self.destroy()

    def reset_defaults(self):
        self.populate_form(self.default_config)
        
    def load_json(self):
        fpath = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                with open(fpath, 'r') as f:
                    loaded = json.load(f)
                filtered = {k: v for k, v in loaded.items() if k in self.current_config}
                temp = self.current_config.copy()
                temp.update(filtered)
                self.populate_form(temp)
            except Exception as e:
                messagebox.showerror("Load Error", str(e))
                
    def save_json(self):
        fpath = filedialog.asksaveasfilename(defaultextension=".json", filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                new_cfg = self.current_config.copy()
                for k, (var, t) in self.entries.items():
                    val = var.get()
                    if t == bool:
                        new_cfg[k] = val.lower() in ['true', '1', 't', 'y', 'yes']
                    else:
                        new_cfg[k] = t(val)
                with open(fpath, 'w') as f:
                    json.dump(new_cfg, f, indent=4)
                messagebox.showinfo("Saved", f"Parameters saved to {os.path.basename(fpath)}")
            except Exception as e:
                messagebox.showerror("Save Error", str(e))

class SpermGUI:
    """
    The primary Tkinter-based graphical user interface for the Sperm Segmentation ROI Tool.

    Two-panel layout:
    - **Left sidebar**: directory loading, Z-slice navigation, tool selection,
      ROI polygon drawing, single and batch analysis, progress bars.
    - **Right canvas**: matplotlib ``FigureCanvasTkAgg`` rendering the current
      Z-slice (raw / overlay) and accepting polygon drawing interactions.

    Interaction modes (``mode_var``)
    ---------------------------------
    - ``'view'``   — raw image inspection and panning.
    - ``'review'`` — displays saved overlay PNG (requires prior batch run).
    - ``'roi'``    — polygon drawing; left-click = add vertex, right-click = undo.

    Thread architecture
    -------------------
    Batch analysis runs in a background thread so the Tkinter event loop stays
    responsive; progress bar updates are posted via ``root.after()``.
    """
    def open_parameter_editor(self):
        """
        Opens the :class:`ParameterEditor` ``Toplevel`` window and defines the
        ``on_apply`` callback that merges edited values into ``CONFIG``.
        """
        def on_apply(new_cfg):
            CONFIG.update(new_cfg)
            self.lbl_roi.config(text="Parameters updated in memory.")
            
        editor = ParameterEditor(self.root, CONFIG, self.default_config, on_apply)

    def __init__(self, root):
        """
        Initialises the main application window, all sidebar controls, the matplotlib
        canvas, and mouse/key event bindings.

        Args:
            root (tk.Tk): The root Tkinter window created by :func:`launch_gui`.
        """
        self.root = root
        self.root.title(f'Sperm Segmentation ROI Tool - v11 (Excel + Unique Folders)')
        self.root.geometry('1450x920')

        self.input_dir = ''
        self.files = []
        self.current_idx = 0
        self.current_img = None
        self.last_out_dir = ""

        self.roi_points = []
        self.drawing = False
        self.roi_active = False
        self._loaded_roi_mask = None

        self.sidebar = tk.Frame(root, width=280, bg='#f0f0f0')
        self.sidebar.pack(side='left', fill='y')

        self.default_config = CONFIG.copy()
        tk.Button(self.sidebar, text='⚙ Configure Parameters', command=self.open_parameter_editor, bg='#e2e3e5').pack(fill='x', padx=6, pady=(10,6))

        tk.Button(self.sidebar, text='Load Directory', command=self.load_directory, height=2).pack(fill='x', padx=6, pady=6)
        self.lbl_status = tk.Label(self.sidebar, text='No directory loaded', wraplength=260, justify='left')
        self.lbl_status.pack(pady=6)

        tk.Label(self.sidebar, text='Z-Slice Navigation').pack(pady=(20, 0))
        self.scale_z = tk.Scale(self.sidebar, from_=0, to=0, orient='horizontal', command=self.on_slide_change)
        self.scale_z.pack(fill='x', padx=10)
        self.lbl_z = tk.Label(self.sidebar, text='Z: 0 / 0')
        self.lbl_z.pack()

        tk.Label(self.sidebar, text='Tools').pack(pady=(20, 5))
        self.mode_var = tk.StringVar(value='view')
        tk.Radiobutton(self.sidebar, text='View/Nav (Raw Image)', variable=self.mode_var, value='view', command=self.render).pack(anchor='w', padx=10)
        tk.Radiobutton(self.sidebar, text='Review Overlays (After Batch)', variable=self.mode_var, value='review', command=self.render).pack(anchor='w', padx=10)
        tk.Radiobutton(self.sidebar, text='Draw ROI (Polygon)', variable=self.mode_var, value='roi', command=self.render).pack(anchor='w', padx=10)
        tk.Label(self.sidebar, text='(Left-click points, Right-click undo)', font=('Arial', 8, 'italic'), fg='dimgray').pack()
        tk.Button(self.sidebar, text='Finalize Polygon', command=self.finalize_roi, bg='#ffeeba').pack(fill='x', padx=20, pady=2)

        tk.Button(self.sidebar, text='Run Analysis on Slice', command=self.run_analysis_slice).pack(fill='x', padx=6, pady=18)
        tk.Button(self.sidebar, text='Run Batch (All Slices + 3D Track)', command=self.run_batch_analysis, bg='#d4edda', font=('Arial', 10, 'bold')).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Reset ROI', command=self.reset_roi).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Save ROI Mask', command=self.save_roi_mask).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Load ROI Mask', command=self.load_roi_mask).pack(fill='x', padx=6, pady=6)

        self.lbl_roi = tk.Label(self.sidebar, text='ROI: none', wraplength=260, justify='left')
        self.lbl_roi.pack(pady=10)

        self.canvas_frame = tk.Frame(root, bg='black')
        self.canvas_frame.pack(side='right', expand=True, fill='both')

        # --- TOP STATUS BAR (For Progress Visibility) ---
        self.top_status_frame = tk.Frame(self.canvas_frame, bg='#f0f0f0', height=80)
        self.top_status_frame.pack(side='top', fill='x')

        # --- Sub-frame for Status (Left) ---
        self.status_sub_frame = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.status_sub_frame.pack(side='left', padx=30, fill='y')

        # --- Dynamic Status Label ---
        self.lbl_batch_op = tk.Label(self.status_sub_frame, text='GUI Ready', font=('Arial', 10, 'bold'), fg='#2c3e50', bg='#f0f0f0')
        self.lbl_batch_op.pack(side='left', pady=5)

        # --- PROGRESS BOX: Sequential Progress Bars (At Top Right) ---
        self.p_container = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.p_container.pack(side='right', padx=20, pady=5)

        # Frame 1: 2D Batch Segmentation
        self.batch_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        self.batch_p_frame.pack(fill='x')
        tk.Label(self.batch_p_frame, text='Batch Progress (2D)', font=('Arial', 9, 'bold'), bg='#f0f0f0').pack(side='left', padx=10)
        self.progress = ttk.Progressbar(self.batch_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress.pack(side='left', padx=5, pady=5)
        self.lbl_progress_val = tk.Label(self.batch_p_frame, text='0%', font=('Arial', 10, 'bold'), fg='blue', bg='#f0f0f0')
        self.lbl_progress_val.pack(side='left', padx=10)

        # Frame 2: Post-Analysis (Initially Hidden)
        self.post_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        # We don't pack self.post_p_frame here
        tk.Label(self.post_p_frame, text='Post-Analysis Progress', font=('Arial', 9, 'bold'), bg="#f0f0f0").pack(side='left', padx=10)
        self.progress_post = ttk.Progressbar(self.post_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress_post.pack(side='left', padx=5, pady=5)
        self.lbl_post_progress_val = tk.Label(self.post_p_frame, text='Waiting...', font=('Arial', 10, 'bold'), fg='dimgray', bg="#f0f0f0")
        self.lbl_post_progress_val.pack(side='left', padx=10)

        self.fig = Figure(figsize=(8, 8), dpi=100)
        self.ax = self.fig.add_subplot(111)
        self.ax.axis('off')
        self.canvas = FigureCanvasTkAgg(self.fig, master=self.canvas_frame)
        self.canvas.get_tk_widget().pack(fill='both', expand=True)

        self.canvas.mpl_connect('button_press_event', self.on_click)
        self.canvas.mpl_connect('key_press_event', self.on_key)

    def load_directory(self):
        """
        Opens a file-picker dialog so the user can select any image in a Z-stack folder.

        Discovers all supported image files (``*.tif``, ``*.tiff``, ``*.png``, ``*.jpg``,
        ``*.jpeg``) in the same directory as the selected file, applies natural sort order
        (so ``z01`` comes before ``z10``), synchronises the Z-slice slider to the
        selected file, and calls :meth:`load_image` to display the first slice.
        """
        initial = CONFIG.get('INPUT_DIR', os.getcwd())
        
        # Use file picker so user can see images
        fpath = filedialog.askopenfilename(
            initialdir=str(initial),
            title="Select any image in the stack",
            filetypes=[("Image files", "*.tif *.tiff *.png *.jpg *.jpeg"), ("All files", "*.*")]
        )
        
        if not fpath:
            return
            
        selected_file = pl.Path(fpath)
        p = selected_file.parent
        self.input_dir = str(p)
        
        exts = ['.tif', '.tiff', '.png', '.jpg', '.jpeg']
        found_files_path = []
        
        # 1. Discover all images in the same folder
        for ext in exts:
            found_files_path.extend(list(p.glob(f"*{ext}")))
            found_files_path.extend(list(p.glob(f"*{ext.upper()}")))
            
        # 2. Recursive fallback (if needed, though file picker implies they are in the right spot)
        if not found_files_path:
            for ext in exts:
                found_files_path.extend(list(p.rglob(f"*{ext}")))
                found_files_path.extend(list(p.rglob(f"*{ext.upper()}")))
        
        if not found_files_path:
            messagebox.showerror('Error', f"No supported images found in:\n{p}")
            return
            
        # Standardize and Natural Sort
        found_files_str = [os.path.abspath(str(f)) for f in found_files_path]
        unique_files = list(set(found_files_str))
        self.files = sorted(unique_files, key=natural_sort_key)
            
        # Sync to the selected file
        try:
            self.current_idx = self.files.index(os.path.abspath(fpath))
        except ValueError:
            self.current_idx = 0
            
        self.scale_z.config(to=len(self.files) - 1)
        self.scale_z.set(self.current_idx)
        self.reset_roi(redraw=False)
        self.load_image()
        self.lbl_status.config(text=f'Opened: {selected_file.name}\n({len(self.files)} slices in folder)', fg='blue')
        self.root.update()

    def load_image(self):
        if not self.files:
            return
        try:
            self.current_img = robust_imread(self.files[self.current_idx])
            self.lbl_z.config(text=f'Z: {self.current_idx} / {len(self.files)-1}')
            self.render()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load image:\n{self.files[self.current_idx]}\n\nError: {e}")

    def on_slide_change(self, val):
        self.current_idx = int(val)
        self.load_image()

    def render(self):
        self.ax.clear()
        self.ax.axis('off')
        
        # --- NEW: Review Overlays Logic ---
        if self.mode_var.get() == 'review':
            try:
                if not hasattr(self, 'last_out_dir') or not self.last_out_dir:
                    self.ax.text(0.5, 0.5, "No Batch Analysis Results Found.\nRun Batch First.", 
                                 ha='center', va='center', color='red', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return
                
                z_idx = extract_z_index(self.files[self.current_idx])
                panel_path = os.path.join(self.last_out_dir, "overlays", f"z{z_idx:02d}_panel.png")
                
                if os.path.exists(panel_path):
                    if _HAVE_CV2:
                        img = _cv2.imread(panel_path)
                        img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
                    else:
                        img = plt.imread(panel_path)
                    self.ax.imshow(img)
                    self.canvas.draw_idle()
                    return
                else:
                    self.ax.text(0.5, 0.5, f"Overlay not found for Z={z_idx:02d}\n{os.path.basename(panel_path)}", 
                                 ha='center', va='center', color='orange', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return
            except Exception as e:
                print(f"Overlay Render Error: {e}")
                
        if self.current_img is not None and isinstance(self.current_img, np.ndarray):
            img = self.current_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            self.ax.imshow(disp, cmap='gray')
            if len(self.roi_points) > 0:
                pts = np.array(self.roi_points)
                # If active (closed), draw it closed. Otherwise, draw the open line and endpoints.
                if self.roi_active:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=2)
                else:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=1.5)
                    self.ax.plot(pts[:,0], pts[:,1], 'ro', markersize=4)
            # If a loaded mask exists, always redraw the red contour so it persists across slices
            elif self._loaded_roi_mask is not None:
                self.ax.contour(self._loaded_roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def on_click(self, event):
        if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
            return
        if self.mode_var.get() == 'roi':
            # Left click = add point
            if event.button == 1:
                # If they start clicking after already having a finished ROI, wipe it to start fresh
                if self.roi_active:
                    self.reset_roi(redraw=False)
                self.roi_points.append([event.xdata, event.ydata])
                self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                self.render()
            
            # Right click = undo last point
            elif event.button == 3:
                if not self.roi_active and len(self.roi_points) > 0:
                    self.roi_points.pop()
                    self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                    self.render()

    def finalize_roi(self):
        if not self.roi_active and len(self.roi_points) > 2:
            self.roi_points.append(self.roi_points[0])
            self.roi_active = True
            self.lbl_roi.config(text=f'ROI: active ({len(self.roi_points)-1} points)')
            self.render()
        elif not self.roi_active and len(self.roi_points) > 0:
            messagebox.showwarning('Draw ROI', 'Please place at least 3 points before finalizing.')

    def on_key(self, event):
        if self.mode_var.get() == 'roi' and event.key == 'enter':
            self.finalize_roi()

    def reset_roi(self, redraw=True):
        self.roi_points = []
        self.roi_active = False
        self._loaded_roi_mask = None
        self.lbl_roi.config(text='ROI: none')
        if redraw:
            self.render()

    def build_roi_mask(self):
        # If a mask was loaded from file, use it directly
        if self._loaded_roi_mask is not None and self.current_img is not None:
            return self._loaded_roi_mask
        if not self.roi_active or len(self.roi_points) < 4 or self.current_img is None:
            return None
        h, w = self.current_img.shape
        yy, xx = np.mgrid[:h, :w]
        pts = np.column_stack((xx.ravel(), yy.ravel()))
        path = Path(self.roi_points)
        return path.contains_points(pts).reshape(h, w)

    def save_roi_mask(self):
        mask = self.build_roi_mask()
        if mask is None:
            messagebox.showinfo('Save ROI', 'No active ROI to save.')
            return
        default = f'roi_z{self.current_idx:02d}.npy'
        path = filedialog.asksaveasfilename(defaultextension='.npy', initialfile=default,
                                            filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        np.save(path, mask.astype(np.uint8))
        messagebox.showinfo('Save ROI', f'Saved ROI mask to:\n{path}')

    def load_roi_mask(self):
        if self.current_img is None:
            messagebox.showinfo('Load ROI', 'Load an image first.')
            return
        path = filedialog.askopenfilename(filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        mask = np.load(path).astype(bool)
        if mask.shape != self.current_img.shape:
            messagebox.showerror('Load ROI', f'Mask shape {mask.shape} does not match image shape {self.current_img.shape}')
            return
        # Convert mask boundary to points for display
        ys, xs = np.where(mask)
        if len(xs) == 0:
            messagebox.showerror('Load ROI', 'Loaded ROI mask is empty.')
            return
        self.roi_points = []
        self.roi_active = True
        self._loaded_roi_mask = mask  # Store the loaded mask for build_roi_mask!
        self.lbl_roi.config(text=f'ROI: loaded mask\n{os.path.basename(path)}')
        self.render()
        self.ax.contour(mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def run_analysis_slice(self):
        if self.current_img is None:
            messagebox.showinfo('Info', 'No image loaded. Load a directory first.')
            return

        self.lbl_roi.config(text='Running analysis...')
        self.root.update_idletasks()

        try:
            import sys, time as _t
            log_lines = []
            def log(msg):
                log_lines.append(msg)
                print(msg)
                sys.stdout.flush()

            log(f"\n--- GUI Analysis: slice {self.current_idx} ---")
            log(f"  File: {os.path.basename(self.files[self.current_idx]) if self.files else 'N/A'}")
            log(f"  Image shape: {self.current_img.shape}, dtype: {self.current_img.dtype}")

            params = CONFIG.copy()
            params['SAVE_DEBUG_IMAGES'] = False
            roi_mask = self.build_roi_mask()

            # Determine image to process
            full_img = self.current_img
            crop_offset_y, crop_offset_x = 0, 0
            H, W = full_img.shape

            # ═══════════════════════════════════════════════════════════════════
            # TWO-PASS HYBRID SEGMENTATION ENGINE
            # Pass 1: Full-frame analysis (correct global percentiles for sparse outer nuclei)
            # Pass 2: Tight ROI crop analysis (local percentiles naturally separate dense clusters)
            # Merge: Spatially deduplicate overlapping centroids, keeping unique detections from both
            # ═══════════════════════════════════════════════════════════════════

            t0 = _t.time()

            # ── PASS 1: Full-frame analysis ──────────────────────────────────
            log(f"  PASS 1: Full-frame analysis...")
            seg1 = segment_slice(full_img, params, roi_mask=roi_mask, z_idx=self.current_idx)
            meas1 = measure_spermatids(seg1, params)
            results1 = meas1['results']
            skel_label1 = meas1['skel_label']

            # Filter Pass 1 results to ROI
            if roi_mask is not None and results1:
                filtered1 = []
                keep1 = []
                for r in results1:
                    cy = min(max(int(round(r['centroid_y'])), 0), roi_mask.shape[0] - 1)
                    cx = min(max(int(round(r['centroid_x'])), 0), roi_mask.shape[1] - 1)
                    if roi_mask[cy, cx]:
                        filtered1.append(r)
                        keep1.append(r['label'])
                results1 = filtered1
                skel_label1 = np.where(np.isin(skel_label1, keep1), skel_label1, 0).astype(np.int32)
            log(f"  PASS 1: {len(results1)} detections (full-frame statistics)")

            # ── PASS 2: Tight ROI crop analysis ──────────────────────────────
            results2_global = []
            skel_label2_full = np.zeros((H, W), dtype=np.int32)

            if roi_mask is not None:
                log(f"  PASS 2: Tight-crop local analysis...")
                ys, xs = np.where(roi_mask)
                y0, y1 = int(ys.min()), int(ys.max()) + 1
                x0, x1 = int(xs.min()), int(xs.max()) + 1
                pad = 10
                y0c = max(0, y0 - pad)
                x0c = max(0, x0 - pad)
                y1c = min(H, y1 + pad)
                x1c = min(W, x1 + pad)

                crop_img = full_img[y0c:y1c, x0c:x1c]
                crop_roi = roi_mask[y0c:y1c, x0c:x1c]

                seg2 = segment_slice(crop_img, params, roi_mask=crop_roi, z_idx=self.current_idx)
                meas2 = measure_spermatids(seg2, params)
                results2 = meas2['results']
                skel_label2_crop = meas2['skel_label']

                # Filter Pass 2 to ROI
                if results2:
                    filtered2 = []
                    keep2 = []
                    for r in results2:
                        cy = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0] - 1)
                        cx = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1] - 1)
                        if crop_roi[cy, cx]:
                            filtered2.append(r)
                            keep2.append(r['label'])
                    results2 = filtered2
                    skel_label2_crop = np.where(np.isin(skel_label2_crop, keep2), skel_label2_crop, 0).astype(np.int32)

                # Map Pass 2 back to full-image coordinates
                ch2, cw2 = skel_label2_crop.shape
                max_label1 = int(skel_label1.max()) if skel_label1.max() > 0 else 0
                skel_label2_crop_shifted = np.where(skel_label2_crop > 0, skel_label2_crop + max_label1, 0).astype(np.int32)
                skel_label2_full[y0c:y0c+ch2, x0c:x0c+cw2] = skel_label2_crop_shifted

                for r in results2:
                    r['centroid_x'] += x0c
                    r['centroid_y'] += y0c
                    r['label'] += max_label1
                results2_global = results2
                log(f"  PASS 2: {len(results2_global)} detections (local-crop statistics)")
            else:
                log(f"  PASS 2: Skipped (no ROI)")

            # ── MERGE: Spatial deduplication ──────────────────────────────────
            # For each Pass 2 detection, check if Pass 1 already found a spermatid
            # within 5 pixels of the same centroid. If not, it's a new unique detection.
            dedup_radius = 5.0  # pixels
            merged_results = list(results1)
            skel_label_full = skel_label1.copy()
            new_from_pass2 = 0

            for r2 in results2_global:
                cx2, cy2 = r2['centroid_x'], r2['centroid_y']
                is_duplicate = False
                for r1 in results1:
                    dx = r1['centroid_x'] - cx2
                    dy = r1['centroid_y'] - cy2
                    if (dx*dx + dy*dy) < dedup_radius * dedup_radius:
                        is_duplicate = True
                        break
                if not is_duplicate:
                    merged_results.append(r2)
                    # Splice the Pass 2 skeleton pixels into the merged map
                    mask2 = (skel_label2_full == r2['label'])
                    skel_label_full[mask2] = r2['label']
                    new_from_pass2 += 1

            results = merged_results
            log(f"  MERGE: {len(results1)} (Pass1) + {new_from_pass2} unique from Pass2 = {len(results)} total")

            elapsed = _t.time() - t0
            log(f"  RESULT: {len(results)} spermatids detected ({elapsed:.1f}s)")

            overlay = make_overlay(full_img, skel_label_full)

            # Write log to file
            try:
                log_path = os.path.join(CONFIG['OUTPUT_DIR'], 'gui_analysis_log.txt')
                ensure_dir(CONFIG['OUTPUT_DIR'])
                with open(log_path, 'a', encoding='utf-8') as f:
                    f.write('\n'.join(log_lines) + '\n\n')
            except Exception:
                pass

            # Show results popup
            top = tk.Toplevel(self.root)
            top.title(f'Results Z={self.current_idx} - {len(results)} spermatids')
            top.geometry('1200x650')

            fig = Figure(figsize=(14, 6))
            ax1 = fig.add_subplot(121)
            ax2 = fig.add_subplot(122)

            img = full_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            ax1.imshow(disp, cmap='gray')
            if roi_mask is not None:
                ax1.contour(roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.2)
            ax1.set_title('Original + ROI')
            ax1.axis('off')

            ax2.imshow(overlay)
            ax2.set_title(f'Overlay (N={len(results)})')
            ax2.axis('off')

            um = params['UM_PER_PX_XY']
            for r in results:
                ax2.text(r['centroid_x'], r['centroid_y'],
                         f"{r['length_px_geodesic'] * um:.1f}",
                         color='white', fontsize=5, ha='center', va='center')

            fig.tight_layout()
            can = FigureCanvasTkAgg(fig, master=top)
            can.get_tk_widget().pack(fill='both', expand=True)
            can.draw()

            if results:
                lengths = [r['length_px_geodesic'] * um for r in results]
                text = f'Found {len(results)} spermatids | median length {np.median(lengths):.2f} um ({elapsed:.1f}s)'
            else:
                text = f'Found 0 spermatids ({elapsed:.1f}s) - see gui_analysis_log.txt for diagnostics'
            lbl_stats = tk.Label(top, text=text, font=('Arial', 11))
            lbl_stats.pack(pady=4)

            lbl_tool = tk.Label(top, text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue', font=('Arial', 10, 'bold'))
            lbl_tool.pack(pady=2)

            self.lbl_roi.config(text=f'Analysis done: {len(results)} spermatids')

            # ── INTERACTIVE MANUAL CORRECTION LOGIC ──
            class ManualCorrector:
                def __init__(self, canvas, ax_overlay, seg_data, prms, crop_oy, crop_ox, fimg):
                    self.canvas = canvas
                    self.ax = ax_overlay
                    self.seg = seg_data
                    self.params = prms
                    self.crop_oy = crop_oy
                    self.crop_ox = crop_ox
                    self.fimg = fimg
                    
                    self.active_tool = None
                    self.cid_press = self.canvas.mpl_connect('button_press_event', self.on_click)
                    self.cid_key = self.canvas.mpl_connect('key_press_event', self.on_key)
                    
                    self.overlay_imshow = None
                    self.text_artists = []
                    
                def on_key(self, event):
                    if event.key == 'e':
                        self.active_tool = 'erase'
                        lbl_tool.config(text="Active Tool: ERASE (Click a colored spermatid to delete it)", fg='red')
                    elif event.key == 's':
                        self.active_tool = 'split'
                        lbl_tool.config(text="Active Tool: SPLIT (Click on a skeleton branch to sever it)", fg='orange')
                    elif event.key == 'escape':
                        self.active_tool = None
                        lbl_tool.config(text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue')

                def on_click(self, event):
                    if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
                        return
                    if not self.active_tool:
                        return

                    # Map global click to cropped coordinate space
                    x_crop = int(round(event.xdata)) - self.crop_ox
                    y_crop = int(round(event.ydata)) - self.crop_oy
                    
                    ch, cw = self.seg['skel_pruned'].shape
                    if not (0 <= x_crop < cw and 0 <= y_crop < ch):
                        return # Clicked outside the ROI working area

                    modified = False
                    if self.active_tool == 'erase':
                        # Find the label at this pixel in the current skeleton
                        lab = self.seg['skel_labeled'][y_crop, x_crop]
                        # If the user clicked slightly off, search a 3x3 neighborhood
                        if lab == 0:
                            for dy in (-1, 0, 1):
                                for dx in (-1, 0, 1):
                                    yy, xx = y_crop+dy, x_crop+dx
                                    if 0 <= yy < ch and 0 <= xx < cw and self.seg['skel_labeled'][yy, xx] > 0:
                                        lab = self.seg['skel_labeled'][yy, xx]
                                        break
                        if lab > 0:
                            self.seg['skel_pruned'][self.seg['skel_labeled'] == lab] = False
                            modified = True
                            
                    elif self.active_tool == 'split':
                        # Draw a small black circle to sever the topological skeleton connection
                        from skimage.draw import disk
                        rr, cc = disk((y_crop, x_crop), radius=2.5, shape=self.seg['skel_pruned'].shape)
                        self.seg['skel_pruned'][rr, cc] = False
                        modified = True
                        
                    if modified:
                        self.recalculate_and_redraw()
                        
                def recalculate_and_redraw(self):
                    # Re-label the modified skeleton
                    self.seg['skel_labeled'] = measure.label(self.seg['skel_pruned'])
                    
                    # Re-run measurement filters
                    new_meas = measure_spermatids(self.seg, self.params)
                    new_results = new_meas['results']
                    
                    # Filter by ROI inside centroid logic
                    if crop_roi is not None:
                        filtered = []
                        keep_labels = []
                        for r in new_results:
                            c_y = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0] - 1)
                            c_x = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1] - 1)
                            if crop_roi[c_y, c_x]:
                                filtered.append(r)
                                keep_labels.append(r['label'])
                        new_results = filtered
                        skel_l = np.where(np.isin(new_meas['skel_label'], keep_labels), new_meas['skel_label'], 0).astype(np.int32)
                    else:
                        skel_l = new_meas['skel_label']

                    # Map back to full image
                    H, W = self.fimg.shape
                    skel_label_full = np.zeros((H, W), dtype=np.int32)
                    ch, cw = skel_l.shape
                    skel_label_full[self.crop_oy:self.crop_oy+ch, self.crop_ox:self.crop_ox+cw] = skel_l

                    # Redraw Overlay
                    new_overlay = make_overlay(self.fimg, skel_label_full)
                    
                    self.ax.clear()
                    self.ax.imshow(new_overlay)
                    self.ax.set_title(f'Overlay (N={len(new_results)}) - Manual Corrections Applied')
                    self.ax.axis('off')

                    # Redraw Text
                    _um = self.params['UM_PER_PX_XY']
                    for r in new_results:
                        self.ax.text(r['centroid_x'] + self.crop_ox, r['centroid_y'] + self.crop_oy,
                                 f"{r['length_px_geodesic'] * _um:.1f}",
                                 color='white', fontsize=5, ha='center', va='center')
                                 
                    self.canvas.draw()
                    
                    # Update Stats Label
                    if new_results:
                        lengths = [r['length_px_geodesic'] * _um for r in new_results]
                        lbl_stats.config(text=f'Corrected: {len(new_results)} spermatids | median length {np.median(lengths):.2f} um')
                    else:
                        lbl_stats.config(text=f'Corrected: 0 spermatids')

            # Attach to popup so it doesn't get garbage collected
            top.corrector = ManualCorrector(can, ax2, seg1, params, crop_offset_y, crop_offset_x, full_img)

        except Exception as e:
            import traceback
            traceback.print_exc()
            sys.stdout.flush()
            self.lbl_roi.config(text=f'Analysis error: {e}')
            messagebox.showerror('Analysis Error', f'{type(e).__name__}: {e}')

    def run_batch_analysis(self):
        if not self.files:
            messagebox.showinfo('Info', 'No directory loaded.')
            return

        # Auto-incremental output directory inside selected folder
        out_dir = get_unique_batch_dir(self.input_dir)
        self.last_out_dir = out_dir
        
        # EXPLICIT CONFIRMATION: Show the user where the data will go
        confirm = messagebox.askokcancel("Confirm Output", 
            f"Results (Excel, PDF, CSV) will be saved to:\n\n{out_dir}\n\nContinue?")
        if not confirm:
            return
            
        ensure_dir(out_dir)
        overlay_dir = os.path.join(out_dir, "overlays")
        ensure_dir(overlay_dir)

        params = CONFIG.copy()
        params['OUTPUT_DIR'] = out_dir
        params['SAVE_DEBUG_IMAGES'] = False
        params['DO_TRACKING'] = True
        
        roi_mask = self.build_roi_mask()

        self.lbl_roi.config(text="Processing... See Top Bar")
        self.root.update_idletasks()
        
        try:
            t_batch = _t.time()
            self.lbl_batch_op.config(text=f"Batch Segmenting: 0 / {len(self.files)} slices...", fg='blue')
            self.root.update()
            
            all_rows = []
            summaries = []
            
            # Z-Projection accumulation
            max_proj_raw = None
            max_proj_ov = None
            
            # Robust initialization
            ts = None
            df_trk = None

            self.progress['value'] = 0
            self.progress['maximum'] = len(self.files)
            self.progress_post['value'] = 0
            self.progress_post['maximum'] = 100
            self.lbl_post_progress_val.config(text="0%", fg='orange')
            
            for idx, fpath in enumerate(self.files):
                z_idx = extract_z_index(fpath, sequence_idx=idx)
                
                pct = int(((idx + 1) / len(self.files)) * 100)
                self.progress['value'] = idx + 1
                self.lbl_progress_val.config(text=f"{pct}%")
                self.root.update()
                
                print(f"[{idx+1}/{len(self.files)}] Processing Z={z_idx:02d}...")
                
                img_raw = robust_imread(fpath)
                
                # Handling ROI cropping to speed up processing
                full_img = img_raw
                crop_oy, crop_ox = 0, 0
                
                # ── BUGFIX: NEVER crop the image before analysis ──
                # Cropping to a tight ROI artificially removes the dark background,
                # which completely destroys the np.percentile() thresholding math that
                # relies on the global image statistics!
                process_img = full_img
                crop_roi = roi_mask
                seg = segment_slice(process_img, params, z_idx=z_idx)
                meas = measure_spermatids(seg, params)
                res = meas['results']
                sl_crop = meas['skel_label']
                
                if crop_roi is not None and res:
                    filtered = []
                    keep = []
                    for r in res:
                        cy = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0]-1)
                        cx = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1]-1)
                        if crop_roi[cy, cx]:
                            filtered.append(r)
                            keep.append(r['label'])
                    res = filtered
                    sl_crop = np.where(np.isin(sl_crop, keep), sl_crop, 0).astype(np.int32)
                
                for r in res:
                    r['centroid_x'] += crop_ox
                    r['centroid_y'] += crop_oy
                    
                H, W = full_img.shape
                sl_full = np.zeros((H, W), dtype=np.int32)
                ch, cw = sl_crop.shape
                sl_full[crop_oy:crop_oy+ch, crop_ox:crop_ox+cw] = sl_crop
                
                um = params['UM_PER_PX_XY']
                ls_um = [r['length_px_geodesic']*um for r in res]
                
                all_rows.extend(rows_from_results(res, z_idx, um))
                summaries.append({
                    "z_slice": z_idx,
                    "n_spermatids": len(res),
                    "median_length_um": round(float(np.median(ls_um)), 3) if ls_um else 0,
                })
                
                if params['SAVE_OVERLAYS']:
                    ov = make_overlay(full_img, sl_full)
                    # Create side-by-side panel
                    orig_rgb = (normalize_display(full_img) * 255).astype(np.uint8)
                    if orig_rgb.ndim == 2:
                        orig_rgb = np.stack([orig_rgb]*3, axis=-1)
                    panel = np.hstack([orig_rgb, ov])
                    _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_panel.png"), panel)
                    
                    # ---- LIVE GUI UPDATE ----
                    # Show the side-by-side segmentation panel during batch execution
                    if hasattr(self, 'ax') and hasattr(self, 'canvas'):
                        try:
                            self.ax.clear()
                            self.ax.axis('off')
                            self.ax.imshow(panel)
                            self.canvas.draw()
                        except Exception:
                            pass
                    
                    # Update Z-Projections
                    if max_proj_raw is None:
                        max_proj_raw = img_raw.copy().astype(np.float32)
                        max_proj_ov = ov.copy().astype(np.float32)
                    else:
                        max_proj_raw = np.maximum(max_proj_raw, img_raw)
                        max_proj_ov = np.maximum(max_proj_ov, ov.astype(np.float32))

                # Update Progress Bar for each slice
                self.lbl_batch_op.config(text=f"Batch Segmenting: {idx+1} / {len(self.files)} slices...", fg='blue')
                self.progress['value'] = idx + 1
                self.root.update()
            
            df = pd.DataFrame(all_rows)
            df_sum = pd.DataFrame(summaries)
            df.to_csv(os.path.join(out_dir, "spermatid_measurements.csv"), index=False)
            df_sum.to_csv(os.path.join(out_dir, "slice_summary.csv"), index=False)
            
            if not df.empty:
                self.lbl_batch_op.config(text='Running 3D Tracking & Morphometrics...', fg='#e67e22') # Orange-ish
                
                # --- Sequential Progress Bar Swap ---
                # Hide 2D progress, Show Post-Analysis progress
                self.batch_p_frame.pack_forget()
                self.post_p_frame.pack(fill='x')
                
                self.progress_post['value'] = 25
                self.lbl_post_progress_val.config(text="25%")
                self.root.update()

                df_trk, ts = track_across_slices(df, params)
                
                # --- Advanced 3D Morphometrics ---
                self.lbl_batch_op.config(text='Calculating Advanced 3D Metrics...', fg='#e67e22')
                # Metrics are natively generated correctly and safely in track_across_slices

                self.progress_post['value'] = 60
                self.lbl_post_progress_val.config(text="60%")
                self.root.update()

                df_trk.to_csv(os.path.join(out_dir, "measurements_with_tracks.csv"), index=False)
                ts.to_csv(os.path.join(out_dir, "track_summary.csv"), index=False)
                
                # Save Global Z-Projection
                if max_proj_raw is not None:
                    self.lbl_batch_op.config(text='Generating Global Z-Projection...', fg='#e67e22')
                    raw_p = (normalize_display(max_proj_raw.astype(np.uint16)) * 255).astype(np.uint8)
                    if raw_p.ndim == 2: raw_p = np.stack([raw_p]*3, axis=-1)
                    ov_p = max_proj_ov.astype(np.uint8)
                    global_panel = np.hstack([raw_p, ov_p])
                    _imwrite(os.path.join(out_dir, "global_z_projection.png"), global_panel)
                    
                self.progress_post['value'] = 80
                self.lbl_post_progress_val.config(text="80%")
                self.root.update()
            
            elapsed = _t.time() - t_batch
            msg = f"Batch complete in {elapsed:.1f}s!\nSaved to: {out_dir}"
            print(msg)
            self.lbl_batch_op.config(text='Batch Analysis Complete.', fg='green')
            self.lbl_progress_val.config(text="100% - Done", fg='green')
            self.root.update()
            
            # Generate High-Res Graphical Report and Excel Audit
            self.lbl_batch_op.config(text='Generating PDF & Excel Reports...', fg='#8e44ad') # Purple
            
            def update_cb(v):
                self.progress_post['value'] = v
                self.lbl_post_progress_val.config(text=f"{v}%")
                self.root.update()

            generate_batch_report(out_dir, df, df_sum, um, ts if not df.empty else None, update_cb)
            generate_excel_report(out_dir, df, df_sum, ts if not df.empty else None)
            
            self.lbl_batch_op.config(text='Generating PowerPoint Dashboard...', fg='#c0392b') # Red-ish
            self.root.update()
            try:
                ok = generate_pptx_report(out_dir, df, df_sum, um, ts if not df.empty else None)
                if not ok:
                    print("WARNING: PPTX generation returned False — check console for traceback.")
            except Exception as pptx_err:
                import traceback as _tb
                print(f"ERROR generating PPTX: {pptx_err}")
                _tb.print_exc()
            
            msg = f"Batch complete in {elapsed:.1f}s!\nSaved to: {out_dir}"
            print(msg)
            self.lbl_batch_op.config(text='ALL OPERATIONS COMPLETE', fg='green')
            self.lbl_progress_val.config(text="100%", fg='green')
            self.lbl_post_progress_val.config(text="100% - DONE", fg='green')
            self.progress_post['value'] = 100
            self.root.update()

            messagebox.showinfo('Batch Complete', msg)
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            self.lbl_batch_op.config(text=f'Batch error: {e}', fg='red')
            messagebox.showerror('Batch Error', str(e))



def launch_gui():
    if not _TK_AVAILABLE:
        raise RuntimeError("Tkinter GUI components are not available in this Python environment.")
    root = tk.Tk()
    app = SpermGUI(root)
    root.mainloop()


# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    print(f"\n{'='*60}")
    print(f" SPERMATID ANALYSIS PIPELINE - v{_VERSION} - INITIALIZING...")
    print(f" Mode: {'GUI' if '--gui' in sys.argv or len(sys.argv)<=1 else 'CLI'}")
    print(f" Features: [Incremental Folders] [Multi-Tab Excel] [3D Stats]")
    print(f"{'='*60}\n")
    
    ap = argparse.ArgumentParser(description=f"Spermatid segmentation {_VERSION} / ROI GUI v11")
    ap.add_argument("--batch",  action="store_true", help="Run batch processing")
    ap.add_argument("--single", action="store_true", help="Run a single-image analysis")
    ap.add_argument("--gui",    action="store_true", help="Launch ROI GUI")
    ap.add_argument("--z",      type=int, default=None, help="Choose z-index in single mode")
    args = ap.parse_args()

    validate_config(CONFIG)

    # Launch GUI by default if no explicit CLI flags are provided
    if args.gui or not (args.batch or args.single or args.z is not None):
        launch_gui()
        raise SystemExit

    if args.batch:
        CONFIG["RUN_MODE"] = "batch"
        # CLI Incremental Folder Logic
        if not os.path.isabs(CONFIG["OUTPUT_DIR"]):
            # If default, make it relative to execution dir
            CONFIG["OUTPUT_DIR"] = get_unique_batch_dir(os.getcwd())
        else:
            # If absolute path provided in config, make unique relative to parent
            CONFIG["OUTPUT_DIR"] = get_unique_batch_dir(os.path.dirname(CONFIG["OUTPUT_DIR"]))
        
        ensure_dir(CONFIG["OUTPUT_DIR"])
        print(f"CLI BATCH MODE: Results will be saved to: {CONFIG['OUTPUT_DIR']}")
    if args.single:
        CONFIG["RUN_MODE"] = "single"
    if args.z is not None:
        CONFIG["SINGLE_IMAGE_SELECTION_MODE"] = "z_index"
        CONFIG["SINGLE_Z_INDEX"] = args.z

    if CONFIG["RUN_MODE"] == "single":
        process_one_image(choose_single_image(CONFIG), CONFIG, CONFIG["OUTPUT_DIR"])
    else:
        process_batch(CONFIG)
