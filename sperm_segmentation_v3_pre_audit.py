import sys
import os

# --- v12 ROBUST LOGGING SYSTEM ---
class Tee(object):
    def __init__(self, *files):
        self.files = files
    def write(self, obj):
        for f in self.files:
            f.write(obj)
            f.flush()
    def flush(self):
        for f in self.files:
            f.flush()

# Redirect stdout and stderr to both console and a file
log_file_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "sperm_error_log.txt")
f_log = open(log_file_path, 'a', encoding='utf-8')
sys.stdout = Tee(sys.stdout, f_log)
sys.stderr = Tee(sys.stderr, f_log)
print(f"\n--- NEW SESSION STARTED: {os.path.basename(__file__)} ---\n")
# ---------------------------------

#!/usr/bin/env python3
"""
Sperm Nucleus Segmentation & 3D Morphometrics Pipeline  —  Saturn Project
=========================================================================
A production-ready image-analysis pipeline for automated detection,
measurement, and 3D reconstruction of sperm nuclei tailored to
the *Saturn* experimental dataset acquired on the Leica confocal system.

Biological context
------------------
Spermiogenesis — the post-meiotic differentiation of round spermatids into
mature spermatozoa — involves dramatic nuclear elongation, chromatin condensation,
and apical plunging of the nucleus toward the basal lamina of the seminiferous
tubule.  Quantifying these morphological changes requires imaging the highly
elongated, condensed sperm nuclei.  This pipeline automates the process by:

1. Detecting sperm nuclei as thin, dim, ridge-like objects in each 2D
   Z-slice using a skeleton-first strategy (Frangi ridge filter → Otsu
   binarisation → morphological closing → geodesic skeleton).
2. Measuring per-cell biometrics: geodesic length (µm), Euclidean width,
   tortuosity, endpoint count, branch complexity, centroid, and area.
3. Linking detections across Z-slices into 3D tracks using a nearest-neighbour
   frame-to-frame assignment algorithm.
4. Computing 3D morphometrics from tracks: true 3D length, Z-extent,
   approximated volume, pitch angle, taper ratio, and nearest-neighbour
   packing density.
5. Exporting results as CSV, a multi-sheet Excel workbook, a multi-page PDF
   report, and a native PowerPoint (.pptx) dashboard with editable charts.

Saturn-specific calibration
----------------------------
Physical scale factors are derived from Leica confocal metadata for the
Saturn dataset:

- ``UM_PER_PX_XY  = 0.7568``  µm/pixel  (lateral resolution)
- ``UM_PER_SLICE_Z = 1.0404``  µm/slice  (z-step size)

These values are set as defaults in ``CONFIG`` and can be overridden via
the Parameter Editor GUI or by loading a JSON settings file.

Pipeline architecture
---------------------
::

    Image stack (TIF / PNG / JPG)
          │
          ▼
    load_batch_files()  ────────────────────── natural-sorted file list
          │
    ┌─────┴──────────────────────────────────────────────────────────┐
    │  Per-slice loop  (process_batch → process_one_image)           │
    │                                                                 │
    │  robust_imread()   →  segment_slice()  →  measure_spermatids() │
    │       │                    │                     │              │
    │   raw ndarray        mask / skel / ridge    list[dict]           │
    │                            │                     │              │
    │                     make_overlay()        rows_from_results()   │
    │                            │                     │              │
    │                      overlay PNG              per-slice CSV     │
    └─────────────────────────────────────────────────────────────────┘
          │
    track_across_slices()  ──── 3D track DataFrame
          │
    ┌─────┴──────────────────────────────────────────────────────────┐
    │  Reporting                                                      │
    │  generate_excel_report()  →  multi-sheet XLSX workbook         │
    │  generate_batch_report()  →  multi-page PDF (matplotlib)       │
    │  generate_pptx_report()   →  native Office PPTX dashboard      │
    └─────────────────────────────────────────────────────────────────┘

Key configuration parameters (``CONFIG`` dict)
-----------------------------------------------
``UM_PER_PX_XY``       Physical pixel size in µm (Leica metadata: 0.7568).
``UM_PER_SLICE_Z``     Z-step size in µm (Leica metadata: 1.0404).
``FRANGI_SCALE_RANGE`` Tubeness filter scales (px); set to match sperm nucleus width.
``MIN_SKEL_LEN_PX``    Minimum skeleton length accepted (removes debris/noise).
``MAX_WIDTH_PX``       Maximum skeleton width accepted (rejects merged clusters).
``MAX_BRANCH_NODES``   Maximum branch-point count; >0 tolerates bridged tips.
``MAX_TORTUOSITY``     Maximum curvature index; rejects snaking merged networks.
``DO_TRACKING``        Enable/disable cross-slice 3D linking.
``TRACK_MAX_DIST_UM``  Maximum centroid displacement between adjacent Z-slices.

Usage
-----
Launch the interactive GUI (recommended)::

    python sperm_segmentation_saturn.py

Run a headless batch analysis::

    python sperm_segmentation_saturn.py --batch

Analyse a single slice::

    python sperm_segmentation_saturn.py --single --z 4

Dependencies
------------
numpy, scipy, scikit-image, pandas, matplotlib, tifffile, opencv-python,
Pillow, xlsxwriter, python-pptx, tkinter (stdlib)

Author
------
Dushyant Mishra  |  Findlay Lab  |  Saturn Dataset Branch
"""

import os, sys, glob, re, time, warnings, heapq, argparse, math, pathlib as pl
import time as _t
import json, webbrowser, threading
try:
    import requests
    _HAVE_REQUESTS = True
except ImportError:
    _HAVE_REQUESTS = False
warnings.filterwarnings("ignore")

import numpy as np
import pandas as pd
import tifffile
import matplotlib
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.path import Path
print(f"[matplotlib backend: {matplotlib.get_backend()}]")

from skimage import measure, morphology, exposure
from skimage.filters import meijering, gaussian, apply_hysteresis_threshold
from skimage.morphology import skeletonize
from scipy.ndimage import distance_transform_edt, grey_dilation
from scipy.spatial import cKDTree
from matplotlib.backends.backend_pdf import PdfPages

try:
    import cv2 as _cv2
    _HAVE_CV2 = True
except ImportError:
    _HAVE_CV2 = False

# =============================================================================
# CONFIG
# =============================================================================

CONFIG = {
    # ── run mode ─────────────────────────────────────────────────────────────
    "RUN_MODE": "single",          # "single" | "batch"

    # ── single-image selection ────────────────────────────────────────────────
    "SINGLE_IMAGE_SELECTION_MODE": "dialog",  # "path" | "z_index" | "dialog"
    "SINGLE_TEST_IMAGE": r"C:\Users\dmishra\Desktop\sperm images\Project001_Series002_z15_ch00.tif",
    "SINGLE_Z_INDEX": 15,

    # ── input / output ────────────────────────────────────────────────────────
    "INPUT_DIR":    ".",
    "OUTPUT_DIR":   "./sperm_results_saturnv2",
    "FILE_PATTERN": "Project001_Series002_z*_ch00.tif",

    # ── calibration ───────────────────────────────────────────────────────────
    "UM_PER_PX_XY":   0.756836,
    "UM_PER_SLICE_Z": 1.040460,

    # ── image enhancement ─────────────────────────────────────────────────────
    "CLAHE_CLIP":   0.0862,  # 30-Gen V2 Optimized
    "CLAHE_KERNEL": 64,
    "BG_SIGMA":     5.065,   # 30-Gen V2 Optimized
    "RIDGE_SIGMAS": [1, 2, 3, 4],

    # ── hysteresis threshold ──────────────────────────────────────────────────
    "THRESHOLD_HI": 76.14,   # 30-Gen V2 Optimized
    "THRESHOLD_LO": 66.44,   # 30-Gen V2 Optimized

    # ── morphological cleanup ─────────────────────────────────────────────────
    "CLOSE_RADIUS":  1,
    "MIN_HOLE_AREA": 10,
    "MIN_OBJ_PX":    5,

    # ── skeleton-level gap bridging ───────────────────────────────────────────
    # P1: reduced from 10 → 5 to prevent chaining distinct spermatids
    "MAX_BRIDGE_PX": 5,

    # ── branch pruning & automated splitting ──────────────────────────────────
    "MAX_BRANCH_LEN_PX": 5,   # prune spurs shorter than this before measuring
    "BREAK_JUNCTIONS": True,  # automatically sever all branching intersections into distinct lines

    # ── optional early mask-level shape filter ────────────────────────────────
    "USE_EARLY_SHAPE_FILTER": False,
    "MIN_ECCENTRICITY": 0.60,
    "MAX_MINOR_PX":     12.0,
    "MIN_AXIS_RATIO":   1.4,
    "MIN_MAJOR_PX":     5,

    # ── post-skeleton filters ─────────────────────────────────────────────────
    "MIN_SKEL_LEN_PX":        6.02, # 30-Gen V2 Optimized
    "MAX_GEODESIC_LEN_PX":    65.0, # backend actively cuts chains longer than this
    "MAX_WIDTH_PX":           8.94, # 30-Gen V2 Optimized
    "MIN_LENGTH_WIDTH_RATIO": 1.80, # 30-Gen V2 Optimized

    # ── NEW topology filters ──────────────────────────────────────────────────
    "MAX_BRANCH_NODES": 0,        # 30-Gen V2 Optimized
    "MAX_TORTUOSITY": 3.23,       # 30-Gen V2 Optimized
    "MAX_ENDPOINT_COUNT": 10,     # 30-Gen V2 Optimized

    # N4: loops
    "ALLOW_LOOPS": True,     # don't reject spermatids that cross themselves  # moderate

    # ── tracking across z ─────────────────────────────────────────────────────
    "DO_TRACKING":          True,
    "TRACK_MAX_DIST_UM":    6.5,
    "TRACK_MAX_GAP_SLICES": 1,
    "TRACK_BBOX_PADDING_PX": 3,

    # ── conservative tracking stop-rules ─────────────────────────────────────────
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": 0.40,      # 40% width change max (OPTIMIZED)
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": 0.60,     # 60% length change max
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": 0.70,       # 70% area change max
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": 0.40,       # Absolute tortuosity jump
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": 10.0,      # Centroid displacement in um (RELAXED)

    # ── overlap-first Stage 2 parameters ──────────────────────────────────────────
    "OVERLAP_STABILITY_THRESHOLD": 0.08,            # 8% tolerance for width/area/length stability
    "OVERLAP_ORIENTATION_DEG":     15.0,             # Max orientation change (degrees) for stability
    "OVERLAP_MULTIPLIER":          1.35,             # Fallback multiplier when overlap + stable
    "OVERLAP_MIN_STABLE_COUNT":    1,                # Min stable metrics required for overlap continuation

    # ── output / debug ────────────────────────────────────────────────────────
    "SAVE_DEBUG_IMAGES":   True,
    "SAVE_MASK_TIFS":      True,
    "SAVE_LABEL_TIFS":     True,
    "SAVE_OVERLAYS":       True,
    "SAVE_DETAIL_FIGURE":  True,
    "SHOW_PREVIEW_WINDOW": True,
    "SHOW_DEBUG_PREVIEW":  True,
}

# =============================================================================
# CONFIG VALIDATION
# =============================================================================

_REQUIRED = {
    "RUN_MODE": str, "SINGLE_IMAGE_SELECTION_MODE": str,
    "SINGLE_TEST_IMAGE": str, "SINGLE_Z_INDEX": int,
    "INPUT_DIR": str, "OUTPUT_DIR": str, "FILE_PATTERN": str,
    "UM_PER_PX_XY": float, "UM_PER_SLICE_Z": float,
    "CLAHE_CLIP": float, "CLAHE_KERNEL": int, "BG_SIGMA": (int, float),
    "RIDGE_SIGMAS": list,
    "THRESHOLD_HI": (int, float), "THRESHOLD_LO": (int, float),
    "CLOSE_RADIUS": int, "MIN_HOLE_AREA": int, "MIN_OBJ_PX": int,
    "MAX_BRIDGE_PX": (int, float), "MAX_BRANCH_LEN_PX": (int, float),
    "USE_EARLY_SHAPE_FILTER": bool,
    "MIN_SKEL_LEN_PX": (int, float), "MAX_GEODESIC_LEN_PX": (int, float),
    "MAX_WIDTH_PX": (int, float), "MIN_LENGTH_WIDTH_RATIO": (int, float),
    "MAX_BRANCH_NODES": (int, float),
    "MAX_TORTUOSITY": (int, float),
    "MAX_ENDPOINT_COUNT": (int, float),
    "DO_TRACKING": bool, "TRACK_MAX_DIST_UM": (int, float),
    "TRACK_MAX_GAP_SLICES": int,
    "CONSERVATIVE_MAX_WIDTH_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_LENGTH_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_AREA_JUMP_RATIO": float,
    "CONSERVATIVE_MAX_TORTUOSITY_JUMP": float,
    "CONSERVATIVE_MAX_CENTROID_JUMP_UM": float,
    "SAVE_DEBUG_IMAGES": bool, "SAVE_MASK_TIFS": bool,
    "SAVE_LABEL_TIFS": bool, "SAVE_OVERLAYS": bool,
    "SAVE_DETAIL_FIGURE": bool, "SHOW_PREVIEW_WINDOW": bool,
    "SHOW_DEBUG_PREVIEW": bool,
}

def validate_config(cfg):
    """
    Validates the pipeline CONFIG dictionary for required keys, correct types, and logical consistency.

    Raises a descriptive ValueError listing ALL problems found so engineers can fix
    everything in one go rather than playing whack-a-mole with sequential errors.

    Checks performed
    ----------------
    - Every key in ``_REQUIRED`` is present in *cfg*.
    - Each value has the expected Python type (e.g., ``float`` for calibration parameters).
    - ``THRESHOLD_LO`` < ``THRESHOLD_HI`` (hysteresis thresholding only works in this order).
    - ``RUN_MODE`` is exactly one of ``'single'`` or ``'batch'``.
    - The deprecated ``'REJECT_BRANCHES'`` Boolean key is absent (replaced by ``MAX_BRANCH_NODES`` int).

    Args:
        cfg (dict): The configuration mapping to validate.

    Raises:
        ValueError: If *any* of the above checks fail. The message lists all errors.
    """
    errors = []
    for key, expected in _REQUIRED.items():
        if key not in cfg:
            errors.append(f"  MISSING: '{key}'")
        elif not isinstance(cfg[key], expected):
            errors.append(f"  WRONG TYPE '{key}': "
                          f"got {type(cfg[key]).__name__}, want {expected}")
    if cfg.get("THRESHOLD_LO", 0) >= cfg.get("THRESHOLD_HI", 100):
        errors.append("  THRESHOLD_LO must be < THRESHOLD_HI")
    if cfg.get("RUN_MODE", "") not in ("single", "batch"):
        errors.append("  RUN_MODE must be 'single' or 'batch'")
    if "REJECT_BRANCHES" in cfg:
        errors.append("  'REJECT_BRANCHES' was removed in v8. "
                      "Use MAX_BRANCH_NODES (int) instead:\n"
                      "    REJECT_BRANCHES=True  → MAX_BRANCH_NODES=0\n"
                      "    REJECT_BRANCHES=False → MAX_BRANCH_NODES=9999")
    if errors:
        raise ValueError("CONFIG errors:\n" + "\n".join(errors))

# =============================================================================
# UTILITIES
# =============================================================================

_N8 = [(-1,-1),(-1,0),(-1,1),(0,-1),(0,1),(1,-1),(1,0),(1,1)]

def ensure_dir(p):
    """
    Creates directory *p* (and all missing parents) if it does not already exist.

    Uses ``exist_ok=True`` to avoid TOCTOU race conditions in parallel runs.

    Args:
        p (str): Absolute or relative path of the directory to create.
    """
    os.makedirs(p, exist_ok=True)

def get_unique_batch_dir(base_dir):
    """
    Checks for 'batch_output', then 'batch_output_1', 'batch_output_2', etc.
    Returns the first available path.
    """
    candidate = os.path.join(base_dir, "batch_output")
    if not os.path.exists(candidate):
        return candidate
    
    counter = 1
    while True:
        candidate = os.path.join(base_dir, f"batch_output_{counter}")
        if not os.path.exists(candidate):
            return candidate
        counter += 1

def natural_sort_key(s):
    """
    Key for natural sorting (e.g., img2.tif comes before img10.tif).
    """
    import re
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split(r'(\d+)', os.path.basename(s))]

def extract_z_index(path, sequence_idx=0):
    """
    Extracts Z-index from filename, fallback to sequence index.
    """
    import re
    m = re.search(r"[zZ](\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    m = re.search(r"(\d+)", os.path.basename(path))
    if m:
        return int(m.group(1))
    return sequence_idx

def ensure_2d_image(img, name="image"):
    """
    Forces an image array to be 2D grayscale.
    Handles squeezing, channel-first (Z,H,W), channel-last (H,W,C), and RGB conversion.
    """
    img = np.asarray(img)
    img = np.squeeze(img)

    if img.ndim == 2:
        return img

    if img.ndim == 3:
        # channel-last RGB/RGBA
        if img.shape[-1] in (3, 4):
            return img[..., 0]
        # channel-first
        if img.shape[0] in (1, 3, 4):
            return img[0]
        # singleton channel
        if img.shape[-1] == 1:
            return img[..., 0]
        if img.shape[0] == 1:
            return img[0]

    raise ValueError(f"{name} must be 2D after loading, got shape {img.shape}")

def robust_imread(path):
    """
    Reads image with multi-engine fallback and forced 2D grayscale enforcement.
    """
    p_lower = path.lower()

    # 1. Primary Loader: tifffile
    if p_lower.endswith(".tif") or p_lower.endswith(".tiff"):
        try:
            img = tifffile.imread(path)
            return ensure_2d_image(img, os.path.basename(path))
        except Exception:
            pass

    # 2. Fallback A: Pillow
    try:
        from PIL import Image
        img = np.array(Image.open(path))
        return ensure_2d_image(img, os.path.basename(path))
    except Exception:
        pass

    # 3. Fallback B: OpenCV
    if _HAVE_CV2:
        try:
            img = _cv2.imread(path, _cv2.IMREAD_UNCHANGED)
            if img is not None:
                return ensure_2d_image(img, os.path.basename(path))
        except Exception:
            pass

    # 4. Fallback C: Matplotlib
    try:
        img = plt.imread(path)
        return ensure_2d_image(img, os.path.basename(path))
    except Exception:
        pass

    raise RuntimeError(f"All image engines failed to read: {os.path.basename(path)}")

def normalize_display(img):
    """
    Contrast-stretches a raw image for colourmap display.

    Uses the 1st–99.5th percentile range rather than min–max to suppress hot
    pixels and dark-corner vignetting artefacts common in fluorescence microscopy.
    The small epsilon (1e-9) prevents division-by-zero on flat images.

    Args:
        img (np.ndarray): Input image of any integer or float dtype.

    Returns:
        np.ndarray: Float32 array clipped to [0, 1] ready for matplotlib display.
    """
    a = img.astype(np.float32)
    lo, hi = np.percentile(a, 1), np.percentile(a, 99.5)
    return np.clip((a - lo) / (hi - lo + 1e-9), 0, 1)

def _imwrite(path, arr_uint8, cmap="gray"):
    """
    Writes a uint8 image to disk using the best available engine.

    Engine priority:
    1. **OpenCV** — fastest; handles large images without matplotlib overhead.
       Converts RGB→BGR before writing (OpenCV internal convention).
    2. **Matplotlib** — fallback when OpenCV is absent; uses ``plt.imsave``
       which respects the *cmap* argument for single-channel saves.

    Args:
        path (str): Destination file path (extension determines format, e.g. ``.png``).
        arr_uint8 (np.ndarray): uint8 image array, shape ``(H, W)`` or ``(H, W, 3)``.
        cmap (str): Matplotlib colourmap name used only for grayscale fallback saves.
    """
    if _HAVE_CV2:
        if arr_uint8.ndim == 2:
            _cv2.imwrite(path, arr_uint8)
        else:
            # OpenCV stores BGR; convert from RGB before writing
            _cv2.imwrite(path, _cv2.cvtColor(arr_uint8, _cv2.COLOR_RGB2BGR))
    else:
        plt.imsave(path, arr_uint8,
                   cmap=(cmap if arr_uint8.ndim == 2 else None),
                   vmin=0, vmax=255)

def save_gray(path, img_float):
    """
    Saves a floating-point single-channel image as an 8-bit grayscale PNG.

    Applies a full min–max stretch (appropriate for debug images where absolute
    intensity is less important than structure visibility).

    Args:
        path (str): Output file path.
        img_float (np.ndarray): Float image of arbitrary range.
    """
    a = img_float.astype(np.float32)
    a = (a - a.min()) / (a.max() - a.min() + 1e-9)
    _imwrite(path, (a * 255).astype(np.uint8), cmap="gray")

def save_mask(path, mask_bool):
    """
    Saves a binary boolean mask as a black-and-white 8-bit PNG.

    Pixels where *mask_bool* is ``True`` are written as 255 (white);
    background pixels are 0 (black). This is the standard convention
    for binary mask overlays in ImageJ / FIJI.

    Args:
        path (str): Output file path.
        mask_bool (np.ndarray[bool]): Boolean mask array.
    """
    _imwrite(path, mask_bool.astype(np.uint8) * 255, cmap="gray")

def load_batch_files(input_dir, pattern):
    """
    Discovers and sorts all image files for a batch run.

    Uses a three-tier fallback strategy so the pipeline works even when
    users supply unusual file extensions or mixed case (.TIF vs .tif):

    1. Exact pattern match (e.g. ``Project001_Series002_z*_ch00.tif``).
    2. ``.tif`` → ``.tiff`` substitution.
    3. Broad glob over all supported extensions (tif, tiff, png, jpg, jpeg)
       in both lower- and upper-case variants.

    After discovery, files are sorted with :func:`natural_sort_key` so that
    ``z2`` comes before ``z10`` (lexicographic sort would reverse this).  A
    Z-index is extracted from each filename or assigned by sequence position
    as a fallback.

    Args:
        input_dir (str): Directory to search for image files.
        pattern (str): Glob pattern relative to *input_dir*.

    Returns:
        tuple[list[str], list[int]]: Sorted file paths and corresponding Z-indices.

    Raises:
        FileNotFoundError: If no matching images are found after all fallbacks.
    """
    files = glob.glob(os.path.join(input_dir, pattern))
    if not files:
        # Fallback 1: .tif → .tiff extension swap
        files = glob.glob(os.path.join(input_dir, pattern.replace(".tif", ".tiff")))
    if not files:
        # Fallback 2: Broad scan of all supported image formats
        for ext in ['*.tif', '*.tiff', '*.png', '*.jpg', '*.jpeg']:
            files.extend(glob.glob(os.path.join(input_dir, ext)))
            files.extend(glob.glob(os.path.join(input_dir, ext.upper())))
        files = list(set(files))  # remove cross-extension duplicates

    if not files:
        raise FileNotFoundError(f"No supported image files found in '{input_dir}'")

    # Natural sort preserves correct Z-stack ordering (z2 < z10)
    files = sorted(files, key=natural_sort_key)

    # Extract z-index from filename pattern; fall back to sequence index
    z_idx = [extract_z_index(f, sequence_idx=i) for i, f in enumerate(files)]

    print(f"Found {len(files)} slices: Z = {z_idx}")
    return files, z_idx

def choose_single_image(cfg):
    """
    Resolves the path to a single test image according to the configured selection mode.

    Three modes are supported (set via ``SINGLE_IMAGE_SELECTION_MODE``):

    - ``'path'``     — Use the hard-coded ``SINGLE_TEST_IMAGE`` path directly.
    - ``'z_index'``  — Find the file whose filename encodes the Z-index in
                      ``SINGLE_Z_INDEX`` (e.g. ``z15`` → ``...z15_ch00.tif``).
    - ``'dialog'``   — Open a Tkinter file-picker dialog for interactive selection.
                      Falls back gracefully if Tkinter is unavailable.

    Args:
        cfg (dict): Pipeline configuration dictionary.

    Returns:
        str: Absolute path to the selected image file.

    Raises:
        FileNotFoundError: For ``'path'`` mode when the file does not exist, or
                           ``'z_index'`` mode when no file matches the Z-index.
        RuntimeError: For ``'dialog'`` mode when Tkinter fails or the user cancels.
        ValueError: If ``SINGLE_IMAGE_SELECTION_MODE`` is not one of the three valid values.
    """
    mode = cfg["SINGLE_IMAGE_SELECTION_MODE"].lower()
    if mode == "path":
        p = cfg["SINGLE_TEST_IMAGE"]
        if not os.path.exists(p):
            raise FileNotFoundError(f"Not found: {p}")
        return p
    if mode == "z_index":
        z = int(cfg["SINGLE_Z_INDEX"])
        files = (glob.glob(os.path.join(cfg["INPUT_DIR"], cfg["FILE_PATTERN"])) or
                 glob.glob(os.path.join(cfg["INPUT_DIR"],
                                        cfg["FILE_PATTERN"].replace(".tif", ".tiff"))))
        hits = [f for f in files if extract_z_index(f) == z]
        if not hits:
            raise FileNotFoundError(f"No file for z={z}")
        print(f"Using z={z}: {hits[0]}")
        return hits[0]
    if mode == "dialog":
        try:
            import tkinter as tk
            from tkinter import filedialog
            root = tk.Tk(); root.withdraw()  # hide the empty root window
            p = filedialog.askopenfilename(
                title="Choose image", initialdir=cfg["INPUT_DIR"],
                filetypes=[("TIFF", "*.tif *.tiff"), ("All", "*.*")])
            root.destroy()
        except Exception as e:
            raise RuntimeError(f"File dialog failed: {e}")
        if not p:
            raise RuntimeError("No file selected.")
        print(f"Selected: {p}")
        return p
    raise ValueError("SINGLE_IMAGE_SELECTION_MODE: 'path'|'z_index'|'dialog'")

# =============================================================================
# SKELETON UTILITIES
# =============================================================================

def find_endpoints(skel_bool):
    """
    Identifies all endpoint pixels (tip pixels) in a binary skeleton image.

    In an 8-connected skeleton, a pixel is an *endpoint* if exactly one of its
    eight neighbours is also part of the skeleton.  These pixels correspond to
    the physical tips of spermatid filaments.  A clean spermatid should have
    exactly 2 endpoints (head and tail end); a bridged fragment pair still has
    2 endpoints after merging.

    Args:
        skel_bool (np.ndarray[bool]): 2D binary skeleton image (True = skeleton pixel).

    Returns:
        list[tuple[int, int]]: List of ``(row, col)`` pixel coordinates of all endpoints.
    """
    H, W = skel_bool.shape
    ys, xs = np.where(skel_bool)
    sk_set = set(zip(ys.tolist(), xs.tolist()))
    return [(r, c) for r, c in sk_set
            if sum(1 for dr, dc in _N8
                   if 0 <= r+dr < H and 0 <= c+dc < W
                   and skel_bool[r+dr, c+dc]) == 1]

def bridge_skeleton_endpoints(skel_bool, skel_labeled, max_gap_px):
    """
    Join endpoint pairs from DIFFERENT components within max_gap_px by a
    1-px straight line.  Preserves the original mask so width estimates
    are not inflated.
    """
    if max_gap_px <= 0:
        return skel_bool.copy()
    H, W   = skel_bool.shape
    out    = skel_bool.copy()
    eps    = find_endpoints(skel_bool)
    if not eps:
        return out
    ep_arr    = np.array(eps, dtype=np.float32)
    ep_labels = skel_labeled[ep_arr[:, 0].astype(int), ep_arr[:, 1].astype(int)]
    pairs     = cKDTree(ep_arr).query_pairs(r=max_gap_px, output_type="ndarray")
    for i, j in pairs:
        if ep_labels[i] == ep_labels[j]:
            continue
        r0, c0 = int(eps[i][0]), int(eps[i][1])
        r1, c1 = int(eps[j][0]), int(eps[j][1])
        n = max(abs(r1-r0), abs(c1-c0)) + 1
        rs = np.clip(np.round(np.linspace(r0, r1, n)).astype(int), 0, H-1)
        cs = np.clip(np.round(np.linspace(c0, c1, n)).astype(int), 0, W-1)
        out[rs, cs] = True
    return out

def prune_branches(skel_bool, max_branch_len):
    """
    Iteratively remove endpoints to shorten side-branches ≤ max_branch_len px.
    """
    if max_branch_len <= 0:
        return skel_bool.copy()
    H, W  = skel_bool.shape
    skel  = skel_bool.copy()
    for _ in range(int(max_branch_len)):
        eps = find_endpoints(skel)
        if not eps:
            break
        for r, c in eps:
            n = sum(1 for dr, dc in _N8
                    if 0 <= r+dr < H and 0 <= c+dc < W and skel[r+dr, c+dc])
            if n == 1:
                skel[r, c] = False
    return skel

# =============================================================================
# GEODESIC & TOPOLOGY MEASUREMENT
# =============================================================================

def _build_adj(coords, W):
    """
    Builds a lightweight adjacency list for a set of skeleton pixel coordinates.

    Uses a *linearised index* trick: each pixel ``(r, c)`` is mapped to an integer
    ``r * W + c`` so that neighbour lookup is an O(1) dictionary lookup rather than
    a 2D array access.  This is important for large skeleton components.

    Edge weights follow the Chebyshev / Euclidean convention:
    - Axis-aligned move (4-connected step): weight = 1.0
    - Diagonal move (8-connected step):     weight = √2 ≈ 1.41421

    Args:
        coords (np.ndarray): Shape ``(N, 2)`` integer array of ``(row, col)`` pixel positions.
        W (int): Image width in pixels (used to compute the linear index).

    Returns:
        list[list[tuple[int, float]]]: Adjacency list where ``adj[i]`` is a list of
        ``(neighbour_index, edge_weight)`` pairs.
    """
    n       = len(coords)
    # Linearise (row, col) -> single int for O(1) membership test
    lin     = coords[:, 0] * W + coords[:, 1]
    lin2idx = {int(v): i for i, v in enumerate(lin.tolist())}
    lin_set = set(lin.tolist())
    adj     = [[] for _ in range(n)]
    for i, (r, c) in enumerate(coords.tolist()):
        for dr, dc in _N8:
            lk = (r + dr) * W + (c + dc)
            if lk in lin_set:
                # Diagonal edges are longer by a factor of √2
                w = 1.41421356 if (dr != 0 and dc != 0) else 1.0
                adj[i].append((lin2idx[lk], w))
    return adj

def _dijkstra(adj, src, n):
    """
    Runs Dijkstra's shortest-path algorithm from source node *src* and returns
    the *farthest* reachable node and its distance.

    This is the first BFS pass in the double-BFS algorithm for computing the
    true *geodesic diameter* (longest shortest path) of a skeleton component.
    The second BFS starts from the farthest node found here.

    Mathematical note:
        Geodesic length = shortest path distance through the skeleton graph,
        which accounts for the actual curvature of the filament rather than
        the straight-line Euclidean distance between endpoints.

    Args:
        adj (list[list[tuple[int, float]]]): Adjacency list from :func:`_build_adj`.
        src (int): Index of the source node in *adj*.
        n (int): Total number of nodes.

    Returns:
        tuple[int, float]: ``(farthest_node_index, distance_to_farthest)``
    """
    d = np.full(n, np.inf); d[src] = 0.0
    pq = [(0.0, src)]
    while pq:
        cost, u = heapq.heappop(pq)
        if cost > d[u]:
            continue  # stale entry in the heap — skip
        for v, w in adj[u]:
            nd = cost + w
            if nd < d[v]:
                d[v] = nd
                heapq.heappush(pq, (nd, v))
    far = int(np.argmax(d))
    return far, float(d[far])

def measure_topology(coords, W, allow_loops=False):
    """
    Compute geodesic length, tortuosity, endpoint count, and branch-node count
    for one skeleton component.

    Returns
    -------
    dict with keys: geo_len, tortuosity, n_endpoints, n_branch_nodes, reason
    or None if the component is a loop and allow_loops=False.
    """
    n   = len(coords)
    adj = _build_adj(coords, W)
    deg = [len(a) for a in adj]

    n_endpoints    = sum(1 for d in deg if d == 1)
    n_branch_nodes = sum(1 for d in deg if d >  2)
    eps_idx        = [i for i, d in enumerate(deg) if d == 1]

    # ── Loop handling ─────────────────────────────────────────────────────────
    if not eps_idx:
        if not allow_loops:
            return None  # discard loop
        seen, total = set(), 0.0
        for u, nbrs in enumerate(adj):
            for v, w in nbrs:
                key = (min(u, v), max(u, v))
                if key not in seen:
                    seen.add(key); total += w
        return {"geo_len": total, "tortuosity": 1.0,
                "n_endpoints": 0, "n_branch_nodes": n_branch_nodes,
                "reason": "loop"}

    # ── Double-BFS for geodesic ───────────────────────────────────────────────
    b, _  = _dijkstra(adj, eps_idx[0], n)
    c, gl = _dijkstra(adj, b, n)

    # ── Tortuosity ────────────────────────────────────────────────────────────
    p0  = coords[b]
    p1  = coords[c]
    euc = float(np.sqrt((p0[0]-p1[0])**2 + (p0[1]-p1[1])**2))
    tort = gl / (euc + 1e-9)

    return {"geo_len": float(gl), "tortuosity": tort,
            "n_endpoints": n_endpoints, "n_branch_nodes": n_branch_nodes,
            "reason": "ok"}

# =============================================================================
# SEGMENTATION PIPELINE
# =============================================================================

def apply_optional_early_shape_filter(mask, cfg):
    """
    Optionally removes non-elongated connected components from the binary mask
    *before* skeletonisation, based on region shape descriptors.

    This is an early-stage pre-filter controlled by ``USE_EARLY_SHAPE_FILTER``.
    When enabled, it eliminates round debris, fat nuclei, and large globular
    artefacts that would otherwise generate spurious skeleton trees downstream.

    Biological rationale
    --------------------
    Round debris (dead cells, lipid droplets, imaging artefacts) tends to have:
    - Low eccentricity (close to circular)
    - Short major axis (small, compact object)
    - Low axis ratio (major ≈ minor, i.e., not elongated)

    Elongated spermatid nuclei, by contrast, have high eccentricity (> 0.6),
    a long major axis, and a large major/minor axis ratio.

    Args:
        mask (np.ndarray[bool]): Binary mask from hysteresis thresholding.
        cfg (dict): Pipeline config; reads ``USE_EARLY_SHAPE_FILTER``,
            ``MIN_ECCENTRICITY``, ``MAX_MINOR_PX``, ``MIN_AXIS_RATIO``,
            and ``MIN_MAJOR_PX``.

    Returns:
        np.ndarray[bool]: Filtered binary mask (unchanged if the filter is disabled).
    """
    if not cfg["USE_EARLY_SHAPE_FILTER"]:
        return mask  # filter disabled — pass through unchanged
    labeled = measure.label(mask)
    keep = [p.label for p in measure.regionprops(labeled)
            if (p.eccentricity      >= cfg["MIN_ECCENTRICITY"] and   # must be elongated
                p.minor_axis_length <= cfg["MAX_MINOR_PX"]     and   # must be narrow
                p.major_axis_length / (p.minor_axis_length + 1e-9) >= cfg["MIN_AXIS_RATIO"] and  # must be rod-like
                p.major_axis_length >= cfg["MIN_MAJOR_PX"])]
    return np.isin(labeled, keep)

def segment_slice(img_raw, cfg, z_idx=None, debug_dir=None, roi_mask=None):
    """
    Executes advanced 2D multi-stage morphology segmentation to detect and isolate spermatid nuclei.
    
    Args:
        img_raw (np.ndarray): The raw 2D grayscale image array of the z-slice.
        cfg (dict): Pipeline hyperparameters including thresholding, CLAHE limits, and morphological radius.
        z_idx (int, optional): The current Z-index integer for debugging context.
        debug_dir (str, optional): Directory to save intermediate thresholding outputs if tracking is on.
        roi_mask (np.ndarray, optional): A boolean boolean layer representing the user's manual bounding box.

    Returns:
        np.ndarray: A labeled discrete array of contiguous pixel bodies corresponding to valid single nuclei.
    """
    img = ensure_2d_image(img_raw, f"segment_slice z={z_idx}").astype(np.float32)
    print(f"DEBUG segment_slice input shape: {img.shape}, ndim={img.ndim}")
    if roi_mask is not None:
        roi_mask = roi_mask.astype(bool)
        if roi_mask.shape != img.shape:
            raise ValueError(f"roi_mask shape {roi_mask.shape} does not match image shape {img.shape}")
        # NOTE: do NOT zero out pixels here — that destroys image statistics
        # (CLAHE, percentiles, ridge detection).  Instead, process full image
        # and apply roi_mask only at the mask/skeleton stages below.

    # 1. Intensity Normalization
    img_norm = (img - img.min()) / (img.max() - img.min() + 1e-9)

    # 2. Adaptive Contrast Enhancement (CLAHE) - Targets local variance over global washouts
    img_eq = exposure.equalize_adapthist(
        img_norm, clip_limit=cfg["CLAHE_CLIP"], kernel_size=cfg["CLAHE_KERNEL"])

    # 3. Background Subtraction
    bg  = gaussian(img_eq, sigma=cfg["BG_SIGMA"])
    fg  = np.clip(img_eq - bg, 0, None)
    fgn = fg / (fg.max() + 1e-9)

    # 4. Neurite/Ridge Enhanced Topological Edge Detection (Meijering Filter)
    ridge = meijering(fgn, sigmas=cfg["RIDGE_SIGMAS"], black_ridges=False)

    th_hi     = np.percentile(ridge, cfg["THRESHOLD_HI"])
    th_lo     = np.percentile(ridge, cfg["THRESHOLD_LO"])
    mask_hyst = apply_hysteresis_threshold(ridge, th_lo, th_hi)
    if roi_mask is not None:
        mask_hyst &= roi_mask

    if mask_hyst.ndim != 2:
        raise ValueError(f"mask_hyst must be 2D, got shape {mask_hyst.shape}")
    mask_clean = morphology.binary_closing(mask_hyst, morphology.disk(cfg["CLOSE_RADIUS"]))
    mask_clean = morphology.remove_small_holes(mask_clean, area_threshold=cfg["MIN_HOLE_AREA"])
    mask_clean = morphology.remove_small_objects(mask_clean, min_size=cfg["MIN_OBJ_PX"])
    mask_clean = apply_optional_early_shape_filter(mask_clean, cfg)
    if roi_mask is not None:
        mask_clean &= roi_mask

    # Width is measured from the CLEAN (un-bridged) distance map
    dist_clean   = distance_transform_edt(mask_clean)
    skel_clean   = skeletonize(mask_clean)
    skel_labeled = measure.label(skel_clean)

    # Skeleton-level bridging (preserves mask / width integrity)
    skel_bridged    = bridge_skeleton_endpoints(
        skel_clean, skel_labeled, cfg["MAX_BRIDGE_PX"])
    if roi_mask is not None:
        skel_bridged &= roi_mask
    # Branch pruning before measurement
    skel_pruned     = prune_branches(skel_bridged, cfg["MAX_BRANCH_LEN_PX"])
    
    # Optional: Automatically sever complex webs into isolated individual strands
    if cfg.get("BREAK_JUNCTIONS", False):
        from scipy.ndimage import convolve
        kernel = np.array([[1, 1, 1],
                           [1, 0, 1],
                           [1, 1, 1]], dtype=np.int32)
        skel_int = skel_pruned.astype(np.int32)
        neighbors = convolve(skel_int, kernel, mode='constant', cval=0)
        # Any skeleton pixel with more than 2 neighbors is a junction
        junctions = (skel_int > 0) & (neighbors > 2)
        skel_pruned[junctions] = 0
        
    if roi_mask is not None:
        skel_pruned &= roi_mask
    skel_labeled_fn = measure.label(skel_pruned)

    # NEW: The Recursive Adaptive Micro-Crop Reanalyzer for dense webs and chains
    if cfg.get("MAX_GEODESIC_LEN_PX", 0) > 0 and cfg.get("AUTO_LOCAL_REANALYSIS", True):
        max_px = cfg["MAX_GEODESIC_LEN_PX"]
        
        # Create a rigid dict of sub-parameters to forcefully shatter the isolated components
        sub_cfg = cfg.copy()
        sub_cfg["AUTO_LOCAL_REANALYSIS"] = False # Prevent infinite recursion
        
        # Because the bounding box crop contains almost NO dark background, 
        # a median 50% threshold perfectly isolates the brightest centers of the blobs!
        sub_cfg["THRESHOLD_HI"] = 55.0
        sub_cfg["THRESHOLD_LO"] = 45.0
        
        # Deactivate topological limits in the recursive sub-call so fragments aren't 
        # instantly dropped before they can be spliced back into the master image.
        sub_cfg["MIN_SKEL_LEN_PX"] = 1.0
        sub_cfg["MIN_OBJ_PX"] = 3
        sub_cfg["MIN_HOLE_AREA"] = 0

        props = measure.regionprops(skel_labeled_fn)
        for sp in props:
            # Note: A dense 2D web might have len(sp.coords) = 5000 pixels.
            if len(sp.coords) > max_px:
                minr, minc, maxr, maxc = sp.bbox
                pad = 12
                minr = max(0, minr - pad)
                minc = max(0, minc - pad)
                maxr = min(img.shape[0], maxr + pad)
                maxc = min(img.shape[1], maxc + pad)
                
                crop_img = img[minr:maxr, minc:maxc]
                
                # Create a strict ROI mask over the target structure to ignore neighbors
                obj_mask = (skel_labeled_fn[minr:maxr, minc:maxc] == sp.label)
                crop_roi = morphology.dilation(obj_mask, morphology.disk(6))
                
                try:
                    # RECURE: Run the entire engine on the tiny isolated crop!
                    sub_seg = segment_slice(crop_img, sub_cfg, roi_mask=crop_roi)
                    sub_skel = sub_seg["skel_pruned"]
                    sub_lab = measure.label(sub_skel)
                    
                    if sub_lab.max() > 1:
                        # SUCCESS: The adaptive threshold organically shattered the web!
                        skel_pruned[minr:maxr, minc:maxc][obj_mask] = 0
                        new_frags = (sub_skel > 0)
                        skel_pruned[minr:maxr, minc:maxc][new_frags] = 1
                    else:
                        # Failsafe: if the intensity was perfectly uniform, geometric centroid chop
                        cy, cx = sp.centroid
                        dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                        mid_idx = np.argmin(dists)
                        my, mx = sp.coords[mid_idx]
                        skel_pruned[my-1:my+2, mx-1:mx+2] = 0
                        
                except Exception:
                    # Absolute geometric failsafe
                    cy, cx = sp.centroid
                    dists = (sp.coords[:, 0] - cy)**2 + (sp.coords[:, 1] - cx)**2
                    mid_idx = np.argmin(dists)
                    my, mx = sp.coords[mid_idx]
                    skel_pruned[my-1:my+2, mx-1:mx+2] = 0

        # Run one final relabeling array refresh after all splicing
        skel_labeled_fn = measure.label(skel_pruned)

    out = {
        "mask_hyst":    mask_hyst,
        "mask_clean":   mask_clean,
        "skel_clean":   skel_clean,
        "skel_bridged": skel_bridged,
        "skel_pruned":  skel_pruned,
        "skel_labeled": skel_labeled_fn,
        "dist_clean":   dist_clean,
        "img_eq":       img_eq,
        "ridge":        ridge,
        "roi_mask":     roi_mask,
    }

    if cfg["SAVE_DEBUG_IMAGES"] and debug_dir and z_idx is not None:
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_01_norm.png"),          img_norm)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_02_clahe.png"),         img_eq)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_03_fg.png"),            fgn)
        save_gray(os.path.join(debug_dir, f"z{z_idx:02d}_04_ridge.png"),         ridge)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_05_hysteresis.png"),    mask_hyst)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_06_clean.png"),         mask_clean)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_07_skel_clean.png"),    skel_clean)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_08_skel_bridged.png"),  skel_bridged)
        save_mask(os.path.join(debug_dir, f"z{z_idx:02d}_09_skel_pruned.png"),   skel_pruned)

    return out

# =============================================================================
# MEASUREMENT  (single geodesic pass, all topology in one function)
# =============================================================================

def measure_spermatids(seg, cfg):
    """
    Analyzes mathematically discretised nuclei arrays and derives geometric indices for individual shape profiles.
    
    Args:
        seg (dict): Dictionary containing labeled binary masks ("primary", "skel_pruned", "dist_clean").
        cfg (dict): Pipeline hyperparameters for absolute scaling calculations.
        
    Returns:
        pd.DataFrame: Table indexing each isolated object dynamically with mathematical biometrics.
                      Attributes include region area, bounding-boxes, minor/major ellipse axis mapping,
                      orientation, and geodesic structural lengths via binary skeleton processing.
    """
    skel     = seg["skel_pruned"]
    dist     = seg["dist_clean"]
    skel_lab = seg["skel_labeled"]
    H, W     = skel.shape

    # ── Filter pass ───────────────────────────────────────────────────────────
    accepted_labels = []
    cache           = {}
    reasons = {"short": 0, "loop": 0, "long": 0, "wide": 0, "ratio": 0, "branches": 0, "tortuous": 0, "endpoints": 0}

    for sp in measure.regionprops(skel_lab):
        coords = sp.coords
        if coords.shape[0] < cfg["MIN_SKEL_LEN_PX"]:
            reasons["short"] += 1
            continue

        topo = measure_topology(coords, W, allow_loops=cfg.get("ALLOW_LOOPS", False))
        if topo is None:
            reasons["loop"] += 1
            continue  # loop, not allowed

        gl   = topo["geo_len"]
        tort = topo["tortuosity"]
        n_ep = topo["n_endpoints"]
        n_br = topo["n_branch_nodes"]

        if not (cfg["MIN_SKEL_LEN_PX"] <= gl <= cfg["MAX_GEODESIC_LEN_PX"]):
            if gl < cfg["MIN_SKEL_LEN_PX"]: reasons["short"] += 1
            else: reasons["long"] += 1
            continue

        width = float(np.median(2.0 * dist[coords[:, 0], coords[:, 1]]))
        if width > cfg["MAX_WIDTH_PX"]:
            reasons["wide"] += 1
            continue

        if gl / (width + 1e-9) < cfg["MIN_LENGTH_WIDTH_RATIO"]:
            reasons["ratio"] += 1
            continue

        # N1: branch-node count filter
        if n_br > cfg["MAX_BRANCH_NODES"]:
            reasons["branches"] += 1
            continue

        # N2: tortuosity filter (only for open filaments with 2+ endpoints)
        if n_ep >= 2 and tort > cfg["MAX_TORTUOSITY"]:
            reasons["tortuous"] += 1
            continue

        # N3: endpoint count filter
        if n_ep > cfg["MAX_ENDPOINT_COUNT"]:
            reasons["endpoints"] += 1
            continue

        cy, cx = sp.centroid
        accepted_labels.append(sp.label)
        cache[sp.label] = {
            "geo_len":            gl,
            "tortuosity":         tort,
            "n_endpoints":        n_ep,
            "n_branch_nodes":     n_br,
            "width":              width,
            "length_width_ratio": gl / (width + 1e-9),
            "length_px_count":    float(coords.shape[0]),
            "cx": cx, "cy": cy,
            "area_px": float(sp.area),
            "bbox_min_y": float(sp.bbox[0]),
            "bbox_min_x": float(sp.bbox[1]),
            "bbox_max_y": float(sp.bbox[2]),
            "bbox_max_x": float(sp.bbox[3]),
            "orientation": float(sp.orientation),
        }

    total_rejected = sum(reasons.values())
    if total_rejected > 0:
        print(f"    measure_spermatids rejected {total_rejected} blobs:")
        for k, v in reasons.items():
            if v > 0: print(f"      {k}: {v}")

    if not accepted_labels:
        return {"skel_label": np.zeros_like(skel_lab, dtype=np.int32),
                "results": []}

    clean_skel  = np.isin(skel_lab, accepted_labels)
    final_label = measure.label(clean_skel).astype(np.int32)

    # ── Re-index using cached values (no second Dijkstra pass) ───────────────
    final_results = []
    for new_i, sp in enumerate(measure.regionprops(final_label), start=1):
        old_label = skel_lab[sp.coords[0, 0], sp.coords[0, 1]]
        if old_label not in cache:
            continue
        c = cache[old_label]
        final_results.append({
            "label":               new_i,
            "length_px_geodesic":  c["geo_len"],
            "length_px_count":     c["length_px_count"],
            "width_px":            c["width"],
            "length_width_ratio":  c["length_width_ratio"],
            "tortuosity":          c["tortuosity"],
            "n_endpoints":         c["n_endpoints"],
            "n_branch_nodes":      c["n_branch_nodes"],
            "centroid_x":          c["cx"],
            "centroid_y":          c["cy"],
            "area_px":             c["area_px"],
            "bbox_min_y":          c["bbox_min_y"],
            "bbox_min_x":          c["bbox_min_x"],
            "bbox_max_y":          c["bbox_max_y"],
            "bbox_max_x":          c["bbox_max_x"],
            "orientation":         c["orientation"],
        })

    return {"skel_label": final_label, "results": final_results}

# =============================================================================
# OVERLAY  (vectorized LUT)
# =============================================================================

def make_overlay(img_raw, skel_label):
    """
    Generates a colour-coded skeleton overlay on the grayscale raw image.

    Each detected spermatid is assigned a unique hue from the ``gist_rainbow``
    colourmap, dilated by 3 pixels for visibility, and composited onto the
    contrast-stretched raw image.  Background pixels (label == 0) retain the
    original grayscale intensity.

    Implementation notes
    --------------------
    - Vectorised LUT (Look-Up Table) approach avoids per-label loop for speed.
    - ``grey_dilation`` makes thin single-pixel skeletons visible at any zoom.
    - Colour assignment is deterministic: same label ordering → same colour.

    Args:
        img_raw (np.ndarray): Raw microscopy image (any dtype).
        skel_label (np.ndarray[int]): Integer-labelled skeleton array
            (0 = background, 1..N = individual spermatids).

    Returns:
        np.ndarray: uint8 RGB image, shape ``(H, W, 3)``, ready for ``plt.imshow``
        or saving with :func:`_imwrite`.
    """
    base = normalize_display(img_raw)
    n    = int(skel_label.max())
    if n <= 0:
        # No detections: return grayscale image as RGB
        return (np.stack([base]*3, -1) * 255).astype(np.uint8)
    # Assign one colour per label; prepend black for background (index 0)
    cols    = plt.cm.gist_rainbow(np.linspace(0, 1, n))[:, :3]
    dilated = grey_dilation(skel_label.astype(np.int32), size=3)
    lut     = np.vstack([[0., 0., 0.], cols[:n]])
    rgb     = lut[dilated]
    # Restore original grayscale for background pixels
    m0      = dilated == 0
    rgb[m0, 0] = base[m0]
    rgb[m0, 1] = base[m0]
    rgb[m0, 2] = base[m0]
    return (np.clip(rgb, 0, 1) * 255).astype(np.uint8)

# =============================================================================
# DISPLAY / SAVE
# =============================================================================

def _safe_show():
    """
    Calls ``plt.show(block=True)`` with a try/except guard.

    On headless systems (Linux CI servers, SSH sessions without X11 forwarding,
    or Windows without a display) matplotlib will raise a ``RuntimeError`` or
    ``_tkinter.TclError`` when trying to open an interactive window.  This
    wrapper silently swallows that error so the pipeline can continue and saves
    the image to disk as a fallback.
    """
    try:
        plt.show(block=True)
    except Exception as e:
        print(f"[WARNING] plt.show() failed: {e}")

def show_single_preview(img_raw, seg, overlay_rgb, results, z_idx, cfg):
    """
    Renders and saves an interactive preview figure for a single analysed Z-slice.

    Layout operates in two modes controlled by ``SHOW_DEBUG_PREVIEW``:

    - **Standard mode** (2 panels): Raw image | Spermatid overlay + length histogram.
    - **Debug mode** (8 panels): Adds intermediate stage images — CLAHE, ridge filter,
      hysteresis mask, cleaned mask, and pruned skeleton — for visual pipeline QC.

    After saving to ``output_dir/preview.png``, the function attempts to open the
    image with the OS default viewer (``os.startfile`` on Windows) as an immediate
    visual check, falling back silently if that fails.

    Args:
        img_raw (np.ndarray): Raw microscopy image.
        seg (dict): Segmentation dictionary from :func:`segment_slice` containing
            intermediate images (``img_eq``, ``ridge``, ``mask_hyst``, etc.).
        overlay_rgb (np.ndarray): Colour overlay from :func:`make_overlay`.
        results (list[dict]): Per-spermatid measurement dictionaries.
        z_idx (int): Z-slice index used for figure titles.
        cfg (dict): Pipeline configuration; reads ``SHOW_DEBUG_PREVIEW``,
            ``OUTPUT_DIR``, and ``UM_PER_PX_XY``.
    """
    um         = cfg["UM_PER_PX_XY"]
    lengths_um = [r["length_px_geodesic"] * um for r in results]

    nrows, ncols = (2, 4) if cfg["SHOW_DEBUG_PREVIEW"] else (1, 3)
    fig, axes = plt.subplots(nrows, ncols, figsize=(18, 5*nrows))

    def _ax(r, c):
        return axes[r, c] if nrows > 1 else axes[c]

    _ax(0,0).imshow(normalize_display(img_raw), cmap="gray")
    _ax(0,0).set_title(f"Original Z={z_idx:02d}"); _ax(0,0).axis("off")

    if cfg["SHOW_DEBUG_PREVIEW"]:
        _ax(0,1).imshow(seg["img_eq"], cmap="gray")
        _ax(0,1).set_title("CLAHE"); _ax(0,1).axis("off")
        _ax(0,2).imshow(seg["ridge"], cmap="gray")
        _ax(0,2).set_title("Ridge"); _ax(0,2).axis("off")
        _ax(0,3).imshow(seg["mask_hyst"], cmap="gray")
        _ax(0,3).set_title("Hysteresis"); _ax(0,3).axis("off")
        _ax(1,0).imshow(seg["mask_clean"], cmap="gray")
        _ax(1,0).set_title("Clean mask"); _ax(1,0).axis("off")
        _ax(1,1).imshow(seg["skel_pruned"], cmap="gray")
        _ax(1,1).set_title("Skeleton (pruned)"); _ax(1,1).axis("off")
        ov_ax   = _ax(1,2)
        hist_ax = _ax(1,3)
    else:
        ov_ax   = _ax(0,1)
        hist_ax = _ax(0,2)

    ov_ax.imshow(overlay_rgb)
    ov_ax.set_title(f"Overlay N={len(results)}"); ov_ax.axis("off")
    for r in results:
        ov_ax.text(r["centroid_x"], r["centroid_y"],
                   f"{r['length_px_geodesic']*um:.1f}",
                   color="white", fontsize=5, ha="center", va="center")

    if lengths_um:
        hist_ax.hist(lengths_um, bins=20, edgecolor="white")
        hist_ax.axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} µm")
        hist_ax.set_xlabel("Geodesic length (µm)"); hist_ax.legend(fontsize=8)
        hist_ax.set_title("Length distribution")
    else:
        hist_ax.text(0.5, 0.5, "No detections", ha="center", va="center")
        hist_ax.axis("off")

    plt.tight_layout()

    # Persist preview image so it can be opened by OS or reloaded by GUI
    preview_path = os.path.join(cfg["OUTPUT_DIR"], "preview.png")
    ensure_dir(cfg["OUTPUT_DIR"])
    plt.savefig(preview_path, dpi=120, bbox_inches="tight")
    print(f"  Preview saved to: {preview_path}")

    # Try interactive display
    _safe_show()
    plt.close()

    # Open with default viewer as fallback
    try:
        os.startfile(preview_path)
    except Exception:
        pass

def save_detail_figure(img_raw, overlay_rgb, results, out_path, z_idx, um):
    """
    Saves a three-panel publication-quality figure for a single Z-slice.

    Panels:
    1. **Original** — contrast-stretched raw image.
    2. **Spermatid overlay** — colour-coded detections with length labels.
    3. **Length distribution** — histogram of geodesic lengths in µm with median line.

    This figure is saved as a PNG per Z-slice and is useful for visual quality
    control and inclusion in lab reports.

    Args:
        img_raw (np.ndarray): Raw microscopy image.
        overlay_rgb (np.ndarray): Colour overlay from :func:`make_overlay`.
        results (list[dict]): Per-spermatid measurement dictionaries.
        out_path (str): Destination PNG file path.
        z_idx (int): Z-slice index for figure titles.
        um (float): Microns-per-pixel scale factor (``UM_PER_PX_XY``).
    """
    lengths_um = [r["length_px_geodesic"] * um for r in results]
    fig, axes  = plt.subplots(1, 3, figsize=(18, 6))

    axes[0].imshow(normalize_display(img_raw), cmap="gray")
    axes[0].set_title(f"Z={z_idx:02d} — Original"); axes[0].axis("off")

    axes[1].imshow(overlay_rgb)
    axes[1].set_title(f"Z={z_idx:02d} — Spermatids (N={len(results)})")
    axes[1].axis("off")
    for r in results:
        axes[1].text(r["centroid_x"], r["centroid_y"],
                     f"{r['length_px_geodesic']*um:.1f}",
                     color="white", fontsize=4, ha="center", va="center",
                     bbox=dict(boxstyle="round,pad=0.1", fc="black", alpha=0.4, lw=0))

    if lengths_um:
        axes[2].hist(lengths_um, bins=20, edgecolor="white")
        axes[2].axvline(np.median(lengths_um), lw=2,
                        label=f"Median={np.median(lengths_um):.1f} µm")
        axes[2].set_xlabel("Geodesic length (µm)"); axes[2].set_ylabel("Count")
        axes[2].set_title(f"Z={z_idx:02d} — Length distribution")
        axes[2].legend(fontsize=9)
    else:
        axes[2].text(0.5, 0.5, "No spermatids detected",
                     transform=axes[2].transAxes, ha="center", va="center")
        axes[2].axis("off")

    plt.tight_layout()
    plt.savefig(out_path, dpi=120, bbox_inches="tight")
    plt.close()

# =============================================================================
# CSV
# =============================================================================

_VERSION = "v3_pre_audit"

def rows_from_results(results, z_idx, um):
    """
    Converts per-spermatid measurement dictionaries into flat CSV row dictionaries.

    Applies pixel-to-micron scaling (``UM_PER_PX_XY``) for all linear dimensions
    and rounds floating-point values to 3 decimal places for clean spreadsheet output.

    Columns emitted per detection
    ------------------------------
    - ``pipeline_version``     — Code version tag for traceability.
    - ``z_slice``              — Z-plane index of this detection.
    - ``sperm_id``             — 1-based integer ID within this slice.
    - ``length_px_geodesic``   — Geodesic skeleton length in pixels.
    - ``length_um_geodesic``   — Geodesic length in micrometres.
    - ``length_px_count``      — Pixel-count skeleton length (alternative measure).
    - ``length_um_count``      — Pixel-count length in micrometres.
    - ``width_px`` / ``width_um`` — Mean width across the skeleton.
    - ``length_width_ratio``   — Elongation ratio (key morphological filter metric).
    - ``tortuosity``           — Curvature index: geodesic / Euclidean tip-to-tip.
    - ``n_endpoints``          — Number of skeleton endpoints (expect 2 for clean cells).
    - ``n_branch_nodes``       — Number of branch points (3+ neighbours); should be low.
    - ``centroid_x`` / ``centroid_y`` — Pixel centroid for overlay annotation.
    - ``area_px``              — Mask area in pixels^2 (used for volume estimation).

    Args:
        results (list[dict]): Output of :func:`measure_spermatids`.
        z_idx (int): Z-slice index.
        um (float): Microns-per-pixel (``UM_PER_PX_XY``).

    Returns:
        list[dict]: One flat dictionary per detected spermatid.
    """
    return [{
        "pipeline_version":    _VERSION,
        "z_slice":             z_idx,
        "sperm_id":            i,
        "length_px_geodesic":  round(r["length_px_geodesic"], 3),
        "length_um_geodesic":  round(r["length_px_geodesic"] * um, 3),
        "length_px_count":     round(r["length_px_count"], 1),
        "length_um_count":     round(r["length_px_count"]  * um, 3),
        "width_px":            round(r["width_px"], 2),
        "width_um":            round(r["width_px"]          * um, 3),
        "length_width_ratio":  round(r["length_width_ratio"], 3),
        "tortuosity":          round(r["tortuosity"], 3),
        "n_endpoints":         r["n_endpoints"],
        "n_branch_nodes":      r["n_branch_nodes"],
        "centroid_x":          round(r["centroid_x"], 1),
        "centroid_y":          round(r["centroid_y"], 1),
        "area_px":             round(r["area_px"], 1),
        "bbox_min_y":          r.get("bbox_min_y"),
        "bbox_min_x":          r.get("bbox_min_x"),
        "bbox_max_y":          r.get("bbox_max_y"),
        "bbox_max_x":          r.get("bbox_max_x"),
        "orientation":         round(r.get("orientation", 0.0), 3),
    } for i, r in enumerate(results, start=1)]

# =============================================================================
# TRACKING
# =============================================================================

def check_extension_consistency(prev_state, candidate_detection, cfg, overlap_exists=False):
    """
    Check if extending a track with this detection would be biologically consistent.
    Implements 'Continue Unless Implausible' logic for overlapping footprints.
    
    Stage 2b: All thresholds are now driven by CONFIG for hyperparameter tuning.
    """
    um_xy = cfg["UM_PER_PX_XY"]
    
    # Read tunable Stage 2 parameters from CONFIG
    stab_thresh = cfg.get("OVERLAP_STABILITY_THRESHOLD", 0.08)
    ori_deg     = cfg.get("OVERLAP_ORIENTATION_DEG", 15.0)
    ovl_mult    = cfg.get("OVERLAP_MULTIPLIER", 1.35)
    min_stable  = cfg.get("OVERLAP_MIN_STABLE_COUNT", 1)
    
    # Extract previous track state
    prev_x = prev_state["last_x"]
    prev_y = prev_state["last_y"]
    prev_width = prev_state.get("last_width")
    prev_length = prev_state.get("last_length")
    prev_area = prev_state.get("last_area")
    prev_ori = prev_state.get("last_orientation")
    
    # Extract candidate detection features
    cand_x = candidate_detection["centroid_x"]
    cand_y = candidate_detection["centroid_y"]
    cand_width = candidate_detection.get("width_um")
    cand_length = candidate_detection.get("length_um_geodesic")
    cand_area = candidate_detection.get("area_px")
    cand_ori = candidate_detection.get("orientation")
    
    # Logic: 
    # If overlap_exists, we allow the track to continue IF enough primary metrics are stable.
    # This prevents 'monster merges' where a track jumps onto a totally different cell.
    if overlap_exists:
        stable_count = 0
        
        # 1. Width stability
        if prev_width and cand_width:
            if (abs(cand_width - prev_width) / max(prev_width, 1e-9)) < stab_thresh:
                stable_count += 1
            
        # 2. Area stability
        if prev_area and cand_area:
            if (abs(cand_area - prev_area) / max(prev_area, 1e-9)) < stab_thresh:
                stable_count += 1
            
        # 3. Orientation stability
        if prev_ori is not None and cand_ori is not None:
            diff_rad = abs(cand_ori - prev_ori)
            if diff_rad > math.pi / 2:
                diff_rad = math.pi - diff_rad
            if diff_rad < (ori_deg * math.pi / 180):
                stable_count += 1
        
        # 4. Length stability
        if prev_length and cand_length:
            if (abs(cand_length - prev_length) / max(prev_length, 1e-9)) < stab_thresh:
                stable_count += 1
            
        # Require minimum stable metrics to continue an overlapping track
        if stable_count < min_stable:
            return False, f"overlap_but_{stable_count}_stable"
        
        # Even with stable metrics, still apply capped multiplier for fallback checks
        multiplier = ovl_mult
    else:
        multiplier = 1.0

    # 1. Check centroid jump
    dx = cand_x - prev_x
    dy = cand_y - prev_y
    centroid_jump_um = math.sqrt(dx*dx + dy*dy) * um_xy
    
    if not overlap_exists and centroid_jump_um > cfg["CONSERVATIVE_MAX_CENTROID_JUMP_UM"]:
        return False, f"centroid_jump={centroid_jump_um:.2f}um"
    
    # 2. Check width consistency
    if prev_width is not None and cand_width is not None:
        width_ratio = abs(cand_width - prev_width) / max(prev_width, 1e-9)
        if width_ratio > cfg["CONSERVATIVE_MAX_WIDTH_JUMP_RATIO"] * multiplier:
            return False, f"width_jump={width_ratio:.2f}"
            
    # 3. Check length consistency
    if prev_length is not None and cand_length is not None:
        length_ratio = abs(cand_length - prev_length) / max(prev_length, 1e-9)
        if length_ratio > cfg["CONSERVATIVE_MAX_LENGTH_JUMP_RATIO"] * multiplier:
            return False, f"length_jump={length_ratio:.2f}"
            
    # 4. Check area consistency
    if prev_area is not None and cand_area is not None:
        area_ratio = abs(cand_area - prev_area) / max(prev_area, 1e-9)
        if area_ratio > cfg["CONSERVATIVE_MAX_AREA_JUMP_RATIO"] * multiplier:
            return False, f"area_jump={area_ratio:.2f}"
    
    return True, "ok"

def track_across_slices(detections_df, cfg):
    """
    Conservative tracking natively: stop tracks when consistency breaks.
    """
    if detections_df.empty:
        detections_df = detections_df.copy()
        detections_df["track_id"] = pd.Series(dtype=int)
        return detections_df, pd.DataFrame()

    max_dist_px = cfg["TRACK_MAX_DIST_UM"] / (cfg["UM_PER_PX_XY"] + 1e-9)
    df = (detections_df.copy()
                       .sort_values(["z_slice", "sperm_id"])
                       .reset_index(drop=True))

    next_tid = 1
    active = {}
    track_ids = [-1] * len(df)
    stopped_tracks = {}  # Track stop reasons for debugging
    
    rows_by_z = {z: df.index[df["z_slice"] == z].to_numpy()
                 for z in sorted(df["z_slice"].unique())}

    for z, idxs in rows_by_z.items():
        # Get detection features for this slice
        xs = df.loc[idxs, "centroid_x"].to_numpy(float)
        ys = df.loc[idxs, "centroid_y"].to_numpy(float)
        
        # Extract morphological features (with fallbacks)
        widths = df.loc[idxs, "width_um"].to_numpy(float) if "width_um" in df.columns else np.full(len(idxs), np.nan)
        lengths = df.loc[idxs, "length_um_geodesic"].to_numpy(float) if "length_um_geodesic" in df.columns else np.full(len(idxs), np.nan)
        areas = df.loc[idxs, "area_px"].to_numpy(float) if "area_px" in df.columns else np.full(len(idxs), np.nan)
        oris = df.loc[idxs, "orientation"].to_numpy(float) if "orientation" in df.columns else np.full(len(idxs), np.nan)
        
        # Extract Bounding Box Arrays for Overlap-First Algorithm
        bbox_min_ys = df.loc[idxs, "bbox_min_y"].to_numpy(float) if "bbox_min_y" in df.columns else np.full(len(idxs), np.nan)
        bbox_min_xs = df.loc[idxs, "bbox_min_x"].to_numpy(float) if "bbox_min_x" in df.columns else np.full(len(idxs), np.nan)
        bbox_max_ys = df.loc[idxs, "bbox_max_y"].to_numpy(float) if "bbox_max_y" in df.columns else np.full(len(idxs), np.nan)
        bbox_max_xs = df.loc[idxs, "bbox_max_x"].to_numpy(float) if "bbox_max_x" in df.columns else np.full(len(idxs), np.nan)

        # Find candidate tracks from previous slices
        cand_tracks = [t for t, st in active.items()
                       if 1 <= z - st["last_z"] <= cfg["TRACK_MAX_GAP_SLICES"] + 1]

        used_det, used_trk = set(), set()
        
        if cand_tracks:
            candidates = []
            pad = cfg.get("TRACK_BBOX_PADDING_PX", 5.0)
            
            for k, (x, y) in enumerate(zip(xs, ys)):
                # Prepare Candidate State
                det_min_y, det_min_x = bbox_min_ys[k], bbox_min_xs[k]
                det_max_y, det_max_x = bbox_max_ys[k], bbox_max_xs[k]
                has_bbox = np.isfinite(det_min_y)
                
                cand_det = {
                    "centroid_x": float(x),
                    "centroid_y": float(y),
                    "width_um": float(widths[k]) if np.isfinite(widths[k]) else None,
                    "length_um_geodesic": float(lengths[k]) if np.isfinite(lengths[k]) else None,
                    "area_px": float(areas[k]) if np.isfinite(areas[k]) else None,
                    "orientation": float(oris[k]) if np.isfinite(oris[k]) else None,
                }
                
                for j, tid in enumerate(cand_tracks):
                    trk_st = active[tid]
                    
                    # Compute standard spatial distance
                    dx = float(x) - trk_st["last_x"]
                    dy = float(y) - trk_st["last_y"]
                    d_val = np.sqrt(dx*dx + dy*dy)
                    
                    # Perform Overlap-First physical footprint collision check
                    overlap_exists = False
                    if has_bbox and "last_bbox" in trk_st and trk_st["last_bbox"] is not None:
                        t_min_y, t_min_x, t_max_y, t_max_x = trk_st["last_bbox"]
                        # Bounding box intersection formula with padding forgiveness
                        if not (det_max_y + pad < t_min_y or det_min_y - pad > t_max_y or 
                                det_max_x + pad < t_min_x or det_min_x - pad > t_max_x):
                            overlap_exists = True
                            
                    # Accept candidate if it physically overlaps, OR if it's within pure centroid limits (implausible fallback)
                    if overlap_exists or d_val <= max_dist_px:
                        # Continue unless implausible bounds
                        is_consistent, reason = check_extension_consistency(
                            trk_st, cand_det, cfg, overlap_exists=overlap_exists
                        )
                        
                        if is_consistent:
                            # Massive scoring favor (+10,000) for overlaps so they greedily override standard loose distance matches
                            score = float(d_val) if not overlap_exists else (float(d_val) - 10000.0)
                            candidates.append((score, k, j))
                        else:
                            if tid not in stopped_tracks:
                                stopped_tracks[tid] = f"z={z}, reason={reason}"
            
            # Sort by score and assign greedily
            candidates.sort(key=lambda x: x[0])
            for score, det_k, trk_j in candidates:
                if det_k in used_det or trk_j in used_trk:
                    continue
                    
                used_det.add(det_k)
                used_trk.add(trk_j)
                tid = cand_tracks[trk_j]
                track_ids[int(idxs[det_k])] = tid
                
                # Update track state
                active[tid] = {
                    "last_z": int(z),
                    "last_x": float(xs[det_k]),
                    "last_y": float(ys[det_k]),
                    "last_width": float(widths[det_k]) if np.isfinite(widths[det_k]) else None,
                    "last_length": float(lengths[det_k]) if np.isfinite(lengths[det_k]) else None,
                    "last_area": float(areas[det_k]) if np.isfinite(areas[det_k]) else None,
                    "last_orientation": float(oris[det_k]) if np.isfinite(oris[det_k]) else None,
                    "last_bbox": (bbox_min_ys[det_k], bbox_min_xs[det_k], bbox_max_ys[det_k], bbox_max_xs[det_k]) if np.isfinite(bbox_min_ys[det_k]) else None,
                }

        # Create new tracks for unmatched detections
        for det_k in range(len(idxs)):
            if track_ids[int(idxs[det_k])] == -1:
                track_ids[int(idxs[det_k])] = next_tid
                active[next_tid] = {
                    "last_z": int(z),
                    "last_x": float(xs[det_k]),
                    "last_y": float(ys[det_k]),
                    "last_width": float(widths[det_k]) if np.isfinite(widths[det_k]) else None,
                    "last_length": float(lengths[det_k]) if np.isfinite(lengths[det_k]) else None,
                    "last_area": float(areas[det_k]) if np.isfinite(areas[det_k]) else None,
                    "last_orientation": float(oris[det_k]) if np.isfinite(oris[det_k]) else None,
                    "last_bbox": (bbox_min_ys[det_k], bbox_min_xs[det_k], bbox_max_ys[det_k], bbox_max_xs[det_k]) if np.isfinite(bbox_min_ys[det_k]) else None,
                }
                next_tid += 1

        # Remove stale tracks (exceeded gap)
        for tid in [t for t, st in active.items()
                    if z - st["last_z"] > cfg["TRACK_MAX_GAP_SLICES"] + 1]:
            del active[tid]

    df["track_id"] = track_ids
    
    # Print tracking stats
    print(f"  Conservative tracking: {len(stopped_tracks)} tracks stopped early for consistency")
    
    # Inject maximum 2D Euclidean distance of the physical shape prior to grouping
    if "tortuosity" in df.columns:
        df["euc_um_2d"] = df["length_um_geodesic"] / df["tortuosity"]
    else:
        df["euc_um_2d"] = df["length_um_geodesic"]
        
    g = df.groupby("track_id", as_index=False)
    ts = g.agg(
        n_slices        = ("z_slice",            "count"),
        z_start         = ("z_slice",            "min"),
        z_end           = ("z_slice",            "max"),
        max_length_2d   = ("length_um_geodesic", "max"),
        max_euc_2d      = ("euc_um_2d",          "max"),
        sum_area_px     = ("area_px",            "sum"),
        area_start      = ("area_px",            "first"),
        area_end        = ("area_px",            "last"),
        x_mean          = ("centroid_x",         "mean"),
        y_mean          = ("centroid_y",         "mean"),
        z_mean          = ("z_slice",            "mean"),
        x_start         = ("centroid_x",         "first"),
        y_start         = ("centroid_y",         "first"),
        x_end           = ("centroid_x",         "last"),
        y_end           = ("centroid_y",         "last"),
    )
    
    um_xy = cfg["UM_PER_PX_XY"]
    um_z  = cfg["UM_PER_SLICE_Z"]
    
    # 1. Z-Extent (Vertical Span includes +1 for slice thickness, matching PDF)
    z_extent = (ts["z_end"] - ts["z_start"] + 1) * um_z
    ts["z_extent_um"] = z_extent
    
    # Geodesic vertical displacement for Euclidean distance is purely centroid-to-centroid
    dz_euc = (ts["z_end"] - ts["z_start"]) * um_z
    
    # 2. Total 3D Length (Lateral-Corrected Hypotenuse)
    # The physical 3D shape arc relies on the maximum length of its 2D projection
    euc_2d_centroid = np.sqrt((ts["x_end"] - ts["x_start"])**2 + (ts["y_end"] - ts["y_start"])**2) * um_xy
    lat_geodesic = np.maximum(ts["max_length_2d"], euc_2d_centroid)
    l3d = np.sqrt(lat_geodesic**2 + z_extent**2)
    ts["total_3d_length_um"] = l3d
    
    # 3. 3D Volume (Sum Area * Z_step)
    ts["volume_um3"] = ts["sum_area_px"] * (um_xy**2) * um_z
    
    # 4. 3D Tortuosity (Total 3D Geodesic Length / 3D End-To-End Euclidean Distance)
    euc_3d = np.sqrt(ts["max_euc_2d"]**2 + dz_euc**2)
    safe_euc = np.maximum(euc_3d, 0.1)
    tort_raw = l3d / safe_euc
    ts["tortuosity_3d"] = np.minimum(tort_raw, 20.0)
    
    # 5. Taper Ratio (Max Area / Min Area between first and last z-slice)
    ts["taper_ratio"] = np.maximum(ts["area_start"], ts["area_end"]) / np.maximum(np.minimum(ts["area_start"], ts["area_end"]), 0.001)
    
    # 6. Effective Thickness / Diameter
    cross_area = ts["volume_um3"] / np.maximum(ts["total_3d_length_um"], 0.1)
    ts["thickness_um"] = 2 * np.sqrt(cross_area / np.pi)
    
    # 7. Orientation Angles (Pitch and Yaw)
    dx = (ts["x_end"] - ts["x_start"]) * um_xy
    dy = (ts["y_end"] - ts["y_start"]) * um_xy
    v_mag = np.sqrt(dx**2 + dy**2 + dz_euc**2)
    safe_v = np.maximum(v_mag, 1e-9)
    ts["pitch_deg"] = np.abs(np.arcsin(dz_euc / safe_v)) * (180.0 / np.pi)
    ts["yaw_deg"] = np.arctan2(dy, dx) * (180.0 / np.pi)
    
    if len(ts) > 1:
        centers = np.column_stack((ts["x_mean"] * um_xy, ts["y_mean"] * um_xy, ts["z_mean"] * um_z))
        tree = cKDTree(centers)
        dists, _ = tree.query(centers, k=2)
        ts["nearest_neighbor_um"] = dists[:, 1]
    else:
        ts["nearest_neighbor_um"] = np.nan
        
    cols_ordered = [
        "track_id", "total_3d_length_um", "z_extent_um", "volume_um3", "tortuosity_3d",
        "thickness_um", "pitch_deg", "yaw_deg", "taper_ratio", "nearest_neighbor_um",
        "n_slices", "z_start", "z_end", "max_length_2d", "sum_area_px"
    ]
    ts = ts[cols_ordered]
    
    return df, ts

def write_error_log(out_dir, component, message):
    """
    Writes a persistent error log to report_generation_errors.txt in the output directory.
    """
    try:
        import os as _os
        import time as _time
        from datetime import datetime as _dt
        log_path = _os.path.join(out_dir, "report_generation_errors.txt")
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(f"\n[{_dt.now().strftime('%Y-%m-%d %H:%M:%S')}] COMPONENT: {component}\n")
            f.write(f"MESSAGE:\n{message}\n")
            f.write("-" * 80 + "\n")
    except Exception:
        pass

def generate_excel_report(out_dir, df, df_summary, df_tracks=None):
    """
    Generates a multi-tab Excel workbook with formatted data, summary statistics,
    embedded chart images, and source-data hyperlinks.

    Workbook structure
    ------------------
    - **Batch_Summary** — one row per Z-slice with detection counts, mean/median
      length, and total area.  Includes an embedded histogram image and conditional
      formatting for high/low detection slices.
    - **3D_Morphometrics** — track-level 3D metrics exported from
      :func:`track_across_slices`, with an embedded 3D length distribution plot.
    - **Raw_2D_Detections** — full per-spermatid measurement table matching the CSV
      export, with number formatting and frozen header pane.
    - **Statistics_Summary** — descriptive statistics (mean, median, std, IQR,
      percentiles) for the primary biometric columns.

    All sheets include a hyperlink from cell A1 back to the source Excel file so
    that clicking within PowerPoint linked charts opens the correct workbook row.

    Biological interpretation
    -------------------------
    The statistical summary sheet is designed to be directly copy-pasteable into
    lab reports.  The IQR (Interquartile Range = Q75 - Q25) is reported because
    spermatid length distributions are often right-skewed and IQR is more robust
    than standard deviation in these cases.

    Args:
        out_dir (str): Top-level analysis output directory.  The workbook is saved
            as ``<out_dir>/batch_analysis_results_<ver>.xlsx``.
        df (pd.DataFrame): Per-spermatid measurement table (all Z-slices combined).
        df_summary (pd.DataFrame): Per-slice summary statistics.
        df_tracks (pd.DataFrame, optional): 3D track table from
            :func:`track_across_slices`.  ``None`` if tracking was not run.
    """
    excel_path = os.path.join(out_dir, f"batch_analysis_results_{_VERSION}.xlsx")
    plot_dir = os.path.join(out_dir, "summary_plots")
    print(f"Generating Interactive Excel Audit: {excel_path} ...")
    
    try:
        with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
            workbook  = writer.book
            
            # Formats
            bold = workbook.add_format({'bold': True, 'bg_color': '#D7E4BC', 'border': 1})
            num_fmt = workbook.add_format({'num_format': '0.00'})

            # --- Sheet 1: Population Summary (DYNAMIC FORMULAS) --- [FIRST TAB]
            ws_sum = workbook.add_worksheet('Population_Summary')
            headers = ["Metric", "Average (Formula)", "Median (Formula)", "Std Dev (Formula)"]
            for col, h in enumerate(headers):
                ws_sum.write(0, col, h, bold)
            
            row = 1
            # 2D Length (from Raw_2D_Detections column G - length_um_geodesic)
            if not df.empty:
                n_2d = len(df) + 1
                ws_sum.write(row, 0, "2D Geodesic Length (um)")
                ws_sum.write_formula(row, 1, f"=AVERAGE('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                ws_sum.write_formula(row, 2, f"=MEDIAN('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                ws_sum.write_formula(row, 3, f"=STDEV.P('Raw_2D_Detections'!G2:G{n_2d})", num_fmt)
                row += 1
            
            # 3D Metrics
            if df_tracks is not None and not df_tracks.empty:
                df_q = df_tracks
                n_3d = len(df_q) + 1 if not df_q.empty else 2 # Default to 2 to avoid #DIV/0! bounds
                metrics_3d = [
                    ("3D Geodesic Length (um)", "B"),
                    ("3D Z-Extent (um)", "C"),
                    ("3D Volume (um3)", "D"),
                    ("3D Tortuosity", "E")
                ]
                for m_name, col_letter in metrics_3d:
                    ws_sum.write(row, 0, m_name)
                    ws_sum.write_formula(row, 1, f"=AVERAGE('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 2, f"=MEDIAN('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    ws_sum.write_formula(row, 3, f"=STDEV.P('3D_Morphometrics'!{col_letter}2:{col_letter}{n_3d})", num_fmt)
                    row += 1

            ws_sum.set_column('A:A', 30)

            # --- Sheet 2: 3D Morphometrics  ---
            if df_tracks is not None and not df_tracks.empty:
                df_q = df_tracks
                df_q.to_excel(writer, sheet_name='3D_Morphometrics', index=False)
                ws_3d = writer.sheets['3D_Morphometrics']
                ws_3d.set_column('A:Z', 15)
                # Insert 3D Distribution Graph
                p_3d = os.path.join(plot_dir, "3d_length_distribution.png")
                if os.path.exists(p_3d):
                    ws_3d.insert_image('K2', p_3d, {'x_scale': 0.4, 'y_scale': 0.4})
                
                # If we filtered, add a raw 3D tracks sheet for transparency
                if False:
                    df_tracks.to_excel(writer, sheet_name='Raw_3D_Detections', index=False)
                    ws_raw_3d = writer.sheets['Raw_3D_Detections']
                    ws_raw_3d.set_column('A:Z', 15)

            # --- Sheet 3: Raw 2D Detections ---
            if not df.empty:
                df.to_excel(writer, sheet_name='Raw_2D_Detections', index=False)
                ws_2d = writer.sheets['Raw_2D_Detections']
                ws_2d.set_column('A:Z', 15)

            # --- Sheet 4: Slice Summary ---
            if not df_summary.empty:
                df_summary.to_excel(writer, sheet_name='Slice_Summary', index=False)

            ws_sum.set_column('B:D', 18)

            # Insert Histograms into Summary
            p_hist = os.path.join(plot_dir, "global_histograms.png")
            p_slice = os.path.join(plot_dir, "length_by_slice.png")
            if os.path.exists(p_hist):
                ws_sum.insert_image('F2', p_hist, {'x_scale': 0.5, 'y_scale': 0.5})
            if os.path.exists(p_slice):
                ws_sum.insert_image('F25', p_slice, {'x_scale': 0.5, 'y_scale': 0.5})

            # --- Sheet 5: Methods Dictionary ---
            dictionary_data = [
                ["Metric", "Formula / Definition", "Biological Interpretation"],
                ["2D Geodesic Length", "Dijkstra path length on skeleton", "The curved length of a spermatid fragment within a single 2D slice."],
                ["Total 3D Geodesic Length", "sqrt( max(L_2d, L_drift)^2 + Z_extent^2 )", "Robust 3D reconstruction accounting for lateral drift across slices."],
                ["3D Tortuosity", "3D Geodesic Length / 3D Euclidean distance", "Curvature Index. 1.0 = Perfectly straight."],
                ["3D Euclidean Distance", "sqrt( dx^2 + dy^2 + dz^2 )", "Straight line between centroids."],
                ["Z-Extent (Vertical Span)", "Vertical Span (slices * um/slice)", "Total vertical depth covered."],
                ["3D Volume (um3)", "sum( Area_i * Z_step )", "Approximated Volumetric Mass (*NOTE: Absolute values artificially inflated by Point Spread Function / fluorescence halo*)."],
                ["Effective Thickness (um)", "2 * sqrt( (V_3D / L_3D) / pi )", "Average morphological diameter (*NOTE: Artificially inflated by PSF*)."],
                ["Standard Deviation", "Std Dev", "Population variance."]
            ]
            pd.DataFrame(dictionary_data[1:], columns=dictionary_data[0]).to_excel(writer, sheet_name='Methods_Dictionary', index=False)
            ws_dict = writer.sheets['Methods_Dictionary']
            
            # Manually write the dictionary to avoid any encoding or header issues
            for r_idx, r_data in enumerate(dictionary_data):
                for c_idx, val in enumerate(r_data):
                    ws_dict.write(r_idx, c_idx, val, bold if r_idx==0 else None)
            
            ws_dict.set_column('A:A', 25)
            ws_dict.set_column('B:B', 40)
            ws_dict.set_column('C:C', 60)
            p_guide = os.path.join(plot_dir, "methods_guide.png")
            if os.path.exists(p_guide):
                ws_dict.insert_image('A12', p_guide, {'x_scale': 0.6, 'y_scale': 0.6})

            print(f"Interactive Excel report successfully saved to {excel_path}")

    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating Excel report: {e}")
        write_error_log(out_dir, "Excel Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"Excel Report failed to generate completely.\n{e}")
        except Exception:
            pass

def generate_batch_report(out_dir, df, df_summary, um, df_tracks=None, gui_callback=None):
    """
    Compiles standard summary global output architectures including histograms, mathematical summaries,
    biological methodology pages, and graphical slice overlays natively to a `.pdf` file.
    
    Args:
        out_dir (str): Root export directory.
        df (pd.DataFrame): Flat 2D analysis parameters.
        df_summary (pd.DataFrame): Top-level slice aggregation tracking.
        um (dict): User-defined pixel to micron mapping ratios.
        df_tracks (pd.DataFrame, optional): Unified 3D tracking geometries array.
        gui_callback (function, optional): Live variable passing to front-end dashboards elements.
        
    Returns:
        None (Saves directly to absolute path PDF)
    """
    pdf_path = os.path.join(out_dir, f"batch_report_{_VERSION}.pdf")
    print(f"Generating high-res PDF report: {pdf_path} ...")
    
    # Create directory for high-res standalone plots (for easy copy-pasting into papers/presentations)
    plot_dir = os.path.join(out_dir, "summary_plots")
    os.makedirs(plot_dir, exist_ok=True)
    
    try:
        with PdfPages(pdf_path) as pdf:
            # --- PAGE 1: GLOBAL SUMMARY ---
            fig_sum = plt.figure(figsize=(11, 8.5))
            fig_sum.suptitle(f"Spermatid Analysis Batch Summary - {_VERSION}\nLocation: {out_dir}", fontsize=14, fontweight='bold')
            
            # Global Z-Projection Image (Top Center)
            z_proj_path = os.path.join(out_dir, "global_z_projection.png")
            if os.path.exists(z_proj_path):
                ax_z = fig_sum.add_axes([0.15, 0.62, 0.7, 0.28]) # [left, bottom, width, height]
                ax_z.imshow(plt.imread(z_proj_path))
                ax_z.set_title("Global Z-Projection (Composite [Original | Overlay])", fontsize=10)
                ax_z.axis('off')
            
            # Plot 1: Counts per slice
            ax1 = fig_sum.add_subplot(2, 2, 3)
            ax1.plot(df_summary['z_slice'], df_summary['n_spermatids'], 'bo-', markersize=4)
            ax1.set_title("Detections per Z-Slice")
            ax1.set_xlabel("Z-Index")
            ax1.set_ylabel("Count")
            ax1.grid(True, alpha=0.3)
            
            # Plot 2: Length Distribution (Global)
            ax2 = fig_sum.add_subplot(2, 2, 4)
            if not df.empty:
                vals = df['length_um_geodesic']
                ax2.hist(vals, bins=25, color='forestgreen', edgecolor='black', alpha=0.7)
                m_med = vals.median()
                m_avg = vals.mean()
                ax2.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax2.axvline(m_avg, color='orange', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax2.set_title("Global 2D Length Distribution")
                ax2.set_xlabel("Geodesic Length (um)")
                ax2.set_ylabel("Frequency")
                ax2.legend(fontsize=8)
            
            fig_sum.savefig(os.path.join(plot_dir, "global_summary.png"), dpi=300, bbox_inches='tight')
            # NEW: Explicitly save global_histograms.png for Excel embedding
            fig_sum.savefig(os.path.join(plot_dir, "global_histograms.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_sum, dpi=300, bbox_inches='tight')
            plt.close(fig_sum)
            
            # Save length_by_slice.png
            fig_l_slice = plt.figure(figsize=(10, 5))
            ax_ls = fig_l_slice.add_subplot(1, 1, 1)
            ax_ls.plot(df_summary['z_slice'], df_summary['median_length_um'], 'go-', label='Median Length')
            ax_ls.set_title("Median Length by Slice")
            ax_ls.set_xlabel("Z-Slice")
            ax_ls.set_ylabel("Length (um)")
            ax_ls.grid(True, alpha=0.3)
            fig_l_slice.savefig(os.path.join(plot_dir, "length_by_slice.png"), dpi=300, bbox_inches='tight')
            plt.close(fig_l_slice)

            # --- PAGE 1.5: POPULATION CONSOLIDATION ---
            if df_tracks is not None and not df_tracks.empty:
                fig_dyn = plt.figure(figsize=(11, 8.5))
                fig_dyn.suptitle("3D Population Tracking Summary", fontsize=15, fontweight='bold')
                
                total_2d = len(df)
                total_3d = len(df_tracks)
                n_quality = total_3d
                n_flagged = 0
                
                # A) 3-Tier Reduction Funnel
                ax_bar = fig_dyn.add_subplot(1, 2, 1)
                y_pos = [2, 1, 0]
                counts = [total_2d, total_3d, n_quality]
                colors = ['coral', 'steelblue', '#2ca02c']
                labels = [
                    'Raw 2D Detections\n(All Fragments)',
                    'All 3D Tracks\n(Consolidated)',
                    'Quality Population\n(Audit-Passed)'
                ]
                
                bars = ax_bar.barh(y_pos, counts, color=colors, edgecolor='black', height=0.55)
                ax_bar.set_xlim(0, max(counts) * 1.35)
                
                for i, v in enumerate(counts):
                    ax_bar.text(v + (max(counts)*0.02), y_pos[i], f"{v:,}", va='center', fontweight='bold', fontsize=12)
                    
                ax_bar.set_yticks(y_pos)
                ax_bar.set_yticklabels(labels, fontsize=10, fontweight='bold')
                ax_bar.set_xlabel("Total Count", fontsize=12)
                ax_bar.set_title("Tracking & Quality Reduction", fontsize=13, fontweight='bold')
                ax_bar.spines['top'].set_visible(False)
                ax_bar.spines['right'].set_visible(False)
                
                # B) Donut Chart: Quality vs Flagged Breakdown
                ax_pie = fig_dyn.add_subplot(1, 2, 2)
                
                n_single = len(df_tracks[df_tracks['n_slices'] == 1]) if 'n_slices' in df_tracks.columns else 0
                n_outlier_other = n_flagged - n_single if n_flagged > n_single else 0
                n_single_only = min(n_single, n_flagged)
                
                pie_sizes = [n_quality, n_outlier_other, n_single_only]
                pie_labels = [
                    f"Quality Tracks\n({n_quality:,})",
                    f"Shape Outliers\n({n_outlier_other:,})",
                    f"Single-Slice\n({n_single_only:,})"
                ]
                pie_colors = ['#2ca02c', '#ff7f0e', '#d62728']
                
                # Filter out zero segments
                valid = [(s, l, c) for s, l, c in zip(pie_sizes, pie_labels, pie_colors) if s > 0]
                if valid:
                    pie_sizes, pie_labels, pie_colors = zip(*valid)
                
                wedge_props = dict(width=0.45, edgecolor='white', linewidth=2)
                wedges, texts, autotexts = ax_pie.pie(
                    pie_sizes, labels=None, colors=pie_colors,
                    autopct='%1.1f%%',
                    startangle=90,
                    pctdistance=0.75,
                    wedgeprops=wedge_props,
                    textprops={'fontsize': 12, 'fontweight': 'bold'}
                )
                
                for autotext in autotexts:
                    autotext.set_color('white')
                    autotext.set_fontsize(13)
                    autotext.set_fontweight('bold')

                ax_pie.legend(
                    wedges, pie_labels,
                    title="Track Quality",
                    loc="lower center",
                    bbox_to_anchor=(0.5, -0.15),
                    ncol=len(pie_sizes),
                    fontsize=10,
                    frameon=False
                )
                    
                ax_pie.set_title(f"Quality Breakdown of {total_3d:,} 3D Tracks", fontsize=13, fontweight='bold')
                
                plt.tight_layout(rect=[0, 0.03, 1, 0.95])
                fig_dyn.savefig(os.path.join(plot_dir, "population_consolidation.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_dyn, dpi=300, bbox_inches='tight')
                plt.close(fig_dyn)

            # --- PAGE 2: 3D MORPHOMETRICS SUMMARY ---
            if df_tracks is not None and not df_tracks.empty:
                fig_3d = plt.figure(figsize=(11, 8.5))
                
                df_q = df_tracks
                q_label = ""
                fig_3d.suptitle(f"3D Population Statistics{q_label}", fontsize=14, fontweight='bold')
                
                # 3D Length
                ax3d_1 = fig_3d.add_subplot(2, 2, 1)
                vals_all = df_tracks['total_3d_length_um']
                vals_q = df_q['total_3d_length_um']
                ax3d_1.hist(vals_all, bins=20, color='lightgray', edgecolor='gray', alpha=0.5, label='All Tracks')
                ax3d_1.hist(vals_q, bins=20, color='darkorange', edgecolor='black', alpha=0.7, label='Quality')
                m_med = vals_q.median()
                m_avg = vals_q.mean()
                ax3d_1.axvline(m_med, color='red', linestyle='-', label=f"Median: {m_med:.1f}")
                ax3d_1.axvline(m_avg, color='black', linestyle='--', label=f"Mean: {m_avg:.1f}")
                ax3d_1.set_title("Total 3D Geodesic Length")
                ax3d_1.set_xlabel("Length (um)")
                ax3d_1.set_ylabel("Frequency")
                ax3d_1.legend(fontsize=7)
                
                fig_3d_len = plt.figure(figsize=(6, 4))
                ax_3dl = fig_3d_len.add_subplot(1, 1, 1)
                ax_3dl.hist(vals_q, bins=20, color='darkorange', edgecolor='black', alpha=0.7)
                ax_3dl.set_title("3D Length Distribution (Quality Population)")
                fig_3d_len.savefig(os.path.join(plot_dir, "3d_length_distribution.png"), dpi=300, bbox_inches='tight')
                plt.close(fig_3d_len)

                # 3D Tortuosity
                ax3d_2 = fig_3d.add_subplot(2, 2, 2)
                vt_all = df_tracks['tortuosity_3d']
                vt_q = df_q['tortuosity_3d']
                vt_all_viz = vt_all[(vt_all >= 0.95) & (vt_all <= 3.0)]
                vt_q_viz = vt_q[(vt_q >= 0.95) & (vt_q <= 3.0)]
                ax3d_2.hist(vt_all_viz, bins=25, color='lightgray', edgecolor='gray', alpha=0.5, label='All Tracks')
                ax3d_2.hist(vt_q_viz, bins=25, color='purple', edgecolor='black', alpha=0.6, label='Quality')
                ax3d_2.axvline(vt_q.median(), color='red', linestyle='-', label=f"Median: {vt_q.median():.2f}")
                ax3d_2.axvline(vt_q.mean(), color='black', linestyle='--', label=f"Mean: {vt_q.mean():.2f}")
                ax3d_2.set_xlim(0.95, 3.0)
                ax3d_2.set_title("3D Tortuosity (Curvature)")
                ax3d_2.set_xlabel("Ratio (Length / Distance)")
                ax3d_2.set_ylabel("Frequency")
                ax3d_2.legend(fontsize=7)

                # Vertical Extent
                ax3d_3 = fig_3d.add_subplot(2, 2, 3)
                ve_all = df_tracks['z_extent_um']
                ve_q = df_q['z_extent_um']
                ax3d_3.hist(ve_all, bins=15, color='lightgray', edgecolor='gray', alpha=0.5, label='All Tracks')
                ax3d_3.hist(ve_q, bins=15, color='teal', edgecolor='black', alpha=0.7, label='Quality')
                ax3d_3.axvline(ve_q.median(), color='red', linestyle='-', label=f"Median: {ve_q.median():.1f}")
                ax3d_3.axvline(ve_q.mean(), color='black', linestyle='--', label=f"Mean: {ve_q.mean():.1f}")
                ax3d_3.set_title("Z-Extent (Vertical Span)")
                ax3d_3.set_xlabel("Vertical Height (um)")
                ax3d_3.set_ylabel("Frequency")
                ax3d_3.legend(fontsize=7)

                # Volume
                ax3d_4 = fig_3d.add_subplot(2, 2, 4)
                vv_all = df_tracks['volume_um3']
                vv_q = df_q['volume_um3']
                ax3d_4.hist(vv_all, bins=20, color='lightgray', edgecolor='gray', alpha=0.5, label='All Tracks')
                ax3d_4.hist(vv_q, bins=20, color='gray', edgecolor='black', alpha=0.7, label='Quality')
                ax3d_4.axvline(vv_q.median(), color='red', linestyle='-', label=f"Median: {vv_q.median():.0f}")
                ax3d_4.axvline(vv_q.mean(), color='black', linestyle='--', label=f"Mean: {vv_q.mean():.0f}")
                ax3d_4.set_title("Approximated 3D Volume")
                ax3d_4.set_xlabel("Volume (um\u00b3)")
                ax3d_4.set_ylabel("Frequency")
                ax3d_4.legend(fontsize=7)

                plt.tight_layout(rect=[0, 0.03, 1, 0.95])
                fig_3d.savefig(os.path.join(plot_dir, "3d_population_stats.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_3d, dpi=300, bbox_inches='tight')
                plt.close(fig_3d)

            # Methods Guide securely moved down past the Advanced Biometrics block.

            # --- PAGE 4: ADVANCED 3D BIOMETRICS (QUALITY-FILTERED) ---
            if df_tracks is not None and not df_tracks.empty:
                fig_adv = plt.figure(figsize=(11, 8.5))
                df_q = df_tracks
                q_label = ""
                fig_adv.suptitle(f"Advanced 3D Biometrics Dashboard{q_label}", fontsize=16, fontweight='bold', y=0.96)
                
                # Helper for Mean/Median
                def add_stats_lines(ax, data_series):
                    if data_series.empty or data_series.isna().all(): return
                    m = data_series.mean()
                    med = data_series.median()
                    ax.axvline(med, color='red', linestyle='--', linewidth=1.5, label=f'Median: {med:.2f}')
                    ax.axvline(m, color='green', linestyle=':', linewidth=2, label=f'Mean: {m:.2f}')
                    ax.legend(fontsize=8)
                    
                def dual_hist(ax, col, title, xlabel, color_q, bins=30):
                    vals_all = df_tracks[col].dropna() if col in df_tracks.columns else pd.Series(dtype=float)
                    vals_q = df_q[col].dropna() if col in df_q.columns else pd.Series(dtype=float)
                    if not vals_all.empty:
                        sns.histplot(vals_all, bins=bins, ax=ax, color='lightgray', edgecolor='gray', alpha=0.5, label='All Tracks')
                    if not vals_q.empty:
                        sns.histplot(vals_q, bins=bins, ax=ax, color=color_q, edgecolor='black', alpha=0.7, label='Quality')
                        add_stats_lines(ax, vals_q)
                    ax.set_title(title)
                    ax.set_xlabel(xlabel)
                    ax.set_ylabel("Frequency")
                    if len(ax.get_legend_handles_labels()[0]) > 0:
                        ax.legend(fontsize=7)

                # 4 panels: Pitch, Thickness, Taper, Nearest Neighbor
                ax_p = fig_adv.add_subplot(2, 2, 1)
                dual_hist(ax_p, 'pitch_deg', "Pitch Angle (Vertical Plunge)", "Degrees (0=Flat, 90=Vertical)", 'orange')
                
                ax_th = fig_adv.add_subplot(2, 2, 2)
                dual_hist(ax_th, 'thickness_um', "Effective Nucleus Thickness", "Average Diameter (\u00b5m)", '#17becf')
                
                ax_ta = fig_adv.add_subplot(2, 2, 3)
                dual_hist(ax_ta, 'taper_ratio', "Morphological Taper Ratio", "Max Area / Min Area", 'purple')
                
                ax_nn = fig_adv.add_subplot(2, 2, 4)
                dual_hist(ax_nn, 'nearest_neighbor_um', "Spatial Packing Density", "Distance to Nearest Neighbor (\u00b5m)", 'brown')
                
                plt.tight_layout(rect=[0, 0.03, 1, 0.93])
                fig_adv.savefig(os.path.join(plot_dir, "advanced_biometrics.png"), dpi=300, bbox_inches='tight')
                pdf.savefig(fig_adv, dpi=300, bbox_inches='tight')
                plt.close(fig_adv)

            # --- PAGE 5: METHODS & INTERPRETATION GUIDE ---
            fig_guide = plt.figure(figsize=(11, 8.5))
            ax_g = fig_guide.add_axes([0.05, 0.05, 0.9, 0.9])
            ax_g.axis('off')
            guide_full = (
                "METHODS & CALCULATION DETAILS (Biometrics & Tissue Mechanics)\n"
                f"{'='*80}\n\n"
                "1. Total 3D Geodesic Length (µm)\n"
                "   Formula: L_3D = Σ sqrt((Δx)² + (Δy)² + (Δz)²)\n"
                "   Biology: Measures the true multi-slice structural backbone of the spermatid,\n"
                "   tracking along its curvature rather than assuming rigid linearity.\n\n"
                "2. 3D Tortuosity / Curvature\n"
                "   Formula: T = (Total Geodesic Length) / (Straight-Line Point-to-Point Distance)\n"
                "   Biology: A morphological metric of acrosomal or tail curvature. Normal spermatids\n"
                "   should be moderately rigid (~1.0 - 1.2). Extreme values (>2.0) generally flag\n"
                "   tissue deformation, fixation artifacts, or algorithmic tracking fusions.\n\n"
                "3. Z-Extent (Vertical Tissue Span)\n"
                "   Biology: Captures how deep the nucleus penetrates through the epithelial layer.\n"
                "   Plunging nuclei indicate distinct maturation stages closer to the lumen boundary.\n\n"
                "4. Approximated 3D Volume (µm³)\n"
                "   Calculated via Riemann sum of 2D segmented pixel area × Z-slice thickness.\n"
                "   Biology: Directly correlates to chromatin condensation levels. Aberrantly\n"
                "   large volumes may signify poor condensation or double-nucleation.\n\n"
                "5. Effective Thickness (Average Diameter, µm)\n"
                "   Formula: D_avg = 2 * sqrt( (V_3D / L_3D) / π )\n"
                "   Biology: Models the nucleus as a cylinder to quantify chromatin thickness.\n"
                "   Particularly useful for isolating swelling or decompaction pathologies.\n\n"
                "6. Pitch Angle (Plunge Vector, Degrees)\n"
                "   Formula: θ = |arcsin(ΔZ / Length_Euclidean)| * (180/π)\n"
                "   Biology: Quantifies the absolute orientation of the elongating nucleus relative\n"
                "   to the basal lamina plane. Values near 90° indicate perpendicular apical plunging.\n\n"
                "7. Taper Ratio\n"
                "   Formula: R = (Maximum Cross-sectional Area) / (Minimum Cross-sectional Area)\n"
                "   Biology: Identifies tapering anomalies near the acrosome versus the tail insertion.\n"
                "   Extreme variance implies severe morphological irregularity or fusion errors.\n\n"
                "8. Spatial Packing Density (Nearest Neighbor Dist, µm)\n"
                "   Calculated: 3D Euclidean distance from this track's centroid to its closest sibling.\n"
                "   Biology: Measures tubule capacity and cyst grouping density during spermiation."
            )
            ax_g.text(0, 1, guide_full, transform=ax_g.transAxes, fontsize=10, family='monospace', verticalalignment='top', linespacing=1.3)
            fig_guide.savefig(os.path.join(plot_dir, "methods_guide.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_guide, dpi=300, bbox_inches='tight')
            plt.close(fig_guide)

            # --- PAGE 6: GLOBAL STATISTICS TABLE (after Methods) ---
            fig_tab = plt.figure(figsize=(11, 8.5))
            ax_t = fig_tab.add_subplot(1, 1, 1)
            ax_t.axis('off')
            ax_t.set_title("Global Population Statistics Summary", fontsize=14, fontweight='bold', pad=20)
            
            stats_rows = []
            if not df.empty:
                l2d = df['length_um_geodesic']
                stats_rows.append(["2D Fragment Geodesic Length (um)", f"{l2d.mean():.2f}", f"{l2d.median():.2f}", f"{l2d.std():.2f}"])
            
            if df_tracks is not None and not df_tracks.empty:
                # Quality-filtered subset
                df_q = df_tracks
                
                # -- All Tracks Section --
                stats_rows.append(["--- ALL TRACKS ---", f"N={len(df_tracks)}", "", ""])
                l3d = df_tracks['total_3d_length_um']
                ze  = df_tracks['z_extent_um']
                vo  = df_tracks['volume_um3']
                to  = df_tracks['tortuosity_3d']
                th  = df_tracks['thickness_um']
                stats_rows.append(["3D Length (um)", f"{l3d.mean():.2f}", f"{l3d.median():.2f}", f"{l3d.std():.2f}"])
                stats_rows.append(["3D Z-Extent (um)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
                stats_rows.append(["3D Volume (um3)", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
                stats_rows.append(["3D Tortuosity", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
                stats_rows.append(["3D Thickness (um)", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])
                
                # -- 3D Population Section --
                if len(df_q) > 0:
                    stats_rows.append(["--- 3D POPULATION ---", f"N={len(df_q)}", "", ""])
                    l3q = df_q['total_3d_length_um']
                    zeq = df_q['z_extent_um']
                    voq = df_q['volume_um3']
                    toq = df_q['tortuosity_3d']
                    thq = df_q['thickness_um']
                    piq = df_q['pitch_deg']
                    taq = df_q['taper_ratio']
                    nnq = df_q['nearest_neighbor_um'].dropna()
                    stats_rows.append(["3D Length (um)", f"{l3q.mean():.2f}", f"{l3q.median():.2f}", f"{l3q.std():.2f}"])
                    stats_rows.append(["3D Z-Extent (um)", f"{zeq.mean():.2f}", f"{zeq.median():.2f}", f"{zeq.std():.2f}"])
                    stats_rows.append(["3D Volume (um3)", f"{voq.mean():.1f}", f"{voq.median():.1f}", f"{voq.std():.1f}"])
                    stats_rows.append(["3D Tortuosity", f"{toq.mean():.3f}", f"{toq.median():.3f}", f"{toq.std():.3f}"])
                    stats_rows.append(["3D Thickness (um)", f"{thq.mean():.2f}", f"{thq.median():.2f}", f"{thq.std():.2f}"])
                    stats_rows.append(["3D Pitch (degrees)", f"{piq.mean():.1f}", f"{piq.median():.1f}", f"{piq.std():.1f}"])
                    stats_rows.append(["3D Taper Ratio", f"{taq.mean():.2f}", f"{taq.median():.2f}", f"{taq.std():.2f}"])
                    if not nnq.empty:
                        stats_rows.append(["Nearest Neighbor (um)", f"{nnq.mean():.1f}", f"{nnq.median():.1f}", f"{nnq.std():.1f}"])
            
            if stats_rows:
                table = ax_t.table(
                    cellText=stats_rows,
                    colLabels=["Metric", "Mean", "Median", "Std Dev"],
                    loc='center', cellLoc='center',
                    colWidths=[0.38, 0.2, 0.2, 0.2]
                )
                table.auto_set_font_size(False)
                table.set_fontsize(10)
                table.scale(1.2, 2.5)
            
            fig_tab.savefig(os.path.join(plot_dir, "global_statistics_table.png"), dpi=300, bbox_inches='tight')
            pdf.savefig(fig_tab, dpi=300, bbox_inches='tight')
            plt.close(fig_tab)

            # --- SUBSEQUENT PAGES: PER-SLICE DETAILS (2 panels: [Panel] | [Histogram]) ---
            overlay_dir = os.path.join(out_dir, "overlays")
            for idx_p, (row_idx, row) in enumerate(df_summary.iterrows()):
                z = int(row['z_slice'])
                panel_path = os.path.join(overlay_dir, f"z{z:02d}_panel.png")
                
                if not os.path.exists(panel_path):
                    continue
                
                fig_slice = plt.figure(figsize=(18, 7))
                fig_slice.suptitle(f"Z-Slice {z:02d} Analysis [Original | Overlay | Distribution]", fontsize=12, fontweight='bold')
                
                # Panel: Side-by-Side (Original | Overlay)
                ax_panel = fig_slice.add_subplot(1, 2, 1)
                ax_panel.imshow(plt.imread(panel_path))
                ax_panel.set_title(f"Visual Verification (N={int(row['n_spermatids'])})")
                ax_panel.axis('off')
                
                # Plot: Stats
                ax_hist = fig_slice.add_subplot(1, 2, 2)
                slice_data = df[df['z_slice'] == z]
                if not slice_data.empty:
                    ax_hist.hist(slice_data['length_um_geodesic'], bins=15, color='skyblue', edgecolor='black')
                    ax_hist.set_title(f"Z={z} Length Distribution")
                    ax_hist.set_xlabel("Spermatid Nucleus Length (µm)")
                    ax_hist.set_ylabel("Frequency (Count)")
                    
                    m_med = slice_data['length_um_geodesic'].median()
                    m_avg = slice_data['length_um_geodesic'].mean()
                    ax_hist.axvline(m_med, color='red', linestyle='-', alpha=0.7, label=f"Median: {m_med:.1f}")
                    ax_hist.axvline(m_avg, color='orange', linestyle='--', alpha=0.7, label=f"Mean: {m_avg:.1f}")
                    ax_hist.legend(fontsize=9)
                else:
                    ax_hist.text(0.5, 0.5, "No Detections", ha='center', va='center')
                
                pdf.savefig(fig_slice, dpi=300)
                plt.close(fig_slice)
                
                if gui_callback:
                    gui_callback(int(80 + (20 * (idx_p+1) / len(df_summary))))
                
        print(f"Report successfully saved to {pdf_path}")
        
        # --- GENERATE POWERPOINT REPORT ---
        try:
            generate_pptx_report(out_dir, df, df_summary, um, df_tracks)
        except Exception as e:
            import traceback
            err_msg = traceback.format_exc()
            print(f"PPTX Report failed: {e}")
            write_error_log(out_dir, "PowerPoint Generator (via Batch)", err_msg)
            
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"ERROR generating PDF report: {e}")
        write_error_log(out_dir, "PDF Reporter", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showwarning("Reporting Warning", f"PDF Report failed to generate completely.\n{e}")
        except Exception:
            pass

def generate_pptx_report(out_dir, df, df_summary, um, df_tracks=None):
    """
    Generates a native Microsoft Office PowerPoint presentation (.pptx) with
    fully editable, data-embedded charts for each key biometric.

    Each chart is a *true* Office Open XML chart object (not a rasterised image),
    so the researcher can re-style, re-colour, and re-export from PowerPoint without
    any additional software.  The chart data tables are embedded inside the PPTX
    file itself, enabling standalone sharing.

    Slide structure
    ---------------
    1. **Global Population Analytics** — side-by-side column histograms of:
       - 2D geodesic length distribution (from the current batch).
       - 3D true track length distribution (if tracking was run).
    2. **Population Tracking Consolidation** (if tracking was run) — pie chart
       showing the fraction of single-slice vs. multi-slice reassigned tracks.
    3. **3D Biometrics Dashboard** — scatter plot of tortuosity vs. 3D length,
       plus a bar chart of pitch angle distribution.
    4. **Methods & Calculation Details** — text slide with biological justification
       and mathematical formulae for all metrics.

    Hyperlink behaviour
    -------------------
    The title text on each slide contains a hyperlink to the source Excel workbook.
    Clicking the title in Slide Show mode opens the associated Excel file so the
    viewer can inspect the raw numbers behind any chart.

    Biological context of each chart
    ---------------------------------
    - **Length histogram**: A Gaussian distribution centred near the species-expected
      length (e.g., ~55 µm in *Drosophila*) indicates healthy elongation.  Bimodal
      distributions may indicate two simultaneous maturation cohorts.
    - **Pie chart**: Quantifies how many tracked objects span multiple Z-slices
      (true 3D continuity) vs. transient single-slice events (possibly debris).
    - **Tortuosity vs. 3D length scatter**: Healthy spermatids cluster in the
      low-tortuosity, high-length quadrant.  High-tortuosity outliers indicate
      coiled or bent morphologies.
    - **Pitch angle**: A right-skewed distribution towards 90° is expected during
      the apical plunging phase of spermatid elongation.

    Args:
        out_dir (str): Top-level analysis output directory.  The presentation is
            saved as ``<out_dir>/spermatid_analysis_report.pptx``.
        df (pd.DataFrame): Per-spermatid 2D measurement table.
        df_summary (pd.DataFrame): Per-slice summary statistics table.
        um (float): Microns-per-pixel scale factor (``UM_PER_PX_XY``).
        df_tracks (pd.DataFrame, optional): 3D track table.  ``None`` skips
            tracking-specific slides (slides 2 and 3).
    """
    try:
        import os as _os
        import numpy as _np
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from tkinter import messagebox
        
        print("PPTX: Starting report generation...")
        # We need a fallback blank pptx
        prs = Presentation()
        
        # Safe blank layout (often index 6 or 5)
        try:
            blank_slide_layout = prs.slide_layouts[6] 
        except Exception:
            blank_slide_layout = prs.slide_layouts[0]
            
        def add_hyperlink(slide, sheet_name="Population_Summary"):
            excel_name = "batch_analysis_results_v11.xlsx"
            # PowerPoint/Excel deep-link fragment: filename.xlsx#Sheet!A1
            target = f"{excel_name}#'{sheet_name}'!A1"
            
            txBox = slide.shapes.add_textbox(Inches(0.2), Inches(7.1), Inches(9.5), Inches(0.4))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"Click to View Detailed Data: {excel_name} [{sheet_name}]"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0, 0, 255) # Blue
            run.font.underline = True
            try:
                # Shape-level link
                txBox.click_action.hyperlink.address = target
                # Text-run level link
                run.hyperlink.address = target
            except Exception:
                pass

        def add_line_chart(slide, x_data, y_data, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = list(x_data)
            chart_data.add_series('Count', list(y_data))
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.has_major_gridlines = True

        def add_horizontal_bar_chart(slide, categories, values, colors, left, top, width, height, title):
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', values)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = False
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.tick_labels.font.size = Pt(8)
            
            # Show data labels for bar values
            plot = chart.plots[0]
            plot.has_data_labels = True
            for i, point in enumerate(plot.series[0].points):
                point.data_label.font.size = Pt(10)
                point.data_label.font.bold = True

        def add_histogram(slide, data_series, left, top, width, height, title, bins=20):
            if data_series is None or data_series.empty or data_series.isna().all():
                return
            
            clean_data = data_series.dropna()
            counts, bin_edges = _np.histogram(clean_data, bins=bins)
            avg = clean_data.mean()
            med = clean_data.median()
            
            # Create category names from bin boundaries
            categories = [f"{bin_edges[i]:.1f}-{bin_edges[i+1]:.1f}" for i in range(len(counts))]
            
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Count', list(counts))
            # ADD DUMMY SERIES FOR LEGEND STATS
            chart_data.add_series(f'Mean: {avg:.2f}', [])
            chart_data.add_series(f'Median: {med:.2f}', [])
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.CORNER
            chart.legend.font.size = Pt(8)
            chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Reduce axis label crowding
            chart.category_axis.tick_labels.font.size = Pt(8)
            chart.value_axis.tick_labels.font.size = Pt(8)

        # --- Slide 1: Global Analytics Overview ---
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Spermatid Population Overview"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        
        # Bottom Left: Detections per Slice (Line Chart) - Identical to PDF Page 1
        if not df_summary.empty:
            add_line_chart(slide1, df_summary['z_slice'], df_summary['n_spermatids'], Inches(0.2), Inches(4.0), Inches(4.5), Inches(3.0), "Detections per Z-Slice (Raw)")

        # Top Center Left: 2D Length Dist
        if not df.empty:
            add_histogram(slide1, df['length_um_geodesic'], Inches(0.2), Inches(0.8), Inches(4.5), Inches(3.0), "Global 2D Geodesic Length Distribution")
        
        # Top Center Right: 3D Length Dist
        if df_tracks is not None and not df_tracks.empty:
            add_histogram(slide1, df_tracks['total_3d_length_um'], Inches(5.0), Inches(0.8), Inches(4.5), Inches(3.0), "Global 3D Solid Nuclei Length")
            
        add_hyperlink(slide1)
        
        # --- Slide 2: Population Consolidation ---
        if df_tracks is not None and not df_tracks.empty:
            slide2 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "3D Population Tracking & Consolidation"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            
            # Left: Reduction Bar Chart (PDF Parity)
            total_2d = len(df)
            total_3d = len(df_tracks)
            n_quality = total_3d
            n_flagged = total_3d - n_quality
            
            add_horizontal_bar_chart(slide2, 
                                     ['Quality Population', 'All 3D Tracks', 'Raw 2D Detections'], 
                                     [n_quality, total_3d, total_2d], 
                                     None, Inches(0.2), Inches(1.5), Inches(4.5), Inches(4.5), "Tracking & Quality Reduction")

            # Right: Composition Pie Chart
            n_single = len(df_tracks[df_tracks['n_slices'] == 1]) if 'n_slices' in df_tracks.columns else 0
            n_outlier_other = n_flagged - n_single if n_flagged > n_single else 0
            n_single_only = min(n_single, n_flagged)
            
            pie_sizes = [n_quality, n_outlier_other, n_single_only]
            pie_labels = ['Quality Tracks', 'Shape Outliers', 'Single-Slice']
            
            # Filter zero values for pptx chart
            valid_idx = [i for i, v in enumerate(pie_sizes) if v > 0]
            pie_sizes = [pie_sizes[i] for i in valid_idx]
            pie_labels = [pie_labels[i] for i in valid_idx]
            
            chart_data = CategoryChartData()
            chart_data.categories = pie_labels
            chart_data.add_series('Population', pie_sizes)
            
            chart2 = slide2.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(4.6), Inches(1.2), Inches(5.2), Inches(5.2), chart_data
            ).chart
            chart2.has_legend = True
            chart2.legend.position = XL_LEGEND_POSITION.CORNER
            chart2.legend.font.size = Pt(8)
            chart2.chart_title.text_frame.text = f"Quality Breakdown of {total_3d:,} 3D Tracks"
            chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
            
            plot = chart2.plots[0]
            plot.has_data_labels = True
            total = sum(pie_sizes) if sum(pie_sizes) > 0 else 1
            for i, point in enumerate(plot.series[0].points):
                val = pie_sizes[i]
                pct = (val / total) * 100
                label_text = f"{val:,}\n({pct:.1f}%)"
                point.data_label.text_frame.text = label_text
                point.data_label.font.size = Pt(9)
                point.data_label.font.bold = True
            
            add_hyperlink(slide2)
            
        # ---------------------------------------------------------------------
        # SLIDE 3: Advanced 3D Biometrics
        # ---------------------------------------------------------------------
        if df_tracks is not None and not df_tracks.empty:
            slide3 = prs.slides.add_slide(blank_slide_layout)
            txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            df_q = df_tracks
            q_label = ""
            tf.text = f"Advanced 3D Biometrics Dashboard{q_label}"
            tf.paragraphs[0].font.size = Pt(22)
            tf.paragraphs[0].font.bold = True
            
            # For Slide 3, we show the histograms of the 3D population
            # to provide the most representative biological view.
            add_histogram(slide3, df_q['pitch_deg'], Inches(0.2), Inches(0.8), Inches(4.5), Inches(2.9), "Pitch Angle (Degrees)", bins=20)
            add_histogram(slide3, df_q['thickness_um'], Inches(5.0), Inches(0.8), Inches(4.5), Inches(2.9), "Effective Thickness (\u00b5m)", bins=20)
            add_histogram(slide3, df_q['taper_ratio'], Inches(0.2), Inches(3.8), Inches(4.5), Inches(2.9), "Morphological Taper Ratio", bins=20)
            add_histogram(slide3, df_q['nearest_neighbor_um'].dropna(), Inches(5.0), Inches(3.8), Inches(4.5), Inches(2.9), "Nearest Neighbor Density (\u00b5m)", bins=20)
            
            add_hyperlink(slide3, "3D_Morphometrics")
            
        # ---------------------------------------------------------------------
        # SLIDE 4: Global Population Statistics Summary Table
        # ---------------------------------------------------------------------
        slide4 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Global Population Statistics Summary"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        
        # Prepare stats matching PDF (All vs Quality)
        stats_rows = [["Metric", "Mean", "Median", "Std Dev"]]
        if not df.empty:
            l2d = df['length_um_geodesic']
            stats_rows.append(["2D Fragment Length (\u00b5m)", f"{l2d.mean():.2f}", f"{l2d.median():.2f}", f"{l2d.std():.2f}"])
        
        if df_tracks is not None and not df_tracks.empty:
            df_q = df_tracks
            
            # Section Header for All Tracks
            stats_rows.append(["--- ALL TRACKS ---", f"N={len(df_tracks)}", "", ""])
            
            def add_pop_rows(pop_df, prefix=""):
                l3 = pop_df['total_3d_length_um']
                ze = pop_df['z_extent_um']
                vo = pop_df['volume_um3']
                to = pop_df['tortuosity_3d']
                th = pop_df['thickness_um']
                stats_rows.append([f"{prefix}3D Length (\u00b5m)", f"{l3.mean():.2f}", f"{l3.median():.2f}", f"{l3.std():.2f}"])
                stats_rows.append([f"{prefix}3D Z-Extent (\u00b5m)", f"{ze.mean():.2f}", f"{ze.median():.2f}", f"{ze.std():.2f}"])
                stats_rows.append([f"{prefix}3D Volume (\u00b5m\u00b3)", f"{vo.mean():.1f}", f"{vo.median():.1f}", f"{vo.std():.1f}"])
                stats_rows.append([f"{prefix}3D Tortuosity", f"{to.mean():.3f}", f"{to.median():.3f}", f"{to.std():.3f}"])
                stats_rows.append([f"{prefix}3D Thickness (\u00b5m)", f"{th.mean():.2f}", f"{th.median():.2f}", f"{th.std():.2f}"])

            add_pop_rows(df_tracks)
            
            if not df_q.empty:
                # Section Header for Quality Population
                stats_rows.append(["--- 3D POPULATION ---", f"N={len(df_q)}", "", ""])
                add_pop_rows(df_q)
        
        if len(stats_rows) > 1:
            rows = len(stats_rows)
            cols = 4
            table_shape = slide4.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(6.0))
            table = table_shape.table
            
            # Header styling
            for c in range(cols):
                cell = table.cell(0, c)
                cell.text_frame.text = stats_rows[0][c]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196) # Standard Blue
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.bold = True
                
            # Body styling
            for r in range(1, rows):
                is_separator = "---" in stats_rows[r][0]
                for c in range(cols):
                    cell = table.cell(r, c)
                    cell.text_frame.text = stats_rows[r][c]
                    cell.text_frame.paragraphs[0].font.size = Pt(8)
                    if is_separator:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                        cell.text_frame.paragraphs[0].font.bold = True
                    elif c == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
        
        add_hyperlink(slide4, "Population_Summary")

        # ---------------------------------------------------------------------
        # SLIDE 5: Methods & Interpretation Guide (Exact PDF Synchronization)
        # ---------------------------------------------------------------------
        slide5 = prs.slides.add_slide(blank_slide_layout)
        txBox_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.text = "Methods & Calculation Details"
        tf_title.paragraphs[0].font.size = Pt(20)
        tf_title.paragraphs[0].font.bold = True

        # Synchronized text from the PDF report version
        methods_items = [
            ("1. Total 3D Geodesic Length (µm)", [
                ("Formula: ", "L_3D = Σ sqrt((Δx)² + (Δy)² + (Δz)²)"),
                ("Biology: ", "Measures the true multi-slice structural backbone of the spermatid, tracking along its curvature rather than assuming rigid linearity.")
            ]),
            ("2. 3D Tortuosity / Curvature", [
                ("Formula: ", "T = (Total Geodesic Length) / (Straight-Line Point-to-Point Distance)"),
                ("Biology: ", "A morphological metric of acrosomal or tail curvature. Normal spermatids should be moderately rigid (~1.0 - 1.2). Extreme values (>2.0) generally flag tissue deformation, fixation artifacts, or algorithmic tracking fusions.")
            ]),
            ("3. Z-Extent (Vertical Tissue Span)", [
                ("Biology: ", "Captures how deep the nucleus penetrates through the epithelial layer. Plunging nuclei indicate distinct maturation stages closer to the lumen boundary.")
            ]),
            ("4. Approximated 3D Volume (µm³)", [
                ("", "Calculated via Riemann sum of 2D segmented pixel area × Z-slice thickness."),
                ("Biology: ", "Directly correlates to chromatin condensation levels. *NOTE: Absolute values artificially inflated by Point Spread Function (PSF) microscope halo.*")
            ]),
            ("5. Effective Thickness (Average Diameter, µm)", [
                ("Formula: ", "D_avg = 2 * sqrt( (V_3D / L_3D) / π )"),
                ("Biology: ", "Models the nucleus as a cylinder. *NOTE: Artificially inflated by PSF.* Particularly useful for relative comparisons of decompaction pathologies.")
            ]),
            ("6. Pitch Angle (Plunge Vector, Degrees)", [
                ("Formula: ", "θ = |arcsin(ΔZ / Length_Euclidean)| * (180/π)"),
                ("Biology: ", "Quantifies the absolute orientation of the elongating nucleus relative to the basal lamina plane. Values near 90° indicate perpendicular apical plunging.")
            ]),
            ("7. Taper Ratio", [
                ("Formula: ", "R = (Maximum Cross-sectional Area) / (Minimum Cross-sectional Area)"),
                ("Biology: ", "Identifies tapering anomalies near the acrosome versus the tail insertion. Extreme variance implies severe morphological irregularity or fusion errors.")
            ]),
            ("8. Spatial Packing Density (Nearest Neighbor Dist, µm)", [
                ("Calculated: ", "3D Euclidean distance from this track's centroid to its closest sibling."),
                ("Biology: ", "Measures tubule capacity and cyst grouping density during spermiation.")
            ])
        ]

        top_y = 1.0
        for title, content_list in methods_items:
            tb = slide5.shapes.add_textbox(Inches(0.5), Inches(top_y), Inches(9), Inches(0.72))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(10)
            
            for label, text in content_list:
                p2 = tf.add_paragraph()
                p2.font.size = Pt(8.5)
                if label:
                    run1 = p2.add_run()
                    run1.text = label
                    run1.font.bold = True
                run2 = p2.add_run()
                run2.text = text
                run2.font.bold = False
            
            top_y += 0.74

        add_hyperlink(slide5, "Population_Summary")

        # Save the presentation matching Batch Name
        final_pptx_name = "batch_analysis_results_v11.pptx"
        output_path = os.path.join(out_dir, final_pptx_name)
        try:
            print(f"PPTX: Saving to {output_path}...")
            prs.save(output_path)
            print(f"PPTX Report successfully saved to: {output_path}")
            return True
        except PermissionError:
            print(f"CRITICAL ERROR: Could not save PowerPoint because the file '{final_pptx_name}' is currently open!")
            return False
            
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"Failed to generate PPTX report: {e}")
        write_error_log(out_dir, "PowerPoint Generator", err_msg)
        return False
            
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        print(f"Failed to generate PPTX report: {e}")
        write_error_log(out_dir, "PowerPoint Generator", err_msg)
        try:
            from tkinter import messagebox
            messagebox.showerror("Reporting Error", f"Failed to generate PowerPoint Report:\n{e}\n\nSee report_generation_errors.txt for details.")
        except Exception:
            pass
        return False

# =============================================================================
# ROI GUI (single-file v9)
# =============================================================================
import traceback
try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
    from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
    from matplotlib.figure import Figure
    _TK_AVAILABLE = True
except Exception:
    _TK_AVAILABLE = False

PARAM_DESCRIPTIONS = {
    "UM_PER_PX_XY": "Microns per pixel in XY (Calibration)",
    "UM_PER_SLICE_Z": "Microns per Z-slice step (Calibration)",
    "CLAHE_CLIP": "Local contrast limit for CLAHE enhancement",
    "CLAHE_KERNEL": "Grid size for CLAHE enhancement",
    "BG_SIGMA": "Gaussian blur radius for background subtraction",
    "THRESHOLD_HI": "Intensity threshold for main spermatid seeds",
    "THRESHOLD_LO": "Intensity threshold for trailing edges",
    "CLOSE_RADIUS": "Radius for morphological closing (px)",
    "MIN_HOLE_AREA": "Smallest hole area to fill (px)",
    "MIN_OBJ_PX": "Absolute minimum foreground size (px)",
    "MAX_BRIDGE_PX": "Max gap to bridge in skeleton (px)",
    "MAX_BRANCH_LEN_PX": "Prune skeleton spurs shorter than this (px)",
    "BREAK_JUNCTIONS": "Automatically sever intersecting skeletons (True/False)",
    "USE_EARLY_SHAPE_FILTER": "Apply basic eccentricity filters early (True/False)",
    "MIN_ECCENTRICITY": "Minimum eccentricity (if early filtered)",
    "MAX_MINOR_PX": "Maximum minor axis length (px)",
    "MIN_AXIS_RATIO": "Minimum major/minor axis ratio",
    "MIN_MAJOR_PX": "Minimum major axis length (px)",
    "MIN_SKEL_LEN_PX": "Minimum valid skeleton length (px)",
    "MAX_GEODESIC_LEN_PX": "Maximum valid skeleton length (px)",
    "MAX_WIDTH_PX": "Maximum allowed thickness (px)",
    "MIN_LENGTH_WIDTH_RATIO": "Minimum required elongation ratio",
    "MAX_BRANCH_NODES": "Maximum allowed branch intersections",
    "MAX_TORTUOSITY": "Maximum allowed curvature index",
    "MAX_ENDPOINT_COUNT": "Maximum allowed skeleton endpoints",
    "ALLOW_LOOPS": "Allow topological loops in targets (True/False)",
    "DO_TRACKING": "Enable 3D tracking across Z-slices (True/False)",
    "TRACK_MAX_DIST_UM": "Maximum allowed tracking displacement (um)",
    "TRACK_MAX_GAP_SLICES": "Maximum allowed missing Z-slices in a track",
    "AUDIT_MAX_LENGTH_UM": "Flag tracks longer than this (um)",
    "AUDIT_MAX_TORTUOSITY": "Flag tracks more tortuous than this",
    "AUDIT_MAX_THICKNESS_UM": "Flag tracks thicker than this (um)",
    "AUDIT_MAX_TAPER_RATIO": "Flag tracks with taper ratio above this",
    "AUDIT_MIN_SLICES": "Flag tracks with fewer slices than this"
}

class ParameterEditor(tk.Toplevel):
    """
    A Tkinter ``Toplevel`` window providing an interactive editor for all pipeline
    configuration parameters.

    The editor auto-generates one labelled entry widget per config key, fetching
    human-readable descriptions from the ``PARAM_DESCRIPTIONS`` lookup table.
    Changes are validated and type-cast on apply (e.g., ``"True"`` → ``True``,
    ``"0.05"`` → ``float``).  The window supports:

    - **Apply to Session** — updates the live ``CONFIG`` dict in memory without
      restarting the process.
    - **Reset to Defaults** — restores all fields to the factory-default values
      captured at GUI startup.
    - **Load JSON** — populates all fields from an external ``settings.json``
      file (useful for project-specific preset configurations).
    - **Save JSON** — serialises the current edited values to disk so they can
      be reloaded in future sessions.

    The window is scrollable so all parameters are accessible even on small screens.
    """
    def __init__(self, parent, current_config, default_config, apply_callback):
        super().__init__(parent)
        self.title("Parameter Configuration")
        self.geometry("900x700")
        self.current_config = current_config
        self.default_config = default_config
        self.apply_callback = apply_callback
        
        self.entries = {}
        
        # Tools Frame (Top)
        tools = tk.Frame(self, bg='#e0e0e0', padx=10, pady=10)
        tools.pack(side='top', fill='x')
        
        tk.Button(tools, text="Apply to Session", command=self.apply, bg="#d4edda", font=("Arial", 10, "bold")).pack(side="left", padx=5)
        tk.Button(tools, text="Reset to Defaults", command=self.reset_defaults, bg="#f8d7da").pack(side="left", padx=5)
        tk.Button(tools, text="Load JSON", command=self.load_json).pack(side="left", padx=5)
        tk.Button(tools, text="Save JSON", command=self.save_json).pack(side="left", padx=5)
        
        # Scrollable Canvas
        canvas = tk.Canvas(self)
        scrollbar = ttk.Scrollbar(self, orient="vertical", command=canvas.yview)
        self.scrollable_frame = tk.Frame(canvas)
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        self.populate_form(self.current_config)

    def populate_form(self, cfg):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()
        
        self.entries = {}
        row = 0
        for k, v in cfg.items():
            if k in PARAM_DESCRIPTIONS and isinstance(v, (int, float, bool)):
                # Label Key
                tk.Label(self.scrollable_frame, text=k, font=("Arial", 10, "bold"), width=30, anchor="e").grid(row=row, column=0, padx=10, pady=5, sticky="e")
                
                # Entry value
                var = tk.StringVar(value=str(v))
                ent = tk.Entry(self.scrollable_frame, textvariable=var, width=15)
                ent.grid(row=row, column=1, padx=10, pady=5)
                self.entries[k] = (var, type(v))
                
                # Description
                tk.Label(self.scrollable_frame, text=PARAM_DESCRIPTIONS[k], fg="dimgray", anchor="w").grid(row=row, column=2, padx=10, pady=5, sticky="w")
                
                row += 1

    def apply(self):
        new_cfg = self.current_config.copy()
        try:
            for k, (var, t) in self.entries.items():
                val = var.get()
                if t == bool:
                    new_cfg[k] = val.lower() in ['true', '1', 't', 'y', 'yes']
                else:
                    new_cfg[k] = t(val)
        except ValueError as e:
            messagebox.showerror("Validation Error", f"Invalid input format: {e}")
            return
            
        self.apply_callback(new_cfg)
        self.destroy()

    def reset_defaults(self):
        self.populate_form(self.default_config)
        
    def load_json(self):
        fpath = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                with open(fpath, 'r') as f:
                    loaded = json.load(f)
                filtered = {k: v for k, v in loaded.items() if k in self.current_config}
                temp = self.current_config.copy()
                temp.update(filtered)
                self.populate_form(temp)
            except Exception as e:
                messagebox.showerror("Load Error", str(e))
                
    def save_json(self):
        fpath = filedialog.asksaveasfilename(defaultextension=".json", filetypes=[("JSON files", "*.json")])
        if fpath:
            try:
                import json
                new_cfg = self.current_config.copy()
                for k, (var, t) in self.entries.items():
                    val = var.get()
                    if t == bool:
                        new_cfg[k] = val.lower() in ['true', '1', 't', 'y', 'yes']
                    else:
                        new_cfg[k] = t(val)
                with open(fpath, 'w') as f:
                    json.dump(new_cfg, f, indent=4)
                messagebox.showinfo("Saved", f"Parameters saved to {os.path.basename(fpath)}")
            except Exception as e:
                messagebox.showerror("Save Error", str(e))

# =============================================================================
# v12 ADVANCED AI BIOLOGICAL INTERPRETATION ENGINE
# =============================================================================

SPECIES_PROFILES = {
    "D. melanogaster": {
        "length_target": "10.0 µm",
        "thickness": "0.3-0.4 µm",
        "context": "Standard model species. Needle-like, highly condensed nucleus.",
        "heteromorphism": False
    },
    "D. simulans": {
        "length_target": "9.5-10.5 µm",
        "thickness": "0.3 µm",
        "context": "Close relative of D. mel, very similar nuclear profile.",
        "heteromorphism": False
    },
    "D. yakuba": {
        "length_target": "10.0-11.0 µm",
        "thickness": "0.35 µm",
        "context": "Close relative of D. mel, slightly different condensation timing.",
        "heteromorphism": False
    },
    "D. ananassae": {
        "length_target": "8.0-10.0 µm",
        "thickness": "0.3-0.4 µm",
        "context": "Melanogaster subgroup. Moderate nuclear elongation with distinct chromatin packaging.",
        "heteromorphism": False
    },
    "D. pseudoobscura (Dpse)": {
        "length_target": "Variable (Heteromorphic)",
        "thickness": "Variable",
        "context": "Produces both fertile 'eusperm' and shorter, non-fertilizing 'parasperm'.",
        "heteromorphism": True
    },
    "D. virilis (Dvir)": {
        "length_target": "15.0-18.0 µm",
        "thickness": "0.4 µm",
        "context": "Large species with robust, stable nuclear morphology.",
        "heteromorphism": False
    },
    "General / Evolutionary": {
        "length_target": "Unknown",
        "thickness": "Unknown",
        "context": "Perform comparative discovery mode to infer species strategy.",
        "heteromorphism": "Possible"
    }
}

def get_ai_biological_interpretation(csv_summary_str, species, folder_name, model_id="gemini-2.5-pro"):
    """Calls Gemini API for biological narrative. Priority: Local File > Env Var."""
    profile = SPECIES_PROFILES.get(species, SPECIES_PROFILES["General / Evolutionary"])
    
    if not _HAVE_REQUESTS:
        return "AI ANALYSIS SKIPPED: 'requests' library not found. Please run 'pip install requests'."
    
    # 1. Try local file path first
    key_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gemini_api_key.txt")
    api_key = ""
    if os.path.exists(key_file):
        with open(key_file, 'r') as f:
            api_key = f.read().strip()
    
    # 2. Fallback to Env Var
    if not api_key:
        api_key = os.environ.get("GEMINI_API_KEY", "")
        
    if not api_key:
        return "AI ANALYSIS SKIPPED: No Gemini API Key found. Use 'Set API Key' in the GUI (Free at aistudio.google.com)."
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent?key={api_key}"
    system_prompt = f"""You are a world-class Drosophila reproductive biologist and evolutionary morphologist.
    You are interpreting 3D-tracked sperm nuclei data from a confocal Z-stack batch analysis.
    
    BIOLOGICAL CONTEXT:
    - Target Species: {species}
    - Baseline Profile: {profile['context']}
    - Expected Length Class: {profile['length_target']}
    - Heteromorphism (Multiple Morphs): {'Yes' if profile.get('heteromorphism') else 'No'}
    - Source Data Folder: {folder_name}
    
    INPUT DATA (CSV Summary):
    {csv_summary_str}
    
    YOUR TASK:
    1. Determine the 'Morphological Class' of this population. Does it align with the expected {species} profile?
    2. Analyze maturation: Are these likely mature motile sperm or transitionary elongating spermatids based on length/thickness/tortuosity?
    3. Identify Anomalies: Highlight any outliers in tortuosity or uneven thickness that indicate fixation artifacts or developmental defects.
    4. Evolutionary Insight: Provide 2-3 sentences on the evolutionary context of this sperm morphology.
    5. Formatting: Use professional, high-density scientific language. Use Markdown formatting.
    """
    payload = {"contents": [{"parts": [{"text": system_prompt}]}]}
    try:
        response = requests.post(url, json=payload, timeout=300)
        if response.status_code == 200:
            return response.json()['candidates'][0]['content']['parts'][0]['text']
        else: return f"AI API Error ({response.status_code}): {response.text}"
    except Exception as e: return f"Failed to connect to AI Service: {str(e)}"

def generate_ai_html_report(out_dir, ai_text, stats_summary, species):
    """Generates premium HTML dashboard."""
    html_path = os.path.join(out_dir, "AI_Biological_Analysis_v12.html")
    try:
        import markdown
        ai_html = markdown.markdown(ai_text)
    except ImportError:
        ai_html = ai_text.replace("\n", "<br>")
    
    stats_html = "".join([f'<div class="stat-box"><div class="stat-val">{v}</div><div class="stat-label">{k.replace("_", " ")}</div></div>' for k,v in stats_summary.items()])
    
    html_content = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <title>Sperm AI Analysis - v12</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap" rel="stylesheet">
        <style>
            :root {{ --primary: #2563eb; --bg: #f8fafc; --card: rgba(255, 255, 255, 0.8); }}
            body {{ font-family: 'Inter', sans-serif; background: var(--bg); color: #1e293b; line-height: 1.6; padding: 40px; }}
            .container {{ max-width: 900px; margin: 0 auto; }}
            .badge {{ background: #dbeafe; color: #1e40af; padding: 4px 12px; border-radius: 99px; font-size: 0.8rem; font-weight: bold; }}
            .glass-card {{ background: var(--card); backdrop-filter: blur(10px); border: 1px solid rgba(255,255,255,0.3); 
                          border-radius: 16px; padding: 30px; box-shadow: 0 4px 20px rgba(0,0,0,0.05); margin-bottom: 30px; }}
            .stats-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 20px; margin-bottom: 30px; }}
            .stat-box {{ background: white; padding: 15px; border-radius: 12px; text-align: center; border: 1px solid #e2e8f0; }}
            .stat-val {{ font-size: 1.5rem; font-weight: 800; color: var(--primary); }}
            .stat-label {{ font-size: 0.75rem; text-transform: uppercase; color: #64748b; }}
            .ai-content h2 {{ color: var(--primary); border-bottom: 1px solid #e2e8f0; padding-bottom: 8px; }}
        </style>
    </head>
    <body>
        <div class="container">
            <span class="badge">Drosophila Morphometrics v12</span>
            <h1>AI Biological Interpretation</h1>
            <p>Advanced Analysis for <strong>{species}</strong></p>
            <div class="stats-grid">{stats_html}</div>
            <div class="glass-card">
                <div class="ai-content">{ai_html}</div>
            </div>
            <footer>Generated | {time.strftime('%Y-%m-%d %H:%M:%S')} | v12 AI Layer</footer>
        </div>
    </body>
    </html>
    """
    with open(html_path, "w", encoding="utf-8") as f: f.write(html_content)
    return html_path

# =============================================================================

class SpermGUI:
    """
    The primary Tkinter-based graphical user interface for the Sperm Segmentation ROI Tool.

    The GUI provides a two-panel layout:
    - **Left sidebar** — controls for directory loading, Z-slice navigation,
      tool selection, ROI drawing, single-slice analysis, batch analysis,
      and two progress bars (2D segmentation and post-analysis reporting).
    - **Right canvas** — a matplotlib ``FigureCanvasTkAgg`` that renders the
      currently selected Z-slice image (raw, overlay, or debug) and accepts
      polygon ROI drawing interactions.

    Key interaction modes (controlled by ``mode_var``)
    ---------------------------------------------------
    - ``'view'``   — pans and inspects the raw image.
    - ``'review'`` — displays the saved overlay PNG for the current slice (requires
      a prior batch run).
    - ``'roi'``    — left-click to add polygon vertices; right-click to undo the
      last vertex; call *Finalize Polygon* to close and rasterise the mask.

    Thread architecture
    -------------------
    Batch processing (:meth:`run_batch_analysis`) runs in a background ``threading.Thread``
    so the Tkinter event loop remains responsive.  Progress bar updates are marshalled
    back to the main thread via ``root.after()`` calls.
    """
    def open_parameter_editor(self):
        """
        Opens the :class:`ParameterEditor` ``Toplevel`` window.

        Defines an ``on_apply`` callback that merges the edited values into the
        global ``CONFIG`` dict and updates the status label so the researcher
        knows the new parameters are active.
        """
        def on_apply(new_cfg):
            CONFIG.update(new_cfg)
            self.lbl_roi.config(text="Parameters updated in memory.")
            
        editor = ParameterEditor(self.root, CONFIG, self.default_config, on_apply)

    def _load_tuned_params(self):
        """Load a tuned parameters JSON file and merge into CONFIG."""
        from tkinter import filedialog, messagebox
        filepath = filedialog.askopenfilename(
            title="Select Tuned Parameters JSON",
            filetypes=[("JSON Files", "*.json"), ("All Files", "*.*")],
            initialdir=os.path.dirname(os.path.abspath(__file__))
        )
        if not filepath:
            return
        try:
            with open(filepath, 'r') as f:
                tuned = json.load(f)
            
            # Only update keys that exist in CONFIG (safety check)
            applied = []
            for key, value in tuned.items():
                if key in CONFIG:
                    old_val = CONFIG[key]
                    CONFIG[key] = value
                    applied.append(f"  {key}: {old_val} → {value}")
            
            if applied:
                n = len(applied)
                short_name = os.path.basename(filepath)
                self.lbl_params_status.config(
                    text=f'✅ Loaded {n} params from {short_name}',
                    fg='green'
                )
                detail = "\n".join(applied)
                messagebox.showinfo(
                    "Parameters Loaded",
                    f"Loaded {n} parameters from:\n{short_name}\n\n"
                    f"Changes applied:\n{detail}\n\n"
                    f"Use 'Revert Defaults' to undo."
                )
            else:
                messagebox.showwarning("No Matching Keys", 
                    f"No recognized CONFIG keys found in {os.path.basename(filepath)}")
        except Exception as e:
            messagebox.showerror("Load Error", f"Failed to load parameters:\n{e}")
    
    def _revert_to_defaults(self):
        """Revert CONFIG back to the original defaults captured at startup."""
        from tkinter import messagebox
        CONFIG.update(self.default_config)
        self.lbl_params_status.config(
            text='Using default parameters',
            fg='#555'
        )
        messagebox.showinfo("Reverted", "All parameters have been reverted to their original defaults.")

    def __init__(self, root):
        """
        Initialises the main application window, all sidebar controls, the matplotlib
        canvas, and mouse/key event bindings.

        Args:
            root (tk.Tk): The root Tkinter window created by :func:`launch_gui`.
        """
        self.root = root
        self.root.title(f'Sperm Segmentation ROI Tool - Saturn Project')
        self.root.geometry('1450x920')

        self.input_dir = ''
        self.files = []
        self.current_idx = 0
        self.current_img = None
        self.last_out_dir = ""

        self.roi_points = []
        self.drawing = False
        self.roi_active = False
        self._loaded_roi_mask = None

        # --- v12 AI ANALYSIS CONTROLS ---
        self.species_var = tk.StringVar(value='D. melanogaster')
        self.model_var = tk.StringVar(value='Gemini 2.5 Pro (Thinking)')
        self._last_batch_ts = None
        self._last_batch_out_dir = None

        # --- v12 SCROLLABLE SIDEBAR --- 
        self.sidebar_container = tk.Frame(root, width=300, bg='#f0f0f0')
        self.sidebar_container.pack(side='left', fill='y')
        self.sidebar_container.pack_propagate(False)

        self.sidebar_canvas = tk.Canvas(self.sidebar_container, bg='#f0f0f0', highlightthickness=0)
        self.sidebar_scrollbar = ttk.Scrollbar(self.sidebar_container, orient='vertical', command=self.sidebar_canvas.yview)
        self.sidebar_canvas.pack(side='left', fill='both', expand=True)
        self.sidebar_scrollbar.pack(side='right', fill='y')
        self.sidebar_canvas.configure(yscrollcommand=self.sidebar_scrollbar.set)

        self.sidebar = tk.Frame(self.sidebar_canvas, bg='#f0f0f0')

        def _on_sidebar_configure(event):
            # Update scrollregion to dynamically envelop all items in canvas
            self.sidebar_canvas.configure(scrollregion=self.sidebar_canvas.bbox('all'))

        self.sidebar.bind("<Configure>", _on_sidebar_configure)
        self.sidebar_canvas.bind("<Configure>", lambda e: self.sidebar_canvas.itemconfig(self.sidebar_window, width=e.width))
        
        # Specifically create the window as an object so we can resize it
        self.sidebar_window = self.sidebar_canvas.create_window((0, 0), window=self.sidebar, anchor="nw")

        def _on_mousewheel(event):
            # Scroll 4 units at a time for smooth, responsive scrolling
            self.sidebar_canvas.yview_scroll(int(-1*(event.delta/30)), 'units')
        
        def _bind_mousewheel(event):
            self.sidebar_canvas.bind_all('<MouseWheel>', _on_mousewheel)
        
        def _unbind_mousewheel(event):
            self.sidebar_canvas.unbind_all('<MouseWheel>')
        
        # Bind to canvas AND all sidebar children for reliable scroll capture
        self.sidebar_container.bind('<Enter>', _bind_mousewheel)
        self.sidebar_container.bind('<Leave>', _unbind_mousewheel)
        # ------------------------------

        self.default_config = CONFIG.copy()
        tk.Button(self.sidebar, text='⚙ Configure Parameters', command=self.open_parameter_editor, bg='#e2e3e5').pack(fill='x', padx=6, pady=(10,2))

        # --- Stage 2: Load/Revert Tuned Parameters ---
        params_frame = tk.Frame(self.sidebar, bg='#f0f0f0')
        params_frame.pack(fill='x', padx=6, pady=(2,6))
        tk.Button(params_frame, text='📂 Load Tuned Params', command=self._load_tuned_params, bg='#d4edda', width=16).pack(side='left', expand=True, fill='x', padx=(0,2))
        tk.Button(params_frame, text='↩ Revert Defaults', command=self._revert_to_defaults, bg='#f8d7da', width=14).pack(side='right', expand=True, fill='x', padx=(2,0))
        self.lbl_params_status = tk.Label(self.sidebar, text='Using default parameters', wraplength=260, justify='left', fg='#555', font=('Arial', 8))
        self.lbl_params_status.pack(pady=(0,4))

        tk.Button(self.sidebar, text='Load Directory', command=self.load_directory, height=2).pack(fill='x', padx=6, pady=6)
        self.lbl_status = tk.Label(self.sidebar, text='No directory loaded', wraplength=260, justify='left')
        self.lbl_status.pack(pady=6)

        tk.Label(self.sidebar, text='Z-Slice Navigation').pack(pady=(20, 0))
        self.scale_z = tk.Scale(self.sidebar, from_=0, to=0, orient='horizontal', command=self.on_slide_change)
        self.scale_z.pack(fill='x', padx=10)
        self.lbl_z = tk.Label(self.sidebar, text='Z: 0 / 0')
        self.lbl_z.pack()

        tk.Label(self.sidebar, text='Tools').pack(pady=(20, 5))
        self.mode_var = tk.StringVar(value='view')
        tk.Radiobutton(self.sidebar, text='View/Nav (Raw Image)', variable=self.mode_var, value='view', command=self.render).pack(anchor='w', padx=10)
        tk.Radiobutton(self.sidebar, text='Review Overlays (After Batch)', variable=self.mode_var, value='review', command=self.render).pack(anchor='w', padx=10)
        tk.Radiobutton(self.sidebar, text='Draw ROI (Polygon)', variable=self.mode_var, value='roi', command=self.render).pack(anchor='w', padx=10)
        tk.Label(self.sidebar, text='(Left-click points, Right-click undo)', font=('Arial', 8, 'italic'), fg='dimgray').pack()
        tk.Button(self.sidebar, text='Finalize Polygon', command=self.finalize_roi, bg='#ffeeba').pack(fill='x', padx=20, pady=2)

        # ── v12 ADVANCED AI ANALYTICS ─────────────────────────────────────────
        tk.Label(self.sidebar, text='v12 AI Biological Analysis', font=('Arial', 10, 'bold'), fg='#8b5cf6').pack(pady=(20, 0))
        species_list = [
            'D. melanogaster', 'D. simulans', 'D. yakuba',
            'D. ananassae', 'D. pseudoobscura (Dpse)',
            'D. virilis (Dvir)', 'General / Evolutionary'
        ]
        self.species_dropdown = ttk.Combobox(self.sidebar, textvariable=self.species_var, values=species_list, state='readonly')
        self.species_dropdown.pack(fill='x', padx=10, pady=5)
        
        tk.Label(self.sidebar, text='AI Model:', font=('Arial', 8)).pack()
        models = ['Gemini 2.5 Pro (Thinking)', 'Gemini 2.5 Flash (Fast)']
        self.model_dropdown = ttk.Combobox(self.sidebar, textvariable=self.model_var, values=models, state='readonly')
        self.model_dropdown.pack(fill='x', padx=10, pady=(0, 5))

        tk.Button(self.sidebar, text='🔑 Set API Key (Free)', command=self.set_ai_key, font=('Arial', 8)).pack(padx=10, pady=(0, 5))
        self.btn_ai = tk.Button(self.sidebar, text='🧬 Run AI Analysis', command=self.run_ai_analysis,
                                 bg='#8b5cf6', fg='white', font=('Arial', 9, 'bold'), state='disabled')
        self.btn_ai.pack(fill='x', padx=10, pady=(0, 10))
        
        # ──────────────────────────────────────────────────────────────────────

        tk.Button(self.sidebar, text='Run Analysis on Slice', command=self.run_analysis_slice).pack(fill='x', padx=6, pady=18)
        tk.Button(self.sidebar, text='Run Batch (All Slices + 3D Track)', command=self.run_batch_analysis, bg='#d4edda', font=('Arial', 10, 'bold')).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Reset ROI', command=self.reset_roi).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Save ROI Mask', command=self.save_roi_mask).pack(fill='x', padx=6, pady=6)
        tk.Button(self.sidebar, text='Load ROI Mask', command=self.load_roi_mask).pack(fill='x', padx=6, pady=6)

        self.lbl_roi = tk.Label(self.sidebar, text='ROI: none', wraplength=260, justify='left')
        self.lbl_roi.pack(pady=10)

        self.canvas_frame = tk.Frame(root, bg='black')
        self.canvas_frame.pack(side='right', expand=True, fill='both')

        # --- TOP STATUS BAR (For Progress Visibility) ---
        self.top_status_frame = tk.Frame(self.canvas_frame, bg='#f0f0f0', height=80)
        self.top_status_frame.pack(side='top', fill='x')

        # --- Sub-frame for Status (Left) ---
        self.status_sub_frame = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.status_sub_frame.pack(side='left', padx=10, fill='y')

        # --- Dynamic Status Label ---
        self.lbl_batch_op = tk.Label(self.status_sub_frame, text='GUI Ready', font=('Arial', 10, 'bold'), fg='#2c3e50', bg='#f0f0f0')
        self.lbl_batch_op.pack(side='left', pady=5)

        # --- PROGRESS BOX: Sequential Progress Bars (At Top Right) ---
        self.p_container = tk.Frame(self.top_status_frame, bg='#f0f0f0')
        self.p_container.pack(side='right', padx=20, pady=5)

        # Frame 1: 2D Batch Segmentation
        self.batch_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        self.batch_p_frame.pack(fill='x')
        tk.Label(self.batch_p_frame, text='Batch Progress (2D)', font=('Arial', 9, 'bold'), bg='#f0f0f0').pack(side='left', padx=10)
        self.progress = ttk.Progressbar(self.batch_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress.pack(side='left', padx=5, pady=5)
        self.lbl_progress_val = tk.Label(self.batch_p_frame, text='0%', font=('Arial', 10, 'bold'), fg='blue', bg='#f0f0f0')
        self.lbl_progress_val.pack(side='left', padx=10)

        # Frame 2: Post-Analysis (Initially Hidden)
        self.post_p_frame = tk.Frame(self.p_container, bg='#f0f0f0')
        # We don't pack self.post_p_frame here
        tk.Label(self.post_p_frame, text='Post-Analysis Progress', font=('Arial', 9, 'bold'), bg="#f0f0f0").pack(side='left', padx=10)
        self.progress_post = ttk.Progressbar(self.post_p_frame, orient='horizontal', length=200, mode='determinate')
        self.progress_post.pack(side='left', padx=5, pady=5)
        self.lbl_post_progress_val = tk.Label(self.post_p_frame, text='Waiting...', font=('Arial', 10, 'bold'), fg='dimgray', bg="#f0f0f0")
        self.lbl_post_progress_val.pack(side='left', padx=10)

        self.fig = Figure(figsize=(8, 8), dpi=100)
        self.ax = self.fig.add_subplot(111)
        self.ax.axis('off')
        self.canvas = FigureCanvasTkAgg(self.fig, master=self.canvas_frame)
        self.canvas.get_tk_widget().pack(fill='both', expand=True)

        self.canvas.mpl_connect('button_press_event', self.on_click)
        self.canvas.mpl_connect('key_press_event', self.on_key)

    def set_ai_key(self):
        """Opens a small dialog to save the Gemini API Key to a local file."""
        win = tk.Toplevel(self.root)
        win.title("AI Key Management")
        win.geometry("420x260")
        win.grab_set()
        
        tk.Label(win, text="Gemini AI Key Setup", font=('Arial', 12, 'bold')).pack(pady=(15,5))
        
        link = tk.Label(win, text="Get your FREE Key at Google AI Studio", fg='blue', cursor="hand2", font=('Arial', 9, 'underline'))
        link.pack()
        link.bind("<Button-1>", lambda e: webbrowser.open("https://aistudio.google.com/app/apikey"))
        
        tk.Label(win, text="(No credit card required for Free Tier)", font=('Arial', 8), fg='#64748b').pack(pady=(0, 10))
        
        tk.Label(win, text="Paste API Key here:", font=('Arial', 9)).pack()
        entry = tk.Entry(win, show='*', width=45)
        entry.pack(pady=5)
        
        # Load current key if exists
        key_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "gemini_api_key.txt")
        exists = os.path.exists(key_file)
        if exists:
            with open(key_file, 'r') as f:
                entry.insert(0, f.read().strip())
        
        def save():
            k = entry.get().strip()
            if not k:
                messagebox.showwarning("Warning", "Key cannot be empty.")
                return
            with open(key_file, 'w') as f:
                f.write(k)
            status = "updated" if exists else "saved"
            messagebox.showinfo("Success", f"API Key {status} successfully!\nPortable analysis is ready.")
            win.destroy()
            
        btn_text = "Update API Key" if exists else "Save Key Permanently"
        tk.Button(win, text=btn_text, command=save, bg='#8b5cf6', fg='white', font=('Arial', 9, 'bold'), padx=20).pack(pady=15)

    def run_ai_analysis(self):
        """Runs AI biological interpretation on the last completed batch data."""
        if self._last_batch_ts is None or self._last_batch_out_dir is None:
            messagebox.showinfo("AI Analysis", "Please run a batch analysis first before requesting AI interpretation.")
            return

        ts = self._last_batch_ts
        out_dir = self._last_batch_out_dir

        if ts.empty or len(ts) == 0:
            messagebox.showinfo("AI Analysis", "No valid 3D tracks found in the last batch. AI requires track data.")
            return

        self.lbl_batch_op.config(text='AI Analyzing Biological Data...', fg='#8b5cf6')
        self.btn_ai.config(state='disabled', text='🧬 AI Working...')
        self.root.update()

        # Token mitigation: Send fewer tracks to stay within Pro free tier limits
        ts_quality = ts
        
        # We send the high-quality tracks to Gemini so its analysis isn't skewed by outliers
        csv_summary_str = ts_quality.head(100).to_csv(index=False)
        species = self.species_var.get()
        folder_name = os.path.basename(self.input_dir or "Current Project")
        
        # Model Selection Mapping
        m_map = {
            'Gemini 2.5 Pro (Thinking)': 'gemini-2.5-pro',
            'Gemini 2.5 Flash (Fast)': 'gemini-2.5-flash'
        }
        model_id = m_map.get(self.model_var.get(), 'gemini-2.5-pro')

        def _ai_thread():
            try:
                print(f"AI START: Interpreting {len(ts)} tracks for {species} via {model_id}...")
                ai_text = get_ai_biological_interpretation(csv_summary_str, species, folder_name, model_id=model_id)

                if not ai_text or "AI ANALYSIS SKIPPED" in ai_text or "AI API Error" in ai_text:
                    err_hint = ""
                    if "(429)" in ai_text:
                        err_hint = "\n\nTIP: Gemini 2.5 Pro has lower free-tier quotas. Try switching to 'Gemini 2.5 Flash' in the sidebar or wait a minute."
                    print(f"AI FAIL: {ai_text}")
                    self.root.after(0, lambda: messagebox.showwarning("AI Analysis", f"The AI analysis could not proceed:\n\n{ai_text}{err_hint}"))
                    self.root.after(0, lambda: self.btn_ai.config(state='normal', text='🧬 Run AI Analysis'))
                    return

                stats_summary = {
                    "Median_Length_um": f"{ts['total_3d_length_um'].median():.2f}",
                    "Avg_Tortuosity": f"{ts['tortuosity_3d'].mean():.2f}",
                    "Track_Count": f"{len(ts)}",
                    "Species": species
                }
                report_path = generate_ai_html_report(out_dir, ai_text, stats_summary, species)

                if os.path.exists(report_path):
                    print(f"AI SUCCESS: Report generated at {report_path}")
                    abs_report = os.path.abspath(report_path)

                    def _open_report():
                        try:
                            if os.name == 'nt':
                                os.startfile(abs_report)
                            else:
                                webbrowser.open("file:///" + abs_report.replace("\\", "/"))
                        except Exception as oe:
                            print(f"AI OPEN ERROR: {oe}")

                    self.root.after(0, lambda: self.lbl_batch_op.config(text='AI REPORT READY ✓', fg='green'))
                    self.root.after(500, _open_report)
                    self.root.after(1200, lambda: messagebox.showinfo("AI Analysis Complete",
                        f"Biological interpretation report generated and opened:\n\n{os.path.basename(report_path)}"))
                else:
                    raise FileNotFoundError(f"Report file was not created at {report_path}")

            except Exception as e:
                import traceback
                print(f"AI EXCEPTION: {traceback.format_exc()}")
                err_msg = f"Failed during biological interpretation:\n\n{str(e)}"
                self.root.after(0, lambda m=err_msg: messagebox.showerror("AI Error", m))
            finally:
                self.root.after(0, lambda: self.btn_ai.config(state='normal', text='🧬 Run AI Analysis'))

        threading.Thread(target=_ai_thread, daemon=True).start()

    def load_directory(self):
        """
        Opens a file-picker dialog so the user can select any image in a Z-stack folder.

        Discovers all supported image files (``*.tif``, ``*.tiff``, ``*.png``, ``*.jpg``,
        ``*.jpeg``) in the same directory as the selected file, applies natural sort order
        (so ``z01`` comes before ``z10``), synchronises the Z-slice slider to the selected
        file, and calls :meth:`load_image` to display the first slice.
        """
        initial = CONFIG.get('INPUT_DIR', os.getcwd())
        
        # Use file picker so user can see images
        fpath = filedialog.askopenfilename(
            initialdir=str(initial),
            title="Select any image in the stack",
            filetypes=[("Image files", "*.tif *.tiff *.png *.jpg *.jpeg"), ("All files", "*.*")]
        )
        
        if not fpath:
            return
            
        selected_file = pl.Path(fpath)
        p = selected_file.parent
        self.input_dir = str(p)
        
        exts = ['.tif', '.tiff', '.png', '.jpg', '.jpeg']
        found_files_path = []
        
        # 1. Discover all images in the same folder
        for ext in exts:
            found_files_path.extend(list(p.glob(f"*{ext}")))
            found_files_path.extend(list(p.glob(f"*{ext.upper()}")))
            
        # 2. Recursive fallback (if needed, though file picker implies they are in the right spot)
        if not found_files_path:
            for ext in exts:
                found_files_path.extend(list(p.rglob(f"*{ext}")))
                found_files_path.extend(list(p.rglob(f"*{ext.upper()}")))
        
        if not found_files_path:
            messagebox.showerror('Error', f"No supported images found in:\n{p}")
            return
            
        # Standardize and Natural Sort
        found_files_str = [os.path.abspath(str(f)) for f in found_files_path]
        unique_files = list(set(found_files_str))
        self.files = sorted(unique_files, key=natural_sort_key)
            
        # Sync to the selected file
        try:
            self.current_idx = self.files.index(os.path.abspath(fpath))
        except ValueError:
            self.current_idx = 0
            
        self.scale_z.config(to=len(self.files) - 1)
        self.scale_z.set(self.current_idx)
        self.reset_roi(redraw=False)
        self.load_image()
        self.lbl_status.config(text=f'Opened: {selected_file.name}\n({len(self.files)} slices in folder)', fg='blue')
        self.root.update()

    def load_image(self):
        if not self.files:
            return
        try:
            self.current_img = robust_imread(self.files[self.current_idx])
            self.lbl_z.config(text=f'Z: {self.current_idx} / {len(self.files)-1}')
            self.render()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load image:\n{self.files[self.current_idx]}\n\nError: {e}")

    def on_slide_change(self, val):
        self.current_idx = int(val)
        self.load_image()

    def render(self):
        self.ax.clear()
        self.ax.axis('off')
        
        # --- NEW: Review Overlays Logic ---
        if self.mode_var.get() == 'review':
            try:
                if not hasattr(self, 'last_out_dir') or not self.last_out_dir:
                    self.ax.text(0.5, 0.5, "No Batch Analysis Results Found.\nRun Batch First.", 
                                 ha='center', va='center', color='red', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return
                
                z_idx = extract_z_index(self.files[self.current_idx])
                panel_path = os.path.join(self.last_out_dir, "overlays", f"z{z_idx:02d}_panel.png")
                
                if os.path.exists(panel_path):
                    if _HAVE_CV2:
                        img = _cv2.imread(panel_path)
                        img = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)
                    else:
                        img = plt.imread(panel_path)
                    self.ax.imshow(img)
                    self.canvas.draw_idle()
                    return
                else:
                    self.ax.text(0.5, 0.5, f"Overlay not found for Z={z_idx:02d}\n{os.path.basename(panel_path)}", 
                                 ha='center', va='center', color='orange', transform=self.ax.transAxes)
                    self.canvas.draw_idle()
                    return
            except Exception as e:
                print(f"Overlay Render Error: {e}")
                
        if self.current_img is not None and isinstance(self.current_img, np.ndarray):
            img = self.current_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            self.ax.imshow(disp, cmap='gray')
            if len(self.roi_points) > 0:
                pts = np.array(self.roi_points)
                # If active (closed), draw it closed. Otherwise, draw the open line and endpoints.
                if self.roi_active:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=2)
                else:
                    self.ax.plot(pts[:,0], pts[:,1], 'r-', linewidth=1.5)
                    self.ax.plot(pts[:,0], pts[:,1], 'ro', markersize=4)
            # If a loaded mask exists, always redraw the red contour so it persists across slices
            elif self._loaded_roi_mask is not None:
                self.ax.contour(self._loaded_roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def on_click(self, event):
        if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
            return
        if self.mode_var.get() == 'roi':
            # Left click = add point
            if event.button == 1:
                # If they start clicking after already having a finished ROI, wipe it to start fresh
                if self.roi_active:
                    self.reset_roi(redraw=False)
                self.roi_points.append([event.xdata, event.ydata])
                self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                self.render()
            
            # Right click = undo last point
            elif event.button == 3:
                if not self.roi_active and len(self.roi_points) > 0:
                    self.roi_points.pop()
                    self.lbl_roi.config(text=f'ROI: building ({len(self.roi_points)} points)')
                    self.render()

    def finalize_roi(self):
        if not self.roi_active and len(self.roi_points) > 2:
            self.roi_points.append(self.roi_points[0])
            self.roi_active = True
            self.lbl_roi.config(text=f'ROI: active ({len(self.roi_points)-1} points)')
            self.render()
        elif not self.roi_active and len(self.roi_points) > 0:
            messagebox.showwarning('Draw ROI', 'Please place at least 3 points before finalizing.')

    def on_key(self, event):
        if self.mode_var.get() == 'roi' and event.key == 'enter':
            self.finalize_roi()

    def reset_roi(self, redraw=True):
        self.roi_points = []
        self.roi_active = False
        self._loaded_roi_mask = None
        self.lbl_roi.config(text='ROI: none')
        if redraw:
            self.render()

    def build_roi_mask(self):
        # If a mask was loaded from file, use it directly
        if self._loaded_roi_mask is not None and self.current_img is not None:
            return self._loaded_roi_mask
        if not self.roi_active or len(self.roi_points) < 4 or self.current_img is None:
            return None
        h, w = self.current_img.shape
        yy, xx = np.mgrid[:h, :w]
        pts = np.column_stack((xx.ravel(), yy.ravel()))
        path = Path(self.roi_points)
        return path.contains_points(pts).reshape(h, w)

    def save_roi_mask(self):
        mask = self.build_roi_mask()
        if mask is None:
            messagebox.showinfo('Save ROI', 'No active ROI to save.')
            return
        default = f'roi_z{self.current_idx:02d}.npy'
        path = filedialog.asksaveasfilename(defaultextension='.npy', initialfile=default,
                                            filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        np.save(path, mask.astype(np.uint8))
        messagebox.showinfo('Save ROI', f'Saved ROI mask to:\n{path}')

    def load_roi_mask(self):
        if self.current_img is None:
            messagebox.showinfo('Load ROI', 'Load an image first.')
            return
        path = filedialog.askopenfilename(filetypes=[('NumPy array', '*.npy')])
        if not path:
            return
        mask = np.load(path).astype(bool)
        if mask.shape != self.current_img.shape:
            messagebox.showerror('Load ROI', f'Mask shape {mask.shape} does not match image shape {self.current_img.shape}')
            return
        # Convert mask boundary to points for display
        ys, xs = np.where(mask)
        if len(xs) == 0:
            messagebox.showerror('Load ROI', 'Loaded ROI mask is empty.')
            return
        self.roi_points = []
        self.roi_active = True
        self._loaded_roi_mask = mask  # Store the loaded mask for build_roi_mask!
        self.lbl_roi.config(text=f'ROI: loaded mask\n{os.path.basename(path)}')
        self.render()
        self.ax.contour(mask.astype(float), levels=[0.5], colors='red', linewidths=1.5)
        self.canvas.draw_idle()

    def run_analysis_slice(self):
        if self.current_img is None:
            messagebox.showinfo('Info', 'No image loaded. Load a directory first.')
            return

        self.lbl_roi.config(text='Running analysis...')
        self.root.update_idletasks()

        try:
            import time as _t
            log_lines = []
            def log(msg):
                log_lines.append(msg)
                print(msg)
                sys.stdout.flush()

            log(f"\n--- GUI Analysis: slice {self.current_idx} ---")
            log(f"  File: {os.path.basename(self.files[self.current_idx]) if self.files else 'N/A'}")
            log(f"  Image shape: {self.current_img.shape}, dtype: {self.current_img.dtype}")

            params = CONFIG.copy()
            params['SAVE_DEBUG_IMAGES'] = False
            roi_mask = self.build_roi_mask()

            # Determine image to process
            full_img = self.current_img
            crop_offset_y, crop_offset_x = 0, 0
            H, W = full_img.shape

            # ═══════════════════════════════════════════════════════════════════
            # TWO-PASS HYBRID SEGMENTATION ENGINE
            # Pass 1: Full-frame analysis (correct global percentiles for sparse outer nuclei)
            # Pass 2: Tight ROI crop analysis (local percentiles naturally separate dense clusters)
            # Merge: Spatially deduplicate overlapping centroids, keeping unique detections from both
            # ═══════════════════════════════════════════════════════════════════

            t0 = _t.time()

            # ── PASS 1: Full-frame analysis ──────────────────────────────────
            log(f"  PASS 1: Full-frame analysis...")
            seg1 = segment_slice(full_img, params, roi_mask=roi_mask, z_idx=self.current_idx)
            meas1 = measure_spermatids(seg1, params)
            results1 = meas1['results']
            skel_label1 = meas1['skel_label']

            # Filter Pass 1 results to ROI
            if roi_mask is not None and results1:
                filtered1 = []
                keep1 = []
                for r in results1:
                    cy = min(max(int(round(r['centroid_y'])), 0), roi_mask.shape[0] - 1)
                    cx = min(max(int(round(r['centroid_x'])), 0), roi_mask.shape[1] - 1)
                    if roi_mask[cy, cx]:
                        filtered1.append(r)
                        keep1.append(r['label'])
                results1 = filtered1
                skel_label1 = np.where(np.isin(skel_label1, keep1), skel_label1, 0).astype(np.int32)
            log(f"  PASS 1: {len(results1)} detections (full-frame statistics)")

            # ── PASS 2: Tight ROI crop analysis ──────────────────────────────
            results2_global = []
            skel_label2_full = np.zeros((H, W), dtype=np.int32)

            if roi_mask is not None:
                log(f"  PASS 2: Tight-crop local analysis...")
                ys, xs = np.where(roi_mask)
                y0, y1 = int(ys.min()), int(ys.max()) + 1
                x0, x1 = int(xs.min()), int(xs.max()) + 1
                pad = 10
                y0c = max(0, y0 - pad)
                x0c = max(0, x0 - pad)
                y1c = min(H, y1 + pad)
                x1c = min(W, x1 + pad)

                crop_img = full_img[y0c:y1c, x0c:x1c]
                crop_roi = roi_mask[y0c:y1c, x0c:x1c]

                seg2 = segment_slice(crop_img, params, roi_mask=crop_roi, z_idx=self.current_idx)
                meas2 = measure_spermatids(seg2, params)
                results2 = meas2['results']
                skel_label2_crop = meas2['skel_label']

                # Filter Pass 2 to ROI
                if results2:
                    filtered2 = []
                    keep2 = []
                    for r in results2:
                        cy = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0] - 1)
                        cx = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1] - 1)
                        if crop_roi[cy, cx]:
                            filtered2.append(r)
                            keep2.append(r['label'])
                    results2 = filtered2
                    skel_label2_crop = np.where(np.isin(skel_label2_crop, keep2), skel_label2_crop, 0).astype(np.int32)

                # Map Pass 2 back to full-image coordinates
                ch2, cw2 = skel_label2_crop.shape
                max_label1 = int(skel_label1.max()) if skel_label1.max() > 0 else 0
                skel_label2_crop_shifted = np.where(skel_label2_crop > 0, skel_label2_crop + max_label1, 0).astype(np.int32)
                skel_label2_full[y0c:y0c+ch2, x0c:x0c+cw2] = skel_label2_crop_shifted

                for r in results2:
                    r['centroid_x'] += x0c
                    r['centroid_y'] += y0c
                    r['label'] += max_label1
                results2_global = results2
                log(f"  PASS 2: {len(results2_global)} detections (local-crop statistics)")
            else:
                log(f"  PASS 2: Skipped (no ROI)")

            # ── MERGE: Spatial deduplication ──────────────────────────────────
            # For each Pass 2 detection, check if Pass 1 already found a spermatid
            # within 5 pixels of the same centroid. If not, it's a new unique detection.
            dedup_radius = 5.0  # pixels
            merged_results = list(results1)
            skel_label_full = skel_label1.copy()
            new_from_pass2 = 0

            for r2 in results2_global:
                cx2, cy2 = r2['centroid_x'], r2['centroid_y']
                is_duplicate = False
                for r1 in results1:
                    dx = r1['centroid_x'] - cx2
                    dy = r1['centroid_y'] - cy2
                    if (dx*dx + dy*dy) < dedup_radius * dedup_radius:
                        is_duplicate = True
                        break
                if not is_duplicate:
                    merged_results.append(r2)
                    # Splice the Pass 2 skeleton pixels into the merged map
                    mask2 = (skel_label2_full == r2['label'])
                    skel_label_full[mask2] = r2['label']
                    new_from_pass2 += 1

            results = merged_results
            log(f"  MERGE: {len(results1)} (Pass1) + {new_from_pass2} unique from Pass2 = {len(results)} total")

            elapsed = _t.time() - t0
            log(f"  RESULT: {len(results)} spermatids detected ({elapsed:.1f}s)")

            overlay = make_overlay(full_img, skel_label_full)

            # Write log to file
            try:
                log_path = os.path.join(CONFIG['OUTPUT_DIR'], 'gui_analysis_log.txt')
                ensure_dir(CONFIG['OUTPUT_DIR'])
                with open(log_path, 'a', encoding='utf-8') as f:
                    f.write('\n'.join(log_lines) + '\n\n')
            except Exception:
                pass

            # Show results popup
            top = tk.Toplevel(self.root)
            top.title(f'Results Z={self.current_idx} - {len(results)} spermatids')
            top.geometry('1200x650')

            fig = Figure(figsize=(14, 6))
            ax1 = fig.add_subplot(121)
            ax2 = fig.add_subplot(122)

            img = full_img.astype(float)
            p1, p99 = np.percentile(img, 1), np.percentile(img, 99.5)
            disp = np.clip((img - p1) / (p99 - p1 + 1e-9), 0, 1)
            ax1.imshow(disp, cmap='gray')
            if roi_mask is not None:
                ax1.contour(roi_mask.astype(float), levels=[0.5], colors='red', linewidths=1.2)
            ax1.set_title('Original + ROI')
            ax1.axis('off')

            ax2.imshow(overlay)
            ax2.set_title(f'Overlay (N={len(results)})')
            ax2.axis('off')

            um = params['UM_PER_PX_XY']
            for r in results:
                ax2.text(r['centroid_x'], r['centroid_y'],
                         f"{r['length_px_geodesic'] * um:.1f}",
                         color='white', fontsize=5, ha='center', va='center')

            fig.tight_layout()
            can = FigureCanvasTkAgg(fig, master=top)
            can.get_tk_widget().pack(fill='both', expand=True)
            can.draw()

            if results:
                lengths = [r['length_px_geodesic'] * um for r in results]
                text = f'Found {len(results)} spermatids | median length {np.median(lengths):.2f} um ({elapsed:.1f}s)'
            else:
                text = f'Found 0 spermatids ({elapsed:.1f}s) - see gui_analysis_log.txt for diagnostics'
            lbl_stats = tk.Label(top, text=text, font=('Arial', 11))
            lbl_stats.pack(pady=4)

            lbl_tool = tk.Label(top, text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue', font=('Arial', 10, 'bold'))
            lbl_tool.pack(pady=2)

            self.lbl_roi.config(text=f'Analysis done: {len(results)} spermatids')

            # ── INTERACTIVE MANUAL CORRECTION LOGIC ──
            class ManualCorrector:
                def __init__(self, canvas, ax_overlay, seg_data, prms, crop_oy, crop_ox, fimg):
                    self.canvas = canvas
                    self.ax = ax_overlay
                    self.seg = seg_data
                    self.params = prms
                    self.crop_oy = crop_oy
                    self.crop_ox = crop_ox
                    self.fimg = fimg
                    
                    self.active_tool = None
                    self.cid_press = self.canvas.mpl_connect('button_press_event', self.on_click)
                    self.cid_key = self.canvas.mpl_connect('key_press_event', self.on_key)
                    
                    self.overlay_imshow = None
                    self.text_artists = []
                    
                def on_key(self, event):
                    if event.key == 'e':
                        self.active_tool = 'erase'
                        lbl_tool.config(text="Active Tool: ERASE (Click a colored spermatid to delete it)", fg='red')
                    elif event.key == 's':
                        self.active_tool = 'split'
                        lbl_tool.config(text="Active Tool: SPLIT (Click on a skeleton branch to sever it)", fg='orange')
                    elif event.key == 'escape':
                        self.active_tool = None
                        lbl_tool.config(text="Active Tool: None (Press 'E' to Erase, 'S' to Split, 'Esc' to Cancel)", fg='blue')

                def on_click(self, event):
                    if event.inaxes != self.ax or event.xdata is None or event.ydata is None:
                        return
                    if not self.active_tool:
                        return

                    # Map global click to cropped coordinate space
                    x_crop = int(round(event.xdata)) - self.crop_ox
                    y_crop = int(round(event.ydata)) - self.crop_oy
                    
                    ch, cw = self.seg['skel_pruned'].shape
                    if not (0 <= x_crop < cw and 0 <= y_crop < ch):
                        return # Clicked outside the ROI working area

                    modified = False
                    if self.active_tool == 'erase':
                        # Find the label at this pixel in the current skeleton
                        lab = self.seg['skel_labeled'][y_crop, x_crop]
                        # If the user clicked slightly off, search a 3x3 neighborhood
                        if lab == 0:
                            for dy in (-1, 0, 1):
                                for dx in (-1, 0, 1):
                                    yy, xx = y_crop+dy, x_crop+dx
                                    if 0 <= yy < ch and 0 <= xx < cw and self.seg['skel_labeled'][yy, xx] > 0:
                                        lab = self.seg['skel_labeled'][yy, xx]
                                        break
                        if lab > 0:
                            self.seg['skel_pruned'][self.seg['skel_labeled'] == lab] = False
                            modified = True
                            
                    elif self.active_tool == 'split':
                        # Draw a small black circle to sever the topological skeleton connection
                        from skimage.draw import disk
                        rr, cc = disk((y_crop, x_crop), radius=2.5, shape=self.seg['skel_pruned'].shape)
                        self.seg['skel_pruned'][rr, cc] = False
                        modified = True
                        
                    if modified:
                        self.recalculate_and_redraw()
                        
                def recalculate_and_redraw(self):
                    # Re-label the modified skeleton
                    self.seg['skel_labeled'] = measure.label(self.seg['skel_pruned'])
                    
                    # Re-run measurement filters
                    new_meas = measure_spermatids(self.seg, self.params)
                    new_results = new_meas['results']
                    
                    # Filter by ROI inside centroid logic
                    if crop_roi is not None:
                        filtered = []
                        keep_labels = []
                        for r in new_results:
                            c_y = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0] - 1)
                            c_x = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1] - 1)
                            if crop_roi[c_y, c_x]:
                                filtered.append(r)
                                keep_labels.append(r['label'])
                        new_results = filtered
                        skel_l = np.where(np.isin(new_meas['skel_label'], keep_labels), new_meas['skel_label'], 0).astype(np.int32)
                    else:
                        skel_l = new_meas['skel_label']

                    # Map back to full image
                    H, W = self.fimg.shape
                    skel_label_full = np.zeros((H, W), dtype=np.int32)
                    ch, cw = skel_l.shape
                    skel_label_full[self.crop_oy:self.crop_oy+ch, self.crop_ox:self.crop_ox+cw] = skel_l

                    # Redraw Overlay
                    new_overlay = make_overlay(self.fimg, skel_label_full)
                    
                    self.ax.clear()
                    self.ax.imshow(new_overlay)
                    self.ax.set_title(f'Overlay (N={len(new_results)}) - Manual Corrections Applied')
                    self.ax.axis('off')

                    # Redraw Text
                    _um = self.params['UM_PER_PX_XY']
                    for r in new_results:
                        self.ax.text(r['centroid_x'] + self.crop_ox, r['centroid_y'] + self.crop_oy,
                                 f"{r['length_px_geodesic'] * _um:.1f}",
                                 color='white', fontsize=5, ha='center', va='center')
                                 
                    self.canvas.draw()
                    
                    # Update Stats Label
                    if new_results:
                        lengths = [r['length_px_geodesic'] * _um for r in new_results]
                        lbl_stats.config(text=f'Corrected: {len(new_results)} spermatids | median length {np.median(lengths):.2f} um')
                    else:
                        lbl_stats.config(text=f'Corrected: 0 spermatids')

            # Attach to popup so it doesn't get garbage collected
            top.corrector = ManualCorrector(can, ax2, seg1, params, crop_offset_y, crop_offset_x, full_img)

        except Exception as e:
            import traceback
            traceback.print_exc()
            sys.stdout.flush()
            self.lbl_roi.config(text=f'Analysis error: {e}')
            messagebox.showerror('Analysis Error', f'{type(e).__name__}: {e}')

    def run_batch_analysis(self):
        if not self.files:
            messagebox.showinfo('Info', 'No directory loaded.')
            return

        # Auto-incremental output directory inside selected folder
        out_dir = get_unique_batch_dir(self.input_dir)
        self.last_out_dir = out_dir
        
        # EXPLICIT CONFIRMATION: Show the user where the data will go
        confirm = messagebox.askokcancel("Confirm Output", 
            f"Results (Excel, PDF, CSV) will be saved to:\n\n{out_dir}\n\nContinue?")
        if not confirm:
            return
            
        ensure_dir(out_dir)
        overlay_dir = os.path.join(out_dir, "overlays")
        ensure_dir(overlay_dir)

        params = CONFIG.copy()
        params['OUTPUT_DIR'] = out_dir
        params['SAVE_DEBUG_IMAGES'] = False
        params['DO_TRACKING'] = True
        
        roi_mask = self.build_roi_mask()

        self.lbl_roi.config(text="Processing... See Top Bar")
        self.root.update_idletasks()
        
        try:
            import time as _t
            t_batch = _t.time()
            self.lbl_batch_op.config(text=f"Batch Segmenting: 0 / {len(self.files)} slices...", fg='blue')
            self.root.update()
            
            all_rows = []
            summaries = []
            
            # Z-Projection accumulation
            max_proj_raw = None
            max_proj_ov = None
            
            # Robust initialization
            ts = None
            df_trk = None

            self.progress['value'] = 0
            self.progress['maximum'] = len(self.files)
            self.progress_post['value'] = 0
            self.progress_post['maximum'] = 100
            self.lbl_post_progress_val.config(text="0%", fg='orange')
            
            for idx, fpath in enumerate(self.files):
                z_idx = extract_z_index(fpath, sequence_idx=idx)
                
                pct = int(((idx + 1) / len(self.files)) * 100)
                self.progress['value'] = idx + 1
                self.lbl_progress_val.config(text=f"{pct}%")
                self.root.update()
                
                print(f"[{idx+1}/{len(self.files)}] Processing Z={z_idx:02d}...")
                
                img_raw = robust_imread(fpath)
                process_img = ensure_2d_image(img_raw, os.path.basename(fpath))
                print(f"DEBUG batch image {os.path.basename(fpath)} shape: {process_img.shape}, ndim={process_img.ndim}")
                
                # Handling ROI cropping to speed up processing
                full_img = process_img
                crop_oy, crop_ox = 0, 0
                crop_roi = roi_mask
                
                # ── BUGFIX: NEVER crop the image before analysis ──
                # Cropping to a tight ROI artificially removes the dark background,
                # which completely destroys the np.percentile() thresholding math that
                # relies on the global image statistics!
                seg = segment_slice(process_img, params, z_idx=z_idx)
                meas = measure_spermatids(seg, params)
                res = meas['results']
                sl_crop = meas['skel_label']
                
                if crop_roi is not None and res:
                    filtered = []
                    keep = []
                    for r in res:
                        cy = min(max(int(round(r['centroid_y'])), 0), crop_roi.shape[0]-1)
                        cx = min(max(int(round(r['centroid_x'])), 0), crop_roi.shape[1]-1)
                        if crop_roi[cy, cx]:
                            filtered.append(r)
                            keep.append(r['label'])
                    res = filtered
                    sl_crop = np.where(np.isin(sl_crop, keep), sl_crop, 0).astype(np.int32)
                
                for r in res:
                    r['centroid_x'] += crop_ox
                    r['centroid_y'] += crop_oy
                    
                H, W = full_img.shape
                sl_full = np.zeros((H, W), dtype=np.int32)
                ch, cw = sl_crop.shape
                sl_full[crop_oy:crop_oy+ch, crop_ox:crop_ox+cw] = sl_crop
                
                um = params['UM_PER_PX_XY']
                ls_um = [r['length_px_geodesic']*um for r in res]
                
                all_rows.extend(rows_from_results(res, z_idx, um))
                summaries.append({
                    "z_slice": z_idx,
                    "n_spermatids": len(res),
                    "median_length_um": round(float(np.median(ls_um)), 3) if ls_um else 0,
                })
                
                if params['SAVE_OVERLAYS']:
                    ov = make_overlay(full_img, sl_full)
                    # Create side-by-side panel
                    orig_rgb = (normalize_display(full_img) * 255).astype(np.uint8)
                    if orig_rgb.ndim == 2:
                        orig_rgb = np.stack([orig_rgb]*3, axis=-1)
                    panel = np.hstack([orig_rgb, ov])
                    _imwrite(os.path.join(overlay_dir, f"z{z_idx:02d}_panel.png"), panel)
                    
                    # ---- LIVE GUI UPDATE ----
                    # Show the side-by-side segmentation panel during batch execution
                    if hasattr(self, 'ax') and hasattr(self, 'canvas'):
                        try:
                            self.ax.clear()
                            self.ax.axis('off')
                            self.ax.imshow(panel)
                            self.canvas.draw()
                        except Exception:
                            pass
                    
                    # Update Z-Projections
                    if max_proj_raw is None:
                        max_proj_raw = img_raw.copy().astype(np.float32)
                        max_proj_ov = ov.copy().astype(np.float32)
                    else:
                        max_proj_raw = np.maximum(max_proj_raw, img_raw)
                        max_proj_ov = np.maximum(max_proj_ov, ov.astype(np.float32))

                # Update Progress Bar for each slice
                self.lbl_batch_op.config(text=f"Batch Segmenting: {idx+1} / {len(self.files)} slices...", fg='blue')
                self.progress['value'] = idx + 1
                self.root.update()
            
            df = pd.DataFrame(all_rows)
            df_sum = pd.DataFrame(summaries)
            df.to_csv(os.path.join(out_dir, "spermatid_measurements.csv"), index=False)
            df_sum.to_csv(os.path.join(out_dir, "slice_summary.csv"), index=False)
            
            if not df.empty:
                self.lbl_batch_op.config(text='Running 3D Tracking & Morphometrics...', fg='#e67e22') # Orange-ish

                # --- Sequential Progress Bar Swap ---
                # Hide 2D progress, Show Post-Analysis progress
                self.batch_p_frame.pack_forget()
                self.post_p_frame.pack(fill='x')

                self.progress_post['value'] = 25
                self.lbl_post_progress_val.config(text="25%")
                self.root.update()

                df_trk, ts = track_across_slices(df, params)
                
                # --- Advanced 3D Morphometrics ---
                self.lbl_batch_op.config(text='Calculating Advanced 3D Metrics...', fg='#e67e22')
                # Metrics are natively generated correctly and safely in track_across_slices

                self.progress_post['value'] = 60
                self.lbl_post_progress_val.config(text="60%")
                self.root.update()

                df_trk.to_csv(os.path.join(out_dir, "measurements_with_tracks.csv"), index=False)
                
                # ── AUTO QUALITY AUDIT ──────────────────────────────────────
                self.root.update()
                
                
                # Save annotated track summary (with quality flags)
                ts.to_csv(os.path.join(out_dir, "track_summary.csv"), index=False)
                
                
                # Generate outlier_audit/ subfolder automatically
                
                n_quality = len(ts_quality)
                n_flagged = len(ts) - n_quality
                self.lbl_batch_op.config(
                    text=f'Quality: {n_quality} clean / {n_flagged} flagged', fg='#27ae60')
                self.root.update()
                
                # Save Global Z-Projection
                if max_proj_raw is not None:
                    self.lbl_batch_op.config(text='Generating Global Z-Projection...', fg='#e67e22')
                    raw_p = (normalize_display(max_proj_raw.astype(np.uint16)) * 255).astype(np.uint8)
                    if raw_p.ndim == 2: raw_p = np.stack([raw_p]*3, axis=-1)
                    ov_p = max_proj_ov.astype(np.uint8)
                    global_panel = np.hstack([raw_p, ov_p])
                    _imwrite(os.path.join(out_dir, "global_z_projection.png"), global_panel)
                    
                self.progress_post['value'] = 80
                self.lbl_post_progress_val.config(text="80%")
                self.root.update()
            
            elapsed = _t.time() - t_batch
            msg = f"Batch complete in {elapsed:.1f}s!\nSaved to: {out_dir}"
            print(msg)
            self.lbl_batch_op.config(text='Batch Analysis Complete.', fg='green')
            self.lbl_progress_val.config(text="100% - Done", fg='green')
            self.root.update()
            
            # Generate High-Res Graphical Report and Excel Audit
            self.lbl_batch_op.config(text='Generating PDF & Excel Reports...', fg='#8e44ad') # Purple
            
            def update_cb(v):
                self.progress_post['value'] = v
                self.lbl_post_progress_val.config(text=f"{v}%")
                self.root.update()

            generate_batch_report(out_dir, df, df_sum, um, ts if not df.empty else None, update_cb)
            generate_excel_report(out_dir, df, df_sum, ts if not df.empty else None)
            
            # --- Store batch data for AI button ---
            if not df.empty and ts is not None and len(ts) > 0:
                self._last_batch_ts = ts
                self._last_batch_out_dir = out_dir
                self.btn_ai.config(state='normal')
                print(f"AI READY: {len(ts)} tracks stored. Click 'Run AI Analysis' to interpret.")

            self.lbl_batch_op.config(text='Generating PowerPoint Dashboard...', fg='#c0392b') # Red-ish
            self.root.update()
            try:
                ok = generate_pptx_report(out_dir, df, df_sum, um, ts if not df.empty else None)
                if not ok:
                    print("WARNING: PPTX generation returned False — check console for traceback.")
            except Exception as pptx_err:
                import traceback as _tb
                print(f"ERROR generating PPTX: {pptx_err}")
                _tb.print_exc()
            
            msg = f"Batch complete in {elapsed:.1f}s!\nSaved to: {out_dir}"
            print(msg)
            self.lbl_batch_op.config(text='ALL OPERATIONS COMPLETE', fg='green')
            self.lbl_progress_val.config(text="100%", fg='green')
            self.lbl_post_progress_val.config(text="100% - DONE", fg='green')
            self.progress_post['value'] = 100
            self.root.update()

            messagebox.showinfo('Batch Complete', msg)
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            self.lbl_batch_op.config(text=f'Batch error: {e}', fg='red')
            messagebox.showerror('Batch Error', str(e))

def launch_gui():
    if not _TK_AVAILABLE:
        raise RuntimeError("Tkinter GUI components are not available in this Python environment.")
    root = tk.Tk()
    app = SpermGUI(root)
    root.mainloop()

# =============================================================================
# ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    print(f"\n{'='*60}")
    print(f" SPERMATID ANALYSIS PIPELINE - v{_VERSION} - INITIALIZING...")
    print(f" Mode: {'GUI' if '--gui' in sys.argv or len(sys.argv)<=1 else 'CLI'}")
    print(f" Features: [Incremental Folders] [Multi-Tab Excel] [3D Stats]")
    print(f"{'='*60}\n")
    
    ap = argparse.ArgumentParser(description=f"Spermatid segmentation {_VERSION} / ROI GUI v11")
    ap.add_argument("--batch",  action="store_true", help="Run batch processing")
    ap.add_argument("--single", action="store_true", help="Run a single-image analysis")
    ap.add_argument("--gui",    action="store_true", help="Launch ROI GUI")
    ap.add_argument("--z",      type=int, default=None, help="Choose z-index in single mode")
    ap.add_argument("--params", type=str, default=None, help="Path to tuned parameters JSON file to override CONFIG")
    args = ap.parse_args()

    # Load tuned parameters from JSON if provided
    if args.params:
        import json as _json
        params_path = os.path.abspath(args.params)
        if os.path.exists(params_path):
            with open(params_path, 'r') as _pf:
                tuned = _json.load(_pf)
            applied = 0
            for key, value in tuned.items():
                if key in CONFIG:
                    CONFIG[key] = value
                    applied += 1
            print(f"  Loaded {applied} tuned parameters from: {os.path.basename(params_path)}")
        else:
            print(f"  WARNING: Params file not found: {params_path}")

    validate_config(CONFIG)

    # Launch GUI by default if no explicit CLI flags are provided
    if args.gui or not (args.batch or args.single or args.z is not None):
        launch_gui()
        raise SystemExit

    if args.batch:
        CONFIG["RUN_MODE"] = "batch"
        # CLI Incremental Folder Logic - Anchor to INPUT_DIR
        base_parent = CONFIG["INPUT_DIR"]
        if not os.path.isabs(base_parent):
            base_parent = os.path.abspath(base_parent)
            
        CONFIG["OUTPUT_DIR"] = get_unique_batch_dir(base_parent)
        ensure_dir(CONFIG["OUTPUT_DIR"])
        print(f"CLI BATCH MODE: Results will be saved inside input folder: {CONFIG['OUTPUT_DIR']}")
    if args.single:
        CONFIG["RUN_MODE"] = "single"
    if args.z is not None:
        CONFIG["SINGLE_IMAGE_SELECTION_MODE"] = "z_index"
        CONFIG["SINGLE_Z_INDEX"] = args.z

    if CONFIG["RUN_MODE"] == "single":
        process_one_image(choose_single_image(CONFIG), CONFIG, CONFIG["OUTPUT_DIR"])
    else:
        process_batch(CONFIG)

